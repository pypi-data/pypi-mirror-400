from typing import Any, Dict, List, Optional, Tuple
from abc import ABC, abstractmethod
import tempfile
from pathlib import Path
from io import BytesIO
import markdown
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
from navconfig.logging import logging


class PowerPointFile(ABC):
    """
    Generic PowerPointFile Interface

    Fully configurable through missing_sections without hardcoded assumptions.
    """

    def __init__(self, *args, **kwargs):
        self.logger = logging.getLogger('Flowtask.PowerPointFile')
        self._layout_cache: Dict[int, Dict[str, Any]] = {}
        self._presentation: Optional[Presentation] = None

        # Default styles - applied when not specified in missing_sections
        self.default_styles = {
            'font_name': 'Calibri',
            'font_size': 14,
            'bold': False,
            'color': '000000',  # Black
            'alignment': 'left'
        }

        super().__init__(*args, **kwargs)

    def load_template(self, template_path: Path) -> bool:
        """Load PowerPoint template and analyze layouts."""
        try:
            self._presentation = Presentation(str(template_path))
            self.logger.info(f"Loaded PowerPoint template: {template_path}")
            return True
        except Exception as e:
            self.logger.error(f"Failed to load PowerPoint template: {e}")
            return False

    def clear_existing_slides(self):
        """Remove existing slides."""
        if not self._presentation:
            return
        while len(self._presentation.slides) > 0:
            rId = self._presentation.slides._sldIdLst[0].rId
            self._presentation.part.drop_rel(rId)
            del self._presentation.slides._sldIdLst[0]

    def convert_markdown_to_html(self, markdown_text: str) -> str:
        """Convert markdown to HTML."""
        try:
            html = markdown.markdown(markdown_text, extensions=['extra', 'tables'])
            return html
        except Exception as e:
            self.logger.error(f"Failed to convert markdown to HTML: {e}")
            return markdown_text

    def extract_text_from_html(self, html: str) -> str:
        """Extract clean text from HTML for PowerPoint."""
        try:
            soup = BeautifulSoup(html, 'html.parser')
            text_parts = []

            for element in soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'li', 'div']):
                text = element.get_text().strip()
                if text:
                    if element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                        text = f"• {text}"
                    elif element.name == 'li':
                        text = f"  - {text}"
                    text_parts.append(text)

            return '\n'.join(text_parts)
        except Exception as e:
            self.logger.error(f"Failed to extract text from HTML: {e}")
            return html

    def save_presentation(self, output_filename, output_path: Path, override: bool = True):
        """Save the presentation to file."""
        if not self._presentation:
            raise RuntimeError("No presentation to save")
        output_filename = output_filename if output_filename.endswith('.pptx') else f"{output_filename}.pptx"
        path = output_path.joinpath(output_filename)

        if not override and path.exists():
            self.logger.warning(f"PowerPoint presentation already exists: {path}")
            return
        if override and path.exists():
            self.logger.info(f"Overriding existing PowerPoint presentation: {path}")
            path.unlink(missing_ok=True)

        self._presentation.save(str(path))
        self.logger.notice(f"PowerPoint presentation saved: {path}")

    def _layout_catalog(self):
        """Scan all layouts and collect {layout_idx: {placeholder_name: placeholder_idx}}"""
        catalog = {}
        for i, ly in enumerate(self._presentation.slide_layouts):
            names = {}
            for ph in getattr(ly, "placeholders", []):
                try:
                    nm = getattr(ph, "name", None)
                    idx = ph.placeholder_format.idx
                except Exception:
                    continue
                if nm:
                    names[nm] = idx
            catalog[i] = names
        return catalog

    def _choose_best_layout_index(self) -> int:
        """Choose the layout that contains the most placeholder names from column_mapping."""
        wanted = {dest for dest in self.column_mapping.values() if isinstance(dest, str)}
        catalog = self._layout_catalog()
        best_idx, best_score = None, -1
        for idx, names in catalog.items():
            score = len(wanted.intersection(set(names.keys())))
            if score > best_score:
                best_idx, best_score = idx, score
        for fb in (getattr(self, 'slide_layout', 1), 1, 6):
            if best_idx is None and fb < len(self._presentation.slide_layouts):
                best_idx = fb
        return best_idx if best_idx is not None else 1

    def _safe_placeholder_idx(self, shape):
        try:
            return shape.placeholder_format.idx
        except Exception:
            return None

    def _find_placeholder_by_name(self, slide, placeholder_name: str):
        # Direct name match first
        for shp in slide.shapes:
            if getattr(shp, "name", None) == placeholder_name:
                return shp

        # Match by placeholder index from the layout
        ly = slide.slide_layout
        for lph in getattr(ly, "placeholders", []):
            if getattr(lph, "name", None) == placeholder_name:
                target_idx = self._safe_placeholder_idx(lph)
                if target_idx is None:
                    continue
                for shp in slide.shapes:
                    if self._safe_placeholder_idx(shp) == target_idx:
                        return shp
        return None

    def _ensure_placeholder(self, slide, section_name: str, autoplace_state: dict):
        """
        Create placeholder using missing_sections configuration.
        Falls back to simple flow layout if no configuration provided.
        """
        try:
            # First try to find existing placeholder
            shp = self._find_placeholder_by_name(slide, section_name)
            if shp:
                if shp.left < 0 or shp.top < 0:
                    sw, sh = self._presentation.slide_width, self._presentation.slide_height
                    shp.left = Inches(0.5)
                    shp.top = sh - Inches(0.6)
                    shp.width = sw - Inches(1.0)
                self.logger.debug(f"Found existing placeholder for '{section_name}'")
                return shp

            # Get section configuration
            spec = getattr(self, "missing_sections", {}).get(section_name, {})

            # Try configured positioning first
            if self._has_positioning_config(spec):
                return self._create_configured_placeholder(slide, section_name, spec, autoplace_state)

            # Fall back to simple flow layout
            return self._create_flow_placeholder(slide, section_name, autoplace_state)

        except Exception as e:
            self.logger.error(f"Error creating placeholder for '{section_name}': {e}")
            return self._create_emergency_placeholder(slide, section_name, autoplace_state)

    def _has_positioning_config(self, spec: dict) -> bool:
        """Check if spec has positioning information."""
        return (
            ("anchor" in spec and "position" in spec) or ("left_in" in spec and "top_in" in spec)
        )

    def _create_configured_placeholder(self, slide, section_name: str, spec: dict, autoplace_state: dict):
        """Create placeholder using missing_sections configuration."""
        # Try anchor-based positioning
        if "anchor" in spec:
            anchor = self._find_placeholder_by_name(slide, spec["anchor"])
            if anchor:
                return self._create_anchored_placeholder(slide, section_name, spec, anchor)
            else:
                self.logger.debug(f"Anchor '{spec['anchor']}' not found for '{section_name}'")

        # Try absolute positioning
        if "left_in" in spec and "top_in" in spec:
            return self._create_absolute_placeholder(slide, section_name, spec)

        # Config exists but is incomplete - fall back to flow
        self.logger.debug(f"Incomplete positioning config for '{section_name}', using flow layout")
        return self._create_flow_placeholder(slide, section_name, autoplace_state)

    def _create_anchored_placeholder(self, slide, section_name: str, spec: dict, anchor):
        """Create placeholder relative to anchor."""
        pos = (spec.get("position", "below")).lower()
        off = spec.get("offset_in", {})
        dx, dy = Inches(off.get("dx", 0.0)), Inches(off.get("dy", 0.0))
        sz = spec.get("size_in", {})
        width, height = Inches(sz.get("width", 3.0)), Inches(sz.get("height", 0.5))

        if pos == "below":
            left, top = anchor.left, anchor.top + anchor.height + dy
        elif pos == "right":
            left, top = anchor.left + anchor.width + dx, anchor.top
        elif pos == "left":
            left, top = anchor.left - width + dx, anchor.top
        else:  # above
            left, top = anchor.left, anchor.top - height + dy

        box = slide.shapes.add_textbox(left, top, width, height)
        box.name = section_name
        self.logger.debug(f"Created anchored section '{section_name}' at ({left}, {top})")
        return box

    def _create_absolute_placeholder(self, slide, section_name: str, spec: dict):
        """Create placeholder at absolute position."""
        left = Inches(spec["left_in"])
        top = Inches(spec["top_in"])
        sz = spec.get("size_in", {})
        width, height = Inches(sz.get("width", 3.0)), Inches(sz.get("height", 0.5))

        box = slide.shapes.add_textbox(left, top, width, height)
        box.name = section_name
        self.logger.debug(f"Created absolute-positioned section '{section_name}' at ({left}, {top})")
        return box

    def _create_flow_placeholder(self, slide, section_name: str, autoplace_state: dict):
        """Create placeholder using simple flow layout."""
        left = Inches(0.5)
        top = autoplace_state.setdefault("flow_top", Inches(0.5))
        width, height = Inches(4.0), Inches(0.5)
        autoplace_state["flow_top"] = top + height + Inches(0.1)

        box = slide.shapes.add_textbox(left, top, width, height)
        box.name = section_name
        self.logger.debug(f"Created flow section '{section_name}' at ({left}, {top})")
        return box

    def _create_emergency_placeholder(self, slide, section_name: str, autoplace_state: dict):
        """Emergency placeholder creation that should never fail."""
        try:
            left = Inches(0.5)
            top = autoplace_state.setdefault("emergency_top", Inches(0.5))
            width, height = Inches(2.0), Inches(0.5)
            autoplace_state["emergency_top"] = top + height + Inches(0.2)

            box = slide.shapes.add_textbox(left, top, width, height)
            box.name = f"{section_name}_emergency"
            self.logger.debug(f"Emergency-created section '{section_name}' at ({left}, {top})")
            return box
        except Exception as e:
            self.logger.error(f"Emergency placeholder creation failed for '{section_name}': {e}")
            return None

    def _normalize_image(self, value, column_name: str = ""):
        """Normalize various image formats to PIL.Image."""
        from PIL import Image as PILImage
        import pandas as pd
        from datetime import datetime, date

        if pd.isna(value) or value is None:
            return None

        if isinstance(value, (str, int, float, bool, datetime, date, pd.Timestamp)):
            return None

        try:
            if isinstance(value, PILImage.Image):
                return value
            elif isinstance(value, BytesIO):
                value.seek(0)
                img = PILImage.open(value)
                if img.mode in ('RGBA', 'LA', 'P'):
                    img = img.convert('RGB')
                return img
            elif isinstance(value, (bytes, bytearray)):
                img = PILImage.open(BytesIO(value))
                if img.mode in ('RGBA', 'LA', 'P'):
                    img = img.convert('RGB')
                return img
            elif hasattr(value, 'read'):
                img = PILImage.open(value)
                if img.mode in ('RGBA', 'LA', 'P'):
                    img = img.convert('RGB')
                return img
            elif "numpy" in str(type(value)):
                return PILImage.fromarray(value)
            else:
                return None
        except Exception:
            return None

    def _format_value(self, col, val):
        """Enhanced value formatting for different data types and contexts."""
        import pandas as pd
        from datetime import datetime, date

        if pd.isna(val) or val is None:
            return getattr(self, 'reference_text', {}).get(col, f"[{col}]")

        if isinstance(val, bool):
            if any(word in col.lower() for word in ['compliant', 'status', 'pass', 'fail']):
                return "✓ Compliant" if val else "✗ Non-Compliant"
            elif any(word in col.lower() for word in ['active', 'enabled', 'valid']):
                return "✓ Yes" if val else "✗ No"
            return "True" if val else "False"

        if isinstance(val, (pd.Timestamp, datetime, date)):
            if hasattr(val, 'strftime'):
                return val.strftime("%Y-%m-%d %H:%M") if hasattr(val, 'hour') else val.strftime("%Y-%m-%d")
            return str(val)

        if isinstance(val, float):
            if any(word in col.lower() for word in ['score', 'rate', 'percent', 'pct']):
                if 0 <= val <= 1:
                    return f"{val:.1%}"
                else:
                    return f"{val:.1f}"
            return f"{val:.2f}" if val != int(val) else str(int(val))

        if 'store_id' in col.lower() or 'id' in col.lower():
            return f"Store ID: {val}" if 'store' in col.lower() else f"ID: {val}"

        return str(val)

    def _apply_text_styles(self, text_frame, section_name: str, custom_styles: dict = None):
        """Apply text styles from missing_sections configuration or defaults."""
        try:
            # Get styles from missing_sections first, then custom_styles, then defaults
            section_spec = getattr(self, "missing_sections", {}).get(section_name, {})
            styles = {}

            # Start with defaults
            styles.update(self.default_styles)

            # Override with custom styles if provided
            if custom_styles:
                styles.update(custom_styles)

            # Override with section-specific config from missing_sections
            if section_spec:
                for style_key in ['font_name', 'font_size', 'bold', 'color', 'alignment']:
                    if style_key in section_spec:
                        styles[style_key] = section_spec[style_key]

            # Apply styles
            for paragraph in text_frame.paragraphs:
                font = paragraph.font
                font.name = styles.get('font_name', 'Calibri')

                if 'font_size' in styles:
                    font.size = Pt(int(styles['font_size']))

                font.bold = bool(styles.get('bold', False))

                if 'color' in styles:
                    color_hex = str(styles['color']).strip('#')
                    if len(color_hex) == 6:
                        try:
                            r = int(color_hex[0:2], 16)
                            g = int(color_hex[2:4], 16)
                            b = int(color_hex[4:6], 16)
                            font.color.rgb = RGBColor(r, g, b)
                        except ValueError:
                            self.logger.warning(f"Invalid color hex: {color_hex}")

                if 'alignment' in styles:
                    alignment = styles['alignment'].lower()
                    if alignment == 'center':
                        paragraph.alignment = PP_ALIGN.CENTER
                    elif alignment == 'right':
                        paragraph.alignment = PP_ALIGN.RIGHT

        except Exception as e:
            self.logger.error(f"Failed to apply text styles to {section_name}: {e}")

    def _replace_placeholder_with_image(self, slide, placeholder, pil_image: Image.Image):
        """Replace placeholder with image."""
        try:
            left, top, width, height = placeholder.left, placeholder.top, placeholder.width, placeholder.height

            try:
                slide.shapes._spTree.remove(placeholder._element)
            except Exception as e:
                self.logger.debug(f"Could not remove placeholder element: {e}")

            if pil_image.mode in ('RGBA', 'LA', 'P'):
                pil_image = pil_image.convert('RGB')

            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                pil_image.save(tmp, format="PNG")
                path = tmp.name

            pic = slide.shapes.add_picture(path, left, top, width, height)

            try:
                Path(path).unlink(missing_ok=True)
            except Exception as e:
                self.logger.debug(f"Could not delete temporary file {path}: {e}")

            self.logger.debug("Successfully replaced placeholder with image")
            return pic

        except Exception as e:
            self.logger.error(f"Failed to replace placeholder with image: {e}")
            return None
