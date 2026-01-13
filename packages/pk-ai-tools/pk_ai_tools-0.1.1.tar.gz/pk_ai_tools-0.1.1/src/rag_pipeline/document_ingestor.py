from __future__ import annotations
import os, json, logging, time, subprocess, shutil, tempfile, io, csv
from dataclasses import dataclass, field
from typing import Iterable, List, Optional, Tuple, Union
from contextlib import contextmanager

from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
    Docx2txtLoader,
    CSVLoader,
)
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain.docstore.document import Document

# Ny importväg (med fallback om paketet inte finns installerat)
try:
    from langchain_ollama import OllamaEmbeddings
except Exception:
    from langchain_community.embeddings import OllamaEmbeddings  # fallback
    logging.getLogger(__name__).warning(
        "Uses deprecated langchain_community.OllamaEmbeddings."
        "Please install langchain-ollama and change the import path."
    )

# Excel (Unstructured) – valfritt
try:
    from langchain_community.document_loaders import UnstructuredExcelLoader
    HAS_UNSTRUCTURED_EXCEL = True
except Exception:
    HAS_UNSTRUCTURED_EXCEL = False

# Word (Unstructured) – redan i din kod
try:
    from langchain_community.document_loaders import UnstructuredWordDocumentLoader
    HAS_UNSTRUCTURED_WORD = True
except Exception:
    HAS_UNSTRUCTURED_WORD = False

import pandas as pd
import requests
import ollama
# Unstructured loaders (valfria – vi faller tillbaka om de saknas)
try:
    from langchain_community.document_loaders import (
        UnstructuredMarkdownLoader,
        UnstructuredHTMLLoader,
        UnstructuredPowerPointLoader,
        UnstructuredODTLoader,
        UnstructuredRTFLoader,
    )
    HAS_UNSTRUCTURED_MD = True
    HAS_UNSTRUCTURED_HTML = True
    HAS_UNSTRUCTURED_PPT = True
    HAS_UNSTRUCTURED_ODT = True
    HAS_UNSTRUCTURED_RTF = True
except Exception:
    HAS_UNSTRUCTURED_MD = HAS_UNSTRUCTURED_HTML = HAS_UNSTRUCTURED_PPT = HAS_UNSTRUCTURED_ODT = HAS_UNSTRUCTURED_RTF = False

# Fallbacks
try:
    from bs4 import BeautifulSoup
    HAS_BS4 = True
except Exception:
    HAS_BS4 = False

try:
    import pptx  # python-pptx
    HAS_PYTHON_PPTX = True
except Exception:
    HAS_PYTHON_PPTX = False


# Filtyper vi snappar upp i mappskanningen
SUPPORTED_EXTS = (
    ".pdf", ".txt", ".docx", ".doc", ".csv", ".xlsx", ".xls", ".ods",
    ".md", ".markdown", ".html", ".htm", ".pptx", ".ppt", ".odt", ".odp",
    ".tsv", ".json", ".jsonl",
)


# Google Sheets-URL detektering
def _is_gsheet_url(s: str) -> bool:
    return isinstance(s, str) and s.startswith(("http://", "https://")) and "docs.google.com/spreadsheets" in s

@contextmanager
def timed(logger: logging.Logger, label: str):
    # logger.info(f"▶️  Start: {label}")
    t0 = time.perf_counter()
    try:
        yield
    finally:
        dt = time.perf_counter() - t0
        # logger.info(f"✅ Klart: {label} (tog {dt:.2f}s)")

@dataclass
class IngestConfig:
    doc_folder: str
    chroma_path: str
    vector_store_name: str
    embedding_model: str = "nomic-embed-text"
    chunk_size: int = 1000
    chunk_overlap: int = 200
    manifest_filename: str = "processed_files.json"
    use_pdf_fallback: bool = False
    verbose: bool = True  # styr logg-nivå
    # Word / .doc konvertering:
    convert_doc_with_soffice: bool = False
    soffice_path: Optional[str] = None  # ex: r"C:\Program Files\LibreOffice\program\soffice.exe"
    convert_ods_with_soffice: bool = False
    max_rows_per_sheet: Optional[int] = None
    convert_odt_with_soffice: bool = False  # odt -> docx fallback
    convert_odp_with_soffice: bool = False  # odp -> pptx fallback
    convert_ppt_with_soffice: bool = False
    error_on_zero_docs: bool = False
    recursive: bool = True
    follow_symlinks: bool = False
    ignore_dirs: Tuple[str, ...] = (".git", "__pycache__", "node_modules", ".venv", "venv",
                                    ".mypy_cache", ".pytest_cache", ".idea", ".vscode")
    log_enabled = False  # stänger AV/PÅ all logg i denna modul
    log_level = logging.DEBUG  # miniminivå för modulen (t.ex. INFO/WARNING)
    console_level = logging.WARNING  # vad som syns i konsolen
    log_to_file = r"./ingest.log"  # eller None
    dedup_ttl = 2.0

@dataclass
class DocumentIngestor:
    cfg: IngestConfig
    logger: logging.Logger = field(default_factory=lambda: logging.getLogger(__name__))

    # -------- Public API --------
    def build_or_update_vector_db(self, files: Optional[Iterable[str]] = None) -> Chroma:
        """
        Läser in nya dokument/URL:er, splittrar och bygger/uppdaterar Chroma.
        - Sparar ENDAST lyckade mål i manifestet.
        - Om inga nya mål finns: öppnar befintlig Chroma om den finns.
        - Om nya mål finns men 0 docs inlästa: raise eller mjuk hantering beroende på cfg.error_on_zero_docs.
        """
        self._configure_logging_once()
        with timed(self.logger, "Förbered kataloger"):
            self._ensure_dirs()

        manifest_path = os.path.join(self.cfg.chroma_path, self.cfg.manifest_filename)

        with timed(self.logger, "Läsa manifest"):
            processed_files = [
                os.path.normpath(p) if not _is_gsheet_url(p) else p
                for p in self._load_manifest(manifest_path)
            ]

        with timed(self.logger, "Lista målfiler"):
            current_targets = self._list_target_items(files)  # paths och/eller gsheet-URL:er
            # self.logger.info(f"Hittade {len(current_targets)} mål (filer/URL:er)")

        with timed(self.logger, "Beräkna nya mål"):
            norm_processed = set(processed_files)
            new_targets = [t for t in current_targets if t not in norm_processed]
            # self.logger.info(f"Nya mål: {len(new_targets)}")

        failed: List[str] = []
        succeeded: List[str] = []

        if not new_targets:
            # self.logger.info("Inga nya mål att processa.")
            if os.path.exists(self.cfg.chroma_path):
                with timed(self.logger, "Öppna befintlig Chroma-databas"):
                    return Chroma(
                        embedding_function=OllamaEmbeddings(model=self.cfg.embedding_model),
                        collection_name=self.cfg.vector_store_name,
                        persist_directory=self.cfg.chroma_path,
                    )
            raise FileNotFoundError(
                f"Ingen befintlig Chroma på {self.cfg.chroma_path} och inga nya filer att indexera."
            )

        # Det finns nya mål → försök ladda dem
        with timed(self.logger, "Ladda dokument"):
            docs, succeeded, failed = self._load_documents(new_targets)
            # self.logger.info(f"Dokumentobjekt (sammanlagt): {len(docs)}")
            self._log_ingest_summary(succeeded, failed, stage="Efter laddning")

            if not docs:
                # Ingen indexering – spara inget nytt i manifestet
                self._log_ingest_summary(succeeded, failed, final=True)
                error_on_zero = getattr(self.cfg, "error_on_zero_docs", True)
                if error_on_zero:
                    raise RuntimeError("Inga dokument kunde läsas in (nya mål fanns, men gav 0 Docs).")
                # Mjuk hantering: öppna ev. Befintlig Chroma, annars samma FileNotFoundError som ovan
                if os.path.exists(self.cfg.chroma_path):
                    # self.logger.warning("0 Docs inlästa; hoppar indexering och återanvänder befintlig Chroma.")
                    with timed(self.logger, "Öppna befintlig Chroma-databas"):
                        return Chroma(
                            embedding_function=OllamaEmbeddings(model=self.cfg.embedding_model),
                            collection_name=self.cfg.vector_store_name,
                            persist_directory=self.cfg.chroma_path,
                        )
                raise FileNotFoundError(
                    f"Ingen befintlig Chroma på {self.cfg.chroma_path} och inga indexerbara dokument."
                )

        with timed(self.logger, "Splitta dokument till chunks"):
            chunks = self._split_documents(docs)
            # self.logger.info(f"Totalt antal chunks: {len(chunks)}")
            if not chunks:
                # self.logger.warning("0 chunks efter split – hoppar indexering.")
                self._log_ingest_summary(succeeded, failed, final=True)
                error_on_zero = getattr(self.cfg, "error_on_zero_docs", True)
                if error_on_zero:
                    raise RuntimeError("0 chunks genererades från inlästa dokument.")
                if os.path.exists(self.cfg.chroma_path):
                    with timed(self.logger, "Öppna befintlig Chroma-databas"):
                        return Chroma(
                            embedding_function=OllamaEmbeddings(model=self.cfg.embedding_model),
                            collection_name=self.cfg.vector_store_name,
                            persist_directory=self.cfg.chroma_path,
                        )
                raise FileNotFoundError(
                    f"Ingen befintlig Chroma på {self.cfg.chroma_path} och 0 chunks att indexera."
                )

        with timed(self.logger, f"Säkerställ embeddingsmodell '{self.cfg.embedding_model}'"):
            self._ensure_model_available(self.cfg.embedding_model)

        with timed(self.logger, "Bygga Chroma-index från dokument"):
            vectordb = Chroma.from_documents(
                documents=chunks,
                embedding=OllamaEmbeddings(model=self.cfg.embedding_model),
                collection_name=self.cfg.vector_store_name,
                persist_directory=self.cfg.chroma_path,
            )

        with timed(self.logger, "Spara uppdaterat manifest (endast lyckade)"):
            # Lägg till ENBART de som faktiskt gav docs
            new_manifest = sorted(set(processed_files) | set(succeeded))
            self._save_manifest(manifest_path, new_manifest)

        self._log_ingest_summary(succeeded, failed, final=True)
        return vectordb

    # -------- Interna hjälpmetoder --------
    def _walk_files(self, start_dir: str) -> List[str]:
        """Gå rekursivt genom start_dir och returnera alla filer med SUPPORTED_EXTS.
        Respekterar ignore_dirs, follow_symlinks och hoppar över chroma_path om den ligger under roten."""
        results: List[str] = []
        start_abs = os.path.abspath(start_dir)
        chroma_abs = os.path.abspath(self.cfg.chroma_path)

        for root, dirs, files in os.walk(start_abs, followlinks=self.cfg.follow_symlinks):
            # Pruna bort ignorerade mappar
            dirs[:] = [d for d in dirs if d not in self.cfg.ignore_dirs]

            # Hoppa över chroma_path om den råkar ligga under doc_folder
            # (och dess underkataloger som _converted)
            pruned = []
            for d in dirs:
                d_abs = os.path.abspath(os.path.join(root, d))
                if d_abs == chroma_abs or d_abs.startswith(chroma_abs + os.sep):
                    continue
                pruned.append(d)
            dirs[:] = pruned

            for f in files:
                if f.lower().endswith(SUPPORTED_EXTS):
                    results.append(os.path.normpath(os.path.join(root, f)))
        return results

    def _configure_logging_once(self):
        if not getattr(self, "_logging_configured", False):
            level = logging.DEBUG if self.cfg.verbose else logging.INFO
            logging.basicConfig(
                level=level,
                format="%(asctime)s | %(levelname)s | %(message)s",
                datefmt="%H:%M:%S",
            )
            self._logging_configured = True
            # self.logger.debug("Logger konfigurerad (verbose=%s)", self.cfg.verbose)

    def _ensure_dirs(self) -> None:
        if not os.path.isdir(self.cfg.doc_folder):
            raise NotADirectoryError(f"Ogiltig katalog: {self.cfg.doc_folder}")
        os.makedirs(self.cfg.chroma_path, exist_ok=True)

    def _list_target_items(self, items: Optional[Iterable[str]] = None) -> List[str]:
        """
        Returnerar en lista av:
          - normaliserade paths (absoluta) för lokala filer med SUPPORTED_EXTS (rekursivt om directory)
          - oförändrade Google Sheets-URL:er
        """
        results: List[str] = []

        # Om användaren anger explicita items (filer, mappar eller URL:er)
        if items:
            for x in items:
                if _is_gsheet_url(x):
                    results.append(x.strip())
                    continue

                # Relativ → absolut från doc_folder
                x_path = x if os.path.isabs(x) else os.path.join(self.cfg.doc_folder, x)
                x_path = os.path.normpath(x_path)

                if os.path.isdir(x_path):
                    # Rekursiv genomsökning av mappen
                    results.extend(self._walk_files(x_path))
                elif os.path.isfile(x_path) and x_path.lower().endswith(SUPPORTED_EXTS):
                    results.append(x_path)
                else:
                    pass
                    # Ignorera okända/icke-existerande paths tyst (eller logga om du vill)
                    # self.logger.debug(f"Ignorerar item (ej fil/URL eller fel ändelse): {x}")
            return sorted(set(results))

        # Annars: skanna doc_folder (rekursivt om så är konfigurerat)
        if self.cfg.recursive:
            return sorted(set(self._walk_files(self.cfg.doc_folder)))

        # Icke-rekursiv fallback (endast toppnivå)
        results = []
        for f in os.listdir(self.cfg.doc_folder):
            if f.lower().endswith(SUPPORTED_EXTS):
                results.append(os.path.normpath(os.path.join(self.cfg.doc_folder, f)))
        return sorted(set(results))

    def _load_documents(self, targets: Iterable[str]) -> Tuple[List[Document], List[str], List[str]]:
        all_docs: List[Document] = []
        succeeded: List[str] = []
        failed: List[str] = []

        for target in targets:
            t0 = time.perf_counter()
            is_url = _is_gsheet_url(target)
            label = target if is_url else os.path.basename(target)
            try:
                if is_url:
                    data = self._load_gsheet(target)
                else:
                    ext = os.path.splitext(target)[1].lower()
                    if ext == ".pdf":
                        data = self._load_pdf(target)
                    elif ext == ".txt":
                        data = self._load_txt(target)
                    elif ext in (".docx", ".doc"):
                        data = self._load_word(target, ext)
                    elif ext in (".csv",):
                        data = self._load_csv(target)
                    elif ext in (".xlsx", ".xls"):
                        data = self._load_excel(target)
                    elif ext == ".ods":
                        data = self._load_ods(target)
                    elif ext in (".md", ".markdown"):
                        data = self._load_markdown(target)
                    elif ext in (".html", ".htm"):
                        data = self._load_html(target)
                    elif ext in (".pptx", ".ppt"):
                        data = self._load_ppt(target)
                    elif ext == ".odt":
                        data = self._load_odt(target)
                    elif ext == ".odp":
                        data = self._load_odp(target)
                    elif ext == ".tsv":
                        data = self._load_tsv(target)
                    elif ext in (".json", ".jsonl"):
                        data = self._load_json(target)
                    else:
                        data = []

                if data:
                    all_docs.extend(data)
                    succeeded.append(target if is_url else os.path.normpath(target))
                else:
                    failed.append(target if is_url else os.path.normpath(target))
            except Exception as e:
                # self.logger.warning(f"Misslyckades att läsa {target}: {e}")
                failed.append(target if is_url else os.path.normpath(target))
            finally:
                dt = time.perf_counter() - t0
                # self.logger.info(f"⏱️  Mål klart: {label} tog {dt:.2f}s")

        return all_docs, succeeded, failed

    # ----- Enskilda loaders -----
    def _load_pdf(self, full_path: str) -> List[Document]:
        try:
            data = PyPDFLoader(file_path=full_path).load()
            # self.logger.info(f"Läste PDF: {full_path} (docs: {len(data)})")
            return data
        except Exception as e:
            # self.logger.warning(f"PyPDFLoader misslyckades för {full_path}: {e}")
            if self.cfg.use_pdf_fallback:
                # TODO: lägg ev. fitz-baserad fallback
                pass
            return []

    def _load_txt(self, full_path: str) -> List[Document]:
        data = TextLoader(file_path=full_path, encoding="utf-8").load()
        # self.logger.info(f"Läste TXT: {full_path} (docs: {len(data)})")
        return data

    def _load_word(self, full_path: str, ext: str) -> List[Document]:
        if HAS_UNSTRUCTURED_WORD:
            try:
                data = UnstructuredWordDocumentLoader(full_path).load()
                # self.logger.info(f"Läste Word via Unstructured: {full_path} (docs: {len(data)})")
                return data
            except Exception as e:
                pass
                # self.logger.warning(f"Unstructured (Word) misslyckades för {full_path}: {e}")

        if ext == ".docx":
            try:
                data = Docx2txtLoader(full_path).load()
                # self.logger.info(f"Läste DOCX via docx2txt: {full_path} (docs: {len(data)})")
                return data
            except Exception as e:
                # self.logger.warning(f"Docx2txt misslyckades för {full_path}: {e}")
                return []

        if ext == ".doc" and self.cfg.convert_doc_with_soffice:
            try:
                converted = self._convert_doc_to_docx_with_soffice(full_path)
                data = Docx2txtLoader(converted).load()
                # self.logger.info(f"Läste DOC (konverterad): {full_path} -> {converted} (docs: {len(data)})")
                return data
            except Exception as e:
                pass
                # self.logger.warning(f"Konvertering .doc → .docx misslyckades för {full_path}: {e}")

        # self.logger.warning(f"Inget fungerande sätt att läsa {full_path}. "f"Tips: installera 'unstructured' eller aktivera convert_doc_with_soffice.")
        return []

    def _load_csv(self, full_path: str) -> List[Document]:
        try:
            loader = CSVLoader(file_path=full_path, autodetect_encoding=True)
            data = loader.load()
            # self.logger.info(f"Läste CSV: {full_path} (docs: {len(data)})")
            return data
        except Exception as e:
            # self.logger.warning(f"CSVLoader misslyckades för {full_path}: {e}")
            # Fallback med pandas
            try:
                df = pd.read_csv(full_path)
                data = self._df_to_documents(df, source=full_path)
                # self.logger.info(f"Läste CSV via pandas: {full_path} (docs: {len(data)})")
                return data
            except Exception as e2:
                # self.logger.warning(f"Pandas CSV-fallback misslyckades för {full_path}: {e2}")
                return []

    def _load_excel(self, full_path: str) -> List[Document]:
        # 1) Försök först med pandas (snabbare, inga extra beroenden)
        docs: List[Document] = []
        try:
            dfs = pd.read_excel(full_path, sheet_name=None, dtype=str)  # läs alla blad som strängar
            for sheet_name, df in dfs.items():
                # Vanligt fall: det finns rader -> konvertera direkt
                if not df.empty:
                    docs.extend(self._df_to_documents(df, source=f"{full_path}#{sheet_name}"))
                    continue

                # Specialfall: blad med bara rubriker / bara en cell
                # Försök igen med header=None så att första raden blir data
                try:
                    df2 = pd.read_excel(full_path, sheet_name=sheet_name, header=None, dtype=str)
                    if not df2.empty:
                        # ge generiska kolumnnamn
                        df2.columns = [f"col{i + 1}" for i in range(df2.shape[1])]
                        docs.extend(self._df_to_documents(df2, source=f"{full_path}#{sheet_name}"))
                        continue
                except Exception as e:
                    pass
                    # self.logger.debug(f"Excel header=None fallback misslyckades ({sheet_name}): {e}")

                # Sista utväg: skapa ett enda Document av “rubrikerna” om något finns
                if len(df.columns) > 0:
                    header_text = " | ".join(map(str, df.columns.tolist()))
                    from langchain.docstore.document import Document
                    docs.append(Document(
                        page_content=header_text,
                        metadata={"source": f"{full_path}#{sheet_name}", "type": "sheet_header_only"},
                    ))
            if docs:
                # self.logger.info(f"Läste Excel via pandas: {full_path} (docs: {len(docs)})")
                return docs
        except Exception as e2:
            pass
            # self.logger.warning(f"Pandas Excel-läsning misslyckades för {full_path}: {e2}")

        # 2) Fallback: Unstructured (kan varna om msoffcrypto saknas – det är okej)
        if HAS_UNSTRUCTURED_EXCEL:
            try:
                data = UnstructuredExcelLoader(full_path).load()
                # self.logger.info(f"Läste Excel via Unstructured: {full_path} (docs: {len(data)})")
                return data
            except Exception as e:
                pass
                # self.logger.warning(f"Unstructured (Excel) misslyckades för {full_path}: {e}")

        return []

    def _load_ods(self, full_path: str) -> List[Document]:
        # Försök Pandas med odfpy först — dtype=str minskar typstrul (t.ex. Datum)
        try:
            dfs = pd.read_excel(full_path, sheet_name=None, engine="odf", dtype=str)
            docs: List[Document] = []
            for sheet_name, df in dfs.items():
                docs.extend(self._df_to_documents(df, source=f"{full_path}#{sheet_name}"))
            # self.logger.info(f"Läste ODS via pandas(odf): {full_path} (docs: {len(docs)})")
            return docs
        except Exception as e:
            pass
            # self.logger.warning(f"ODS-läsning via pandas(odf) misslyckades för {full_path}: {e}")

        # Fallback: konvertera ODS → XLSX med LibreOffice och läs igen
        if self.cfg.convert_ods_with_soffice:
            try:
                converted = self._convert_ods_to_xlsx_with_soffice(full_path)
                dfs = pd.read_excel(converted, sheet_name=None)  # openpyxl/xlrd auto
                docs: List[Document] = []
                for sheet_name, df in dfs.items():
                    docs.extend(self._df_to_documents(df, source=f"{converted}#{sheet_name}"))
                # self.logger.info(f"Läste ODS (konverterad till XLSX): {full_path} -> {converted} (docs: {len(docs)})")
                return docs
            except Exception as e2:
                pass
                # self.logger.warning(f"ODS-konvertering/läsning misslyckades för {full_path}: {e2}")

        # Gav upp
        return []

    def _load_markdown(self, full_path: str):
        if HAS_UNSTRUCTURED_MD:
            try:
                data = UnstructuredMarkdownLoader(full_path).load()
                # self.logger.info(f"Läste Markdown via Unstructured: {full_path} (docs: {len(data)})")
                return data
            except Exception as e:
                pass
                # self.logger.warning(f"Unstructured (MD) misslyckades för {full_path}: {e}")
        # Fallback: vanlig text
        try:
            data = TextLoader(file_path=full_path, encoding="utf-8").load()
            # self.logger.info(f"Läste Markdown som text: {full_path} (docs: {len(data)})")
            return data
        except Exception as e:
            # self.logger.warning(f"MD-fallback misslyckades för {full_path}: {e}")
            return []

    def _load_html(self, full_path: str):
        if HAS_UNSTRUCTURED_HTML:
            try:
                data = UnstructuredHTMLLoader(full_path).load()
                # self.logger.info(f"Läste HTML via Unstructured: {full_path} (docs: {len(data)})")
                return data
            except Exception as e:
                pass
                # self.logger.warning(f"Unstructured (HTML) misslyckades för {full_path}: {e}")
        # Fallback: BeautifulSoup -> ren text
        if HAS_BS4:
            try:
                with open(full_path, "r", encoding="utf-8", errors="ignore") as f:
                    html = f.read()
                text = BeautifulSoup(html, "lxml").get_text(separator="\n")
                from langchain.docstore.document import Document
                data = [Document(page_content=text, metadata={"source": full_path, "type": "html"})]
                # self.logger.info(f"Läste HTML via BeautifulSoup: {full_path} (docs: {len(data)})")
                return data
            except Exception as e:
                pass
                # self.logger.warning(f"HTML BeautifulSoup-fallback misslyckades för {full_path}: {e}")
        return []

    def _load_ppt(self, full_path: str) -> List[Document]:
        """
        Läser PPT(X):
        - För .pptx: försök Unstructured → python-pptx.
        - För .ppt: konvertera med soffice till .pptx och läs den.
        """
        ext = os.path.splitext(full_path)[1].lower()

        def _load_pptx_inner(pptx_path: str) -> List[Document]:
            # 1) Unstructured (brukar funka bra på pptx)
            if HAS_UNSTRUCTURED_PPT:
                try:
                    data = UnstructuredPowerPointLoader(pptx_path).load()
                    # self.logger.info(f"Läste PPT via Unstructured: {pptx_path} (docs: {len(data)})")
                    return data
                except Exception as e:
                    pass
                    # self.logger.warning(f"Unstructured (PPTX) misslyckades för {pptx_path}: {e}")

            # 2) Fallback: python-pptx – extrahera text från shapes
            if HAS_PYTHON_PPTX:
                try:
                    import pptx
                    prs = pptx.Presentation(pptx_path)
                    chunks = []
                    for i, slide in enumerate(prs.slides, start=1):
                        texts = []
                        for shp in slide.shapes:
                            if hasattr(shp, "text") and shp.text:
                                texts.append(shp.text)
                        if texts:
                            chunks.append(Document(
                                page_content="\n".join(texts),
                                metadata={"source": pptx_path, "type": "slide", "slide": i},
                            ))
                    # self.logger.info(f"Läste PPT via python-pptx: {pptx_path} (docs: {len(chunks)})")
                    return chunks
                except Exception as e:
                    pass
                    # self.logger.warning(f"python-pptx-fallback misslyckades för {pptx_path}: {e}")

            return []

        # ----- Grenar per filändelse -----
        if ext == ".pptx":
            return _load_pptx_inner(full_path)

        if ext == ".ppt":
            # Konvertera till pptx först (default: på). Kräver LibreOffice.
            if getattr(self.cfg, "convert_ppt_with_soffice", True):
                try:
                    converted = self._convert_any_with_soffice(full_path, "pptx")
                    # self.logger.info(f"Konverterad .ppt → .pptx: {full_path} -> {converted}")
                    data = _load_pptx_inner(converted)
                    if data:
                        return data
                except Exception as e:
                    pass
                    # self.logger.warning(f"PPT-konvertering misslyckades för {full_path}: {e}")

            # Sista utväg: prova Unstructured direkt på .ppt (kräver soffice i PATH)
            if HAS_UNSTRUCTURED_PPT:
                try:
                    data = UnstructuredPowerPointLoader(full_path).load()
                    # self.logger.info(f"Läste PPT via Unstructured (direkt): {full_path} (docs: {len(data)})")
                    return data
                except Exception as e:
                    pass
                    # self.logger.warning(f"Unstructured (PPT) misslyckades för {full_path}: {e}")

            return []

        # Okänt tillägg (ska inte hända, vi routar bara .ppt/.pptx hit)
        # self.logger.warning(f"_load_ppt anropad med okänt tillägg: {full_path}")
        return []

    def _load_odt(self, full_path: str):
        # 1) Unstructured ODT
        if HAS_UNSTRUCTURED_ODT:
            try:
                data = UnstructuredODTLoader(full_path).load()
                # self.logger.info(f"Läste ODT via Unstructured: {full_path} (docs: {len(data)})")
                return data
            except Exception as e:
                pass
                # self.logger.warning(f"Unstructured (ODT) misslyckades för {full_path}: {e}")
        # 2) Fallback: konvertera ODT -> DOCX -> Docx2txt
        if self.cfg.convert_odt_with_soffice:
            try:
                converted = self._convert_any_with_soffice(full_path, "docx")
                data = Docx2txtLoader(converted).load()
                # self.logger.info(f"Läste ODT (konverterad till DOCX): {full_path} -> {converted} (docs: {len(data)})")
                return data
            except Exception as e:
                pass
                # self.logger.warning(f"ODT-konvertering misslyckades för {full_path}: {e}")
        return []

    def _load_odp(self, full_path: str):
        # 1) Unstructured PPT loader brukar klara .odp också, annars konvertera
        if HAS_UNSTRUCTURED_PPT:
            try:
                data = UnstructuredPowerPointLoader(full_path).load()
                # self.logger.info(f"Läste ODP via Unstructured: {full_path} (docs: {len(data)})")
                return data
            except Exception as e:
                pass
                # self.logger.warning(f"Unstructured (ODP) misslyckades för {full_path}: {e}")
        # 2) Fallback: konvertera ODP -> PPTX -> python-pptx
        if self.cfg.convert_odp_with_soffice:
            try:
                converted = self._convert_any_with_soffice(full_path, "pptx")
                return self._load_ppt(converted)
            except Exception as e:
                pass
                # self.logger.warning(f"ODP-konvertering misslyckades för {full_path}: {e}")
        return []

    def _load_tsv(self, full_path: str):
        try:
            import pandas as pd
            df = pd.read_csv(full_path, sep="\t")
            data = self._df_to_documents(df, source=full_path)
            # self.logger.info(f"Läste TSV via pandas: {full_path} (docs: {len(data)})")
            return data
        except Exception as e:
            pass
            # self.logger.warning(f"TSV-läsning misslyckades för {full_path}: {e}")
            return []

    def _load_json(self, full_path: str):
        try:
            # För JSONL/NDJSON (en rad = ett objekt)
            if full_path.lower().endswith(".jsonl"):
                import pandas as pd
                df = pd.read_json(full_path, lines=True)
                data = self._df_to_documents(df, source=full_path)
                # self.logger.info(f"Läste JSONL via pandas: {full_path} (docs: {len(data)})")
                return data

            # Vanlig JSON
            import json
            with open(full_path, "r", encoding="utf-8") as f:
                obj = json.load(f)

            from langchain.docstore.document import Document
            if isinstance(obj, list):
                # lista av objekt -> DataFrame -> Documents
                import pandas as pd
                df = pd.DataFrame(obj)
                data = self._df_to_documents(df, source=full_path)
                # self.logger.info(f"Läste JSON-lista som tabell: {full_path} (docs: {len(data)})")
                return data
            elif isinstance(obj, dict):
                # Försök hitta en "records"-nyckel eller platta ut nycklar
                import pandas as pd
                try:
                    df = pd.json_normalize(obj, max_level=1)
                    data = self._df_to_documents(df, source=full_path)
                    # self.logger.info(f"Läste JSON-dict via normalize: {full_path} (docs: {len(data)})")
                    return data
                except Exception:
                    # annars 1 dokument med pretty JSON
                    text = json.dumps(obj, ensure_ascii=False, indent=2)
                    return [Document(page_content=text, metadata={"source": full_path, "type": "json"})]
            else:
                text = str(obj)
                return [Document(page_content=text, metadata={"source": full_path, "type": "json"})]
        except Exception as e:
            # self.logger.warning(f"JSON-läsning misslyckades för {full_path}: {e}")
            return []

    def _load_gsheet(self, url: str) -> List[Document]:
        """
        Simpel public/export-hämtning:
         - Tar en vanlig Google Sheets-URL och hämtar CSV-exporten.
         - För privata ark krävs OAuth + GoogleDriveLoader (se README/inställningar).
        """
        try:
            sheet_id = None
            gid = None
            # URL-format: https://docs.google.com/spreadsheets/d/<id>/edit#gid=<gid>
            parts = url.split("/d/")
            if len(parts) >= 2:
                rest = parts[1]
                sheet_id = rest.split("/")[0]
            if "#gid=" in url:
                gid = url.split("#gid=")[-1].split("&")[0]

            if not sheet_id:
                raise ValueError("Kunde inte parsa Google Sheets-ID från URL.")

            export_url = f"https://docs.google.com/spreadsheets/d/{sheet_id}/export?format=csv"
            if gid:
                export_url += f"&gid={gid}"

            resp = requests.get(export_url, timeout=30)
            if resp.status_code != 200:
                raise RuntimeError(f"Sheets export gav HTTP {resp.status_code} – kräver sannolikt delning eller OAuth.")

            # Spara som temp .csv och låt CSVLoader skapa Documents (radvis)
            with tempfile.NamedTemporaryFile("wb", suffix=".csv", delete=False) as tmp:
                tmp.write(resp.content)
                temp_path = tmp.name

            try:
                data = CSVLoader(file_path=temp_path, autodetect_encoding=True).load()
                # self.logger.info(f"Läste Google Sheet via export: {url} (docs: {len(data)})")
                return data
            finally:
                try:
                    os.remove(temp_path)
                except Exception:
                    pass

        except Exception as e:
            # self.logger.warning(f"Google Sheets-läsning misslyckades för URL {url}: {e}. "f"För privata ark: använd langchain_google_community.GoogleDriveLoader.")
            return []

    # ----- Hjälpare -----
    def _df_to_documents(self, df: pd.DataFrame, source: str) -> List[Document]:
        """
        Gör om en DataFrame till Documents (en per rad) på ett robust sätt.
        - Tvingar kolumnnamn -> str för att undvika KeyError vid numeriska etiketter (t.ex. 2008).
        - Fyller NaN med "".
        """
        # Gör en kopia så vi inte muterar originalet
        df = df.copy()

        # 1) Se till att kolumnnamn är strängar (fixar KeyError '2008')
        df.columns = ["" if c is None else str(c) for c in df.columns]

        # 2) Fyll NaN -> ""
        df = df.fillna("")

        # 3) Konvertera radvis till text
        cols = list(df.columns)
        docs: List[Document] = []
        if isinstance(getattr(self, "cfg", None), IngestConfig) and self.cfg.max_rows_per_sheet:
            if len(df) > self.cfg.max_rows_per_sheet:
                # self.logger.warning(f"Trunkerar {source}: {len(df)} rader -> {self.cfg.max_rows_per_sheet} rader för prestanda.")
                df = df.iloc[: self.cfg.max_rows_per_sheet]
        for rec in df.to_dict(orient="records"):
            parts = [f"{c}: {rec.get(c, '')}" for c in cols]
            text = " | ".join(map(str, parts))
            docs.append(Document(page_content=text, metadata={"source": source, "type": "table_row"}))
        return docs

    def _convert_doc_to_docx_with_soffice(self, full_path: str) -> str:
        soffice = self.cfg.soffice_path or shutil.which("soffice") or shutil.which("soffice.exe")
        if not soffice:
            raise RuntimeError("LibreOffice 'soffice' hittas inte. Ange soffice_path i config eller lägg i PATH.")
        outdir = os.path.join(self.cfg.chroma_path, "_converted")
        os.makedirs(outdir, exist_ok=True)
        cmd = [soffice, "--headless", "--convert-to", "docx", "--outdir", outdir, full_path]
        # self.logger.info(f"Konverterar .doc → .docx med: {' '.join(cmd)}")
        res = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, check=True)
        # self.logger.debug(f"LibreOffice stdout: {res.stdout.strip()}")
        # self.logger.debug(f"LibreOffice stderr: {res.stderr.strip()}")
        base = os.path.splitext(os.path.basename(full_path))[0] + ".docx"
        new_path = os.path.join(outdir, base)
        if not os.path.exists(new_path):
            raise RuntimeError(f"Förväntad konverterad fil saknas: {new_path}")
        return new_path

    def _convert_ods_to_xlsx_with_soffice(self, full_path: str) -> str:
        soffice = self.cfg.soffice_path or shutil.which("soffice") or shutil.which("soffice.exe")
        if not soffice:
            raise RuntimeError("LibreOffice 'soffice' hittas inte. Ange soffice_path i config eller lägg i PATH.")
        outdir = os.path.join(self.cfg.chroma_path, "_converted")
        os.makedirs(outdir, exist_ok=True)
        cmd = [soffice, "--headless", "--convert-to", "xlsx", "--outdir", outdir, full_path]
        # self.logger.info(f"Konverterar .ods → .xlsx med: {' '.join(cmd)}")
        res = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, check=True)
        # self.logger.debug(f"LibreOffice stdout: {res.stdout.strip()}")
        # self.logger.debug(f"LibreOffice stderr: {res.stderr.strip()}")
        base = os.path.splitext(os.path.basename(full_path))[0] + ".xlsx"
        new_path = os.path.join(outdir, base)
        if not os.path.exists(new_path):
            raise RuntimeError(f"Förväntad konverterad fil saknas: {new_path}")
        return new_path

    def _convert_any_with_soffice(self, full_path: str, target_ext: str) -> str:
        soffice = self.cfg.soffice_path or shutil.which("soffice") or shutil.which("soffice.exe")
        if not soffice:
            raise RuntimeError("LibreOffice 'soffice' hittas inte. Ange soffice_path i config eller lägg i PATH.")
        outdir = os.path.join(self.cfg.chroma_path, "_converted")
        os.makedirs(outdir, exist_ok=True)
        cmd = [soffice, "--headless", "--convert-to", target_ext, "--outdir", outdir, full_path]
        # self.logger.info(f"Konverterar med soffice: {' '.join(cmd)}")
        res = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, check=True)
        # self.logger.debug(f"LibreOffice stdout: {res.stdout.strip()}")
        # self.logger.debug(f"LibreOffice stderr: {res.stderr.strip()}")
        base = os.path.splitext(os.path.basename(full_path))[0] + f".{target_ext}"
        new_path = os.path.join(outdir, base)
        if not os.path.exists(new_path):
            raise RuntimeError(f"Förväntad konverterad fil saknas: {new_path}")
        return new_path

    def _split_documents(self, documents: List[Document]):
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.cfg.chunk_size,
            chunk_overlap=self.cfg.chunk_overlap
        )
        return splitter.split_documents(documents)

    def _load_manifest(self, manifest_path: str) -> List[str]:
        if os.path.exists(manifest_path):
            try:
                with open(manifest_path, "r", encoding="utf-8") as f:
                    return json.load(f)
            except Exception as e:
                pass
                # self.logger.warning(f"Kunde inte läsa manifest {manifest_path}: {e}")
        return []

    def _save_manifest(self, manifest_path: str, items: List[str]) -> None:
        try:
            # behåll URL:er som är Sheets; normalisera bara paths
            normed = [i if _is_gsheet_url(i) else os.path.normpath(i) for i in items]
            with open(manifest_path, "w", encoding="utf-8") as f:
                json.dump(normed, f, ensure_ascii=False, indent=2)
        except Exception as e:
            pass
            # self.logger.warning(f"Kunde inte spara manifest {manifest_path}: {e}")

    def _ensure_model_available(self, model: str) -> None:
        try:
            local = {m["model"] for m in ollama.list().get("models", [])}
        except Exception:
            local = set()
        if model in local:
            # self.logger.info(f"Modellen '{model}' finns redan lokalt – hoppar över pull.")
            return
        # self.logger.info(f"Modellen '{model}' saknas lokalt – kör ollama.pull (kan ta tid första gången).")
        ollama.pull(model)

    def _log_ingest_summary(self, succeeded: List[str], failed: List[str], stage: str = "", final: bool = False):
        prefix = "SLUTRAPPORT" if final else f"Sammanfattning ({stage})"
        # self.logger.info("─" * 60)
        # self.logger.info(f"{prefix}: {len(succeeded)} lyckade, {len(failed)} misslyckade")
        if succeeded:
            # self.logger.info("Lyckade:")
            for p in succeeded:
                pass
                # self.logger.info(f"  ✓ {p}")
        if failed:
            # self.logger.warning("Misslyckade (ej sparade i manifest):")
            for p in failed:
                pass
                # self.logger.warning(f"  ✗ {p}")
        #self.logger.info("─" * 60)