from typing import Any, Dict, List
from collections.abc import Callable
import asyncio
import tempfile
from pathlib import Path
import pandas as pd
from PIL import Image
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from .flow import FlowComponent
from ..interfaces.ppt import PowerPointFile
from ..exceptions import ComponentError, ConfigError, DataNotFound
from ..utils.functions import is_empty


class PandasToPowerPoint(PowerPointFile, FlowComponent):
    """
    PandasToPowerPoint

    Overview:
        FlowTask component for generating PowerPoint presentations from DataFrame rows.
        Each row becomes a slide using hybrid template/programmatic approach.
        Uses PowerPointFile interface for reusable functionality.

    Properties:
        template_file (str): Path to PowerPoint template file (.pptx or .potx)
        column_mapping (dict): Maps DataFrame columns to PowerPoint sections
        reference_text (dict): Default text for missing DataFrame columns
        output_file (str): Output PowerPoint file path
        slide_layout (int): Preferred slide layout index (fallback if auto-selection fails)
        markdown_columns (list): Columns containing markdown to convert
        image_size (dict): Default image dimensions
        text_styles (dict): Default text styling
        missing_sections (dict): Configuration for creating missing placeholders
        overwrite (bool): Whether to overwrite existing presentation files

    Hybrid Approach:
        1. Analyze template layouts to find best match for required sections
        2. If good match found: use template placeholders
        3. If no good match: create layout programmatically
        4. Fall back to specified slide_layout if both fail

    Returns:
        DataFrame: Original + 'powerpoint_slide_created' + 'layout_used' columns

    |---|---|---|
    | version | No | version of component |


        Example:

        | Name | Required | Summary |
    |---|---|---|
    | version | No | version of component |


        Example:

        ```yaml
          PandasToPowerPoint:
          # attributes here
        ```
    """
    _version = "1.0.0"

    def __init__(
        self,
        loop: asyncio.AbstractEventLoop = None,
        job: Callable = None,
        stat: Callable = None,
        **kwargs,
    ):
        # Component-specific configuration
        self.template_file: str = kwargs.pop('template_file', 'powerpoint_template.pptx')
        self.column_mapping: Dict[str, str] = kwargs.pop('column_mapping', {})
        self.reference_text: Dict[str, str] = kwargs.pop('reference_text', {})
        self.output_file: str = kwargs.pop('output_file', 'presentation.pptx')
        self.slide_layout: int = kwargs.pop('slide_layout', 1)  # Fallback layout
        self.markdown_columns: List[str] = kwargs.pop('markdown_columns', [])
        self.slide_layout_name: str | None = kwargs.pop('slide_layout_name', None)
        self.missing_sections: Dict[str, Any] = kwargs.pop('missing_sections', {})
        self._overwrite: bool = kwargs.pop('overwrite', True)
        self.custom_text_styles: Dict[str, Any] = kwargs.pop('text_styles', {})

        # Validate required parameters
        if not self.template_file:
            raise ConfigError("template_file is required")
        if not self.column_mapping:
            raise ConfigError("column_mapping is required")

        # Initialize parent classes
        super().__init__(loop=loop, job=job, stat=stat, **kwargs)

        # Internal state
        self._slides_created = 0
        self._layout_stats = {
            'template_used': 0,
            'programmatic_used': 0,
            'fallback_used': 0,
            'failed': 0
        }

    def _is_image_column(self, column_name: str, section_name: str) -> bool:
        """
        Determine if a column should be treated as containing image data.
        Uses column names and section names as hints.
        """
        # Check column name for image indicators
        image_keywords = ['image', 'photo', 'picture', 'pic', 'img', 'overlay']
        if any(keyword in column_name.lower() for keyword in image_keywords):
            return True

        # Check section name for image indicators
        if any(keyword in section_name.lower() for keyword in image_keywords):
            return True

        # Specific known image placeholders
        image_sections = {'main_photo', 'photo', 'picture', 'image'}
        if section_name.lower() in image_sections:
            return True

        return False

    def _choose_best_layout_index(self) -> int:
        """
        Enhanced layout selection with slide_layout_name support.
        """
        # First, try to find layout by name if specified
        if self.slide_layout_name:
            for i, ly in enumerate(self._presentation.slide_layouts):
                if (ly.name or '').strip().lower() == self.slide_layout_name.strip().lower():
                    self._logger.debug(f"Found layout by name: '{self.slide_layout_name}' at index {i}")
                    return i
            self._logger.warning(f"Layout name '{self.slide_layout_name}' not found, using best match")

        # Use parent class method for best match
        return super()._choose_best_layout_index()

    async def start(self, **kwargs):
        """Initialize the component."""
        if self.previous:
            self.data = self.input
        else:
            if is_empty(self.data):
                raise DataNotFound("Previous Data was Not Found")

        await super().start(**kwargs)

        # Setup paths
        self.template_directory = self._taskstore.get_path().joinpath(self._program).joinpath('templates')
        template_path = self.template_directory.joinpath(self.template_file)

        if not template_path.exists():
            raise ComponentError(f"Template file not found: {self.template_file}")

        self.output_file = self.mask_replacement(self.output_file)
        self.output_path = self._filestore.get_directory('').joinpath('presentations')
        if not self.output_path.exists():
            self.output_path.mkdir(parents=True, exist_ok=True)

        # Load template using interface
        if not self.load_template(template_path):
            raise ComponentError(f"Failed to load PowerPoint template: {template_path}")

        self._log_template_info()
        return True

    def _log_template_info(self):
        """Log template information for debugging."""
        self._logger.info(f"Template has {len(self._presentation.slide_layouts)} slide layouts")

        # Log layout details
        catalog = self._layout_catalog()
        for i, (layout, placeholders) in enumerate(zip(self._presentation.slide_layouts, catalog.values())):
            layout_name = getattr(layout, 'name', f'Layout {i}')
            placeholder_names = list(placeholders.keys()) if placeholders else []
            self._logger.debug(f"Layout {i}: '{layout_name}' - Placeholders: {placeholder_names}")

        # Log column mapping
        self._logger.debug(f"Column mapping: {self.column_mapping}")

        # Check for potential matches
        wanted_placeholders = set(self.column_mapping.values())
        available_placeholders = set()
        for placeholders in catalog.values():
            available_placeholders.update(placeholders.keys())

        matching = wanted_placeholders.intersection(available_placeholders)
        missing = wanted_placeholders - available_placeholders

        if matching:
            self._logger.info(f"Template contains matching placeholders: {matching}")
        if missing:
            self._logger.info(f"Missing placeholders (will be created): {missing}")

    def _try_template_layout(self, row: pd.Series, row_index: int) -> tuple[bool, str]:
        """Try to create slide using template layouts."""
        attempts = []
        best = self._choose_best_layout_index()
        attempts.append(best)

        # Add fallback layout options
        for fallback in (self.slide_layout, 1, 6, 0):
            if fallback not in attempts and fallback < len(self._presentation.slide_layouts):
                attempts.append(fallback)

        for layout_idx in attempts:
            try:
                layout = self._presentation.slide_layouts[layout_idx]
                slide = self._presentation.slides.add_slide(layout)

                if self._populate_template_slide(slide, row, row_index):
                    self._slides_created += 1
                    self._layout_stats['template_used'] += 1
                    layout_name = getattr(layout, 'name', f'Layout {layout_idx}')
                    return True, f"template_{layout_idx}_{layout_name}"
                else:
                    # Remove the failed slide
                    self._remove_last_slide()

            except Exception as e:
                self._logger.debug(f"Template layout {layout_idx} failed: {e}")
                # Remove the failed slide if it was created
                try:
                    self._remove_last_slide()
                except Exception as e:
                    self._logger.warning(f"Failed to remove last slide: {e}")
                continue

        return False, "template_failed"

    def _remove_last_slide(self):
        """Remove the last slide from the presentation."""
        if len(self._presentation.slides) > 0:
            rId = self._presentation.slides._sldIdLst[-1].rId
            self._presentation.part.drop_rel(rId)
            del self._presentation.slides._sldIdLst[-1]

    def _populate_template_slide(self, slide, row: pd.Series, row_index: int) -> bool:
        """Populate slide using template placeholders and missing sections."""
        try:
            placeholders_found = 0
            autoplace_state = {}  # for missing sections without specs

            for column_name, section_name in self.column_mapping.items():
                # Get value from row
                val = row.get(column_name) if column_name in row and not pd.isna(row[column_name]) else None

                # Apply default text if value is missing
                if val is None:
                    val = self.reference_text.get(column_name, f"[{column_name}]")

                # Find or create placeholder first
                placeholder = self._ensure_placeholder(slide, section_name, autoplace_state)
                if not placeholder:
                    self._logger.warning(f"Could not create placeholder for '{section_name}'")
                    continue

                # Check if this column should be treated as an image
                is_image_column = self._is_image_column(column_name, section_name)

                if is_image_column:
                    # Try to process as image
                    img = self._normalize_image(val, column_name)
                    if img is not None:
                        result = self._replace_placeholder_with_image(slide, placeholder, img)
                        if result is not None:
                            placeholders_found += 1
                            self._logger.debug(f"Added image for {section_name} from column '{column_name}'")
                        else:
                            self._logger.warning(f"Failed to add image for {section_name}")
                        continue
                    else:
                        self._logger.warning(f"Expected image in column '{column_name}' but could not process as image")

                # Process markdown if specified (for text content)
                if column_name in self.markdown_columns and isinstance(val, str):
                    html_content = self.convert_markdown_to_html(val)
                    val = self.extract_text_from_html(html_content)

                # Handle text content
                if hasattr(placeholder, "text_frame"):
                    formatted_text = self._format_value(column_name, val)
                    placeholder.text_frame.text = formatted_text
                    self._apply_text_styles(placeholder.text_frame, section_name, self.custom_text_styles)
                    placeholders_found += 1
                    self._logger.debug(f"Set {section_name} = '{formatted_text}' from column '{column_name}'")
                else:
                    self._logger.warning(f"Placeholder {section_name} has no text_frame")

            success = placeholders_found > 0
            if success:
                self._logger.debug(f"Successfully populated slide with {placeholders_found} placeholders")
            else:
                self._logger.warning("No placeholders were populated on slide")

            return success

        except Exception as e:
            self._logger.error(f"Error populating template slide: {e}")
            return False

    def _create_programmatic_slide(self, row: pd.Series, row_index: int) -> tuple[bool, str]:
        """Create slide programmatically as fallback."""
        try:
            # Find blank layout (typically index 6, but fallback to 0)
            blank_idx = 6 if len(self._presentation.slide_layouts) > 6 else 0
            blank_layout = self._presentation.slide_layouts[blank_idx]
            slide = self._presentation.slides.add_slide(blank_layout)

            # Create basic layout programmatically
            if self._create_basic_layout(slide, row):
                self._slides_created += 1
                self._layout_stats['programmatic_used'] += 1
                return True, f"programmatic_{blank_idx}"
            else:
                # Remove failed slide
                self._remove_last_slide()
                return False, "programmatic_failed"

        except Exception as e:
            self._logger.error(f"Programmatic slide creation failed: {e}")
            return False, "programmatic_error"

    def _create_basic_layout(self, slide, row: pd.Series) -> bool:
        """Create basic layout programmatically."""
        try:
            spec = self.get_programmatic_layout_spec()["placeholders"]
            created_sections = 0

            for col, section in self.column_mapping.items():
                # Get and process value
                val = row.get(col) if col in row and not pd.isna(row[col]) else None
                if val is None:
                    val = self.reference_text.get(col, f"[{col}]")

                # Process markdown
                if col in self.markdown_columns and isinstance(val, str):
                    val = self.extract_text_from_html(self.convert_markdown_to_html(val))

                # Ensure placeholder exists with programmatic specs
                if section not in self.missing_sections and section in spec:
                    s = spec[section]
                    self.missing_sections[section] = {
                        "left_in": s["left"],
                        "top_in": s["top"],
                        "size_in": {"width": s["width"], "height": s["height"]}
                    }

                # Create autoplace state for this slide
                autoplace_state = {}
                shp = self._ensure_placeholder(slide, section, autoplace_state)
                if not shp:
                    continue

                # Check if this should be treated as an image
                is_image_column = self._is_image_column(col, section)

                if is_image_column:
                    # Handle images
                    img = self._normalize_image(val, col)
                    if img is not None:
                        if self._replace_placeholder_with_image(slide, shp, img):
                            created_sections += 1
                        continue

                # Handle text
                if hasattr(shp, "text_frame"):
                    formatted_text = self._format_value(col, val)
                    shp.text_frame.text = formatted_text
                    self._apply_text_styles(shp.text_frame, section, self.custom_text_styles)
                    created_sections += 1

            return created_sections > 0

        except Exception as e:
            self._logger.error(f"Error creating basic layout: {e}")
            return False

    def get_programmatic_layout_spec(self) -> Dict[str, Any]:
        """Return layout specification for programmatic slide creation."""
        return {
            "slide_width": 10.0,
            "slide_height": 7.5,
            "margin": 0.3,
            "header_height": 0.6,
            "footer_height": 0.5,
            "photo_width_percent": 0.75,
            "analysis_width_percent": 0.25,

            "placeholders": {
                "store_id_text": {
                    "left": 0.3, "top": 0.3, "width": 4.0, "height": 0.6,
                    "font_size": 16, "font_bold": True, "font_color": "003366"
                },
                "created_on_text": {
                    "left": 0.3, "top": 0.9, "width": 4.0, "height": 0.4,
                    "font_size": 11, "font_color": "666666"
                },
                "who_text": {
                    "left": 4.5, "top": 0.3, "width": 5.2, "height": 0.6,
                    "font_size": 12, "font_color": "666666"
                },
                "main_photo": {
                    "left": 0.3, "top": 1.4, "width": 7.2, "height": 4.8,
                    "type": "image"
                },
                "analysis_text": {
                    "left": 7.7, "top": 1.4, "width": 2.0, "height": 3.0,
                    "font_size": 10, "font_color": "333333"
                },
                "score_text": {
                    "left": 7.7, "top": 4.6, "width": 2.0, "height": 0.6,
                    "font_size": 14, "font_color": "008000"
                },
                "status_text": {
                    "left": 7.7, "top": 5.3, "width": 2.0, "height": 0.6,
                    "font_size": 12, "font_color": "008000"
                }
            }
        }

    def _create_slide_from_row(self, row: pd.Series, row_index: int) -> tuple[bool, str]:
        """Create slide from row with fallback strategy."""
        # Log row being processed
        self._logger.debug(f"Processing row {row_index + 1}/{len(self.data)}")

        # First try: Use template layout
        success, layout_type = self._try_template_layout(row, row_index)
        if success:
            return success, layout_type

        # Second try: Programmatic creation
        self._logger.debug(f"Template failed for row {row_index}, trying programmatic approach")
        success, layout_type = self._create_programmatic_slide(row, row_index)
        if success:
            return success, layout_type

        # Final fallback failed
        self._layout_stats['failed'] += 1
        self._logger.error(f"All slide creation methods failed for row {row_index}")
        return False, "failed"

    async def run(self):
        """
        Generate PowerPoint presentation from DataFrame rows.

        Returns:
            DataFrame: Original + 'powerpoint_slide_created' + 'layout_used' columns
        """
        if self.data is None or self.data.empty:
            raise ComponentError("No data available for processing")

        self._logger.info(f"Starting PowerPoint generation for {len(self.data)} rows")

        # Clear existing slides
        self.clear_existing_slides()

        # Create result DataFrame
        result_df = self.data.copy()
        result_df['powerpoint_slide_created'] = False
        result_df['layout_used'] = 'none'
        result_df['error_message'] = ''

        # Process each row
        for idx, row in self.data.iterrows():
            try:
                success, layout_type = self._create_slide_from_row(row, idx)
                result_df.at[idx, 'powerpoint_slide_created'] = success
                result_df.at[idx, 'layout_used'] = layout_type

                if success:
                    self._logger.debug(f"✓ Created slide {idx+1} using {layout_type}")
                else:
                    self._logger.warning(f"✗ Failed to create slide for row {idx+1}")

            except Exception as e:
                self._logger.error(f"Error processing row {idx+1}: {e}")
                result_df.at[idx, 'powerpoint_slide_created'] = False
                result_df.at[idx, 'layout_used'] = 'error'
                result_df.at[idx, 'error_message'] = str(e)
                self._layout_stats['failed'] += 1

        # Save presentation
        try:
            self.save_presentation(
                self.output_file,
                self.output_path,
                override=self._overwrite
            )
        except Exception as e:
            raise ComponentError(f"Failed to save PowerPoint presentation: {e}")

        # Log comprehensive statistics
        self._log_completion_stats(result_df)

        # Store result
        self._result = result_df
        return self._result

    def _log_completion_stats(self, result_df: pd.DataFrame):
        """Log comprehensive completion statistics."""
        successful_slides = result_df['powerpoint_slide_created'].sum()
        total_rows = len(result_df)

        self._logger.notice(f"PowerPoint generation completed: {successful_slides}/{total_rows} slides created")

        # Layout usage statistics
        if self._slides_created > 0:
            self._logger.info("Layout usage statistics:")
            for layout_type, count in self._layout_stats.items():
                if count > 0:
                    percentage = (count / self._slides_created) * 100
                    self._logger.info(f"  {layout_type}: {count} slides ({percentage:.1f}%)")

        # Error summary
        failed_count = self._layout_stats['failed']
        if failed_count > 0:
            self._logger.warning(f"Failed to create {failed_count} slides - check error_message column for details")

        # Column mapping success rate
        if hasattr(self, 'column_mapping'):
            self._logger.debug(f"Processed {len(self.column_mapping)} column mappings per slide")

    async def close(self):
        """Clean up resources."""
        self._presentation = None
        self._slides_created = 0
        self._layout_stats = {'template_used': 0, 'programmatic_used': 0, 'fallback_used': 0, 'failed': 0}
        return True
