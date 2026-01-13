from io import BytesIO
from abc import ABC
from pathlib import Path
from typing import Iterator, List, Dict, Any, Tuple, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
from navconfig.logging import logging


logging.getLogger("PIL").setLevel(logging.INFO)


class PowerPointClient(ABC):
    def create_presentation_from_template(
        self,
        template_path: Path,
        slide_contents: Iterator[Dict[str, Any]],
        file_path: Path,
        default_master_index: int = 0,
        default_layout_index: int = 1,
    ):
        """
        Create a presentation from a template, adding texts or images per slide.

        :param template_path: Path to the PowerPoint template file.
        :param slide_contents: An iterator where each item is a dict with content for a slide.
        :param file_path: Path where the new PowerPoint file will be saved.
        :param default_layout_index: Default slide layout index to use if not specified per slide.
        """
        prs = Presentation(template_path)
        prs = self.remove_existing_slides(prs)

        for content in slide_contents:
            master_index = content.get("master_index", default_master_index)
            layout_index = content.get("layout_index", default_layout_index)

            if master_index < 0 or master_index >= len(prs.slide_masters):
                raise ValueError(f"Invalid master index {master_index}.")
            master = prs.slide_masters[master_index]

            # Use the layout_index from content if specified, else use default
            if layout_index < 0 or layout_index >= len(master.slide_layouts):
                raise ValueError(
                    f"Invalid layout index {layout_index} for master {master_index}."
                )
            slide_layout = master.slide_layouts[layout_index]

            # Create slide with selected layout
            slide = prs.slides.add_slide(slide_layout)

            self._add_content_to_slide(
                slide, prs, content.get("text_content", {}), content.get("image")
            )
        prs.save(file_path)

        return file_path

    def _add_content_to_slide(
        self,
        slide,
        prs,  # Presentation Object
        content: List[dict],
        image_infos: Optional[List[Dict[str, Any]]],
    ):
        """
        Internal method to add content to a slide.

        :param slide: The slide object to add content to.
        :param prs: The Presentation object.
        :param content: List of dictionaries with text content.
        :param image_infos: List of dictionaries with image details.
        """
        # Create a placeholder map by name
        placeholder_map = {shape.name: shape for shape in slide.placeholders}

        # Process text content first.
        for text_entry in content:
            for placeholder_id, text_details in text_entry.items():
                text_content = text_details.get("text")
                if placeholder_id in placeholder_map:
                    try:
                        placeholder = placeholder_map[placeholder_id]
                        if placeholder.has_text_frame:
                            text_frame = placeholder.text_frame
                            new_paragraph = text_details.get("new_paragraph", False)

                            if new_paragraph:
                                p = text_frame.add_paragraph() if text_frame.paragraphs else text_frame.paragraphs[0]
                            else:
                                p = text_frame.paragraphs[0]

                            run = p.add_run()
                            run.text = text_content
                            font = run.font

                            if "font_name" in text_details:
                                font.name = text_details["font_name"]

                            if "font_size" in text_details:
                                font.size = Pt(text_details["font_size"])

                            if "font_color" in text_details:
                                font.color.rgb = RGBColor.from_string(text_details["font_color"])

                            if "bold" in text_details:
                                font.bold = text_details["bold"]

                            if "italic" in text_details:
                                font.italic = text_details["italic"]

                    except KeyError:
                        # Fallback: add a textbox manually if placeholder lookup fails
                        left = Inches(1)
                        top = Inches(1)
                        width = Inches(8)
                        height = Inches(2)
                        textbox = slide.shapes.add_textbox(left, top, width, height)
                        text_frame = textbox.text_frame
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = text_content
                        font = run.font
                        font.name = text_details.get("font_name", "Calibri")
                        font.size = text_details.get("font_size", Pt(18))
                        font.color.rgb = RGBColor.from_string(
                            text_details.get("font_color", "000000")
                        )
                        font.bold = text_details.get("bold", False)
                        font.italic = text_details.get("italic", False)

        # Process each image info in the list.

        if image_infos:
            for image in image_infos:
                # Step 1: Get source (BytesIO or file path)
                if "data" in image and isinstance(image["data"], BytesIO):
                    image_source = image["data"]
                    image_source.seek(0)
                elif "path" in image:
                    image_source = image["path"]
                else:
                    raise ValueError("Image must have either 'path' or 'data' field")

                # Step 2: Load image with PIL to check size and convert MPO if needed
                with Image.open(image_source) as img:
                    if img.format == "MPO":
                        img.seek(0)
                        img = img.convert("RGB")

                        if "path" in image:
                            # Overwrite the file if file-based
                            img.save(image["path"], format="JPEG")
                            image_source = image["path"]
                        else:
                            # Replace in-memory BytesIO with JPEG-converted version
                            buffer = BytesIO()
                            img.save(buffer, format="JPEG")
                            buffer.seek(0)
                            # image["data"] = buffer
                            image_source = buffer

                    img_width, img_height = img.size

                img_width = self.to_emu(img_width)
                img_height = self.to_emu(img_height)

                # Step 3: Try inserting into placeholder
                img_id = image.get("placeholder_id")
                if img_id in placeholder_map:
                    slide_holder = placeholder_map[img_id]
                    try:
                        picture = slide_holder.insert_picture(image_source)
                        print(picture, picture.shape_type)
                        continue
                    except Exception:
                        pass  # fallback to manual below

                # Step 4: Manual fallback or no placeholder
                slide_width = prs.slide_width
                slide_height = prs.slide_height

                left = image.get("left", Inches(1))
                top = image.get("top", Inches(2))

                if left == "center":
                    effective_width = slide_width - Inches(1)
                else:
                    effective_width = slide_width - left

                if top == "center":
                    effective_height = slide_height - Inches(1.5)
                else:
                    effective_height = slide_height - top

                scale_factor = image.get("scale_factor")
                if not scale_factor:
                    scale_factor = min(
                        effective_width / img_width,
                        effective_height / img_height,
                        1,
                    )

                width = int(img_width * scale_factor)
                height = int(img_height * scale_factor)

                if left == "center":
                    left = (slide_width - width) // 2
                if top == "center":
                    top = (slide_height - height) // 2

                width = image.get("width", width)
                height = image.get("height", height)

                slide.shapes.add_picture(
                    image_source, left, top, width=width, height=height
                )

    @staticmethod
    def create_empty_presentation(file_path: Path):
        """
        Create an empty PowerPoint presentation.

        :param file_path: Path where the new PowerPoint file will be saved.
        """
        prs = Presentation()
        prs.save(file_path)

    @staticmethod
    def create_presentation_from_text_list(text_list: List[str], file_path: Path):
        """
        Create a PowerPoint file where each slide contains text from a list.

        :param text_list: List of texts to add to slides.
        :param file_path: Path where the new PowerPoint file will be saved.
        """
        prs = Presentation()
        for text in text_list:
            slide = prs.slides.add_slide(
                prs.slide_layouts[1]
            )  # Title and Content layout
            textbox = slide.shapes.placeholders[1]
            textbox.text = text
        prs.save(file_path)

    @staticmethod
    def append_slides_to_presentation(
        file_path: Path, slide_contents: Iterator[Dict[str, Any]]
    ):
        """
        Edit a PowerPoint file to append more slides from an iterator.

        :param file_path: Path to the existing PowerPoint file.
        :param slide_contents: An iterator where each item is a dict with content for a slide.
        """
        prs = Presentation(file_path)
        for content in slide_contents:
            slide = prs.slides.add_slide(
                prs.slide_layouts[1]
            )  # Adjust layout as needed
            if "text" in content:
                textbox = slide.shapes.placeholders[1]
                textbox.text = content["text"]
            if "image" in content:
                left = Inches(1)
                top = Inches(2)
                slide.shapes.add_picture(content["image"], left, top, width=Inches(5))
        prs.save(file_path)

    @staticmethod
    def create_slide(
        text_content: Optional[str] = None,
        image_path: Optional[str] = None,
        image_size: Optional[Tuple[float, float, float]] = None,
        font_size: float = 18,
        font_color: str = "000000",
        font_name: str = "Calibri",
        layout_index: Optional[int] = None,
    ) -> Dict[str, Any]:
        """
        Create a slide content dictionary.

        :param text_content: Text to include in the slide.
        :param image_path: Path to the image file to include in the slide.
        :param image_size: Tuple of three elements (left, top, width) for image positioning.
        :param font_size: Font size for the text.
        :param font_color: Font color for the text in hex format (e.g., 'FF0000' for red).
        :param font_name: Font name for the text.
        :param layout_index: Optional index of the slide layout to use for this slide.
        :return: A dictionary representing the slide content.
        """
        slide_content = {}

        if text_content:
            slide_content["text"] = {
                "content": text_content,
                "font_size": Pt(font_size),
                "font_color": font_color,
                "font_name": font_name,
            }

        if image_path and image_size:
            left, top, width = image_size
            slide_content["image"] = {
                "path": image_path,
                "left": Inches(left),
                "top": Inches(top),
                "width": Inches(width),
            }

        if layout_index is not None:
            slide_content["layout_index"] = layout_index

        return slide_content

    @staticmethod
    def to_emu(size):
        return size * 9525

    @staticmethod
    def list_placeholder_ids(template_path: Path):
        prs = Presentation(template_path)
        for master_index, master in enumerate(prs.slide_masters):
            print(f"\nMaster Slide {master_index}:")

            for layout_index, layout in enumerate(master.slide_layouts):
                print(f"  Layout {layout_index}: {layout.name}")

                # Create a slide in a temporary presentation to inspect placeholder IDs
                temp_prs = Presentation()  # Temporary Presentation object
                slide = temp_prs.slides.add_slide(layout)

                for placeholder in slide.placeholders:
                    # Access idx directly from the XML element
                    idx = placeholder.element.get(
                        "idx"
                    )  # Retrieve 'idx' attribute directly from XML
                    name = placeholder.name
                    print(f"    Placeholder IDX: {idx}, Name: {name}")

    @staticmethod
    def remove_existing_slides(prs):
        for slide in list(prs.slides):
            # Get the slide ID and relationship ID
            slide_id = slide.slide_id
            rId = None

            # Get the list of slide IDs
            sldIdLst = prs.slides._sldIdLst

            # Remove the slide ID and get the relationship ID
            for sldId in sldIdLst:
                if sldId.id == slide_id:
                    rId = sldId.rId
                    sldIdLst.remove(sldId)
                    break

            if rId is not None:
                # Remove the relationship
                prs.part.drop_rel(rId)
            else:
                print(f"Slide ID {slide_id} not found.")

        return prs
