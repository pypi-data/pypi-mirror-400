import copy
import os
import re

import toml

from chgksuite.common import log_wrap, replace_escaped, tryint
from chgksuite.composer.composer_common import (
    BaseExporter,
    backtick_replace,
    parseimg,
    remove_accents_standalone,
)
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.util import Inches as PptxInches
from pptx.util import Pt as PptxPt


class PptxExporter(BaseExporter):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.config_path = os.path.abspath(self.args.pptx_config)
        with open(self.config_path, encoding="utf8") as f:
            self.c = toml.load(f)
        self.qcount = 0
        hs = self.labels["question_labels"]["handout"]
        self.re_handout_1 = re.compile(
            "\\[" + hs + ".(?P<body>.+?)\\]", flags=re.DOTALL
        )
        self.re_handout_2 = re.compile("^" + hs + ".(?P<body>.+?)$")

    def get_textbox_qnumber(self, slide):
        kwargs = {}
        for param in ("left", "top", "width", "height"):
            try:
                kwargs[param] = PptxInches(self.c["number_textbox"][param])
            except KeyError:
                pass

        return self.get_textbox(slide, **kwargs)

    def get_textbox(self, slide, left=None, top=None, width=None, height=None):
        if left is None:
            left = PptxInches(self.c["textbox"]["left"])
        if top is None:
            top = PptxInches(self.c["textbox"]["top"])
        if width is None:
            width = PptxInches(self.c["textbox"]["width"])
        if height is None:
            height = PptxInches(self.c["textbox"]["height"])
        textbox = slide.shapes.add_textbox(left, top, width, height)
        return textbox

    def add_run(self, para, text, color=None):
        r = para.add_run()
        r.text = text
        if color is None:
            color = self.c["textbox"].get("color")
        if color:
            r.font.color.rgb = RGBColor(*color)
        if self.args.language == "ru":
            r.font.language_id = MSO_LANGUAGE_ID.RUSSIAN
        return r

    def pptx_format(self, el, para, tf, slide, replace_spaces=True):
        def r_sp(text):
            if replace_spaces:
                return self._replace_no_break(text)
            return text

        if isinstance(el, list):
            if len(el) > 1 and isinstance(el[1], list):
                self.pptx_format(el[0], para, tf, slide)
                licount = 0
                for li in el[1]:
                    licount += 1
                    self.add_run(para, "\n{}. ".format(licount))
                    self.pptx_format(li, para, tf, slide)
            else:
                licount = 0
                for li in el:
                    licount += 1
                    self.add_run(para, "\n{}. ".format(licount))
                    self.pptx_format(li, para, tf, slide)

        if isinstance(el, str):
            self.logger.debug("parsing element {}:".format(log_wrap(el)))
            el = backtick_replace(el)

            for run in self.parse_4s_elem(el):
                if run[0] == "screen":
                    self.add_run(para, r_sp(run[1]["for_screen"]))

                elif run[0] == "linebreak":
                    self.add_run(para, "\n")

                elif run[0] == "strike":
                    r = self.add_run(para, r_sp(run[1]))
                    r.font.strike = True  # TODO: doesn't work as of 2023-12-24, cf. https://github.com/scanny/python-pptx/issues/339

                elif run[0] == "img":
                    pass  # image processing is moved to other places

                else:
                    r = self.add_run(para, r_sp(run[1]))
                    if "italic" in run[0]:
                        r.font.italic = True
                    if "bold" in run[0]:
                        r.font.bold = True
                    if "underline" in run[0]:
                        r.font.underline = True

    def pptx_process_text(
        self,
        s,
        image=None,
        strip_brackets=True,
        replace_spaces=True,
        do_not_remove_accents=False,
    ):
        hs = self.regexes["handout_short"]
        if isinstance(s, list):
            for i in range(len(s)):
                s[i] = self.pptx_process_text(s[i], image=image)
            return s
        if not (self.args.do_not_remove_accents or do_not_remove_accents):
            s = remove_accents_standalone(s, self.regexes)
        if strip_brackets:
            s = self.remove_square_brackets(s)
            s = s.replace("]\n", "]\n\n")
        else:
            s = replace_escaped(s)
        if image:
            s = re.sub("\\[" + hs + "(.+?)\\]", "", s, flags=re.DOTALL)
            s = s.strip()
        elif hs in s:
            re_hs = re.search("\\[" + hs + ".+?: ?(.+)\\]", s, flags=re.DOTALL)
            if re_hs:
                s = s.replace(re_hs.group(0), re_hs.group(1))
        s = re.sub(" +", " ", s)
        for punct in (".", ",", "!", "?", ":"):
            s = s.replace(" " + punct, punct)
        if replace_spaces:
            s = self._replace_no_break(s)
        s = s.strip()
        return s

    def apply_vertical_alignment_if_needed(self, text_frame):
        align = self.c["textbox"].get("vertical_align")
        if align:
            text_frame.auto_size = MSO_AUTO_SIZE.NONE
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0
            text_frame.vertical_anchor = getattr(MSO_VERTICAL_ANCHOR, align.upper())

    def _process_block(self, block):
        section = [x for x in block if x[0] == "section"]
        editor = [x for x in block if x[0] == "editor"]
        meta = [x for x in block if x[0] == "meta"]
        if not section and not editor and not meta:
            return
        slide = self.prs.slides.add_slide(self.BLANK_SLIDE)
        textbox = self.get_textbox(slide)
        tf = textbox.text_frame
        self.apply_vertical_alignment_if_needed(tf)
        tf.word_wrap = True
        text_for_size = (
            (self.recursive_join([x[1] for x in section]) or "")
            + "\n"
            + (self.recursive_join([x[1] for x in editor]) or "")
            + "\n"
            + (self.recursive_join([x[1] for x in meta]) or "")
        )
        p = self.init_paragraph(tf, text=text_for_size)
        add_line_break = False
        if section:
            if self.c.get("tour_as_question_number"):
                txt = self.pptx_process_text(section[0][1])
                if self.c.get("tour_as_question_number") == "caps":
                    txt = txt.upper()
                self.set_question_number(slide, number=txt)
            else:
                r = self.add_run(
                    p, self._replace_no_break(self.pptx_process_text(section[0][1]))
                )
                r.font.size = PptxPt(self.c["text_size_grid"]["section"])
                add_line_break = True
        if editor:
            r = self.add_run(
                p,
                self._replace_no_break(
                    ("\n\n" if add_line_break else "")
                    + self.pptx_process_text(editor[0][1])
                ),
            )
            add_line_break = True
        if meta:
            for element in meta:
                r = self.add_run(
                    p,
                    self._replace_no_break(
                        ("\n\n" if add_line_break else "")
                        + self.pptx_process_text(element[1])
                    ),
                )
                add_line_break = True

    def process_buffer(self, buffer):
        heading_block = []
        editor_block = []
        section_block = []
        block = heading_block
        for element in buffer:
            if element[0] == "section":
                block = section_block
            if element[0] == "editor" and not section_block:
                block = editor_block
            block.append(element)
        heading = [x for x in heading_block if x[0] == "heading"]
        ljheading = [x for x in heading_block if x[0] == "ljheading"]
        title_text = ljheading or heading
        date_text = [x for x in heading_block if x[0] == "date"]
        if title_text:
            if len(self.prs.slides) == 1:
                slide = self.prs.slides[0]
            else:
                slide = self.prs.slides.add_slide(self.TITLE_SLIDE)
            title = slide.shapes.title
            title.text = title_text[0][1]
            if date_text:
                try:
                    subtitle = slide.placeholders[1]
                    subtitle.text = date_text[0][1]
                except KeyError:
                    pass
        for block in (editor_block, section_block):
            self._process_block(block)

    def set_question_number(self, slide, number):
        if self.args.disable_numbers:
            return
        qntextbox = self.get_textbox_qnumber(slide)
        qtf = qntextbox.text_frame
        qtf_p = self.init_paragraph(qtf)
        if self.c["number_textbox"].get("align"):
            qtf_p.alignment = getattr(
                PP_ALIGN, self.c["number_textbox"]["align"].upper()
            )
        if self.c.get("question_number_format") == "caps" and tryint(number):
            number = f"ВОПРОС {number}"
        qtf_r = self.add_run(qtf_p, number)
        if self.c["number_textbox"].get("bold"):
            qtf_r.font.bold = True
        if self.c["number_textbox"].get("color"):
            qtf_r.font.color.rgb = RGBColor(*self.c["number_textbox"]["color"])
        if self.c["number_textbox"].get("font_size"):
            qtf_r.font.size = PptxPt(self.c["number_textbox"]["font_size"])

    def _get_handout_from_4s(self, text):
        if isinstance(text, list):
            for el in text:
                handout = self._get_handout_from_4s(el)
                if handout:
                    return handout
        elif isinstance(text, str):
            match_ = self.re_handout_1.search(text)
            if match_:
                return match_.group("body")
            else:
                lines = text.split("\n")
                for line in lines:
                    match_ = self.re_handout_2.search(line)
                    if match_:
                        return match_.group("body")

    def _get_image_from_4s(self, text):
        if isinstance(text, list):
            for el in text:
                image = self._get_image_from_4s(el)
                if image:
                    return image
        elif isinstance(text, str):
            for run in self.parse_4s_elem(text):
                if run[0] == "img":
                    parsed_image = parseimg(
                        run[1],
                        dimensions="inches",
                        tmp_dir=self.dir_kwargs.get("tmp_dir"),
                        targetdir=self.dir_kwargs.get("targetdir"),
                    )
                    return parsed_image

    def make_slide_layout(self, image, slide, allowbigimage=True):
        if image:
            ratio = image["width"] / image["height"]
            img_base_width = PptxInches(image["width"])
            img_base_height = PptxInches(image["height"])
            base_left = PptxInches(self.c["textbox"]["left"])
            base_top = PptxInches(self.c["textbox"]["top"])
            base_width = PptxInches(self.c["textbox"]["width"])
            base_height = PptxInches(self.c["textbox"]["height"])
            if self.c.get("disable_autolayout"):
                slide.shapes.add_picture(
                    image["imgfile"],
                    left=base_left,
                    top=base_top,
                    width=img_base_width,
                    height=img_base_height,
                )
                return self.get_textbox(slide), 1
            big_mode = (
                image["big"] and not self.c.get("text_is_duplicated") and allowbigimage
            )
            if ratio < 1:  # vertical image
                max_width = base_width // 3
                if big_mode:
                    max_width *= 2
                if img_base_width > max_width or big_mode:
                    img_width = max_width
                    img_height = int(img_base_height * (max_width / img_base_width))
                else:
                    img_width = img_base_width
                    img_height = img_base_height
                left = base_left + img_width
                top = base_top
                width = base_width - img_width
                height = base_height
                img_left = base_left
                img_top = int(base_top + 0.5 * (base_height - img_height))
            else:  # horizontal/square image
                max_height = base_height // 3
                if big_mode:
                    max_height *= 2
                if img_base_height > max_height or big_mode:
                    img_height = max_height
                    img_width = int(img_base_width * (max_height / img_base_height))
                else:
                    img_width = img_base_width
                    img_height = img_base_height
                left = base_left
                top = base_top + img_height
                width = base_width
                height = base_height - img_height
                img_top = base_top
                img_left = int(base_left + 0.5 * (base_width - img_width))
            slide.shapes.add_picture(
                image["imgfile"],
                left=img_left,
                top=img_top,
                width=img_width,
                height=img_height,
            )
            textbox = slide.shapes.add_textbox(left, top, width, height)
            return textbox, (width * height) / (base_width * base_height)
        else:
            return self.get_textbox(slide), 1

    def add_slide_with_image(self, image, number=None):
        slide = self.prs.slides.add_slide(self.BLANK_SLIDE)
        if number:
            self.set_question_number(slide, number)
        img_width = PptxInches(image["width"])
        img_height = PptxInches(image["height"])
        base_left = PptxInches(self.c["textbox"]["left"])
        base_top = PptxInches(self.c["textbox"]["top"])
        base_width = PptxInches(self.c["textbox"]["width"])
        base_height = PptxInches(self.c["textbox"]["height"])
        if image["big"] or img_width > base_width:
            img_width, img_height = (
                base_width,
                int(img_height * (base_width / img_width)),
            )
        if img_height > base_height:
            img_width, img_height = (
                int(img_width * (base_height / img_height)),
                base_height,
            )
        img_left = int(base_left + 0.5 * (base_width - img_width))
        img_top = int(base_top + 0.5 * (base_height - img_height))
        slide.shapes.add_picture(
            image["imgfile"],
            left=img_left,
            top=img_top,
            width=img_width,
            height=img_height,
        )

    def put_question_on_slide(self, image, slide, q, allowbigimage=True):
        textbox, coeff = self.make_slide_layout(
            image, slide, allowbigimage=allowbigimage
        )
        tf = textbox.text_frame
        self.apply_vertical_alignment_if_needed(tf)
        tf.word_wrap = True
        self.set_question_number(slide, self.number)
        question_text = self.pptx_process_text(q["question"], image=image)
        if self.c.get("force_text_size_question"):
            p = self.init_paragraph(tf, size=self.c["force_text_size_question"])
        else:
            p = self.init_paragraph(tf, text=question_text, coeff=coeff)
        self.pptx_format(question_text, p, tf, slide)

    def recursive_join(self, s):
        if isinstance(s, str):
            return s
        if isinstance(s, list):
            return "\n".join(self.recursive_join(x) for x in s)

    def add_slide_with_handout(self, handout, number=None):
        slide = self.prs.slides.add_slide(self.BLANK_SLIDE)
        textbox = self.get_textbox(slide)
        tf = textbox.text_frame
        self.apply_vertical_alignment_if_needed(tf)
        tf.word_wrap = True
        if number is not None:
            self.set_question_number(slide, number)
        p = self.init_paragraph(tf, text=handout)
        self.pptx_format(
            self.pptx_process_text(handout, do_not_remove_accents=True), p, tf, slide
        )

    def process_question_text(self, q):
        image = self._get_image_from_4s(q["question"])
        handout = self._get_handout_from_4s(q["question"])
        add_handout_on_separate_slide = self.c.get("add_handout_on_separate_slide")
        add_handout_on_separate_slide = (
            add_handout_on_separate_slide is None or add_handout_on_separate_slide
        )
        if image and add_handout_on_separate_slide:
            self.add_slide_with_image(image, number=self.number)
        elif handout and add_handout_on_separate_slide:
            self.add_slide_with_handout(handout, number=self.number)
        slide = self.prs.slides.add_slide(self.BLANK_SLIDE)
        text_is_duplicated = bool(self.c.get("text_is_duplicated"))
        self.put_question_on_slide(
            image, slide, q, allowbigimage=not text_is_duplicated
        )
        if image and image["big"] and text_is_duplicated:
            self.add_slide_with_image(image, number=self.number)

    def add_answer_slide(self, q):
        slide = self.prs.slides.add_slide(self.BLANK_SLIDE)
        if self.c.get("override_answer_caption"):
            self.set_question_number(slide, self.c["override_answer_caption"])
        else:
            self.set_question_number(slide, self.number)
        fields = ["answer"]
        if q.get("zachet") and self.c.get("add_zachet"):
            fields.append("zachet")
        if q.get("nezachet") and self.c.get("add_zachet"):
            fields.append("nezachet")
        if self.c["add_comment"] and "comment" in q:
            fields.append("comment")
        if self.c.get("add_source") and "source" in q:
            fields.append("source")
        textbox = None
        coeff = 1
        for field in fields:
            image = self._get_image_from_4s(q[field])
            if image:
                textbox, coeff = self.make_slide_layout(image, slide)
                break
        if not textbox:
            textbox = self.get_textbox(slide)
        tf = textbox.text_frame
        self.apply_vertical_alignment_if_needed(tf)
        tf.word_wrap = True

        text_for_size = self.recursive_join(
            self.pptx_process_text(q["answer"], strip_brackets=False)
        )
        if q.get("zachet") and self.c.get("add_zachet"):
            text_for_size += "\n" + self.recursive_join(
                self.pptx_process_text(q["zachet"], strip_brackets=False)
            )
        if q.get("nezachet") and self.c.get("add_zachet"):
            text_for_size += "\n" + self.recursive_join(
                self.pptx_process_text(q["nezachet"], strip_brackets=False)
            )
        if q.get("comment") and self.c.get("add_comment"):
            text_for_size += "\n" + self.recursive_join(
                self.pptx_process_text(q["comment"])
            )
        if q.get("source") and self.c.get("add_source"):
            text_for_size += "\n" + self.recursive_join(
                self.pptx_process_text(q["source"])
            )
        if q.get("author") and self.c.get("add_author"):
            text_for_size += "\n" + self.recursive_join(
                self.pptx_process_text(q["author"])
            )
        if self.c.get("force_text_size_answer"):
            p = self.init_paragraph(tf, size=self.c["force_text_size_answer"])
        else:
            p = self.init_paragraph(tf, text=text_for_size, coeff=coeff)
        r = self.add_run(p, f"{self.get_label(q, 'answer')}: ")
        r.font.bold = True
        self.pptx_format(
            self.pptx_process_text(q["answer"], strip_brackets=False), p, tf, slide
        )
        if q.get("zachet") and self.c.get("add_zachet"):
            zachet_text = self.pptx_process_text(q["zachet"], strip_brackets=False)
            r = self.add_run(p, f"\n{self.get_label(q, 'zachet')}: ")
            r.font.bold = True
            self.pptx_format(zachet_text, p, tf, slide)
        if q.get("nezachet") and self.c.get("add_zachet"):
            nezachet_text = self.pptx_process_text(q["nezachet"], strip_brackets=False)
            r = self.add_run(p, f"\n{self.get_label(q, 'nezachet')}: ")
            r.font.bold = True
            self.pptx_format(nezachet_text, p, tf, slide)
        if self.c["add_comment"] and "comment" in q:
            comment_text = self.pptx_process_text(q["comment"])
            r = self.add_run(p, f"\n{self.get_label(q, 'comment')}: ")
            r.font.bold = True
            self.pptx_format(comment_text, p, tf, slide)
        if self.c.get("add_source") and "source" in q:
            source_text = self.pptx_process_text(q["source"])
            r = self.add_run(p, f"\n{self.get_label(q, 'source')}: ")
            r.font.bold = True
            self.pptx_format(source_text, p, tf, slide)
        if self.c.get("add_author") and "author" in q:
            author_text = self.pptx_process_text(q["author"])
            r = self.add_run(p, f"\n{self.get_label(q, 'author')}: ")
            r.font.bold = True
            self.pptx_format(author_text, p, tf, slide)

    def process_question(self, q):
        if "number" not in q:
            self.qcount += 1
        if "setcounter" in q:
            self.qcount = int(q["setcounter"])
        self.number = str(self.qcount if "number" not in q else q["number"])

        if isinstance(q["question"], list):
            for i in range(len(q["question"][1])):
                qn = copy.deepcopy(q)
                qn["question"][1] = q["question"][1][: i + 1]
                self.process_question_text(qn)
        else:
            self.process_question_text(q)

        if self.c["add_plug"]:
            slide = self.prs.slides.add_slide(self.BLANK_SLIDE)
            self.set_question_number(slide, self.number)
        self.add_answer_slide(q)

    def determine_size(self, text, coeff=1):
        text = self.recursive_join(text)
        len_for_size = round((len(text) + 50 * text.count("\n")) / coeff)
        for element in self.c["text_size_grid"]["elements"]:
            if len_for_size <= element["length"]:
                return element["size"]
        return self.c["text_size_grid"]["smallest"]

    def init_paragraph(self, text_frame, text=None, coeff=1, size=None, color=None):
        p = text_frame.paragraphs[0]
        p.font.name = self.c["font"]["name"]
        if size:
            _size = size
        else:
            _size = self.c["text_size_grid"]["default"]
            if text:
                _size = self.determine_size(text, coeff=coeff)
        p.font.size = PptxPt(_size)
        return p

    def export(self, outfilename):
        self.outfilename = outfilename
        wd = os.getcwd()
        os.chdir(os.path.dirname(self.config_path))
        template = os.path.abspath(self.c["template_path"])
        os.chdir(wd)
        self.prs = Presentation(template)
        self.TITLE_SLIDE = self.prs.slide_layouts[0]
        self.BLANK_SLIDE = self.prs.slide_layouts[6]
        buffer = []
        for element in self.structure:
            if element[0] != "Question":
                buffer.append(element)
                continue
            if element[0] == "Question":
                if buffer:
                    self.process_buffer(buffer)
                    buffer = []
                self.process_question(element[1])
        self.prs.save(outfilename)
        self.logger.info("Output: {}".format(outfilename))
