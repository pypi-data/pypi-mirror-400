"""
Author: HECE - University of Liege, Pierre Archambeau
Date: 2024

Copyright (c) 2024 University of Liege. All rights reserved.

This script and its content are protected by copyright law. Unauthorized
copying or distribution of this file, via any medium, is strictly prohibited.
"""
try:
    from osgeo import gdal
except ImportError as e:
    print(e)
    raise ImportError("I can't find the 'gdal' package. You should get it from https://www.lfd.uci.edu/~gohlke/pythonlibs/")

try:
    import numpy as np
    from wx import dataview, TreeCtrl
    import wx
    import wx.propgrid as pg
    # from wxasync import AsyncBind
    from wx.core import VERTICAL, BoxSizer, Height, ListCtrl, StaticText, TextCtrl, Width
    from wx.glcanvas import GLCanvas, GLContext
    from wx.dataview import TreeListCtrl
    import wx.lib.ogl as ogl
    from PIL import Image, ImageOps
    from PIL.PngImagePlugin import PngInfo
    import io
    import json
    import glob
    import traceback
    from datetime import datetime
    from sklearn import linear_model, datasets

except ImportError as e:
    print(e)
    raise ImportError("Error importing wxPython, numpy, PIL, json, glob, traceback, sklearn. Please check your installation.")

try:
    from time import sleep
    from datetime import timedelta
    from multiprocessing import Pool
    from pathlib import Path
except ImportError as e:
    print(e)
    raise ImportError("Error importing time, datetime, multiprocessing, pathlib. Please check your installation.")

try:
    from OpenGL.GL import *
    from OpenGL.GLUT import *
except ImportError as e:
    msg=_('Error importing OpenGL library')
    msg+=_('   Python version : ' + sys.version)
    msg+=_('   Please check your version of opengl32.dll -- conflict may exist between different files present on your desktop')
    raise Exception(msg)

try:
    import matplotlib.pyplot as plt
    from matplotlib.widgets import Button as mplButton
    from matplotlib.ticker import FormatStrFormatter
    from os import scandir, listdir
    from os.path import exists, join, normpath
    from pptx import Presentation
    import threading
    from enum import Enum
    from typing import Literal, Union
    import logging
except ImportError as e:
    print(e)
    raise ImportError("Error importing matplotlib, os, threading, enum, typing, logging. Please check your installation.")

try:
    from .wolf_texture import genericImagetexture,imagetexture,Text_Image_Texture
    from .xyz_file import xyz_scandir, XYZFile
    from .mesh2d import wolf2dprev
    from .PyPalette import wolfpalette
    from .wolfresults_2D import Wolfresults_2D, views_2D, Extractable_results
    from .PyTranslate import _
    from .PyVertex import cloud_vertices, getIfromRGB, getRGBfromI
    from .RatingCurve import SPWMIGaugingStations, SPWDCENNGaugingStations
    from .wolf_array import WOLF_ARRAY_MB, SelectionData, WolfArray, WolfArrayMB, CropDialog, header_wolf, WolfArrayMNAP, WOLF_ARRAY_FULL_SINGLE, WOLF_ARRAY_FULL_INTEGER8, WOLF_ARRAY_FULL_INTEGER16, WOLF_ARRAY_FULL_DOUBLE, WOLF_ARRAY_FULL_INTEGER
    from .PyParams import Wolf_Param, key_Param, Type_Param
    from .mesh2d.bc_manager import BcManager
    from .PyVertexvectors import *
    from .Results2DGPU import wolfres2DGPU
    from .PyCrosssections import crosssections, profile, Interpolator, Interpolators
    from .GraphNotebook import PlotNotebook
    from .lazviewer.laz_viewer import myviewer, read_laz, clip_data_xyz, xyz_laz_grids, choices_laz_colormap, Classification_LAZ, Wolf_LAZ_Data, viewer as viewerlaz
    from . import Lidar2002
    from .picc import Picc_data, Cadaster_data
    from .wolf_zi_db import ZI_Databse_Elt, PlansTerrier, Ouvrages, Particularites, Enquetes, Profils
    from .math_parser.calculator import Calculator
    from .wintab.wintab import Wintab
    from .images_tiles import ImagesTiles
    from .PyWMS import Alaro_Navigator, get_Alaro_legend
    from .PyPictures import PictureCollection
    from .irm_alaro import IRM_Alaro, GribFiles, _convert_col2date_str

except ImportError as e:
    print(e)
    raise ImportError("Error importing wolf_texture, xyz_file, mesh2d, PyPalette, wolfresults_2D, PyTranslate, PyVertex, RatingCurve, wolf_array, PyParams, mesh2d.bc_manager, PyVertexvectors, Results2DGPU, PyCrosssections, GraphNotebook, lazviewer, picc, wolf_zi_db, math_parser.calculator, wintab. Please check your installation.")

try:
    from .dike import DikeWolf, InjectorWolfDike as InjectorDike
    WOLFPYDIKE_AVAILABLE = True
except:
    logging.warning(_("Missing package. Install wolfpydike module via pip."))
    WOLFPYDIKE_AVAILABLE = False

try:
    from .hydrometry.kiwis_wolfgui import hydrometry_wolfgui
except ImportError as e:
    print(e)
    raise ImportError("Error importing hydrometry.kiwis_wolfgui. Please check your installation.")

try:
    from .pyshields import get_d_cr
    from .pyviews import WolfViews
    from .PyConfig import handle_configuration_dialog, WolfConfiguration, ConfigurationKeys
    from .GraphProfile import ProfileNotebook
    from .pybridges import Bridges, Bridge, Weirs, Weir
    from .tools_mpl import *
    from .wolf_tiles import Tiles
    from .lagrangian.particle_system_ui import Particle_system_to_draw as Particle_system
    from .opengl.py3d import Wolf_Viewer3D
    from .pyGui1D import GuiNotebook1D
    from .matplotlib_fig import Matplotlib_Figure as MplFig, PRESET_LAYOUTS
except ImportError as e:
    print(e)
    raise ImportError("Error importing pyshields, pyviews, PyConfig, GraphProfile, pybridges, tools_mpl, wolf_tiles, lagrangian.particle_system_ui, opengl.py3d, pyGui1D. Please check your installation.")

try:
    from .apps.curvedigitizer import Digitizer
except ImportError as e:
    print(e)
    raise ImportError("Error importing apps.curvedigitizer. Please check your installation.")

try:
    from .drowning_victims.drowning_class import Drowning_victim_Viewer
except ImportError as e:
    print(e)
    raise ImportError("Error importing Drowning_victims.Class. Please check your installation.")

ID_SELECTCS = 1000
ID_SORTALONG = 1001
ID_LOCMINMAX = 1002
ID_PLOTCS = 1003   #Manageactions ID for profile plots

LIST_1TO9 = [wx.WXK_NUMPAD1, wx.WXK_NUMPAD2, wx.WXK_NUMPAD3, wx.WXK_NUMPAD4, wx.WXK_NUMPAD5, wx.WXK_NUMPAD6, wx.WXK_NUMPAD7, wx.WXK_NUMPAD8, wx.WXK_NUMPAD9 ] + [ord(str(cur)) for cur in range(1,10)]

PROJECT_ACTION = 'action'
PROJECT_CS     = 'cross_sections'
PROJECT_VECTOR = 'vector'
PROJECT_ARRAY  = 'array'
PROJECT_TILES  = 'tiles'
PROJECT_LAZ    = 'laz_grid'
PROJECT_CLOUD  = 'cloud'
PROJECT_WOLF2D = 'wolf2d'
PROJECT_GPU2D  = 'gpu2d'
PROJECT_PALETTE = 'palette'
PROJECT_PALETTE_ARRAY = 'palette-array'
PROJECT_LINK_CS = 'cross_sections_link'
PROJECT_LINK_VEC_ARRAY = 'vector_array_link'

PROJECT_GROUP_KEYS = {PROJECT_ACTION : {'which': 'compare_arrays'},
                      PROJECT_CS: {'id - file': 'id to use - full or relative path to CS file',
                                          'format': '(mandatory) 2000, 2022, vecz, sxy',
                                          'dirlaz': 'Path to LAZ data (prepro Numpy)'},
                      PROJECT_VECTOR: {'id - file': 'id to use - full or relative path to vector file (.vec, .vecz, .shp)'},
                      PROJECT_ARRAY: {'id - file': 'id to use - full or relative path to array file (.bin, .tif, .npy, .npz)'},
                      PROJECT_TILES : {'id': '(mandatory) id to use',
                                  'tiles_file': '(mandatory) Path to tiles file',
                                  'data_dir': '(mandatory) Path to data directory',
                                  'comp_dir': 'Path to comparison directory'},
                      PROJECT_LAZ: {'data_dir': '(mandatory) Path to data directory (prepro Numpy)',
                                    'classification': 'Color classification for LAZ data - default SPW-Geofit 2023',},
                      PROJECT_CLOUD: {'id - file': 'id to use - full or relative path to cloud file (.xyz, .txt)'},
                      PROJECT_WOLF2D: {'id - dir': 'id to use - full or relative path to wolf2d simulation directory'},
                      PROJECT_GPU2D: {'id - dir': 'id to use - full or relative path to gpu2d simulation directory'},
                      PROJECT_PALETTE : {'id - file': 'id to use - full or relative path to palette file (.pal)'},
                      PROJECT_PALETTE_ARRAY : {'idarray - idpal': 'id of array - id of palette to link'},
                      PROJECT_LINK_CS : {'linkzones' : '(mandatory) id of vector to link to cross sections',
                                                'sortzone' : '(mandatory) id of the zone to use for sorting',
                                                'sortname' : '(mandatory) id of the polyline to use for sorting',
                                                'downfirst' : 'is the first vertex downstream or upstream? (1 is True, 0 is False - default is False)'},
                      PROJECT_LINK_VEC_ARRAY : {'id - id vector': 'id of array/wolf2d/gpu2d - id of vector to link (only 1 vector in 1 zone)'},
                    }

class MplFigViewer(MplFig):

    def __init__(self, layout = None, idx:str='',
                 mapviewer:"WolfMapViewer" = None,
                 caption:str = '', size:tuple = (800, 600),
                 style:int = wx.DEFAULT_FRAME_STYLE ^ wx.RESIZE_BORDER):

        super().__init__(layout)
        self._mapviewer = mapviewer
        self._idx = idx

        self.SetTitle(caption)

        self.SetSize(size)

        dpi = self.fig.get_dpi()
        size_x = (size[0]-16)/dpi
        size_y = (size[1]-240)/dpi

        self.fig.set_size_inches(size_x, size_y)


        self.SetWindowStyle(style)

        self.Bind(wx.EVT_CLOSE, self.OnClose)

    def OnClose(self, event):
        """ Close the window """

        if self.mapviewer is not None:
            self.mapviewer.destroy_fig_by_id(self.idx)
        else:
            self.Destroy()

    @property
    def mapviewer(self):
        return self._mapviewer

    @property
    def idx(self):
        return self._idx

class Memory_View():
    """ Memory view """

    def __init__(self, screen_width, screen_height, xmin, xmax, ymin, ymax):
        """ Constructor """

        self.screen_width = screen_width
        self.screen_height = screen_height
        self.xmin = xmin
        self.xmax = xmax
        self.ymin = ymin
        self.ymax = ymax

    @property
    def width(self):
        """ Width of the view """
        return self.xmax - self.xmin

    @property
    def height(self):
        """ Height of the view """
        return self.ymax - self.ymin

    def __str__(self):
        """ String representation of the view """

        return f"Memory view : {self.screen_width}x{self.screen_height} ({self.xmin},{self.ymin})-({self.xmax},{self.ymax})"

    def serialize(self):
        """ Serialize the view """

        return {
            "screen_width": self.screen_width,
            "screen_height": self.screen_height,
            "xmin": self.xmin,
            "xmax": self.xmax,
            "ymin": self.ymin,
            "ymax": self.ymax
        }

    @staticmethod
    def deserialize(data:dict):
        """ Deserialize the view """

        return Memory_View(data["screen_width"], data["screen_height"], data["xmin"], data["xmax"], data["ymin"], data["ymax"])

class Memory_View_encoder(json.JSONEncoder):
    """ Memory view encoder """

    def default(self, o):
        """ Default method """

        if isinstance(o, Memory_View):
            return o.serialize()
        else:
            return super().default(o)

class Memory_View_decoder(json.JSONDecoder):
    """ Memory view decoder """

    def __init__(self, *args, **kwargs):
        """ Constructor """

        super().__init__(object_hook=self.object_hook, *args, **kwargs)

    def object_hook(self, obj):
        """ Decode the object """

        if "screen_width" in obj and "screen_height" in obj and "xmin" in obj and "xmax" in obj and "ymin" in obj and "ymax" in obj:
            return Memory_View.deserialize(obj)
        else:
            return obj

class Memory_Views():
    """ Memory views """

    def __init__(self):

        self.views:dict[str,Memory_View] = {}

    def __len__(self):
        """ Number of views """

        return len(self.views)

    def add_view(self, name:str, screen_width:int, screen_height:int, xmin:float, xmax:float, ymin:float, ymax:float):
        """ Add a new view to the memory views """
        self.views[name] = Memory_View(screen_width, screen_height, xmin, xmax, ymin, ymax)

    def remove_view(self, name:str):
        """ Remove a view from the memory views """

        if name in self.views:
            self.views.pop(name)

    def reset(self):
        """ Reset the memory views """

        self.views = {}

    def __getitem__(self, name:str) -> Memory_View:
        """ Get a view from the memory views """

        if name in self.views:
            return self.views[name]
        else:
            return None

    def zoom_on(self, name:str, mapviewer:"WolfMapViewer"):
        """ Zoom on a view """

        if name not in self.views:
            return

        view = self.views[name]

        mapviewer.zoom_on(width= view.width, height=view.height, xll=view.xmin, yll=view.ymin, canvas_height=view.screen_height)

    def save(self, filename:str):
        """ Save the memory views """

        with open(filename, 'w') as f:
            json.dump(self.views, f,
                      cls=Memory_View_encoder,
                      indent=4)

    def load(self, filename:str):
        """ Load the memory views """

        with open(filename, 'r') as f:
            self.views = json.load(f, cls=Memory_View_decoder)

        # self.views = {k:Memory_View.deserialize(v) for k,v in tmp_views.items()}

class Memory_Views_GUI(wx.Frame):
    """ Memory views GUI """

    def __init__(self, parent, title, memory_views:Memory_Views, mapviewer:"WolfMapViewer"):
        """ Constructor """

        super(Memory_Views_GUI, self).__init__(parent, title=title, size=(200, 400), style = wx.DEFAULT_FRAME_STYLE & ~ (wx.RESIZE_BORDER | wx.MAXIMIZE_BOX | wx.MINIMIZE_BOX))
        self.mapviewer = mapviewer
        self._memory_views = memory_views


        panel = wx.Panel(self)
        panel.SetBackgroundColour(wx.Colour(255, 255, 255))

        sizer = wx.BoxSizer(wx.VERTICAL)

        self._views = wx.ListBox(panel, choices= list(memory_views.views.keys()), style=wx.LB_SINGLE)

        self._cmdZoom = wx.Button(panel, wx.ID_ANY, _('Zoom on'))
        self._cmdZoom.SetToolTip(_('Zoom on the selected view'))

        self._cmdAdd = wx.Button(panel, wx.ID_ADD, _('+'))
        self._cmdAdd.SetToolTip(_('Add a view based on the current zoom and shape of the canvas'))

        self._cmdDelete = wx.Button(panel, wx.ID_DELETE, _('-'))
        self._cmdDelete.SetToolTip(_('Delete the selected view'))

        self._cmdReset = wx.Button(panel, wx.ID_RESET, _('Reset'))
        self._cmdReset.SetToolTip(_('Reset the views'))

        sizer.Add(self._views, 5, wx.EXPAND | wx.ALL, 2)

        sizer_but = wx.BoxSizer(wx.HORIZONTAL)
        sizer_but.Add(self._cmdAdd, 1, wx.EXPAND)
        sizer_but.Add(self._cmdDelete, 1, wx.EXPAND)

        sizer.Add(self._cmdZoom, 1, wx.EXPAND, 2)
        sizer.Add(sizer_but, 1,  wx.EXPAND, 2)
        sizer.Add(self._cmdReset, 1, wx.EXPAND , 2)

        sizer_manual = wx.BoxSizer(wx.VERTICAL)

        sizer_xmin = wx.BoxSizer(wx.HORIZONTAL)
        self._label_xmin = wx.StaticText(panel, label=_('X min'))
        self._xmin = wx.TextCtrl(panel, value=str(mapviewer.xmin), style=wx.ALIGN_CENTER_HORIZONTAL)
        sizer_xmin.Add(self._label_xmin, 1, wx.ALL, 2)
        sizer_xmin.Add(self._xmin, 1, wx.ALL, 2)

        sizer_xmax = wx.BoxSizer(wx.HORIZONTAL)
        self._label_xmax = wx.StaticText(panel, label=_('X max'))
        self._xmax = wx.TextCtrl(panel, value=str(mapviewer.xmax), style=wx.ALIGN_CENTER_HORIZONTAL)
        sizer_xmax.Add(self._label_xmax, 1, wx.ALL, 2)
        sizer_xmax.Add(self._xmax, 1, wx.ALL, 2)

        sizer_ymin = wx.BoxSizer(wx.HORIZONTAL)
        self._label_ymin = wx.StaticText(panel, label=_('Y min'))
        self._ymin = wx.TextCtrl(panel, value=str(mapviewer.ymin), style=wx.ALIGN_CENTER_HORIZONTAL)
        sizer_ymin.Add(self._label_ymin, 1, wx.ALL, 2)
        sizer_ymin.Add(self._ymin, 1, wx.ALL, 2)

        sizer_ymax = wx.BoxSizer(wx.HORIZONTAL)
        self._label_ymax = wx.StaticText(panel, label=_('Y max'))
        self._ymax = wx.TextCtrl(panel, value=str(mapviewer.ymax), style=wx.ALIGN_CENTER_HORIZONTAL)
        sizer_ymax.Add(self._label_ymax, 1, wx.ALL, 2)
        sizer_ymax.Add(self._ymax, 1, wx.ALL, 2)

        sizer_canvas_height = wx.BoxSizer(wx.HORIZONTAL)
        self._label_canvas_height = wx.StaticText(panel, label=_('Canvas height'))
        self._canvas_height = wx.TextCtrl(panel, value=str(mapviewer.canvasheight), style=wx.ALIGN_CENTER_HORIZONTAL)

        self._label_canvas_height.SetToolTip(_('Height of the canvas in pixels'))
        self._canvas_height.SetToolTip(_('Height of the canvas in pixels'))

        sizer_canvas_height.Add(self._label_canvas_height, 1, wx.ALL, 2)
        sizer_canvas_height.Add(self._canvas_height, 1, wx.ALL, 2)

        sizer_manual.Add(sizer_xmin, 1, wx.ALL, 2)
        sizer_manual.Add(sizer_xmax, 1, wx.ALL, 2)
        sizer_manual.Add(sizer_ymin, 1, wx.ALL, 2)
        sizer_manual.Add(sizer_ymax, 1, wx.ALL, 2)
        sizer_manual.Add(sizer_canvas_height, 1, wx.ALL, 2)

        sizer_but2 = wx.BoxSizer(wx.HORIZONTAL)
        self._cmdGet = wx.Button(panel, wx.ID_ANY, _('Get'))
        self._cmdGet.SetToolTip(_('Get the current bounds of the canvas and fill the fields'))

        self._cmdApply = wx.Button(panel, wx.ID_APPLY, _('Apply'))
        self._cmdApply.SetToolTip(_('Apply the values to the selected view'))

        self._cmdApply.Bind(wx.EVT_BUTTON, self.OnApply)
        self._cmdGet.Bind(wx.EVT_BUTTON, self.OnGet)

        sizer_but2.Add(self._cmdGet, 1, wx.ALL, 2)
        sizer_but2.Add(self._cmdApply, 1, wx.ALL, 2)

        sizer_manual.Add(sizer_but2, 1, wx.EXPAND, 2)
        sizer.Add(sizer_manual, 1, wx.ALL, 2)

        sizer_save_load = wx.BoxSizer(wx.HORIZONTAL)
        self._cmdSave = wx.Button(panel, wx.ID_SAVE, _('Save'))
        self._cmdSave.SetToolTip(_('Save the memory views to a json file'))

        self._cmdLoad = wx.Button(panel, wx.ID_OPEN, _('Load'))
        self._cmdLoad.SetToolTip(_('Load the memory views from a json file'))

        sizer_save_load.Add(self._cmdSave, 1, wx.ALL, 2)
        sizer_save_load.Add(self._cmdLoad, 1, wx.ALL, 2)

        sizer.Add(sizer_save_load, 1, wx.EXPAND, 2)

        self._views.Bind(wx.EVT_LISTBOX, self.OnSelectView)
        self.Bind(wx.EVT_BUTTON, self.OnAdd, self._cmdAdd)
        self.Bind(wx.EVT_BUTTON, self.OnDelete, self._cmdDelete)
        self.Bind(wx.EVT_BUTTON, self.OnReset, self._cmdReset)
        self.Bind(wx.EVT_CLOSE, self.OnClose)
        self.Bind(wx.EVT_BUTTON, self.OnSave, self._cmdSave)
        self.Bind(wx.EVT_BUTTON, self.OnLoad, self._cmdLoad)
        self.Bind(wx.EVT_BUTTON, self.OnZoom, self._cmdZoom)

        panel.SetSizer(sizer)

        self.CenterOnScreen()

        icon = wx.Icon()
        icon_path = Path(__file__).parent / "apps/wolf.ico"
        icon.CopyFromBitmap(wx.Bitmap(str(icon_path), wx.BITMAP_TYPE_ANY))
        self.SetIcon(icon)

        self.Show()

    def OnSave(self, event):
        """ Save the memory views """

        with wx.FileDialog(self, _('Save the memory views'), wildcard="JSON files (*.json)|*.json",
                            style=wx.FD_SAVE | wx.FD_OVERWRITE_PROMPT) as fileDialog:
            if fileDialog.ShowModal() == wx.ID_CANCEL:
                return

            self._memory_views.save(fileDialog.GetPath())

    def OnLoad(self, event):
        """ Load the memory views """

        with wx.FileDialog(self, _('Load the memory views'), wildcard="JSON files (*.json)|*.json",
                            style=wx.FD_OPEN | wx.FD_FILE_MUST_EXIST) as fileDialog:
            if fileDialog.ShowModal() == wx.ID_CANCEL:
                return

            self._memory_views.load(fileDialog.GetPath())

        if self._memory_views is None:
            logging.error(_('Error while loading the memory views'))
            return

        if self._memory_views.views is None:
            logging.error(_('Error while loading the memory views'))
            return

        self._views.Set(list(self._memory_views.views.keys()))

    def OnClose(self, event):
        """ Close the memory views GUI """

        self.mapviewer._memory_views_gui = None
        self.Destroy()

    def OnSelectView(self, event):
        """ Select a view """

        view = self._views.GetStringSelection()
        locview = self._memory_views[view]

        if locview is None:
            return

        self._xmax.SetValue(str(self._memory_views[view].xmax))
        self._xmin.SetValue(str(self._memory_views[view].xmin))
        self._ymax.SetValue(str(self._memory_views[view].ymax))
        self._ymin.SetValue(str(self._memory_views[view].ymin))
        self._canvas_height.SetValue(str(self._memory_views[view].screen_height))

    def OnZoom(self, event):
        """ Zoom on the current view """

        view = self._views.GetStringSelection()
        self._memory_views.zoom_on(view, self.mapviewer)

    def OnAdd(self, event):
        """ Add a view """

        name = "View " + str(len(self._memory_views.views) + 1)
        with wx.TextEntryDialog(self, _('Enter the name of the view'), _('Add a view'), name) as dlg:
            if dlg.ShowModal() == wx.ID_OK:
                name = dlg.GetValue()

        self._memory_views.add_view(name, self.mapviewer.canvaswidth, self.mapviewer.canvasheight, self.mapviewer.xmin, self.mapviewer.xmax, self.mapviewer.ymin, self.mapviewer.ymax)

        self._views.Set(list(self._memory_views.views.keys()))

    def OnDelete(self, event):
        """ Delete a view """

        view = self._views.GetStringSelection()
        self._memory_views.remove_view(view)

        self._views.Set(list(self._memory_views.views.keys()))

    def OnReset(self, event):
        """ Reset the views """

        self._memory_views.reset()

    def OnApply(self, event):
        """ Apply the changes """

        view = self._views.GetStringSelection()

        try:
            xmin = float(self._xmin.GetValue())
            xmax = float(self._xmax.GetValue())
            ymin = float(self._ymin.GetValue())
            ymax = float(self._ymax.GetValue())
            canvas_height = int(self._canvas_height.GetValue())

            self._memory_views.add_view(view, self.mapviewer.canvaswidth, canvas_height, xmin, xmax, ymin, ymax)
        except:
            logging.error(_('Error while applying the changes'))

    def OnGet(self, event):
        """ Get the values """

        self._xmin.SetValue(str(self.mapviewer.xmin))
        self._xmax.SetValue(str(self.mapviewer.xmax))
        self._ymin.SetValue(str(self.mapviewer.ymin))
        self._ymax.SetValue(str(self.mapviewer.ymax))
        self._canvas_height.SetValue(str(self.mapviewer.canvasheight))

class draw_type(Enum):
    # FIXME: change this to be more robust -> Done !
    # Be careful with the enum name, it must be the same than the one used to create the tree list elements, but in lower case
    # see : self.treelist.AppendItem in __init__
    ARRAYS = 'arrays'
    BRIDGES= 'bridges'
    WEIRS = 'weirs'
    VECTORS = 'vectors'
    CLOUD = 'clouds'
    TRIANGULATION = 'triangulations'
    PARTICLE_SYSTEM = 'particle systems'
    CROSS_SECTIONS = 'cross_sections'
    OTHER = 'others'
    VIEWS = 'views'
    RES2D = 'wolf2d'
    WMSBACK = 'wms-background'
    WMSFORE = 'wms-foreground'
    TILES = 'tiles'
    IMAGESTILES = 'imagestiles'
    LAZ = 'laz'
    DROWNING = 'drowning'
    DIKE = 'dike'
    PICTURECOLLECTION = 'picture_collection'
    INJECTOR = 'injector'

class Colors_1to9(wx.Frame):

    def __init__(self, parent):

        self._parent = parent
        self.colors1to9 = [(0, 0, 255, 255),
                            (0, 255, 0, 255),
                            (0, 128, 255, 255),
                            (255, 255, 0, 255),
                            (255, 165, 0, 255),
                            (128, 0, 128, 255),
                            (255, 192, 203, 255),
                            (165, 42, 42, 255),
                            (128, 128, 128, 255)]

        if self.file.exists():
            self.OnLoad(None)

    @property
    def directory(self):
        tmp = Path(__file__).parent / 'data'
        if not tmp.exists():
            tmp.mkdir()

        return tmp

    @property
    def file(self):
        return self.directory / 'colors1to9.json'

    def __getitem__(self, i):
        return self.colors1to9[i]

    def change_colors(self, e):

        super(Colors_1to9, self).__init__(self._parent, title=_('Colors 1 to 9'), size=(200, 400), style = wx.DEFAULT_FRAME_STYLE & ~ (wx.RESIZE_BORDER | wx.MAXIMIZE_BOX | wx.MINIMIZE_BOX | wx.CLOSE_BOX))
        panel = wx.Panel(self)
        panel.SetBackgroundColour(wx.Colour(255, 255, 255))

        sizer = wx.BoxSizer(wx.VERTICAL)

        self.pickers={}
        for i in range(9):
            horsizer = wx.BoxSizer(wx.HORIZONTAL)
            horsizer.Add(wx.StaticText(panel, label=_('Color ') + str(i + 1)), 0, wx.ALL, 2)
            color = self.colors1to9[i]
            color = wx.Colour(color[0], color[1], color[2], color[3])
            self.pickers[i] = wx.ColourPickerCtrl(panel, colour=color)
            horsizer.Add(self.pickers[i], 0, wx.ALL, 2)
            sizer.Add(horsizer, 0, wx.ALL, 2)

        cmdOK = wx.Button(panel, wx.ID_OK, _('OK'))
        cmdCancel = wx.Button(panel, wx.ID_CANCEL, _('Cancel'))
        cdmSetDefault = wx.Button(panel, wx.ID_APPLY, _('Default'))
        cmdSave = wx.Button(panel, wx.ID_SAVE, _('Save'))

        horsizer = wx.BoxSizer(wx.HORIZONTAL)
        horsizer2 = wx.BoxSizer(wx.HORIZONTAL)
        horsizer.Add(cmdOK, 0, wx.ALL, 2)
        horsizer.Add(cmdCancel, 0, wx.ALL, 2)
        horsizer2.Add(cdmSetDefault, 0, wx.ALL, 2)
        horsizer2.Add(cmdSave, 0, wx.ALL, 2)

        sizer.Add(horsizer, 0, wx.ALL, 2)
        sizer.Add(horsizer2, 0, wx.ALL, 2)

        panel.SetSizer(sizer)

        self.Bind(wx.EVT_BUTTON, self.OnOK, cmdOK)
        self.Bind(wx.EVT_BUTTON, self.OnCancel, cmdCancel)
        self.Bind(wx.EVT_BUTTON, self.OnSetDefault, cdmSetDefault)
        self.Bind(wx.EVT_BUTTON, self.OnSave, cmdSave)

        icon = wx.Icon()
        icon_path = Path(__file__).parent / "apps/wolf.ico"
        icon.CopyFromBitmap(wx.Bitmap(str(icon_path), wx.BITMAP_TYPE_ANY))
        self.SetIcon(icon)

        self.CenterOnScreen()

        self.Show()

    def Apply(self):

        for i in range(9):
            color = self.pickers[i].GetColour()
            self.colors1to9[i] = (color.Red(), color.Green(), color.Blue(), color.Alpha())

    def OnOK(self, event):

        self.Apply()
        self.Close()

    def OnCancel(self, event):
        self.Close()

    def OnSetDefault(self, event):
        self.colors1to9 = [(0, 0, 255, 255),
                            (0, 255, 0, 255),
                            (0, 128, 255, 255),
                            (255, 255, 0, 255),
                            (255, 165, 0, 255),
                            (128, 0, 128, 255),
                            (255, 192, 203, 255),
                            (165, 42, 42, 255),
                            (128, 128, 128, 255)]
        for i in range(9):
            color = self.colors1to9[i]
            color = wx.Colour(color[0], color[1], color[2], color[3])
            self.pickers[i].SetColour(color)

    def OnSave(self, event):
        self.Apply()
        with open(self.file, 'w') as f:
            json.dump(self.colors1to9, f)

    def OnLoad(self, event):
        with open(self.file, 'r') as f:
            self.colors1to9 = json.load(f)
class DragdropFileTarget(wx.FileDropTarget):
    def __init__(self, window:"WolfMapViewer"):
        wx.FileDropTarget.__init__(self)
        self.window = window

    def OnDropFiles(self, x, y, filenames):

        def test_if_array(filename):

            ext = Path(filename).suffix

            if ext.lower() in ['.bin', '.npy', '.hbin', '.qxin','.qybin', '.top',
                               '.kbin', '.epsbin', '.tif', '.tiff', '.frot', '.topini_fine']:
                return True
            else:
                return False

        def test_if_arrayMB(filename):

            ext = Path(filename).suffix

            if ext.lower() in ['.hbinb', '.qxbinb','.qybinb', '.kbinb',
                               '.epsbinb', '.topini', '.frotini']:
                return True
            else:
                return False

        def test_if_vector(filename):
            ext = Path(filename).suffix

            if ext.lower() in ['.vec', '.vecz', '.shp', '.dxf']:
                return True
            else:
                return False

        def test_if_cloud(filename):
            ext = Path(filename).suffix

            if ext.lower() in ['.xyz']:
                return True
            else:
                return False

        pgbar = wx.ProgressDialog(_('Loading files'), _('Loading files'), maximum=len(filenames), parent=self.window, style=wx.PD_APP_MODAL | wx.PD_AUTO_HIDE)

        for name in filenames:

            if Path(name).is_dir():
                for file in scandir(name):
                    if file.is_file():
                        self.OnDropFiles(x, y, [file.path])
                continue

            if test_if_array(name):
                ids = self.window.get_list_keys(draw_type.ARRAYS, checked_state=None)
                id = Path(name).stem
                while id in ids:
                    id = id + '_1'

                try:
                    h = header_wolf.read_header(name)

                    if h.nb_blocks>0:
                        newobj = WolfArrayMB(fname=name, mapviewer= self.window)
                    else:
                        newobj = WolfArray(fname=name, mapviewer= self.window)
                    self.window.add_object('array', newobj = newobj, id  = id)
                except:
                    logging.error(_('Error while loading array : ') + name)

            elif test_if_arrayMB(name):
                ids = self.window.get_list_keys(draw_type.ARRAYS, checked_state=None)
                id = Path(name).stem
                while id in ids:
                    id = id + '_1'

                try:
                    newobj = WolfArrayMB(fname=name, mapviewer= self.window)
                    self.window.add_object('array', newobj = newobj, id  = id)
                except:
                    logging.error(_('Error while loading array : ') + name)

            elif test_if_vector(name):
                ids = self.window.get_list_keys(draw_type.VECTORS, checked_state=None)
                id = Path(name).stem
                while id in ids:
                    id = id + '_1'

                try:
                    newobj = Zones(filename=name, parent=self.window, mapviewer=self.window)
                    self.window.add_object('vector', newobj = newobj, id  = id)
                except:
                    logging.error(_('Error while loading vector : ') + name)

            elif test_if_cloud(name):
                ids = self.window.get_list_keys(draw_type.CLOUD, checked_state=None)
                id = Path(name).stem
                while id in ids:
                    id = id + '_1'

                try:
                    newobj = cloud_vertices(fname=name, mapviewer=self.window)
                    self.window.add_object('cloud', newobj = newobj, id  = id)
                except:
                    logging.error(_('Error while loading cloud : ') + name)

            pgbar.Update(pgbar.GetValue() + 1)

        pgbar.Destroy()

        return True


class Sim_Explorer(wx.Frame):

    def __init__(self, parent, title, mapviewer:"WolfMapViewer", sim:Wolfresults_2D):

        super(Sim_Explorer, self).__init__(parent, title=title, size=(150, 250), style = wx.DEFAULT_FRAME_STYLE & ~ (wx.MAXIMIZE_BOX | wx.MINIMIZE_BOX))

        self._panel = wx.Panel(self)

        self.mapviewer = mapviewer
        self.active_res2d:Wolfresults_2D = sim

        main_sizer = wx.BoxSizer(wx.HORIZONTAL)

        left_bar = wx.BoxSizer(wx.VERTICAL)
        right_bar = wx.BoxSizer(wx.VERTICAL)

        self._all_times_steps = self.active_res2d.get_times_steps()

        # Right bar
        # ---------

        # Slider
        self._slider_steps = wx.Slider(self._panel, minValue=1, maxValue=sim.get_nbresults(), style=wx.SL_HORIZONTAL | wx.SL_AUTOTICKS | wx.SL_MIN_MAX_LABELS | wx.SL_LABELS)
        self._slider_steps.Bind(wx.EVT_SLIDER, self.OnSliderSteps)
        right_bar.Add(self._slider_steps, 1, wx.EXPAND | wx.ALL, 2)

        # Explore by index
        self._label_idx = wx.StaticText(self._panel, label=_('Index'))
        right_bar.Add(self._label_idx, 1, wx.EXPAND | wx.ALL, 2)

        self._step_idx = wx.ListBox(self._panel, choices=[str(i) for i in range(1, sim.get_nbresults()+1)], style=wx.LB_SINGLE)
        self._step_idx.Bind(wx.EVT_LISTBOX, self.OnSelectIdxStep)
        right_bar.Add(self._step_idx, 1, wx.EXPAND | wx.ALL, 5)

        # Explore by time
        self._label_time = wx.StaticText(self._panel, label=_('Time [s]'))
        right_bar.Add(self._label_time, 1, wx.EXPAND | wx.ALL, 2)

        _now = datetime.now()
        self._starting_date = datetime(year=_now.year, month=_now.month, day=_now.day, hour=0, minute=0, second=0)
        self._texttime = wx.TextCtrl(self._panel, value=self._starting_date.strftime('%Y-%m-%d %H:%M:%S'))
        right_bar.Add(self._texttime, 1, wx.EXPAND | wx.ALL, 5)
        self._texttime.Bind(wx.EVT_TEXT, self.OnTextTime)

        self._step_time = wx.ListBox(self._panel, choices=['{:.3f} - {}'.format(i, datetime.strftime(self._starting_date + timedelta(seconds=float(i)), '%Y-%m-%d %H:%M:%S')) for i in self._all_times_steps[0]], style=wx.LB_SINGLE)
        self._step_time.Bind(wx.EVT_LISTBOX, self.OnSelectNumStep)
        right_bar.Add(self._step_time, 1, wx.EXPAND | wx.ALL, 5)

        # Explore by time step
        self._label_steps = wx.StaticText(self._panel, label=_('Time step [-]'))
        right_bar.Add(self._label_steps, 1, wx.EXPAND | wx.ALL, 2)

        self._step_num = wx.ListBox(self._panel, choices=[str(i) for i in self._all_times_steps[1]], style=wx.LB_SINGLE)
        self._step_num.Bind(wx.EVT_LISTBOX, self.OnSelectCurTime)
        right_bar.Add(self._step_num, 1, wx.EXPAND | wx.ALL, 5)

        # Left bar
        # --------

        # Apply selected step
        self._cmd_apply = wx.Button(self._panel, wx.ID_APPLY, _('Apply'))
        self._cmd_apply.SetToolTip(_('Apply the selected parameters to the map'))
        self._cmd_apply.Bind(wx.EVT_BUTTON, self.OnApply)
        left_bar.Add(self._cmd_apply, 1, wx.EXPAND | wx.ALL, 5)

        # Update listbox from files on disk
        self._cmd_update = wx.Button(self._panel, wx.ID_REFRESH, _('Update'))
        self._cmd_update.SetToolTip(_('Update the list of available results based on the files on disk'))
        self._cmd_update.Bind(wx.EVT_BUTTON, self.OnUpdate)
        left_bar.Add(self._cmd_update, 1, wx.EXPAND | wx.ALL, 5)

        #Plot
        self._cmd_plot = wx.Button(self._panel, wx.ID_PREVIEW, _('Plot simulation informations'))
        self._cmd_plot.SetToolTip(_('Plot synthesis of the simulation (computation time, time step, clock time, mostly dry mesh...)'))
        self._cmd_plot.Bind(wx.EVT_BUTTON, self.OnPlot)
        left_bar.Add(self._cmd_plot, 1, wx.EXPAND | wx.ALL, 5)

        # Next step
        self._cmd_next = wx.Button(self._panel, wx.ID_FORWARD, _('Next'))
        self._cmd_next.SetToolTip(_('Go to the next step -- using the selected mode'))
        self._cmd_next.Bind(wx.EVT_BUTTON, self.OnNext)
        left_bar.Add(self._cmd_next, 1, wx.EXPAND | wx.ALL, 5)

        # Previous step
        self._cmd_prev = wx.Button(self._panel, wx.ID_BACKWARD, _('Previous'))
        self._cmd_prev.SetToolTip(_('Go to the previous step -- using the selected mode'))
        self._cmd_prev.Bind(wx.EVT_BUTTON, self.OnPrev)
        left_bar.Add(self._cmd_prev, 1, wx.EXPAND | wx.ALL, 5)

        # Check Mode movement
        self._mode = wx.ListBox(self._panel, choices=['by time [s]', 'by time [hour]', 'by index', 'by time step'], style=wx.LB_SINGLE)
        self._mode.SetToolTip(_('Select the mode to move through the simulation'))
        self._mode.SetSelection(2)
        self._interval = wx.TextCtrl(self._panel, value='1', style=wx.ALIGN_CENTER_HORIZONTAL)
        self._interval.SetToolTip(_('Interval for the mode selected -- unit depends on the mode'))
        self._interval.Bind(wx.EVT_TEXT, self.OnInterval)

        left_bar.Add(self._mode, 1, wx.EXPAND | wx.ALL, 5)
        left_bar.Add(self._interval, 0, wx.EXPAND | wx.ALL, 5)

        self.Bind(wx.EVT_CLOSE, self.OnClose)

        main_sizer.Add(left_bar, 1, wx.EXPAND | wx.ALL, 2)
        main_sizer.Add(right_bar, 1, wx.EXPAND | wx.ALL, 2)

        self._panel.SetSizer(main_sizer)
        self._panel.SetAutoLayout(True)

        self.MinSize = (450, 500)

        self.Fit()
        self.Show()

        self.SetIcon(wx.Icon(str(Path(__file__).parent / "apps/wolf.ico")))

        self._set_all(0)

    def OnPlot(self, event):
        """ Create a scatter plot of all steps.

        Major x_axis is time in seconds, Minor X-axis is time by date.

        Plots:
            - Computation time step (Dt)
            - Computation steps (N)
            - Clock time (s)
            - Mostly dry mesh (N)

        """

        main_x = self._all_times_steps[0]
        second_x = [self._starting_date + timedelta(seconds=i) for i in main_x]

        if isinstance(self.active_res2d, wolfres2DGPU):

            ax:list[Axes]
            fig, ax= plt.subplots(5, 1, figsize=(10, 8))

            ax[0].plot(main_x, self._all_times_steps[1], 'o-')
            ax[0].set_ylabel(_('Computation\ntime step (N)'), fontsize=8)
            ax[0].ticklabel_format(axis='y', style='sci', scilimits=(0,0))

            ax[0].grid(which='both')
            ax[0].set_xticks(main_x)
            ax[0].set_xticklabels([])

            secax:Axes = ax[0].secondary_xaxis('top')
            secax.set_xlabel(_('Real date\n[Y-M-D H:M:S]'), fontsize=8)
            secax.set_xticks(main_x)
            secax.set_xticklabels([datetime.strftime(i, '%Y-%m-%d %H:%M') for i in second_x], fontsize=8)
            secax.tick_params(axis='x', rotation=30)

            ax[1].plot(main_x, self.active_res2d.all_dt, 'o-')
            ax[1].set_ylabel(_(r'$\Delta t$ [s]'), fontsize=8)
            ax[1].grid(which='both')
            ax[1].set_xticks(main_x)
            ax[1].set_xticklabels([])
            ax[1].ticklabel_format(axis='y', style='sci', scilimits=(0,0))

            ctime = self.active_res2d.all_clock_time
            ax[2].plot(main_x, self.active_res2d.all_clock_time, 'o-')
            ax[2].set_ylabel(_('Clock time [s]'), fontsize=8)
            ax[2].grid(which='both')
            ax[2].set_xticks([])
            ax[2].set_xticks(main_x)
            ax[2].set_xticklabels([])
            ax[2].ticklabel_format(axis='y', style='sci', scilimits=(0,0))

            # Fit a line on the (main_x - clock time) plot
            # This will give a mean acceleration factor
            # The inverse of the slope of the line is the accelaration factor
            # The line is y = slope * x + intercept
            from scipy.stats import linregress
            slope, intercept, r_value, p_value, std_err = linregress(main_x, ctime)

            # Plot the info on the ax[2]
            msg = _('Acceleration factor:')
            ax[2].text(0.5, 0.5, f'{msg} {1/slope:.2f}', transform=ax[2].transAxes, fontsize=12, verticalalignment='top')

            ax[3].plot(main_x, self.active_res2d.all_wet_meshes, 'o-', color='blue')
            ax[3].plot(main_x, self.active_res2d.all_mostly_dry_mesh, 'o-', color='green')
            ax[3].set_ylabel(_('Wet and Mostly dry\nmeshes [N]'), fontsize=8)
            ax[3].grid(which='both')
            ax[3].set_xticks(main_x)
            ax[3].set_xlabel([])
            ax[3].ticklabel_format(axis='y', style='sci', scilimits=(0,0))

            ax[4].plot(main_x, [i/j*100 if j>0 else 0 for i, j in zip(self.active_res2d.all_mostly_dry_mesh, self.active_res2d.all_wet_meshes)], 'o-', color='red')
            ax[4].set_ylabel(_('Wet/Mostly dry\nmeshes [%]'), fontsize=8)
            ax[4].grid(which='both')
            ax[4].set_xticks(main_x)
            ax[4].set_xlabel(_('Simulated time [s]'), fontsize=8)
            ax[4].ticklabel_format(axis='y', style='sci', scilimits=(0,0))

            fig.suptitle('Simulation {}'.format(self.active_res2d.idx), fontsize=10)

            fig.tight_layout()
            fig.show()

    def OnInterval(self, event):
        """ Change the interval """

        try:
            interv = float(self._interval.GetValue())
            if interv <= 0:
                interv = 1.
                self._interval.SetValue('1')
        except:
            interv = 1
            self._interval.SetValue('1')

    def _find_next(self, idx:int):
        """ Find the next step based on the mode and interval """

        mode = int(self._mode.GetSelection())

        if mode == 0:
            # By time [s]
            next_time = self._all_times_steps[0][idx] + float(self._interval.GetValue())
            diff = [abs(next_time - i) for i in self._all_times_steps[0][idx:]]
            next_idx = diff.index(min(diff)) + idx

            return next_idx

        elif mode == 1:
            # By time [hour]
            next_time = self._all_times_steps[0][idx] + float(self._interval.GetValue())*3600
            diff = [abs(next_time - i) for i in self._all_times_steps[0][idx:]]
            next_idx = diff.index(min(diff)) + idx

            return next_idx

        elif mode == 2:
            # By index
            next_idx = min(idx + int(self._interval.GetValue()), len(self._all_times_steps[0])-1)

            return next_idx

        elif mode == 3:
            # By time step
            next_idx = self._all_times_steps[1].index(self._all_times_steps[1][idx] + int(self._interval.GetValue()))
            diff = [abs(next_idx - i) for i in self._all_times_steps[1][idx:]]
            next_idx = diff.index(min(diff)) + idx

            return next_idx

    def _find_prev(self, idx:int):
        """ Find the previous step based on the mode and interval """

        mode = int(self._mode.GetSelection())

        if mode == 0:
            # By time [s]
            prev_time = self._all_times_steps[0][idx] - float(self._interval.GetValue())
            diff = [abs(prev_time - i) for i in self._all_times_steps[0][:idx]]
            prev_idx = diff.index(min(diff))

            return prev_idx

        elif mode == 1:
            # By time [hour]
            prev_time = self._all_times_steps[0][idx] - float(self._interval.GetValue())*3600
            diff = [abs(prev_time - i) for i in self._all_times_steps[0][:idx]]
            prev_idx = diff.index(min(diff))

            return prev_idx

        elif mode == 2:
            # By index
            prev_idx = max(idx - int(self._interval.GetValue()), 0)

            return prev_idx

        elif mode == 3:
            # By time step
            prev_idx = self._all_times_steps[1].index(self._all_times_steps[1][idx] - int(self._interval.GetValue()))
            diff = [abs(prev_idx - i) for i in self._all_times_steps[1][:idx]]
            prev_idx = diff.index(min(diff))

            return prev_idx

    def OnNext(self, event):
        """ Go to the next step """

        selected_step = self._slider_steps.GetValue()-1
        next_idx = self._find_next(selected_step)

        if next_idx != selected_step:
            self._set_all(next_idx)
            self.Refresh(next_idx)

    def OnPrev(self, event):
        """ Go to the previous step """

        selected_step = self._slider_steps.GetValue()-1
        prev_idx = self._find_prev(selected_step)

        if prev_idx != selected_step:
            self._set_all(prev_idx)
            self.Refresh(prev_idx)

    def OnTextTime(self, event):
        try:
            self._starting_date = datetime.strptime(self._texttime.GetValue(), '%Y-%m-%d %H:%M:%S')
            self._step_time.Set(['{:.3f} - {}'.format(i, datetime.strftime(self._starting_date + timedelta(seconds=i), '%Y-%m-%d %H:%M:%S')) for i in self._all_times_steps[0]])
        except:
            logging.info('Error while parsing the date')
            pass

    def OnClose(self, event):
        """ Close the simulation explorer """

        self.mapviewer._pop_sim_explorer(self.active_res2d)
        self.Destroy()

    def OnUpdate(self, event):
        self._update()

    def OnApply(self, event):
        selected_step = self._slider_steps.GetValue()-1

        self._cmd_apply.SetBackgroundColour(wx.Colour(255, 0, 0))  # Set button color to red
        self._cmd_apply.Refresh()  # Refresh the button to apply the color change

        self.Refresh(selected_step)

        self._cmd_apply.SetBackgroundColour(wx.NullColour)  # Reset button color to default
        self._cmd_apply.Refresh()  # Refresh the button to apply the color change

    def _set_all(self, idx:int):

        # test if idx is in range
        if idx < 0 :
            logging.error(_('Index out of range'))
            return
        if idx >= len(self._all_times_steps[0]):
            self._update()
            if idx >= len(self._all_times_steps[0]):
                logging.error(_('Index out of range'))
                return

        try:
            self._slider_steps.SetValue(idx+1)
            self._step_idx.SetSelection(idx)
            self._step_time.SetSelection(idx)
            self._step_num.SetSelection(idx)
        except:
            logging.error(_('Error while setting the step selection'))

    def Refresh(self, idx:int):
        self.active_res2d.read_oneresult(idx)
        self.active_res2d.set_currentview()
        self.mapviewer.Refresh()

    def OnSliderSteps(self, event):
        selected_step = self._slider_steps.GetValue()
        self._set_all(selected_step-1)

    def OnSelectCurTime(self, event):
        selected_time = self._step_num.GetSelection()
        self._set_all(selected_time)

    def OnSelectNumStep(self, event):
        selected_step = self._step_time.GetSelection()
        self._set_all(selected_step)

    def OnSelectIdxStep(self, event):
        selected_step = self._step_idx.GetSelection()
        self._set_all(selected_step)

    def _update(self):
        nb = self.active_res2d.get_nbresults()
        self._all_times_steps = self.active_res2d.get_times_steps()

        self._slider_steps.SetMax(nb)
        self._step_idx.Set([str(i) for i in range(1,nb+1)])
        self._step_time.Set(['{:.3f} - {}'.format(i, datetime.strftime(self._starting_date + timedelta(seconds=float(i)), '%Y-%m-%d %H:%M:%S')) for i in self._all_times_steps[0]])
        self._step_num.Set([str(i) for i in self._all_times_steps[1]])


class Sim_VideoCreation(wx.Dialog):

    def __init__(self, parent, title, mapviewer:"WolfMapViewer", sim:Wolfresults_2D):
        super(Sim_VideoCreation, self).__init__(parent, title=title, size=(350, 250), style = wx.DEFAULT_DIALOG_STYLE & ~ (wx.RESIZE_BORDER | wx.MAXIMIZE_BOX | wx.MINIMIZE_BOX | wx.CLOSE_BOX))

        self.mapviewer = mapviewer
        self.active_res2d:Wolfresults_2D = sim

        self._framerate = 25
        self._start_step = 1
        self._end_step = self.active_res2d.get_nbresults()
        self._interval = 1
        self._fn = str(Path(self.active_res2d.filename).parent / f'{Path(self.active_res2d.filename).stem}.avi')

        self._fontsize = 16
        self._fontcolor = (255, 255, 255, 255)
        self._timeposition = 'top-left' # 'top-right', 'bottom-left', 'bottom-right', 'top-center', 'bottom-center'

        panel = wx.Panel(self)
        vbox = wx.BoxSizer(wx.VERTICAL)

        hbox1 = wx.BoxSizer(wx.HORIZONTAL)
        st1 = wx.StaticText(panel, -1, _('File name'))
        hbox1.Add(st1, 1, wx.EXPAND|wx.ALIGN_LEFT|wx.ALL, 5)
        tc1 = wx.TextCtrl(panel, -1, self._fn)
        hbox1.Add(tc1,1,wx.EXPAND|wx.ALIGN_LEFT|wx.ALL,5)
        vbox.Add(hbox1,0,wx.EXPAND|wx.ALIGN_LEFT|wx.ALL,5)
        self.tc1 = tc1

        # Add a button to choose the file name
        btn_browse = wx.Button(panel, -1, _('Browse'))
        hbox1.Add(btn_browse, 0, wx.EXPAND|wx.ALIGN_LEFT|wx.ALL,5)
        self.btn_browse = btn_browse
        self.btn_browse.Bind(wx.EVT_BUTTON, self.OnBrowse)

        hbox2 = wx.BoxSizer(wx.HORIZONTAL)
        st2 = wx.StaticText(panel, -1, _('Frame rate [nb_images/second]'), style=wx.ALIGN_CENTER)
        hbox2.Add(st2, 1, wx.EXPAND|wx.ALIGN_LEFT|wx.ALL, 5)
        tc2 = wx.TextCtrl(panel, -1, str(self._framerate))
        hbox2.Add(tc2,1,wx.EXPAND|wx.ALIGN_LEFT|wx.ALL,5)
        vbox.Add(hbox2,0,wx.EXPAND|wx.ALIGN_LEFT|wx.ALL,5)
        self.tc2 = tc2

        hbox3 = wx.BoxSizer(wx.HORIZONTAL)
        st3 = wx.StaticText(panel, -1, _('First step'), style=wx.ALIGN_CENTER)
        hbox3.Add(st3, 1, wx.EXPAND|wx.ALIGN_LEFT|wx.ALL, 5)
        tc3 = wx.TextCtrl(panel, -1, str(self._start_step))
        hbox3.Add(tc3,1,wx.EXPAND|wx.ALIGN_LEFT|wx.ALL,5)
        vbox.Add(hbox3,0,wx.EXPAND|wx.ALIGN_LEFT|wx.ALL,5)
        self.tc3 = tc3

        hbox4 = wx.BoxSizer(wx.HORIZONTAL)
        st4 = wx.StaticText(panel, -1, _('Final step'), style=wx.ALIGN_CENTER)
        hbox4.Add(st4, 1, wx.EXPAND|wx.ALIGN_LEFT|wx.ALL, 5)
        tc4 = wx.TextCtrl(panel, -1, str(self._end_step))
        hbox4.Add(tc4,1,wx.EXPAND|wx.ALIGN_LEFT|wx.ALL,5)
        vbox.Add(hbox4,0,wx.EXPAND|wx.ALIGN_LEFT|wx.ALL,5)
        self.tc4 = tc4

        hbox5 = wx.BoxSizer(wx.HORIZONTAL)
        st5 = wx.StaticText(panel, -1, _('Interval'), style=wx.ALIGN_CENTER)
        hbox5.Add(st5, 1, wx.EXPAND|wx.ALIGN_LEFT|wx.ALL, 5)
        tc5 = wx.TextCtrl(panel, -1, str(self._interval))
        hbox5.Add(tc5,1,wx.EXPAND|wx.ALIGN_LEFT|wx.ALL,5)
        vbox.Add(hbox5,0,wx.EXPAND|wx.ALIGN_LEFT|wx.ALL,5)
        self.tc5 = tc5

        hbox7 = wx.BoxSizer(wx.HORIZONTAL)
        st7 = wx.StaticText(panel, -1, _("Font size (for time stamp)"), style=wx.ALIGN_CENTER)
        hbox7.Add(st7, 1, wx.EXPAND|wx.ALIGN_LEFT|wx.ALL, 5)
        tc7 = wx.TextCtrl(panel, -1, str(self._fontsize))
        hbox7.Add(tc7,1,wx.EXPAND|wx.ALIGN_LEFT|wx.ALL,5)
        vbox.Add(hbox7,0,wx.EXPAND|wx.ALIGN_LEFT|wx.ALL,5)
        self.tc7 = tc7

        hbox8 = wx.BoxSizer(wx.HORIZONTAL)
        st8 = wx.StaticText(panel, -1, _("Font color (R,G,B,A)"), style=wx.ALIGN_CENTER)
        hbox8.Add(st8, 1, wx.EXPAND|wx.ALIGN_LEFT|wx.ALL, 5)
        tc8 = wx.ColourPickerCtrl(panel, -1, wx.Colour(*self._fontcolor))
        hbox8.Add(tc8,1,wx.EXPAND|wx.ALIGN_LEFT|wx.ALL,5)
        vbox.Add(hbox8,0,wx.EXPAND|wx.ALIGN_LEFT|wx.ALL,5)
        self.tc8 = tc8

        hbox9 = wx.BoxSizer(wx.HORIZONTAL)
        st9 = wx.StaticText(panel, -1, _("Time position"), style=wx.ALIGN_CENTER)
        hbox9.Add(st9, 1, wx.EXPAND|wx.ALIGN_LEFT|wx.ALL, 5)
        choices = ['top-left', 'top-right', 'bottom-left', 'bottom-right', 'top-center', 'bottom-center']
        tc9 = wx.Choice(panel, -1, choices=choices)
        hbox9.Add(tc9, 1, wx.EXPAND|wx.ALIGN_LEFT|wx.ALL, 5)
        vbox.Add(hbox9, 0, wx.EXPAND|wx.ALIGN_LEFT|wx.ALL, 5)
        self.tc9 = tc9

        hbox6 = wx.BoxSizer(wx.HORIZONTAL)
        btn1 = wx.Button(panel, -1, _('Ok'))
        hbox6.Add(btn1,1,wx.EXPAND|wx.ALIGN_LEFT|wx.ALL,5)
        self.btn1 = btn1

        btn2 = wx.Button(panel, -1, _('Cancel'))
        hbox6.Add(btn2,1,wx.EXPAND|wx.ALIGN_LEFT|wx.ALL,5)
        self.btn2 = btn2
        vbox.Add(hbox6,0,wx.EXPAND|wx.ALIGN_LEFT|wx.ALL,5)
        panel.SetSizer(vbox)
        vbox.Fit(self)

        self.btn1.Bind(wx.EVT_BUTTON, self.OnOk)
        self.btn2.Bind(wx.EVT_BUTTON, self.OnCancel)

        # Add validation to the text controls
        self.tc2.Bind(wx.EVT_TEXT, self.OnValidate)
        self.tc3.Bind(wx.EVT_TEXT, self.OnValidate)
        self.tc4.Bind(wx.EVT_TEXT, self.OnValidate)
        self.tc5.Bind(wx.EVT_TEXT, self.OnValidate)
        self.tc7.Bind(wx.EVT_TEXT, self.OnValidate)
        self.tc8.Bind(wx.EVT_COLOURPICKER_CHANGED, self.OnValidate)
        self.tc9.SetSelection(0)
        self.tc9.Bind(wx.EVT_CHOICE, self.OnValidate)

        icon = wx.Icon()
        icon_path = Path(__file__).parent / "apps/wolf.ico"
        icon.CopyFromBitmap(wx.Bitmap(str(icon_path), wx.BITMAP_TYPE_ANY))
        self.SetIcon(icon)
        self.CenterOnScreen()
        self.Show()

    def OnBrowse(self, event):
        """ Browse a file name to save the video """
        with wx.FileDialog(self, _("Save video file"), wildcard="MP4 files (*.mp4)|*.mp4|AVI files (*.avi)|*.avi|All files (*.*)|*.*",
                           style=wx.FD_SAVE | wx.FD_OVERWRITE_PROMPT) as fileDialog:
            if fileDialog.ShowModal() == wx.ID_CANCEL:
                return     # the user changed idea...

            self._fn = fileDialog.GetPath()
            self.tc1.SetValue(self._fn)

    def get_values(self):
        """ Return the values set in the dialog """
        return self._fn, self._framerate, self._start_step, self._end_step, self._interval, self._fontsize, self._fontcolor, self._timeposition

    def OnValidate(self, event):
        """ Validate the text controls to be sure that the values are correct """
        try:
            framerate = int(self.tc2.GetValue())
            if framerate <= 0:
                framerate = 25
                self.tc2.SetValue(str(framerate))
            self._framerate = framerate
        except:
            self.tc2.SetValue(str(self._framerate))
        try:
            start_step = int(self.tc3.GetValue())
            if start_step < 1:
                start_step = 1
                self.tc3.SetValue(str(start_step))
            if start_step > self.active_res2d.get_nbresults():
                start_step = self.active_res2d.get_nbresults()
                self.tc3.SetValue(str(start_step))
            self._start_step = start_step

        except:
            self.tc3.SetValue(str(self._start_step))
        try:
            end_step = int(self.tc4.GetValue())
            if end_step < 1:
                end_step = 1
                self.tc4.SetValue(str(end_step))
            if end_step > self.active_res2d.get_nbresults():
                end_step = self.active_res2d.get_nbresults()
                self.tc4.SetValue(str(end_step))
            self._end_step = end_step
        except:
            self.tc4.SetValue(str(self._end_step))
        try:
            interval = int(self.tc5.GetValue())
            if interval < 1:
                interval = 1
                self.tc5.SetValue(str(interval))
            self._interval = interval
        except:
            self.tc5.SetValue(str(self._interval))

        try:
            fontsize = int(self.tc7.GetValue())
            if fontsize < 1:
                fontsize = 16
                self.tc7.SetValue(str(fontsize))
            self._fontsize = fontsize
        except:
            self.tc7.SetValue(str(self._fontsize))

        try:
            color = self.tc8.GetColour()
            self._fontcolor = (color.Red(), color.Green(), color.Blue(), color.Alpha())
        except:
            self.tc8.SetColour(wx.Colour(*self._fontcolor))

        try:
            pos_idx = self.tc9.GetSelection()
            choices = ['top-left', 'top-right', 'bottom-left', 'bottom-right', 'top-center', 'bottom-center']
            if pos_idx < 0 or pos_idx >= len(choices):
                pos_idx = 0
                self.tc9.SetSelection(pos_idx)
            self._timeposition = choices[pos_idx]
        except:
            self.tc9.SetSelection(0)

    def OnOk(self, event):
        """ Create the video file """
        self.OnValidate(None)
        self.EndModal(wx.ID_OK)

    def OnCancel(self, event):
        """ Cancel the video creation """
        self.EndModal(wx.ID_CANCEL)


class Drowning_Explorer(wx.Frame):

    def __init__(self, parent, title, mapviewer:any, sim:Drowning_victim_Viewer):

        super().__init__(parent, title=title, size=(150, 250), style = wx.DEFAULT_FRAME_STYLE & ~ (wx.MAXIMIZE_BOX | wx.MINIMIZE_BOX))

        self._panel = wx.Panel(self)

        self.mapviewer = mapviewer
        self.active_drowning:Drowning_victim_Viewer = sim

        main_sizer = wx.BoxSizer(wx.HORIZONTAL)

        left_bar = wx.BoxSizer(wx.VERTICAL)
        right_bar = wx.BoxSizer(wx.VERTICAL)

        self._all_times_steps = self.active_drowning.wanted_time

        # Right bar
        # ---------

        # Slider
        self._slider_steps = wx.Slider(self._panel, minValue=1, maxValue=len(self.active_drowning.wanted_time)-1, style=wx.SL_HORIZONTAL | wx.SL_AUTOTICKS | wx.SL_MIN_MAX_LABELS | wx.SL_LABELS)
        self._slider_steps.Bind(wx.EVT_SLIDER, self.OnSliderSteps)
        right_bar.Add(self._slider_steps, 1, wx.EXPAND | wx.ALL, 2)

        self._time_drowning = wx.TextCtrl(self._panel, value=f"Drowning at 0 days, 0 hours,\n0 minutes and 0 seconds", style=wx.TE_MULTILINE | wx.TE_READONLY | wx.TE_CENTER)
        right_bar.Add(self._time_drowning, 0, wx.EXPAND | wx.ALL, 2)

        # Explore by time
        self._label_time = wx.StaticText(self._panel, label=_('Time [s]'))
        right_bar.Add(self._label_time, 1, wx.EXPAND | wx.ALL, 2)

        _now = datetime.now()
        self._starting_date = datetime(year=_now.year, month=_now.month, day=_now.day, hour=0, minute=0, second=0)
        self._texttime = wx.TextCtrl(self._panel, value=self._starting_date.strftime('%Y-%m-%d %H:%M:%S'))
        right_bar.Add(self._texttime, 1, wx.EXPAND | wx.ALL, 5)
        self._texttime.Bind(wx.EVT_TEXT, self.OnTextTime)

        self._step_time = wx.ListBox(self._panel, choices=['{:.3f} - {}'.format(i, datetime.strftime(self._starting_date + timedelta(seconds=i), '%Y-%m-%d %H:%M:%S')) for i in self._all_times_steps[:-1]], style=wx.LB_SINGLE)
        self._step_time.Bind(wx.EVT_LISTBOX, self.OnSelectNumStep)
        right_bar.Add(self._step_time, 1, wx.EXPAND | wx.ALL, 5)

        self._step_idx = wx.ListBox(self._panel, choices=[str(i) for i in range(1, len(self.active_drowning.wanted_time)-1+1)], style=wx.LB_SINGLE)
        self._step_idx.Bind(wx.EVT_LISTBOX, self.OnSelectIdxStep)
        right_bar.Add(self._step_idx, 1, wx.EXPAND | wx.ALL, 5)

        # Left bar
        # --------

        # Apply selected step
        self._cmd_apply = wx.Button(self._panel, wx.ID_APPLY, _('Apply'))
        self._cmd_apply.SetToolTip(_('Apply the selected parameters to the map'))
        self._cmd_apply.Bind(wx.EVT_BUTTON, self.OnApply)
        left_bar.Add(self._cmd_apply, 1, wx.EXPAND | wx.ALL, 5)

        # Next step
        self._cmd_next = wx.Button(self._panel, wx.ID_FORWARD, _('Next'))
        self._cmd_next.SetToolTip(_('Go to the next step -- using the selected mode'))
        self._cmd_next.Bind(wx.EVT_BUTTON, self.OnNext)
        left_bar.Add(self._cmd_next, 1, wx.EXPAND | wx.ALL, 5)

        # Previous step
        self._cmd_prev = wx.Button(self._panel, wx.ID_BACKWARD, _('Previous'))
        self._cmd_prev.SetToolTip(_('Go to the previous step -- using the selected mode'))
        self._cmd_prev.Bind(wx.EVT_BUTTON, self.OnPrev)
        left_bar.Add(self._cmd_prev, 1, wx.EXPAND | wx.ALL, 5)

        self.Bind(wx.EVT_CLOSE, self.OnClose)

        main_sizer.Add(left_bar, 1, wx.EXPAND | wx.ALL, 2)
        main_sizer.Add(right_bar, 1, wx.EXPAND | wx.ALL, 2)

        self._panel.SetSizer(main_sizer)
        self._panel.SetAutoLayout(True)

        self.MinSize = (450, 500)

        self.Fit()
        self.Show()

        self.SetIcon(wx.Icon(str(Path(__file__).parent / "apps/wolf.ico")))

        self._set_all(0)

    def _find_next(self, idx:int):
        """ Find the next step based on the mode and interval """

        mode = 2
        if mode == 0:
            # By time [s]
            next_time = self._all_times_steps[idx] + float(self._interval.GetValue())
            diff = [abs(next_time - i) for i in self._all_times_steps[idx:]]
            next_idx = diff.index(min(diff)) + idx

            return next_idx

        elif mode == 1:
            # By time [hour]
            next_time = self._all_times_steps[idx] + float(self._interval.GetValue())*3600
            diff = [abs(next_time - i) for i in self._all_times_steps[idx:]]
            next_idx = diff.index(min(diff)) + idx

            return next_idx

        elif mode == 2:
            # By index
            next_idx = min(idx + int(1), len(self._all_times_steps)-1)

            return next_idx

        elif mode == 3:
            # By time step
            next_idx = self._all_times_steps[1].index(self._all_times_steps[idx] + int(1))
            diff = [abs(next_idx - i) for i in self._all_times_steps[idx:]]
            next_idx = diff.index(min(diff)) + idx

            return next_idx

    def _find_prev(self, idx:int):
        """ Find the previous step based on the mode and interval """

        mode = 2

        if mode == 0:
            # By time [s]
            prev_time = self._all_times_steps[idx] - float(1)
            diff = [abs(prev_time - i) for i in self._all_times_steps[:idx]]
            prev_idx = diff.index(min(diff))

            return prev_idx

        elif mode == 1:
            # By time [hour]
            prev_time = self._all_times_steps[idx] - float(1)*3600
            diff = [abs(prev_time - i) for i in self._all_times_steps[:idx]]
            prev_idx = diff.index(min(diff))

            return prev_idx

        elif mode == 2:
            # By index
            prev_idx = max(idx - int(1), 0)

            return prev_idx

        elif mode == 3:
            # By time step
            prev_idx = self._all_times_steps[1].index(self._all_times_steps[idx] - int(1))
            diff = [abs(prev_idx - i) for i in self._all_times_steps[:idx]]
            prev_idx = diff.index(min(diff))

            return prev_idx

    def OnNext(self, event):
        """ Go to the next step """

        selected_step = self._slider_steps.GetValue()+1
        next_idx = self._find_next(selected_step)

        if next_idx != selected_step:
            self._set_all(next_idx)
            self.Refresh(next_idx)

    def OnTextTime(self, event):
        try:
            self._starting_date = datetime.strptime(self._texttime.GetValue(), '%Y-%m-%d %H:%M:%S')
            self._step_time.Set(['{:.3f} - {}'.format(int(i/3600/24), datetime.strftime(self._starting_date + timedelta(seconds=i), '%Y-%m-%d %H:%M:%S')) for i in self._all_times_steps[:-1]])
        except:
            pass

    def OnPrev(self, event):
        """ Go to the previous step """

        selected_step = self._slider_steps.GetValue()-1
        prev_idx = self._find_prev(selected_step)

        if prev_idx != selected_step:
            self._set_all(prev_idx)
            self.Refresh(prev_idx)

    def OnClose(self, event):
        """ Close the simulation explorer """

        self.mapviewer._pop_sim_explorer(self.active_drowning)
        self.Destroy()

    def OnUpdate(self, event):
        self._update()

    def OnApply(self, event):
        selected_step = self._slider_steps.GetValue()-1

        self._cmd_apply.SetBackgroundColour(wx.Colour(255, 0, 0))  # Set button color to red
        self._cmd_apply.Refresh()  # Refresh the button to apply the color change

        self.Refresh(selected_step)

        self._cmd_apply.SetBackgroundColour(wx.NullColour)  # Reset button color to default
        self._cmd_apply.Refresh()  # Refresh the button to apply the color change

    def _set_all(self, idx:int):
        self._slider_steps.SetValue(idx+1)
        self._step_idx.SetSelection(idx)

    def Refresh(self, idx:int):
        self.active_drowning.read_oneresult(idx)
        self.mapviewer.Refresh()

    def OnSliderSteps(self, event):
        selected_step = self._slider_steps.GetValue()-1
        self.active_drowning.time_id = selected_step
        time_id = self._slider_steps.GetValue()-1
        time_value = self.active_drowning.wanted_time[time_id]

        days = np.floor(time_value // 86400)
        hours = np.floor((time_value % 86400) / 3600)
        minutes = np.floor(((time_value % 86400) % 3600) / 60)
        seconds = np.floor(((time_value % 86400) % 3600) % 60)

        self._time_drowning.SetValue(
            f"Drowning at {int(days)} days, {int(hours)} hours,\n"
            f"{int(minutes)} minutes and {int(seconds)} seconds"
        )
        self._set_all(selected_step)

    def OnSelectNumStep(self, event):
        selected_step = self._step_time.GetSelection()
        self.active_drowning.time_id = selected_step
        time_id = selected_step
        time_value = self.active_drowning.wanted_time[time_id]

        days = np.floor(time_value // 86400)
        hours = np.floor((time_value % 86400) / 3600)
        minutes = np.floor(((time_value % 86400) % 3600) / 60)
        seconds = np.floor(((time_value % 86400) % 3600) % 60)

        self._time_drowning.SetValue(
            f"Drowning at {int(days)} days, {int(hours)} hours,\n"
            f"{int(minutes)} minutes and {int(seconds)} seconds"
        )
        self._set_all(selected_step)

    def OnSelectIdxStep(self, event):
        selected_step = self._step_idx.GetSelection()
        self.active_drowning.time_id = selected_step
        time_id = selected_step
        time_value = self.active_drowning.wanted_time[time_id]

        days = np.floor(time_value // 86400)
        hours = np.floor((time_value % 86400) / 3600)
        minutes = np.floor(((time_value % 86400) % 3600) / 60)
        seconds = np.floor(((time_value % 86400) % 3600) % 60)

        self._time_drowning.SetValue(
            f"Drowning at {int(days)} days, {int(hours)} hours,\n"
            f"{int(minutes)} minutes and {int(seconds)} seconds"
        )
        self._set_all(selected_step)

    def _update(self):
        nb = len(self.active_drowning.wanted_time)
        self._all_times_steps = self.active_drowning.wanted_time

        self._slider_steps.SetMax(nb)
        self._step_idx.Set([str(i) for i in range(1,nb+1)])

class Select_Begin_end_interval_step(wx.Dialog):
    """ wx.frame to select the begin and end of the interval to extract """

    def __init__(self, parent, title, sim:Wolfresults_2D, checkbox:bool = False):

        super(Select_Begin_end_interval_step, self).__init__(parent, title=title, size=(500, 350), style = wx.DEFAULT_FRAME_STYLE & ~ (wx.MAXIMIZE_BOX | wx.MINIMIZE_BOX))

        # ajout d'un slider pour choisir le début et la fin de l'intervalle -> selrange
        # ajout d'un slider pour choisir le pas de l'intervalle


        # + les mêmes informations mais sous forme de TextCtrl

        # ajout d'un bouton pour valider
        # ajout d'un bouton pour annuler

        self._panel = wx.Panel(self)

        sizer = wx.BoxSizer(wx.VERTICAL)

        self.begin = 1
        self.end = sim.get_nbresults(True)
        self.step = 1
        self.check_all = True
        self.check_violin = False

        self._slider_begin = wx.Slider(self._panel, minValue=self.begin, maxValue=self.end, style=wx.SL_HORIZONTAL | wx.SL_AUTOTICKS | wx.SL_MIN_MAX_LABELS | wx.SL_LABELS)
        self._slider_begin.SetToolTip(_('Select the first result to export'))
        self._slider_begin.SetValue(self.begin)

        self._slider_begin.Bind(wx.EVT_SLIDER, self.OnSliderBegin)
        sizer.Add(self._slider_begin, 1, wx.EXPAND | wx.ALL, 2)

        self._slider_end = wx.Slider(self._panel, minValue=self.begin, maxValue=self.end, style=wx.SL_HORIZONTAL | wx.SL_AUTOTICKS | wx.SL_MIN_MAX_LABELS | wx.SL_LABELS)
        self._slider_end.SetToolTip(_('Select the last result to export - If step is > 1, this value will be forced if not already captured'))
        self._slider_end.SetValue(self.end)
        self._slider_end.Bind(wx.EVT_SLIDER, self.OnSliderEnd)

        sizer.Add(self._slider_end, 1, wx.EXPAND | wx.ALL, 2)

        sizer_txt1 = wx.BoxSizer(wx.HORIZONTAL)
        self._label_range = wx.StaticText(self._panel, label=_('Range'))
        self._text_range = wx.TextCtrl(self._panel, value='1 - {}'.format(sim.get_nbresults(True)))

        sizer_txt1.Add(self._label_range, 0, wx.EXPAND | wx.ALL, 2)
        sizer_txt1.Add(self._text_range, 1, wx.EXPAND | wx.ALL, 2)

        sizer.Add(sizer_txt1, 0, wx.EXPAND | wx.ALL, 2)

        self._slider_step = wx.Slider(self._panel, minValue=1, maxValue=sim.get_nbresults(True), style=wx.SL_HORIZONTAL | wx.SL_AUTOTICKS | wx.SL_MIN_MAX_LABELS | wx.SL_LABELS)
        self._slider_step.SetToolTip(_('Export one result every N steps'))
        self._slider_step.Bind(wx.EVT_SLIDER, self.OnSliderStep)
        sizer.Add(self._slider_step, 1, wx.EXPAND | wx.ALL, 2)

        sizer_txt2 = wx.BoxSizer(wx.HORIZONTAL)
        self._label_step = wx.StaticText(self._panel, label=_('Step'))
        self._text_step = wx.TextCtrl(self._panel, value='1')

        sizer_txt2.Add(self._label_step, 0, wx.EXPAND | wx.ALL, 2)
        sizer_txt2.Add(self._text_step, 0, wx.EXPAND | wx.ALL, 2)

        sizer.Add(sizer_txt2, 0, wx.EXPAND | wx.ALL, 2)

        sizer_but = wx.BoxSizer(wx.HORIZONTAL)
        self._cmd_apply = wx.Button(self._panel, wx.ID_APPLY, _('Apply'))
        self._cmd_apply.Bind(wx.EVT_BUTTON, self.OnApply)

        self._cmd_ok = wx.Button(self._panel, wx.ID_OK, _('OK'))
        self._cmd_ok.Bind(wx.EVT_BUTTON, self.OnOK)

        self._cmd_cancel = wx.Button(self._panel, wx.ID_CANCEL, _('Cancel'))
        self._cmd_cancel.Bind(wx.EVT_BUTTON, self.OnCancel)

        sizer_but.Add(self._cmd_apply, 1, wx.EXPAND | wx.ALL, 2)
        sizer_but.Add(self._cmd_ok, 1, wx.EXPAND | wx.ALL, 2)
        sizer_but.Add(self._cmd_cancel, 1, wx.EXPAND | wx.ALL, 2)

        sizer.Add(sizer_but, 1, wx.EXPAND | wx.ALL, 2)

        if checkbox:
            sizer_check = wx.BoxSizer(wx.HORIZONTAL)
            self._check_all = wx.CheckBox(self._panel, label=_('Statistics and values'), style=wx.CHK_2STATE)
            self._check_all.SetToolTip(_('If checked, export statistics and all values for each step'))
            self._check_all.SetValue(True)
            self._check_all.Bind(wx.EVT_CHECKBOX, self.OnCheckAll)
            sizer_check.Add(self._check_all, 1, wx.EXPAND | wx.ALL, 2)
            sizer.Add(sizer_check, 1, wx.EXPAND | wx.ALL, 2)

            self._check_violin= wx.CheckBox(self._panel, label=_('Violin plot (experimental)'), style=wx.CHK_2STATE)
            self._check_violin.SetToolTip(_('If checked, create a violin plot for each step'))
            self._check_violin.SetValue(False)
            self._check_violin.Bind(wx.EVT_CHECKBOX, self.OnCheckViolin)
            sizer_check.Add(self._check_violin, 1, wx.EXPAND | wx.ALL, 2)

        self._panel.SetSizer(sizer)

        self.CenterOnScreen()

        self.SetIcon(wx.Icon(str(Path(__file__).parent / "apps/wolf.ico")))

        self.Show()


    def OnCheckAll(self, event):
        self.check_all = self._check_all.IsChecked()

    def OnCheckViolin(self, event):
        self.check_violin = self._check_violin.IsChecked()

    def OnSliderBegin(self, event):

        self.begin = min(self._slider_begin.GetValue(), self.end)
        self._slider_begin.SetValue(self.begin)
        self._text_range.SetValue('{} - {}'.format(self.begin, self.end))

    def OnSliderEnd(self, event):

        self.end = max(self._slider_end.GetValue(), self.begin)
        self._slider_end.SetValue(self.end)
        self._text_range.SetValue('{} - {}'.format(self.begin, self.end))

    def OnSliderStep(self, event):

        self.step = self._slider_step.GetValue()
        self._text_step.SetValue(str(self.step))

    def OnApply(self, event):

        try:
            txt_begin, txt_end = self._text_range.GetValue().split('-')
        except:
            self._text_range.SetValue('{} - {}'.format(self.begin, self.end))

        txt_step = self._text_step.GetValue()

        try:
            if self.step != int(txt_step):
                self._slider_step.SetValue(int(txt_step))
        except:
            logging.error('Error while parsing the step')
            return

        try:
            if int(txt_begin) != self.begin or int(txt_end) != self.end:
                self._slider_begin.SetRange(int(txt_begin), int(txt_end))
        except:
            logging.error('Error while parsing the range')
            return

    def OnOK(self, event):
        self.Hide()

    def OnCancel(self, event):
        self.begin = -1
        self.end = -1
        self.step = -1

        self.Hide()


class PrecomputedDEM_DTM(Enum):
    """ Enum for Precomputed DEM/DTM array """

    DEMDTM_50cm = "AllData.vrt"
    DEMDTM_1m_average = "Combine_1m_average.vrt"
    DEMDTM_1m_min = "Combine_1m_minimum.vrt"
    DEMDTM_1m_max = "Combine_1m_maximum.vrt"
    DEMDTM_2m_average = "Combine_2m_average.vrt"
    DEMDTM_2m_min = "Combine_2m_minimum.vrt"
    DEMDTM_2m_max = "Combine_2m_maximum.vrt"
    DEMDTM_5m_average = "Combine_5m_average.vrt"
    DEMDTM_5m_min = "Combine_5m_minimum.vrt"
    DEMDTM_5m_max = "Combine_5m_maximum.vrt"
    DEMDTM_10m_average = "Combine_10m_average.vrt"
    DEMDTM_10m_min = "Combine_10m_minimum.vrt"
    DEMDTM_10m_max = "Combine_10m_maximum.vrt"

class Precomputed_DEM_DTM_Dialog(wx.Dialog):
    """ wx.Dialog to select Precomputed DEM/DTM array

    Resolutions are 50cm, 1m, 2m, 5m, 10m
    Operators are average, min, max
    """

    def __init__(self, parent, title, directory:Path | str, mapviewer:"WolfMapViewer"):

        super(Precomputed_DEM_DTM_Dialog, self).__init__(parent, title=title, size=(500, 350), style = wx.DEFAULT_FRAME_STYLE & ~ (wx.MAXIMIZE_BOX | wx.MINIMIZE_BOX))

        self._dir = Path(directory)
        self._header = None
        self._vrt = None
        self._mapviewer = mapviewer

        self.available_vrt()

        self._panel = wx.Panel(self)
        sizer = wx.BoxSizer(wx.VERTICAL)

        # Listbox with all available operators
        self._res = ['50cm', '1m', '2m', '5m', '10m']
        self._ops = ['average', 'minimum', 'maximum']
        self._resolution = wx.ListBox(self._panel, choices=self._res, style=wx.LB_SINGLE)
        self._resolution.Bind(wx.EVT_LISTBOX, self.OnSelectResolution)

        self._operations = wx.ListBox(self._panel, choices=[], style=wx.LB_SINGLE)

        sizer.Add(self._resolution, 1, wx.EXPAND | wx.ALL, 2)
        sizer.Add(self._operations, 1, wx.EXPAND | wx.ALL, 2)

        sizer_btns = wx.BoxSizer(wx.HORIZONTAL)
        self._cmd_sameactive = wx.Button(self._panel, wx.ID_APPLY, _('Same as active array...'))
        self._cmd_sameas = wx.Button(self._panel, wx.ID_APPLY, _('Same as file...'))
        self._cmd_zoom = wx.Button(self._panel, wx.ID_APPLY, _('On current zoom...'))

        self._cmd_sameas.Bind(wx.EVT_BUTTON, self.OnSameAs)
        self._cmd_zoom.Bind(wx.EVT_BUTTON, self.OnZoom)
        self._cmd_sameactive.Bind(wx.EVT_BUTTON, self.OnSameActive)

        sizer_btns.Add(self._cmd_sameactive, 1, wx.EXPAND | wx.ALL, 2)
        sizer_btns.Add(self._cmd_sameas, 1, wx.EXPAND | wx.ALL, 2)
        sizer_btns.Add(self._cmd_zoom, 1, wx.EXPAND | wx.ALL, 2)

        sizer.Add(sizer_btns, 1, wx.EXPAND | wx.ALL, 2)

        self._panel.SetSizer(sizer)

        self.CenterOnScreen()

        self.SetIcon(wx.Icon(str(Path(__file__).parent / "apps/wolf.ico")))

        self.Show()

    def OnSameAs(self, event):
        """ Set the Precomputed DEM/DTM array to the same bounds as an existing array """

        dlg = wx.FileDialog(self, _('Select a file'), str(self._dir), '', "All supported formats|*.bin;*.tif;*.tiff;*.top;*.flt;*.npy;*.npz;*.vrt|bin (*.bin)|*.bin|Elevation WOLF2D (*.top)|*.top|Geotif (*.tif)|*.tif|Float ESRI (*.flt)|*.flt|Numpy (*.npy)|*.npy|Numpy named arrays(*.npz)|*.npz|all (*.*)|*.*", wx.FD_OPEN | wx.FD_FILE_MUST_EXIST)
        ret = dlg.ShowModal()

        if ret == wx.ID_OK:
            fname = Path(dlg.GetPath())

            self._header = header_wolf()
            self._header.read_txt_header(fname)

            res = self._resolution.GetStringSelection()
            if res == '50cm':
                res = 0.5
            elif res in ['1m', '2m', '5m', '10m']:
                res = float(res[:-1])
            else:
                logging.error('Resolution not found')
                return

            if self._header.dx != res or self._header.dy != res:
                logging.warning(_('Resolution not the same'))
                logging.warning(_('Forcing resolution to {}m').format(res))

                self._header.origx = float(int(self._header.origx / res) * res)
                self._header.origy = float(int(self._header.origy / res) * res)
                self._header.nbx = int(np.ceil(float(self._header.nbx) * self._header.dx / res))
                self._header.nby = int(np.ceil(float(self._header.nby) * self._header.dy / res))
                self._header.dx = res
                self._header.dy = res

                logging.info(_('New header:'))
                logging.info(self._header)

            self.add_array()
            self.Hide()

    def OnSameActive(self, event):
        """ Set the Precomputed DEM/DTM array to the same bounds as the active array """

        if self._mapviewer is None:
            logging.error('No mapviewer to get the active array')
            return

        active = self._mapviewer.active_array
        if active is None:
            logging.error('No active array to get the bounds')
            return

        self._header = active.get_header()

        res = self._resolution.GetStringSelection()
        if res == '50cm':
            res = 0.5
        elif res in ['1m', '2m', '5m', '10m']:
            res = float(res[:-1])
        else:
            logging.error('Resolution not found')
            return

        if self._header.dx != res or self._header.dy != res:
            logging.warning(_('Resolution not the same'))
            logging.warning(_('Forcing resolution to {}m').format(res))

            self._header.origx = float(int(self._header.origx / res) * res)
            self._header.origy = float(int(self._header.origy / res) * res)
            self._header.nbx = int(np.ceil(float(self._header.nbx) * self._header.dx / res))
            self._header.nby = int(np.ceil(float(self._header.nby) * self._header.dy / res))
            self._header.dx = res
            self._header.dy = res

            logging.info(_('New header:'))
            logging.info(self._header)

        newarray = self.add_array()

        #copy palette
        if newarray is not None:
            newarray.mypal.automatic = False
            newarray.mypal.values = active.mypal.values

        self.Hide()
        self._mapviewer.Refresh()

    def OnZoom(self, event):
        """ Set the Precomputed DEM/DTM array to the current zoom """

        if self._mapviewer is None:
            logging.error('No mapviewer to get the current zoom')
            return

        onzoom = [self._mapviewer.xmin, self._mapviewer.xmax, self._mapviewer.ymin, self._mapviewer.ymax]

        self._header = header_wolf()

        # round to the nearest resolution
        self._header.origx = float(int(onzoom[0]))
        self._header.origy = float(int(onzoom[2]))

        res = self._resolution.GetStringSelection()
        if res == '50cm':
            res = 0.5
        elif res in ['1m', '2m', '5m', '10m']:
            res = float(res[:-1])
        else:
            logging.error('Resolution not found')
            return

        self._header.dx = res
        self._header.dy = res

        self._header.nbx = int(float(np.ceil(onzoom[1]) - int(onzoom[0])) / res)
        self._header.nby = int(float(np.ceil(onzoom[3]) - int(onzoom[2])) / res)

        self.add_array()
        self.Hide()
        self._mapviewer.Refresh()

    @property
    def selected_vrt(self):

        res = self._resolution.GetStringSelection()
        op = self._operations.GetStringSelection()

        vrt_names = [cur.name for cur in self._vrt]
        if res == '50cm':
            if PrecomputedDEM_DTM.DEMDTM_50cm.value in vrt_names:
                return self._dir / PrecomputedDEM_DTM.DEMDTM_50cm.value
        elif res in ['1m', '2m', '5m', '10m']:
            to_test = 'Combine_{}_{}.vrt'.format(res, op)
            if to_test in vrt_names:
                return self._dir / to_test
            else:
                logging.error(_('Operator not found - Did you select one?'))
                return None
        else:
            logging.error(_('Resolution not found - Did you select one?'))
            return None

    def add_array(self):
        """ Add a new array to the viewer """

        if self._mapviewer is None:
            logging.error(_('No mapviewer to add the array'))
            return

        if self._header is None:
            logging.error(_('No header defined'))
            return

        vrt = self.selected_vrt
        if vrt is None:
            logging.error(_('No vrt selected'))
            return

        newarray = WolfArray(vrt, crop= [[self._header.origx, self._header.origx + self._header.nbx * self._header.dx], [self._header.origy, self._header.origy + self._header.nby * self._header.dy]])
        self._mapviewer.add_object(newobj = newarray, id = vrt.stem)

        return newarray

    def OnSelectResolution(self, event):
        """ Select the resolution """

        res = self._resolution.GetStringSelection()

        vrt_names = [i.name for i in self._vrt]
        if res == '50cm':
            if PrecomputedDEM_DTM.DEMDTM_50cm.value in vrt_names:
                self._operations.Set(['No operator to choose - 50cm resolution'])
        elif res in ['1m', '2m', '5m', '10m']:
            to_test = {i: 'Combine_{}_{}.vrt'.format(res, i) for i in self._ops}
            self._operations.Set([i for i, val in to_test.items() if val in vrt_names])
        else:
            self._operations.Set([''])


    def available_vrt(self):
        """ List all available vrt files in the directory """

        self._vrt = [i for i in self._dir.iterdir() if i.suffix == '.vrt']
        # test if vrt are in PrecomputedDEM_DTM
        self._vrt = [i for i in self._vrt if i.name in [j.value for j in PrecomputedDEM_DTM]]


class WolfMapViewer(wx.Frame):
    """
    Fenêtre de visualisation de données WOLF grâce aux WxWidgets
    """

    TIMER_ID = 100  # délai d'attente avant action

    mybc: list[BcManager]  # Gestionnaire de CL
    myarrays: list  # matrices ajoutées
    myvectors: list[Zones]  # zones vectorielles ajoutées
    myclouds: list[cloud_vertices]  # nuages de vertices
    mytri: list[Triangulation]  # triangulations
    myothers: list
    myviews:list[views_2D]
    mywmsback: list
    mywmsfore: list
    myres2D: list
    mytiles: list[Tiles]
    myimagestiles: list[ImagesTiles]
    mypartsystems: list[Particle_system]
    myviewers3d:list[Wolf_Viewer3D]
    mylazdata:list[Wolf_LAZ_Data]
    mydrownings: list[Drowning_victim_Viewer]
    mypicturecollections: list[PictureCollection]

    if WOLFPYDIKE_AVAILABLE:
        mydikes: list[DikeWolf]
        myinjectors: list[InjectorDike]

    mymplfigs:list[MplFigViewer]

    sim_explorers: dict[Wolfresults_2D:Sim_Explorer]

    canvas: GLCanvas  # canvas OpenGL
    context: GLContext  # context OpenGL
    mytooltip: Wolf_Param  # Objet WOLF permettant l'analyse de ce qui est sous la souris
    treelist: TreeListCtrl  # Gestion des éléments sous forme d'arbre
    _lbl_selecteditem: StaticText
    leftbox: BoxSizer

    # DEPRECEATED
    # added: dict  # dictionnaire des éléments ajoutés

    active_vector: vector
    active_zone: zone
    active_zones: Zones
    active_array: WolfArray
    active_bc: BcManager
    active_view: WolfViews
    active_vertex: wolfvertex
    active_cs: crosssections
    active_tri: Triangulation
    active_tile: Tiles
    active_imagestiles: ImagesTiles
    active_particle_system: Particle_system
    active_viewer3d: Wolf_Viewer3D
    active_viewerlaz: viewerlaz
    active_bridges: Bridges
    active_bridge: Bridge
    active_weirs : Weirs
    active_weir : Weir
    active_laz : Wolf_LAZ_Data
    active_drowning: Drowning_victim_Viewer

    if WOLFPYDIKE_AVAILABLE:
        active_dike : DikeWolf
        active_injector : InjectorDike

    active_picturecollection: PictureCollection
    active_alaro: IRM_Alaro

    active_fig: MplFigViewer

    alaro_navigator: Alaro_Navigator


    # def check_user_activity(self, *args):
    #     while True:
    #         sleep(1)
    #         if datetime.now() - self._last_activity_time > timedelta(seconds=3): # 5 secondes d'inactivité
    #             args[0]._user_active = False

    # def _user_activity_true(self):
    #     self._user_active = True
    #     self._last_activity_time = datetime.now()

    # def background_task(self, *args):

    #     while args[0]._user_active:
    #         sleep(1)

    #     args[0]._update_background()
    #     args[0].Paint()

    #     args[0]._thread_update_background = None

    def __init__(self,
                 wxparent = None,
                 title:str = _('Default Wolf Map Viewer'),
                 w:int=500,
                 h:int=500,
                 treewidth:int=200,
                 wolfparent=None,
                 wxlogging=None):

        """
        Create a Viewer for WOLF data/simulation

        :params wxparent: wx parent - set to None if main window
        :params title: title of the window
        :params w: width of the window in pixels
        :params h: height of the window in pixels
        :params treewidth: width of the tree list in pixels
        :params wolfparent: WOLF object parent -- see PyGui.py
        :params wxlogging: wx logging object

        """

        self._show_dialog_wx = True  # Show dialog boxes

        # self._user_active = True # True if the user is active in the viewer
        # self._last_activity_time  = datetime.now() # last time the user was active in the viewer
        # self._check_activity_thread = threading.Thread(target=self.check_user_activity, args=[self]) # thread to check user activity
        # self._check_activity_thread.start() # start the thread
        # self._thread_update_background = None # thread to update the background

        self.treewidth = treewidth
        super(WolfMapViewer, self).__init__(wxparent, title=title, size=(w + self.treewidth, h))

        self._wxlogging = wxlogging
        self.action = None  # Action à entreprendre
        self.update_absolute_minmax = False  # Force la MAJ de la palette
        self.copyfrom = None  # aucun élément pointé par CTRL+C

        self.wolfparent = wolfparent

        self.regular = True  # Gestion de la taille de fenêtre d'affichage, y compris l'arbre de gestion
        self.sx = 1  # facteur d'échelle selon X = largeur en pixels/largeur réelle
        self.sy = 1  # facteur d'échelle selon Y = hauteur en pixels/hauteur réelle
        self.samescale = True  # force le même facteur d'échelle

        self.dynapar_dist = 1.

        # emprise initiale
        self.xmin = 0.
        self.ymin = 0.
        self.xmax = 40.
        self.ymax = 40.
        self.width = self.xmax - self.xmin  # largeur de la zone d'affichage en coordonnées réelles
        self.height = self.ymax - self.ymin  # hauteur de la zone d'affichage en coordonnées réelles
        self.canvaswidth = 100
        self.canvasheight = 100

        # position de la caméra
        self.mousex = self.width / 2.
        self.mousey = self.height / 2.
        self._last_mouse_pos = (0, 0, (0,0))

        self.bordersize = 0  # zone réservée au contour
        self.titlesize = 0  # zone réservée au titre
        self.treewidth = 200  # largeur de la zone d'arbre "treelist"

        self.backcolor = wx.Colour(255, 255, 255)  # couleur de fond
        self.mousedown = (-99999., -99999.)  # position initiale du bouton position bas
        # self.mouseup   = (-99999., -99999.)  # position initiale du bouton position haut
        self.oneclick = True  # détection d'un simple click ou d'un double-click
        # self.move = False  # la souris est-elle en train de bouger?

        self.linked = False
        self.link_shareopsvect = True
        self.linkedList = None
        self.link_params = None

        self.project_pal = None

        self.forcemimic = True
        self.currently_readresults = False

        self.mylazgrid:xyz_laz_grids = None # LAZ grid preprocessed by Numpy

        self.colors1to9 = Colors_1to9(self)

        self._dragdrop = DragdropFileTarget(self)
        self.SetDropTarget(self._dragdrop)

        # Gestion des menus
        self.popupmenu = wx.Menu()
        self.popupmenu.Bind(wx.EVT_MENU, self.OnPopupItemSelected)

        for text in [_('Save'), _('Save as'), _('Rename'), _('Duplicate'), _('Delete'), _('Up'), _('Down'), _('Check/Uncheck'), _('Properties'), _('Reload')]:
            item = self.popupmenu.Append(-1, text)

        self.menubar = wx.MenuBar()

        self.menuwolf2d = None
        self.menu_landmap = None
        self.menu_bridge = None
        self.menu_weir = None
        self.menu2d_cache_setup = None
        self.menuparticlesystem = None
        self.menu2dGPU = None
        self.menuLandUseLandCover = None
        self.timer_ps = None
        self.menusim2D  = None
        self.menusim2D_GPU = None
        self.menulaz    = None
        self.menutiles  = None
        self.menuimagestiles = None
        self.menudrowning = None
        self.menudike = None
        self.menupicturecollections = None
        self.menuqdfidf = None
        self.menualaro = None

        self.alaro_navigator = None

        self.filemenu = wx.Menu()
        openitem = self.filemenu.Append(wx.ID_OPEN, _('Open/Add project'), _('Open a full project from file'))
        saveproject = self.filemenu.Append(wx.ID_ANY, _('Save project as...'), _('Save the current project to file'))
        self.filemenu.AppendSeparator()
        saveitem = self.filemenu.Append(wx.ID_SAVE, _('Save'), _('Save all checked arrays or vectors to files'))
        saveasitem = self.filemenu.Append(wx.ID_SAVEAS, _('Save as...'), _('Save all checked arrays or vectors to new files --> one file dialog per data'))
        savecanvas = self.filemenu.Append(wx.ID_ANY, _('Save to image...'), _('Save the canvas to image file on disk'))
        copycanvas = self.filemenu.Append(wx.ID_ANY, _('Copy image...'), _('Copy the canvas to image file to the clipboard'))

        self.filemenu.AppendSeparator()
        # --- GLTF
        self.menugltf = wx.Menu()
        self.filemenu.Append(wx.ID_ANY,_('Gltf2...'), self.menugltf)

        exportgltf = self.menugltf.Append(wx.ID_ANY, _('Export...'), _('Save data to gltf files'))
        importgltf = self.menugltf.Append(wx.ID_ANY, _('Import...'), _('Import data from gltf files'))
        compareitem = self.menugltf.Append(wx.ID_ANY, _('Compare...'), _('Create new frames to compare sculpting'))
        updategltf = self.menugltf.Append(wx.ID_ANY, _('Update...'), _('Update data from gltf files'))

        self.filemenu.AppendSeparator()

        # SIMULATION 2D

        self.menu_sim2d = wx.Menu()
        self.menu_sim2d_cpu = wx.Menu()
        self.menu_sim2d_gpu = wx.Menu()
        self.menu_sim1d = wx.Menu()

        sim2d = self.menu_sim2d_cpu.Append(wx.ID_ANY, _('Create/Open multiblock model'), _('Create or open a multiblock model in the viewer --> CPU/Fortran Wolf2D model'))
        check2D = self.menu_sim2d_cpu.Append(wx.ID_ANY, _('Check headers'), _('Check the header .txt files from an existing 2D CPU simulation'))

        sim2dgpu = self.menu_sim2d_gpu.Append(wx.ID_ANY, _('Create/Open GPU model'), _('Create or open a GPU model in the viewer --> GPU Wolf2D model'))

        create1Dmodel = self.menu_sim1d.Append(wx.ID_ANY, _('Create Wolf1D...'),('Create a 1D model using crossections, vectors and arrays...'))


        self.menu_sim2d.Append(wx.ID_ANY,_('2D GPU'),self.menu_sim2d_gpu)
        self.menu_sim2d.Append(wx.ID_ANY,_('2D CPU'),self.menu_sim2d_cpu)
        self.filemenu.Append(wx.ID_ANY,_('2D Model'),self.menu_sim2d)
        self.filemenu.Append(wx.ID_ANY,_('1D Model'),self.menu_sim1d)

        # self.filemenu.AppendSeparator()

        # SIMULATION Hydrologique

        self.menu_hydrology = wx.Menu()
        hydrol = self.menu_hydrology.Append(wx.ID_ANY, _('Create/Open Hydrological model'), _('Hydrological simulation'))
        self.filemenu.Append(wx.ID_ANY,_('Hydrology'),self.menu_hydrology)

        self.filemenu.AppendSeparator()

        # MULTIVIEWER

        compareitem = self.filemenu.Append(wx.ID_ANY, _('Set comparison'), _('Set comparison'))
        multiview = self.filemenu.Append(wx.ID_ANY, _('Multiviewer'), _('Multiviewer'))
        viewer3d = self.filemenu.Append(wx.ID_ANY, _('3D viewer'), _('3D viewer'))
        self.filemenu.AppendSeparator()


        # ---
        self.menucreateobj = wx.Menu()
        self.filemenu.Append(wx.ID_ANY,_('Create...'),self.menucreateobj)

        createarray = self.menucreateobj.Append(wx.ID_FILE6, _('Create array...'), _('New array (binary file - real)'))
        createarray2002 = self.menucreateobj.Append(wx.ID_ANY, _('Create array from Lidar 2002...'),
                                               _('Create array from Lidar 2002 (binary file - real)'))
        createarrayxyz = self.menucreateobj.Append(wx.ID_ANY, _('Create array from bathymetry file...'),
                                              _('Create array from XYZ (ascii file - real)'))
        createvector = self.menucreateobj.Append(wx.ID_FILE7, _('Create vectors...'), _('New vectors'))
        createview = self.menucreateobj.Append(wx.ID_ANY, _('Create view...'), _('New view'))
        createcloud = self.menucreateobj.Append(wx.ID_FILE8, _('Create cloud...'), _('New cloud'))
        createmanager2D = self.menucreateobj.Append(wx.ID_ANY, _('Create Wolf2D manager ...'), _('New manager 2D'))
        createscenario2D = self.menucreateobj.Append(wx.ID_ANY, _('Create scenarios manager ...'), _('New scenarios manager 2D'))
        createbcmanager2D = self.menucreateobj.Append(wx.ID_ANY, _('Create BC manager Wolf2D...'), _('New BC manager 2D'))
        createpartsystem = self.menucreateobj.Append(wx.ID_ANY, _('Create particle system...'), _('Create a particle system - Lagrangian view'))
        create_acceptability = self.menucreateobj.Append(wx.ID_ANY, _('Create acceptability manager...'), _('Create acceptability manager'))
        create_inbe = self.menucreateobj.Append(wx.ID_ANY, _('Create INBE manager...'), _('Create INBE manager'))
        createdrowning = self.menucreateobj.Append(wx.ID_ANY, _('Create a drowning...'),_('Create a drowning'))

        if WOLFPYDIKE_AVAILABLE:
            createdike = self.menucreateobj.Append(wx.ID_ANY, _('Create dike...'), _('New dike'))

        self.filemenu.AppendSeparator()


        # -----
        self.menuaddobj = wx.Menu()
        self.filemenu.Append(wx.ID_ANY,_('Add...'),self.menuaddobj)

        addarray = self.menuaddobj.Append(wx.ID_FILE1, _('Add array...'), _('Add array (binary file - real)'))
        addarraycrop = self.menuaddobj.Append(wx.ID_ANY, _('Add array and crop...'),
                                            _('Add array and crop (binary file - real)'))
        addvector = self.menuaddobj.Append(wx.ID_FILE2, _('Add vectors...'), _('Add vectors'))
        addpictcollection = self.menuaddobj.Append(wx.ID_ANY, _('Add picture collection...'), _('Add a collection of pictures'))
        addtiles = self.menuaddobj.Append(wx.ID_ANY, _('Add tiles...'), _('Add tiles'))
        addimagestiles = self.menuaddobj.Append(wx.ID_ANY, _('Add images tiles...'), _('Add georeferenced images tiles'))
        addtilescomp = self.menuaddobj.Append(wx.ID_ANY, _('Add tiles comparator...'), _('Add tiles comparator'))
        addtilesgpu = self.menuaddobj.Append(wx.ID_ANY, _('Add tiles GPU...'), _('Add tiles from 2D GPU model -- 2 arrays will be added'))
        addcloud = self.menuaddobj.Append(wx.ID_FILE3, _('Add cloud...'), _('Add cloud'))
        addtri = self.menuaddobj.Append(wx.ID_ANY, _('Add triangulation...'), _('Add triangulation'))
        addprofiles = self.menuaddobj.Append(wx.ID_FILE4, _('Add cross sections...'), _('Add cross sections'))
        addres2D = self.menuaddobj.Append(wx.ID_ANY, _('Add Wolf2D results...'), _('Add Wolf 2D results'))
        addres2Dgpu = self.menuaddobj.Append(wx.ID_ANY, _('Add Wolf2D GPU results...'), _('Add Wolf 2D GPU results'))
        addpartsystem = self.menuaddobj.Append(wx.ID_ANY, _('Add particle system...'), _('Add a particle system - Lagrangian view'))
        addbridges = self.menuaddobj.Append(wx.ID_ANY, _('Add bridges...'), _('Add bridges from directory'))
        addweirs = self.menuaddobj.Append(wx.ID_ANY, _('Add weirs...'), _('Add bridges from directory'))
        addview = self.menuaddobj.Append(wx.ID_ANY, _('Add view...'), _('Add view from project file'))
        adddrowning = self.menuaddobj.Append(wx.ID_ANY, _('Add a drowning result...'),_('Add a drowning result'))

        if WOLFPYDIKE_AVAILABLE:
            adddike = self.menuaddobj.Append(wx.ID_ANY, _('Add dike...'), _('Add dike'))

        self.precomputed_menu = None
        if self.default_dem != "":
            self.filemenu.AppendSeparator()
            self.precomputed_menu = wx.Menu()
            self.precomputed_menu.Append(wx.ID_ANY,_('Precomputed DEM'))
            self.filemenu.Append(wx.ID_ANY, _('Precomputed...'), self.precomputed_menu)
        if self.default_dtm != "":
            if self.precomputed_menu is None:
                self.filemenu.AppendSeparator()
                self.precomputed_menu = wx.Menu()
                self.filemenu.Append(wx.ID_ANY, _('Precomputed...'), self.precomputed_menu)

            self.precomputed_menu.Append(wx.ID_ANY,_('Precomputed DTM'))

        self.filemenu.AppendSeparator()
        addscan = self.filemenu.Append(wx.ID_FILE5, _('Recursive scan...'), _('Add recursively'))

        # Tools
        # ----------------

        self.tools_menu = wx.Menu()

        self.menu_contour_from_arrays = self.tools_menu.Append(wx.ID_ANY, _("Create contour from checked arrays..."), _("Create contour"))
        self.menu_calculator = self.tools_menu.Append(wx.ID_ANY, _("Calculator..."), _("Calculator"))
        self.menu_views = self.tools_menu.Append(wx.ID_ANY, _("Memory views..."), _("Memory views"))

        self.menu_distances = self.tools_menu.Append(wx.ID_ANY, _("Memory distances..."), _("Memory distances"))
        self.menu_distances_add = self.tools_menu.Append(wx.ID_ANY, _("Add distances to viewer..."), _("Add memory distances"))

        self.menu_digitizer = self.tools_menu.Append(wx.ID_ANY, _("Image digitizer..."), _("Image Digitizer"))
        self.calculator = None
        self.memory_views = None
        self._memory_views_gui = None

        # Cross sections
        # ----------------

        self.cs_menu = wx.Menu()
        self.link_cs_zones = self.cs_menu.Append(wx.ID_ANY, _("Link cross sections to active zones"),
                                                   _("Link cross section"))
        self.sortalong = self.cs_menu.Append(ID_SORTALONG, _("Sort along..."),
                                               _("Sort cross sections along support vector"))
        self.select_cs = self.cs_menu.Append(ID_SELECTCS, _("Pick one cross section"), _("Select cross section"),
                                               kind=wx.ITEM_CHECK)
        self.menumanagebanks = self.cs_menu.Append(wx.ID_ANY, _("Manage banks..."), _("Manage banks"))
        self.menucreatenewbanks = self.cs_menu.Append(wx.ID_ANY, _("Create banks from vertices..."),
                                                        _("Manage banks"))
        self.renamecs = self.cs_menu.Append(wx.ID_ANY, _("Rename cross sections..."), _("Rename"))
        self.menutrianglecs = self.cs_menu.Append(wx.ID_ANY, _("Triangulate cross sections..."), _("Triangulate"))
        self.menuexportgltfonebyone = self.cs_menu.Append(wx.ID_ANY, _("Export cross sections to gltf..."),
                                                            _("Export gltf"))
        self.menupontgltfonebyone = self.cs_menu.Append(wx.ID_ANY, _("Create bridge and export gltf..."),
                                                          _("Bridge gltf"))
        # self.menuimport3dfaces_from_DXF = self.toolsmenu.Append(wx.ID_ANY, _("Import triangulation..."), _("DXF"))

        #Profile plots
        #The action for plotting cross section's profile is initialised.
        self.plot_cs = self.cs_menu.Append(ID_PLOTCS, _("Plot cross section"),_("Plot cross section"),kind=wx.ITEM_CHECK)

        self.menuviewerinterpcs = None
        self.menuinterpcs = None

        # COLORMAP  menu
        self.minmaxmenu = wx.Menu()
        self.locminmax = self.minmaxmenu.Append(ID_LOCMINMAX, _("Local minmax"), _("Adapt colormap on current zoom"),
                                                kind=wx.ITEM_CHECK)
        paluniform= self.minmaxmenu.Append(wx.ID_ANY, _("Compute and apply unique colormap on all..."),
                                           _("Unique colormap"))
        paluniform_fomfile= self.minmaxmenu.Append(wx.ID_ANY, _("Load and apply unique colormap on all..."),
                                           _("Unique colormap"))
        paluniform_inparts= self.minmaxmenu.Append(wx.ID_ANY, _("Force uniform in parts on all..."),
                                           _("Uniform in parts"))
        pallinear= self.minmaxmenu.Append(wx.ID_ANY, _("Force linear interpolation on all..."),
                                           _("Linear colormap"))

        self.analyzemenu = wx.Menu()

        self.analyzeplot = wx.Menu()
        self.analyzeexport = wx.Menu()
        self.analyzeinpaint = wx.Menu()
        self.analyzesimsheet = wx.Menu()

        plotvect = self.analyzeplot.Append(wx.ID_ANY, _("Plot active vector..."),
                                           _("Plot the active vector and linked arrays"))
        plotpoly = self.analyzeplot.Append(wx.ID_ANY, _("Plot active polygons..."),
                                           _("Plot the active polygons and linked arrays"))

        self.analyzeplot.AppendSeparator()

        self.analyzemenu.Append(wx.ID_ANY,_('Plot...'), self.analyzeplot)
        self.analyzemenu.Append(wx.ID_ANY,_('Export...'), self.analyzeexport)
        self.analyzemenu.Append(wx.ID_ANY,_('Inpaint...'), self.analyzeinpaint)

        self.analyzemenu.AppendSeparator()

        self.analyzemenu.Append(wx.ID_ANY,_('Report...'), self.analyzesimsheet)

        self.analyzesimsheet.Append(wx.ID_ANY, _("Active simulation..."), _("Generate a summary PDF report for the active simulation"))
        self.analyzesimsheet.Append(wx.ID_ANY, _("All checked simulations..."), _("Generate a summary PDF report for all checked simulations"))
        self.analyzesimsheet.AppendSeparator()
        self.analyzesimsheet.Append(wx.ID_ANY, _("One simulation from disk..."), _("Generate a summary PDF report for one simulation"))
        self.analyzesimsheet.Append(wx.ID_ANY, _("All simulations in directory..."), _("Generate a summary PDF report for all simulations in the current directory"))
        self.analyzesimsheet.AppendSeparator()
        self.analyzesimsheet.Append(wx.ID_ANY, _("Compare checked simulations..."), _("Generate a summary PDF report for all the loaded simulations"))
        self.analyzesimsheet.Append(wx.ID_ANY, _("Compare all simulations in a directory..."), _("Generate a summary PDF report for all the simulations in a directory"))
        self.analyzesimsheet.AppendSeparator()
        self.analyzesimsheet.Append(wx.ID_ANY, _("Compare arrays..."), _("Generate a summary PDF report for two loaded arrays"))
        self.analyzesimsheet.Append(wx.ID_ANY, _("Compare arrays from files..."), _("Generate a summary PDF report for two arrays from files"))


        self.analyzeinpaint.Append(wx.ID_ANY, _("Inpaint active array..."), _("Inpaint active array"))
        self.analyzeinpaint.Append(wx.ID_ANY, _("Inpaint waterlevel..."), _("Inpaint a waterlevel result array based on sepcified dem and dtm data"))
        self.analyzeinpaint.Append(wx.ID_ANY, _("Inpaint array with mask..."), _("Inpaint an array based on sepcified mask and test data"))

        self.analyzemenu.AppendSeparator()

        masksimul = self.analyzemenu.Append(wx.ID_ANY, _("Load and apply mask (nap)..."), _("Apply mask from sim2D"))
        filterinund = self.analyzemenu.Append(wx.ID_ANY, _("Filter inundation arrays..."), _("Filter arrays"))

        # Plot hydrographs

        plotqvect = self.analyzeplot.Append(wx.ID_ANY, _("Plot integrated Q along active vector..."), _("Integrate Q along the active vector and plot"))
        plotqvect = self.analyzeplot.Append(wx.ID_ANY, _("Plot integrated Q along active zone..."), _("Integrate Q along the active zone and plot"))

        self.analyzeplot.AppendSeparator()


        exportqvect = self.analyzeexport.Append(wx.ID_ANY, _("Export integrated Q along active vector..."), _("Integrate Q along the active vector and export"))
        exportqvect = self.analyzeexport.Append(wx.ID_ANY, _("Export integrated Q along all vectors in active zone..."), _("Integrate Q along ALL VECTORS of the active zone and export"))

        self.analyzeexport.AppendSeparator()

        plothselect = self.analyzeplot.Append(wx.ID_ANY, _("Plot stats unknown (selected nodes)..."), _("Compute stats and plot on the selected nodes"))
        plothvector = self.analyzeplot.Append(wx.ID_ANY, _("Plot stats unknown (inside active vector)..."), _("Compute stats and plot on nodes inside the active vector"))
        plothzone = self.analyzeplot.Append(wx.ID_ANY, _("Plot stats unknown (inside active zone)..."), _("Compute stats and plot on nodes inside the active zone"))

        exporthselect = self.analyzeexport.Append(wx.ID_ANY, _("Export stats unknown (selected nodes)..."), _("Compute stats and export on the selected nodes"))
        exporthvector = self.analyzeexport.Append(wx.ID_ANY, _("Export stats unknown (inside active vector)..."), _("Compute stats and export on nodes inside the active vector"))
        exporthzone = self.analyzeexport.Append(wx.ID_ANY, _("Export stats unknown (inside active zone)..."), _("Compute stats and export on nodes inside the active zone"))

        self.filemenu.AppendSeparator()
        menuquit = self.filemenu.Append(wx.ID_EXIT, _('&Quit\tCTRL+Q'), _('Quit application'))

        # If one uses the accelerator key then it is tied to the
        # wx.ID_EXIT. Moreover the accelerator key will be shadowed by
        # EVT_CHAR_HOOK if one is not careful.  Note that using
        # accelerators on anything else then wx.EVT_MENU is reported
        # hackish at best on the WWW.

        accel_tbl = wx.AcceleratorTable([(wx.ACCEL_CTRL,  ord('Q'), menuquit.GetId() )])
        self.SetAcceleratorTable(accel_tbl)

        # Gestion des outils --> Utile pour ManageActions
        self.tools = {}
        curtool = self.tools[ID_SELECTCS] = {}
        curtool['menu'] = self.select_cs
        curtool['name'] = 'Select nearest profile'

        curtool = self.tools[ID_PLOTCS] = {}
        curtool['menu'] = self.plot_cs
        curtool['name'] = 'Plot cross section'

        self.mybc = []

        self.active_vector = None
        self.active_zone = None
        self.active_zones = None
        self.active_vertex = None
        self.active_array = None
        self.active_bc = None
        self.active_tri = None
        self.active_cloud = None
        self.active_view = None
        self.active_cs = None
        self.active_profile = None
        self.active_res2d = None
        self.active_particle_system = None
        self.active_viewer3d = None
        self.active_viewerlaz = None
        self.active_landmap:PlansTerrier = None
        self.active_tile = None
        self.selected_treeitem = None
        self.selected_object = None
        self.active_bridges = None
        self.active_bridge = None
        self.active_weirs = None
        self.active_weir = None
        self.active_laz = None
        self.active_dike = None
        self.active_injector = None
        self.active_picturecollection = None
        self.active_qdfidf = None
        self.active_alaro = None

        self.active_fig = None
        self.active_drowning = None

        curtool = self.tools[ID_SORTALONG] = {}
        curtool['menu'] = self.sortalong
        curtool['name'] = 'Sort along vector'

        curtool = self.tools[ID_LOCMINMAX] = {}
        curtool['menu'] = self.locminmax
        curtool['name'] = None

        self.menubar.Append(self.filemenu, _('&File'))

        # Help
        self.helpmenu = wx.Menu()
        self.helpmenu.Append(wx.ID_ANY, _('Shortcuts'), _('Shortcuts'))
        self.helpmenu.Append(wx.ID_ANY, _('Project .proj'), _('A project file ".proj", what is it?'))
        self.helpmenu.Append(wx.ID_ANY, _('Show logs/informations'), _('Logs'))
        self.helpmenu.Append(wx.ID_ANY, _('Show values'), _('Data/Values'))
        self.helpmenu.Append(wx.ID_ANY, _('About'), _('About'))
        self.helpmenu.Append(wx.ID_ANY, _('Check for updates'), _('Update?'))

        self.menubar.Append(self.helpmenu, _('&Help'))

        # ajout du menu pour les données LAZ
        self.menu_laz()

        self.menubar.Append(self.tools_menu, _('&Tools'))
        self.menubar.Append(self.cs_menu, _('&Cross sections'))

        self.menubar.Append(self.minmaxmenu, _('&Colormap'))
        self.menubar.Append(self.analyzemenu, _('&Analyze'))
        self.SetMenuBar(self.menubar)
        self.Bind(wx.EVT_MENU, self.OnMenubar)
        self.Bind(wx.EVT_MENU_HIGHLIGHT, self.OnMenuHighlight)

        # Ajout du conteneur OpenGL
        self.canvas = GLCanvas(self)
        self.canvas.SetDropTarget(self._dragdrop)

        self.context = GLContext(self.canvas)
        self.mybackisloaded = False
        self.myfrontisloaded = False

        # ajout d'une liste en arbre des objets
        self.treelist = TreeListCtrl(self, style= wx.dataview.TL_CHECKBOX | wx.LC_EDIT_LABELS | wx.TR_FULL_ROW_HIGHLIGHT)
        self._lbl_selecteditem = StaticText(self, style=wx.ALIGN_CENTER_HORIZONTAL | wx.ALIGN_CENTER_VERTICAL)

        self.root = self.treelist.GetRootItem()
        self.treelist.AppendColumn(_('Objects to plot'))
        self.myitemsarray  = self.treelist.AppendItem(self.root, _("Arrays"))
        self.myitemsvector = self.treelist.AppendItem(self.root, _("Vectors"))
        self.myitemscloud  = self.treelist.AppendItem(self.root, _("Clouds"))
        self.myitemslaz    = self.treelist.AppendItem(self.root, _("Laz"))
        self.myitemstri    = self.treelist.AppendItem(self.root, _("Triangulations"))
        self.myitemsres2d  = self.treelist.AppendItem(self.root, _("Wolf2D"))
        self.myitemsps     = self.treelist.AppendItem(self.root, _("Particle systems"))
        self.myitemsothers = self.treelist.AppendItem(self.root, _("Others"))
        self.myitemsviews  = self.treelist.AppendItem(self.root, _("Views"))
        self.myitemswmsback = self.treelist.AppendItem(self.root, _("WMS-background"))
        self.myitemswmsfore = self.treelist.AppendItem(self.root, _("WMS-foreground"))
        self.myitemsdrowning = self.treelist.AppendItem(self.root,_("Drowning"))
        self.myitemsdike    = self.treelist.AppendItem(self.root, _("Dikes"))
        self.myitemsinjector= self.treelist.AppendItem(self.root, _("Injectors"))
        self.myitemspictcollection = self.treelist.AppendItem(self.root, _("Pictures"))

        width, height = self.GetClientSize()
        self.bordersize = int((w - width + self.treewidth) / 2)
        self.titlesize = h - height - self.bordersize
        self.SetSize(w + 2 * self.bordersize + self.treewidth, h + self.bordersize + self.titlesize)

        # dimensionnement et positionnement de la fenêtre OpenGL
        self.canvas.SetSize(width - self.treewidth, height)
        self.canvas.SetPosition((self.treewidth, 0))

        self.setbounds()

        # dimensionnement et positionnement de l'arbre
        self.leftbox = BoxSizer(orient=wx.VERTICAL)
        self.leftbox.Add(self.treelist, 1, wx.LEFT)
        self.leftbox.Add(self._lbl_selecteditem, 0, wx.LEFT)
        self.treelist.SetSize(self.treewidth, height)


        self.CreateStatusBar(1)

        self.SetSizer(self.leftbox)

        # self.treelist.SetPosition((0,0))

        # fenêtre ToolTip
        self.mytooltip = Wolf_Param(self, _("Data/Results"), to_read=False, withbuttons=False, toolbar=False, DestroyAtClosing=False)
        self.mytooltip.SetSize(300, 400)
        self.mytooltip.prop.SetDescBoxHeight(20) # Hauteur de la zone de description
        self.mytooltip.Show(True)
        self._oldpos_tooltip = None

        #Notebooks
        self.notebookcs = None
        self.notebookprof = None
        self.notebookbanks = None

        #Axes
        self.myaxcs = None
        self.myaxprof = None

        #Figures
        self.myfigcs = None
        self.myfigprof = None

        self.cloudmenu=None
        self.trianglesmenu = None

        self._configuration = None

        self.compare_results = None

        self.InitUI()

        self._tmp_vector_distance = None  # distance computation the vector
        self._distances = Zones(mapviewer=self, idx=_('Distances/Areas'), parent = self)
        self._distances.add_zone(zone(name='memory distances', parent=self._distances))  # distances memory

        self.menu_alaro_forecasts()

        # self._wintab = Wintab(self.GetHandle())

        # if self._wintab:

        #     import win32gui
        #     import win32con

        #     # self.oldWndProc = win32gui.SetWindowLong(self.GetHandle(), win32con.GWL_WNDPROC, self.MyWndProc)

    # def MyWndProc(self, hWnd, msg, wParam, lParam):
    #     import win32con

    #     # Intercept a specific Windows message (for example, WM_KEYDOWN)
    #     # if msg == 0x7FF0:
    #     #     key_code = wParam
    #     #     print(f"Key pressed: {key_code}")

    #     #     # Process the message or do something custom
    #     #     if key_code == win32con.VK_ESCAPE:
    #     #         print("Escape key pressed, intercepting the event.")

    #             # # You can return 0 to indicate the message has been processed
    #             # return 0
    #     # print(msg)
    #     return 0

    def _check_id_for_fig(self, idx:str):
        """ Check if an ID is already used for a figure """

        ids = [cur.idx for cur in self.mymplfigs]

        if idx in ids:
            return True

        return False

    def _create_id_for_fig(self):

        idx = 'Figure'
        while not self._check_id_for_fig(idx):
            idx += '_'

        return idx

    def _validate_id_for_fig(self, idx:str):
        """ Validate an ID for a figure """

        if idx is None:
            return self._create_id_for_fig()

        while self._check_id_for_fig(idx):
            idx += '_'

        return idx

    def new_fig(self, caption:str, idx:str = None, layout = PRESET_LAYOUTS.DEFAULT, size = (800,600), show:bool = True) -> MplFigViewer:
        """ Create a new figure """

        if idx is None:
            with wx.TextEntryDialog(self, _('Enter an id for the figure'), _('Figure id'), _('Figure')) as dlg:
                dlg: wx.TextEntryDialog
                if dlg.ShowModal() == wx.ID_CANCEL:
                    return None

                idx = dlg.GetValue()

                idx = self._validate_id_for_fig(idx)
        else:
            idx = self._validate_id_for_fig(idx)

        logging.info(f'Figure ID: {idx}')

        added_fig = MplFigViewer(layout, idx= idx, mapviewer = self, caption = caption, size= size)

        if show:
            added_fig.Show()
        else:
            added_fig.Hide()

        self.mymplfigs.append(added_fig)

        return added_fig

    def destroy_fig_by_id(self, idx:str) -> bool:
        """ Destroy a figure by its ID """

        for id, fig in enumerate(self.mymplfigs):
            if fig.idx == idx:
                if self.active_fig is fig:
                    self.active_fig = None
                fig.Destroy()
                self.mymplfigs.pop(id)
                return True

        return False

    def get_fig(self, idx:str) -> MplFigViewer:
        """ Get a figure by its ID """

        for cur in self.mymplfigs:
            if cur.idx == idx:
                return cur

        return None

    def list_ids_figs(self) -> list[str]:
        """ List all IDs of figures """

        return [cur.idx for cur in self.mymplfigs]

    @property
    def viewer_name(self):
        return self.GetTitle()

    @viewer_name.setter
    def viewer_name(self, value):
        self.SetTitle(value)

    @property
    def wxlogging(self):
        return self._wxlogging

    @wxlogging.setter
    def wxlogging(self, value):
        self._wxlogging = value

    def check_logging(self):
        """ Check if logging window is shown """

        if self._wxlogging is None:
            logging.info(_('No logging window'))
            return

        self._wxlogging.Show()

    def check_tooltip(self):
        """ Check if tooltip window is shown """

        if self.mytooltip is None:
            logging.info(_('No tooltip window'))
            return

        self.mytooltip.Show()


    def open_hydrological_model(self):
        """ Open a hydrological model """

        from .PyGui import HydrologyModel

        newview = HydrologyModel(splash = False)

    def create_2D_MB_model(self):
        """ Create a 2D model """

        from .PyGui import Wolf2DModel

        newview = Wolf2DModel(splash = False)

    def create_2D_GPU_model(self):
        """ Create a 2D GPU model """

        from .PyGui import Wolf2DGPUModel

        newview = Wolf2DGPUModel(splash = False)

    def check_2D_MB_headers(self):
        """ Check headers of a 2D simulation without opening viewer"""

        # Check 2D simulation
        dlg = wx.FileDialog(self, _("Choose 2D simulation file"), wildcard="all (*.*)|*.*", style=wx.FD_OPEN)
        if dlg.ShowModal() == wx.ID_CANCEL:
            dlg.Destroy()
            return

        filename = dlg.GetPath()
        dlg.Destroy()

        from .mesh2d.wolf2dprev import prev_sim2D

        sim = prev_sim2D(filename)
        sim.verify_files()


    def get_mapviewer(self):
        """ Retourne une instance WolfMapViewer """
        return self

    def do_quit(self):
        pass

    def create_triangles_menu(self):
        """ Menu for triangulations """

        if self.trianglesmenu is None:

            self.trianglesmenu = wx.Menu()
            self.menubar.Append(self.trianglesmenu, _('&Triangulations'))

            self._menuinteractptri = self.trianglesmenu.Append(wx.ID_ANY, _("Interpolate on active triangulation..."),
                                                               _('Interpolate active array on active triangulation'))

            self._menuinteractptri_above = self.trianglesmenu.Append(wx.ID_ANY, _("Interpolate on active triangulation (keep only above)..."),
                                                                     _('Interpolate active array on active triangulation but keep only values above the array'))

            self._menuinteractptri_below = self.trianglesmenu.Append(wx.ID_ANY, _("Interpolate on active triangulation (keep only below)..."),
                                                                     _('Interpolate active array on active triangulation but keep only values below the array'))

            self._menucomparetri = self.trianglesmenu.Append(wx.ID_ANY, _("Compare triangles to array..."), _("Comparison"))
            self._menumovetri = self.trianglesmenu.Append(wx.ID_ANY, _("Move triangles..."), _("Move triangles"))
            self._menurotatetri = self.trianglesmenu.Append(wx.ID_ANY, _("Rotate triangles..."), _("Rotate triangles"))

    def create_cloud_menu(self):
        """ Menu for cloud points """

        if self.cloudmenu is None:
            self.cloudmenu = wx.Menu()
            self.menubar.Append(self.cloudmenu, _('Cloud'))

            interpcloudonarray = self.cloudmenu.Append(wx.ID_ANY, _("Interpolate active cloud on active array..."),_("Interpolation"))
            self._menucomparecloud = self.cloudmenu.Append(wx.ID_ANY, _("Compare cloud to array..."), _("Comparison"))
            split_cloud = self.cloudmenu.Append(wx.ID_ANY, _("Split cloud..."), _("Split cloud"))

    def split_cloud_by_vector(self):
        """ Split cloud by vector """

        inside_cloud, outside_cloud = self.active_vector.split_cloud(self.active_cloud)

        self.add_object('cloud', newobj = inside_cloud, id = inside_cloud.idx)
        self.add_object('cloud', newobj = outside_cloud, id = outside_cloud.idx)


    def get_choices_arrays(self):
        """Boîte de dialogue permettant de choisir une ou plusieurs matrices parmi celles chargées"""

        dlg = wx.MultiChoiceDialog(self,_('Choose one or multiple arrays'),
                                    _('Choose'),choices=[cur.idx for cur in self.myarrays])
        ret = dlg.ShowModal()

        if ret == wx.ID_CANCEL:
            return None

        mychoices = self.myarrays[dlg.GetSelections()]
        return mychoices

    def menu_tiles(self):
        """ Menu for tiles """
        if self.menutiles is None:
            self.menutiles = wx.Menu()
            self.menubar.Append(self.menutiles, _('&Tiles'))

            picktiles = self.menutiles.Append(wx.ID_ANY, _('Pick a tile and load data'), _('Right click to pick a tile'))
            data_active_polygon_tiles = self.menutiles.Append(wx.ID_ANY, _('Select data within the active polygonal area'), _('Select data within the bouding box of the active polygonal area'))
            data_tmpvec_tiles = self.menutiles.Append(wx.ID_ANY, _('Select data within a temporary polygonal area'), _('Right click to add points + Enter'))

            self.Bind(wx.EVT_MENU, self.pick_tile, picktiles)
            self.Bind(wx.EVT_MENU, self.create_data_from_tiles_activevec, data_active_polygon_tiles)
            self.Bind(wx.EVT_MENU, self.create_data_from_tiles_tmpvec, data_tmpvec_tiles)

    def pîck_image_tile(self, event: wx.Event):

        if self.active_imagestiles is None:
            logging.warning(_('No active image tile -- Please load data first'))
            return

        self.action = 'select active image tile'
        logging.info(_('Select active image tile'))

    def menu_pictcollection(self):
        """ Menu for picture collections """

        if self.menupicturecollections is None:
            self.menupicturecollections = wx.Menu()
            self.menubar.Append(self.menupicturecollections, _('&Pictures'))

            scaleall = self.menupicturecollections.Append(wx.ID_ANY, _('Scale all pictures'), _('Scale all pictures in the collection'))
            pick = self.menupicturecollections.Append(wx.ID_ANY, _('Pick a picture'), _('Right click to pick a picture'))
            hide = self.menupicturecollections.Append(wx.ID_ANY, _('Hide all pictures'), _('Reset the picture collection'))
            allvisible = self.menupicturecollections.Append(wx.ID_ANY, _('Show all pictures'), _('Set all pictures in the collection visible'))
            extract = self.menupicturecollections.Append(wx.ID_ANY, _('Extract pictures'), _('Extract all visible pictures from the collection'))

            self.Bind(wx.EVT_MENU, self.action_pictcollections, pick)
            self.Bind(wx.EVT_MENU, self.action_pictcollections, scaleall)
            self.Bind(wx.EVT_MENU, self.action_pictcollections, hide)
            self.Bind(wx.EVT_MENU, self.action_pictcollections, allvisible)
            self.Bind(wx.EVT_MENU, self.action_pictcollections, extract)

    def menu_qdfidf(self):
        """ Menu for QDF/IDF files """

        if self.menuqdfidf is None:
            self.menuqdfidf = wx.Menu()
            self.menubar.Append(self.menuqdfidf, _('QDF/IDF'))

            download = self.menuqdfidf.Append(wx.ID_ANY, _('Download data'), _('Download all QDF/IDF data for Belgium and preprocess it !'))
            load = self.menuqdfidf.Append(wx.ID_ANY, _('Load data'), _('Load a QDF/IDF for Belgium'))
            pick = self.menuqdfidf.Append(wx.ID_ANY, _('Pick municipality'), _('Pick a municipality and show its data'))
            show_tables = self.menuqdfidf.Append(wx.ID_ANY, _('Show tables'), _('Show tables on viewer'))
            show_plots = self.menuqdfidf.Append(wx.ID_ANY, _('Show plots'), _('Show plots on viewer'))
            scale = self.menuqdfidf.Append(wx.ID_ANY, _('Scale images'), _('Scale the QDF/IDF images'))

            self.Bind(wx.EVT_MENU, self.action_qdfidf, download)
            self.Bind(wx.EVT_MENU, self.action_qdfidf, load)
            self.Bind(wx.EVT_MENU, self.action_qdfidf, pick)
            self.Bind(wx.EVT_MENU, self.action_qdfidf, show_tables)
            self.Bind(wx.EVT_MENU, self.action_qdfidf, show_plots)
            self.Bind(wx.EVT_MENU, self.action_qdfidf, scale)

    def action_qdfidf(self, event: wx.Event):
        """ Action for QDF/IDF files """

        item = event.GetEventObject().FindItemById(event.GetId())
        itemlabel = item.ItemLabel

        if self.active_qdfidf is None and itemlabel != _('Load data') and itemlabel != _('Download data'):
            logging.warning(_('No active QDF/IDF -- Please load data first'))
            return

        if itemlabel == _('Load data'):

            from .irm_qdf import QDF_Hydrology_Draw

            dlg_dir = wx.DirDialog(self, _('Choose directory with QDF/IDF files'), style=wx.DD_DEFAULT_STYLE)
            if dlg_dir.ShowModal() == wx.ID_CANCEL:
                dlg_dir.Destroy()
                return
            dirpath = dlg_dir.GetPath()
            dlg_dir.Destroy()

            pgbar = wx.ProgressDialog(_('Loading QDF/IDF data'), _('Loading data...'), maximum=100, parent=self, style=wx.PD_APP_MODAL | wx.PD_AUTO_HIDE)
            pgbar.Pulse(_('Loading QDF/IDF data... (estimated time 10 min.)'))

            self.active_qdfidf = QDF_Hydrology_Draw(dirpath, idx = 'QDF/IDF', mapviewer = self)

            pgbar.Update(100, _('QDF/IDF data loaded'))
            pgbar.Destroy()

            self.treelist.SetFocus()

        elif itemlabel == _('Download data'):

            from .irm_qdf import QDF_Hydrology_Draw

            dlg_dir = wx.DirDialog(self, _('Choose an empty directory to store QDF/IDF files'), style=wx.DD_DEFAULT_STYLE | wx.DD_NEW_DIR_BUTTON)
            if dlg_dir.ShowModal() == wx.ID_CANCEL:
                dlg_dir.Destroy()
                return

            dirpath = dlg_dir.GetPath()
            dlg_dir.Destroy()

            # Test if the directory is empty
            if os.listdir(dirpath):
                logging.error(_('The directory {} is not empty. Please choose an empty directory.').format(dirpath))
                wx.MessageBox(_('The directory {} is not empty. Please choose an empty directory.').format(dirpath), _('Error'), wx.OK | wx.ICON_ERROR)
                return

            pgbar = wx.ProgressDialog(_('Downloading QDF/IDF data'), _('Downloading data...'), maximum=100, parent=self, style=wx.PD_APP_MODAL | wx.PD_AUTO_HIDE)
            pgbar.Pulse(_('Downloading QDF/IDF data...'))

            try:
                self.active_qdfidf = QDF_Hydrology_Draw(dirpath, idx = 'QDF/IDF', mapviewer = self)
            except Exception as e:
                logging.error(_('Error downloading QDF/IDF data: {}').format(e))
                wx.MessageBox(_('Error downloading QDF/IDF data: {}').format(e), _('Error'), wx.OK | wx.ICON_ERROR)

            pgbar.Update(100, _('QDF/IDF data loaded and processed'))
            pgbar.Destroy()

            self.treelist.SetFocus()

        elif itemlabel == _('Pick municipality'):
            self.action = 'pick municipality'
            logging.info(_('Pick a municipality from the QDF/IDF data'))

        elif itemlabel == _('Show tables'):

            if self.active_qdfidf is not None:

                pgbar = wx.ProgressDialog(_('Loading QDF/IDF data'), _('Loading data...'), maximum=100, parent=self, style=wx.PD_APP_MODAL | wx.PD_AUTO_HIDE)
                pgbar.Pulse(_('Loading QDF/IDF data...'))

                self.active_qdfidf.show_plot = False
                self.active_qdfidf.show_table = True
                logging.info(_('Show tables for QDF/IDF data'))
                self.Refresh()

                pgbar.Update(100, _('QDF/IDF data loaded'))
                pgbar.Destroy()
                self.treelist.SetFocus()


        elif itemlabel == _('Show plots'):

            if self.active_qdfidf is not None:

                pgbar = wx.ProgressDialog(_('Loading QDF/IDF data'), _('Loading data...'), maximum=100, parent=self, style=wx.PD_APP_MODAL | wx.PD_AUTO_HIDE)
                pgbar.Pulse(_('Loading QDF/IDF data...'))

                self.active_qdfidf.show_table = False
                self.active_qdfidf.show_plot = True
                logging.info(_('Show tables for QDF/IDF data'))
                self.Refresh()

                pgbar.Update(100, _('QDF/IDF data loaded'))
                pgbar.Destroy()
                self.treelist.SetFocus()


        elif itemlabel == _('Scale images'):
            if self.active_qdfidf is not None:
                scalefactor = wx.GetTextFromUser(_('Enter the scale factor (default is 1.0)'), _('Scale factor'), '1.0', self)
                if scalefactor == '':
                    scalefactor = '1.0'
                try:
                    scalefactor = float(scalefactor)
                except ValueError:
                    logging.error(_('Invalid scale factor: {}').format(scalefactor))
                    wx.MessageBox(_('Invalid scale factor: {}').format(scalefactor), _('Error'), wx.OK | wx.ICON_ERROR)
                    return

            self.active_qdfidf.scale_images(scalefactor)

            self.Refresh()

    def action_pictcollections(self, event: wx.Event):
        """ Action for picture collections """

        if self.active_picturecollection is None:
            logging.warning(_('No active picture collection -- Please load data first'))
            return

        item = event.GetEventObject().FindItemById(event.GetId())
        itemlabel = item.ItemLabel

        if itemlabel == _('Pick a picture'):
            self.action = 'pick a picture'
            logging.info(_('Pick a picture from the collection'))

        elif itemlabel == _('Scale all pictures'):

            scalefactor = wx.GetTextFromUser(_('Enter the scale factor (default is 1.0)'), _('Scale factor'), '1.0', self)
            if scalefactor == '':
                scalefactor = '1.0'
            try:
                scalefactor = float(scalefactor)
            except ValueError:
                logging.error(_('Invalid scale factor: {}').format(scalefactor))
                wx.MessageBox(_('Invalid scale factor: {}').format(scalefactor), _('Error'), wx.OK | wx.ICON_ERROR)
                return

            self.active_picturecollection.scale_all_pictures(scalefactor)

            self.Refresh()

        elif itemlabel == _('Hide all pictures'):
            self.active_picturecollection.hide_all_pictures()
            self.Refresh()

        elif itemlabel == _('Show all pictures'):
            self.active_picturecollection.show_all_pictures()
            self.Refresh()

        elif itemlabel == _('Extract pictures'):

            dlg = wx.DirDialog(self, _('Choose directory to extract picture collection'), style= wx.DD_DEFAULT_STYLE)
            if dlg.ShowModal() == wx.ID_CANCEL:
                dlg.Destroy()

            else:
                dirpath = dlg.GetPath()
                dlg.Destroy()
                self.active_picturecollection.extract_pictures(dirpath)
                logging.info(_('Pictures extracted to {}').format(dirpath))

    def menu_imagestiles(self):
        """ Menu for image tiles """
        if self.menuimagestiles is None:
            self.menuimagestiles = wx.Menu()
            self.menubar.Append(self.menuimagestiles, _('&Image tiles'))

            picktiles = self.menuimagestiles.Append(wx.ID_ANY, _('Pick a tile and (un)load data'), _('Right click to pick a tile'))
            self.Bind(wx.EVT_MENU, self.pîck_image_tile, picktiles)


    def pick_tile(self, event: wx.Event):

        if self.active_tile is None:
            logging.warning(_('No active tile -- Please load data first'))
            return

        self.action = 'select active tile'
        logging.info(_('Select active tile'))

    def create_data_from_tiles_activevec(self, event: wx.Event):

        if self.active_tile is None:
            logging.warning(_('No active tile -- Please load data first'))
            return

        if self.active_vector is None:
            logging.warning(_('No active vector -- Please activate a vector first'))
            return

        self._create_data_from_tiles_common()

    def _create_data_from_tiles_common(self):

        from .wolf_vrt import create_vrt, crop_vrt

        dirdata = self.active_tile.linked_data_dir

        glob_vrt = glob.glob(join(dirdata,'*.vrt'))

        if len(glob_vrt) == 0:
            file_vrt = r'tmp.vrt'
            create_vrt(dirdata, fout=file_vrt)
            glob_vrt = file_vrt
        else:
            glob_vrt = glob_vrt[0]

        dlg = wx.FileDialog(None, _('Choose filename'), wildcard='tif (*.tif)|*.tif', defaultDir=dirdata, defaultFile='{}_crop.tif'.format(self.active_vector.myname), style=wx.FD_SAVE)
        ret = dlg.ShowModal()
        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return
        fout = dlg.GetPath()
        dlg.Destroy()

        bbox = self.active_vector.get_bounds_xx_yy()

        crop_vrt(glob_vrt, bbox, fout=fout)
        logging.info(_('File {} created').format(fout))

        dlg = wx.MessageDialog(self, _('Do you want to load the created file ?'), _('Load file'), wx.YES_NO | wx.ICON_QUESTION)
        ret = dlg.ShowModal()
        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return
        elif ret == wx.ID_YES:
            self.add_object('array', filename=fout, id = self.active_vector.myname)

    def create_data_from_tiles_tmpvec(self, event: wx.Event):

        if self.active_tile is None:
            logging.warning(_('No active tile -- Please load data first'))
            return

        self.start_action('create polygon - tiles', _('Extract data from tiles inside polygon'))

        self.active_vector = vector()
        self.active_vector.myname = 'crop_tiles'
        self.active_vector.add_vertex(wolfvertex(0.,0.))

    def menu_laz(self):
        """ Menu for LAZ Data """
        if self.menulaz is None:

            self.menulaz = wx.Menu()
            self.menulazdata = wx.Menu()
            self.menulazgrid = wx.Menu()

            self.menubar.Append(self.menulaz, _('&LAZ'))
            self.menulaz.AppendSubMenu(self.menulazdata, _('LAZ data'))
            self.menulaz.AppendSubMenu(self.menulazgrid, _('LAZ grid'))

            readlaz = self.menulazdata.Append(wx.ID_ANY, _('Initialize from laz, las or npz'), _('LAZ data from one specific file (type laz, las or npz)'))
            readlaz_gridinfo = self.menulazgrid.Append(wx.ID_ANY, _('Initialize from directory'), _('LAZ GRID stored in a directory - subgridding of LAZ files'), kind=wx.ITEM_CHECK)
            copylaz_gridinfo = self.menulazgrid.Append(wx.ID_ANY, _('Copy from current zoom'), _('Make a copy of the used files'))
            updatecolors_laz = self.menulazgrid.Append(wx.ID_ANY, _('Change colors - Classification'), _('Change color map associated to the current classification'),)

            bridgelaz = self.menulazdata.Append(wx.ID_ANY, _('Create cloud points from bridges'), _('Extract bridge from LAZ data as cloud points (class 10)'))
            buildinglaz = self.menulazdata.Append(wx.ID_ANY, _('Create cloud points from buildings'), _('Extract buildings from LAZ data as cloud points (class 1)'))
            classlaz = self.menulazdata.Append(wx.ID_ANY, _('Create cloud points from specified classes'), _('Extract cloud points from LAZ data (class to specify)'))

            croplaz = self.menulaz.Append(wx.ID_ANY, _('Clip LAZ grid on current zoom'), _('Select LAZ data based on the visible screen extent'),)
            viewlaz = self.menulaz.Append(wx.ID_ANY, _('Create LAZ viewer'), _('Create a LAZ Viewer based on loaded data'))
            filterlaz = self.menulaz.Append(wx.ID_ANY, _('Filter data based on codes'), _('Filter LAZ data based on codes'),)
            descimate_laz = self.menulaz.Append(wx.ID_ANY, _('Descimate LAZ data'), _('Descimate LAZ data'),)
            aroundlaz = self.menulaz.Append(wx.ID_ANY, _('Plot LAZ around active vector'), _('Display a Matplotlib plot with the LAZ values around the active vector/polyline'),)
            pick_aroundlaz = self.menulaz.Append(wx.ID_ANY, _('Plot LAZ around temporary vector'), _('Display a Matplotlib plot with the LAZ values around a temporary vector/polyline -- Right clicks to add points + Enter'),)
            fillarray_laz = self.menulaz.Append(wx.ID_ANY, _('Fill active array from LAZ data'), _('Fill an array from the LAZ data'),)
            selectarray_laz = self.menulaz.Append(wx.ID_ANY, _('Select cells in array from LAZ data'), _('Select nodes in active array from the LAZ data'),)
            countarray_laz = self.menulaz.Append(wx.ID_ANY, _('Count LAZ data in cells'), _('Count the number of LAZ data in each cell of the matrix'),)

    def menu_wolf2d(self):

        if self.menuwolf2d is None:
            self.menuwolf2d = wx.Menu()

            self.menu2d_explore_results = self.menuwolf2d.Append(wx.ID_ANY, _("Explore time/index results"), _("Open a dialog to explore time/index results"))

            self.menuwolf2d.AppendSeparator()

            self.menu2d_curentview = self.menuwolf2d.Append(wx.ID_ANY, _("Change current view"), _("Current view"))
            self.menu2d_lastres = self.menuwolf2d.Append(wx.ID_ANY, _("Read last result"), _("Current view"))
            self.menu2d_epsilon = self.menuwolf2d.Append(wx.ID_ANY, _("Set epsilon water depth"), _("Set the epsilon used in the mask"))

            self.menu_filter_independent = self.menuwolf2d.Append(wx.ID_ANY, _("Filter independent"), _("Filter independent"), kind=wx.ITEM_CHECK)

            # self.menu2d_bc = self.menuwolf2d.Append(wx.ID_ANY, _("Manage boundary conditions..."), _("BC manager"))
            self.menu2d_video = self.menuwolf2d.Append(wx.ID_ANY, _("Create video..."), _("Video/Movie"))

            self.menuwolf2d.AppendSeparator()

            self.menu2d_dangermap = self.menuwolf2d.Append(wx.ID_ANY, _("Danger map"), _("Compute the danger map"))
            self.menu2d_dangermap_tiled = self.menuwolf2d.Append(wx.ID_ANY, _("Danger map tiled"), _("Compute the danger map in tiled mode -- Need wolfgpu >= 1.4.0"))
            self.menu2d_dangermap_mp = self.menuwolf2d.Append(wx.ID_ANY, _("Danger map (multiprocess)"), _("Compute the danger map using multiprocessing -- Need to duplicate the model, the memory usage can be very high for large model"))
            self.menu2d_dangermaph = self.menuwolf2d.Append(wx.ID_ANY, _("Danger map - only h"), _("Compute the danger map - only waterdepth"))

            self.menuwolf2d.AppendSeparator()

            self.menu2d_export_as = self.menuwolf2d.Append(wx.ID_ANY, _("Export results as..."), _("Export results as Geotif, Shapefile or Numpy arrays"))

            self.menuwolf2d.AppendSeparator()
            # Possible cache entries will be added after this separator

            self.menubar.Append(self.menuwolf2d, _('Results 2D'))

            self.menuwolf2d.Bind(wx.EVT_MENU, self.Onmenuwolf2d)

    def menu_alaro_forecasts(self):
        """ Menu for Alaro forecasts """
        if self.menualaro is None:

            self.menualaro = wx.Menu()
            self.menubar.Append(self.menualaro, _('Alaro forecasts'))
            self.menualaro_ftp = wx.Menu()
            self.menualaro.AppendSubMenu(self.menualaro_ftp, _('FTP server'))

            self.menualaro.Bind(wx.EVT_MENU, self.Onmenualaro)

            self.alaro_show_grid = self.menualaro.Append(wx.ID_ANY, _("Show grid"), _("Show the Alaro grid on the viewer"))
            self.alaro_list_runs = self.menualaro_ftp.Append(wx.ID_ANY, _("List runs on FTP server"), _("List available Alaro runs on the FTP server"))
            self.alaro_download_runs = self.menualaro_ftp.Append(wx.ID_ANY, _("Download runs"), _("Download Alaro runs"))
            self.alaro_download_runs_only_rain = self.menualaro_ftp.Append(wx.ID_ANY, _("Download runs - Only Rain and Temperature"), _("Download Alaro runs"))
            self.alaro_load_run = self.menualaro.Append(wx.ID_ANY, _("Load run"), _("Load an available Alaro run"))
            self.alaro_show_run = self.menualaro.Append(wx.ID_ANY, _("Add forecast(s) as array(s)"), _("Show Alaro forecasts on the viewer as WolfArray - Total precipitation [mm]"))
            self.alaro_show_all = self.menualaro.Append(wx.ID_ANY, _("Add all forecasts as arrays"), _("Show all available Alaro runs on the viewer as WolfArrays - Total precipitation [mm]"))
            self.alaro_plot_xy = self.menualaro.Append(wx.ID_ANY, _("Plot for XY"), _("Plot the next selected point"))

            self.menualaro_videos = wx.Menu()
            self.menualaro.AppendSubMenu(self.menualaro_videos, _('Videos'))

            self.alaro_videos_totprec = self.menualaro_videos.Append(wx.ID_ANY, _("Total precip. [mm] - One run"), _("Create a video of the Total Precipitation forecasts for a specific run"))
            self.alaro_videos_totprec_all = self.menualaro_videos.Append(wx.ID_ANY, _("Total precip. [mm] - Multiple runs"), _("Create videos of the Total Precipitation forecasts for all runs"))
            self.alaro_videos_comparison = self.menualaro_videos.Append(wx.ID_ANY, _("Rain intensity [mm/h] - Comparison") , _("Create a video comparing the Rain Intensity forecasts of different runs"))

    def Onmenualaro(self, event: wx.Event):
        """ Action for Alaro forecasts """

        from datetime import timezone as tz

        item = event.GetEventObject().FindItemById(event.GetId())
        itemlabel = item.ItemLabel

        if self.active_alaro is None:
            self.active_alaro = IRM_Alaro()

        if itemlabel == _("List runs on FTP server"):

            runs = self.active_alaro.run_dates_str

            for run in runs:
                logging.info(f'Available Alaro run: {run}')

            dlg = wx.MessageDialog(self, _('Available Alaro runs:\n\n') + '\n'.join(runs), _('Alaro runs'), wx.OK | wx.ICON_INFORMATION)
            dlg.ShowModal()

        elif itemlabel == _("Rain intensity [mm/h] - Comparison"):

            data_dir = self.active_alaro.data_directory
            dates = self.active_alaro.list_run_dates_cached()

            dates_str = [datetime.strptime(date, '%Y%m%d%H').replace(tzinfo=tz.utc) for date in dates]
            dates_str = [date.strftime('%Y-%m-%d %H+00') for date in dates_str]

            dlg = wx.MultiChoiceDialog(self, _('Choose one or multiple Alaro runs to treat'), _('Choose Alaro runs'), choices=dates_str)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return
            choices = dlg.GetSelections()
            dlg.Destroy()

            if len(choices) == 0:
                logging.warning(_('No Alaro run selected'))
                return

            dates = [dates[choice] for choice in choices]

            dlg = wx.FileDialog(self, _('Choose output video filename'), wildcard='mp4 (*.mp4)|*.mp4', defaultDir=str(data_dir), defaultFile='Alaro_RainIntensity_Comparison.mp4', style=wx.FD_SAVE)
            ret = dlg.ShowModal()
            if ret == wx.ID_OK:
                output_dir = dlg.GetPath()
            dlg.Destroy()

            self.active_alaro.reset_gdf()

            pgbar = wx.ProgressDialog(_('Loading Alaro runs'), _('Loading data...'), maximum=100, parent=self, style=wx.PD_APP_MODAL | wx.PD_AUTO_HIDE)

            pgbar.Pulse(_('Creating video...'))
            outfile = self.active_alaro.video_gradient_cumulated_rain_compare(Path(output_dir), run_dates = dates)
            pgbar.Update(100, _('Video created'))
            pgbar.Destroy()

            logging.info(f'Videos: {dates}')

            # Propose to lauch the video
            dlg = wx.MessageDialog(self, _('Video created: {}\n\nDo you want to launch it ?').format(outfile), _('Video created'), wx.YES_NO | wx.ICON_QUESTION)
            ret = dlg.ShowModal()
            if ret == wx.ID_YES:
                if sys.platform == "win32":
                    os.startfile(outfile)
                # elif sys.platform == "darwin":
                #     subprocess.call(["open", outfile])
                # elif sys.platform == "linux":
                #     subprocess.call(["xdg-open", outfile])
            dlg.Destroy()

        elif itemlabel == _("Total precip. [mm] - One run"):

            data_dir = self.active_alaro.data_directory
            dates = self.active_alaro.list_run_dates_cached()

            dates_str = [datetime.strptime(date, '%Y%m%d%H').replace(tzinfo=tz.utc) for date in dates]
            dates_str = [date.strftime('%Y-%m-%d %H+00') for date in dates_str]

            dlg = wx.SingleChoiceDialog(self, _('Choose one Alaro run to treat'), _('Choose Alaro run'), choices=dates_str)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return
            choice = dlg.GetSelection()
            dlg.Destroy()

            if choice == wx.NOT_FOUND:
                logging.warning(_('No Alaro run selected'))
                return

            dates = [dates[choice]]

            dlg = wx.FileDialog(self, _('Choose output video filename'), wildcard='mp4 (*.mp4)|*.mp4', defaultDir=str(data_dir), defaultFile=f"Alaro_cumulated_rain_{dates[0]}.mp4", style=wx.FD_SAVE)
            ret = dlg.ShowModal()
            if ret == wx.ID_OK:
                output_dir = dlg.GetPath()
            dlg.Destroy()

            self.active_alaro.reset_gdf()

            pgbar = wx.ProgressDialog(_('Loading Alaro runs'), _('Loading data...'), maximum=100, parent=self, style=wx.PD_APP_MODAL | wx.PD_AUTO_HIDE)

            pgbar.Pulse(_('Creating video...'))
            outfile = self.active_alaro.video_cumulated_rain(dates[0], Path(output_dir))
            pgbar.Update(100, _('Video created'))
            pgbar.Destroy()

            # Propose to lauch the video
            dlg = wx.MessageDialog(self, _('Video created: {}\n\nDo you want to launch it ?').format(outfile), _('Video created'), wx.YES_NO | wx.ICON_QUESTION)
            ret = dlg.ShowModal()
            if ret == wx.ID_YES:
                if sys.platform == "win32":
                    os.startfile(outfile)
                # elif sys.platform == "darwin":
                #     subprocess.call(["open", outfile])
                # elif sys.platform == "linux":
                #     subprocess.call(["xdg-open", outfile])
            dlg.Destroy()

        elif itemlabel == _("Total precip. [mm] - Multiple runs"):

            data_dir = self.active_alaro.data_directory
            dates = self.active_alaro.list_run_dates_cached()

            dates_str = [datetime.strptime(date, '%Y%m%d%H').replace(tzinfo=tz.utc) for date in dates]
            dates_str = [date.strftime('%Y-%m-%d %H+00') for date in dates_str]

            dlg = wx.MultiChoiceDialog(self, _('Choose multiple Alaro runs to treat'), _('Choose Alaro runs'), choices=dates_str)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return
            choices = dlg.GetSelections()
            dlg.Destroy()

            if len(choices) == 0:
                logging.warning(_('No Alaro run selected'))
                return

            dates = [dates[choice] for choice in choices]

            dlg = wx.DirDialog(self, _('Choose output video directory'), defaultPath=str(data_dir), style=wx.DD_DEFAULT_STYLE | wx.DD_DIR_MUST_EXIST)
            ret = dlg.ShowModal()
            if ret == wx.ID_OK:
                output_dir = dlg.GetPath()
            dlg.Destroy()

            self.active_alaro.reset_gdf()

            pgbar = wx.ProgressDialog(_('Loading Alaro runs'), _('Loading data...'), maximum=100, parent=self, style=wx.PD_APP_MODAL | wx.PD_AUTO_HIDE)

            pgbar.Pulse(_('Creating video...'))
            outdir = self.active_alaro.videos_cumulated_rain_allforecasts(output_dir, run_dates=dates)
            pgbar.Update(100, _('Video created'))
            pgbar.Destroy()

        elif itemlabel == _("Download runs"):

            pgbar = wx.ProgressDialog(_('Downloading Alaro runs'), _('Downloading data...'), maximum=100, parent=self, style=wx.PD_APP_MODAL | wx.PD_AUTO_HIDE)
            pgbar.Pulse(_('Downloading Alaro runs...'))
            ret = self.active_alaro.download_all_available_files()
            pgbar.Update(100, _('Alaro runs downloaded'))
            pgbar.Destroy()

            for r in ret:
                logging.info(f'Downloaded Alaro run: {r}')

        elif itemlabel == _("Download runs - Only Rain and Temperature"):

            pgbar = wx.ProgressDialog(_('Downloading Alaro runs'), _('Downloading data...'), maximum=100, parent=self, style=wx.PD_APP_MODAL | wx.PD_AUTO_HIDE)
            pgbar.Pulse(_('Downloading Alaro runs...'))
            ret = self.active_alaro.download_TotalPrecipitations_available_files()
            pgbar.Update(100, _('Alaro runs downloaded'))
            pgbar.Destroy()

            for r in ret:
                logging.info(f'Downloaded Alaro run: {r}')

        elif itemlabel == _("Load run"):

            data_dir = self.active_alaro.data_directory
            dates = self.active_alaro.list_run_dates_cached()

            dates_str = [datetime.strptime(date, '%Y%m%d%H').replace(tzinfo=tz.utc) for date in dates]
            dates_str = [date.strftime('%Y-%m-%d %H+00') for date in dates_str]

            dlg = wx.MultiChoiceDialog(self, _('Choose one or multiple Alaro runs to load'), _('Choose Alaro runs'), choices=dates_str)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return
            choices = dlg.GetSelections()
            dlg.Destroy()

            if len(choices) == 0:
                logging.warning(_('No Alaro run selected'))
                return

            dates = [dates[choice] for choice in choices]

            self.active_alaro.reset_gdf()

            pgbar = wx.ProgressDialog(_('Loading Alaro runs'), _('Loading data...'), maximum=100, parent=self, style=wx.PD_APP_MODAL | wx.PD_AUTO_HIDE)
            pgbar.Pulse(_('Loading Alaro runs...'))
            self.active_alaro.load_grib_data_to_gdf(GribFiles.FILE_TotPrecip, dates)
            pgbar.Update(100, _('Alaro runs loaded'))
            pgbar.Destroy()
            logging.info(f'Loaded Alaro run: {dates}')

        elif itemlabel == _("Add forecast(s) as array(s)"):

            if self.active_alaro._gdf is None:
                logging.warning(_('No Alaro run loaded -- Please load a run first'))
                return

            columns = self.active_alaro.get_forecast_columns()

            columns_str = [_convert_col2date_str(col) for col in columns]

            dlg = wx.MultiChoiceDialog(self, _('Choose one or multiple Alaro forecast to show'), _('Choose Alaro forecasts'), choices=columns_str)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            choices = dlg.GetSelections()
            dlg.Destroy()

            if len(choices) == 0:
                logging.warning(_('No Alaro forecast selected'))
                return

            columns = [columns[choice] for choice in choices]

            pgbar = wx.ProgressDialog(_('Loading Alaro forecasts'), _('Loading data...'), maximum=100, parent=self, style=wx.PD_APP_MODAL | wx.PD_AUTO_HIDE)
            pgbar.Pulse(_('Loading Alaro forecasts...'))
            arrays = self.active_alaro.forecasts_to_arrays(columns)
            pgbar.Update(100, _('Alaro forecasts loaded'))
            pgbar.Destroy()

            if len(arrays) == 0:
                logging.warning(_('No Alaro forecast to show'))
                return

            for array in arrays:
                array.array *= 1000.0  # Convert from m to mm
                self.add_object('array', newobj=array, id=array.idx)

            self.Refresh()

            logging.info(_('Alaro forecasts added to the viewer'))

        elif itemlabel == _("Show grid"):

            if self.active_alaro._zones is None:
                self.active_alaro._prepare_Zones_from_grib(GribFiles.FILE_TotPrecip, self.active_alaro.run_dates[0])

            self.active_alaro._zones.prep_listogl()
            self.add_object('vector', newobj=self.active_alaro._zones, id='Alaro_grid')
            self.Refresh()

        elif itemlabel == _("Add all forecasts as arrays"):
            if self.active_alaro._gdf is None:
                logging.warning(_('No Alaro run loaded -- Please load a run first'))
                return

            pgbar = wx.ProgressDialog(_('Loading Alaro forecasts'), _('Loading data...'), maximum=100, parent=self, style=wx.PD_APP_MODAL | wx.PD_AUTO_HIDE)
            pgbar.Pulse(_('Loading Alaro forecasts...'))
            arrays = self.active_alaro.forecasts_to_arrays()
            pgbar.Update(100, _('Alaro forecasts loaded'))
            pgbar.Destroy()

            if len(arrays) == 0:
                logging.warning(_('No Alaro forecast to show'))
                return

            for array in arrays:
                array.array *= 1000.0  # Convert from m to mm
                self.add_object('array', newobj=array, id=array.idx, ToCheck= False)

            self.Refresh()
            logging.info(_('Alaro forecasts added to the viewer'))

        elif itemlabel == _("Plot for XY"):
            if self.active_alaro._gdf is None:
                logging.warning(_('No Alaro run loaded -- Please load a run first'))
                return

            self.start_action('plot alaro xy', _('Plot Alaro forecasts for a specific point -- Right click to select points'))

    def menu_landuse_landcover(self):

        if self.menuLandUseLandCover is None:

            self.menuLandUseLandCover = wx.Menu()

            self._importFromFile_menu = wx.Menu()

            self._vectorImport = self._importFromFile_menu.Append(wx.ID_ANY, _("Import vector LU-LC data"), _("Import vector landuse/landcover data and map to Manning's n for active array"))
            self._rasterImport = self._importFromFile_menu.Append(wx.ID_ANY, _("Import raster LU-LC data"), _("Import raster landuse/landcover data and map to Manning's n for active array"))

            self._menuwalous = wx.Menu()
            self._uts_menu = wx.Menu()
            self._ocs_menu = wx.Menu()

            self._uts_crop = self._uts_menu.Append(wx.ID_ANY, _("Crop on active array"), _("Crop UTS data (vectorized - shp, gpkg) on the active array extent"))
            self._uts_cropscreen = self._uts_menu.Append(wx.ID_ANY, _("Crop on screen"), _("Crop UTS data (vectorized - shp, gpkg) on the current screen extent"))
            self._uts_map = self._uts_menu.Append(wx.ID_ANY, _("Map active array (WAL_UTS -> Manning)"), _("Map Walous UTS active array to Manning's n"))
            self._uts_legend = self._uts_menu.Append(wx.ID_ANY, _("Legend"), _("Legend"))

            self._ocs_crop_10m_2023 = self._ocs_menu.Append(wx.ID_ANY, _("Crop on active array (prepared data 10 m - 2023)"), _("Crop OCS data (matrix - geotif) on the active array extent using prepared 10m data"))
            self._ocs_cropscreen_10m_2023 = self._ocs_menu.Append(wx.ID_ANY, _("Crop on screen (prepared data 10 m - 2023)"), _("Crop OCS data (matrix - geotif) on the current screen extent using prepared 10m data"))
            self._ocs_menu.AppendSeparator()
            self._ocs_crop_10m_2020 = self._ocs_menu.Append(wx.ID_ANY, _("Crop on active array (prepared data 10 m - 2020)"), _("Crop OCS data (matrix - geotif) on the active array extent using prepared 10m data"))
            self._ocs_cropscreen_10m_2020 = self._ocs_menu.Append(wx.ID_ANY, _("Crop on screen (prepared data 10 m - 2020)"), _("Crop OCS data (matrix - geotif) on the current screen extent using prepared 10m data"))
            self._ocs_menu.AppendSeparator()
            self._ocs_crop = self._ocs_menu.Append(wx.ID_ANY, _("Crop on active array"), _("Crop OCS data (matrix - geotif) on the active array extent"))
            self._ocs_cropscreen = self._ocs_menu.Append(wx.ID_ANY, _("Crop on screen"), _("Crop OCS data (matrix - geotif) on the current screen extent"))
            self._ocs_menu.AppendSeparator()
            self._ocs_map = self._ocs_menu.Append(wx.ID_ANY, _("Map active array (WAL_OCS -> Hydrology)"), _("Map Walous OCS active array to Hydrology's landuse classification"))
            self._ocs_map = self._ocs_menu.Append(wx.ID_ANY, _("Map active array (WAL_OCS -> Manning)"), _("Map Walous OCS active array to Manning's n"))
            self._ocs_legend = self._ocs_menu.Append(wx.ID_ANY, _("Legend"), _("Legend"))

            self._walous_UTS_filepath = None
            self._walous_OCS_filepath = None
            self._walous_layer = None
            self._walous_map = None

            self._menuwalous.AppendSubMenu(self._uts_menu, _('UTS'))
            self._menuwalous.AppendSubMenu(self._ocs_menu, _('OCS'))

            self.menuLandUseLandCover.AppendSubMenu(self._importFromFile_menu, _('Map imported data to Manning'))
            self.menuLandUseLandCover.AppendSubMenu(self._menuwalous, _('Walous'))
            self.menubar.Append(self.menuLandUseLandCover, _('Landuse/Landcover'))

            self._importFromFile_menu.Bind(wx.EVT_MENU, self.Onmenu_landuse_landcover_importfromfile)
            self._uts_menu.Bind(wx.EVT_MENU, self.Onmenuwalous_uts)
            self._ocs_menu.Bind(wx.EVT_MENU, self.Onmenuwalous_ocs)

    def menu_drowning(self):

        if self.menudrowning is None:

            self.menudrowning = wx.Menu()

            self.menudrowning_plot = wx.Menu()

            self.menudrowning_explore_results = self.menudrowning.Append(wx.ID_ANY, _("Explore time/index results"), _("Open a dialog to explore time/index results"))

            self.menudrowning.AppendSeparator()

            self.menudrowning.Append(wx.ID_ANY, _("Plot..."),self.menudrowning_plot)
            self.plot_runs = self.menudrowning_plot.Append(wx.ID_ANY, _("Plot runs positions"), _("Plot runs positions"),kind=wx.ITEM_CHECK)
            self.plot_cells = self.menudrowning_plot.Append(wx.ID_ANY, _("Plot cells positions"), _("Plot the cells where bodies colored as a function of bodies in it"),kind=wx.ITEM_CHECK)
            self.plot_KDE = self.menudrowning_plot.Append(wx.ID_ANY, _("Plot KDE"), _("Plot Kernel Density Estimation - Map of probability of presence"),kind=wx.ITEM_CHECK)

            self.menudrowning_lastres = self.menudrowning.Append(wx.ID_ANY, _("Read last result"), _("Current view"))

            self.menudrowning_zoom = self.menudrowning.Append(wx.ID_ANY, _("Zoom on hotspots"), _("Zoom on areas where you have the highest probability of presence"))

            self.menudrowning.AppendSeparator()

            self.menudrowning_get_bodies = self.menudrowning.Append(wx.ID_ANY, _("Get bodies characteristics"), _("Get a table of the characteristics of all simulated bodies"))
            self.menudrowning_get_perc = self.menudrowning.Append(wx.ID_ANY, _("Vertical position proportion"), _("Get proportion of runs at the surface and at the bottom at selected time"))

            self.menudrowning.AppendSeparator()

            self.menudrowning_video = self.menudrowning.Append(wx.ID_ANY, _("Create video..."), _("Video/Movie"))

            self.menubar.Append(self.menudrowning, _('Drowning'))

            self.menudrowning.Bind(wx.EVT_MENU, self.Onmenudrowning)

    def menu_dike(self):

        if self.menudike is None:
            self.menudike = wx.Menu()

            self.menudike_launchsimu = self.menudike.Append(wx.ID_ANY, _("Launch lumped simulation"), _("Launch lumped simulation"))

            self.menudike.AppendSeparator()

            self.menudike_setinjector = self.menudike.Append(wx.ID_ANY, _("Set injector"), _("Set injector"))
            self.menudike_launch2Dsimu = self.menudike.Append(wx.ID_ANY, _("Launch 2D-coupled simulation"), _("Launch 2D-coupled simulation"))

            self.menudike.AppendSeparator()

            self.menudike_showtri = self.menudike.Append(wx.ID_ANY, _("Show triangulation"), _("Show triangulation"))
            self.menudike_plotQ = self.menudike.Append(wx.ID_ANY, _("Plot discharges"), _("Plot discharges"))
            self.menudike_plotz = self.menudike.Append(wx.ID_ANY, _("Plot water levels/breach bottom"), _("Plot water levels/breach bottom"))
            self.menudike_plotB = self.menudike.Append(wx.ID_ANY, _("Plot breach width"), _("Plot breach width"))

            self.menudike.AppendSeparator()

            self.menudike_showparam = self.menudike.Append(wx.ID_ANY, _("Show parameters"), _("Show parameters"))

            self.menubar.Append(self.menudike, _('Dike'))

            self.menudike.Bind(wx.EVT_MENU, self.Onmenudike)


    def get_canvas_bounds(self, gridsize:float = None):
        """
        Retourne les limites de la zone d'affichage

        :return: [xmin, ymin, xmax, ymax]

        """

        if gridsize is None:

            return [self.xmin, self.ymin, self.xmax, self.ymax]

        else:

            xmin = float(np.rint(self.xmin / gridsize) * gridsize)
            ymin = float(np.rint(self.ymin / gridsize) * gridsize)
            xmax = float(np.rint(self.xmax / gridsize) * gridsize)
            ymax = float(np.rint(self.ymax / gridsize) * gridsize)

            return [xmin, ymin, xmax, ymax]

    def get_bounds(self, gridsize:float = None) -> tuple:
        """
        Retourne les limites de la zone d'affichage, voir aussi get_canvas_bounds

        :return: ([xmin, xmax], [ymin, ymax])
        """
        xmin, ymin, xmax, ymax = self.get_canvas_bounds(gridsize=gridsize)

        return ([xmin, xmax], [ymin, ymax])

    def get_bounds_as_polygon(self, gridsize:float = None) -> vector:
        """
        Retourne les limites de la zone d'affichage sous forme de polygone
        :return: vector
        """
        xmin, ymin, xmax, ymax = self.get_canvas_bounds(gridsize=gridsize)
        poly = vector()
        poly.add_vertex(wolfvertex(xmin, ymin))
        poly.add_vertex(wolfvertex(xmax, ymin))
        poly.add_vertex(wolfvertex(xmax, ymax))
        poly.add_vertex(wolfvertex(xmin, ymax))
        poly.force_to_close()
        return poly

    def Onmenu_landuse_landcover_importfromfile(self, event: wx.MenuEvent):
        """ Handle the Landuse/Landcover import from file menu events """
        id = event.GetId()
        item = self.menubar.FindItemById(event.GetId())

        if item is None:
            return

        itemlabel = item.ItemLabel

        if itemlabel == _("Import vector LU-LC data"):
            logging.info("Import vector LU-LC data selected")

            dlg = wx.FileDialog(self, _("Choose the vector LU-LC file"), wildcard="Shapefile (*.shp)|*.shp", style=wx.FD_OPEN)
            if dlg.ShowModal() == wx.ID_CANCEL:
                dlg.Destroy()
                return

            self._vectorLU_LC_filepath = Path(dlg.GetPath())
            dlg.Destroy()

            # Import vector LU-LC data
            try:
                self._vectorLU_LC_data = gpd.read_file(self._vectorLU_LC_filepath)
            except Exception as e:
                logging.error(f"Error reading shapefile: {e}")
                wx.MessageBox(
                    _("Failed to read the shapefile.\nPlease ensure it is valid."),
                    _("Error"),
                    wx.OK | wx.ICON_ERROR
                )
                return

            if self._vectorLU_LC_data.empty:
                logging.error(f"Failed to read vector LU-LC data from {self._vectorLU_LC_filepath}")
                wx.MessageBox(
                    _("The shapefile appears to be empty."),
                    _("Error"),
                    wx.OK | wx.ICON_ERROR
                )
                return
            else:
                logging.info(f"Vector LU-LC data imported from {self._vectorLU_LC_filepath}")

            headers = list(self._vectorLU_LC_data.columns)

            dlg_attr = wx.SingleChoiceDialog(
                self,
                message=_("Select the attribute column to use for Manning's mapping:"),
                caption=_("Select LU-LC Attribute"),
                choices=headers,
                style=wx.CHOICEDLG_STYLE
            )

            if dlg_attr.ShowModal() == wx.ID_OK:
                self._selected_LULC_column = dlg_attr.GetStringSelection()
                logging.info(f"Selected LU-LC attribute column: {self._selected_LULC_column}")
            else:
                logging.info("User cancelled attribute column selection.")
                dlg_attr.Destroy()
                return

            dlg_attr.Destroy()

            unique_values = sorted(self._vectorLU_LC_data[self._selected_LULC_column].dropna().unique().tolist())

            if not unique_values:
                wx.MessageBox(
                    _("No valid attribute values found in the selected column."),
                    _("Error"),
                    wx.OK | wx.ICON_ERROR
                )
                return

            dlg_map = wx.Dialog(self, title=_("Map LU-LC Classes to Manning's Coefficients"), size=(500, 400))

            panel = wx.Panel(dlg_map)
            vbox = wx.BoxSizer(wx.VERTICAL)

            info_text = wx.StaticText(
                panel,
                label=_("Enter Manning's n value for each LU/LC class:")
            )
            vbox.Add(info_text, 0, wx.ALL, 10)

            # Create table for mapping
            grid = wx.grid.Grid(panel)
            grid.CreateGrid(len(unique_values), 2)
            grid.SetColLabelValue(0, _("LU/LC Class"))
            grid.SetColLabelValue(1, _("Manning's n"))
            grid.EnableEditing(True)
            grid.SetColSize(0, 250)
            grid.SetColSize(1, 150)

            for i, val in enumerate(unique_values):
                grid.SetCellValue(i, 0, str(val))
                grid.SetReadOnly(i, 0, True)  # LU/LC column is not editable

            vbox.Add(grid, 1, wx.EXPAND | wx.ALL, 10)

            hbox = wx.BoxSizer(wx.HORIZONTAL)
            btn_ok = wx.Button(panel, wx.ID_OK, _("OK"))
            btn_cancel = wx.Button(panel, wx.ID_CANCEL, _("Cancel"))
            hbox.Add(btn_ok, 0, wx.ALL, 5)
            hbox.Add(btn_cancel, 0, wx.ALL, 5)

            vbox.Add(hbox, 0, wx.ALIGN_CENTER | wx.BOTTOM, 10)
            panel.SetSizer(vbox)

            dlg_map.Layout()

            while True:
                if dlg_map.ShowModal() != wx.ID_OK:
                    logging.info("User cancelled LU-LC to Manning's mapping.")
                    return

                mapping = {}
                validation_failed = False
                for i, val in enumerate(unique_values):
                    n_str = grid.GetCellValue(i, 1).strip()
                    if not n_str:
                        wx.MessageBox(
                            _(f"Missing Manning's n value for class '{val}'."),
                            _("Error"),
                            wx.OK | wx.ICON_ERROR
                        )
                        validation_failed = True
                        break

                    try:
                        n_value = float(n_str)
                        if n_value <= 0:
                            wx.MessageBox(
                                _(f"Invalid Manning's n value for class '{val}'. Must be positive."),
                                _("Error"),
                                wx.OK | wx.ICON_ERROR
                            )
                            validation_failed = True
                            break
                    except ValueError:
                        wx.MessageBox(
                            _(f"Invalid Manning's n value for class '{val}'. Must be numeric."),
                            _("Error"),
                            wx.OK | wx.ICON_ERROR
                        )
                        validation_failed = True
                        break

                    mapping[val] = n_value

                if validation_failed:
                    continue

                self._lulc_manning_mapping = mapping
                logging.info("LU-LC to Manning's mapping completed:")
                for k, v in mapping.items():
                    logging.info(f"  {k}: {v}")
                break

            dlg_map.Destroy()

            if hasattr(self, "_lulc_manning_mapping") and self._lulc_manning_mapping:
                manning_col_name = "MANNING_N"

                self._vectorLU_LC_data[manning_col_name] = self._vectorLU_LC_data[self._selected_LULC_column].map(self._lulc_manning_mapping)

                mapped_count = self._vectorLU_LC_data[manning_col_name].notna().sum()
                total_count = len(self._vectorLU_LC_data)

                logging.info(
                    f"Applied Manning's coefficients to {mapped_count}/{total_count} features "
                    f"using attribute '{self._selected_LULC_column}'."
                )

                wx.MessageBox(
                    _(
                        f"Manning's n values assigned to {mapped_count} out of {total_count} features.\n"
                        f"A new column '{manning_col_name}' has been added."
                    ),
                    _("Mapping Complete"),
                    wx.OK | wx.ICON_INFORMATION
                )
            else:
                logging.warning("No Manning's mapping found or user cancelled mapping process.")

            export_prompt = wx.MessageDialog(
                self,
                _(
                    "Do you want to export the updated LU/LC data (with Manning's n values) "
                    "to a new file?"
                ),
                _("Export Updated Layer"),
                style=wx.YES_NO | wx.ICON_QUESTION
            )

            if export_prompt.ShowModal() == wx.ID_YES:
                export_prompt.Destroy()

                save_dlg = wx.FileDialog(self,_("Save updated LU/LC file as"),wildcard="Shapefile (*.shp)|*.shp",style=wx.FD_SAVE | wx.FD_OVERWRITE_PROMPT)

                if save_dlg.ShowModal() == wx.ID_CANCEL:
                    save_dlg.Destroy()
                    logging.info("User cancelled export.")
                    return

                export_path = Path(save_dlg.GetPath())
                save_dlg.Destroy()

                try:
                    driver = "ESRI Shapefile"

                    self._vectorLU_LC_data.to_file(export_path, driver=driver)

                    logging.info(f"Exported updated LU/LC data to {export_path}")
                    wx.MessageBox(_(f"Successfully exported updated LU/LC data to:\n{export_path}"),_("Export Successful"),wx.OK | wx.ICON_INFORMATION)

                except Exception as e:
                    logging.error(f"Failed to export LU/LC data: {e}")
                    wx.MessageBox(_(f"Error while exporting:\n{e}"), _("Export Error"), wx.OK | wx.ICON_ERROR)

            else:
                export_prompt.Destroy()
                logging.info("User chose not to export updated LU/LC data.")

            raster_prompt = wx.MessageDialog(self, _("Do you want to create a raster file from the Manning's n data?"), _("Create Raster"), style=wx.YES_NO | wx.ICON_QUESTION)

            if raster_prompt.ShowModal() == wx.ID_YES:
                raster_prompt.Destroy()

                choices = [_("Use vector bounds"), _("Crop to active array")]
                bounds_dlg = wx.SingleChoiceDialog(self, _("Select bounding extent:"), _("Raster Bounds"), choices)

                if bounds_dlg.ShowModal() == wx.ID_CANCEL:
                    bounds_dlg.Destroy()
                    return

                bounds_choice = bounds_dlg.GetStringSelection()
                bounds_dlg.Destroy()

                if bounds_choice == _("Crop to active array"):
                    if self.active_array is None:
                        logging.warning(_("No active array -- Please activate data first"))
                        wx.MessageBox(_("No active array available.\nCannot crop raster."), _("Error"), wx.OK | wx.ICON_ERROR)
                        return
                    bounds_x, bounds_y = self.active_array.get_bounds()
                    xmin, xmax = bounds_x
                    ymin, ymax = bounds_y
                    spatial_res = self.active_array.dx
                    def_outdir = Path(self.active_array.filename).parent

                else:
                    xmin, ymin, xmax, ymax = self._vectorLU_LC_data.total_bounds
                    def_outdir = Path(self._vectorLU_LC_filepath).parent

                    res_dlg = wx.TextEntryDialog(self, _("Enter spatial resolution (m):"), _("Spatial Resolution"), value="10.0")
                    if res_dlg.ShowModal() == wx.ID_OK:
                        try:
                            spatial_res = float(res_dlg.GetValue())
                        except ValueError:
                            wx.MessageBox(_("Invalid resolution entered. Must be numeric."), _("Error"), wx.OK | wx.ICON_ERROR)
                            res_dlg.Destroy()
                            return
                    else:
                        res_dlg.Destroy()
                        return
                    res_dlg.Destroy()

                save_dlg = wx.FileDialog(
                    self,
                    _("Save raster file as"),
                    defaultDir=str(def_outdir),
                    wildcard="GeoTIFF (*.tif)|*.tif",
                    style=wx.FD_SAVE | wx.FD_OVERWRITE_PROMPT
                )

                if save_dlg.ShowModal() == wx.ID_CANCEL:
                    save_dlg.Destroy()
                    return

                raster_path = Path(save_dlg.GetPath())
                save_dlg.Destroy()

                # Rasterization
                import rasterio
                from rasterio import features
                try:
                    logging.info("Starting rasterization of Manning's n data...")

                    width = int((xmax - xmin) / spatial_res)
                    height = int((ymax - ymin) / spatial_res)

                    transform = rasterio.transform.from_origin(xmin, ymax, spatial_res, spatial_res)

                    shapes = zip(
                        self._vectorLU_LC_data.geometry,
                        self._vectorLU_LC_data["MANNING_N"]
                    )

                    raster_data = features.rasterize(
                        shapes=shapes,
                        out_shape=(height, width),
                        transform=transform,
                        fill=np.nan,
                        dtype="float32"
                    )

                    new_dataset = rasterio.open(
                        raster_path,
                        "w",
                        driver="GTiff",
                        height=height,
                        width=width,
                        count=1,
                        dtype="float32",
                        crs=self._vectorLU_LC_data.crs,
                        transform=transform,
                        nodata=np.nan
                    )
                    new_dataset.write(raster_data, 1)
                    new_dataset.close()

                    logging.info(f"Raster created successfully: {raster_path}")
                    wx.MessageBox(_(f"Raster file created successfully:\n{raster_path}"), _("Rasterization Complete"), wx.OK | wx.ICON_INFORMATION)

                    dlg = wx.MessageDialog(self, _('Do you want to load the created file ?'), _('Load file'), wx.YES_NO | wx.ICON_QUESTION)
                    ret = dlg.ShowModal()
                    dlg.Destroy()

                    if ret == wx.ID_CANCEL:
                        return

                    if ret == wx.ID_YES:
                        manning_array = WolfArray(fname=raster_path)
                        self.add_object('array', newobj=manning_array, id='manning_array')

                except Exception as e:
                    logging.error(f"Error creating raster: {e}")
                    wx.MessageBox(_(f"Error creating raster:\n{e}"),_("Rasterization Error"),wx.OK | wx.ICON_ERROR)

            else:
                raster_prompt.Destroy()
                logging.info("User skipped raster creation.")

            self.Autoscale()

        elif itemlabel == _("Import raster LU-LC data"):
            logging.info("Import raster LU-LC data selected")
            dlg = wx.FileDialog(self, _("Choose the raster LU-LC file"), wildcard="GeoTIFF (*.tif)|*.tif|all (*.*)|*.*", style=wx.FD_OPEN)
            if dlg.ShowModal() == wx.ID_CANCEL:
                dlg.Destroy()
                return
            raster_path = Path(dlg.GetPath())
            dlg.Destroy()
            try:
                lu_lc_array = WolfArray(fname=raster_path)

                unique_values = sorted(lu_lc_array.get_unique_values())

                if not unique_values:
                    wx.MessageBox(
                        _("Something went wrong reading unique values from the raster. Please check the file."),
                        _("Error"),
                        wx.OK | wx.ICON_ERROR
                    )
                    return

                dlg_map = wx.Dialog(self, title=_("Map LU-LC Classes to Manning's Coefficients"), size=(500, 400))

                panel = wx.Panel(dlg_map)
                vbox = wx.BoxSizer(wx.VERTICAL)

                info_text = wx.StaticText(
                    panel,
                    label=_("Enter Manning's n value for each LU/LC class. \nCheck the file's documentation for class meanings:")
                )
                vbox.Add(info_text, 0, wx.ALL, 10)

                # Create table for mapping
                grid = wx.grid.Grid(panel)
                grid.CreateGrid(len(unique_values), 2)
                grid.SetColLabelValue(0, _("LU/LC Class"))
                grid.SetColLabelValue(1, _("Manning's n"))
                grid.EnableEditing(True)
                grid.SetColSize(0, 150)
                grid.SetColSize(1, 150)

                for i, val in enumerate(unique_values):
                    grid.SetCellValue(i, 0, str(val))
                    grid.SetReadOnly(i, 0, True)  # LU/LC column is not editable

                vbox.Add(grid, 1, wx.EXPAND | wx.ALL, 10)

                hbox = wx.BoxSizer(wx.HORIZONTAL)
                btn_ok = wx.Button(panel, wx.ID_OK, _("OK"))
                btn_cancel = wx.Button(panel, wx.ID_CANCEL, _("Cancel"))
                hbox.Add(btn_ok, 0, wx.ALL, 5)
                hbox.Add(btn_cancel, 0, wx.ALL, 5)

                vbox.Add(hbox, 0, wx.ALIGN_CENTER | wx.BOTTOM, 10)
                panel.SetSizer(vbox)

                dlg_map.Layout()

                while True:
                    if dlg_map.ShowModal() != wx.ID_OK:
                        logging.info("User cancelled LU-LC to Manning's mapping.")
                        return

                    mapping = {}
                    validation_failed = False
                    for i, val in enumerate(unique_values):
                        n_str = grid.GetCellValue(i, 1).strip()
                        if not n_str:
                            wx.MessageBox(
                                _(f"Missing Manning's n value for class '{val}'."),
                                _("Error"),
                                wx.OK | wx.ICON_ERROR
                            )
                            validation_failed = True
                            break

                        try:
                            n_value = float(n_str)
                            if n_value <= 0:
                                wx.MessageBox(
                                    _(f"Invalid Manning's n value for class '{val}'. Must be positive."),
                                    _("Error"),
                                    wx.OK | wx.ICON_ERROR
                                )
                                validation_failed = True
                                break
                        except ValueError:
                            wx.MessageBox(
                                _(f"Invalid Manning's n value for class '{val}'. Must be numeric."),
                                _("Error"),
                                wx.OK | wx.ICON_ERROR
                            )
                            validation_failed = True
                            break

                        mapping[val] = n_value

                    if validation_failed:
                        continue

                    self._lulc_manning_mapping = mapping
                    logging.info("LU-LC to Manning's mapping completed:")
                    for k, v in self._lulc_manning_mapping.items():
                        logging.info(f"  {k}: {v}")
                    break

                dlg_map.Destroy()

                if hasattr(self, "_lulc_manning_mapping") and self._lulc_manning_mapping:
                    lu_lc_array_copy = lu_lc_array.array.data.copy()
                    for lu_lc_value, manning_n in self._lulc_manning_mapping.items():
                        lu_lc_array.array.data[lu_lc_array_copy == lu_lc_value] = manning_n

                    if self.active_array is not None:
                        dlg_crop = wx.MessageDialog(
                            self, _("Do you want to crop the raster to the active array extent?"), _('Crop Raster'), wx.YES_NO | wx.ICON_QUESTION
                        )

                        if dlg_crop.ShowModal() == wx.ID_YES:
                            try:
                                lu_lc_array = WolfArray(mold=lu_lc_array, crop=self.active_array.get_bounds())
                                logging.info("Raster cropped to active array extent.")
                            except Exception as e:
                                logging.error(f"Error cropping raster: {e}")
                        dlg_crop.Destroy()

                        dlg_resample = wx.TextEntryDialog(
                            self, _(f"Resample the raster to active array resolution ({self.active_array.dx} m) \n or enter desired resolution"),
                            _("Resample Raster"), value=str(self.active_array.dx)
                        )
                        try:
                            if dlg_resample.ShowModal() == wx.ID_OK:
                                try:
                                    new_resolution = float(dlg_resample.GetValue())
                                    if new_resolution <= 0:
                                        raise ValueError("Resolution must be positive.")
                                    lu_lc_array.rebin(new_resolution / lu_lc_array.dx, operation='min')
                                    logging.info(f"Raster resampled to {new_resolution} m.")

                                    save_dlg = wx.FileDialog(
                                        self,
                                        _("Save raster file as"),
                                        defaultDir=str(raster_path),
                                        wildcard="GeoTIFF (*.tif)|*.tif",
                                        style=wx.FD_SAVE | wx.FD_OVERWRITE_PROMPT
                                    )

                                    if save_dlg.ShowModal() == wx.ID_CANCEL:
                                        save_dlg.Destroy()
                                        return

                                    raster_path = Path(save_dlg.GetPath())
                                    save_dlg.Destroy()

                                    lu_lc_array.write_all(raster_path)
                                    logging.info(f"Raster file saved: {raster_path}")

                                    dlg = wx.MessageDialog(self, _('Do you want to load the created array ?'), _('Load file'), wx.YES_NO | wx.ICON_QUESTION)
                                    ret = dlg.ShowModal()
                                    dlg.Destroy()

                                    if ret == wx.ID_CANCEL:
                                        return

                                    if ret == wx.ID_YES:
                                        manning_array = WolfArray(fname=raster_path)
                                        self.add_object('array', newobj=manning_array, id=f'{raster_path.stem}')

                                except ValueError as e:
                                    wx.MessageBox(
                                        _(f"Invalid resolution value: {e}"),
                                        _("Error"),
                                        wx.OK | wx.ICON_ERROR
                                    )
                                    return
                        finally:
                            dlg_resample.Destroy()
                    else:
                        logging.info("No active array to resample raster to.")
                        dlg_resample = wx.TextEntryDialog(
                            self, _(f"Enter desired spatial resolution (m) for the raster \n Current: {lu_lc_array.dx}"), _("Resample Raster"), value=""
                        )

                        try:
                            if dlg_resample.ShowModal() == wx.ID_OK:
                                try:
                                    new_resolution = float(dlg_resample.GetValue())
                                    if new_resolution <= 0:
                                        raise ValueError("Resolution must be positive.")
                                    lu_lc_array.rebin(new_resolution / lu_lc_array.dx, operation='min')
                                    logging.info(f"Raster resampled to {new_resolution} m.")

                                    save_dlg = wx.FileDialog(
                                        self,
                                        _("Save raster file as"),
                                        defaultDir=str(raster_path),
                                        wildcard="GeoTIFF (*.tif)|*.tif",
                                        style=wx.FD_SAVE | wx.FD_OVERWRITE_PROMPT
                                    )

                                    if save_dlg.ShowModal() == wx.ID_CANCEL:
                                        save_dlg.Destroy()
                                        return

                                    raster_path = Path(save_dlg.GetPath())
                                    save_dlg.Destroy()

                                    lu_lc_array.write_all(raster_path)
                                    logging.info(f"Raster file saved: {raster_path}")

                                    dlg = wx.MessageDialog(self, _('Do you want to load the created array ?'), _('Load file'), wx.YES_NO | wx.ICON_QUESTION)
                                    ret = dlg.ShowModal()
                                    dlg.Destroy()

                                    if ret == wx.ID_CANCEL:
                                        return

                                    if ret == wx.ID_YES:
                                        manning_array = WolfArray(fname=raster_path)
                                        self.add_object('array', newobj=manning_array, id=f'{raster_path.stem}')

                                except ValueError as e:
                                    wx.MessageBox(
                                        _(f"Invalid resolution value: {e}"),
                                        _("Error"),
                                        wx.OK | wx.ICON_ERROR
                                    )
                                    return
                        finally:
                            dlg_resample.Destroy()

                else:
                    logging.info("No LU-LC to Manning's mapping available.")

            except Exception as e:
                logging.error(f"Error: {e}")

            self.Autoscale()

    def Onmenuwalous_ocs(self, event: wx.MenuEvent):
        """ Handle the Walous OCS menu events """
        id = event.GetId()
        item = self.menubar.FindItemById(event.GetId())

        if item is None:
            return

        itemlabel = item.ItemLabel

        if _("Crop on active array") in itemlabel or _("Crop on screen") in itemlabel:

            if _("Crop on screen") in itemlabel:

                bounds = self.get_canvas_bounds(gridsize=1.)

                def_outdrir = ''
                spatial_res = 1.

                if self.active_array is not None:
                    spatial_res = self.active_array.dx

                dlg = wx.TextEntryDialog(None,_("Spatial resolution [m] ?"), value = str(spatial_res))

                dlg.ShowModal()
                try:
                    spatial_res = float(dlg.GetValue())
                    dlg.Destroy()
                except:
                    dlg.Destroy()
                    logging.warning(_("Bad value -- Rety"))
                    return

            else:

                if self.active_array is None:
                    logging.warning(_('No active array -- Please activate data first'))
                    return

                bounds = self.active_array.get_bounds()
                def_outdrir = Path(self.active_array.filename).parent
                spatial_res = self.active_array.dx

            from .pywalous import update_palette_walous_ocs

            if self._walous_OCS_filepath is None:

                if itemlabel in [_("Crop on active array (prepared data 10 m - 2023)"), _("Crop on screen (prepared data 10 m - 2023)")]:
                    self._walous_OCS_filepath = toys_dataset('Walous_OCS', 'WALOUS_2023_lbt72_10m.tif')
                elif itemlabel in [_("Crop on active array (prepared data 10 m - 2020)"), _("Crop on screen (prepared data 10 m - 2020)")]:
                    self._walous_OCS_filepath = toys_dataset('Walous_OCS', 'WALOUS_2020_lbt72_10m.tif')

                else:
                    dlg = wx.FileDialog(self, _("Choose the Walous OCS Tif file"), wildcard="Tif file (*.tif)|*.tif|all (*.*)|*.*", style=wx.FD_OPEN)
                    if dlg.ShowModal() == wx.ID_CANCEL:
                        dlg.Destroy()
                        return

                    self._walous_OCS_filepath = Path(dlg.GetPath())
                    dlg.Destroy()

            if self._walous_OCS_filepath is None or not Path(self._walous_OCS_filepath).exists():
                logging.error(_('No Walous OCS file -- Please set it'))
                return

            dlg = wx.FileDialog(self, _("Choose the output file"), wildcard="Geotif (*.tif)|*.tif|all (*.*)|*.*", style=wx.FD_SAVE, defaultDir=str(def_outdrir))
            if dlg.ShowModal() == wx.ID_CANCEL:
                dlg.Destroy()
                return

            output = Path(dlg.GetPath())
            dlg.Destroy()

            header_OCS = header_wolf.read_header(self._walous_OCS_filepath)
            if header_OCS.dx != spatial_res:
                # Adapt bounds to ensure that the rebin will be correct.
                # If spatial_res is a multiple of header_OCS.dx, no need to change bounds.
                # If not, change the bounds to ensure that the crop will include data in all footprints.

                # Convert bvounds to list to be able to modify it
                bounds = [list(bounds[0]), list(bounds[1])]

                if (bounds[0][0] - header_OCS.origx) % header_OCS.dx != 0:
                    bounds[0][0] = header_OCS.origx + ((bounds[0][0] - header_OCS.origx) // header_OCS.dx) * header_OCS.dx
                if (bounds[0][1] - bounds[0][0]) % header_OCS.dx != 0:
                    bounds[0][1] = bounds[0][0] + ((bounds[0][1] - bounds[0][0]) // header_OCS.dx + 1) * header_OCS.dx
                if (bounds[1][0] - header_OCS.origy) % header_OCS.dy != 0:
                    bounds[1][0] = header_OCS.origy + ((bounds[1][0] - header_OCS.origy) // header_OCS.dy) * header_OCS.dy
                if (bounds[1][1] - bounds[1][0]) % header_OCS.dy != 0:
                    bounds[1][1] = bounds[1][0] + ((bounds[1][1] - bounds[1][0]) // header_OCS.dy + 1) * header_OCS.dy

            locwalous = WolfArray(fname=self._walous_OCS_filepath,
                                  crop = [bounds[0][0], bounds[0][1], bounds[1][0], bounds[1][1]])

            if locwalous.dx != spatial_res:
                locwalous.rebin(spatial_res / locwalous.dx, operation='min')
                logging.info(_('Rebin to {} m because original data are {} m').format(spatial_res, locwalous.dx))
                locwalous = WolfArray(mold=locwalous, crop = self.active_array.get_bounds())

            locwalous.write_all(output)

            if Path(output).exists():
                logging.info(_('File {} created').format(output))
            else:
                logging.error(_('File {} not created').format(output))
                return

            dlg = wx.MessageDialog(self, _('Do you want to load the created file ?'), _('Load file'), wx.YES_NO | wx.ICON_QUESTION)
            ret = dlg.ShowModal()

            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            elif ret == wx.ID_YES:
                walousarray = WolfArray(fname=output)
                update_palette_walous_ocs(walousarray.mypal)
                walousarray.reset_plot()
                self.add_object('array', newobj=walousarray, id = 'walous_ocs_crop')
                dlg.Destroy()

        elif itemlabel == _("Legend"):
            from .pywalous import Walous_OCS_Legend
            newlegend = Walous_OCS_Legend(self)
            newlegend.Show()

        elif itemlabel == _("Map active array (WAL_OCS -> Hydrology)"):

            from .pywalous import DlgMapWalous2Hydrology

            if self.active_array is None:
                logging.warning(_('No active array -- Please activate data first'))
                return

            if self.active_array.wolftype != WOLF_ARRAY_FULL_SINGLE:
                logging.error(_('Active array is not a Float32 array -- Please change it to Float32 before mapping'))
                return

            vals = self.active_array.get_unique_values()

            if self._walous_layer is None:

                if vals[0] > 11:
                    logging.error(_('You have values greater than 11 -- Please check your data'))
                    return

            dlg = DlgMapWalous2Hydrology(self,)

            ret = dlg.ShowModal()

            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            mapvals = dlg.get_mapping()
            dlg.Destroy()

            if mapvals == -1:
                logging.error(_('Bad values -- retry'))
                return

            self.active_array.map_values(mapvals)

            self.active_array.reset_plot()

        elif itemlabel == _("Map active array (WAL_OCS -> Manning)"):

            from .pywalous import DlgMapWalousOCS2Manning

            if self.active_array is None:
                logging.warning(_('No active array -- Please activate data first'))
                return

            if self.active_array.wolftype != WOLF_ARRAY_FULL_SINGLE:
                logging.error(_('Active array is not a Float32 array -- Please change it to Float32 before mapping'))
                return

            vals = self.active_array.get_unique_values()

            if self._walous_layer is None:

                if vals[0] > 11:
                    logging.error(_('You have values greater than 11 -- Please check your data'))
                    return

            dlg = DlgMapWalousOCS2Manning(self,)

            ret = dlg.ShowModal()

            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            mapvals = dlg.get_mapping()
            dlg.Destroy()

            if mapvals == -1:
                logging.error(_('Bad values -- retry'))
                return

            self.active_array.map_values(mapvals)

            self.active_array.reset_plot()

    def Onmenuwalous_uts(self, event: wx.MenuEvent):
        """ Handle the Walous UTS menu events """
        id = event.GetId()
        item = self.menubar.FindItemById(event.GetId())

        if item is None:
            return

        itemlabel = item.ItemLabel

        if itemlabel in [_("Crop on active array"), _("Crop on screen")]:

            if itemlabel == _("Crop on screen"):

                bounds = self.get_canvas_bounds(gridsize=1.)

                def_outdrir = ''
                spatial_res = 1.

                if self.active_array is not None:
                    spatial_res = self.active_array.dx

                dlg = wx.TextEntryDialog(None,_("Spatial resolution [m] ?"), value = str(spatial_res))

                dlg.ShowModal()
                try:
                    spatial_res = float(dlg.GetValue())
                    dlg.Destroy()
                except:
                    dlg.Destroy()
                    logging.warning(_("Bad value -- Rety"))
                    return

            else:

                if self.active_array is None:
                    logging.warning(_('No active array -- Please activate data first'))
                    return

                bounds = self.active_array.get_bounds()
                def_outdrir = Path(self.active_array.filename).parent
                spatial_res = self.active_array.dx

            from .pywalous import Walous_data, WALOUS_UTS2MANNING_MAJ_NIV1, WALOUS_UTS2MANNING_MAJ_NIV2, update_palette_walous_uts

            if self._walous_UTS_filepath is None:
                dlg = wx.FileDialog(self, _("Choose the Walous shape file"), wildcard="Geopackage (*.gpkg)|*.gpkg|Shapefile (*.shp)|*.shp|all (*.*)|*.*", style=wx.FD_OPEN)
                if dlg.ShowModal() == wx.ID_CANCEL:
                    dlg.Destroy()
                    return

                self._walous_UTS_filepath = Path(dlg.GetPath())
                dlg.Destroy()


            dlg = wx.FileDialog(self, _("Choose the output file"), wildcard="Geotif (*.tif)|*.tif|all (*.*)|*.*", style=wx.FD_SAVE, defaultDir=str(def_outdrir))
            if dlg.ShowModal() == wx.ID_CANCEL:
                dlg.Destroy()
                return

            output = Path(dlg.GetPath())
            dlg.Destroy()

            # choix de la couche entre MAJ_NIV1 et MAJ_NIV2
            dlg = wx.SingleChoiceDialog(None, _("Choose a layer"), "Choices", ['MAJ_NIV1', 'MAJ_NIV2'])
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            self._walous_layer = 'UTS_' + dlg.GetStringSelection()

            locwalous = Walous_data(self._walous_UTS_filepath.parent, self._walous_UTS_filepath.name)
            ret = locwalous.rasterize(bounds=bounds,
                                      layer=self._walous_layer,
                                      fn_out=output,
                                      pixel_size=spatial_res)

            if isinstance(ret, int):
                logging.error(_('Error {}').format(ret))
                return

            if Path(output).exists():
                logging.info(_('File {} created').format(output))
            else:
                logging.error(_('File {} not created').format(output))
                return

            dlg = wx.MessageDialog(self, _('Do you want to load the created file ?'), _('Load file'), wx.YES_NO | wx.ICON_QUESTION)
            ret = dlg.ShowModal()

            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            elif ret == wx.ID_YES:
                walousarray = WolfArray(fname=output)
                update_palette_walous_uts(self._walous_layer, walousarray.mypal)
                self.add_object('array', newobj=walousarray, id = 'walous_uts_crop')
                dlg.Destroy()

        elif itemlabel == _("Legend"):

            from .pywalous import Walous_UTS_Legend

            newlegend = Walous_UTS_Legend(self)
            newlegend.Show()

        elif itemlabel == _("Map active array (WAL_UTS -> Manning)"):

            from .pywalous import DlgMapWalous2Manning

            if self.active_array is None:
                logging.warning(_('No active array -- Please activate data first'))
                return

            if self.active_array.wolftype != WOLF_ARRAY_FULL_SINGLE:
                logging.error(_('Active array is not a Float32 array -- Please change it to Float32 before mapping'))
                return

            vals = self.active_array.get_unique_values()

            if self._walous_layer is None:

                if vals[0] > 10:
                    self._walous_layer = 'UTS_MAJ_NIV2'
                else:
                    self._walous_layer = 'UTS_MAJ_NIV1'

            dlg = DlgMapWalous2Manning(self, which=self._walous_layer)

            ret = dlg.ShowModal()

            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            mapvals = dlg.get_mapping()
            dlg.Destroy()

            if mapvals == -1:
                logging.error(_('Bad values -- retry'))
                return

            self.active_array.map_values(mapvals)

            self.active_array.reset_plot()

    def _add_sim_explorer(self, which:Wolfresults_2D):
        """ Add a step chooser """

        if which in self.sim_explorers:
            logging.warning(_('Step chooser already exists for this result'))
            self.sim_explorers[which].Show()
            self.sim_explorers[which].Raise()
            self.sim_explorers[which].SetFocus()
            self.sim_explorers[which].Center()
            return

        self.sim_explorers[which] = Sim_Explorer(self, which.idx, self, which)
        self.sim_explorers[which]._set_all(which.current_result)

    def _pop_sim_explorer(self, which:Wolfresults_2D):
        """ Pop a step chooser """

        if which in self.sim_explorers:
            self.sim_explorers.pop(which)
            logging.debug(_('Pop step chooser for result {}'.format(which.idx)))
        else:
            logging.warning(_('No step chooser for this result'))

    def _update_sim_explorer(self, which:Wolfresults_2D = None):

        if which is None:
            if self.active_res2d is None:
                logging.warning(_('No active 2D result -- Please activate a 2D result first'))
                return

            which = self.active_res2d

        if which in self.sim_explorers:
            self.sim_explorers[which]._set_all(which.current_result)

    def Onmenudrowning(self, event: wx.MenuEvent):

        id = event.GetId()
        item = self.menubar.FindItemById(event.GetId())

        if item is None:
            return

        itemlabel = item.ItemLabel


        if itemlabel == _("Read last result"):

            self.active_drowning.read_last_result()
            self.Refresh()
            self._update_mytooltip()

        elif itemlabel == _("Explore time/index results"):
            if self.active_drowning is None:
                logging.warning(_('No active drowning ! -- Please activate a drowning first'))
                return
            which = self.active_drowning
            self.sim_explorers[which] = Drowning_Explorer(self,title=f'Explore drowning results: {self.active_drowning.idx}',mapviewer=self, sim=which)

        elif itemlabel == _("Plot runs positions"):
            if self.plot_runs.IsChecked():
                self.active_drowning.prepare_plot_runs_positions()
            else:
                self.active_drowning.reset_plot_runs_positions()
            self._update_mytooltip()
            self.Refresh()


        elif itemlabel == _("Plot cells positions"):
            if self.plot_cells.IsChecked():
                self._update_mytooltip()
                self.active_drowning.prepare_plot_cells_positions()
            else:
                self.active_drowning.reset_plot_cells_positions()
            self._update_mytooltip()
            self.Refresh()

        elif itemlabel == _("Plot KDE"):
            if self.plot_KDE.IsChecked():
                self.active_drowning.prepare_plot_kde()
            else:
                self.active_drowning.reset_plot_kde()
            self._update_mytooltip()
            self.Refresh()

        elif itemlabel == _("Zoom on hotspots"):
            self.memory_views = Memory_Views()
            self.active_drowning.zoom_on_hotspots(self.memory_views)
            self._memory_views_gui = Memory_Views_GUI(self, _('Memory view manager'), self.memory_views, mapviewer = self)

        elif itemlabel == _("Get bodies characteristics"):
            self.active_drowning.get_bodies_characteristics()

        elif itemlabel == _("Vertical position proportion"):
            self.active_drowning.get_vertical_position_proportion()

        elif itemlabel == _("Create video..."):
            logging.info(_("Not yet implemeted"))
            return


    def Onmenudike(self, event: wx.MenuEvent):

        if not WOLFPYDIKE_AVAILABLE:
            raise ImportError('WolfPyDike not installed -- Please install WolfPyDike from pip to use this feature')

        id = event.GetId()
        item = self.menubar.FindItemById(event.GetId())

        if item is None:
            return

        itemlabel = item.ItemLabel

        if self.active_dike is None:
            logging.warning(_('No active dike -- Please activate a dike first'))
            return

        class plot_types(Enum):
            """
            Enum class for plot types.
            """
            DISCHARGES = 0
            LEVELS = 1
            BREACHWIDTH = 2

        if itemlabel == _("Launch lumped simulation"):

            self.active_dike.run_lumped()

        elif itemlabel == _("Set injector"):

            self.active_dike.set_injector()

        elif itemlabel == _("Launch 2D-coupled simulation"):

            self.active_dike.run_2Dcoupled()

        elif itemlabel == _("Show triangulation"):

            self.active_dike.show_triangulation()

        elif itemlabel == _("Plot discharges"):
            self.active_dike.plot_mainOutputs(plot_types.DISCHARGES.value)

        elif itemlabel == _("Plot water levels/breach bottom"):
            self.active_dike.plot_mainOutputs(plot_types.LEVELS.value)

        elif itemlabel == _("Plot breach width"):
            self.active_dike.plot_mainOutputs(plot_types.BREACHWIDTH.value)

        elif itemlabel == _("Show parameters"):
            self.active_dike.show_properties()


    def Onmenuwolf2d(self, event: wx.MenuEvent):

        id = event.GetId()
        item = self.menubar.FindItemById(event.GetId())

        if item is None:
            return

        itemlabel = item.ItemLabel


        if itemlabel == _("Read last result"):

            self.read_last_result()

        elif itemlabel == _("Export results as..."):

            self.export_results_as()

        elif itemlabel == _("Explore time/index results"):
            if self.active_res2d is None:
                logging.warning(_('No active 2D result ! -- Please activate a 2D result first'))
                return

            self._add_sim_explorer(self.active_res2d)

        elif itemlabel == _("Change current view"):

            # Change view for results

            autoscale = False
            choices = [cur.value for cur in views_2D]
            dlg = wx.SingleChoiceDialog(None, _("Pick a view"), "Choices", choices)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            method = dlg.GetStringSelection()

            method = list(views_2D)[choices.index(method)]

            dlg.Destroy()

            diamsize = None
            if method == views_2D.SHIELDS_NUMBER :

                if self.active_res2d is not None:
                    sediment_diam = self.active_res2d.sediment_diameter
                    sediment_density = self.active_res2d.sediment_density
                elif self.compare_results is not None:
                    sediment_diam = 0.001
                    sediment_density = 2.650
                else:
                    logging.warning(_('No active 2D result or comparison !'))
                    return

                dlg = wx.TextEntryDialog(None,_("Diameter grain size [m] ?"), value = str(sediment_diam))
                ret = dlg.ShowModal()
                if ret == wx.ID_CANCEL:
                    dlg.Destroy()
                    return
                try:
                    diamsize = float(dlg.GetValue())
                except:
                    dlg.Destroy()
                    logging.warning(_("Bad value -- Rety"))
                    return

                dlg = wx.TextEntryDialog(None,_("Density grain [-] ?"), value = str(sediment_density))
                ret = dlg.ShowModal()
                if ret == wx.ID_CANCEL:
                    dlg.Destroy()
                    return
                try:
                    density = float(dlg.GetValue())
                except:
                    dlg.Destroy()
                    logging.warning(_("Bad value -- Rety"))
                    return

            if len(self.myres2D)>1:

                dlg = wx.MessageDialog(None, _('Apply to all results?'), style=wx.YES_NO)
                ret = dlg.ShowModal()
                if ret == wx.ID_NO:
                    if diamsize is not None:
                        self.active_res2d.sediment_diameter = diamsize
                        self.active_res2d.sediment_density = density
                        self.active_res2d.load_default_colormap('shields_cst')

                    self.active_res2d.set_currentview(method, force_wx = True, force_updatepal = True)
                else:
                    for curarray in self.iterator_over_objects(draw_type.RES2D):
                        curarray:Wolfresults_2D
                        if diamsize is not None:
                            curarray.sediment_diameter = diamsize
                            curarray.sediment_density  = density
                            curarray.load_default_colormap('shields_cst')

                        curarray.set_currentview(method, force_wx = True, force_updatepal = True)

            else:
                if self.active_res2d is not None:
                    if diamsize is not None:
                        self.active_res2d.sediment_diameter = diamsize
                        self.active_res2d.sediment_density = density
                        self.active_res2d.load_default_colormap('shields_cst')
                    self.active_res2d.set_currentview(method, force_wx = True, force_updatepal = True)

            if self.compare_results is not None:
                # update compare results
                if diamsize is not None:
                    self.compare_results.set_shields_param(diamsize, density)
                self.compare_results.update_type_result(method)

        elif itemlabel == _("Set epsilon water depth"):

            dlg = wx.TextEntryDialog(self, _('Enter an epsilon [m]'),value='0.0')

            ret = dlg.ShowModal()

            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            try:
                neweps = float(dlg.GetValue())
                dlg.Destroy()
            except:
                logging.error(_('Bad value -- retry !'))
                dlg.Destroy()
                return

            pgbar = wx.ProgressDialog(_('Setting epsilon'), _('Setting epsilon'), maximum = len(self.myres2D), parent=self, style = wx.PD_APP_MODAL | wx.PD_AUTO_HIDE)

            for id, curmodel in enumerate(self.iterator_over_objects(draw_type.RES2D)):
                curmodel: Wolfresults_2D
                curmodel.epsilon = neweps
                curmodel._epsilon_default = neweps
                curmodel.read_oneresult(curmodel.current_result)
                curmodel.set_currentview()

                pgbar.Update(id, _('Setting epsilon for result {}'.format(curmodel.idx)))

            pgbar.Destroy()

        elif itemlabel == _("Filter independent"):

            self.menu_filter_independent.IsChecked = not self.menu_filter_independent.IsChecked

            pgbar = wx.ProgressDialog(_('Filtering independent zones'), _('Filtering independent zones'), maximum = len(self.myres2D), parent=self, style = wx.PD_APP_MODAL | wx.PD_AUTO_HIDE)

            for id, curmodel in enumerate(self.iterator_over_objects(draw_type.RES2D)):
                curmodel: Wolfresults_2D
                curmodel.to_filter_independent = not self.menu_filter_independent.IsChecked

                pgbar.Update(id, _('Filtering independent zones for result {}'.format(curmodel.idx)))

            pgbar.Destroy()

        # elif itemlabel == _("Manage boundary conditions..."):

        #     if self.active_res2d is not None:
        #         self.active_res2d.myparams.editing_bc(self.myres2D)

        elif itemlabel ==_("Create video..."):
            if self.active_res2d is not None:
                self.create_video()

        elif itemlabel == _("Danger map - only h"):
            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            with wx.NumberEntryDialog(None, _('Danger map'), _('From step'), _('Danger map'), 1, 1, self.active_res2d.get_nbresults()) as dlg:

                if dlg.ShowModal() == wx.ID_CANCEL:
                    return

                start_step = dlg.GetValue()

            with wx.NumberEntryDialog(None, _('Danger map'), _('To step'), _('Danger map'), self.active_res2d.get_nbresults(), start_step, self.active_res2d.get_nbresults()) as dlg:

                if dlg.ShowModal() == wx.ID_CANCEL:
                    return

                end_step = dlg.GetValue()

            with wx.NumberEntryDialog(None, _('Danger map'), _('Every'), _('Danger map'), 1, 1, 60) as dlg:

                if dlg.ShowModal() == wx.ID_CANCEL:
                    return

                every = dlg.GetValue()

            with wx.NumberEntryDialog(None, _('Danger map'),
                                      _('Minimum water depth [mm]'),
                                      _('Danger map'),
                                      int(self.active_res2d.epsilon * 1000), 1, 1000) as dlg:

                if dlg.ShowModal() == wx.ID_CANCEL:
                    return

                hmin = float(dlg.GetValue()) / 1000.

            danger_map = self.active_res2d.danger_map_only_h(start_step-1, end_step-1, every, hmin)

            with wx.DirDialog(None, _('Choose a directory'), style=wx.DD_DEFAULT_STYLE) as dlg:

                if dlg.ShowModal() == wx.ID_CANCEL:
                    return

                outdir = dlg.GetPath()

            danger_map.write_all(Path(outdir) / 'danger_h.tif')

        elif itemlabel in [_("Danger map"), _("Danger map (multiprocess)"), _("Danger map tiled")]:

            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            if itemlabel == _("Danger map tiled") and not isinstance(self.active_res2d, wolfres2DGPU):
                logging.error(_('Tiled danger map is only available for wolfres2DGPU results !'))
                return

            with wx.NumberEntryDialog(None, _('Danger map'), _('From step'), _('Danger map'), 1, 1, self.active_res2d.get_nbresults()) as dlg:

                if dlg.ShowModal() == wx.ID_CANCEL:
                    return

                start_step = dlg.GetValue()

            with wx.NumberEntryDialog(None, _('Danger map'), _('To step'), _('Danger map'), self.active_res2d.get_nbresults(), start_step, self.active_res2d.get_nbresults()) as dlg:

                if dlg.ShowModal() == wx.ID_CANCEL:
                    return

                end_step = dlg.GetValue()

            with wx.NumberEntryDialog(None, _('Danger map'), _('Every'), _('Danger map'), 1, 1, 60) as dlg:

                if dlg.ShowModal() == wx.ID_CANCEL:
                    return

                every = dlg.GetValue()

            with wx.NumberEntryDialog(None, _('Danger map'),
                                      _('Minimum water depth [mm]'),
                                      _('Danger map'),
                                      int(self.active_res2d.epsilon * 1000), 1, 1000) as dlg:

                if dlg.ShowModal() == wx.ID_CANCEL:
                    return

                hmin = float(dlg.GetValue()) / 1000.

            if itemlabel == _("Danger map"):
                logging.info(_('Danger map -- Be patient !'))
                pgbar = wx.ProgressDialog(_('Danger map'), _('Danger map'), maximum = end_step-1, parent=self, style = wx.PD_APP_MODAL | wx.PD_AUTO_HIDE)
                def _callback(id, msg):
                    pgbar.Update(id+1-start_step, msg)

                danger_maps = self.active_res2d.danger_map(start_step-1, end_step-1, every, _callback, hmin)

                pgbar.Hide()
                pgbar.Destroy()
                logging.info(_('Danger map done !'))

            elif itemlabel == _("Danger map (multiprocess)"):

                logging.info(_('Multiprocess danger map -- Be patient !'))
                danger_maps = self.active_res2d.danger_map_multiprocess(start_step-1, end_step-1, every, hmin=hmin)
                logging.info(_('Multiprocess danger map done !'))

            elif itemlabel == _("Danger map tiled"):
                logging.info(_('Tiled danger map -- Be patient !'))
                assert isinstance(self.active_res2d, wolfres2DGPU)
                pgbar = wx.ProgressDialog(_('Danger map'), _('Danger map'), maximum = end_step-1, parent=self, style = wx.PD_APP_MODAL | wx.PD_AUTO_HIDE)
                def _callback(id, msg):
                    pgbar.Update(id+1-start_step, msg)

                danger_maps = self.active_res2d.danger_map_gpu_tiled(start_step-1, end_step-1, every, _callback, hmin=hmin)
                logging.info(_('Tiled danger map done !'))

            with wx.DirDialog(None, _('Choose a directory to store results'), style=wx.DD_DEFAULT_STYLE) as dlg:

                if dlg.ShowModal() == wx.ID_CANCEL:
                    return

                outdir = dlg.GetPath()

            names = ['danger_h', 'danger_u',
                     'danger_q', 'danger_Z',
                     'danger_head',
                     'danger_toa', 'danger_tom',
                     'danger_doi', 'danger_toe']

            for name, danger_map in zip(names, danger_maps):

                if isinstance(danger_map, WolfArrayMB):
                    name = name + '.bin'
                    logging.info(_('Saving danger map {}').format(name))
                    danger_map.write_all(Path(outdir) / name)

                elif isinstance(danger_map, WolfArray):
                    name = name + '.tif'
                    logging.info(_('Saving danger map {}').format(name))
                    danger_map.write_all(Path(outdir) / name)

                else:
                    logging.error(_('Bad type for danger map {} -- not saved !').format(name))
                    continue

        elif itemlabel == _("Setup cache..."):

            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            dlg = wx.MessageDialog(None, _('Cache only water depth results ?'), style=wx.YES_NO)
            ret = dlg.ShowModal()
            if ret == wx.ID_NO:
                only_h = False
            else:
                only_h = True
            dlg.Destroy()

            dlg = wx.MessageDialog(None, _('Cache all results ?'), style=wx.YES_NO)
            ret = dlg.ShowModal()
            if ret == wx.ID_NO:

                dlg_start = wx.SingleChoiceDialog(None, _('Choosing the start index'),
                                                  _('Choices'),
                                                  [str(cur) for cur in range(1,self.active_res2d.get_nbresults()+1)])
                ret = dlg_start.ShowModal()
                if ret == wx.ID_CANCEL:
                    dlg_start.Destroy()
                    return

                start_idx = int(dlg_start.GetStringSelection())
                dlg_start.Destroy()

                dlg_end   = wx.SingleChoiceDialog(None, _('Choosing the end index'),
                                                  _('Choices'),
                                                  [str(cur) for cur in range(start_idx + 1,self.active_res2d.get_nbresults()+1)])

                ret = dlg_end.ShowModal()
                if ret == wx.ID_CANCEL:
                    dlg_end.Destroy()
                    return

                dlg_end.Destroy()

                end_idx = int(dlg_end.GetStringSelection())

                logging.info(_('Caching from {} to {} - Be patient !').format(start_idx, end_idx))
                self.active_res2d.setup_cache(start_idx = start_idx-1, end_idx = end_idx-1, only_h=only_h)
                logging.info(_('Caching done !'))
            else:
                logging.info(_('Caching all results - Be patient !'))
                self.active_res2d.setup_cache(only_h=only_h)
                logging.info(_('Caching done !'))

            dlg.Destroy()

        elif itemlabel == _("Clear cache..."):

            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            self.active_res2d.clear_cache()
            logging.info(_('Cache cleared !'))

        elif itemlabel == _("Show tiles..."):

            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            self.active_res2d.show_tiles()


    def menu_2dgpu(self):

        if self.menuwolf2d is not None:
            if self.menu2d_cache_setup  is None:
                self.menu2d_cache_setup = self.menuwolf2d.Append(wx.ID_ANY, _("Setup cache..."), _("Set up cache for 2D GPU model"))
                self.menu2d_cache_reset = self.menuwolf2d.Append(wx.ID_ANY, _("Clear cache..."), _("Clear cache for 2D GPU model"))
                self.menu2d_show_tiles = self.menuwolf2d.Append(wx.ID_ANY, _("Show tiles..."), _("Show a grid of tiles for 2D GPU model"))

    def menu_landmaps(self):

        if self.menu_landmap is None:
            self.menu_landmap = wx.Menu()
            self.menubar.Append(self.menu_landmap, _('&Landmap'))

            self.menupick_landmap_full = self.menu_landmap.Append(wx.ID_ANY, _("Pick landmap full..."), _("Pick landmap full resolution"))
            self.menupick_landmap_low = self.menu_landmap.Append(wx.ID_ANY, _("Pick landmap low..."), _("Pick landmap low resolution"))
            self.menu_landmap.AppendSeparator()

            self.menu_colortransparent_landmap = self.menu_landmap.Append(wx.ID_ANY, _("Transparent color "), _("Change transparent color associated to the landmap"))
            self.menu_tolerance_landmap = self.menu_landmap.Append(wx.ID_ANY, _("Set tolerance"), _("Set tolerance for the transparent color landmap"))
            self.menu_color_landmap = self.menu_landmap.Append(wx.ID_ANY, _("Change colors"), _("Change color map associated to the landmap"))

            self.menu_landmap.Bind(wx.EVT_MENU, self.pick_landmap_full, self.menupick_landmap_full)
            self.menu_landmap.Bind(wx.EVT_MENU, self.pick_landmap_low, self.menupick_landmap_low)
            self.menu_landmap.Bind(wx.EVT_MENU, self.change_colors_landmap, self.menu_color_landmap)
            self.menu_landmap.Bind(wx.EVT_MENU, self.change_transparent_color_landmap, self.menu_colortransparent_landmap)
            self.menu_landmap.Bind(wx.EVT_MENU, self.set_tolerance_landmap, self.menu_tolerance_landmap)

    def change_transparent_color_landmap(self, event: wx.Event):

        if self.active_landmap is None:
            logging.warning(_('No active landmap -- Please load data first'))

        data = wx.ColourData()
        data.SetColour(self.active_landmap.transparent_color)
        with wx.ColourDialog(self, data) as dlg:
            if dlg.ShowModal() == wx.ID_OK:
                data = dlg.GetColourData()
                color = data.GetColour()
                self.active_landmap.set_transparent_color([color.Red(), color.Green(), color.Blue()])

    def set_tolerance_landmap(self, event: wx.Event):

        if self.active_landmap is None:
            logging.warning(_('No active landmap -- Please load data first'))

        dlg = wx.TextEntryDialog(self, _('Set the tolerance for the transparent color'), _('Tolerance'), str(self.active_landmap.tolerance))
        ret = dlg.ShowModal()
        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return
        tol = int(dlg.GetValue())
        tol = max(0, tol)
        dlg.Destroy()

        self.active_landmap.set_tolerance(tol)


    def change_colors_landmap(self, event: wx.Event):

        if self.active_landmap is None:
            logging.warning(_('No active landmap -- Please load data first'))

        data = wx.ColourData()
        data.SetColour(self.active_landmap.color)
        with wx.ColourDialog(self, data) as dlg:
            if dlg.ShowModal() == wx.ID_OK:
                data = dlg.GetColourData()
                color = data.GetColour()
                self.active_landmap.set_color(color)

    def pick_landmap_full(self, event: wx.Event):

        self.action = 'pick landmap full'
        logging.info(_('Pick landmap - Full resolution'))

    def pick_landmap_low(self, event: wx.Event):

        self.action = 'pick landmap low'
        logging.info(_('Pick landmap - Low resolution'))

    def menu_particlesystem(self):
        if self.menuparticlesystem is None:

            self.menuparticlesystem = wx.Menu()
            self.menuparticlesystem_load = wx.Menu()

            self.menuparticlesystem.Append(wx.ID_ANY, _("Set..."), _("Set arrays as the domain/uv of the particle system -- Must be a 2D array - Mask will be used to separate water and land"))
            self.menuparticlesystem.Append(wx.ID_ANY, _("Set emitter from selected nodes"), _("Set the selected nodes as emitters of the particle system"))
            self.menuparticlesystem.AppendSubMenu(self.menuparticlesystem_load, _("Load..."),  _('Load data for the particle system in the UI'))

            self.menuparticlesystem_load.Append(wx.ID_ANY, _("Load domain..."), _("Loading the domain in the UI"))
            self.menuparticlesystem_load.Append(wx.ID_ANY, _("Load ~domain..."), _("Loading the negative of the domain in the UI"))
            self.menuparticlesystem_load.Append(wx.ID_ANY, _("Load emitters..."), _("Loading the emitters in the UI"))
            self.menuparticlesystem_load.Append(wx.ID_ANY, _("Load uv..."), _("Loading the UV velocity field in the UI"))
            self.menuparticlesystem_load.Append(wx.ID_ANY, _("Load uv norm..."), _("Loading the norm of the velocity field in the UI"))
            # self.menuparticlesystem.Append(wx.ID_ANY, _("Set emitters..."), _("Set active zones as the emitters of the particle system -- Each checked zone will be used as an emitter"))
            # self.menuparticlesystem.Append(wx.ID_ANY, _("Set emitter..."), _("Set only the active vector as an emitters of the particle system"))
            # self.menuparticlesystem.Append(wx.ID_ANY, _("Set uv..."), _("Choose U and V arrays for the particle system -- Must be 2D arrays"))
            self.menuparticlesystem.AppendSeparator()
            self.menuparticlesystem.Append(wx.ID_ANY, _("Check"), _("Check if the particle system is ready to be computed"))
            self.menuparticlesystem.Append(wx.ID_ANY, _("Bake"), _("Compute the particle system"))
            self.menuparticlesystem.Append(wx.ID_ANY, _("Reset"), _("Clear all results but keep the particle system settings"))
            self.menuparticlesystem.AppendSeparator()
            # self.menuparticlesystem.AppendSeparator()
            self.menuparticlesystem.Append(wx.ID_ANY, _("Start"), _("Run all steps"))
            self.menuparticlesystem.Append(wx.ID_ANY, _("Stop"), _("Stop the current animation"))
            self.menuparticlesystem.Append(wx.ID_ANY, _("Resume"), _("Resume animation"))

            self.timer_ps = wx.Timer(self)

            self.menuparticlesystem.Bind(wx.EVT_MENU, self.action_menu_particlesystem)
            self.Bind(wx.EVT_TIMER, self.update_particlesystem, self.timer_ps)

            self.menubar.Append(self.menuparticlesystem, _('Particle system'))

    def action_menu_particlesystem(self, event: wx.Event):
        """ Action to perform when the timer is triggered """

        if self.active_particle_system is not None:

            itemlabel = self.menuparticlesystem.FindItemById(event.GetId()).GetItemLabelText()


            if itemlabel == _("Start"):

                if self.active_particle_system is not None:
                    self.active_particle_system.current_step = 0
                    self.active_particle_system.current_step_idx = 0
                    self.timer_ps.Start(1000. / self.active_particle_system.fps)

            elif itemlabel == _("Stop"):

                self.timer_ps.Stop()

            elif itemlabel == _("Resume"):

                self.timer_ps.Start(1000. / self.active_particle_system.fps)

            elif itemlabel == _("Load domain..."):
                domain = self.active_particle_system.get_domain(output_type='wolf')
                self.add_object('array', id=domain.idx, newobj=domain, ToCheck=True)
                self.Refresh()

            elif itemlabel == _("Load ~domain..."):
                domain:WolfArray = self.active_particle_system.get_domain(output_type='wolf')
                domain.idx = domain.idx + '_neg'
                domain.mask_reset()

                ones = np.where(domain.array.data == 1)
                domain.array[:,:] = 1
                domain.array[ones] = 0

                domain.mask_data(domain.nullvalue)
                self.add_object('array', id=domain.idx, newobj=domain, ToCheck=True)
                self.Refresh()

            elif itemlabel == _("Load emitters..."):
                emitters = self.active_particle_system.get_emitters(output_type='wolf')
                self.add_object('vector', id=emitters.idx, newobj=emitters, ToCheck=True)
                self.Refresh()

            elif itemlabel == _("Load uv..."):
                u = self.active_particle_system.get_u(output_type='wolf')
                v = self.active_particle_system.get_v(output_type='wolf')
                self.add_object('array', id=u.idx, newobj=u, ToCheck=True)
                self.add_object('array', id=v.idx, newobj=v, ToCheck=True)
                self.Refresh()

            elif itemlabel == _("Load uv norm..."):
                uvnorm = self.active_particle_system.get_uv_absolute(output_type='wolf')
                self.add_object('array', id=uvnorm.idx, newobj=uvnorm, ToCheck=True)
                self.Refresh()

            elif itemlabel == _("Bake"):
                check, msg = self.active_particle_system.bake()

                if not check:
                    dlg = wx.MessageDialog(self, msg, _('Error'), wx.OK | wx.ICON_ERROR)
                    dlg.ShowModal()
                    dlg.Destroy()
                    return
            elif itemlabel == _("Reset"):
                self.active_particle_system.reset()

            elif itemlabel == _("Set..."):
                from .ui.wolf_multiselection_collapsiblepane import Wolf_MultipleSelection

                setter = Wolf_MultipleSelection(self,
                                                title=_("Set particle system"),
                                                message=_("Choose arrays/emitters for the particle system"),
                                                values_dict={'domain': self.get_list_keys(draw_type.ARRAYS),
                                                             'u': self.get_list_keys(draw_type.ARRAYS),
                                                             'v': self.get_list_keys(draw_type.ARRAYS),
                                                             'emitters': self.get_list_keys(draw_type.VECTORS)},
                                                info='Set : \n - domain (1 value)\n - u and v (multiple values)\n - emitters (multiple values)',
                                                styles=[wx.LB_SINGLE, wx.LB_EXTENDED, wx.LB_EXTENDED, wx.LB_EXTENDED],
                                                max_choices=[1, None, None, None],
                                                delete_if_transfer = [True, False, False, True],
                                                destroyOK=False)
                setter.ShowModal()
                ret_dict = setter.get_values()
                setter.Destroy()

                if 'domain' in ret_dict:
                    if len(ret_dict['domain']) == 1:
                        domain = self.getobj_from_id(ret_dict['domain'][0])
                        self.active_particle_system.set_domain(domain)
                if 'u' in ret_dict and 'v' in ret_dict:
                    if len(ret_dict['u']) >0:
                        assert len(ret_dict['u']) == len(ret_dict['v']), _('Please select the same number of u and v arrays')

                        time = 0.
                        for u,v in zip(ret_dict['u'], ret_dict['v']):
                            u = self.getobj_from_id(u)
                            v = self.getobj_from_id(v)
                            u:WolfArray
                            v:WolfArray
                            assert u.array.shape == v.array.shape, _('Please select arrays with the same shape')
                            assert u.origx == v.origx and u.origy == v.origy, _('Please select arrays with the same origin')
                            assert u.dx == v.dx and u.dy == v.dy, _('Please select arrays with the same resolution')
                            self.active_particle_system.set_uv((u, v),
                                                            (u.origx, u.origy, u.dx, u.dy),
                                                            time = time)
                            time += 1.

                if 'emitters' in ret_dict:
                    if len(ret_dict['emitters'])>0:
                        emitters = [self.getobj_from_id(cur) for cur in ret_dict['emitters']]
                        self.active_particle_system.set_emitters(emitters)

                if self.active_particle_system._ui is not None:
                    self.active_particle_system.show_properties()

            elif itemlabel == _("Set emitter from selected nodes"):
                if self.active_array is None:
                    logging.warning(_('No active array -- Please activate an array first'))
                    return
                if len(self.active_array.SelectionData.myselection) == 0 and len(self.active_array.SelectionData.selections) ==0:
                    logging.warning(_('No selection -- Please select some nodes first'))
                    return

                from .lagrangian.emitter import Emitter

                newemitters=[]
                if len(self.active_array.SelectionData.myselection) > 0:
                    indices = [self.active_array.get_ij_from_xy(cur[0], cur[1]) for cur in self.active_array.SelectionData.myselection]
                    newemitters = [Emitter(indices,
                                         header = (self.active_array.origx, self.active_array.origy, self.active_array.dx, self.active_array.dy))]

                if len(self.active_array.SelectionData.selections) > 0:

                    for cursel in self.active_array.SelectionData.selections.values():
                        indices = [self.active_array.get_ij_from_xy(cur[0], cur[1]) for cur in cursel['select']]
                        newemitters += [Emitter(indices, header = (self.active_array.origx, self.active_array.origy, self.active_array.dx, self.active_array.dy))]

                self.active_particle_system.set_emitters(newemitters)

                if self.active_particle_system._ui is not None:
                    self.active_particle_system.show_properties()

            # elif itemlabel == _("Set emitters..."):

            #     if self.active_zones is None:
            #         logging.warning(_('No active zones -- Please activate zones first'))
            #         return

            #     self.active_particle_system.set_emitters(self.active_zones)

            # elif itemlabel == _("Set emitter..."):

            #     if self.active_vector is None:
            #         logging.warning(_('No active vector -- Please activate a vector first'))
            #         return

            #     self.active_particle_system.set_emitter(self.active_vector)

            # elif itemlabel == _("Set uv..."):

            #     list_arrays = self.multiple_choice_object(draw_type.ARRAYS, message=_('Choose U and V arrays for the particle system -- first == u ; second == v'), titel='UV choice' )

            #     if len(list_arrays) != 2:
            #         logging.error(_('Please select two arrays and ONLY two arrays'))
            #         return

            #     self.active_particle_system.set_uv(tuple(list_arrays))

            elif itemlabel == _("Check"):
                check, msg = self.active_particle_system.check()

                if not check:
                    dlg = wx.MessageDialog(self, msg, _('Error'), wx.OK | wx.ICON_ERROR)
                    dlg.ShowModal()
                    dlg.Destroy()
                    return
                else:
                    dlg = wx.MessageDialog(self, _('All is fine -- You can bake you system !'), _('Chesk particle system'), wx.OK | wx.ICON_INFORMATION)
                    dlg.ShowModal()
                    dlg.Destroy()
                    return

    def update_particlesystem(self, event: wx.Event):
        """ Animation of the particle system """

        if self.active_particle_system is not None:

            nb = self.active_particle_system.nb_steps
            self.active_particle_system.current_step_idx += 1
            self.Paint()
            self._update_mytooltip()

            if self.active_particle_system.current_step_idx == nb-1:
                self.timer_ps.Stop()

    def menu_sim2D(self):
        """ Menu for 2D simulations """

        if self.menusim2D is None:
            self.menusim2D = wx.Menu()
            self.menubar.Append(self.menusim2D, _('Tools 2D'))

            menu2d_options = self.menusim2D.Append(wx.ID_ANY, _("Parameters..."), _("Parameters"))
            menu2d_zbin2hbin = self.menusim2D.Append(wx.ID_ANY, _("Convert zbin to hbin"), _("Convert zbin to hbin"))
            menu2d_hbin2zbin = self.menusim2D.Append(wx.ID_ANY, _("Convert hbin to zbin"), _("Convert hbin to zbin"))
            menu2D_zbinb2hbinb = self.menusim2D.Append(wx.ID_ANY, _("Convert zbinb to hbinb"), _("Convert zbinb to hbinb"))
            menu2d_hbinb2zbinb = self.menusim2D.Append(wx.ID_ANY, _("Convert hbinb to zbinb"), _("Convert hbinb to zbinb"))
            menu2d_forcemask = self.menusim2D.Append(wx.ID_ANY, _("Reset mask of all arrays"), _("Reset mask"))

            # update = self.menusim2D.Append(wx.ID_ANY, _('Update model from current mask'), _('Update model'))
            # updateblocfile = self.menusim2D.Append(wx.ID_ANY, _('Update .bloc file'), _('Update bloc'))
            # updatefreesurface = self.menusim2D.Append(wx.ID_ANY, _('Update free surface elevation - IC'), _('Update free surface elevation'))
            # updaterough = self.menusim2D.Append(wx.ID_ANY, _('Update roughness coeff'), _('Update roughness coefficient'))
            # updateic = self.menusim2D.Append(wx.ID_ANY, _('Update IC reading mode'), _('Update IC'))
            # menu2d_tft_ic = self.menusim2D.Append(wx.ID_ANY,_("Transfer initial conditions..."),_("Transfer IC"))

            self.menusim2D.Bind(wx.EVT_MENU, self.Onmenusim2D)

    def menu_sim2DGPU(self):
        """ Menu for 2D GPU simulations """

        if self.menusim2D_GPU is None:
            self.menusim2D_GPU = wx.Menu()
            self.menubar.Append(self.menusim2D_GPU, _('Tools 2D GPU'))

            menu2dGPU_options = self.menusim2D_GPU.Append(wx.ID_ANY, _("Parameters..."), _("Parameters"))

            self.menusim2D_GPU.Bind(wx.EVT_MENU, self.Onmenusim2DGPU)

    def Onmenusim2DGPU(self, event: wx.MenuEvent):
        """ Action to perform when menu 2D GPU entry is selected """

        id = event.GetId()
        item = self.menubar.FindItemById(event.GetId())

        if item is None:
            return

        itemlabel = item.ItemLabel

        from .PyGui import Wolf2DGPUModel

        if not isinstance(self.wolfparent, Wolf2DGPUModel):
            logging.error(_('This is not a 2D GPUmodel'))
            return

        self.wolfparent:Wolf2DGPUModel

        if itemlabel == _("Parameters..."):
            self.wolfparent.show_properties()

    def Onmenusim2D(self, event: wx.MenuEvent):
        """ Action to perform when menu 2D entry is selected """

        id = event.GetId()
        item = self.menubar.FindItemById(event.GetId())

        if item is None:
            return

        itemlabel = item.ItemLabel

        from .PyGui import Wolf2DModel

        if not isinstance(self.wolfparent, Wolf2DModel):
            logging.error(_('This is not a 2D model'))
            return

        self.wolfparent:Wolf2DModel

        if itemlabel == _('Update .bloc file'):

            msg = _('If you continue the .bloc file will be relpaced !')+'\n'
            msg += '\n'
            msg += _('Continue ?')+'\n'

            dlg = wx.MessageDialog(self,msg,caption = _('Attention'), style = wx.YES_NO)
            ret = dlg.ShowModal()
            dlg.Destroy()
            if ret == wx.ID_NO:
                return

            self.wolfparent.write_bloc_file()

        elif itemlabel == _('Reset mask of all arrays'):

            self.wolfparent.sim.force_mask()

        elif itemlabel == _('Convert zbin to hbin'):

            if self.wolfparent.sim._zbin is not None:
                self.wolfparent.sim.zbin2hbin()
                self.wolfparent.sim.hbin.reset_plot()

        elif itemlabel == _('Convert hbin to zbin'):

            if self.wolfparent.sim._hbin is not None:
                self.wolfparent.sim.hbin2zbin()
                self.wolfparent.sim.zbin.reset_plot()

        elif itemlabel == _('Convert zbinb to hbinb'):

            if self.wolfparent.sim._zbinb is not None:
                self.wolfparent.sim.zbinb2hbinb()
                self.wolfparent.sim.hbinb.reset_plot()

        elif itemlabel == _('Convert hbinb to zbinb'):

            if self.wolfparent.sim._hbinb is not None:
                self.wolfparent.sim.hbinb2zbinb()
                self.wolfparent.sim.zbinb.reset_plot()

        elif itemlabel == _("Transfer initial conditions..."):

            if self.active_array is not None:
                from .PyGui import Wolf2DModel
                if isinstance(self.wolfparent,Wolf2DModel):
                    self.wolfparent.transfer_ic(self.active_vector)

        elif itemlabel == _("Parameters..."):
            self.wolfparent.show_properties()

        elif itemlabel == _('Update free surface elevation - IC'):

            if len(self.active_array.SelectionData.myselection)==0:

                msg = _('There is none selected nodes in the active array !')+'\n'
                msg += '\n'
                msg += _('Please select the desired zone and retry !')+'\n'

                logging.warning(msg)
                return

            self.wolfparent.extend_freesurface_elevation(self.active_array.SelectionData.myselection)

        elif itemlabel== _('Update roughness coeff'):

            if len(self.active_array.SelectionData.myselection)==0:

                msg = _('There is none selected nodes in the active array !')+'\n'
                msg += '\n'
                msg += _('Please select the desired zone and retry !')+'\n'

                logging.warning(msg)
                return

            self.wolfparent.extend_roughness(self.active_array.SelectionData.myselection)

        # elif itemlabel == _('Update IC reading mode'):

        #     self.wolfparent.set_type_ic()

        elif itemlabel == _('Update model from current mask'):

            if type(self.active_array) not in [WolfArray]:
                msg = _('Please select a mono-block array !')+'\n'
                dlg=wx.MessageBox(msg,style=wx.OK)
                return

            msg = _('If you continue, the mask of all arrays will be replaced by the current mask !')+'\n'
            msg += _('The external contour in the .bloc file will also be relpaced.')+'\n'
            msg += '\n'
            msg += _('Continue ?')+'\n'

            dlg = wx.MessageDialog(self,msg,caption = _('Attention'), style = wx.YES_NO)
            ret = dlg.ShowModal()
            dlg.Destroy()
            if ret == wx.ID_NO:
                return

            with wx.BusyInfo(_('Updating 2D model')):
                wait = wx.BusyCursor()

                sux,suy,cont,interior = self.active_array.suxsuy_contour(self.wolfparent.filenamegen,True)

                self.wolfparent.mimic_mask(self.active_array)
                self.wolfparent.replace_external_contour(cont,interior)

                del wait

            self.wolfparent.extend_bed_elevation()

    def get_configuration(self) -> Union[WolfConfiguration, None]:
        """ Get global configuration parameters """

        # At this point, I'm not too sure about
        # which window/frame does what. So to be on
        # the safe side, I make sure that the configuration
        # menu is active only on the "first" window.
        # Moreover, I try to go up the frame/window
        # hierarchy to get the configuration (which will therefore
        # be treated as a singleton)
        if self.wolfparent:
            return self.wolfparent.get_configuration()
        else:
            return None

    @property
    def epsg(self) -> int:
        """ Return the EPSG code from configs """
        config = self.get_configuration()
        if config is None:
            logging.debug(_('No configuration found -- Using default EPSG:31370'))
            return 31370  # Default EPSG code - Lambert 1970
        else:
            strcode = config[ConfigurationKeys.EPSG_CODE]
            try:
                code = int(strcode.lower().replace('epsg:', ''))
                return code
            except:
                logging.error(_('Bad EPSG code in configuration -- Using default EPSG:31370'))
                return 31370  # Default EPSG code - Lambert 1970

    @property
    def active_vector_color(self) -> list[int]:
        """ Return the active vector color from configs """
        config = self.get_configuration()
        if config is None:
            return [0, 0, 0, 255]  # Default black color
        else:
            return config[ConfigurationKeys.ACTIVE_VECTOR_COLOR]

    @property
    def active_vector_square_size(self) -> list[int]:
        """ Return the active vector square size from configs """
        config = self.get_configuration()
        if config is None:
            return 0
        else:
            return config[ConfigurationKeys.ACTIVE_VECTOR_SIZE_SQUARE]

    @property
    def default_dem(self) -> Path:
        """ Return the default DEM file from configs """
        config = self.get_configuration()
        if config is None:
            return Path('')
        else:
            return Path(config[ConfigurationKeys.DIRECTORY_DEM])

    @property
    def default_dtm(self) -> Path:
        """ Return the default DTM file from configs """
        config = self.get_configuration()
        if config is None:
            return Path('')
        else:
            return Path(config[ConfigurationKeys.DIRECTORY_DTM])

    @property
    def default_laz(self):
        """ Return the default LAZ file from configs """
        config = self.get_configuration()
        if config is None:
            return Path('')
        else:
            return Path(config[ConfigurationKeys.DIRECTORY_LAZ])

    @property
    def default_hece_database(self) -> Path:
        """ Return the default HECE database file from configs """
        config = self.get_configuration()
        if config is None:
            return Path('')
        else:
            return Path(config[ConfigurationKeys.XLSX_HECE_DATABASE])

    @property
    def bkg_color(self):
        """ Return the background color from configs """
        config = self.get_configuration()
        if config is None:
            return [255.,255.,255.,255.]
        else:
            return config[ConfigurationKeys.COLOR_BACKGROUND]

    @property
    def ticks_size(self) -> float:
        """ Return the ticks spacing from configs """

        config = self.get_configuration()
        if config is None:
            return 100.
        else:
            return config[ConfigurationKeys.TICKS_SIZE]

    @property
    def ticks_xrotation(self) -> float:
        """ Return the ticks x rotation from configs """

        config = self.get_configuration()
        if config is None:
            return 30.
        else:
            return config[ConfigurationKeys.TICKS_XROTATION]

    @property
    def ticks_fontsize(self) -> int:
        """ Return the ticks font size from configs """

        config = self.get_configuration()
        if config is None:
            return 14
        else:
            return config[ConfigurationKeys.TICKS_FONTSIZE]

    @property
    def assembly_mode(self) -> str:
        """ Return the assembly mode from configs """

        config = self.get_configuration()
        if config is None:
            return 'square'
        else:
            return config[ConfigurationKeys.ASSEMBLY_IMAGES]

    @property
    def ticks_bounds(self) -> bool:
        """ Return the ticks bounds from configs """

        config = self.get_configuration()
        if config is None:
            return True
        else:
            return config[ConfigurationKeys.TICKS_BOUNDS]

    @property
    def palette_for_copy(self) -> wolfpalette:
        """ Return the palette for copy from configs """

        config = self.get_configuration()
        if config is None:
            if self.active_array is not None:
                return self.active_array.palette
            elif self.active_res2d is not None:
                return self.active_res2d.palette
            else:
                return wolfpalette()
        else:
            act_array = config[ConfigurationKeys.ACTIVE_ARRAY_PALETTE_FOR_IMAGE]
            act_res2d = config[ConfigurationKeys.ACTIVE_RES2D_PALETTE_FOR_IMAGE]

            if act_array:
                if self.active_array is not None:
                    return self.active_array.mypal
                else:
                    if self.active_res2d is not None:
                        logging.warning(_('No active array -- Using active 2D result palette instead'))
                        return self.active_res2d.mypal
                    else:
                        return wolfpalette()
            elif act_res2d:
                if self.active_res2d is not None:
                    return self.active_res2d.mypal
                else:
                    if self.active_array is not None:
                        logging.warning(_('No active 2D result -- Using active array palette instead'))
                        return self.active_array.mypal
                    else:
                        return wolfpalette()
            else:
                return wolfpalette()

    def GlobalOptionsDialog(self, event):
        handle_configuration_dialog(self, self.get_configuration())

    # def import_3dfaces(self):

    #     dlg = wx.FileDialog(None, _('Choose filename'),
    #                         wildcard='dxf (*.dxf)|*.dxf|gltf (*.gltf)|*.gltf|gltf binary (*.glb)|*.glb|All (*.*)|*.*', style=wx.FD_OPEN)
    #     ret = dlg.ShowModal()
    #     if ret == wx.ID_CANCEL:
    #         dlg.Destroy()
    #         return

    #     fn = dlg.GetPath()
    #     dlg.Destroy()

    #     mytri = Triangulation(plotted=True,mapviewer=self)

    #     if fn.endswith('.dxf'):
    #         mytri.import_dxf(fn)
    #     elif fn.endswith('.gltf') or fn.endswith('.glb'):
    #         mytri.import_from_gltf(fn)

    #     self.add_object('triangulation',newobj=mytri,id=fn)
    #     self.active_tri = mytri

    def triangulate_cs(self):
        """ Triangulate the active cross sections """

        msg = ''
        if self.active_zones is None:
            msg += _(' The active zones is None. Please activate the desired object !\n')
        if self.active_cs is None:
            msg += _(' The is no cross section. Please active the desired object or load file !')

        if msg != '':
            logging.warning(msg)
            dlg = wx.MessageBox(msg, 'Required action')
            return

        dlg = wx.NumberEntryDialog(None, _('What is the desired size [cm] ?'), 'ds', 'ds size', 100, 1, 10000)
        ret = dlg.ShowModal()
        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return

        ds = float(dlg.GetValue()) / 100.
        dlg.Destroy()

        self.set_interp_cs(Interpolators(self.active_zones, self.active_cs, ds))

    def set_interp_cs(self, obj:Interpolators, add_zones:bool = True):
        """ Set the active cross-sections interpolator """

        assert isinstance(obj, Interpolators), _('Please provide an Interpolators object')

        self.myinterp = obj

        if add_zones:
            self.add_object('vector', newobj=self.myinterp.myzones, ToCheck=False, id='Interp_mesh')

        if self.menuviewerinterpcs is None:
            self.menuviewerinterpcs = self.cs_menu.Append(wx.ID_ANY, _("New cloud Viewer..."),
                                                            _("Cloud viewer Interpolate"))
        if self.menuinterpcs is None:
            self.menuinterpcs = self.cs_menu.Append(wx.ID_ANY, _("Interpolate on active array..."), _("Interpolate"))

        self.Refresh()

    def interpolate_cloud(self):
        """
        Interpolation d'un nuage de point sur une matrice

        Il est possible d'utiliser une autre valeur que la coordonnées Z des vertices
        """
        if self.active_cloud is not None and self.active_array is not None:

            keyvalue='z'
            if self.active_cloud.has_values:
                choices = list(self.active_cloud.myvertices[0].keys())
                dlg = wx.SingleChoiceDialog(None, "Pick the value to interpolate", "Choices", choices)
                ret = dlg.ShowModal()
                if ret == wx.ID_CANCEL:
                    dlg.Destroy()
                    return

                keyvalue = dlg.GetStringSelection()
                dlg.Destroy()

            choices = ["nearest", "linear", "cubic"]
            dlg = wx.SingleChoiceDialog(None, "Pick an interpolate method", "Choices", choices)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            method = dlg.GetStringSelection()
            dlg.Destroy()

            self.active_cloud.interp_on_array(self.active_array,keyvalue,method)

    def interpolate_cs(self):
        """ Interpolate the active cross sections by interpolators """

        if self.active_array is None:
            logging.warning(_('No active array -- Please activate an array first'))
            return

        if self.myinterp is None:
            logging.warning(_('No active interpolator -- Please create an interpolator first'))
            return

        choices = ["nearest", "linear", "cubic"]
        dlg = wx.SingleChoiceDialog(None, "Pick an interpolate method", "Choices", choices)
        ret = dlg.ShowModal()
        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return

        method = dlg.GetStringSelection()
        dlg.Destroy()

        self.myinterp.interp_on_array(self.active_array, method)

    def interpolate_triangulation(self, keep:Literal['all', 'above', 'below'] = 'all'):
        """ Alias to interpolate on triangulation

        :param keep: 'all' to keep all points, 'above' to keep only points above the current array's value, 'below' to keep only points below the current array's value
        """

        if self.active_array is None:
            logging.warning(_('No active array -- Please activate an array first'))
            return

        if self.active_tri is None:
            logging.warning(_('No active triangulation -- Please activate a triangulation first'))
            return

        self.active_array.interpolate_on_triangulation(self.active_tri.pts, self.active_tri.tri, keep=keep)

    def compare_cloud2array(self):
        """
        Compare the active cloud points to the active array

        """

        if self.active_array is None :
            logging.warning(_('No active array -- Please activate an array first'))
            return

        if self.active_cloud is None:
            logging.warning(_('No active cloud -- Please activate a cloud first'))
            return

        self.active_array.compare_cloud(self.active_cloud)

    def compare_tri2array(self):

        if self.active_array is not None and self.active_tri is not None:

            self.active_array.compare_tri(self.active_tri)

    def move_triangles(self):
        """ Move the active triangles """
        if self.active_tri is None:
            logging.warning(_('No active triangles -- Please activate triangles first'))
            return
        self.start_action('move triangles', 'Move the current triangulation -- Please select 2 points to define the translation vector')

    def rotate_triangles(self):
        """ Rotate the active triangles """
        if self.active_tri is None:
            logging.warning(_('No active triangles -- Please activate triangles first'))
            return
        self.start_action('rotate triangles', 'Rotate the current triangulation -- Please select 1 point for the center')


    def display_canvasogl(self, mpl =True,
                          ds=0., fig: Figure = None, ax: Axes = None,
                          clear = True, redraw =True, palette=False, title=''):
        """
        This method takes a matplotlib figure and axe and,
        returns a clear screenshot of the information displayed in the wolfpy GUI.
        """

        self.Paint()
        myax = ax
        if redraw:
            if clear:
                myax.clear()


        if self.SetCurrentContext():
            glPixelStorei(GL_PACK_ALIGNMENT, 1)
            data = glReadPixels(0,0,self.canvaswidth, self.canvasheight, GL_RGBA,GL_UNSIGNED_BYTE)
            myimage: Image.Image
            myimage =  Image.frombuffer("RGBA",(self.canvaswidth,self.canvasheight),data)
            myimage = myimage.transpose(1)

            if mpl:
                if ds ==0.:
                    ds = self.ticks_size

                extent = (self.xmin, self.xmax, self.ymin, self.ymax)

                myax.imshow(myimage, origin ='upper', extent=extent)

                x1 = np.ceil((self.xmin//ds)*ds)
                if x1 < self.xmin:
                    x1 += ds
                x2 = int((self.xmax//ds)*ds)
                if x2 >self.xmax:
                    x2 -= ds
                y1 = np.ceil((self.ymin//ds)*ds)
                if y1 < self.ymin:
                    y1 += ds
                y2 = int((self.ymax // ds) * ds)
                if y2 > self.ymax:
                    y2 -= ds

                x_label_list = np.linspace(x1,x2, int((x2-x1)/ds) +1, True)
                if self.ticks_bounds:
                    x_label_list = np.insert(x_label_list,0,self.xmin)
                    x_label_list = np.insert(x_label_list,-1, self.xmax)
                    x_label_list = np.unique(x_label_list)

                y_label_list = np.linspace(y1, y2, int((y2 - y1) / ds) + 1, True)
                if self.ticks_bounds:
                    y_label_list = np.insert(y_label_list, 0, self.ymin)
                    y_label_list = np.insert(y_label_list, -1, self.ymax)
                    y_label_list = np.unique(y_label_list)

                myax.set_xticks(x_label_list)
                myax.set_yticks(y_label_list)

                myax.set_xticklabels(FormatStrFormatter('%.1f').format_ticks(x_label_list),
                                     fontsize = self.ticks_fontsize, rotation = self.ticks_xrotation)
                myax.set_yticklabels(FormatStrFormatter('%.1f').format_ticks(y_label_list),
                                     fontsize = self.ticks_fontsize)
                myax.xaxis.set_ticks_position('top')
                myax.xaxis.set_label_position('top')

                myax.set_xlabel('X ($m$)')
                myax.set_ylabel('Y ($m$)')
                myax.xaxis.set_ticks_position('bottom')
                myax.xaxis.set_label_position('bottom')

                if title!='':
                    myax.set_title(title)

                fig.tight_layout()
                fig.canvas.draw()
                fig.canvas.flush_events()

        else:
            logging.warning( "Can't open the clipboard", "Error")

    def get_mpl_plot(self, center = [0., 0.], width = 500., height = 500., title='', toshow=True) -> tuple[Figure, Axes]:
        """
        Récupère un graphique matplotlib sur base de la fenêtre OpenGL et de la palette de la matrice/résultat actif.
        """

        self.zoom_on(center=center, width=width, height= height, canvas_height=self.canvasheight, forceupdate=True)

        fig,axes = plt.subplots(1,2, gridspec_kw={'width_ratios': [20, 1]})
        self.display_canvasogl(fig=fig,ax=axes[0])

        palette = self.palette_for_copy
        palette.export_image(None, h_or_v='v', figax=(fig,axes[1]))

        # if self.active_array is not None:
        #     self.active_array.mypal.export_image(None, h_or_v='v', figax=(fig,axes[1]))
        # elif self.active_res2d is not None:
        #     self.active_res2d.mypal.export_image(None, h_or_v='v', figax=(fig,axes[1]))

        axes[0].xaxis.set_ticks_position('bottom')
        axes[0].xaxis.set_label_position('bottom')

        fig.set_size_inches(12,10)

        fontsize(axes[0], 12)
        fontsize(axes[1], 12)

        if title!='':
            axes[0].set_title(title)

        fig.tight_layout()
        if toshow:
            fig.show()

        return fig, axes

    def create_video(self, fn:str = '', framerate:int = 0, start_step:int = 0, end_step:int = 0, every:int = 0):
        """
        Création d'une vidéo sur base des résultats
        """
        try:
            import cv2
        except:
            logging.error(_('Please install opencv-python'))
            return

        dlg = Sim_VideoCreation(None, title = _('Video creation'), sim= self.active_res2d, mapviewer=self)

        ret = dlg.ShowModal()
        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return

        fn, framerate, start_step, end_step, interval, fontsize, fontcolor, timeposition = dlg.get_values()
        dlg.Destroy()

        times,steps = self.active_res2d.get_times_steps()

        framesize = (int(self.canvaswidth), int(self.canvasheight))
        video = cv2.VideoWriter(fn, cv2.VideoWriter_fourcc(*'XVID'), framerate, framesize)

        el_time = str(timedelta(seconds=int(times[self.active_res2d.current_result])))

        zones_time = Zones(mapviewer=self)
        self.add_object('vector', newobj=zones_time, ToCheck=True, id='__VideoTime__')
        zone_time = zone(name='Time')
        vec_time = vector(name='Time')

        zones_time.add_zone(zone_time, forceparent=True)
        zone_time.add_vector(vec_time, forceparent=True)

        if timeposition == 'top-center':
            x = (self.xmax+self.xmin)/2.
            y = self.ymax - 0.05*(self.ymax-self.ymin)
        elif timeposition == 'bottom-center':
            x = (self.xmax+self.xmin)/2.
            y = self.ymin + 0.05*(self.ymax-self.ymin)
        elif timeposition == 'top-left':
            x = self.xmin + 0.15*(self.xmax-self.xmin)
            y = self.ymax - 0.05*(self.ymax-self.ymin)
        elif timeposition == 'bottom-left':
            x = self.xmin + 0.15*(self.xmax-self.xmin)
            y = self.ymin + 0.05*(self.ymax-self.ymin)
        elif timeposition == 'top-right':
            x = self.xmax - 0.15*(self.xmax-self.xmin)
            y = self.ymax - 0.05*(self.ymax-self.ymin)
        elif timeposition == 'bottom-right':
            x = self.xmax - 0.15*(self.xmax-self.xmin)
            y = self.ymin + 0.05*(self.ymax-self.ymin)
        else:
            x = (self.xmax+self.xmin)/2.
            y = self.ymax - 0.05*(self.ymax-self.ymin)

        vec_time.add_vertex(wolfvertex(x,y))
        vec_time.set_legend_position(x,y)
        vec_time.set_legend_text('Time {:0>8} s'.format(el_time))
        vec_time.set_legend_visible(True)
        vec_time.myprop.legendfontsize = fontsize
        vec_time.myprop.legendcolor = getIfromRGB(fontcolor)

        for curmodel in self.iterator_over_objects(draw_type.RES2D):
            curmodel: Wolfresults_2D
            curmodel.step_interval_results = interval

        all_steps = range(0, int((end_step-start_step) // interval) + 1)

        pgbar = wx.ProgressDialog(_('Video creation'),
                                  _('Creating video...'),
                                    maximum = len(all_steps),
                                    parent = None,
                                    style = wx.PD_APP_MODAL | wx.PD_ELAPSED_TIME | wx.PD_ESTIMATED_TIME | wx.PD_REMAINING_TIME)
        pgbar.Show()

        self.read_one_result(start_step-1)
        for idx in tqdm(all_steps, desc=_('Creating video')):

            image = self.get_canvas_as_image()

            video.write(cv2.cvtColor(np.asarray(image),cv2.COLOR_RGB2BGR))

            self.simul_next_step()

            el_time = str(timedelta(seconds=int(times[self.active_res2d.current_result])))
            vec_time.set_legend_text('Time {:0>8} s'.format(el_time))

            if not pgbar.WasCancelled():
                pgbar.Update(idx)
            else:
                break

        pgbar.Destroy()
        video.release()

        for curmodel in self.iterator_over_objects(draw_type.RES2D):
            curmodel: Wolfresults_2D
            curmodel.step_interval_results = 1

        self.removeobj_from_id('__VideoTime__')

    def get_canvas_as_image(self) -> Image.Image:
        """
        Récupère la fenêtre OpenGL sous forme d'image
        """

        self.Paint()

        if self.SetCurrentContext():
            glPixelStorei(GL_PACK_ALIGNMENT, 1)
            data = glReadPixels(0, 0, self.canvaswidth, self.canvasheight, GL_RGBA, GL_UNSIGNED_BYTE)
            myimage: Image.Image
            myimage = Image.frombuffer("RGBA", (self.canvaswidth, self.canvasheight), data)
            myimage = myimage.transpose(1)

            return myimage

    def copy_canvasogl(self,
                       mpl:bool= True,
                       ds:float= 0.,
                       figsizes= [10.,10.],
                       palette:wolfpalette = None):
        """
        Generate image based on UI context and copy to the Clipboard

        :param mpl: Using Matplolib as renderer. Defaults to True.
        :type mpl: bool, optional
        :parem ds: Ticks size. Defaults to 0..
        :type ds: float, optional
        :parem figsizes: fig size in inches
        :type figsizes: list, optional
        """

        if wx.TheClipboard.Open():
            self.Paint()

            if self.SetCurrentContext():

                myimage = self.get_canvas_as_image()

                metadata = PngInfo()
                metadata.add_text('xmin', str(self.xmin))
                metadata.add_text('ymin', str(self.ymin))
                metadata.add_text('xmax', str(self.xmax))
                metadata.add_text('ymax', str(self.ymax))

                if mpl:
                    if ds == 0.:
                        ds = self.ticks_size # Global parameters

                        if ds == 0.:
                            ds = 100.

                        nb_ticks_x = (self.xmax - self.xmin) // ds
                        nb_ticks_y = (self.ymax - self.ymin) // ds

                        if nb_ticks_x > 10 or nb_ticks_y > 10:
                            logging.error(_('Too many ticks for the image. Please raise the ticks size in the global options.'))
                            dlg = wx.MessageDialog(None, _('Too many ticks for the image. Please raise the ticks size in the global options.'), _('Error'), style=wx.OK)
                            dlg.ShowModal()
                            dlg.Destroy()
                            wx.TheClipboard.Close()
                            return

                    # Création d'un graphique Matplotlib
                    extent = (self.xmin, self.xmax, self.ymin, self.ymax)
                    fig, ax = plt.subplots(1, 1)

                    w, h = [self.width, self.height]

                    neww = figsizes[0]
                    newh = h/w * figsizes[0]

                    fig.set_size_inches(neww, newh)

                    pos = ax.imshow(myimage,
                                    origin='upper',
                                    extent=extent)

                    x1 = np.ceil((self.xmin // ds) * ds)
                    if x1 < self.xmin:
                        x1 += ds
                    x2 = int((self.xmax // ds) * ds)
                    if x2 > self.xmax:
                        x2 -= ds
                    y1 = np.ceil((self.ymin // ds) * ds)
                    if y1 < self.ymin:
                        y1 += ds
                    y2 = int((self.ymax // ds) * ds)
                    if y2 > self.ymax:
                        y2 -= ds

                    x_label_list = np.linspace(x1, x2, int((x2 - x1) / ds) + 1, True)
                    if self.ticks_bounds:
                        x_label_list = np.insert(x_label_list, 0, self.xmin)
                        x_label_list = np.insert(x_label_list, -1, self.xmax)
                        x_label_list = np.unique(x_label_list)

                    y_label_list = np.linspace(y1, y2, int((y2 - y1) / ds) + 1, True)
                    if self.ticks_bounds:
                        y_label_list = np.insert(y_label_list, 0, self.ymin)
                        y_label_list = np.insert(y_label_list, -1, self.ymax)
                        y_label_list = np.unique(y_label_list)

                    ax.set_xticks(x_label_list)
                    ax.set_yticks(y_label_list)

                    ax.set_xticklabels(plt.FormatStrFormatter('%.1f').format_ticks(x_label_list),
                                       fontsize = self.ticks_fontsize, rotation = self.ticks_xrotation)
                    ax.set_yticklabels(plt.FormatStrFormatter('%.1f').format_ticks(y_label_list),
                                       fontsize = self.ticks_fontsize)

                    ax.set_xlabel('X ($m$)')
                    ax.set_ylabel('Y ($m$)')

                    fig.tight_layout()

                    #création d'un'buffers
                    buf = io.BytesIO()
                    #sauvegarde de la figure au format png
                    fig.savefig(buf, format='png')

                    #déplacement au début du buffer
                    buf.seek(0)
                    #lecture du buffer et conversion en image avec PIL
                    im = Image.open(buf)

                    if palette is None:
                        palette = self.palette_for_copy
                        # if self.active_array is not None:
                        #     palette = self.active_array.mypal
                        # elif self.active_res2d is not None:
                        #     palette = self.active_res2d.mypal

                    if palette is not None:
                        if palette.values is not None:
                            bufpal = io.BytesIO()
                            palette.export_image(bufpal,'v')

                            try:
                                bufpal.seek(0)

                                #lecture du buffer et conversion en image avec PIL
                                impal = Image.open(bufpal)

                            except Exception as e:

                                text = _('Error while creating the colormap/palette image !')
                                text += '\n'
                                text += _('Please check if an array or a 2D result is active !')

                                logging.error(text)
                                dlg = wx.MessageDialog(None, text, _('Error'), style=wx.OK)
                                dlg.ShowModal()
                                dlg.Destroy()
                                return

                            impal = impal.resize((int(impal.size[0]*im.size[1]*.8/impal.size[1]),int(im.size[1]*.8)))

                            imnew = Image.new('RGB',(im.size[0]+impal.size[0], im.size[1]), (255,255,255))

                            # On colle l'image du buffer et la palette pour ne former qu'une seul image à copier dans le clipboard
                            imnew.paste(im.convert('RGB'),(0,0))
                            imnew.paste(impal.convert('RGB'),(im.size[0]-10, int((im.size[1]-impal.size[1])/3)))
                            im=imnew
                            bufpal.close()
                        else:
                            imnew = Image.new('RGB', (im.size[0], im.size[1]), (255,255,255))

                            # On colle l'image du buffer et la palette pour ne former qu'une seul image à copier dans le clipboard
                            imnew.paste(im.convert('RGB'),(0,0))
                            im=imnew
                    else:
                        imnew = Image.new('RGB', (im.size[0], im.size[1]), (255,255,255))

                        # On colle l'image du buffer et la palette pour ne former qu'une seul image à copier dans le clipboard
                        imnew.paste(im.convert('RGB'),(0,0))
                        im=imnew

                    #création d'un objet bitmap wx
                    wxbitmap = wx.Bitmap().FromBuffer(im.width,im.height,im.tobytes())

                    # objet wx exportable via le clipboard
                    dataobj = wx.BitmapDataObject()
                    dataobj.SetBitmap(wxbitmap)

                    wx.TheClipboard.SetData(dataobj)
                    wx.TheClipboard.Close()

                    fig.set_visible(False)

                    buf.close()

                    return fig, ax, im

                else:
                    """ Création d'un objet bitmap wx sur base du canvas
                    et copie dans le clipboard
                    """
                    # wxbitmap = wx.Bitmap().FromBuffer(myimage.width,myimage.height,myimage.tobytes())
                    wxbitmap = wx.Bitmap().FromBufferRGBA(myimage.width,myimage.height,myimage.tobytes())

                    # objet wx exportable via le clipboard
                    dataobj = wx.BitmapDataObject()
                    dataobj.SetBitmap(wxbitmap)

                    wx.TheClipboard.SetData(dataobj)
                    wx.TheClipboard.Close()

                    return myimage

        else:
            wx.MessageBox("Can't open the clipboard", "Error")

    def save_canvasogl(self,
                       fn:str='',
                       mpl:bool=True,
                       ds:float=0.,
                       dpi:int= 300,
                       add_title:bool = False,
                       figsizes= [10.,10.],
                       arrayid_as_title:bool = False,
                       resid_as_title:bool = False):
        """
        Sauvegarde de la fenêtre d'affichage dans un fichier

        :param fn: File name (.png or .jpg file)
        :param mpl: Using Matplotlib as renderer
        :param ds: Ticks interval
        """

        # FIXME : SHOULD BE MERGEd WITH copy_canvasogl
        fn = str(fn)

        if fn == '':
            dlg = wx.FileDialog(None, _('Choose file name'), wildcard='PNG (*.png)|*.png|JPG (*.jpg)|*.jpg',
                                style=wx.FD_SAVE)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return
            fn = dlg.GetPath()
            dlg.Destroy()
        elif not fn.endswith('.png'):
            fn += '.png'

        if self.SetCurrentContext():
            self.Paint()

            if mpl:
                if ds == 0.:
                    dlg = wx.NumberEntryDialog(self,
                                               _("xmin : {:.3f} \nxmax : {:.3f} \nymin : {:.3f} \nymax : {:.3f} \n\n  dx : {:.3f}\n  dy : {:.3f}").format(
                                                   self.xmin, self.xmax, self.ymin, self.ymax, self.xmax - self.xmin,
                                                   self.ymax - self.ymin),
                                               _("Interval [m]"), _("Ticks interval ?"), 500, 1, 10000)
                    ret = dlg.ShowModal()

                    if ret == wx.ID_CANCEL:
                        dlg.Destroy()
                        return

                    ds = float(dlg.GetValue())
                    dlg.Destroy()

                # Création d'un graphique Matplotlib
                extent = (self.xmin, self.xmax, self.ymin, self.ymax)
                fig, ax = plt.subplots(1, 1)
                w, h = [self.width, self.height]

                neww = figsizes[0]
                newh = h/w * figsizes[0]

                fig.set_size_inches(neww, newh)

                pot_title = self.viewer_name
                if arrayid_as_title:
                    pot_title = self.active_array.idx
                if resid_as_title:
                    pot_title = self.active_res2d.idx

                self.display_canvasogl(fig=fig,
                                       ax=ax,
                                       title=pot_title if add_title else '',
                                       ds = ds)

                #création d'un'buffers
                buf = io.BytesIO()
                #sauvegarde de la figure au format png
                fig.savefig(buf, format='png')

                #déplacement au début du buffer
                buf.seek(0)
                #lecture du buffer et conversion en image avec PIL
                im = Image.open(buf)

                if self.palette_for_copy.values is not None:
                    bufpal = io.BytesIO()
                    self.palette_for_copy.export_image(bufpal,'v')
                    bufpal.seek(0)

                    #lecture du buffer et conversion en image avec PIL
                    impal = Image.open(bufpal)
                    impal = impal.resize((int(impal.size[0]*im.size[1]*.8/impal.size[1]),int(im.size[1]*.8)))

                    imnew = Image.new('RGB',(im.size[0]+impal.size[0], im.size[1]), (255,255,255))

                    # On colle l'image du buffer et la palette pour ne former qu'une seul image à copier dans le clipboard
                    imnew.paste(im.convert('RGB'),(0,0))
                    imnew.paste(impal.convert('RGB'),(im.size[0]-10, int((im.size[1]-impal.size[1])/3)))
                    im=imnew
                    bufpal.close()
                else:
                    imnew = Image.new('RGB', (im.size[0], im.size[1]), (255,255,255))

                    # On colle l'image du buffer et la palette pour ne former qu'une seul image à copier dans le clipboard
                    imnew.paste(im.convert('RGB'),(0,0))
                    im=imnew

                im.save(fn, dpi=(dpi, dpi))
                fig.set_visible(False)
                buf.close()

            else:
                metadata = PngInfo()
                metadata.add_text('xmin', str(self.xmin))
                metadata.add_text('ymin', str(self.ymin))
                metadata.add_text('xmax', str(self.xmax))
                metadata.add_text('ymax', str(self.ymax))

                myimage = self.get_canvas_as_image()
                myimage.save(fn, pnginfo=metadata)

            return fn, ds
        else:
            raise NameError(
                'Opengl setcurrent -- maybe a conflict with an existing opengl32.dll file - please rename the opengl32.dll in the libs directory and retry')

    def reporting(self, dir=''):
        """ First attempt to create a reporting.
        !! Must be improved !!
        """
        if dir == '':
            dlg = wx.DirDialog(None, "Choose directory to store reporting", style=wx.FD_SAVE)
            ret = dlg.ShowModal()

            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            dir = dlg.GetPath()
            dlg.Destroy()

        myppt = Presentation(__file__)
        slide = myppt.slides.add_slide(0)

        for curzone in self.myzones:
            for curvec in curzone.myvectors:
                curvec: vector
                if curvec.nbvertices > 1:
                    oldwidth = curvec.myprop.width
                    curvec.myprop.width = 4
                    myname = curvec.myname

                    self.Activate_vector(curvec)

                    if self.linked:
                        for curview in self.linkedList:
                            title = curview.GetTitle()
                            curview.zoomon_activevector()
                            fn = path.join(dir, title + '_' + myname + '.png')
                            curview.save_canvasogl(fn)
                    else:
                        self.zoomon_activevector()
                        fn = path.join(dir, myname + '.png')
                        self.save_canvasogl(fn)

                        fn = path.join(dir, 'palette_v_' + myname + '.png')
                        self.active_array.mypal.export_image(fn, 'v')
                        fn = path.join(dir, 'palette_h_' + myname + '.png')
                        self.active_array.mypal.export_image(fn, 'h')

                    curvec.myprop.width = oldwidth

    def InitUI(self):
        """ Initialisation de l'interface utilisateur """

        self.Bind(wx.EVT_SIZE, self.OnSize)
        self.Bind(wx.EVT_CLOSE, self.OnClose)

        # self.canvas.Bind(wx.EVT_CONTEXT_MENU, self.OnShowPopup)
        self.canvas.Bind(wx.EVT_PAINT, self.OnPaint)

        self.canvas.Bind(wx.EVT_CHAR_HOOK, self.OnHotKey)

        self.canvas.Bind(wx.EVT_BUTTON, self.On_Mouse_Button)
        self.canvas.Bind(wx.EVT_RIGHT_DCLICK, self.On_Right_Double_Clicks)
        self.canvas.Bind(wx.EVT_LEFT_DCLICK, self.On_Left_Double_Clicks)
        self.canvas.Bind(wx.EVT_LEFT_DOWN, self.On_Mouse_Left_Down)
        self.canvas.Bind(wx.EVT_LEFT_UP, self.On_Mouse_Left_Up)
        self.canvas.Bind(wx.EVT_MIDDLE_DOWN, self.On_Mouse_Left_Down)
        self.canvas.Bind(wx.EVT_RIGHT_DOWN, self.On_Mouse_Right_Down)
        self.canvas.Bind(wx.EVT_RIGHT_UP, self.On_Mouse_Right_Up)
        self.canvas.Bind(wx.EVT_MOTION, self.On_Mouse_Motion)
        self.canvas.Bind(wx.EVT_LEAVE_WINDOW, self.OnLeave)
        self.canvas.Bind(wx.EVT_MOUSEWHEEL, self.On_Mouse_Button)

        self.treelist.Bind(dataview.EVT_TREELIST_ITEM_CHECKED, self.OnCheckItem)
        self.treelist.Bind(dataview.EVT_TREELIST_ITEM_ACTIVATED, self.OnActivateTreeElem)
        self.treelist.Bind(dataview.EVT_TREELIST_ITEM_CONTEXT_MENU, self.OntreeRight)
        self.treelist.Bind(wx.EVT_CHAR_HOOK, self.OnHotKey)
        self.treelist.Bind(dataview.EVT_TREELIST_SELECTION_CHANGED,self.OnSelectItem)
        # dispo dans wxpython 4.1 self.Bind(wx.EVT_GESTURE_ZOOM,self.OnZoomGesture)

        self.Centre()

        self.mybc = []
        self.myarrays = []
        self.mypartsystems = []
        self.myvectors = []
        self.mytiles = []
        self.myimagestiles = []
        self.myclouds = []
        self.mytri = []
        self.myothers = []
        self.myviews = []
        self.mywmsback = []
        self.mywmsfore = []
        self.myres2D = []
        self.myviewers3d = []
        self.myviewerslaz = []
        self.mylazdata = []
        self.mydrownings = []
        self.mypicturecollections = []

        self.mydikes = []
        self.myinjectors = []


        self.mymplfigs = []

        self.sim_explorers = {}

        # liste des éléments modifiable dans l'arbre
        self.all_lists = [self.myarrays, self.myvectors, self.myclouds, self.mytri, self.myothers, self.myviews, self.myres2D, self.mytiles, self.myimagestiles, self.mypartsystems, self.myviewers3d, self.myviewerslaz, self.mydikes, self.mydrownings, self.myinjectors]

        self.menu_options = wx.Menu()
        self._change_title = self.menu_options.Append(wx.ID_ANY, _('Change title'), _('Change title of the window'))
        self.Bind(wx.EVT_MENU, self.OnChangeTitle, self._change_title)

        if self.get_configuration() is not None:
            # see PyGui.py if necessary

            self.menubar.Append(self.menu_options, _('Options'))
            self.option_global = self.menu_options.Append(wx.ID_ANY,_("Global"),_("Modify global options"))
            self.Bind(wx.EVT_MENU, self.GlobalOptionsDialog, self.option_global)

        self.menu_1to9 =self.menu_options.Append(wx.ID_ANY, _('Colors for selections 1->9'), _('Selections'))
        self.Bind(wx.EVT_MENU, self.colors1to9.change_colors, self.menu_1to9)

        self.menu_qdfidf()

        self.Show(True)

    def OnChangeTitle(self, e):
        """
        Change the title of the window
        """

        dlg = wx.TextEntryDialog(None, _('Enter the new title'), _('Change title'), self.GetTitle())
        if dlg.ShowModal() == wx.ID_OK:
            self.SetTitle(dlg.GetValue())
        dlg.Destroy()

    def OnSize(self, e):
        """
        Redimensionnement de la fenêtre
        """
        if self.regular:
            # retrouve la taille de la fenêtre
            width, height = self.GetClientSize()
            # enlève la barre d'arbre
            width -= self.treewidth
            # définit la taille de la fenêtre graphique OpenGL et sa position (à droite de l'arbre)
            self.canvas.SetSize(width, height)
            self.canvas.SetPosition((self.treewidth, 0))
            # calcule les limites visibles sur base de la taille de la fenêtre et des coefficients sx sy
            self.setbounds()
            # fixe la taille de l'arbre (notamment la hauteur)
            # self.treelist.SetSize(self.treewidth,height)
            e.Skip()

    def ManageActions(self, id):
        """
        Gestion des actions via les menus

        TODO : A généraliser?
        """
        curmenu = self.tools[id]['menu']

        if curmenu.IsCheckable():
            if not curmenu.IsChecked():
                curmenu.Check(False)
                self.action = None

                if id == ID_LOCMINMAX:
                    self.update_absolute_minmax = True
            else:
                curmenu.Check()
                if not self.tools[id]['name'] is None:
                    self.action = self.tools[id]['name']

        else:
            if id == ID_SORTALONG:
                # Tri le long d'un vecteur
                if not self.active_cs is None and not self.active_vector is None:
                    self.active_cs: crosssections
                    self.active_vector: vector
                    self.active_cs.sort_along(self.active_vector.asshapely_ls(), self.active_vector.myname, False)
                else:
                    msg = ''
                    if self.active_cs is None:
                        msg += _('Please select the active cross sections \n')
                    if self.active_vector is None:
                        msg += _('Please select the active supprt vector')
                    mydiag = wx.MessageDialog(self, msg, _('Sort along'))
                    mydiag.ShowModal()

    def center_view_on(self, cx, cy):
        """
        Center the view on the point of (map) coordinates (x,y)
        """

        self.mousex, self.mousey = cx, cy

        # retrouve la taille de la fenêtre OpenGL
        width, height = self.canvas.GetSize()

        # calcule la taille selon X et Y en coordonnées réelles
        width = width / self.sx
        height = height / self.sy

        # retrouve les bornes min et max sur base de la valeur centrale qui est censée ne pas bouger
        self.xmin = self.mousex - width / 2.
        self.xmax = self.xmin + width
        self.ymin = self.mousey - height / 2.
        self.ymax = self.ymin + height

    def setbounds(self, updatescale=True):
        """
        Calcule les limites visibles de la fenêtre graphique sur base des
        facteurs d'échelle courants
        """

        if updatescale:
            self.updatescalefactors()

            # retrouve la taille de la fenêtre OpenGL
            width, height = self.canvas.GetSize()
            self.canvaswidth = width
            self.canvasheight = height

            # calcule la taille selon X et Y en coordonnées réelles
            width = width / self.sx
            height = height / self.sy

            # retrouve les bornes min et max sur base de la valeur centrale qui est censée ne pas bouger
            self.xmin = self.mousex - width / 2.
            self.xmax = self.xmin + width
            self.ymin = self.mousey - height / 2.
            self.ymax = self.ymin + height

            self.width = width
            self.height = height

            self.mousex = self.xmin + width / 2.
            self.mousey = self.ymin + height / 2.

            self.updatescalefactors()

        else:
            # retrouve les bornes min et max sur base de la valeur centrale qui est censée ne pas bouger
            self.xmin = self.mousex - self.width / 2.
            self.xmax = self.xmin + self.width
            self.ymin = self.mousey - self.height / 2.
            self.ymax = self.ymin + self.height


        self.mybackisloaded = False
        self.myfrontisloaded = False

        self.Refresh()
        self.mimicme()


    def setsizecanvas(self,width,height):
        """ Redimensionne la fenêtre graphique """
        self.canvas.SetClientSize(width, height)

    def updatescalefactors(self):
        """ Mise à jour des facteurs d'échelle
            This one updates the scale factors based on the relative sizes
            of the GLCanvas and the footprint that should fit in it.
        """

        width, height = self.canvas.GetSize()

        self.sx = 1
        self.sy = 1
        if self.width > 0 and width >0 :
            self.sx = float(width) / self.width
        if self.height > 0 and height > 0 :
            self.sy = float(height) / self.height

        self.sx = min(self.sx, self.sy)
        self.sy = self.sx

    def add_viewer_and_link(self):
        """ Ajout d'une nouvelle fenêtre de visualisation et liaison avec la fenêtre courante """

        dlg = wx.TextEntryDialog(self, _('Enter a caption for the new window'))

        ret = dlg.ShowModal()

        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return

        newcap = dlg.GetValue()
        dlg.Destroy()
        newview = WolfMapViewer(None, newcap, w=600, h=600, wxlogging=self.wxlogging, wolfparent=self.wolfparent)
        newview.add_grid()
        newview.add_WMS()

        if self.linkedList is None:
            self.linkedList = [self]

        self.linkedList.append(newview)

        for curview in self.linkedList:
            curview.linked = True
            curview.linkedList = self.linkedList
            curview.link_shareopsvect = False

        logging.info(_('New viewer added and linked'))

    def add_grid(self):
        """ Ajout d'une grille """

        mygrid = Grid(1000.)
        self.add_object('vector', newobj=mygrid, ToCheck=False, id='Grid')

    def add_WMS(self):
        """ Ajout de couches WMS """
        xmin = 0
        xmax = 0
        ymin = 0
        ymax = 0
        orthos = {'IMAGERIE': {'1971': 'ORTHO_1971', '1994-2000': 'ORTHO_1994_2000',
                               '2006-2007': 'ORTHO_2006_2007',
                               '2009-2010': 'ORTHO_2009_2010',
                               '2012-2013': 'ORTHO_2012_2013',
                               '2015': 'ORTHO_2015',
                               '2016': 'ORTHO_2016',
                               '2017': 'ORTHO_2017',
                               '2018': 'ORTHO_2018',
                               '2019': 'ORTHO_2019',
                               '2020': 'ORTHO_2020',
                               '2021': 'ORTHO_2021',
                               '2022 printemps': 'ORTHO_2022_PRINTEMPS',
                               '2022 été': 'ORTHO_2022_ETE',
                               '2023 été': 'ORTHO_2023_ETE',
                               'Last one': 'ORTHO_LAST',
                               }}
        data_2021 = {'EAU': {'IDW': 'ZONES_INONDEES_IDW',
                             'Emprise': 'ZONES_INONDEES',
                             'Emprise wo Alea': 'ZONES_INONDEES_wo_alea'}}

        lifewatch = {'LW_ecotopes_lc_hr_raster': {'2006': '2006',
                                     '2010': '2010',
                                     '2015': '2015',
                                     '2018': '2018',
                                     '2019': '2019',
                                     '2020': '2020',
                                     '2021': '2021',
                                     '2022': '2022',
                                     }}

        """cat:Literal['orthoimage_coverage',
                       'orthoimage_coverage_2016',
                       'orthoimage_coverage_2017',
                       'orthoimage_coverage_2018',
                       'orthoimage_coverage_2019',
                       'orthoimage_coverage_2020',
                       'orthoimage_coverage_2021',
                       'orthoimage_coverage_2022']
        """
        ign_belgique = {'Orthophotos': {'Last': 'orthoimage_coverage',
                                        '2016': 'orthoimage_coverage_2016',
                                        '2017': 'orthoimage_coverage_2017',
                                        '2018': 'orthoimage_coverage_2018',
                                        '2019': 'orthoimage_coverage_2019',
                                        '2020': 'orthoimage_coverage_2020',
                                        '2021': 'orthoimage_coverage_2021',
                                        '2022': 'orthoimage_coverage_2022',}}

        """ ['crossborder',
                            'crossborder_grey',
                            'overlay',
                            'topo',
                            'topo_grey']"""
        ign_cartoweb = {'CartoWeb': {'Crossborder': 'crossborder',
                                    'Crossborder Grey': 'crossborder_grey',
                                    'Overlay': 'overlay',
                                    'Topographic': 'topo',
                                    'Topographic Grey': 'topo_grey',}}

        ign_postflood2021 = {'PostFlood2021': {'Flood 2021': 'orthoimage_flood'}}

        """
        ['10_m_u__wind_component',
        '10_m_v__wind_component',
        '2_m_Max_temp_since_ppp',
        '2_m_Min_temp_since_ppp',
        '2_m_dewpoint_temperature',
        '2_m_temperature',
        '2m_Relative_humidity',
        'Convective_rain',
        'Convective_snow',
        'Geopotential',
        'Inst_flx_Conv_Cld_Cover',
        'Inst_flx_High_Cld_Cover',
        'Inst_flx_Low_Cld_Cover',
        'Inst_flx_Medium_Cld_Cover',
        'Inst_flx_Tot_Cld_cover',
        'Large_scale_rain',
        'Large_scale_snow',
        'Mean_sea_level_pressure',
        'Relative_humidity',
        'Relative_humidity_isobaric',
        'SBL_Meridian_gust',
        'SBL_Zonal_gust',
        'Specific_humidity',
        'Surf_Solar_radiation',
        'Surf_Thermal_radiation',
        'Surface_CAPE',
        'Surface_Temperature',
        'Surface_orography',
        'Temperature',
        'Total_precipitation',
        'U-velocity',
        'V-velocity',
        'Vertical_velocity',
        'Wet_Bulb_Poten_Temper',
        'freezing_level_zeroDegC_isotherm'],
        """
        # alaro = {'ALARO': {'10m_u_wind_component': '10_m_u__wind_component',
        #                   '10m_v_wind_component': '10_m_v__wind_component',
        #                   '2m_Max_temp_since_ppp': '2_m_Max_temp_since_ppp',
        #                   '2m_Min_temp_since_ppp': '2_m_Min_temp_since_ppp',
        #                   '2m_dewpoint_temperature': '2_m_dewpoint_temperature',
        #                   '2m_temperature': '2_m_temperature',
        #                   '2m_Relative_humidity': '2m_Relative_humidity',
        #                   'Convective_rain': 'Convective_rain',
        #                   'Convective_snow': 'Convective_snow',
        #                   'Geopotential': 'Geopotential',
        #                   'Inst_flx_Conv_Cld_Cover': 'Inst_flx_Conv_Cld_Cover',
        #                   'Inst_flx_High_Cld_Cover': 'Inst_flx_High_Cld_Cover',
        #                   'Inst_flx_Low_Cld_Cover': 'Inst_flx_Low_Cld_Cover',
        #                   'Inst_flx_Medium_Cld_Cover': 'Inst_flx_Medium_Cld_Cover',
        #                   'Inst_flx_Tot_Cld_cover': 'Inst_flx_Tot_Cld_cover',
        #                   'Large_scale_rain': 'Large_scale_rain',
        #                   'Large_scale_snow': 'Large_scale_snow',
        #                   'Mean_sea_level_pressure': 'Mean_sea_level_pressure',
        #                   'Relative_humidity': 'Relative_humidity',
        #                   'Relative_humidity_isobaric': 'Relative_humidity_isobaric',
        #                   'SBL_Meridian_gust': 'SBL_Meridian_gust',
        #                   'SBL_Zonal_gust': 'SBL_Zonal_gust',
        #                   'Specific_humidity': 'Specific_humidity',
        #                   'Surf_Solar_radiation': 'Surf_Solar_radiation',
        #                   'Surf_Thermal_radiation': 'Surf_Thermal_radiation',
        #                   'Surface_CAPE': 'Surface_CAPE',
        #                   'Surface_Temperature': 'Surface_Temperature',
        #                   'Surface_orography': 'Surface_orography',
        #                   'Temperature': 'Temperature',
        #                   'Total_precipitation': 'Total_precipitation',
        #                   'U-velocity': 'U-velocity',
        #                   'V-velocity': 'V-velocity',
        #                   'Vertical_velocity': 'Vertical_velocity',
        #                   'Wet_Bulb_Poten_Temper': 'Wet_Bulb_Poten_Temper',
        #                   'freezing_level_zeroDegC_isotherm': 'freezing_level_zeroDegC_isotherm',}}
        alaro = {'ALARO': {'2m_temperature': '2_m_temperature',
                          'Convective_rain': 'Convective_rain',
                          'Convective_snow': 'Convective_snow',
                          'Large_scale_rain': 'Large_scale_rain',
                          'Large_scale_snow': 'Large_scale_snow',
                          'Surface_Temperature': 'Surface_Temperature',
                          'Total_precipitation': 'Total_precipitation',}}


        for idx, (k, item) in enumerate(orthos.items()):
            for kdx, (m, subitem) in enumerate(item.items()):
                self.add_object(which='wmsback',
                                newobj=imagetexture('PPNC', m, k, subitem,
                                                    self, xmin, xmax, ymin, ymax, -99999, 1024),
                                ToCheck=False, id='PPNC ' + m)

        for idx, (k, item) in enumerate(data_2021.items()):
            for kdx, (m, subitem) in enumerate(item.items()):
                self.add_object(which='wmsback',
                                newobj=imagetexture('PPNC', m, k, subitem,
                                                    self, xmin, xmax, ymin, ymax, -99999, 1024),
                                ToCheck=False, id='Data 2021 ' + m)

        for idx, (k, item) in enumerate(lifewatch.items()):
            for kdx, (m, subitem) in enumerate(item.items()):
                self.add_object(which='wmsback',
                                newobj=imagetexture('LanCover', m, k, subitem,
                                                    self, xmin, xmax, ymin, ymax, -99999, 1024, LifeWatch=True),
                                ToCheck=False, id='LifeWatch LC' + m)

        for idx, (k, item) in enumerate(ign_belgique.items()):
            for kdx, (m, subitem) in enumerate(item.items()):
                self.add_object(which='wmsback',
                                newobj=imagetexture('Orthos IGN', m, k, subitem,
                                                    self, xmin, xmax, ymin, ymax, -99999, 1024, IGN_Belgium=True),
                                ToCheck=False, id='IGN ' + m)

        for idx, (k, item) in enumerate(ign_cartoweb.items()):
            for kdx, (m, subitem) in enumerate(item.items()):
                self.add_object(which='wmsback',
                                newobj=imagetexture('Cartoweb IGN', m, k, subitem,
                                                    self, xmin, xmax, ymin, ymax, -99999, 1024, IGN_Cartoweb=True),
                                ToCheck=False, id='IGN ' + m)

        for idx, (k, item) in enumerate(ign_postflood2021.items()):
            for kdx, (m, subitem) in enumerate(item.items()):
                self.add_object(which='wmsback',
                                newobj=imagetexture('IGN 2021', m, k, subitem,
                                                    self, xmin, xmax, ymin, ymax, -99999, 1024, postFlood2021=True),
                                ToCheck=False, id='orthos post2021')

        self.add_object(which='wmsback',
                        newobj=imagetexture('PPNC', 'Orthos France', 'OI.OrthoimageCoverage.HR', '',
                                            self, xmin, xmax, ymin, ymax, -99999, 1024, France=True, epsg='EPSG:2154'),
                        ToCheck=False, id='Orthos France')


        forelist = {'EAU': {'Aqualim': 'RES_LIMNI_DGARNE', 'Alea': 'ALEA_INOND', 'Lidaxes': 'LIDAXES'},
                    'LIMITES': {'Secteurs Statistiques': 'LIMITES_QS_STATBEL',
                                'Limites administratives': 'LIMITES_ADMINISTRATIVES'},
                    'R3C': {'Limites Communes': 'Municipalities'},
                    # 'INSPIRE': {'Limites administratives': 'AU_wms'},
                    'PLAN_REGLEMENT': {'Plan Parcellaire 2021': 'CADMAP_2021_PARCELLES',
                                       'Plan Parcellaire 2022': 'CADMAP_2022_PARCELLES',
                                       'Plan Parcellaire 2023': 'CADMAP_2023_PARCELLES',
                                       'Plan Parcellaire 2024': 'CADMAP_2024_PARCELLES'}}

        for idx, (k, item) in enumerate(forelist.items()):
            for kdx, (m, subitem) in enumerate(item.items()):
                self.add_object(which='wmsfore',
                                newobj=imagetexture('PPNC', m, k, subitem,
                                                    self, xmin, xmax, ymin, ymax, -99999, 1024),
                                ToCheck=False, id=m)

        for idx, (k, item) in enumerate(alaro.items()):
            for kdx, (m, subitem) in enumerate(item.items()):
                self.add_object(which='wmsfore',
                                newobj=imagetexture('ALARO', m, k, subitem,
                                                    self, xmin, xmax, ymin, ymax, -99999, 1024, Alaro=True),
                                ToCheck=False, id='ALARO ' + m)

        for idx, (k, item) in enumerate(ign_cartoweb.items()):
            for kdx, (m, subitem) in enumerate(item.items()):
                self.add_object(which='wmsfore',
                                newobj=imagetexture('Cartoweb', m, k, subitem,
                                                    self, xmin, xmax, ymin, ymax, -99999, 1024, IGN_Cartoweb=True),
                                ToCheck=False, id='IGN_f ' + m)

        # self.add_object(which='wmsfore',
        #                 newobj=imagetexture('Cadastre Flandres', 'Plan Parcellaire 2024 (Flandres)', 'Adpf', '',
        #                                     self, xmin, xmax, ymin, ymax, -99999, 1024, Vlaanderen=True),
        #                 ToCheck=False, id='Plan Parcellaire 2024 (Flandres)')


    def set_compare(self, ListArrays:list[WolfArray]=None, share_colormap:bool=True):
        """
        Comparison of 2 arrays

        :param ListArrays: List of 2 arrays to compare
        :param share_colormap: Share the colormap between the 2 arrays
        """

        # assert len(ListArrays) == 2, _('List of arrays must contain 2 and only 2 arrays - Here, you have provided {} arrays'.format(len(ListArrays)))

        # Création de 3 fenêtres de visualisation basées sur la classe "WolfMapViewer"
        first = self
        second = WolfMapViewer(None, 'Comparison', w=600, h=600, wxlogging=self.wxlogging, wolfparent = self.wolfparent)
        third  = WolfMapViewer(None, 'Difference', w=600, h=600, wxlogging=self.wxlogging, wolfparent = self.wolfparent)

        second.add_grid()
        third.add_grid()
        second.add_WMS()
        third.add_WMS()

        # Création d'une liste contenant les 3 instances d'objet "WolfMapViewer"
        mylist:list[WolfMapViewer] = []
        mylist.append(first)
        mylist.append(second)
        mylist.append(third)

        # On indique que les objets sont liés en activant le Booléen et en pointant la liste précédente
        for curlist in mylist:
            curlist.linked = True
            curlist.linkedList = mylist

        if ListArrays is not None:
            if len(ListArrays) == 2:
                mnt = ListArrays[0]
                mns = ListArrays[1]

                if not mnt.is_like(mns):
                    logging.warning(_('The 2 arrays must have the same shape - Here, the 2 arrays have different shapes'))
                    return

            else:
                logging.warning(_('List of arrays must contain 2 and only 2 arrays - Here, you have provided {} arrays'.format(len(ListArrays))))
                return
        else:
            logging.warning(_('You must fill the List of arrays with 2 and only 2 arrays - Here, the list is void'))
            return

        mns: WolfArray
        mnt: WolfArray
        diff: WolfArray

        # Recherche d'un masque union des masques partiels
        mns.mask_union(mnt)

        # Création du différentiel -- Les opérateurs mathématiques sont surchargés
        diff = mns - mnt

        # on attribue une matrice par interface graphique
        mnt.change_gui(first)
        mns.change_gui(second)
        diff.change_gui(third)

        path = os.path.dirname(__file__)
        fn = join(path, 'models\\diff16.pal')

        # on partage la palette de couleurs
        if share_colormap:
            mns.add_crosslinked_array(mnt)
            mns.share_palette()

        # on dissocie la palette de la différence
        diff.mypal = wolfpalette()
        if isinstance(diff, WolfArrayMB):
            diff.link_palette()

        diff.mypal.readfile(fn)
        diff.mypal.automatic = False
        diff.myops.palauto.SetValue(0)

        mnt.mypal.automatic = False
        mnt.myops.palauto.SetValue(0)
        if not share_colormap:
            mns.mypal.automatic = False
            mns.myops.palauto.SetValue(0)
            mns.mypal.updatefrompalette(mnt.mypal)

        # Ajout des matrices dans les fenêtres de visualisation
        first.add_object('array', newobj=mnt, ToCheck=True, id='source')
        second.add_object('array', newobj=mns, ToCheck=True, id='comp')
        third.add_object('array', newobj=diff, ToCheck=True, id='diff=comp-source')

        # Partage des vecteurs de la classe d'opérations
        mnt.myops.myzones = mns.myops.myzones
        diff.myops.myzones = mns.myops.myzones

        first.active_array = mnt
        second.active_array = mns
        third.active_array = diff

        mnt.reset_plot()
        mns.reset_plot()
        diff.reset_plot()

    def set_compare_all(self, ListArrays=None, names:list[str] = None):
        """ Comparison of 2 or 3 arrays """

        assert len(ListArrays) == 2 or len(ListArrays) == 3, _('List of arrays must contain 2 or 3 arrays - Here, you have provided {} arrays'.format(len(ListArrays)))
        if names is not None:
            assert len(names) == len(ListArrays)-1, _('List of names must contain the number of names as arrays minus one - Here, you have provided {} names for {} arrays'.format(len(names), len(ListArrays)))
        else:
            names = ['comp1', 'comp2']

        # Création de 3 fenêtres de visualisation basées sur la classe "WolfMapViewer"
        first = self
        second = WolfMapViewer(None, 'Comparison {}'.format(names[0]), w=600, h=600, wxlogging=self.wxlogging, wolfparent=self.wolfparent)
        third  = WolfMapViewer(None, 'Difference {}'.format(names[0]), w=600, h=600, wxlogging=self.wxlogging, wolfparent=self.wolfparent)
        if len(ListArrays) == 3:
            fourth = WolfMapViewer(None, 'Comparison {}'.format(names[1]), w=600, h=600, wxlogging=self.wxlogging, wolfparent=self.wolfparent)
            fifth  = WolfMapViewer(None, 'Difference {}'.format(names[1]), w=600, h=600, wxlogging=self.wxlogging, wolfparent=self.wolfparent)

        # Création d'une liste contenant les multiples instances d'objet "WolfMapViewer"
        list = []
        list.append(first)
        list.append(second)
        list.append(third)
        if len(ListArrays) == 3:
            list.append(fourth)
            list.append(fifth)

        for curview in list:
            if curview is not self:
                curview.add_grid()
                curview.add_grid()

        # On indique que les objets sont liés en actiavt le Booléen et en pointant la liste précédente
        for curview in list:
            curview.linked = True
            curview.linkedList = list

        comp2 = None
        if ListArrays is not None:
            if len(ListArrays) == 2:
                src = ListArrays[0]
                comp1 = ListArrays[1]
            elif len(ListArrays) == 3:
                src = ListArrays[0]
                comp1 = ListArrays[1]
                comp2 = ListArrays[2]
            else:
                return
        else:
            return

        src: WolfArray
        comp1: WolfArray
        diff1: WolfArray
        comp2: WolfArray
        diff2: WolfArray

        # Création du différentiel -- Les opérateurs mathématiques sont surchargés
        diff1 = comp1 - src

        comp1.copy_mask(src, True)
        diff1.copy_mask(src, True)

        src.change_gui(first)
        comp1.change_gui(second)
        diff1.change_gui(third)

        src.mypal.automatic = False
        comp1.mypal.automatic = False
        src.myops.palauto.SetValue(0)
        comp1.myops.palauto.SetValue(0)

        src.mypal.isopop(src.array, src.nbnotnull)
        comp1.mypal.updatefrompalette(src.mypal)

        # Ajout des matrices dans les fenêtres de visualisation
        first.add_object('array', newobj=src, ToCheck=True, id='source')
        second.add_object('array', newobj=comp1, ToCheck=True, id='comp')
        third.add_object('array', newobj=diff1, ToCheck=True, id='diff=comp-source')

        comp1.myops.myzones = src.myops.myzones
        diff1.myops.myzones = src.myops.myzones

        first.active_array = src
        second.active_array = comp1
        third.active_array = diff1

        if comp2 is not None:
            diff2 = comp2 - src
            comp2.copy_mask(src, True)
            diff2.copy_mask(src, True)

            comp2.change_gui(fourth)
            diff2.change_gui(fifth)

            comp2.mypal.automatic = False
            comp2.myops.palauto.SetValue(0)

            comp2.mypal.updatefrompalette(src.mypal)

            # Ajout des matrices dans les fenêtres de visualisation
            fourth.add_object('array', newobj=comp2, ToCheck=True, id='comp2')
            fifth.add_object('array', newobj=diff2, ToCheck=True, id='diff2=comp2-source')

            comp2.myops.myzones = src.myops.myzones
            diff2.myops.myzones = src.myops.myzones

            fourth.active_array = comp2
            fifth.active_array = diff2

    def set_blender_sculpting(self):
        """
        Mise en place de la structure nécessaire pour comparer la donnée de base avec la donnée sculptée sous Blender

        La donnée de base est la matrice contenue dans la fenêtre actuelle

        Fenêtres additionnelles :
            - information sur les volumes de déblai/remblai et bilan
            - matrice sculptée
            - différentiel entre scultage - source
            - gradient
            - laplacien
            - masque de modification
        """
        myframe = wx.Frame(None, title=_('Excavation and backfill'))
        sizergen = wx.BoxSizer(wx.VERTICAL)
        sizer1 = wx.BoxSizer(wx.HORIZONTAL)
        sizer2 = wx.BoxSizer(wx.HORIZONTAL)
        sizer3 = wx.BoxSizer(wx.HORIZONTAL)
        sizergen.Add(sizer1)
        sizergen.Add(sizer2)
        sizergen.Add(sizer3)

        labexc = wx.StaticText(myframe, label=_('Excavation : '))
        labback = wx.StaticText(myframe, label=_('Backfill   : '))
        labbal = wx.StaticText(myframe, label=_('Balance   : '))
        sizer1.Add(labexc)
        sizer2.Add(labback)
        sizer3.Add(labbal)

        font = wx.Font(18, wx.DECORATIVE, wx.NORMAL, wx.NORMAL)

        Exc = wx.StaticText(myframe, label=' [m³]')
        Back = wx.StaticText(myframe, label=' [m³]')
        Bal = wx.StaticText(myframe, label=' [m³]')

        labexc.SetFont(font)
        labback.SetFont(font)
        labbal.SetFont(font)
        Exc.SetFont(font)
        Back.SetFont(font)
        Bal.SetFont(font)

        sizer1.Add(Exc)
        sizer2.Add(Back)
        sizer3.Add(Bal)

        myframe.SetSizer(sizergen)
        myframe.Layout()
        myframe.Centre(wx.BOTH)
        myframe.Show()

        if self.link_params is None:
            self.link_params = {}

        self.link_params['ExcavationBackfill'] = myframe
        self.link_params['Excavation'] = Exc
        self.link_params['Backfill'] = Back
        self.link_params['Balance'] = Bal

        # Création de fenêtres de visualisation basées sur la classe "WolfMapViewer"
        first = self
        second = WolfMapViewer(None, 'Sculpting', w=600, h=600, wxlogging=self.wxlogging, wolfparent=self.wolfparent)
        third  = WolfMapViewer(None, 'Difference', w=600, h=600, wxlogging=self.wxlogging, wolfparent=self.wolfparent)
        fourth = WolfMapViewer(None, 'Gradient', w=600, h=600, wxlogging=self.wxlogging, wolfparent=self.wolfparent)
        fifth  = WolfMapViewer(None, 'Laplace', w=600, h=600, wxlogging=self.wxlogging, wolfparent=self.wolfparent)
        sixth  = WolfMapViewer(None, 'Unitary Mask', w=600, h=600, wxlogging=self.wxlogging, wolfparent=self.wolfparent)

        # Création d'une liste contenant les 3 instances d'objet "WolfMapViewer"
        list = []
        list.append(first)
        list.append(second)
        list.append(third)
        list.append(fourth)
        list.append(fifth)
        list.append(sixth)

        for curlist in list:
            curlist.add_grid()
            curlist.add_WMS()

        # On indique que les objets sont liés en actiavt le Booléen et en pointant la liste précédente
        for curlist in list:
            curlist.linked = True
            curlist.linkedList = list

        source: WolfArray
        sourcenew: WolfArray
        diff: WolfArray
        grad: WolfArray
        lap: WolfArray
        unimask: WolfArray

        source = self.active_array
        sourcenew = WolfArray(mold=source)

        # Création du différentiel -- Les opérateurs mathématiques sont surchargés
        diff = source - source
        grad = source.get_gradient_norm()
        lap = source.get_laplace()
        unimask = WolfArray(mold=diff)

        np.divide(diff.array.data, abs(diff.array.data), out=unimask.array.data, where=diff.array.data != 0.)

        grad.copy_mask(source, True)
        lap.copy_mask(source, True)
        diff.copy_mask(source, True)
        unimask.copy_mask(source, True)

        sourcenew.change_gui(second)
        diff.change_gui(third)
        grad.change_gui(fourth)
        lap.change_gui(fifth)
        unimask.change_gui(sixth)

        path = os.path.dirname(__file__)
        fn=join(path,'models\\diff16.pal')

        if exists(fn):
            diff.mypal.readfile(fn)
            diff.mypal.automatic=False
            diff.myops.palauto.SetValue(0)

        fn=join(path,'models\\diff3.pal')
        if exists(fn):
            unimask.mypal.readfile(fn)
            unimask.mypal.automatic=False
            unimask.myops.palauto.SetValue(0)

        # Ajout des matrices dans les fenêtres de visualisation
        second.add_object('array', newobj=sourcenew, ToCheck=True, id='source_new')
        third.add_object('array', newobj=diff, ToCheck=True, id='diff=comp-source')
        fourth.add_object('array', newobj=grad, ToCheck=True, id='gradient')
        fifth.add_object('array', newobj=lap, ToCheck=True, id='laplace')
        sixth.add_object('array', newobj=unimask, ToCheck=True, id='unimask')

        #pointage des vecteurs attachés à chaque matrice dans chaque GUI de façon à c que les modifications se répercutent  partout
        sourcenew.myops.myzones = source.myops.myzones
        diff.myops.myzones = source.myops.myzones
        grad.myops.myzones = source.myops.myzones
        lap.myops.myzones = source.myops.myzones
        unimask.myops.myzones = source.myops.myzones

        second.active_array = sourcenew
        third.active_array = diff
        fourth.active_array = grad
        fifth.active_array = lap
        sixth.active_array = unimask

        self.mimicme()

    def update_blender_sculpting(self):
        """ Mise à jour des fenêtres de visualisation pour la comparaison avec Blender """
        if not self.linked:
            return
        if len(self.linkedList) != 6:
            return

        # Création de fenêtres de visualisation basées sur la classe "WolfMapViewer"
        first = self.linkedList[0]
        second = self.linkedList[1]
        third = self.linkedList[2]
        fourth = self.linkedList[3]
        fifth = self.linkedList[4]
        sixth = self.linkedList[5]

        source = first.active_array
        sourcenew = second.active_array
        diff = third.active_array
        grad = fourth.active_array
        lap = fifth.active_array
        unimask = sixth.active_array

        fn = ''
        if self.link_params is not None:
            if 'gltf file' in self.link_params.keys():
                fn = self.link_params['gltf file']
                fnpos = self.link_params['gltf pos']

        if fn == '':
            for curgui in self.linkedList:
                if curgui.link_params is not None:
                    if 'gltf file' in curgui.link_params.keys():
                        fn = self.link_params['gltf file']
                        fnpos = self.link_params['gltf pos']
                        break

        with wx.BusyInfo(_('Importing gltf/glb')):
            wait = wx.BusyCursor()
            sourcenew.import_from_gltf(fn, fnpos, 'scipy')
            del wait

        with wx.BusyInfo(_('Update plots')):
            # Création du différentiel -- Les opérateurs mathématiques sont surchargés
            diff.array = (sourcenew - source).array
            grad.array = sourcenew.get_gradient_norm().array
            lap.array = sourcenew.get_laplace().array
            np.divide(diff.array.data, abs(diff.array.data), out=unimask.array.data, where=diff.array.data != 0.)

            diff.copy_mask(sourcenew, True)
            lap.copy_mask(sourcenew, True)
            grad.copy_mask(sourcenew, True)
            unimask.copy_mask(sourcenew, True)

            first.Paint()
            second.Paint()
            third.Paint()
            fourth.Paint()
            fifth.Paint()
            sixth.Paint()

            Exc: wx.StaticText
            Back: wx.StaticText
            Bal: wx.StaticText
            if not 'ExcavationBackfill' in self.link_params.keys():
                for curgui in self.linkedList:
                    if curgui.link_params is not None:
                        if 'ExcavationBackfill' in curgui.link_params.keys():
                            myframe = curgui.link_params['ExcavationBackfill']
                            Exc = curgui.link_params['Excavation']
                            Back = curgui.link_params['Backfill']
                            Bal = curgui.link_params['Balance']
            else:
                myframe = self.link_params['ExcavationBackfill']
                Exc = self.link_params['Excavation']
                Back = self.link_params['Backfill']
                Bal = self.link_params['Balance']

            Exc.SetLabel("{:.2f}".format(np.sum(diff.array[diff.array < 0.])) + ' [m³]')
            Back.SetLabel("{:.2f}".format(np.sum(diff.array[diff.array > 0.])) + ' [m³]')
            Bal.SetLabel("{:.2f}".format(np.sum(diff.array)) + ' [m³]')

    def zoomon_activevector(self, size:float=500., forceupdate:bool=True):
        """
        Zoom on active vector

        :param size: size of the zoomed window
        :param forceupdate: force the update of the window
        """

        if self.active_vector is None:
            logging.warning(_('No active vector'))
            return

        curvec = self.active_vector
        if curvec.xmin == -99999:
            curvec.find_minmax()

        bounds = [curvec.xmin, curvec.xmax, curvec.ymin, curvec.ymax]

        dx = bounds[1] - bounds[0]
        dy = bounds[3] - bounds[2]

        self.mousex = bounds[0] + dx / 2.
        self.mousey = bounds[2] + dy / 2.
        self.width = max(size, dx)
        self.height = max(size, dy)

        self.updatescalefactors()
        self.setbounds()
        self.mimicme()

        if forceupdate:
            self.update()
            if self.linked:
                for cur in self.linkedList:
                    if cur is not self:
                        cur.update()


    def zoomon_active_vertex(self, size:float = 20, forceupdate:bool = True):
        """
        Zoom on active vertex.

        :param size: size of the zoomed window
        :param forceupdate: force the update of the window
        """
        if self.active_vector is None:
            logging.warning(_('No active vector'))
            return
        curvec = self.active_vector
        if curvec.xmin == -99999:
            curvec.find_minmax()

        if self.active_vector is None:
            logging.warning(_('No active vector'))
            return

        grid = self.active_zones.xls
        row = grid.GetGridCursorRow()

        x = float(grid.GetCellValue(row, 0))
        y = float(grid.GetCellValue(row, 1))
        z = float(grid.GetCellValue(row, 2))
        curvert = wolfvertex(x, y, z)
        self.mousex = curvert.x
        self.mousey = curvert.y
        self.width = size
        self.height = size

        self.updatescalefactors()
        self.setbounds()
        self.mimicme()

        if forceupdate:
            self.update()
            if self.linked:
                for cur in self.linkedList:
                    if cur is not self:
                        cur.update()

    def zoom_on_id(self, id:str, drawtype:draw_type = draw_type.ARRAYS, forceupdate=True, canvas_height=1024):
        """
        Zoom on id

        :param id: id of the object to zoom on
        :param drawtype: type of object to zoom on - Different types elements can have the same id

        """

        if drawtype not in [draw_type.ARRAYS, draw_type.VECTORS]:
            logging.warning(_('Draw type must be either ARRAYS or VECTORS'))
            return

        obj = self.get_obj_from_id(id, drawtype)

        if obj is None:
            logging.warning(_('No object found with id {} and drawtype {}'.format(id, drawtype)))
            return

        if drawtype == draw_type.ARRAYS:
            self.zoom_on_array(obj, forceupdate=forceupdate, canvas_height=canvas_height)
        elif drawtype == draw_type.VECTORS:
            self.zoom_on_vector(obj, forceupdate=forceupdate, canvas_height=canvas_height)

    def zoom_on_array(self, array:WolfArray, forceupdate=True, canvas_height=1024):
        """ Zoom on array """

        bounds = array.get_bounds()

        center = [(bounds[0][1] + bounds[0][0]) / 2., (bounds[1][1] + bounds[1][0]) / 2.]
        width  = bounds[0][1] - bounds[0][0]
        height = bounds[1][1] - bounds[1][0]

        self.zoom_on({'center':center, 'width':width, 'height':height}, forceupdate=forceupdate, canvas_height=canvas_height)

    def zoom_on_vector(self, vector:vector, forceupdate=True, canvas_height=1024):
        """ Zoom on vector """

        if vector.xmin == -99999:
            vector.find_minmax()

        bounds = vector.get_bounds_xx_yy()

        center = [(bounds[0][1] + bounds[0][0]) / 2., (bounds[1][1] + bounds[1][0]) / 2.]
        width  = bounds[0][1] - bounds[0][0]
        height = bounds[1][1] - bounds[1][0]

        self.zoom_on({'center':center, 'width':width, 'height':height}, forceupdate=forceupdate, canvas_height=canvas_height)

    def create_Zones_from_arrays(self, arrays:list[WolfArray], id:str = None, add_legend:bool=True) -> Zones:
        """
        Create a Zones instance from list of WolfArrays

        One zone per array.

        One vector per zone with the masked contour.

        :param arrays: list of WolfArrays
        :param id: id of the Zones instance
        :param add_legend: add legend to the vector -- centroid of the contour

        """

        msg = _('This function will force a null border of 1 cell on each array to avoid issues with the contouring algorithm')
        dlg = wx.MessageDialog(None, msg, _('Warning'), wx.OK | wx.CANCEL | wx.ICON_WARNING)
        result = dlg.ShowModal()
        dlg.Destroy()
        if result != wx.ID_OK:
            logging.info(_('Operation cancelled by user'))
            return None

        # création de l'instance de Zones
        new_zones = Zones(idx = 'contour' if id is None else id.lower(), mapviewer=self)

        for curarray in arrays:

            if isinstance(curarray, WolfArray):

                curarray.nullify_border(1)
                curarray.reset_plot()

                new_zone = zone(name = curarray.idx)
                new_zones.add_zone(new_zone, forceparent=True)

                sux, sux, curvect, interior = curarray.suxsuy_contour()
                new_zone.add_vector(curvect, forceparent=True)

                curvect.set_legend_to_centroid(curarray.idx)
                curvect.myprop.width = 2

                rectvect = vector(name = 'rect_boundary')
                new_zone.add_vector(rectvect, forceparent=True)

                bounds = curarray.get_bounds()

                rectvect.add_vertex(wolfvertex(bounds[0][0], bounds[1][0]))
                rectvect.add_vertex(wolfvertex(bounds[0][1], bounds[1][0]))
                rectvect.add_vertex(wolfvertex(bounds[0][1], bounds[1][1]))
                rectvect.add_vertex(wolfvertex(bounds[0][0], bounds[1][1]))
                rectvect.close_force()

                rectvect.myprop.color = getIfromRGB([255,0,0])
                rectvect.myprop.width = 2

                logging.info(_('{} treated'.format(curarray.idx)))
            else:
                logging.warning(_('All elements in the list must be of type WolfArray'))

        new_zones.find_minmax(update=True)

        return new_zones


    def zoom_on(self, zoom_dict = None, width = 500, height = 500, center = None, xll = None, yll = None, forceupdate=True, canvas_height=1024):
        """
        Zoom on a specific area

        It is possible to zoom on a specific area by providing the zoom parameters in :
          - a dictionnary
          - width and height of the zoomed window and the lower left corner coordinates
          - width and height of the zoomed window and the center coordinates

        :param zoom_dict: dictionnary containing the zoom parameters - possible keys : 'width', 'height', 'center', 'xmin', 'ymin', 'xmax', 'ymax'
        :param width: width of the zoomed window [m]
        :param height: height of the zoomed window [m]
        :param center: center of the zoomed window [m] - tuple (x,y)
        :param xll: lower left X coordinate of the zoomed window [m]
        :param yll: lower left Y coordinate of the zoomed window [m]
        :param forceupdate: force the update of the window
        :param canvas_height: height of the canvas [pixels]


        Examples :

          - zoom_on(zoom_dict = {'center':(500,500), 'width':1000, 'height':1000})
          - zoom_on(width = 1000, height = 1000, xll = 500, yll = 500)
          - zoom_on(width = 1000, height = 1000, center = (500,500))
        """
        if zoom_dict is not None:
            width  = 99999
            height = 99999
            xll  = 99999
            yll = 99999
            xmax  = 99999
            ymax = 99999
            if 'center' in zoom_dict.keys():
                center = zoom_dict['center']
            if 'width' in zoom_dict.keys():
                width = zoom_dict['width']
            if 'height' in zoom_dict.keys():
                height = zoom_dict['height']
            if 'xmin' in zoom_dict.keys():
                xll = zoom_dict['xmin']
            if 'ymin' in zoom_dict.keys():
                yll = zoom_dict['ymin']
            if 'xmax' in zoom_dict.keys():
                xmax = zoom_dict['xmax']
            if 'ymax' in zoom_dict.keys():
                ymax = zoom_dict['ymax']

            if width == 99999:
                width = xmax-xll
            if height == 99999:
                height = ymax-yll

        if center is not None and len(center)==2:
            self.mousex = center[0]
            self.mousey = center[1]
            self.width = width
            self.height = height
        elif (xll is not None) and (yll is not None):
            self.mousex = xll + width/2
            self.mousey = yll + height/2
            self.width = width
            self.height = height

        # fixe la taille de la fenêtre
        v_height = canvas_height
        v_width = int(v_height*(float(width)/float(height)))

        self.SetClientSize(v_width + self.treewidth, v_height)

        self.updatescalefactors()
        self.mimicme()

        if forceupdate:
            self.update()
            if self.linked:
                for cur in self.linkedList:
                    if cur is not self:
                        cur.update()

    def zoom_on_active_profile(self, size:float=500., forceupdate:bool=True):
        """ Zoom on active profile """

        curvec = self.active_profile
        if curvec.xmin == -99999:
            curvec.find_minmax()

        bounds = [curvec.xmin, curvec.xmax, curvec.ymin, curvec.ymax]

        dx = bounds[1] - bounds[0]
        dy = bounds[3] - bounds[2]

        self.mousex = bounds[0] + dx / 2.
        self.mousey = bounds[2] + dy / 2.
        self.width = max(size, dx)
        self.height = max(size, dy)

        self.updatescalefactors()
        self.setbounds()
        self.mimicme()

        if forceupdate:
            self.update()
            if self.linked:
                for cur in self.linkedList:
                    if cur is not self:
                        cur.update()

    def read_project(self, fn):
        """
        Projet WOLF GUI

        Fichier de paramètres contenant les types et chemins d'accès aux données à ajouter

        A compléter...

        """
        curdir = Path(os.getcwd())

        real_ids = {}

        myproject = Wolf_Param(None, filename=fn, toShow=False)

        def check_params(myproject, curgroup) -> bool:

            check = True
            pot_keys = list(PROJECT_GROUP_KEYS[curgroup].keys())

            for curkey in pot_keys:
                if 'mandatory' in PROJECT_GROUP_KEYS[curgroup][curkey]:
                    if not myproject.is_in(curgroup, curkey):
                        logging.warning(_('Missing key : ')+ curkey)
                        check = False
            return check

        def sanit_id(id:str, drawtype:draw_type) -> str:
            existing_id = self.get_list_keys(drawtype, None)

            while id in existing_id:
                logging.warning(_('ID already exists - Changing it...'))
                id = id + '_'

            return id

        # COMPLEX ACTIONS
        curgroup = PROJECT_ACTION
        if myproject.is_in(curgroup):

            pot_keys = list(PROJECT_GROUP_KEYS[curgroup].keys())

            for curkey in pot_keys:
                which = myproject[(curgroup, curkey)]

                pot_val = list(PROJECT_GROUP_KEYS[curgroup][curkey].keys())

                if which in pot_val:
                    if which == 'compare_arrays':

                        # Comparaison de plusieurs matrices

                        logging.info(_('Compare action - Searching for arrays to compare...'))
                        ListCompare = []
                        if myproject.is_in('array'):
                            for curval in myproject.get_group('array').values():
                                curid = curval[key_Param.NAME]
                                logging.info(_('Array to compare : ')+ curid)
                                ListCompare.append(WolfArray(Path(myproject[('array', curid)])))
                        else:
                            logging.warning(_('No array to compare - Aborting !'))
                            return

                        logging.info(_('Setting compare...'))
                        self.set_compare(ListCompare)
                        logging.info(_('Compare set !'))
                        return
                else:
                    logging.error(_('Bad parameter in project file - action : ')+ which)

        # CROSS SECTIONS
        curgroup = PROJECT_CS
        if myproject.is_in(curgroup):
            if check_params(myproject, curgroup):
                for curval in myproject.get_group(curgroup).values():
                    curid = curval[key_Param.NAME]
                    if curid != 'format' and curid != 'dirlaz':
                        mycs = crosssections(myproject[(curgroup, curid)],
                                            format = myproject[(curgroup, 'format')],
                                            dirlaz = myproject[(curgroup, 'dirlaz')],
                                            mapviewer = self)

                        locid = real_ids[(draw_type.VECTORS, curid)] = sanit_id(curid, draw_type.VECTORS)

                        self.add_object(curgroup, newobj=mycs, id=locid)
            else:
                logging.warning(_('Bad parameter in project file - cross_sections'))

        # TILES
        curgroup = PROJECT_TILES
        if myproject.is_in(curgroup):
            if check_params(myproject, curgroup):
                curid      = myproject.get_param(curgroup, 'id')
                curfile    = myproject.get_param(curgroup, 'tiles_file')
                curdatadir = myproject.get_param(curgroup, 'data_dir')
                curcompdir = myproject.get_param(curgroup, 'comp_dir')

                if exists(curfile):
                    try:
                        mytiles = Tiles(filename= curfile, parent=self, linked_data_dir=curdatadir)
                        mytiles.set_comp_dir(curcompdir)

                        locid = real_ids[(draw_type.TILES, curid)] = sanit_id(curid, draw_type.TILES)
                        self.add_object(curgroup, newobj=mytiles, id=locid)
                    except Exception as e:
                        logging.error(_('Error in tiles import : ')+ str(e))
                else:
                    logging.warning(_('File does not exist : ')+ str(curfile))
            else:
                logging.warning(_('Bad parameter in project file - tiles'))

        # LAZ GRID
        curgroup = PROJECT_LAZ
        if myproject.is_in(curgroup):
            if check_params(myproject, curgroup):
                try:
                    self.init_laz_from_gridinfos(curdir / myproject[curgroup, 'data_dir'], myproject[(curgroup, 'classification')])
                except Exception as e:
                    logging.error(_('Error in laz_grid import : ')+ str(e))
            else:
                logging.warning(_('Bad parameter in project file - laz_grid'))

        # VECTOR DATA
        curgroup = PROJECT_VECTOR
        if myproject.is_in(curgroup):
            if check_params(myproject, curgroup):
                for curval in myproject.get_group(curgroup).values():
                    curid = curval[key_Param.NAME]
                    name  = curval[key_Param.VALUE]
                    if exists(name):
                        try:
                            myvec = Zones(name, parent = self, mapviewer = self)

                            locid = real_ids[(draw_type.VECTORS, curid)] = sanit_id(curid, draw_type.VECTORS)
                            self.add_object(curgroup, newobj = myvec, id = locid)
                        except Exception as e:
                            logging.error(_('Error in vector import : ')+ str(e))
                    else:
                        logging.info(_('File does not exist : ') + str(name))
            else:
                logging.warning(_('Bad parameter in project file - vector'))

        # ARRAY DATA
        curgroup = PROJECT_ARRAY
        if myproject.is_in(curgroup):
            if check_params(myproject, curgroup):
                for curval in myproject.get_group(curgroup).values():
                    curid = curval[key_Param.NAME]
                    name  = curdir / Path(curval[key_Param.VALUE])
                    if exists(name):
                        try:
                            curarray = WolfArray(name, mapviewer = self)

                            locid = real_ids[(draw_type.ARRAYS, curid)] = sanit_id(curid, draw_type.ARRAYS)
                            self.add_object('array', newobj=curarray, id = locid)
                        except Exception as e:
                            logging.error(_('Error in array import : ')+ str(e))
                    else:
                        logging.info(_('File does not exist : ') + str(name))
            else:
                logging.warning(_('Bad parameter in project file - array'))

        # CLOUD DATA
        curgroup = PROJECT_CLOUD
        if myproject.is_in(curgroup):
            if check_params(myproject, curgroup):
                for curval in myproject.get_group(curgroup).values():
                    curid = curval[key_Param.NAME]
                    name  = curval[key_Param.VALUE]
                    if exists(name):
                        try:
                            mycloud = cloud_vertices(name, mapviewer = self)

                            locid = real_ids[(draw_type.CLOUD, curid)] = sanit_id(curid, draw_type.CLOUD)
                            self.add_object('cloud', newobj = mycloud, id = locid)
                        except Exception as e:
                            logging.error(_('Error in cloud import : ') + str(e))
                    else:
                        logging.info(_('File does not exist : ') + str(name))
            else:
                logging.warning(_('Bad parameter in project file - cloud'))

        # 2D RESULTS

        # CPU code
        curgroup = PROJECT_WOLF2D
        if myproject.is_in(curgroup):
            if check_params(myproject, curgroup):
                for curval in myproject.get_group(curgroup).values():
                    curid  = curval[key_Param.NAME]
                    simdir = Path(curval[key_Param.VALUE])
                    if simdir.exists():
                        try:
                            curwolf = Wolfresults_2D(simdir, mapviewer = self)

                            locid = real_ids[(draw_type.RES2D, curid)] = sanit_id(curid, draw_type.RES2D)
                            self.add_object('res2d', newobj = curwolf, id = locid)
                        except Exception as e:
                            logging.error(_('Error in wolf2d import : ')+ str(e))
                    else:
                        logging.info(_('Directory does not exist ')) + str(simdir)

                self.menu_wolf2d()
            else:
                logging.warning(_('Bad parameter in project file - wolf2d'))

        # GPU code
        curgroup = PROJECT_GPU2D
        if myproject.is_in(curgroup):
            if check_params(myproject, curgroup):

                pgbar = wx.ProgressDialog(_('Loading GPU results'), _('Loading GPU results'), maximum=len(myproject.myparams[curgroup].keys()), parent=self, style=wx.PD_APP_MODAL | wx.PD_AUTO_HIDE)

                for curval in myproject.get_group(curgroup).values():
                    curid  = curval[key_Param.NAME]
                    simdir = Path(curval[key_Param.VALUE])
                    if simdir.exists():
                        try:
                            curwolf = wolfres2DGPU(curdir / simdir, mapviewer = self)

                            locid = real_ids[(draw_type.RES2D, curid)] = sanit_id(curid, draw_type.RES2D)
                            self.add_object('res2d', newobj = curwolf, id = locid)
                        except Exception as e:
                            logging.error(_('Error in gpu2d import : ')+ str(e))
                    else:
                        logging.info(_('Bad directory : ') + str(simdir))

                    pgbar.Update(pgbar.GetValue() + 1)

                pgbar.Destroy()

                self.menu_wolf2d()
                self.menu_2dgpu()
            else:
                logging.warning(_('Bad parameter in project file - gpu2d'))

        # PALETTE/COLORMAP
        curgroup = PROJECT_PALETTE
        if myproject.is_in(curgroup):
            if check_params(myproject, curgroup):
                self.project_pal = {}
                for curval in myproject.get_group(curgroup).values():
                    curid = curval[key_Param.NAME]
                    name  = Path(curval[key_Param.VALUE])
                    if name.exists():
                        if name.suffix == '.pal':
                            mypal = wolfpalette(None, '')
                            mypal.readfile(name)
                            mypal.automatic = False

                            self.project_pal[curid] = mypal
                        else:
                            logging.warning(_('Bad palette file : ')+ str(name))
                    else:
                        logging.info(_('Bad parameter in project file - palette : ')+ str(name))
            else:
                logging.warning(_('Bad parameter in project file - palette'))

        # LINKS
        curgroup = PROJECT_PALETTE_ARRAY
        if myproject.is_in(curgroup):
            if check_params(myproject, curgroup):
                curarray: WolfArray
                if self.project_pal is not None:

                    for curval in myproject.get_group(curgroup).keys():
                        id_array = curval[key_Param.NAME]
                        id_pal = curval[key_Param.VALUE]
                        if id_pal in self.project_pal.keys():
                            try:
                                curarray = self.getobj_from_id(real_ids[(draw_type.ARRAYS, id_array)])
                                if curarray is not None:
                                    mypal:wolfpalette
                                    mypal = self.project_pal[id_pal]
                                    curarray.mypal = mypal
                                    if mypal.automatic:
                                        curarray.myops.palauto.SetValue(1)
                                    else:
                                        curarray.myops.palauto.SetValue(0)
                                    curarray.updatepalette(0)
                                    curarray.reset_plot()
                                else:
                                    logging.warning(_('Bad parameter in project file - palette-array : ')+ str(id_array))
                            except Exception as e:
                                logging.error(_('Error in palette-array link : ')+ str(e))
                        else:
                            logging.warning(_('Bad parameter in project file - palette-array : ')+ str(id_pal))
                else:
                    logging.warning(_('No palettes found in project file ! -- Add palette group in the .proj'))
            else:
                logging.warning(_('Bad parameter in project file - palette-array'))

        curgroup = PROJECT_LINK_CS
        if myproject.is_in(curgroup):
            if self.active_cs is not None:
                if check_params(myproject, curgroup):

                    idx = real_ids[(draw_type.VECTORS, myproject[(curgroup, 'linkzones')])]
                    curzones = self.get_obj_from_id(idx, draw_type.VECTORS)

                    if curzones is not None:
                        self.active_cs.link_external_zones(curzones)

                    zonename = myproject[(curgroup, 'sortzone')]
                    vecname  = myproject[(curgroup, 'sortname')]
                    downfirst = myproject[(curgroup, 'downfirst')]

                    downfirst = False
                    if downfirst == 1 or str(downfirst).lower() == 'true':
                        downfirst = True

                    if zonename != '' and vecname != '':
                        curvec = curzones[(zonename, vecname)]
                        if curvec is not None:
                            try:
                                self.active_cs.sort_along(curvec.asshapely_ls(), curvec.myname, downfirst)
                            except Exception as e:
                                logging.error(_('Error in cross_sections_link sorting : ')+ str(e))
                        else:
                            logging.warning(_('Bad id for sorting vector in project file - cross_sections_link'))
                else:
                    logging.warning(_('Bad parameter in project file - cross_sections_link'))
            else:
                logging.warning(_('No active cross section to link !'))

        curgroup = PROJECT_LINK_VEC_ARRAY
        # Useful to mask data outside of the linked contour
        if myproject.is_in(curgroup):
            if check_params(myproject, curgroup):
                for curval in myproject.get_group(curgroup).keys():

                    id_array = real_ids[(draw_type.ARRAYS, curval[key_Param.NAME])]
                    id_zones = real_ids[(draw_type.VECTORS, curval[key_Param.VALUE])]

                    locarray:WolfArray
                    locvec:Zones

                    locarray = self.get_obj_from_id(id_array, draw_type.ARRAYS)
                    if locarray is None:
                        locarray = self.get_obj_from_id(id_array, draw_type.RES2D)

                    locvec   = self.get_obj_from_id(id_zones, draw_type.VECTORS)

                    if locvec is not None and locarray is not None:
                        try:
                            if locvec.nbzones == 1:
                                if locvec.myzones[0].nbvectors == 1:
                                    locarray.linkedvec = locvec.myzones[0].myvectors[0]
                                else:
                                    logging.warning(_('In vec-array association, You must have only 1 zone and 1 polyline !'))
                            else:
                                logging.warning(_('In vec-array association, You must have only 1 zone and 1 polyline !'))

                        except Exception as e:
                            logging.error(_('Error in vector_array_link : ')+ str(e))
                    else:
                        logging.warning(_('Bad vec-array association in project file !'))
            else:
                logging.warning(_('Bad parameter in project file - vector_array_link'))

    def save_project(self, fn, absolute:bool = True):
        """ Save project file """

        dirproj = Path(fn).parent

        def new_path(drawtype:draw_type, id:str) -> str:
            logging.info(_('Empty path but I need a path !'))

            path = ''

            ext = 'All files (*.*)|*.*'
            if drawtype == draw_type.ARRAYS:
                ext += '|Binary files (*.bin)|*.bin|Tiff files (*.tif)|*.tif|Numpy files (*.npy)|*.npy'
            elif drawtype == draw_type.VECTORS:
                ext += '|VecZ files (*.vecz)|*.vecz|Vec files (*.vec)|*.vec'
            elif drawtype == draw_type.CLOUD:
                ext += '|Cloud files (*.xyz)|*.xyz'

            dlg = wx.FileDialog(None, _('Choose a filename for ') + id, str(dirproj), '', ext, wx.FD_SAVE | wx.FD_OVERWRITE_PROMPT)
            ret = dlg.ShowModal()

            if ret == wx.ID_OK:
                path = Path(dlg.GetPath())

            return path

        def sanit_path(path:Path, absolute:bool, drawtype:draw_type) -> str:

            path = Path(path)

            if not path.exists():
                logging.info(_('Path does not exist : ')+ str(path))

            if absolute:
                return str(path)
            else:
                try:
                    return os.path.relpath(path, dirproj)
                except:
                    logging.error(_('Error in relative path : ')+ str(path) + " - " + str(dirproj))
                    logging.info(_('Returning absolute path instead !'))
                    return str(path.absolute())

        myproject = Wolf_Param(None, toShow=False, to_read=False, filename=fn, init_GUI=False)

        # matrices
        try:
            curgroup = PROJECT_ARRAY
            for curel in self.iterator_over_objects(draw_type.ARRAYS):
                curel:WolfArray
                if curel.filename == '':
                    newpath =  new_path(draw_type.ARRAYS, curel.idx)
                    if newpath == '':
                        logging.warning(_('No path for array : ')+ curel.idx  + _(' - Ignoring it !'))
                        continue
                    curel.write_all(newpath)

                curpath = sanit_path(curel.filename, absolute, draw_type.ARRAYS)

                myproject.add_param(curgroup, curel.idx, curpath)
        except:
            logging.error(_('Error in saving arrays'))

        # résultats 2D
        try:
            curgroup = PROJECT_WOLF2D
            for curel in self.iterator_over_objects(draw_type.RES2D):
                if type(curel) == Wolfresults_2D:
                    myproject.add_param(curgroup, curel.idx, sanit_path(curel.filename, absolute, draw_type.RES2D))

            curgroup = PROJECT_GPU2D
            for curel in self.iterator_over_objects(draw_type.RES2D):
                if type(curel) == wolfres2DGPU:
                    myproject.add_param(curgroup, curel.idx, sanit_path(curel.filename, absolute, draw_type.RES2D))
        except:
            logging.error(_('Error in saving 2D results'))

        # vecteurs
        try:
            curgroup = PROJECT_VECTOR
            for curel in self.iterator_over_objects(draw_type.VECTORS):
                if isinstance(curel, crosssections):
                    continue

                curel:Zones
                if curel.filename == '':
                    newpath =  new_path(draw_type.VECTORS, curel.idx)
                    if newpath == '':
                        logging.warning(_('No path for vector : ')+ curel.idx  + _(' - Ignoring it !'))
                        continue
                    curel.saveas(newpath)

                myproject.add_param(curgroup, curel.idx, sanit_path(curel.filename, absolute, draw_type.VECTORS))
        except:
            logging.error(_('Error in saving vectors'))

        # cross sections
        try:
            curgroup = PROJECT_CS
            for curel in self.iterator_over_objects(draw_type.VECTORS):
                if isinstance(curel, crosssections):
                    myproject.add_param(curgroup, curel.idx, sanit_path(curel.filename, absolute, draw_type.VECTORS))
        except:
            logging.error(_('Error in saving cross sections'))

        # nuages de points
        try:
            curgroup = PROJECT_CLOUD
            for curel in self.iterator_over_objects(draw_type.CLOUD):
                myproject.add_param(curgroup, curel.idx, sanit_path(curel.filename, absolute, draw_type.CLOUD))
        except:
            logging.error(_('Error in saving clouds'))

        # palettes
        try:
            if self.project_pal is not None:
                curgroup = PROJECT_PALETTE
                for curel in self.project_pal.keys():
                    myproject.add_param(curgroup, curel, sanit_path(self.project_pal[curel].filename, absolute, draw_type.OTHER))
        except:
            logging.error(_('Error in saving palettes'))

        # tiles
        try:
            curgroup = PROJECT_TILES
            for curel in self.iterator_over_objects(draw_type.TILES):
                myproject.add_param(curgroup, curel.idx, sanit_path(curel.filename, absolute, draw_type.OTHER))
                myproject.add_param(curgroup, 'data_dir', sanit_path(curel.linked_data_dir, absolute, draw_type.OTHER))
                myproject.add_param(curgroup, 'comp_dir', sanit_path(curel.linked_data_dir_comp, absolute, draw_type.OTHER))
        except:
            logging.error(_('Error in saving tiles'))

        # LAZ GRID
        try:
            if self.mylazgrid is not None:
                curgroup = PROJECT_LAZ
                myproject.add_param(curgroup, 'data_dir', sanit_path(self.mylazgrid.dir, absolute, draw_type.OTHER))
                myproject.add_param(curgroup, 'classification', self.mylazgrid.colors.class_name)
        except:
            logging.error(_('Error in saving laz grid'))

        myproject.Save(fn)


    def help_project(self):
        """ Help for project file.

        Define which elements can be saved in a project file.

        """

        logging.info(_('Project file help'))

        logging.info(_('Project file is a file containing some information about the current project.'))

        logging.info(_('It can contain the following informations :'))

        logging.info(_(' - Arrays :'))
        logging.info(_('   - id'))
        logging.info(_('   - filename in relative or absolute path'))

        logging.info(_(' - Cross sections :'))
        logging.info(_('   - id'))
        logging.info(_('   - filename in relative or absolute path'))

        logging.info(_(' - Vectors :'))
        logging.info(_('   - id'))
        logging.info(_('   - filename in relative or absolute path'))

        logging.info(_(' - Clouds :'))
        logging.info(_('   - id'))
        logging.info(_('   - filename in relative or absolute path'))

        logging.info(_(' - Tiles :'))
        logging.info(_('   - id'))
        logging.info(_('   - filename in relative or absolute path'))

        logging.info(_(' - LAZ grid :'))
        logging.info(_('   - data_dir : directory containing the NUMPY grid'))
        logging.info(_('   - classification : classification of the laz files'))

        logging.info(_(' - Palettes :'))
        logging.info(_('   - id'))
        logging.info(_('   - filename in relative or absolute path'))

        logging.info(_(' - Wolf2D CPU results :'))
        logging.info(_('   - id'))
        logging.info(_('   - filename in relative or absolute path'))

        logging.info(_(' - Wolf2D GPU results :'))
        logging.info(_('   - id'))
        logging.info(_('   - filename in relative or absolute path'))

        logging.info(_(' - Palette-Array links :'))
        logging.info(_('   - id of the array'))
        logging.info(_('   - id of the palette'))

        logging.info(_(' - Vector-Array links :'))
        logging.info(_('   - id of the array'))
        logging.info(_('   - id of the vector (containing only 1 zone and 1 vector)'))

        logging.info(_(' - Cross section links :'))
        logging.info(_('   - id of the cross section'))
        logging.info(_('   - id of the vector to sort along'))
        logging.info(_('   - id of the zone to link'))
        logging.info(_('   - downfirst : True or False'))

        logging.info('')
        logging.info(_('A tabulation is used to separate the value and the key.'))
        logging.info('')

        logging.info(_('Exemple :'))

        logging.info('')

        logging.info('array:')
        logging.info('myid1\tmyfilename_array1')
        logging.info('myid2\tmy../filename_array2')
        logging.info('vector:')
        logging.info('myvec1\tmy../../filename_vecz1')
        logging.info('myvec2\tmyfilename_vecz2')
        logging.info('laz_grid:')
        logging.info('data_dir\tD:\\MODREC-Vesdre\\LAZ_Vesdre\\2023\\grids_flt32')
        logging.info('classification\tSPW-Geofit 2023')

    def plot_laz_around_active_vec(self):
        """ Plot laz data around active vector """

        if self.active_vector is None:
            logging.warning(_('Please activate a vector'))
            return

        if self.mylazgrid is None:
            logging.warning(_('No laz grid'))
            return

        dlg = wx.NumberEntryDialog(None, _('Enter the size of the window around the active vector [cm]'), _('Window size'),_('Window size'), 500, 0, 2000)
        ret = dlg.ShowModal()
        if ret != wx.ID_OK:
            dlg.Destroy()
            return

        value = dlg.GetValue()/100.
        dlg.Destroy()

        fig = self.mylazgrid.plot_laz_wx(self.active_vector.asshapely_ls(), length_buffer=value, show=True)

        if self.active_array is not None:
            copy_vec = vector()
            copy_vec.myvertices = self.active_vector.myvertices.copy()
            copy_vec.split(abs(self.active_array.dx)/2., False)
            copy_vec.get_values_on_vertices(self.active_array)
            s,z = copy_vec.get_sz()
            notmasked = np.where(z != -99999.)

            fig.plot(s[notmasked],z[notmasked], c='black', linewidth=2.0)


    def clip_laz_gridded(self):
        """ Clip laz grid on current zoom """

        if self.mylazgrid is None:
            logging.warning(_('No laz grid -- Please initialize it !'))
            return

        curbounds = [[self.xmin, self.xmin + self.width], [self.ymin, self.ymin + self.height]]

        if self.active_laz is None:
            newobj = Wolf_LAZ_Data()
            newobj.classification = self.mylazgrid.colors
            newobj.from_grid(self.mylazgrid, curbounds)

            self.add_object('laz', newobj= newobj)

        else:
            dlg = wx.MessageDialog(None, _('Do you want to keep the current data ?'), _('Keep data ?'), wx.YES_NO | wx.ICON_QUESTION)
            ret = dlg.ShowModal()

            if ret == wx.ID_YES:
                newobj = Wolf_LAZ_Data()
                newobj.classification = self.mylazgrid.colors
                newobj.from_grid(self.mylazgrid, curbounds)

                self.add_object('laz', newobj= newobj)

            else:

                self.active_laz.from_grid(self.mylazgrid, curbounds)

        logging.info(_('Clip LAZ grid on current zoom'))
        logging.info(_('Bounds {}-{}  {}-{}').format(curbounds[0][0],curbounds[0][1],curbounds[1][0],curbounds[1][1]))
        logging.info(_('Nb points : {:_}').format(self.active_laz.num_points))

    def filter_active_laz(self):
        """ Filter active laz data """

        if self.active_laz is None:
            logging.warning(_('No laz data'))
            return

        codes = self.active_laz.codes_unique()

        names = [self.active_laz.classification.classification[curcode][0] for curcode in codes]

        with wx.MultiChoiceDialog(None, _('Choose the codes to keep'), _('Codes'), names) as dlg:
            if dlg.ShowModal() == wx.ID_OK:
                used_codes = dlg.GetSelections()
                used_codes = [codes[cur] for cur in used_codes]
                self.active_laz.filter_data(used_codes)

                logging.info(_('Filter done - Nb points : {:_}').format(self.active_laz.num_points))
            else:
                logging.info(_('Filter cancelled'))

    def descimate_laz_data(self, factor:int = 10):
        """ Descimate data """

        if self.active_laz is None:
            logging.warning(_('No laz data'))
            return

        self.active_laz.descimate(factor)

    def select_active_array_from_laz(self, array:WolfArray = None, used_codes:list = None, chunk_size:float = 500.):
        """ select some nodes from laz data

        :param array: array to fill
        :param used_codes: codes to use
        """
        if self.mylazgrid is None:
            logging.info(_('No laz grid - Aborting !'))
            return

        if array is None:
            logging.error(_('No array'))
            return

        if used_codes is None:
            keycode = [key for key,val in self.mylazgrid.colors.classification.items()]
            names = [val[0] for key,val in self.mylazgrid.colors.classification.items()]

            with wx.MultiChoiceDialog(None, _('Choose the codes to use'), _('Codes'), names) as dlg:
                if dlg.ShowModal() == wx.ID_OK:
                    used_codes = dlg.GetSelections()
                    used_codes = [float(keycode[cur]) for cur in used_codes]
                else:
                    return

        curbounds = array.get_bounds()

        # align bounds on chunk_size
        curbounds[0][0] = curbounds[0][0] - curbounds[0][0] % chunk_size
        curbounds[0][1] = curbounds[0][1] + chunk_size - curbounds[0][1] % chunk_size
        curbounds[1][0] = curbounds[1][0] - curbounds[1][0] % chunk_size
        curbounds[1][1] = curbounds[1][1] + chunk_size - curbounds[1][1] % chunk_size

        chunck_x = np.arange(curbounds[0][0], curbounds[0][1], chunk_size)
        chunck_y = np.arange(curbounds[1][0], curbounds[1][1], chunk_size)

        for curx in tqdm(chunck_x, 'Chunks'):
            for cury in chunck_y:
                curbounds = [[curx, curx + chunk_size], [cury, cury + chunk_size]]

                logging.info(_('Scan {}-{}  {}-{}').format(curbounds[0][0],curbounds[0][1],curbounds[1][0],curbounds[1][1]))
                mylazdata = self.mylazgrid.scan(curbounds)
                # logging.info(_('Scan done'))

                data = {}
                for curcode in used_codes:
                    data[curcode] = mylazdata[mylazdata[:, 3] == curcode]

                for curdata in data.values():

                    if curdata.shape[0] == 0:
                        continue

                    i,j = array.get_ij_from_xy(curdata[:, 0], curdata[:, 1])

                    keys = np.vstack((i,j)).T

                    # unique keys
                    keys = np.unique(keys, axis=0)

                    array.SelectionData._add_nodes_to_selectionij(keys, verif = False)

        array.SelectionData.update_nb_nodes_selection()
        self.Paint()

        logging.info(_('Selection done'))

    def fill_active_array_from_laz(self, array:WolfArray = None, used_codes:list = [], operator:int = -1, chunk_size:float = 500.):
        """ Fill active array with laz data

        :param array: array to fill
        :param used_codes: codes to use
        :param operator: operator to use
        """

        if self.mylazgrid is None:
            logging.info(_('No laz grid - Aborting !'))
            return

        if array is None:
            logging.error(_('No array'))
            return

        if len(used_codes) == 0 :
            keycode = [key for key,val in self.mylazgrid.colors.classification.items()]
            names = [val[0] for key,val in self.mylazgrid.colors.classification.items()]

            with wx.MultiChoiceDialog(None, _('Choose the codes to use'), _('Codes'), names) as dlg:
                if dlg.ShowModal() == wx.ID_OK:
                    data = {}
                    used_codes = dlg.GetSelections()
                    used_codes = [float(keycode[cur]) for cur in used_codes]
                else:
                    return

        if operator == -1:
            with wx.SingleChoiceDialog(None, _('Choose the operator'), _('Operator'), ['max', 'percentile 95', 'percentile 5', 'min', 'mean', 'median', 'sum']) as dlg:
                if dlg.ShowModal() == wx.ID_OK:
                    if dlg.GetStringSelection() == 'max':
                        operator = np.max
                    elif dlg.GetStringSelection() == 'min':
                        operator = np.min
                    elif dlg.GetStringSelection() == 'mean':
                        operator = np.mean
                    elif dlg.GetStringSelection() == 'median':
                        operator = np.median
                    elif dlg.GetStringSelection() == 'sum':
                        operator = np.sum
                    elif dlg.GetStringSelection() == 'percentile 95':
                        operator = lambda x: np.percentile(x, 95)
                    elif dlg.GetStringSelection() == 'percentile 5':
                        operator = lambda x: np.percentile(x, 5)
                else:
                    return

        with wx.NumberEntryDialog(None, _('Minimum number of points to operate'), _('Minimum'), _('Minimum points'), 1, 1, 20) as dlg:
            if dlg.ShowModal() == wx.ID_OK:
                minpoints = dlg.GetValue()
            else:
                return

        logging.info(_('This could take some time for large area...\n Take a coffee and relax!'))

        bounds = array.get_bounds()

        # align bounds on chunk_size
        bounds[0][0] = bounds[0][0] - bounds[0][0] % chunk_size
        bounds[0][1] = bounds[0][1] + chunk_size - bounds[0][1] % chunk_size
        bounds[1][0] = bounds[1][0] - bounds[1][0] % chunk_size
        bounds[1][1] = bounds[1][1] + chunk_size - bounds[1][1] % chunk_size

        chunks_x = np.arange(bounds[0][0], bounds[0][1], chunk_size)
        chunks_y = np.arange(bounds[1][0], bounds[1][1], chunk_size)

        for curx in tqdm(chunks_x, 'Chunks'):
            for cury in chunks_y:

                curbounds = [[curx, curx + chunk_size], [cury, cury + chunk_size]]

                logging.info(_('Scan {}-{}  {}-{}').format(curbounds[0][0],curbounds[0][1],curbounds[1][0],curbounds[1][1]))
                mylazdata = self.mylazgrid.scan(curbounds)
                # logging.info(_('Scan done'))

                if len(mylazdata) == 0:
                    continue

                # Test codes
                data = {}
                for curcode in used_codes:
                    data[curcode] = mylazdata[mylazdata[:, 3] == curcode]

                # Treat data for each code
                for curdata in data.values():

                    if curdata.shape[0] == 0:
                        continue
                    else:
                        logging.info(_('Code {} : {} points'.format(curdata[0,3], curdata.shape[0])))

                    # get i,j from x,y
                    i,j = array.get_ij_from_xy(curdata[:, 0], curdata[:, 1])

                    # keep only valid points -- inside the array
                    used = np.where((i >=0) & (i < array.nbx) & (j >=0) & (j < array.nby))[0]

                    if len(used) == 0:
                        continue

                    i = i[used]
                    j = j[used]
                    z = curdata[used, 2]

                    # create a key array
                    keys = np.vstack((i,j)).T
                    # find unique keys
                    keys = np.unique(keys, axis=0)

                    # create a ijz array
                    ijz = np.vstack((i, j, z)).T

                    # sort ijz array according to keys
                    #
                    # the most important indice is the last one enumerated in lexsort
                    # see : https://numpy.org/doc/stable/reference/generated/numpy.lexsort.html
                    ijz = ijz[np.lexsort((ijz[:,1], ijz[:,0]))]

                    # find first element of each key
                    idx = np.where(np.abs(np.diff(ijz[:,0])) + np.abs(np.diff(ijz[:,1])) != 0)[0]

                    # add last element
                    idx = np.concatenate((idx, [ijz.shape[0]]))

                    assert len(idx) == keys.shape[0], 'Error in filling'

                    logging.info(_('Cells to fill : {}'.format(len(idx))))

                    # apply operator
                    vals = {}
                    start_ii = 0
                    for ii, key in enumerate(keys):
                        end_ii = idx[ii]+1

                        if end_ii - start_ii >= minpoints:
                            vals[(key[0], key[1])] = operator(ijz[start_ii:end_ii,2])

                        start_ii = end_ii

                    if len(vals) > 0:
                        # create a new ijz array
                        newijz = np.asarray([[key[0], key[1], val] for key, val in vals.items()], dtype = np.float32)

                        array.fillin_from_ijz(newijz)

        array.reset_plot()
        self.Paint()

        logging.info(_('Filling done !'))

    def count_active_array_from_laz(self, array:WolfArray = None, used_codes:list = [], chunk_size:float = 500.):
        """ Fill active array with laz data

        :param array: array to fill
        :param used_codes: codes to use
        :param operator: operator to use
        """

        if self.mylazgrid is None:
            logging.info(_('No laz grid - Aborting !'))
            return

        if array is None:
            logging.error(_('No array'))
            return

        if len(used_codes) == 0 :
            keycode = [key for key,val in self.mylazgrid.colors.classification.items()]
            names = [val[0] for key,val in self.mylazgrid.colors.classification.items()]

            with wx.MultiChoiceDialog(None, _('Choose the codes to use'), _('Codes'), names) as dlg:
                if dlg.ShowModal() == wx.ID_OK:
                    data = {}
                    used_codes = dlg.GetSelections()
                    used_codes = [float(keycode[cur]) for cur in used_codes]
                else:
                    return

        bounds = array.get_bounds()

        # align bounds on chunk_size
        bounds[0][0] = bounds[0][0] - bounds[0][0] % chunk_size
        bounds[0][1] = bounds[0][1] + chunk_size - bounds[0][1] % chunk_size
        bounds[1][0] = bounds[1][0] - bounds[1][0] % chunk_size
        bounds[1][1] = bounds[1][1] + chunk_size - bounds[1][1] % chunk_size

        chunks_x = np.arange(bounds[0][0], bounds[0][1], chunk_size)
        chunks_y = np.arange(bounds[1][0], bounds[1][1], chunk_size)

        for curx in tqdm(chunks_x, 'Chunks'):
            for cury in chunks_y:

                curbounds = [[curx, curx + chunk_size], [cury, cury + chunk_size]]

                logging.info(_('Scan {}-{}  {}-{}').format(curbounds[0][0],curbounds[0][1],curbounds[1][0],curbounds[1][1]))
                mylazdata = self.mylazgrid.scan(curbounds)

                if len(mylazdata) == 0:
                    continue

                # Test codes
                data = {}
                for curcode in used_codes:
                    data[curcode] = mylazdata[mylazdata[:, 3] == curcode]

                # Treat data for each code
                for curdata in data.values():

                    if curdata.shape[0] == 0:
                        continue
                    else:
                        logging.info(_('Code {} : {} points'.format(curdata[0,3], curdata.shape[0])))

                    # get i,j from x,y
                    i,j = array.get_ij_from_xy(curdata[:, 0], curdata[:, 1])

                    # keep only valid points -- inside the array
                    used = np.where((i >=0) & (i < array.nbx) & (j >=0) & (j < array.nby))[0]

                    if len(used) == 0:
                        continue

                    i = i[used]
                    j = j[used]
                    z = curdata[used, 2]

                    # create a key array
                    keys = np.vstack((i,j)).T
                    # find unique keys
                    keys = np.unique(keys, axis=0)

                    # create a ijz array
                    ijz = np.vstack((i, j, z)).T

                    # sort ijz array according to keys
                    #
                    # the most important indice is the last one enumerated in lexsort
                    # see : https://numpy.org/doc/stable/reference/generated/numpy.lexsort.html
                    ijz = ijz[np.lexsort((ijz[:,1], ijz[:,0]))]

                    # find first element of each key
                    idx = np.where(np.abs(np.diff(ijz[:,0])) + np.abs(np.diff(ijz[:,1])) != 0)[0]

                    # add last element
                    idx = np.concatenate((idx, [ijz.shape[0]]))

                    assert len(idx) == keys.shape[0], 'Error in filling'

                    logging.info(_('Cells to fill : {}'.format(len(idx))))

                    # apply operator
                    vals = {}
                    start_ii = 0
                    for ii, key in enumerate(keys):
                        end_ii = idx[ii]+1

                        vals[(key[0], key[1])] = end_ii - start_ii

                        start_ii = end_ii

                    if len(vals) > 0:
                        # create a new ijz array
                        newijz = np.asarray([[key[0], key[1], val] for key, val in vals.items()], dtype = np.float32)

                        array.fillin_from_ijz(newijz)

        array.reset_plot()
        self.Paint()

        logging.info(_('Counting done'))

    def init_laz_from_lazlasnpz(self, fn=None):
        """ Read LAZ data stored in one file

        :param fn: filename (extension .laz, .las, .npz)
        """

        if fn is None:
            filternpz = "LAZ (*.laz)|*.laz|LAS (*.las)|*.las|npz (*.npz)|*.npz|all (*.*)|*.*"
            dlg = wx.FileDialog(None, _('Choose a file containing LAS data'), wildcard=filternpz)
            ret = dlg.ShowModal()
            if ret != wx.ID_OK:
                dlg.Destroy()
                return

            fn = dlg.GetPath()
            dlg.Destroy()

        lazobj = Wolf_LAZ_Data()
        lazobj.from_file(fn)

        self.add_object('laz', newobj= lazobj)

        logging.info(_('LAZ data read from file : ')+ fn)
        logging.info(_('Stored in internal variable'))
        logging.info(_('Nb points : {:_}').format(self.active_laz.num_points))

        if self.linked:
            if len(self.linkedList) > 0:
                for curframe in self.linkedList:
                    if not curframe is self:
                        curframe.mylazdata.append(self.active_laz)

    def _choice_laz_classification(self):

        dlg = wx.SingleChoiceDialog(None, _('Choose the classification'), _('Classification'), ['SPW-Geofit 2023', 'SPW 2013-2014'], wx.CHOICEDLG_STYLE)
        ret = dlg.ShowModal()
        if ret != wx.ID_OK:
            dlg.Destroy()
            return None

        classification = dlg.GetStringSelection()
        dlg.Destroy()

        return classification

    def init_laz_from_gridinfos(self, dirlaz:str = None, classification:Literal['SPW-Geofit 2023', 'SPW 2013-2014', 'SPW 2021-2022'] = 'SPW-Geofit 2023'):

        if dirlaz is None:
            dlg = wx.DirDialog(None, _('Choose directory where LAZ data/gridinfo are stored'), defaultPath= str(self.default_laz))
            ret = dlg.ShowModal()
            if ret != wx.ID_OK:
                return

            dirlaz = dlg.GetPath()

        self.mylazgrid = xyz_laz_grids(dirlaz)

        if classification not in ['SPW-Geofit 2023', 'SPW 2013-2014']:

            classification = self._choice_laz_classification()

        if classification is None:
            logging.warning(_('No classification chosen - Abort !'))
            return
        elif classification == 'SPW 2013-2014':
            self.mylazgrid.colors.init_2013()
        elif classification == "SPW 2021-2022":
            self.mylazgrid.colors.init_2021_2022()
        else:
            self.mylazgrid.colors.init_2023()

        if self.linked:
            if len(self.linkedList) > 0:
                for curframe in self.linkedList:
                    curframe.mylazgrid = self.mylazgrid


    def managebanks(self):
        if self.notebookbanks is None:
            self.notebookbanks = PlotNotebook(self)
            self.mypagebanks = self.notebookbanks.add(_("Manager banks interpolator"), "ManagerInterp")

        msg = ''
        if self.active_cs is None:
            msg += _(' The is no cross section. Please activate the desired object !')

        if msg != '':
            dlg = wx.MessageBox(msg, 'Required action')
            return

        if self.active_cs.linked_zones is None:
            msg += _(' The active zones is None. Please link the desired object to the cross sections !\n')
        # if self.active_zone is None:
        #     msg+=_(' The active zone is None. Please activate the desired object !\n')

        if msg != '':
            dlg = wx.MessageBox(msg, 'Required action')
            return

        self.mypagebanks.pointing(self, self.active_cs, self.active_vector)
        self.notebookbanks.Show(True)

    def _set_fn_fnpos_gltf(self):
        """
        Définition du nom de fichier GLTF/GLB à lire pour réaliser la comparaison
        Utilisation d'une fenêtre de dialogue WX

        Cette fonction n'est a priori appelée que depuis set_fn_fnpos_gltf
        """
        dlg = wx.FileDialog(None, _('Choose filename'),
                            wildcard='glb (*.glb)|*.glb|gltf2 (*.gltf)|*.gltf|All (*.*)|*.*', style=wx.FD_OPEN)
        ret = dlg.ShowModal()
        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return

        fn = dlg.GetPath()
        dlg.Destroy()

        dlg = wx.FileDialog(None, _('Choose pos filename'), wildcard='pos (*.pos)|*.pos|All (*.*)|*.*',
                            style=wx.FD_OPEN)
        ret = dlg.ShowModal()
        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return

        fnpos = dlg.GetPath()
        dlg.Destroy()

        if self.link_params is None:
            self.link_params = {}

        self.link_params['gltf file'] = fn
        self.link_params['gltf pos'] = fnpos

        return fn

    def set_fn_fnpos_gltf(self):
        """
        Définition ou récupération du nom de fichier GLTF/GLB à lire pour réaliser la comparaison

        Le nom de fichier est stocké dans la liste des paramètres partagés de façon à ce que l'appel de mise à jour puisse s'effectuer dans n'importe quel frame
        """
        fn = ''
        fnpos = ''
        if self.linked:
            for curgui in self.linkedList:
                if curgui.link_params is not None:
                    if 'gltf file' in curgui.link_params.keys():
                        fn = curgui.link_params['gltf file']
                        fnpos = curgui.link_params['gltf pos']
                        break
        elif self.link_params is None:
            self.link_params = {}
            fn = self._set_fn_fnpos_gltf()

        if fn == '':
            self._set_fn_fnpos_gltf()

    def read_last_result(self):
        """Lecture du dernier résultat pour les modèles ajoutés et plottés"""

        self.currently_readresults = True

        pgbar = wx.ProgressDialog(_('Reading results'), _('Reading results'), maximum=len(self.myres2D), parent=self, style=wx.PD_APP_MODAL | wx.PD_AUTO_HIDE)

        for id, curmodel in enumerate(self.iterator_over_objects(draw_type.RES2D)):
            curmodel: Wolfresults_2D
            logging.info(_('Updating {} - Last result'.format(curmodel.idx)))

            curmodel.read_oneresult()
            curmodel.set_currentview()
            self._update_sim_explorer(curmodel)

            pgbar.Update(id + 1, _('Reading results') + ' - ' + curmodel.idx)

        pgbar.Destroy()

        self.Refresh()
        self.currently_readresults = False
        self._update_mytooltip()

    def read_one_result(self, which:int):
        """
        Lecture d'un résultat spécific pour les modèles ajoutés et plottés

        :param which: result index (0-based) -- -1 for last result
        0 = first result
        """
        self.currently_readresults = True

        for curmodel in self.iterator_over_objects(draw_type.RES2D):
            curmodel: Wolfresults_2D
            if curmodel.checked:
                logging.info(_('Updating {} - Specific result {}'.format(curmodel.idx, which)))

                curmodel.read_oneresult(which)
                curmodel.set_currentview()
                self._update_sim_explorer(curmodel)


        self.Refresh()
        self.currently_readresults = False
        self._update_mytooltip()


    def simul_previous_step(self):
        """
        Mise à jour au pas précédent
        """
        self.currently_readresults = True

        for curmodel in self.iterator_over_objects(draw_type.RES2D):
            curmodel: Wolfresults_2D
            logging.info(_('Updating {} - Previous result'.format(curmodel.idx)))

            curmodel.read_previous()
            curmodel.set_currentview()
            self._update_sim_explorer(curmodel)

        self.Refresh()
        self.currently_readresults = False
        self._update_mytooltip()

    def particle_next_step(self):
        """ Mise à jour au pas suivant """

        for curps in self.iterator_over_objects(draw_type.PARTICLE_SYSTEM):
            curps: Particle_system
            logging.info(_('Updating {} - Next result'.format(curps.idx)))

            curps.next_step()

        self._update_mytooltip()
        self.Refresh()

    def particle_previous_step(self):
        """ Mise à jour au pas précédent """

        for curps in self.iterator_over_objects(draw_type.PARTICLE_SYSTEM):
            curps: Particle_system
            logging.info(_('Updating {} - Next result'.format(curps.idx)))

            curps.previous_step()

        self._update_mytooltip()
        self.Refresh()

    def simul_next_step(self):
        """
        Mise à jour au pas suivant
        """
        self.currently_readresults = True

        for curmodel in self.iterator_over_objects(draw_type.RES2D):
            curmodel: Wolfresults_2D
            logging.info(_('Updating {} - Next result'.format(curmodel.idx)))

            curmodel.read_next()
            curmodel.set_currentview()
            self._update_sim_explorer(curmodel)


        self.Refresh()
        self.currently_readresults = False
        self._update_mytooltip()

    def OnMenuHighlight(self, event:wx.MenuEvent):

        id = event.GetId()
        item:wx.MenuItem
        item = self.menubar.FindItemById(event.GetId())

        if item is not None:
            self.set_statusbar_text(item.GetHelp())

    def _select_laz_source(self):
        """ Select laz source """

        if self.active_laz is None and self.mylazgrid is None:
            logging.warning(_('No LAZ data loaded/initialized !'))
            return None
        elif self.active_laz is None:
            # No active laz data
            laz_source = self.mylazgrid
        elif self.mylazgrid is None:
            # No laz grid
            laz_source = self.active_laz
        else:
            # We have both
            choices = [_('From active LAZ data'), _('From newly extracted data')]

            keys = self.get_list_keys(draw_type.LAZ, None)
            if len(keys) > 1:
                choices.append(_('From multiple LAZ data'))

            dlg = wx.SingleChoiceDialog(None, _("Pick a data source"), "Choices", choices)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return None

            source = dlg.GetStringSelection()
            idx = choices.index(source)
            dlg.Destroy()

            if idx == 0:
                laz_source = self.active_laz
            elif idx == 1:
                laz_source = self.mylazgrid
            else:
                dlg = wx.MultiChoiceDialog(None, _('Choose the LAZ data to use\n\nIf multiple, a new one will be created !'), _('LAZ data'), keys)
                if dlg.ShowModal() == wx.ID_OK:
                    used_keys = dlg.GetSelections()
                    used_keys = [keys[cur] for cur in used_keys]
                    laz_source = Wolf_LAZ_Data()
                    for curkey in used_keys:
                        laz_source.merge(self.get_obj_from_id(curkey, draw_type.LAZ))

                    self.add_object('laz', newobj=laz_source, id = 'Merged LAZ data')
                    dlg.Destroy()
                else:
                    dlg.Destroy()
                    return None

        return laz_source

    def _choice_laz_colormap(self) -> int:

        choices, ass_values = choices_laz_colormap()
        dlg = wx.SingleChoiceDialog(None, _("Pick a colormap"), "Choices", choices)

        if self.active_laz is not None:
            if self.active_laz.associated_color is not None:
                dlg.SetSelection(ass_values.index(self.active_laz.associated_color))

        ret = dlg.ShowModal()
        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return self.active_laz.associated_color

        colormap = dlg.GetStringSelection()
        idx = choices.index(colormap)
        dlg.Destroy()

        return ass_values[idx]

    def newdrowning(self,itemlabel):

        if itemlabel == ('Create a drowning...'):
            new_drowning = Drowning_victim_Viewer(mapviewer = self)
            self.add_object(which='drowning',newobj=new_drowning,ToCheck=True)

        elif itemlabel == ('Add a drowning result...'):
            dialog = wx.DirDialog(None,_("Select the folder containing your drowning"), style=wx.DD_DEFAULT_STYLE)
            if dialog.ShowModal() == wx.ID_OK:
                # Récupérer le chemin sélectionné
                filedir = dialog.GetPath()
                dialog.Destroy()

            if not os.path.exists(os.path.join(filedir, "Results.npz")):
                logging.error(_("The selected folder does not contain any Results.npz"))
                return

            new_drowning = Drowning_victim_Viewer(mapviewer = self,filedir=filedir)
            new_drowning.file_drowning = 1
            self.add_object(which='drowning',newobj=new_drowning,ToCheck=True)
            new_drowning.load_results()
            new_drowning.time_id = len(new_drowning.wanted_time)-2
            new_drowning.init_plot()
            self.menu_drowning()

    def new_dike(self, itemlabel):
        """
        Called when 'Create dike...' or 'Add dike...' are selected in the viewer
        """
        newdike = DikeWolf(mapviewer = self)

        if _('Add dike...') in itemlabel:
            newdike.load_results()

        self.add_object(which='dike', newobj=newdike, ToCheck = True)#, filename=newdike.filename, id=newdike.filename)

        self.menu_dike()
        autoscale = True

    def _run_compare_arrays(self, dlg):
        """ Run the comparison of two arrays"""

        from .ui.wolf_multiselection_collapsiblepane import Wolf_CompareArrays_Selection

        assert isinstance(dlg, Wolf_CompareArrays_Selection), 'Dialog must be a wx.Dialog instance'

        dlg: Wolf_CompareArrays_Selection

        vals = dlg.get_values()
        id1 = vals[_('Reference array')][0]
        id2 = vals[_('Comparison array')][0]
        min_area = dlg.get_min_area()
        threshold = dlg.get_threshold()
        nb_patches = dlg.get_max_patches()

        ref:WolfArray
        comp:WolfArray
        ref = self.get_obj_from_id(id1, draw_type.ARRAYS)
        comp = self.get_obj_from_id(id2, draw_type.ARRAYS)

        if ref is None or comp is None:
            logging.warning(_('You must select two arrays to compare !'))
            return

        assert isinstance(ref, WolfArray), 'Reference object must be a WolfArray instance'
        assert isinstance(comp, WolfArray), 'Comparison object must be a WolfArray instance'

        if not ref.is_like(comp):
            logging.error(_('The two arrays must have the same shape and type !'))
            return

        # if only 2 arrays, we can use the CompareArrays_wx directly
        from .report.compare_arrays import CompareArrays_wx

        try:
            newsheet = CompareArrays_wx(ref, comp,
                                        size=(800, 600),
                                        ignored_patche_area= min_area,
                                        threshold=threshold,
                                        nb_max_patches = nb_patches,)
            newsheet.Show()

            self.add_object('vector', newobj = newsheet.get_zones(), ToCheck = True, id = 'compare_arrays_{}'.format(ref.idx + comp.idx))
        except:
            logging.error(_('Error in comparing arrays\n'))
            dlg.Destroy()

    def _compare_arrays(self):
        """ Compare two arrays """
        arrays = self.get_list_keys(draw_type.ARRAYS, checked_state = None)

        if len(arrays) == 0:
            logging.warning(_('No arrays to compare !'))
            return
        elif len(arrays) == 1:
            logging.warning(_('Only one array to compare - Nothing to do !'))
            return

        from .ui.wolf_multiselection_collapsiblepane import Wolf_CompareArrays_Selection

        dlg = Wolf_CompareArrays_Selection(parent = self,
                                            title = _("Choose the arrays to compare"),
                                            info = _("Select the reference and comparison arrays"),
                                            values_dict = {_('Reference array'): arrays,
                                                            _('Comparison array'): arrays},
                                            callback= self._run_compare_arrays,
                                            destroyOK = True,
                                            styles = [wx.LB_SINGLE, wx.LB_SINGLE]
                                            )
        dlg.ShowModal()

    def OnMenubar(self, event: wx.MenuEvent):
        """
        Gestion des clicks sur le menu quel que soit le niveau

        Idée générale :
            - récupérer le label du menu sur base de l'id de l'event WX passé en argument --> itemlabel
            - tester le label du menu sur base de la chaîne traduite
            - a priori appeler une autre routine spécifique au traitement choisi
            - éviter autant que possible de coder des fonctions directement dans cette routine ce qui la rendrait complexe à lire

        AUTRE POSSIBILITE:
            - mettre en place un dictionnaire avec key==label, value==action qui se contenterait de tester la présence du label dans les clés et d'appeler l'action
            - dans ce dernier cas, il faudrait que les routines possèdent idéalement une interface unique
        """
        id = event.GetId()
        item = self.menubar.FindItemById(event.GetId())

        if item is None:
            return

        itemlabel = item.ItemLabel

        autoscale = False

        if id == wx.ID_OPEN:
            autoscale = True
            filterProject = "proj (*.proj)|*.proj|param (*.param)|*.param|all (*.*)|*.*"
            file = wx.FileDialog(self, "Choose file", wildcard=filterProject)
            ret = file.ShowModal()
            if ret == wx.ID_CANCEL:
                file.Destroy()
                return
            else:
                # récuparétaion du nom de fichier avec chemin d'accès
                filename = file.GetPath()
                file.Destroy()

            old_dir = os.getcwd()
            os.chdir(os.path.dirname(filename))
            self.read_project(filename)
            os.chdir(old_dir)

        elif itemlabel == _('Shortcuts'):
            # show shortcuts in log
            self.print_shortcuts(True)

        elif itemlabel == _('Project .proj'):
            # show shortcuts in log
            self.help_project()

        elif itemlabel == _('Show logs/informations'):
            self.check_logging()

        elif itemlabel == _('Show values'):
            self.check_tooltip()

        elif itemlabel == _('About'):
            #print About Frame
            self.print_About()

        elif itemlabel == _('Check for updates'):
            # check for new version

            self.check_for_updates()

        elif itemlabel == _("Plot integrated Q along active vector..."):
            """ Integrate Q along active vector """

            if self.active_vector is None:
                logging.warning(_('No active vector !'))
                return

            if self.active_vector.closed:
                logging.error(_('The active vector is closed ! - You can only plot Q along a cross section not a polygon !'))
                return

            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            fig = self.new_fig(_('Q along active vector'), 'Q_along_active_vector', show=False, size=(800, 600))
            self.active_res2d.plot_q_wx(self.active_vector, 'border', toshow=True, fig= fig)

        elif itemlabel == _("Plot integrated Q along active zone..."):
            """ Integrate Q along active zone """

            if self.active_zone is None:
                logging.warning(_('No active zone !'))
                return

            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            fig = self.new_fig(_('Q along active zone'), 'Q_along_active_zone', show=False, size=(800, 600))
            self.active_res2d.plot_q_wx(self.active_zone.myvectors, ['border'] * self.active_zone.nbvectors, toshow=True, fig = fig)

        elif itemlabel == _("Export integrated Q along active vector..."):

            if self.active_vector is None:
                logging.warning(_('No active vector !'))
                return

            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            filterArray = ".csv (*.csv)|*.csv|all (*.*)|*.*"
            fdlg = wx.FileDialog(self, "Choose file name : ", wildcard=filterArray,
                                    style=wx.FD_SAVE)
            ret = fdlg.ShowModal()
            hydrographCSVPath = None
            if ret == wx.ID_OK:
                curfil = fdlg.GetFilterIndex()

                hydrographCSVPath = fdlg.GetPath()

            fdlg.Destroy()

            if hydrographCSVPath is not None:
                # Create a progress dialog
                progress_dialog = wx.ProgressDialog(
                    _("Export Progress"),
                    _("Exporting hydrographs..."),
                    maximum=100,
                    parent=self,
                    style= wx.PD_AUTO_HIDE | wx.PD_APP_MODAL | wx.PD_ELAPSED_TIME
                )

                def update_progress(progress):
                    progress_dialog.Update(progress)

                try:
                    # Call the export function, passing the progress callback
                    status = self.active_res2d.export_hydrographs(
                        vect=self.active_vector,
                        filename=hydrographCSVPath,
                        progress_callback=update_progress
                    )
                finally:
                    progress_dialog.Destroy()  # Ensure dialog is destroyed even if an error occurs

                # Inform the user about the result
                if status:
                    wx.MessageBox(_("Hydrographs exported successfully"), _("Export Hydrographs"), wx.OK | wx.ICON_INFORMATION)
                    logging.info(_('Hydrographs exported successfully'))
                else:
                    wx.MessageBox(_("Error exporting hydrographs"), _("Export Hydrographs"), wx.OK | wx.ICON_ERROR)
                    logging.error(_('Error exporting hydrographs'))


        elif itemlabel == _("Export integrated Q along all vectors in active zone..."):

            if self.active_zone is None:
                logging.warning(_('No active zone !'))
                return

            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            filterArray = ".csv (*.csv)|*.csv|all (*.*)|*.*"
            fdlg = wx.FileDialog(self, "Choose file name : ", wildcard=filterArray,
                                    style=wx.FD_SAVE)
            ret = fdlg.ShowModal()
            hydrographCSVPath = None
            if ret == wx.ID_OK:
                curfil = fdlg.GetFilterIndex()

                hydrographCSVPath = fdlg.GetPath()

            fdlg.Destroy()

            if hydrographCSVPath is not None:
                # Create a progress dialog
                progress_dialog = wx.ProgressDialog(
                    _("Export Progress"),
                    _("Exporting hydrographs..."),
                    maximum=100,
                    parent=self,
                    style= wx.PD_AUTO_HIDE | wx.PD_APP_MODAL | wx.PD_ELAPSED_TIME
                )

                def update_progress(progress):
                    progress_dialog.Update(progress)

                try:
                    # Call the export function, passing the progress callback
                    status = self.active_res2d.export_hydrographs(
                        vect=self.active_zone,
                        filename=hydrographCSVPath,
                        progress_callback=update_progress
                    )
                finally:
                    progress_dialog.Destroy()  # Ensure dialog is destroyed even if an error occurs

                # Inform the user about the result
                if status:
                    wx.MessageBox(_("Hydrographs exported successfully"), _("Export Hydrographs"), wx.OK | wx.ICON_INFORMATION)
                    logging.info(_('Hydrographs exported successfully'))
                else:
                    wx.MessageBox(_("Error exporting hydrographs"), _("Export Hydrographs"), wx.OK | wx.ICON_ERROR)
                    logging.error(_('Error exporting hydrographs'))


        elif itemlabel == _("Plot stats unknown (selected nodes)..."):

            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            all_selected = []
            for curblock in self.active_res2d.myblocks.values():
                if curblock.SelectionData.nb > 0:
                    all_selected += curblock.SelectionData.myselection

            if len(all_selected) == 0:
                logging.warning(_('No selected nodes - Nothing to do !'))
                return

            keys = Extractable_results.get_list()
            dlg = wx.SingleChoiceDialog(None, _('Choose the unknown/variable to plot'), _('Unknown'), keys)

            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                logging.info(_('No unknown chosen - Aborting !'))
                dlg.Destroy()
                return

            which = Extractable_results.get_from_key(dlg.GetStringSelection())
            dlg.Destroy()

            try:
                choice_bes = Select_Begin_end_interval_step(self, _('Choose the interval and step'), self.active_res2d, checkbox=True)
                ret = choice_bes.ShowModal()

                begin = choice_bes.begin
                end = choice_bes.end
                interval = choice_bes.step

            finally:
                choice_bes.Destroy()

            if begin == -1:
                logging.info(_('No interval chosen - Aborting !'))
                return

            newfig = self.new_fig(_('Series of {} - {} (nodes)').format(which.value[0], self.active_res2d.idx), 'series_'+self.active_res2d.idx, show=False, size= (800,600))

            figax = self.active_res2d.plot_some_values(all_selected, which, toshow=False, figax=newfig, for_steps= (begin-1, end-1, interval))

            newfig.Show()

        elif itemlabel == _("Export stats unknown (selected nodes)..."):

            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            all_selected = []
            for curblock in self.active_res2d.myblocks.values():
                if curblock.SelectionData.nb > 0:
                    all_selected += curblock.SelectionData.myselection

            if len(all_selected) == 0:
                logging.warning(_('No selected nodes - Nothing to do !'))
                return

            keys = Extractable_results.get_list()
            dlg = wx.SingleChoiceDialog(None, _('Choose the unknown/variable to plot'), _('Unknown'), keys)

            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                logging.info(_('No unknown chosen - Aborting !'))
                dlg.Destroy()
                return

            which = Extractable_results.get_from_key(dlg.GetStringSelection())
            dlg.Destroy()

            dlg = wx.FileDialog(None, _('Choose the file to export'), wildcard='csv (*.csv)|*.csv', style=wx.FD_SAVE | wx.FD_OVERWRITE_PROMPT)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                logging.info(_('No file chosen - Aborting !'))
                dlg.Destroy()
                return

            filename = Path(dlg.GetPath())
            dlg.Destroy()


            try:
                choice_bes = Select_Begin_end_interval_step(self, _('Choose the interval and step'), self.active_res2d, checkbox=True)
                ret = choice_bes.ShowModal()

                begin = choice_bes.begin
                end = choice_bes.end
                interval = choice_bes.step
                all = choice_bes.check_all

            finally:
                choice_bes.Destroy()

            if begin == -1:
                logging.info(_('No interval chosen - Aborting !'))
                return

            ret = self.active_res2d.export_some_values_to_csv(all_selected, which, filename, for_steps= (begin-1, end-1, interval), all_values=all)

            if not ret:
                logging.error(_('Error in exporting values !'))

        elif itemlabel == _("Plot stats unknown (inside active vector)..."):

            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            if self.active_vector is None:
                logging.warning(_('No active vector !'))
                return

            keys = Extractable_results.get_list()
            dlg = wx.SingleChoiceDialog(None, _('Choose the unknown/variable to plot'), _('Unknown'), keys)

            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                logging.info(_('No unknown chosen - Aborting !'))
                dlg.Destroy()
                return

            which = Extractable_results.get_from_key(dlg.GetStringSelection())
            dlg.Destroy()

            try:
                choice_bes = Select_Begin_end_interval_step(self, _('Choose the interval and step'), self.active_res2d, checkbox=True)
                ret = choice_bes.ShowModal()

                begin = choice_bes.begin
                end = choice_bes.end
                interval = choice_bes.step
                violin = choice_bes.check_violin

            finally:
                choice_bes.Destroy()

            if begin == -1:
                logging.info(_('No interval chosen - Aborting !'))
                return

            newfig = self.new_fig(_('Series of {} - {} (polygon)').format(which.value[0], self.active_res2d.idx), 'series_'+self.active_res2d.idx, show=False, size= (800,600))

            if violin:
                figax = self.active_res2d.plot_violin_values(self.active_vector, which, toshow=False, figax=newfig, for_steps= (begin-1, end-1, interval))
            else:
                figax = self.active_res2d.plot_some_values(self.active_vector, which, toshow=False, figax=newfig, for_steps= (begin-1, end-1, interval))

            newfig.Show()

        elif itemlabel == _("Export stats unknown (inside active vector)..."):

            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            if self.active_vector is None:
                logging.warning(_('No active vector !'))
                return

            keys = Extractable_results.get_list()
            dlg = wx.SingleChoiceDialog(None, _('Choose the unknown/variable to plot'), _('Unknown'), keys)

            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                logging.info(_('No unknown chosen - Aborting !'))
                dlg.Destroy()
                return

            which = Extractable_results.get_from_key(dlg.GetStringSelection())
            dlg.Destroy()

            dlg = wx.FileDialog(None, _('Choose the file to export'), wildcard='csv (*.csv)|*.csv', style=wx.FD_SAVE | wx.FD_OVERWRITE_PROMPT)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                logging.info(_('No file chosen - Aborting !'))
                dlg.Destroy()
                return

            filename = Path(dlg.GetPath())
            dlg.Destroy()

            try:
                choice_bes = Select_Begin_end_interval_step(self, _('Choose the interval and step'), self.active_res2d, checkbox=True)
                ret = choice_bes.ShowModal()

                begin = choice_bes.begin
                end = choice_bes.end
                interval = choice_bes.step

            finally:
                choice_bes.Destroy()

            if begin == -1:
                logging.info(_('No interval chosen - Aborting !'))
                return

            ret = self.active_res2d.export_some_values_to_csv(self.active_vector, which, filename=filename, for_steps= (begin-1, end-1, interval), all_values=all)

        elif itemlabel == _("Plot stats unknown (inside active zone)..."):

            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            if self.active_zone is None:
                logging.warning(_('No active zone !'))
                return

            keys = Extractable_results.get_list()
            dlg = wx.SingleChoiceDialog(None, _('Choose the unknown/variable to plot'), _('Unknown'), keys)

            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                logging.info(_('No unknown chosen - Aborting !'))
                dlg.Destroy()
                return

            which = Extractable_results.get_from_key(dlg.GetStringSelection())
            dlg.Destroy()

            try:
                choice_bes = Select_Begin_end_interval_step(self, _('Choose the interval and step'), self.active_res2d, checkbox=True)
                ret = choice_bes.ShowModal()

                begin = choice_bes.begin
                end = choice_bes.end
                interval = choice_bes.step
                violin = choice_bes.check_violin

            finally:
                choice_bes.Destroy()

            if begin == -1:
                logging.info(_('No interval chosen - Aborting !'))
                return

            for idx, curvect in enumerate(self.active_zone.myvectors):
                logging.info(_('Plotting {} / {}'.format(idx, self.active_zone.nbvectors)))
                if idx ==0:
                    newfig = self.new_fig(_('Series of {} - {} (zone)').format(which.value[0], self.active_res2d.idx), 'series_'+self.active_res2d.idx, show=False, size= (800,600))

                if violin:
                    figax = self.active_res2d.plot_violin_values(curvect, which, toshow=False, figax=newfig, for_steps= (begin-1, end-1, interval))
                else:
                    figax = self.active_res2d.plot_some_values(curvect, which, toshow=False, figax=newfig, for_steps= (begin-1, end-1, interval))

            newfig.Show()

        elif itemlabel == _("Export stats unknown (inside active zone)..."):

            if self.active_res2d is None:
                logging.warning(_('No active 2D result !'))
                return

            if self.active_zone is None:
                logging.warning(_('No active zone !'))
                return

            keys = Extractable_results.get_list()
            dlg = wx.SingleChoiceDialog(None, _('Choose the unknown/variable to plot'), _('Unknown'), keys)

            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                logging.info(_('No unknown chosen - Aborting !'))
                dlg.Destroy()
                return

            which = Extractable_results.get_from_key(dlg.GetStringSelection())
            dlg.Destroy()

            dlg = wx.DirDialog(None, _('Choose the directory where to export'), style= wx.DD_DEFAULT_STYLE)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                logging.info(_('No file chosen - Aborting !'))
                dlg.Destroy()
                return

            directory = Path(dlg.GetPath())
            dlg.Destroy()

            try:
                choice_bes = Select_Begin_end_interval_step(self, _('Choose the interval and step'), self.active_res2d, checkbox=True)
                ret = choice_bes.ShowModal()

                begin = choice_bes.begin
                end = choice_bes.end
                interval = choice_bes.step

            finally:
                choice_bes.Destroy()

            if begin == -1:
                logging.info(_('No interval chosen - Aborting !'))
                return

            allnames = [curvect.myname for curvect in self.active_zone.myvectors]
            if len(set(allnames)) != len(allnames):
                logging.warning(_('Some vectors have the same name !'))

                # create new unique names
                unique_name = []
                for curvect in self.active_zone.myvectors:
                    if curvect.myname in unique_name:
                        unique_name.append(curvect.myname + '_' + str(unique_name.count(curvect.myname)))
                    else:
                        unique_name.append(curvect.myname)


            for idx, (curvect, name) in enumerate(zip(self.active_zone.myvectors, unique_name)):
                self.active_res2d.export_some_values_to_csv(curvect, which, filename=directory / name, for_steps= (begin-1, end-1, interval), all_values=all)

        elif itemlabel == _("Plot active vector..."):
            """ Plot data along active vector """

            if self.active_vector is None:
                logging.warning(_('No active vector !'))
                return

            add_cloud = False
            if self.active_cloud is not None:
                dlg = wx.MessageDialog(self, _('Do you want to plot the cloud ?'), style=wx.YES_NO)

                if dlg.ShowModal() == wx.ID_YES:
                    add_cloud = True

                    prox = wx.TextEntryDialog(None,_('Proximity [m] ?'), value = '5.0')
                    ret = prox.ShowModal()
                    if ret == wx.ID_CANCEL:
                        prox.Destroy()
                        return
                    try:
                        proxval = float(prox.GetValue())
                    except:
                        prox.Destroy()
                        logging.warning(_('Bad value -- Rety'))
                        return

                    tol = wx.TextEntryDialog(None,_('Tolerance [m] ?'), value = '0.5')
                    ret = tol.ShowModal()
                    if ret == wx.ID_CANCEL:
                        tol.Destroy()
                        return
                    try:
                        tolval = float(tol.GetValue())
                    except:
                        tol.Destroy()
                        logging.warning(_('Bad value -- Rety'))
                        return

                else:
                    add_cloud = False

                dlg.Destroy()

            # Création d'un graphique matplotlib sous wx
            lab  = _('Plot of active vector') + ' - ' + self.active_vector.myname
            figmpl = self.new_fig(lab, lab, show=False, size= (800,600))

            linkedarrays = self.get_linked_arrays()

            with wx.MultiChoiceDialog(None, _('Choose the arrays to plot'), _('Arrays'), [curarray for curarray in list(linkedarrays.keys())]) as dlg:
                dlg:wx.MultiChoiceDialog
                dlg.SetSelections(range(len(linkedarrays)))

                if dlg.ShowModal() == wx.ID_CANCEL:
                    dlg.Destroy()
                    return

                selected = dlg.GetSelections()
                keys = list(linkedarrays.keys())
                selected = [keys[cur] for cur in selected]
                dlg.Destroy()

            linkedarrays = {curkey:curval for curkey, curval in linkedarrays.items() if curkey in selected}

            self.active_vector.plot_linked_wx(figmpl, linkedarrays)

            if add_cloud:
                s, z = self.active_cloud.projectontrace(self.active_vector, return_cloud=False, proximity= proxval)

                figmpl.plot( s, z, c='black', s=1.0, marker='x')

                for curs, curz in zip(s,z):
                    figmpl.plot([curs, curs], [curz-tolval, curz+tolval], 'k--', linewidth=0.5)
                    figmpl.plot([curs-.1, curs+.1], [curz+tolval, curz+tolval], c='black', linewidth=0.5)
                    figmpl.plot([curs-.1, curs+.1], [curz-tolval, curz-tolval], c='black', linewidth=0.5)

            figmpl.Show()

        elif itemlabel == _("Compute and apply unique colormap on all..."):
            self.uniquecolormap()

        elif itemlabel == _("Load and apply unique colormap on all..."):
            self.uniquecolormap(True)

        elif itemlabel == _("Force uniform in parts on all..."):
            self.uniforminparts_all(True)

        elif itemlabel == _("Force linear interpolation on all..."):
            self.uniforminparts_all(False)

        elif itemlabel == _("Load and apply mask (nap)..."):
            self.loadnap_and_apply()

        elif itemlabel == _("Active simulation..."):
            if self.active_res2d is None:
                logging.warning(_('No active simulation !'))
                return

            from .report.simplesimgpu import SimpleSimGPU_Report_wx

            if isinstance(self.active_res2d, wolfres2DGPU):
                newsheet = SimpleSimGPU_Report_wx(Path(self.active_res2d.filename).parent, size=(800, 600))
                newsheet.Show()
            else:
                logging.warning(_('Active simulation is not a GPU simulation - Not yet implemented for CPU simulations !'))

        elif itemlabel == _("All checked simulations..."):

            sims = self.get_list_keys(draw_type.RES2D, checked_state=True)
            if len(sims) == 0:
                logging.warning(_('No checked simulation !'))
                return

            from .report.simplesimgpu import SimpleSimGPU_Report_wx

            for curkey in sims:
                curmodel = self.get_obj_from_id(curkey, draw_type.RES2D)
                if isinstance(curmodel, wolfres2DGPU):
                    newsheet = SimpleSimGPU_Report_wx(Path(curmodel.filename).parent, size=(800, 600))
                    newsheet.Show()
                else:
                    logging.warning(_('Simulation {} is not a GPU simulation - Not yet implemented for CPU simulations !').format(curmodel.idx))

        elif itemlabel == _("One simulation from disk..."):
            dlg = wx.DirDialog(None, _('Choose the directory containing the simulation'), style=wx.DD_DEFAULT_STYLE)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return
            directory = Path(dlg.GetPath())
            dlg.Destroy()
            if not directory.exists():
                logging.error(_('Directory {} does not exist !').format(directory))
                wx.MessageBox(_('Directory {} does not exist !').format(directory), _('Error'), wx.OK | wx.ICON_ERROR)
                return
            if not directory.is_dir():
                logging.error(_('Path {} is not a directory !').format(directory))
                wx.MessageBox(_('Path {} is not a directory !').format(directory), _('Error'), wx.OK | wx.ICON_ERROR)
                return

            from .report.simplesimgpu import SimpleSimGPU_Report_wx

            # check if we want to show all wx reports
            newsheet = SimpleSimGPU_Report_wx(directory, size=(800, 600), show=True)


        elif itemlabel == _("All simulations in directory..."):
            dlg = wx.DirDialog(None, _('Choose the directory containing the simulations'), style=wx.DD_DEFAULT_STYLE)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return
            directory = Path(dlg.GetPath())
            dlg.Destroy()
            if not directory.exists():
                logging.error(_('Directory {} does not exist !').format(directory))
                wx.MessageBox(_('Directory {} does not exist !').format(directory), _('Error'), wx.OK | wx.ICON_ERROR)
                return
            if not directory.is_dir():
                logging.error(_('Path {} is not a directory !').format(directory))
                wx.MessageBox(_('Path {} is not a directory !').format(directory), _('Error'), wx.OK | wx.ICON_ERROR)
                return
            from .report.simplesimgpu import SimpleSimGPU_Reports_wx

            # check if we want to show all wx reports
            dlg = wx.MessageDialog(None, _('Do you want to show all reports ?'), _('Show all reports'), style=wx.YES_NO | wx.YES_DEFAULT)
            ret = dlg.ShowModal()
            dlg.Destroy()

            newsheets = SimpleSimGPU_Reports_wx(directory, show = ret == wx.ID_YES, size=(800, 600))

        elif itemlabel == _("Compare arrays..."):

            self._compare_arrays()

        elif itemlabel == _("Compare arrays from files..."):

            from .report.compare_arrays import CompareArrays_wx

            dlg = wx.FileDialog(None, _('Choose the reference file'), wildcard='*.tif, *.bin, *.npy|*.tif;*.bin;*.npy|all (*.*)|*.*', style=wx.FD_OPEN | wx.FD_FILE_MUST_EXIST)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return
            ref_filename = dlg.GetPath()
            dlg.Destroy()

            dlg = wx.FileDialog(None, _('Choose the comparison file'), wildcard='*.tif, *.bin, *.npy|*.tif;*.bin;*.npy|all (*.*)|*.*', style=wx.FD_OPEN | wx.FD_FILE_MUST_EXIST)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return
            comp_filename = dlg.GetPath()
            dlg.Destroy()

            try:
                wa_ref = WolfArray(ref_filename)
                wa_comp = WolfArray(comp_filename)

                if not (wa_ref.loaded and wa_comp.loaded):
                    logging.error(_('Error in loading arrays from files'))
                    wx.MessageBox(_('Error in loading arrays from files'), _('Error'), wx.OK | wx.ICON_ERROR)
                    return

                if wa_ref.is_like(wa_comp):
                    newsheet = CompareArrays_wx(wa_ref, wa_comp, size=(800, 600))
                    newsheet.Show()
                else:
                    logging.error(_('The two arrays are not compatible - Cannot compare !'))
                    wx.MessageBox(_('The two arrays are not compatible - Cannot compare !'), _('Error'), wx.OK | wx.ICON_ERROR)

                logging.info(_('Arrays {} and {} compared successfully').format(ref_filename, comp_filename))
            except Exception as e:
                logging.error(_('Error in comparing arrays from files\n{}'.format(e)))

        elif itemlabel == _("Compare checked simulations..."):

            sims = self.get_list_keys(draw_type.RES2D, checked_state=True)
            if len(sims) == 0:
                logging.warning(_('No checked simulation !'))
                return

            elif len(sims) == 1:
                logging.warning(_('Only one checked simulation - Nothing to compare !'))
                return

            from .report.simplesimgpu import SimpleSimGPU_Report_Compare_wx

            sims = [self.get_obj_from_id(curkey, draw_type.RES2D) for curkey in sims]
            # conserve only GPU simulations
            sims = [Path(curmodel.filename) for curmodel in sims if isinstance(curmodel, wolfres2DGPU)]

            # take parent if "simul_gpu_results" in sims
            sims = [sim.parent for sim in sims if 'simul_gpu_results' in str(sim)]

            if len(sims) == 0:
                logging.warning(_('No GPU simulation to compare !'))
                return
            elif len(sims) == 1:
                logging.warning(_('Only one GPU simulation - Nothing to compare !'))
                return
            elif len(sims) > 1:
                try:
                    newsheet = SimpleSimGPU_Report_Compare_wx(sims, size=(800, 600))
                    newsheet.Show()
                except Exception as e:
                    logging.error(_('Error in comparing simulations\n{}'.format(e)))

        elif itemlabel == _("Compare all simulations in a directory..."):

            dlg = wx.DirDialog(None, _('Choose the directory containing the simulations'), style=wx.DD_DEFAULT_STYLE)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            directory = Path(dlg.GetPath())
            dlg.Destroy()

            if not directory.exists():
                logging.error(_('Directory {} does not exist !').format(directory))
                wx.MessageBox(_('Directory {} does not exist !').format(directory), _('Error'), wx.OK | wx.ICON_ERROR)
                return
            if not directory.is_dir():
                logging.error(_('Path {} is not a directory !').format(directory))
                wx.MessageBox(_('Path {} is not a directory !').format(directory), _('Error'), wx.OK | wx.ICON_ERROR)
                return
            from .report.simplesimgpu import SimpleSimGPU_Report_Compare_wx
            try:
                newsheet = SimpleSimGPU_Report_Compare_wx(directory, size= (800,600))
                newsheet.Show()
            except Exception as e:
                logging.error(_('Error in comparing simulations in directory\n{}'.format(e)))

        elif itemlabel == _("Inpaint active array..."):

            if self.active_array is None:
                logging.warning(_('No active array !'))
                return

            nb = self.active_array.count_holes()
            if nb == 0:
                logging.warning(_('No hole in the array !'))
                return

            dlg = wx.SingleChoiceDialog(None, _('Ignore the last ones ?'), _('Holes'), [str(i) for i in range(10)], style=wx.CHOICEDLG_STYLE)
            if dlg.ShowModal() == wx.ID_CANCEL:
                dlg.Destroy()
                return

            nb = dlg.GetSelection()
            dlg.Destroy()

            self.active_array.inpaint(ignore_last=nb)

        elif itemlabel == _("Inpaint waterlevel..."):

            dlg = InPaint_waterlevel(None, _("Choose the array to inpaint"), (400,400), self)

        elif itemlabel == _("Inpaint array with mask..."):

            dlg = InPaint_array(None, _("Choose the array to inpaint"), (400,400), self)

        elif itemlabel == _("Filter inundation arrays..."):
            self.filter_inundation()

        elif itemlabel == _("Plot active polygons..."):

            if self.active_zone is None:
                logging.warning(_('No active zone ! -- please select a zone containing polygons !'))
                return

            try:
                # plotzone:list[zone]
                plotzone = []
                zonename = self.active_zone.myname
                if '_left_' in zonename or '_right_' in zonename:

                    logging.info(_('Left and Right polygons are detected'))

                    testname = zonename.replace('_left_', '')
                    testname = testname.replace('_right_', '')

                    for curzone in self.active_zones.myzones:
                        if testname == curzone.myname.replace('_left_', '').replace('_right_', ''):
                            plotzone.append(curzone)

                    msg = wx.MessageDialog(self,
                                        _('Left and Right polygons are detected \nDo you want like to plot left and right polygons on the same plot ?'),
                                        style=wx.YES_NO | wx.YES_DEFAULT)
                    ret = msg.ShowModal()
                    msg.Destroy()
                    if ret == wx.ID_NO:
                        plotzone = [self.active_zone]
                else:
                    logging.info(_('Sole polygon detected'))
                    plotzone = [self.active_zone]

                # Création d'un graphique matplotlib sous wx
                figmpl = self.new_fig(_('Plot of active polygons'), 'plot_active_polygons', show=False, size= (800,600))

                linkedarrays = {}

                # Matrices 2D
                for curarray in self.iterator_over_objects(draw_type.ARRAYS):
                    curarray: WolfArray
                    logging.info(_('Plotting array {}').format(curarray.idx))
                    linkedarrays[curarray.idx] = curarray

                # Résultats 2D
                for curarray in self.iterator_over_objects(draw_type.RES2D):
                    curarray: Wolfresults_2D
                    logging.info(_('Plotting results {}').format(curarray.idx))
                    linkedarrays[curarray.idx] = curarray

                linkedvecs={}
                for curvect in self.iterator_over_objects(draw_type.VECTORS):
                    curvect: Zones
                    logging.info(_('Plotting vector {}').format(curvect.idx))
                    linkedvecs[curvect.idx] = curvect

                if len(plotzone) > 1:
                    # left and right polygons
                    for curzone in plotzone:
                        if '_left_' in curzone.myname:
                            locarrays = {}
                            for curkey, curarray in linkedarrays.items():
                                locarrays[curkey+ '_left'] = curarray

                            curzone.plot_linked_polygons_wx(figmpl, locarrays, linked_vec=linkedvecs, linestyle= '--')
                        elif '_right_' in curzone.myname:
                            locarrays = {}
                            for curkey, curarray in linkedarrays.items():
                                locarrays[curkey+ '_right'] = curarray

                            curzone.plot_linked_polygons_wx(figmpl, locarrays, linked_vec=linkedvecs, linestyle= '-.')
                else:
                    # sole polygon
                    plotzone[0].plot_linked_polygons_wx(figmpl, linkedarrays, linked_vec=linkedvecs)

                figmpl.Show()

            except Exception as e:
                logging.error(_('Error in plotting active polygons\n{}'.format(e)))
                logging.warning(_('Are you sure the active zone contains polygons ?'))

        elif itemlabel == _("Manage banks..."):

            if self.active_vector is None:
                msg = _('Active vector is None\nPlease activate the one desired')
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            self.managebanks()

        elif itemlabel == _("Create banks from vertices..."):

            self.active_cs.create_zone_from_banksbed()
            self.active_cs.linked_zones.showstructure()

        elif itemlabel == _("Link cross sections to active zones"):

            if self.active_cs is None:
                msg = _('Active cross sections is None\nPlease activate the one desired')
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            if self.active_zones is None:
                msg = _('Active zone is None\nPlease activate the one desired')
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            self.active_cs.link_external_zones(self.active_zones)

        elif itemlabel == _("Rename cross sections..."):

            dlg = wx.TextEntryDialog(None, _('Which starting point?'))
            ret = dlg.ShowModal()

            idxstart = dlg.GetValue()

            self.active_cs.rename(int(idxstart))

        elif itemlabel == _("Triangulate cross sections..."):
            self.triangulate_cs()

        # elif itemlabel == _("Import triangulation..."):
        #     self.import_3dfaces()

        elif itemlabel == _("Interpolate on active triangulation..."):
            self.interpolate_triangulation(keep='all')

        elif itemlabel == _("Interpolate on active triangulation (keep only above)..."):
            self.interpolate_triangulation(keep='above')

        elif itemlabel == _("Interpolate on active triangulation (keep only below)..."):
            self.interpolate_triangulation(keep='below')

        elif itemlabel==_("Compare cloud to array..."):
            self.compare_cloud2array()

        elif itemlabel==_("Split cloud..."):

            if self.active_cloud is None:
                logging.warning(_('No active cloud !'))
                return

            if self.active_vector is None:
                logging.warning(_('No active vector !'))
                return

            self.split_cloud_by_vector()

        elif itemlabel==_("Compare triangles to array..."):
            self.compare_tri2array()

        elif itemlabel ==  _("Move triangles..."):
            self.move_triangles()

        elif itemlabel == _("Rotate triangles..."):
            self.rotate_triangles()

        elif itemlabel ==  _("Create contour from checked arrays..."):

            # Create contour from checked arrays and add it to the list of objects
            newzones = self.create_Zones_from_arrays(self.get_list_objects(draw_type.ARRAYS, checked_state=True))
            self.add_object('vector', newobj=newzones, ToCheck=True, id='Contours from arrays')

        elif itemlabel == _("Calculator..."):

            if self.calculator is None:
                self.calculator = Calculator(mapviewer = self)
            else:
                try:
                    self.calculator.Show()
                except:
                    self.calculator = Calculator(mapviewer = self)

        elif itemlabel == _('Image digitizer...'):

            new_digitizer = Digitizer()

        elif itemlabel == _("Memory views..."):

            if self.memory_views is None:
                self.memory_views = Memory_Views()
                self._memory_views_gui = Memory_Views_GUI(self, _('Memory view manager'), self.memory_views, mapviewer = self)
            else:
                if self._memory_views_gui is None:
                    self._memory_views_gui = Memory_Views_GUI(self, _('Memory view manager'), self.memory_views, mapviewer = self)

                self._memory_views_gui.Show()

        elif itemlabel == _("Memory distances..."):

            if self._distances is not None:
                if self._distances[-1].nbvectors == 0:
                    logging.warning(_('No vector to show !'))
                    return

                self._distances.showstructure(self)

        elif itemlabel == _("Add distances to viewer..."):

            if self._distances is not None:
                self.add_object('vector', newobj=self._distances, ToCheck=True, id='Distances',)

        elif itemlabel == _("Create bridge and export gltf..."):

            if self.active_cs is None:
                msg = _('Active cross sections is None\nPlease activate the one desired')
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            self.start_action('bridge gltf', _('Create bridge and export gltf...'))

        elif itemlabel == _("Export cross sections to gltf..."):

            if self.active_cs is None:
                msg = _('Active cross sections is None\nPlease activate the one desired')
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            dlg = wx.TextEntryDialog(self, 'Z minimum ?', 'Choose an elevation as base')
            dlg.SetValue('')

            zmin = 0.
            if dlg.ShowModal() == wx.ID_OK:
                zmin = float(dlg.GetValue())
            dlg.Destroy()

            dlg = wx.FileDialog(None, _('Choose filename'),
                                wildcard='glb (*.glb)|*.glb|gltf2 (*.gltf)|*.gltf|All (*.*)|*.*', style=wx.FD_SAVE)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            fn = dlg.GetPath()
            dlg.Destroy()

            self.active_cs.export_gltf(zmin, fn)

        elif itemlabel == _("New cloud Viewer..."):

            if self.myinterp is not None:
                self.myinterp.viewer_interpolator()

        elif itemlabel == _("Interpolate on active array..."):

            if self.myinterp is not None:
                self.interpolate_cs()

        elif itemlabel == _("Interpolate active cloud on active array..."):
            self.interpolate_cloud()

        elif itemlabel == _('Save project as...'):

            filterProject = "proj (*.proj)|*.proj|param (*.param)|*.param|all (*.*)|*.*"
            file = wx.FileDialog(self, "Name your file", wildcard=filterProject, style=wx.FD_SAVE)
            if file.ShowModal() == wx.ID_CANCEL:
                file.Destroy()
                return
            else:
                # récuparétaion du nom de fichier avec chemin d'accès
                filename = file.GetPath()
                file.Destroy()

            abspath = True
            dlg = wx.MessageDialog(None, _('Do you want to save the paths in absolute mode ?'), _('Relative paths'), style=wx.YES_NO)
            ret = dlg.ShowModal()
            if ret == wx.ID_NO:
                abspath = False

            self.save_project(filename, absolute= abspath)

        elif itemlabel == _('Initialize from laz, las or npz'):
            self.init_laz_from_lazlasnpz()

        elif itemlabel == _('Initialize from directory'):
            self.init_laz_from_gridinfos()

        elif itemlabel == _('Copy from current zoom'):

            if self.mylazgrid is None:
                logging.warning(_('No gridded LAZ data loaded !'))
                return

            dlg = wx.DirDialog(None, _('Choose a directory to copy the files'), _('Copy files'), style=wx.DD_DEFAULT_STYLE)
            if dlg.ShowModal() == wx.ID_CANCEL:
                dlg.Destroy()

            dirout = dlg.GetPath()
            dlg.Destroy()

            curbounds = [[self.xmin, self.xmin + self.width], [self.ymin, self.ymin + self.height]]

            self.mylazgrid.copy_files_in_bounds(curbounds, dirout)

        elif itemlabel == _('Create cloud points from bridges'):

            if self.active_laz is None:
                self.init_laz_from_lazlasnpz()

            mybridges = self.active_laz.get_data_class(10)
            mycloud = cloud_vertices()

            mycloud.init_from_nparray(mybridges)
            mycloud.myprop.style = 2
            mycloud.myprop.color = getIfromRGB([255, 102, 102])
            mycloud.myprop.width = .5

            if self.linked:
                if len(self.linkedList) > 0:
                    for curframe in self.linkedList:
                        curframe.add_object('cloud', newobj=mycloud, ToCheck=True, id='Bridges')
            else:
                self.add_object('cloud', newobj=mycloud, ToCheck=True, id='Bridges')

        elif itemlabel == _('Create cloud points from buildings'):

            if self.active_laz is None:
                self.init_laz_from_lazlasnpz()

            mybuildings = self.active_laz.get_data_class(1)
            mycloud = cloud_vertices()

            mycloud.init_from_nparray(mybuildings)
            mycloud.myprop.style = 2
            mycloud.myprop.color = getIfromRGB([102, 102, 102])
            mycloud.myprop.width = .5
            if self.linked:
                if len(self.linkedList) > 0:
                    for curframe in self.linkedList:
                        curframe.add_object('cloud', newobj=mycloud, ToCheck=True, id='Buildings')
            else:
                self.add_object('cloud', newobj=mycloud, ToCheck=True, id='Buildings')

        elif itemlabel == _('Create cloud points from specified classes'):

            if self.active_laz is None:
                self.init_laz_from_lazlasnpz()

            codes = self.active_laz.codes_unique

            dlg = wx.MultiChoiceDialog(None, _('Choose the classes to plot'), _('Classes'), codes)
            if dlg.ShowModal() == wx.ID_CANCEL:
                dlg.Destroy()

            selected = dlg.GetSelections()
            selected = [codes[cur] for cur in selected]
            dlg.Destroy()

            for curcode in selected:
                mycloud = cloud_vertices()
                mydata = self.active_laz.get_data_class(curcode)
                mycloud.init_from_nparray(mydata)
                mycloud.myprop.style = 2
                mycloud.myprop.color = getIfromRGB([102, 102, 102])
                mycloud.myprop.width = .5
                if self.linked:
                    if len(self.linkedList) > 0:
                        for curframe in self.linkedList:
                            curframe.add_object('cloud', newobj=mycloud, ToCheck=True, id='Class {}'.format(curcode))
                else:
                    self.add_object('cloud', newobj=mycloud, ToCheck=True, id='Class {}'.format(curcode))

        elif itemlabel == _('Create LAZ viewer'):

            laz_source = self._select_laz_source()
            if laz_source is None:
                logging.warning(_('No LAZ data loaded !'))
                return

            if laz_source is self.mylazgrid:
                if self.mylazgrid is None:
                    logging.warning(_('No gridded LAZ data loaded !'))
                    return
                self.clip_laz_gridded()

                if self.active_laz.nb_points ==0:
                    logging.warning(_('No points in the active LAZ object -- Aborting !'))
                    return

                self.active_laz.create_viewer(self._choice_laz_colormap(), self.mylazgrid.colors)
                self.myviewerslaz.append(self.active_laz.viewer)
                self.active_viewerlaz = self.myviewerslaz[-1]

                # self.myviewer = myviewer(self.active_laz.data, ass_values[idx], palette_classif = self.mylazgrid.colors)
            else:
                if self.active_laz.nb_points ==0:
                    logging.warning(_('No points in the active LAZ object -- Aborting !'))
                    return

                self.active_laz.create_viewer()
                self.myviewerslaz.append(self.active_laz.viewer)
                self.active_viewerlaz = self.myviewerslaz[-1]

                # self.myviewer = myviewer(laz_source.data, ass_values[idx], palette_classif= laz_source.classification)

        elif itemlabel == _('Filter data based on codes'):

            self.filter_active_laz()

        elif itemlabel == _('Descimate LAZ data'):

            if self.active_laz is None:
                return

            # Choose a decimation factor - integer
            dlg = wx.NumberEntryDialog(None, _('Your dataset contains {} points.\nWould you like to descimate?').format(self.active_laz.num_points),
                                        _('Decaimate factor'), _('Decimation'), 0, 0, 100)

            ret = dlg.ShowModal()

            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            descimate_fact = dlg.GetValue()
            dlg.Destroy()

            if descimate_fact > 0:
                self.active_laz.descimate(descimate_fact)
                logging.info(_('New count : {}').format(self.active_laz.num_points))

        elif itemlabel == _('Clip LAZ grid on current zoom'):

            if self.mylazgrid is None:
                logging.warning(_('No gridded LAZ data loaded !'))
                return

            self.clip_laz_gridded()

            if self.active_laz is None:
                logging.error(_('No data found'))
                return

            if self.active_laz.data.shape[0] > 100_000_000:

                # Choose a decimation factor - integer
                dlg = wx.NumberEntryDialog(None, _('Your data selection is very large (>100 M)\nWould you like to descimate?\n\n{} points').format(self.active_laz.data.shape[0]),
                                           _('Descimate factor'), _('Decimation'), 0, 0, 100)

                ret = dlg.ShowModal()

                if ret == wx.ID_CANCEL:
                    dlg.Destroy()
                    return

                descimate_fact = dlg.GetValue()
                dlg.Destroy()

                if descimate_fact > 0:
                    self.descimate_laz_data(descimate_fact)

        elif itemlabel == _('Fill active array from LAZ data'):

            if self.mylazgrid is None:
                logging.warning('')
                return
            if self.active_array is None:
                logging.warning(_('No active array -- select an array first and retry!'))
                return

            self.fill_active_array_from_laz(self.active_array)

        elif itemlabel == _('Count LAZ data in cells'):

            if self.mylazgrid is None:
                logging.warning('')
                return
            if self.active_array is None:
                logging.warning(_('No active array -- select an array first and retry!'))
                return

            self.count_active_array_from_laz(self.active_array)

        elif itemlabel == _('Select cells in array from LAZ data'):
            if self.mylazgrid is None:
                logging.warning('')
                return
            if self.active_array is None:
                logging.warning(_('No active array -- select an array first and retry!'))
                return

            self.select_active_array_from_laz(self.active_array)

        elif itemlabel == _('Plot LAZ around active vector'):
            self.plot_laz_around_active_vec()

        elif itemlabel == _('Plot LAZ around temporary vector'):

            self.active_vector = vector()
            self.active_vector.add_vertex(wolfvertex(0.,0.))
            self.mimicme()

            self.start_action('laz tmp vector', _('LAZ tmp'))

        elif itemlabel == _('Change colors - Classification'):

            if self.mylazgrid is not None:
                self.mylazgrid.colors.interactive_update_colors()

        elif itemlabel == _('Multiviewer'):

            dlg = wx.NumberEntryDialog(self, _("Additional viewers"), _("How many?"), _("How many additional viewers?"),1, 0, 5)
            ret = dlg.ShowModal()

            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            nb = dlg.GetValue()
            dlg.Destroy()

            if nb > 0:

                # Renaming the current viewer
                dlg = wx.TextEntryDialog(None, _('New name for the current viewer'), _('Rename'), self.viewer_name, style=wx.OK | wx.CANCEL)
                ret = dlg.ShowModal()
                if ret == wx.ID_OK:
                    self.viewer_name = dlg.GetValue()
                    self.SetName(self.viewer_name)
                dlg.Destroy()

                for i in range(nb):
                    self.add_viewer_and_link()
            else:
                logging.warning(_('No additional viewer !'))

        elif itemlabel == _('3D viewer'):

            self.active_viewer3d = Wolf_Viewer3D(self, _("3D Viewer"))
            self.active_viewer3d.Show()
            self.myviewers3d.append(self.active_viewer3d)

            for curarray in self.iterator_over_objects(draw_type.ARRAYS):
                curarray:WolfArray
                if curarray.checked:
                    if curarray._array3d is None:
                        curarray.prepare_3D()

                    if self.active_viewer3d not in curarray.viewers3d:
                        curarray.viewers3d.append(self.active_viewer3d)

                    self.active_viewer3d.add_array(curarray.idx, curarray._array3d)
                    self.active_viewer3d.autoscale()

        elif itemlabel == _('Create/Open multiblock model'):

            self.create_2D_MB_model()

        elif itemlabel == _('Create/Open GPU model'):

            self.create_2D_GPU_model()

        elif itemlabel == _('Create/Open Hydrological model'):

            self.open_hydrological_model()

        elif itemlabel == _('Check headers'):

            self.check_2D_MB_headers()

        elif itemlabel == _('Set comparison'):

            autoscale = True

            # Comparaison de deux résultats ou de deux matrices

            self.compare_results = Compare_Arrays_Results(self, share_cmap_array= True, share_cmap_diff= True)

            add_elt = True
            while add_elt:
                add_elt = self.compare_results.add()

            if len(self.compare_results.paths) < 2 :
                logging.warning(_('Not enough elements to compare !'))
                self.compare_results = None
                return

            self.compare_results.bake()

        elif id == wx.ID_EXIT:

            dlg = wx.MessageDialog(None,_('Do you really want to quit?'), style = wx.YES_NO|wx.NO_DEFAULT)
            ret=dlg.ShowModal()
            if ret == wx.ID_YES:
                wx.Exit()
            else:
                dlg.Destroy()

        elif id == wx.ID_FILE1:
            self.add_object(which='array', ToCheck=True)

        elif itemlabel == _('Add view...'):
            self.add_object(which='views', ToCheck=True)

        elif itemlabel == _('Add tiles GPU...'):
            self.add_object(which='array_tiles', ToCheck=True)

        elif itemlabel == _('Add tiles...'):
            self.add_object(which='tiles', ToCheck=True)

        elif itemlabel ==_('Add images tiles...'):
            self.add_object(which='imagestiles', ToCheck=True)

        elif itemlabel == _('Add tiles comparator...'):
            self.add_object(which='tilescomp', ToCheck=True)

        elif id == wx.ID_FILE2:
            self.add_object(which='vector', ToCheck=True)

        elif id == wx.ID_FILE3:
            self.add_object(which='cloud', ToCheck=True)

        elif itemlabel == _('Add triangulation...'):
            self.add_object(which='triangulation', ToCheck=True)

        elif itemlabel == _('Add particle system...'):
            self.add_object(which = 'particlesystem', ToCheck = True)
            self.menu_particlesystem()

        elif itemlabel == _('Create particle system...'):
            self.active_particle_system = newpart = Particle_system()
            self.add_object(which='particlesystem', newobj=newpart, ToCheck=True)
            self.menu_particlesystem()

        elif id == wx.ID_FILE4:
            self.add_object(which='cross_sections', ToCheck=True)

        elif itemlabel == _('Add Wolf2D results...'):
            self.add_object(which='res2d', ToCheck=True)
            self.menu_wolf2d()

        elif itemlabel == _('Add Wolf2D GPU results...'):
            self.add_object(which='res2d_gpu', ToCheck=True)
            self.menu_wolf2d()
            self.menu_2dgpu()

        elif itemlabel == _('Add bridges...'):
            self.add_object(which='bridges', ToCheck=True)

        elif itemlabel == _('Add weirs...'):
            self.add_object(which='weirs', ToCheck=True)

        elif itemlabel == _('Add array and crop...'):
            self.add_object(which='array_crop', ToCheck=True)

        elif itemlabel == _('Precomputed DEM'):

            dlg = Precomputed_DEM_DTM_Dialog(self, _('Precomputed DEM'), self.default_dem, self)
            ret = dlg.ShowModal()

        elif itemlabel == _('Precomputed DTM'):

            dlg = Precomputed_DEM_DTM_Dialog(self, _('Precomputed DTM'), self.default_dtm, self)
            ret = dlg.ShowModal()

        elif itemlabel == _('Create array from bathymetry file...'):

            self.add_object(which='array_xyz', ToCheck=True)

        elif itemlabel == _('Create array from Lidar 2002...'):

            dlg = wx.SingleChoiceDialog(None, _('What source of data?'), _('Lidar 2002'),
                                        [_('First echo'), _('Second echo')])

            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                return

            sel = dlg.GetStringSelection()

            if sel == _('First echo'):
                self.add_object(which='array_lidar_first', ToCheck=True)
            elif sel == _('Second echo'):
                self.add_object(which='array_lidar_second', ToCheck=True)

        elif id == wx.ID_FILE5:

            def addscandir(mydir):
                for entry in scandir(mydir):
                    if entry.is_dir():
                        addscandir(entry)
                    elif entry.is_file():
                        if entry.name.endswith('.vec') or entry.name.endswith('.vecz'):

                            msg = wx.MessageDialog(self,
                                                   _(entry.name + ' found in ' + mydir + '\n\n Is it a "cross sections" file?'),
                                                   style=wx.YES_NO | wx.NO_DEFAULT)
                            ret = msg.ShowModal()
                            if ret == wx.ID_YES:
                                self.add_object(which='vector',
                                                filename=join(mydir, entry.name),
                                                ToCheck=True,
                                                id=join(mydir, entry.name))
                            else:
                                self.add_object(which='cross_sections',
                                                filename=join(mydir, entry.name),
                                                ToCheck=True,
                                                id=join(mydir, entry.name))

                        elif entry.name.endswith(('.bin', '.tif', '.npy')):
                            self.add_object(which='array',
                                            filename=join(mydir, entry.name),
                                            ToCheck=True,
                                            id=join(mydir, entry.name))

            mydialog = wx.DirDialog(self, _("Choose directory to scan"))
            if mydialog.ShowModal() == wx.ID_CANCEL:
                mydialog.Destroy()
                return
            else:
                # récupération du nom de fichier avec chemin d'accès
                mydir = mydialog.GetPath()
                mydialog.Destroy()

            if exists(mydir):
                addscandir(mydir)

        elif id == wx.ID_FILE6:
            # Création d'une nouvelle matrice
            newarray = WolfArray(create=True, mapviewer=self)
            self.add_object('array', newobj=newarray)

        elif itemlabel == _('Create view...'):

            # Création d'une nouvelle vue
            newview = WolfViews(mapviewer=self)
            self.add_object('views', newobj=newview)

        elif itemlabel==_('Create Wolf2D manager ...'):

            from .mesh2d.config_manager import config_manager_2D
            newmanager = config_manager_2D(mapviewer=self)

        elif itemlabel==_('Create scenarios manager ...'):

            from .scenario.config_manager import Config_Manager_2D_GPU
            newmanager = Config_Manager_2D_GPU(mapviewer=self, create_ui_if_wx=True)

        elif itemlabel == _('Create acceptability manager...'):

            from .acceptability.acceptability_gui import AcceptabilityGui
            newmanager = AcceptabilityGui()
            newmanager.mapviewer = self
            newmanager.Show()

        elif itemlabel == _('Create INBE manager...'):

            from .insyde_be.INBE_gui import INBEGui
            newmanager = INBEGui()
            newmanager.mapviewer = self
            newmanager.Show()

        elif itemlabel==_('Create BC manager Wolf2D...'):

            if self.active_array is not None:

                choices = {'WOLF prev':1, 'WOLF OO':2, 'GPU':3}

                dlg = wx.SingleChoiceDialog(None,
                                            _("Which version of BC Manager"),
                                            _("Version"),
                                            ['WOLF prev', 'WOLF OO', 'GPU'])
                ret = dlg.ShowModal()
                if ret == wx.ID_CANCEL:
                    dlg.Destroy()
                    return

                method = dlg.GetStringSelection()
                dlg.Destroy()

                which_version = choices[method]

                self.mybc.append(BcManager(self,
                                           linked_array=self.active_array,
                                           version = which_version,
                                           DestroyAtClosing=False,
                                           Callback=self.pop_boundary_manager,
                                           mapviewer=self))
                ret = self.mybc[-1].FindBorders()
                if ret == -1:
                    self.mybc.pop(-1)
                    return
                self.active_bc = self.mybc[-1]

        elif itemlabel == _('Create Wolf1D...'):

            self.frame_create1Dfrom2D = GuiNotebook1D(mapviewer= self)
            logging.info(_(f'New window available - Wolf1D.'))

        elif itemlabel in [_('Create dike...'), _('Add dike...')]:
            self.new_dike(itemlabel)

        elif itemlabel in [_('Add picture collection...')]:
            # Création d'une nouvelle collection de photos

            dlg = wx.SingleChoiceDialog(None, _('Choose the type of picture collection'), _('Picture Collection'),
                                        [_('Pictures + shapefile'),
                                         _('Wolf vec format'),
                                         _('Georeferenced pictures'),
                                         _('Pictures + Excel'),
                                         _('URL zip file')])
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            itemlabel = dlg.GetStringSelection()
            dlg.Destroy()

            if itemlabel in [_('Pictures + shapefile'), _('Georeferenced pictures'), _('Pictures + Excel')]:
                dlgdir = wx.DirDialog(self, _('Choose directory to scan for pictures'))
                if dlgdir.ShowModal() == wx.ID_CANCEL:
                    dlgdir.Destroy()
                    return
                mydir = dlgdir.GetPath()
                dlgdir.Destroy()
            elif itemlabel == _('URL zip file'):
                # Demande de l'URL du fichier zip
                dlgurl = wx.TextEntryDialog(self, _('Enter the URL of the zip file containing the pictures'), _('URL zip file'))
                if dlgurl.ShowModal() == wx.ID_CANCEL:
                    dlgurl.Destroy()
                    return
                mydir = dlgurl.GetValue()
                dlgurl.Destroy()
            else:
                dlgfile = wx.FileDialog(self, _('Choose shapefile'), style=wx.FD_OPEN | wx.FD_FILE_MUST_EXIST,
                                        wildcard='Wolf vec (*.vec)| *.vec|Wolf vecz (*.vecz)|*.vecz|All files (*.*)|*.*')
                if dlgfile.ShowModal() == wx.ID_CANCEL:
                    dlgfile.Destroy()
                    return
                mydir = dlgfile.GetPath()
                dlgfile.Destroy()

            if itemlabel == _('Pictures + shapefile'):
                # Création d'une nouvelle collection de photos avec shapefile
                newcollection = PictureCollection(parent=self, mapviewer=self)
                newcollection.load_from_directory_with_shapefile(mydir)
            elif itemlabel == _('Georeferenced pictures'):
                # Création d'une nouvelle collection de photos géoréférencées
                newcollection = PictureCollection(parent=self, mapviewer=self)
                newcollection.load_from_directory_georef_pictures(mydir)
            elif itemlabel == _('Pictures + Excel'):
                # Création d'une nouvelle collection de photos avec Excel
                newcollection = PictureCollection(parent=self, mapviewer=self)
                newcollection.load_from_directory_with_excel(mydir)
            elif itemlabel == _('Wolf vec format'):
                # Création d'une nouvelle collection de photos avec format Wolf vec
                newcollection = PictureCollection(filename = mydir, parent=self, mapviewer=self)
            elif itemlabel == _('URL zip file'):
                # Création d'une nouvelle collection de photos à partir d'un fichier zip
                newcollection = PictureCollection(parent=self, mapviewer=self)
                newcollection.load_from_url_zipfile(mydir)

            count = 0
            for zone in newcollection.myzones:
                count += zone.nbvectors

            if count == 0:
                logging.warning(_('No usable pictures found in the collection !'))
                return

            self.add_object('picture_collection', newobj=newcollection, ToCheck=True)

        elif itemlabel in [_('Create a drowning...'), _('Add a drowning result...')]:

            self.newdrowning(itemlabel)


        elif id == wx.ID_FILE7:
            autoscale = False
            # Création de nouveaux vecteurs
            newzones = Zones(parent=self)
            self.add_object('vector', newobj=newzones)

        elif id == wx.ID_FILE8:
            # Création d'un nouveau nuage de point
            newcloud = cloud_vertices()
            self.add_object('cloud', newobj=newcloud)

        elif id in self.tools.keys():
            # gestion des actions
            self.ManageActions(id)

        elif id == wx.ID_SAVE:

            for obj in self.iterator_over_objects(draw_type.ARRAYS):
                obj: WolfArray

                if obj.filename == '':
                    filterArray = "bin (*.bin)|*.bin|Geotif (*.tif)|*.tif|Numpy (*.npy)|*.npy|all (*.*)|*.*"
                    fdlg = wx.FileDialog(self, "Choose file", wildcard=filterArray, style=wx.FD_SAVE)
                    fdlg.ShowModal()
                    if fdlg.ShowModal() == wx.ID_OK:
                        obj.filename = fdlg.GetPath()

                obj.write_all()

            for obj in self.iterator_over_objects(draw_type.VECTORS):
                obj:Zones
                obj.saveas()

        elif itemlabel == 'Save to image...':

            fn, ds = self.save_canvasogl(mpl=True)

            all_images = self.save_linked_canvas(fn[:-4], mpl= True, ds= ds, add_title= True)
            self.assembly_images(all_images, mode= self.assembly_mode)

        elif itemlabel == _('Copy image...'):

            self.copy_canvasogl()

        elif itemlabel == _('Export...'):

            curarray: WolfArray
            curvec: vector

            msg = ''
            if self.active_array is None:
                msg += _('Active array is None\n')
            if self.active_vector is None:
                msg += _('Active vector is None\n')

            if msg != '':
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            curarray = self.active_array
            curvec = self.active_vector

            curvec.find_minmax()

            i1, j1 = curarray.get_ij_from_xy(curvec.xmin, curvec.ymin)
            x1, y1 = curarray.get_xy_from_ij(i1, j1)
            x1 -= curarray.dx / 2.
            y1 -= curarray.dy / 2.

            i2, j2 = curarray.get_ij_from_xy(curvec.xmax, curvec.ymax)
            x2, y2 = curarray.get_xy_from_ij(i2, j2)
            x2 += curarray.dx / 2.
            y2 += curarray.dy / 2.
            mybounds = [[x1, x2], [y1, y2]]

            dlg = wx.FileDialog(None, _('Choose filename'),
                                wildcard='glb (*.glb)|*.glb|gltf2 (*.gltf)|*.gltf|All (*.*)|*.*', style=wx.FD_SAVE)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            fn = dlg.GetPath()
            dlg.Destroy()

            with wx.BusyInfo(_('Export to gltf/glb')):
                wait = wx.BusyCursor()
                curarray.export_to_gltf(mybounds, fn)
                del wait

        elif itemlabel == _('Import...'):

            curarray: WolfArray

            msg = ''
            if self.active_array is None:
                msg += _('Active array is None\n')

            if msg != '':
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            curarray = self.active_array

            dlg = wx.FileDialog(None, _('Choose filename'),
                                wildcard='glb (*.glb)|*.glb|gltf2 (*.gltf)|*.gltf|All (*.*)|*.*', style=wx.FD_OPEN)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            fn = dlg.GetPath()
            dlg.Destroy()

            dlg = wx.FileDialog(None, _('Choose pos filename'), wildcard='pos (*.pos)|*.pos|All (*.*)|*.*',
                                style=wx.FD_OPEN)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            fnpos = dlg.GetPath()
            dlg.Destroy()

            choices = ["matplotlib", "scipy"] #, "pyvista"]
            dlg = wx.SingleChoiceDialog(None, _("Pick an interpolation method"), _("Choices"), choices)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            method = dlg.GetStringSelection()
            dlg.Destroy()

            with wx.BusyInfo(_('Importing gltf/glb')):
                wait = wx.BusyCursor()
                try:
                    curarray.import_from_gltf(fn, fnpos, method)
                except:
                    pass
                del wait

        elif itemlabel == _('Compare...'):

            msg = ''
            if self.active_array is None:
                msg += _('Active array is None\n')

            if msg != '':
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            self.set_blender_sculpting()
            autoscale = False
            self.set_fn_fnpos_gltf()
            self.update_blender_sculpting()

        elif itemlabel == _('Update...'):

            msg = ''
            if self.active_array is None:
                msg += _('Active array is None\n')

            if msg != '':
                msg += _('\n')
                msg += _('Retry !\n')
                wx.MessageBox(msg)
                return

            self.set_fn_fnpos_gltf()
            self.update_blender_sculpting()

        elif id == wx.ID_SAVEAS:

            for obj in self.iterator_over_objects(draw_type.ARRAYS):
                obj: WolfArray

                filterArray = "bin (*.bin)|*.bin|Geotif (*.tif)|*.tif|Numpy (*.npy)|*.npy|all (*.*)|*.*"
                fdlg = wx.FileDialog(self, "Choose file name for Array : " + obj.idx, wildcard=filterArray,
                                        style=wx.FD_SAVE)
                ret = fdlg.ShowModal()
                if ret == wx.ID_OK:
                    obj.filename = fdlg.GetPath()
                    obj.write_all()

            for obj in self.iterator_over_objects(draw_type.VECTORS):
                obj:Zones
                if obj.idx=='grid':
                    pass
                else:
                    filterArray = "vec (*.vec)|*.vec|vecz (*.vecz)|*.vecz|Shapefile (*.shp)|*.shp|all (*.*)|*.*"
                    fdlg = wx.FileDialog(self, "Choose file name for Vector :" + obj.idx, wildcard=filterArray,
                                         style=wx.FD_SAVE)
                    ret = fdlg.ShowModal()
                    if ret == wx.ID_OK:
                        obj.saveas(fdlg.GetPath())

        if len(self.myarrays) + len(self.myvectors) + len(self.myclouds) + len(self.mytri) + len(self.myres2D) + len(self.mytiles) + len(self.myimagestiles) + len(self.mypartsystems) + len(self.mydikes) + len(self.mydrownings) + len(self.myinjectors) == 2 or autoscale:
            # Trouve les bornes si un seul élément est présent, sinon on conserve l'état du zoom
            self.Autoscale()

    def pop_boundary_manager(self, which:BcManager):
        """ Pop a boundary condition manager after Destroying """

        idx = self.mybc.index(which)
        if self.active_bc is which:
            self.active_bc = None
        self.mybc.pop(idx)

        self.Refresh()


    def get_boundary_manager(self, which:WolfArray):
        """ Get a boundary manager """

        for curbc in self.mybc:
            if curbc.linked_array is which:
                return curbc

        return None

    def uniquecolormap(self, loadfromfile = False):
        """ Compute unique colormap from all (arrays, 2D results) and apply it to all """

        workingarray=[]
        nbnotnull=0

        newpal = wolfpalette(self)

        if loadfromfile :
            newpal.readfile()

            if not newpal.is_valid():
                logging.warning(_('Palette not valid !'))
                return
        else:
            nb  = len(self.myarrays) + len(self.myres2D)
            pgbar = wx.ProgressDialog(_('Compute unique colormap'), _('Compute unique colormap from all arrays'), maximum=nb, parent=self, style=wx.PD_APP_MODAL|wx.PD_AUTO_HIDE)

            curarray:WolfArray
            curres2d:Wolfresults_2D

            for curarray in self.myarrays:
                if curarray.plotted:
                    workingarray.append(curarray.get_working_array())
                    nbnotnull+=curarray.nbnotnull
                pgbar.Update(pgbar.GetValue() + 1, _('Compute unique colormap from array : ') + curarray.idx)

            for curres2d in self.myres2D:
                if curres2d.plotted:
                    workingarray.append(curres2d.get_working_array())
                    nbnotnull+=curres2d.nbnotnull
                pgbar.Update(pgbar.GetValue() + 1, _('Compute unique colormap from 2D result : ') + curres2d.idx)

            pgbar.Destroy()

            workingarray = np.concatenate(workingarray)

            newpal.default16()
            newpal.isopop(workingarray, nbnotnull)

        nb = len(self.myarrays) + len(self.myres2D)
        pgbar = wx.ProgressDialog(_('Applying colormap'), _('Applying colormap to all arrays'), maximum=nb, parent=self, style=wx.PD_APP_MODAL|wx.PD_AUTO_HIDE)
        for curarray in self.myarrays:
            if curarray.plotted:
                curarray.mypal.automatic = False
                curarray.myops.palauto.SetValue(0)
                curarray.mypal.values = newpal.values.copy()
                curarray.mypal.colors = newpal.colors.copy()
                curarray.mypal.fill_segmentdata()
                curarray.reset_plot()
            pgbar.Update(pgbar.GetValue() + 1, _('Applying colormap to array : ') + curarray.idx)

        for curres2d in self.myres2D:
            if curres2d.plotted:
                curres2d.mypal.automatic = False
                curres2d.mypal.nb     = newpal.nb
                curres2d.mypal.values = newpal.values.copy()
                curres2d.mypal.colors = newpal.colors.copy()
                curres2d.mypal.fill_segmentdata()
                curres2d.reset_plot()
            pgbar.Update(pgbar.GetValue() + 1, _('Applying colormap to 2D result : ') + curres2d.idx)

        pgbar.Destroy()

    def loadnap_and_apply(self):

        dlg = wx.MessageDialog(self,_('Load mask for all?'),style=wx.YES_NO|wx.YES_DEFAULT)
        ret=dlg.ShowModal()

        if ret == wx.ID_NO:
            dlg.Destroy()
            return

        with wx.BusyInfo(_('Loading masks')):
            wait = wx.BusyCursor()
            curarray:WolfArray
            for curarray in self.myarrays:
                if curarray.plotted:
                    curarray.loadnap_and_apply()
            del wait

    def uniforminparts_all(self, TrueOrFalse:bool):

        for curarray in self.myarrays:
            curarray:WolfArray
            if curarray.plotted:
                curarray.mypal.interval_cst = TrueOrFalse
                curarray.reset_plot()

        for curarray in self.myres2D:
            curarray:Wolfresults_2D
            if curarray.plotted:
                curarray.mypal.interval_cst = TrueOrFalse
                curarray.link_palette()
                curarray.reset_plot()

    def filter_inundation(self):

        dlg = wx.TextEntryDialog(self,_('Upper bound \n\n All values strictly lower than the bound will not be extracted !'),value='.0005')
        ret=dlg.ShowModal()

        if ret == wx.ID_CANCEL:
            dlg.Destroy()
            return

        bound = float(dlg.GetValue())
        dlg.Destroy()

        logging.info(_('Filtering results'))

        curarray:WolfArray
        for curarray in self.myarrays:
            if curarray.plotted:
                curarray.filter_inundation(epsilon = bound)
                curarray.filter_independent_zones(n_largest = 1)

        curarray:Wolfresults_2D
        for curarray in self.myres2D:
            if curarray.plotted:
                curarray.filter_inundation(eps = bound)
                curarray.filter_independent_zones(n_largest = 1)

        logging.info(_('Filtering done !'))

    def export_results_as(self, which:Literal['geotiff','shape','numpy'] = None, multiband:bool = None):
        """
        Export des résultats WOLF2D vers différents formats.
        Au moins un résultat doit être chargé pour pouvoir être exporté.
        """

        dlg = wx.DirDialog(self,_('Choose output directory'), style = wx.DD_DIR_MUST_EXIST)
        ret=dlg.ShowModal()

        if ret == wx.ID_CANCEL:
            logging.warning(_('Abort!'))
            dlg.Destroy()
            return

        outdir = dlg.GetPath()
        dlg.Destroy()

        if which not in ['geotiff','shape','numpy']:
            dlg = wx.SingleChoiceDialog(self,_('Choose output format'), _('Format'), ['Geotiff','Shape file','Numpy array'])
            ret = dlg.ShowModal()

            if ret == wx.ID_CANCEL:
                logging.warning(_('Abort!'))
                dlg.Destroy()
                return

            sel = dlg.GetSelection()

            if sel == 0:
                which = 'geotiff'
            elif sel == 1:
                which = 'shape'
            else:
                which = 'numpy'

            dlg.Destroy()

        if which == 'geotiff':
            if multiband is None:
                dlg = wx.SingleChoiceDialog(self,_('Choose output format'), _('Format'), ['Multiband (single file)',
                                                                                          'Single band (multiple files)'])
                dlg.ShowModal()

                sel = dlg.GetSelection()
                if sel == 1:
                    multiband = False
                else:
                    multiband = True
                dlg.Destroy()

        logging.info(_('Exporting results -- Be patient !'))

        loaded_res = self.get_list_keys(drawing_type= draw_type.RES2D, checked_state=None)

        dlg = wx.MultiChoiceDialog(self,_('Choose results to export'), _('Results'), choices=loaded_res)
        dlg.SetSelections([idx for idx, res in enumerate(loaded_res) if self.get_obj_from_id(res, drawtype=draw_type.RES2D).plotted])
        ret = dlg.ShowModal()
        if ret == wx.ID_CANCEL:
            logging.warning(_('Abort!'))
            dlg.Destroy()
            return

        sel = dlg.GetSelections() # Get a list if integers
        sel_res = [self.get_obj_from_id(loaded_res[cursel], drawtype=draw_type.RES2D) for cursel in sel] # convert to list of objects

        dlg.Destroy()

        if len(sel) == 0:
            logging.warning(_('No results selected for export'))
            return

        fields = [(views_2D.TOPOGRAPHY, True),
                  (views_2D.WATERDEPTH, True),
                  (views_2D.QX, True),
                  (views_2D.QY, True),
                  (views_2D.UNORM, True),
                  (views_2D.FROUDE, True),
                  (views_2D.HEAD, True),
                  (views_2D.CRITICAL_DIAMETER_SHIELDS, False),
                  (views_2D.CRITICAL_DIAMETER_IZBACH, False),
                  (views_2D.QNORM, False),
                  (views_2D.WATERLEVEL, False),
                  (views_2D.CRITICAL_DIAMETER_SUSPENSION_50, False),
                  (views_2D.CRITICAL_DIAMETER_SUSPENSION_100, False),]

        dlg = wx.MultiChoiceDialog(self,_('Choose fields to export'), _('Fields'), choices= [str(field[0]) for field in fields])
        dlg.SetSelections([idx for idx, field in enumerate(fields) if field[1]])
        ret = dlg.ShowModal()

        if ret == wx.ID_CANCEL:
            logging.warning(_('Abort!'))
            dlg.Destroy()
            return

        sel_fields = dlg.GetSelections() # Get a list if integers
        dlg.Destroy()

        if len(sel_fields) == 0:
            logging.warning(_('No fields selected for export'))
            return

        # Get the views_2D values associated with the selected field names
        fields = [fields[cursel][0] for cursel in sel_fields]

        for cur_res in tqdm(sel_res):
            cur_res:Wolfresults_2D
            cur_res.export_as(outdir, fields, which, multiband)

        logging.info(_('Export done -- Thanks for your patience !'))

    def export_shape(self, outdir:str= '', fn:str = '', myarrays:list[WolfArray]= [], descr:list[str]= [], mask:WolfArray=None):
        """ Export multiple arrays to shapefile

        :param outdir: output directory
        :param fn: filename -- .shp will be added if not present
        :param myarrays: list of Wolfarrays to export
        :param descr: list of descriptions
        :param mask: mask array -- export only where mask > 0
        """

        if len(myarrays)==0:
            logging.warning(_('No arrays provided for shapefile export'))
            return

        if mask is None:
            logging.warning(_('No mask provided for shapefile export'))
            return

        from osgeo import gdal, osr, gdalconst,ogr

        # create the spatial reference system, Lambert72
        srs = osr.SpatialReference()
        srs.ImportFromEPSG(self.epsg)

        # create the data source
        driver: ogr.Driver
        driver = ogr.GetDriverByName("ESRI Shapefile")

        # create the data source
        filename = join(outdir,fn)
        if not filename.endswith('.shp'):
            filename+='.shp'

        ds = driver.CreateDataSource(filename)

        # create one layer
        layer = ds.CreateLayer("results", srs, ogr.wkbPolygon)

        # Add ID fields
        idFields=[]
        for curlab in descr:
            idFields.append(ogr.FieldDefn(curlab, ogr.OFTReal))
            layer.CreateField(idFields[-1])

        # Create the feature and set values
        featureDefn = layer.GetLayerDefn()
        feature = ogr.Feature(featureDefn)

        usednodes = np.argwhere(mask.array>0.)
        for i,j in usednodes:

            x,y = mask.get_xy_from_ij(i,j)
            # Creating a line geometry
            ring = ogr.Geometry(ogr.wkbLinearRing)
            ring.AddPoint(x-mask.dx/2,y-mask.dy/2)
            ring.AddPoint(x+mask.dx/2,y-mask.dy/2)
            ring.AddPoint(x+mask.dx/2,y+mask.dy/2)
            ring.AddPoint(x-mask.dx/2,y+mask.dy/2)
            ring.AddPoint(x-mask.dx/2,y-mask.dy/2)

            # Create polygon
            poly = ogr.Geometry(ogr.wkbPolygon)
            poly.AddGeometry(ring)

            feature.SetGeometry(poly)

            for arr, id in zip(myarrays,descr):

                feature.SetField(id, float(arr.array[i,j]))

            layer.CreateFeature(feature)

        feature = None

        # Save and close DataSource
        ds = None

    def export_geotif(self, outdir:str= '', fn:str = '', myarrays:list[WolfArray]= [], descr:list[str]= [], multiband:bool= True):
        """ Export multiple arrays to geotiff

        :param outdir: output directory
        :param fn: filename -- .tif will be added if not present
        :param myarrays: list of Wolfarrays to export
        :param descr: list of descriptions -- Bands names

        """

        if len(myarrays)==0:
            logging.warning(_('No arrays provided for geotiff export'))
            return

        from osgeo import gdal, osr, gdalconst

        srs = osr.SpatialReference()
        srs.ImportFromEPSG(self.epsg)

        driver: gdal.Driver
        out_ds: gdal.Dataset
        band: gdal.Band
        driver = gdal.GetDriverByName("GTiff")

        if multiband:
            filename = join(outdir,fn)
            if not filename.endswith('.tif'):
                filename+='.tif'

            arr = myarrays[0]
            out_ds = driver.Create(filename, arr.shape[0], arr.shape[1], len(myarrays), arr.dtype_gdal, options=['COMPRESS=LZW'])
            out_ds.SetProjection(srs.ExportToWkt())
            out_ds.SetGeoTransform([myarrays[0].origx+myarrays[0].translx,
                                    myarrays[0].dx,
                                    0.,
                                    myarrays[0].origy+myarrays[0].transly,
                                    0.,
                                    myarrays[0].dy])

            k=1
            for arr, name in zip(myarrays,descr):
                band = out_ds.GetRasterBand(k)
                band.SetNoDataValue(0.)
                band.SetDescription(name)
                band.WriteArray(arr.array.transpose())
                band.FlushCache()
                band.ComputeStatistics(True)
                k+=1

            out_ds = None

        else:
            for arr, name in zip(myarrays,descr):

                if filename.endswith('.tif'):
                    filename = filename[:-4]
                filename = join(outdir,fn+'_'+name)
                filename += '.tif'

                out_ds = driver.Create(filename, arr.shape[0], arr.shape[1], 1, arr.dtype_gdal, options=['COMPRESS=LZW'])
                out_ds.SetProjection(srs.ExportToWkt())
                out_ds.SetGeoTransform([myarrays[0].origx+myarrays[0].translx,
                                        myarrays[0].dx,
                                        0.,
                                        myarrays[0].origy+myarrays[0].transly,
                                        0.,
                                        myarrays[0].dy])

                band = out_ds.GetRasterBand(1)
                band.SetNoDataValue(0.)
                band.SetDescription(name)
                band.WriteArray(arr.array.transpose())
                band.FlushCache()
                band.ComputeStatistics(True)
                out_ds = None

    def get_linked_arrays(self, linked:bool = True) -> dict:
        """ Get all arrays in the viewer and linked viewers """

        linkedarrays = {}

        if self.linked and linked:
            all_dicts = [curviewer.get_linked_arrays(linked = False) for curviewer in self.linkedList]
            for curdict in all_dicts:
                linkedarrays.update(curdict)
        else:
            for locarray in self.iterator_over_objects(draw_type.ARRAYS):
                linkedarrays[locarray.idx] = locarray

            for locarray in self.iterator_over_objects(draw_type.RES2D):
                linkedarrays[locarray.idx] = locarray

        return linkedarrays

    def save_linked_canvas(self, fn:str, mpl:bool= True, ds:float= 0., add_title:bool= True) -> tuple[(str, float), str]:
        """ Save canvas of all linked viewers

        :param fn: filename without extension -- '.png' will be added
        :param mpl: save as matplotlib image
        :param ds: Ticks size for matplotlib image
        :return: list of tuple ((filename, ds), viewer_name)
        """

        fn = str(fn)
        ret = []
        if self.linked:
            for idx, curel in enumerate(self.linkedList):
                ret.append((curel.save_canvasogl(fn + '_' + str(idx) + '.png', mpl, ds, add_title= add_title), self.viewer_name))

        return ret

    def save_arrays_indep(self, fn:str, mpl:bool= True, ds:float= 0., add_title:bool= True) -> tuple[(str, float), str]:
        """ Save each array in a separate file

        :param fn: filename without extension -- '.png' will be added
        :param mpl: save as matplotlib image
        :param ds: Ticks size for matplotlib image
        :return: list of tuple ((filename, ds), viewer_name)
        """

        # Get all checked arrays
        checked_arrays = self.get_list_keys(drawing_type= draw_type.ARRAYS, checked_state= True)
        checked_results = self.get_list_keys(drawing_type= draw_type.RES2D, checked_state= True)

        old_active = self.active_array
        old_res2d = self.active_res2d

        if len(checked_arrays) + len(checked_results) == 0:
            logging.warning(_('No arrays checked for export'))
            return []

        def uncheck_all():
            # uncheck arrays
            for curarray in checked_arrays:
                self.uncheck_id(curarray, unload= False, forceresetOGL= False)
            for curres in checked_results:
                self.uncheck_id(curres, unload= False, forceresetOGL= False)

        fn = str(fn)
        ret = []
        for idx, curel in enumerate(checked_arrays):
            uncheck_all()
            self.check_id(curel)
            self.active_array = self.get_obj_from_id(curel, drawtype= draw_type.ARRAYS)
            ret.append((self.save_canvasogl(fn + '_' + str(idx) + '.png', mpl, ds, add_title= add_title, arrayid_as_title=True), curel))

        for idx, curel in enumerate(checked_results):
            uncheck_all()
            self.check_id(curel)
            self.active_res2d = self.get_obj_from_id(curel, drawtype= draw_type.RES2D)
            ret.append((self.save_canvasogl(fn + '_' + str(idx + len(checked_arrays)) + '.png', mpl, ds, add_title= add_title, resid_as_title=True), curel))

        self.active_array = old_active
        self.active_res2d = old_res2d

        for curarray in checked_arrays:
            self.check_id(curarray)

        for curres in checked_results:
            self.check_id(curres)

        self.Refresh()

        return ret

    def assembly_images(self, all_images, mode:Literal['horizontal', 'vertical', 'square']= 'square'):
        """ Assembly images

        Every image has the same size (width, height)

        :param all_images: list of tuple (filename, viewer_name)
        :param mode: 'horizontal', 'vertical', 'square'
        """

        assert mode in ['horizontal', 'vertical', 'square', 0, 1, 2], 'Mode not recognized'

        from PIL import Image

        images = [Image.open(fn) for (fn, ds), viewername in all_images]

        if len(images) in [1,2] and (mode == 'square' or mode == 2):
            mode = 'horizontal'

        widths, heights = zip(*(i.size for i in images))

        if mode == 'horizontal' or mode==0:

            total_width = sum(widths)
            max_height = max(heights)

            new_im = Image.new('RGB', (total_width, max_height), color=(255,255,255))

            x_offset = 0
            for im in images:
                new_im.paste(im, (x_offset,0))
                x_offset += im.size[0]

            new_im.save(all_images[0][0][0][:-4] + '_assembly.png')

        elif mode == 'vertical' or mode==1:

            total_height = sum(heights)
            max_width = max(widths)

            new_im = Image.new('RGB', (max_width, total_height), color=(255,255,255))

            y_offset = 0
            for im in images:
                new_im.paste(im, (0, y_offset))
                y_offset += im.size[1]

            new_im.save(all_images[0][0][0][:-4] + '_assembly.png')

        elif mode == 'square' or mode==2:

            max_width = max(widths)
            max_height = max(heights)

            nb_hor = int(np.ceil(np.sqrt(len(images))))

            new_im = Image.new('RGB', (max_width*nb_hor, max_height*nb_hor), color=(255,255,255))

            x_offset = 0
            y_offset = 0
            for idx, im in enumerate(images):
                new_im.paste(im, (x_offset, y_offset))
                x_offset += im.size[0]
                if (idx+1) % nb_hor == 0:
                    y_offset += im.size[1]
                    x_offset = 0

            new_im.save(all_images[0][0][0][:-4] + '_assembly.png')

        return new_im

    def thread_update_blender(self):
        print("Update blender")
        if self.SetCurrentContext():
            self.update_blender_sculpting()
            t = threading.Timer(10.0, self.thread_update_blender)
            t.start()

    def add_object(self,
                   which:Literal['array','array_lidar_first','array_lidar_second','array_xyz','array_tiles',
                                 'bridges',
                                 'weirs',
                                 'vector',
                                 'tiles', 'tilescomp'
                                 'cloud', 'laz',
                                 'triangulation',
                                 'cross_sections',
                                 'other',
                                 'views',
                                 'res2d',
                                 'res2d_gpu',
                                 'particlesystem',
                                 'wmsback',
                                 'wmsfore',
                                 'drowning',
                                 'imagestiles',
                                 'dike',
                                 'injector',
                                 'picture_collection'] = 'array',
                   filename='',
                   newobj=None,
                   ToCheck=True,
                   id=''):

        """
        Add object to current Frame/Drawing area
        """

        filterArray = "All supported formats|*.bin;*.tif;*.tiff;*.top;*.flt;*.npy;*.npz;*.vrt|bin (*.bin)|*.bin|Elevation WOLF2D (*.top)|*.top|Geotif (*.tif)|*.tif|Float ESRI (*.flt)|*.flt|Numpy (*.npy)|*.npy|Numpy named arrays(*.npz)|*.npz|all (*.*)|*.*"
        filterjson = "json (*.json)|*.json|all (*.*)|*.*"
        filterall = "all (*.*)|*.*"
        filterres2d = "all (*.*)|*.*"
        filterVector = "All supported formats|*.vec;*.vecz;*.dxf;*.shp|vec (*.vec)|*.vec|vecz (*.vecz)|*.vecz|dxf (*.dxf)|*.dxf|shp (*.shp)|*.shp|all (*.*)|*.*"
        filterCloud = "xyz (*.xyz)|*.xyz|dxf (*.dxf)|*.dxf|text (*.txt)|*.txt|shp (*.shp)|*.shp|all (*.*)|*.*"
        filterlaz = "laz (*.laz)|*.laz|las (*.las)|*.las|Numpy (*.npz)|*.npz|all (*.*)|*.*"
        filtertri = "tri (*.tri)|*.tri|text (*.txt)|*.txt|dxf (*.dxf)|*.dxf|gltf (*.gltf)|*.gltf|gltf binary (*.glb)|*.glb|*.*'all (*.*)|*.*"
        filterCs = "vecz WOLF (*.vecz)|*.vecz|txt 2022 (*.txt)|*.txt|WOLF (*.sxy)|*.sxy|text 2000 (*.txt)|*.txt|xlsx 2025 (*.xlsx)|*.xlsx|all (*.*)|*.*"
        filterimage = "Geotif (*.tif)|*.tif|all (*.*)|*.*"

        if filename == '' and newobj is None:
            # ouverture d'une boîte de dialogue
            if which.lower() == 'array' or which.lower() == 'array_crop':
                file = wx.FileDialog(self, "Choose file", wildcard=filterArray)
            elif which.lower() == 'imagestiles':
                file = wx.DirDialog(self, "Choose directory containing images")
            elif which.lower() == 'particlesystem':
                file = wx.FileDialog(self, "Choose file", wildcard=filterjson)
            elif which.lower() == 'array_lidar_first' or which.lower() == 'array_lidar_second':
                file = wx.DirDialog(self, "Choose directory containing Lidar data")
            elif which.lower() == 'array_xyz':
                file = wx.DirDialog(self, "Choose directory containing XYZ files")
            elif which.lower() == 'array_tiles':
                file = wx.DirDialog(self, "Choose directory containing GPU results")
            elif which.lower() == 'bridges':
                file = wx.DirDialog(self, "Choose directory containing bridges")
            elif which.lower() == 'weirs':
                file = wx.DirDialog(self, "Choose directory containing weirs")
            elif which.lower() in ['vector', 'tiles', 'tilescomp']:
                file = wx.FileDialog(self, "Choose file", wildcard=filterVector)
            elif which.lower() == 'cloud':
                file = wx.FileDialog(self, "Choose file", wildcard=filterCloud)
            elif which.lower() == 'laz':
                file = wx.FileDialog(self, "Choose file", wildcard=filterlaz)
            elif which.lower() == 'triangulation':
                file = wx.FileDialog(self, "Choose file", wildcard=filtertri)
            elif which.lower() == 'cross_sections':
                file = wx.FileDialog(self, "Choose file", wildcard=filterCs)
            elif which.lower() == 'other':
                file = wx.FileDialog(self, "Choose file", wildcard=filterall)
            elif which.lower() == 'views':
                file = wx.FileDialog(self, "Choose file", wildcard=filterall)
            elif which.lower() == 'res2d':
                file = wx.FileDialog(self, "Choose file", wildcard=filterres2d)
            elif which.lower() == 'res2d_gpu':
                file = wx.DirDialog(self, "Choose directory containing WolfGPU results")
            elif which.lower() == 'drowning':
                file = wx.DirDialog(self, "Choose directory containing the drowning")
            elif which.lower() == 'dike':
                file = wx.DirDialog(self, "Choose directory", wildcard=filterall)
            elif which.lower() == 'picture_collection':
                file = wx.DirDialog(self, "Choose directory containing pictures")

            # FIXME : particularize filters for wmsback and wmsfore
            elif which.lower() == 'wmsback':
                file = wx.FileDialog(self, "Choose file", wildcard=filterimage)
            elif which.lower() == 'wmsfore':
                file = wx.FileDialog(self, "Choose file", wildcard=filterimage)

            if file.ShowModal() == wx.ID_CANCEL:
                file.Destroy()
                return -1
            else:
                # récuparétaion du nom de fichier avec chemin d'accès
                filename = file.GetPath()
                try:
                    curfilter = file.GetFilterIndex()
                except:
                    pass
                file.Destroy()

        if filename != '':
            if (not (os.path.exists(filename))):
                logging.warning("Warning : the following file is not present here : " + filename)
                return -1

        all_ids = self.get_list_keys(None, checked_state=None)

        curtree = None
        if which.lower() == 'array' or which.lower() == 'array_crop':

            curdict = self.myarrays
            curtree = self.myitemsarray

            if newobj is None:

                if str(filename).endswith('.npz'):

                    wait = wx.BusyCursor()
                    logging.info(_('Start of importing arrays from npz file'))

                    with np.load(filename) as data:
                        if 'header' in data.keys():
                            header = data['header']

                            if len(header) == 6:
                                logging.info(_('Header found in npz file'))

                                origx, origy, dx, dy, nbx, nby = header

                                logging.info(_('Origin X : ') + str(origx))
                                logging.info(_('Origin Y : ') + str(origy))
                                logging.info(_('dx : ') + str(dx))
                                logging.info(_('dy : ') + str(dy))
                                logging.info(_('nbx : ') + str(nbx))
                                logging.info(_('nby : ') + str(nby))
                                nbx, nby = int(nbx), int(nby)
                            else:
                                logging.warning(_('Header found in npz file but not complete -- Only {} values found - Must be 6').format(len(header)))

                            for key, curarray in data.items():
                                if isinstance(curarray, np.ndarray):
                                    if curarray.shape == (nby, nbx):
                                        logging.info("Importing array : " + key)
                                        curhead = header_wolf()
                                        curhead.origx, curhead.origy, curhead.dx, curhead.dy, curhead.nbx, curhead.nby = origx, origy, dx, dy, nbx, nby
                                        newobj = WolfArray(srcheader=curhead, idx = key)
                                        newobj.set_array_from_numpy(curarray)
                                        self.add_object('array', newobj= newobj, id= key)
                        else:
                            origx, origy, dx, dy, nbx, nby = 0.,0.,1,1.,1,1
                            for key, curarray in data.items():
                                if isinstance(curarray, np.ndarray):
                                    logging.info(_('No header found in npz file - Using default values for header'))
                                    logging.info("Importing array : " + key)
                                    curhead = header_wolf()
                                    curhead.origx, curhead.origy, curhead.dx, curhead.dy, curhead.nbx, curhead.nby = 0., 0., 1., 1., curarray.shape[0], curarray.shape[1]
                                    newobj = WolfArray(srcheader=curhead, idx = key)
                                    newobj.set_array_from_numpy(curarray)
                                    self.add_object('array', newobj= newobj, id= key)

                    logging.info(_('End of importing arrays from npz file'))
                    del wait
                    return -1
                else:
                    testobj = WolfArray()
                    testobj.filename = filename
                    testobj.read_txt_header()

                    if testobj.wolftype in WOLF_ARRAY_MB:
                        # with wx.BusyInfo(_('Importing array')):
                        #     wait = wx.BusyCursor()
                        #     newobj = WolfArrayMB(filename, mapviewer=self)
                        #     del wait
                        newobj = WolfArrayMB(filename, mapviewer=self)
                    else:
                        if which.lower() == 'array_crop':
                            newobj = WolfArray(filename, mapviewer=self, crop='newcrop')
                        else:
                            # with wx.BusyInfo(_('Importing array')):
                            #     wait = wx.BusyCursor()
                            #     newobj = WolfArray(filename, mapviewer=self)
                            #     del wait
                            newobj = WolfArray(filename, mapviewer=self)

                if newobj is not None:
                    if newobj.dx==0. or newobj.dy==0.:
                        dlg_pos = CropDialog(None)
                        dlg_pos.SetTitle(_('Choose informations'))

                        dlg_pos.ox.SetValue('99999.')
                        dlg_pos.oy.SetValue('99999.')

                        dlg_pos.ex.Hide()
                        dlg_pos.ey.Hide()

                        badvalues = True
                        while badvalues:
                            badvalues = False

                            ret = dlg_pos.ShowModal()
                            if ret == wx.ID_CANCEL:
                                newcrop.Destroy()
                                return -1
                            else:
                                cropini = [[float(dlg_pos.ox.Value), float(dlg_pos.ex.Value)],
                                                [float(dlg_pos.oy.Value), float(dlg_pos.ey.Value)]]
                                tmpdx = float(dlg_pos.dx.Value)
                                tmpdy = float(dlg_pos.dy.Value)

                            if tmpdx ==0. or tmpdy==0.:
                                badvalues = True

                        dlg_pos.Destroy()

                        newobj.dx = tmpdx
                        newobj.dy = tmpdy

                        # if newobj.SelectionData is not None:
                        #     newobj.SelectionData.dx = tmpdx
                        #     newobj.SelectionData.dy = tmpdy

                        if cropini[0][0] != 99999. and cropini[1][0]!=99999.:
                            newobj.origx = cropini[0][0]
                            newobj.origy = cropini[1][0]

            if newobj.epsg is None:
                logging.info(_('Array EPSG not defined -- Setting it to viewer EPSG ({})').format(self.epsg))
                newobj.epsg = self.epsg
            else:
                if newobj.epsg != self.epsg:
                    logging.error(_('Array EPSG ({}) different from viewer EPSG ({}) -- Reproject the array before adding it to the viewer').format(newobj.epsg, self.epsg))
                    if self._show_dialog_wx:
                        dlg = wx.MessageDialog(self, _('Array EPSG ({}) different from viewer EPSG ({}) -- Reproject the array before adding it to the viewer').format(newobj.epsg, self.epsg), style=wx.OK | wx.ICON_ERROR)
                        dlg.ShowModal()
                    return -1

            newobj.updatepalette(0)
            self.myarrays.append(newobj)
            newobj.change_gui(self)
            self.active_array = newobj
            self._set_active_bc()

        elif which.lower() == 'picture_collection':

            curdict = self.mypicturecollections
            curtree = self.myitemspictcollection

            if newobj is None:
                newobj = PictureCollection(parent=self, mapviewer=self)
                newobj.load_from_directory_with_shapefile(filename)

            curdict.append(newobj)
            self.active_picturecollection = newobj
            self.menu_pictcollection()

        elif which.lower() == 'array_tiles':

            res = wolfres2DGPU(filename, plotted=False)

            tilesmap = res._result_store._tile_packer.tile_indirection_map()
            if tilesmap is None:
                logging.warning(_('No tile map found in the simulation'))
                return

            header = header_wolf()
            res_header = res[0].get_header()

            header.origx = res_header.origx
            header.origy = res_header.origy
            header.dx = res_header.dx * 16.
            header.dy = res_header.dy * 16.
            header.nbx = tilesmap.shape[1]
            header.nby = tilesmap.shape[0]

            newobj_i = WolfArray(mapviewer=self, srcheader=header, idx = 'tils_i')
            newobj_j = WolfArray(mapviewer=self, srcheader=header, idx = 'tils_j')

            newobj_i.array = np.ma.asarray(tilesmap[:,:,0].T.astype(np.float32))
            newobj_j.array = np.ma.asarray(tilesmap[:,:,1].T.astype(np.float32))

            newobj_i.mask_data(0.)
            newobj_j.mask_data(0.)

            self.add_object('array', newobj=newobj_i, id=newobj_i.idx)
            self.add_object('array', newobj=newobj_j, id=newobj_j.idx)

            return


        elif which.lower() == 'imagestiles':

            curdict = self.myimagestiles
            curtree = self.myitemsvector

            if newobj is None:

                newobj = ImagesTiles('', parent=self, mapviewer=self)
                newobj.scan_dir(Path(filename))

            self.myimagestiles.append(newobj)
            self.active_imagestiles = newobj
            self.menu_imagestiles()

        elif which.lower() == 'bridges':
            curdict = self.myvectors
            curtree = self.myitemsvector

            if newobj is None:
                with wx.BusyInfo(_('Importing files')):
                    wait = wx.BusyCursor()
                    newobj = Bridges(filename, mapviewer=self)
                    del wait

            self.myvectors.append(newobj)

            self.active_bridges = newobj

            self.menu_bridges()

        elif which.lower() == 'weirs':
            curdict = self.myvectors
            curtree = self.myitemsvector

            if newobj is None:
                with wx.BusyInfo(_('Importing files')):
                    wait = wx.BusyCursor()
                    newobj = Weirs(filename, mapviewer=self)
                    del wait
            self.myvectors.append(newobj)

            self.active_weirs = newobj

            self.menu_weirs()

        elif which.lower() in ['tiles', 'tilescomp']:
            curdict = self.mytiles
            curtree = self.myitemsvector

            if newobj is None:

                file = wx.DirDialog(self, "Choose directory containing data")
                if file.ShowModal() == wx.ID_CANCEL:
                    file.Destroy()
                    return -1
                else:
                    # récuparétaion du nom de fichier avec chemin d'accès
                    dirname = file.GetPath()
                    file.Destroy()

                if which.lower() == 'tilescomp':
                    file = wx.DirDialog(self, "Choose directory containing comparison data")
                    if file.ShowModal() == wx.ID_CANCEL:
                        file.Destroy()
                        return -1
                    else:
                        # récuparétaion du nom de fichier avec chemin d'accès
                        dirname_comp = file.GetPath()
                        file.Destroy()

                with wx.BusyInfo(_('Importing files')):
                    wait = wx.BusyCursor()
                    newobj = Tiles(filename, parent=self, linked_data_dir=dirname, mapviewer=self)
                    del wait

                    if which.lower() == 'tilescomp':
                        newobj.linked_data_dir_comp = dirname_comp

            self.mytiles.append(newobj)
            self.active_tile = newobj
            self.menu_tiles()

        elif which.lower() == 'array_xyz':

            curdict = self.myarrays
            curtree = self.myitemsarray

            msg = wx.MessageDialog(self, _('Do you want to crop the data?'), style=wx.YES_NO | wx.YES_DEFAULT)
            ret = msg.ShowModal()
            msg.Destroy()

            if ret == wx.ID_YES:

                newcrop = CropDialog(None)

                badvalues = True
                while badvalues:
                    badvalues = False

                    ret = newcrop.ShowModal()
                    if ret == wx.ID_CANCEL:
                        newcrop.Destroy()
                        return -1
                    else:
                        cropini = [[float(newcrop.ox.Value), float(newcrop.ex.Value)],
                                   [float(newcrop.oy.Value), float(newcrop.ey.Value)]]
                        tmpdx = float(newcrop.dx.Value)
                        tmpdy = float(newcrop.dy.Value)

                newcrop.Destroy()

                myxyz = xyz_scandir(filename, cropini)

                myhead = newcrop.get_header()
                # if min(myhead.dx, myhead.dy) != 1.:
                #     myhead.nbx = int(myhead.nbx * myhead.dx)
                #     myhead.nby = int(myhead.nby * myhead.dy)
                #     myhead.dx = 1.
                #     myhead.dy = 1.

            else:

                dlg = wx.TextEntryDialog(self,_('Spatial step size (assuming dx == dy) ?'), value='1')
                ret=dlg.ShowModal()

                if ret == wx.ID_CANCEL:
                    dlg.Destroy()
                    return -1

                tmpdx = float(dlg.GetValue())
                dlg.Destroy()

                dy = dx

                myxyz = xyz_scandir(filename, None)
                myhead = header_wolf()

                myhead.origx = np.min(myxyz[:, 0]) - dx/2.
                myhead.origy = np.min(myxyz[:, 1]) - dy/2.

                myhead.dx = dx
                myhead.dy = dy

                myhead.nbx = int(np.max(myxyz[:, 0]) - myhead.origx) + 1
                myhead.nby = int(np.max(myxyz[:, 1]) - myhead.origy) + 1

            if len(myxyz) == 0:
                return -1

            newobj = WolfArray()

            newobj.init_from_header(myhead)
            newobj.nullvalue = -99999.
            newobj.array.data[:, :] = -99999.

            newobj.fillin_from_xyz(myxyz)

            newobj.mask_data(newobj.nullvalue)

            newobj.change_gui(self)
            newobj.updatepalette(0)
            self.myarrays.append(newobj)
            self.active_array = newobj
            self._set_active_bc()

        elif which.lower() == 'array_lidar_first' or which.lower() == 'array_lidar_second':

            curdict = self.myarrays
            curtree = self.myitemsarray

            newcrop = CropDialog(None)

            badvalues = True
            while badvalues:
                badvalues = False

                ret = newcrop.ShowModal()
                if ret == wx.ID_CANCEL:
                    newcrop.Destroy()
                    return -1
                else:
                    cropini = [[float(newcrop.ox.Value), float(newcrop.ex.Value)],
                               [float(newcrop.oy.Value), float(newcrop.ey.Value)]]
                    tmpdx = float(newcrop.dx.Value)
                    tmpdy = float(newcrop.dy.Value)

            newcrop.Destroy()

            first, sec = Lidar2002.lidar_scandir(filename, cropini)

            if which.lower() == 'array_lidar_first':
                if len(first) == 0:
                    return -1

                newobj = Lidar2002.create_wolfarray(first, bounds=cropini)

                if min(tmpdx, tmpdy) != 1.:
                    newobj.rebin(min(tmpdx, tmpdy))

                newobj.change_gui(self)
                newobj.updatepalette(0)
                self.myarrays.append(newobj)
                self.active_array = newobj
                self._set_active_bc()

                id = 'lidar2002_firstecho'
            else:
                if len(sec) == 0:
                    return -1
                newobj = Lidar2002.create_wolfarray(sec, bounds=cropini)
                if min(tmpdx, tmpdy) != 1.:
                    newobj.rebin(min(tmpdx, tmpdy))

                newobj.change_gui(self)
                newobj.updatepalette(0)
                self.myarrays.append(newobj)
                self.active_array = newobj
                self._set_active_bc()
                id = 'lidar2002_secondecho'

        elif which.lower() == 'res2d':

            curdict = self.myres2D
            curtree = self.myitemsres2d

            if newobj is None:
                with wx.BusyInfo(_('Importing 2D model')):
                    wait = wx.BusyCursor()
                    newobj = Wolfresults_2D(filename, mapviewer=self)
                    del wait

            newobj.get_nbresults(True)
            newobj.updatepalette()
            self.myres2D.append(newobj)
            self.active_res2d = newobj

        elif which.lower() == 'res2d_gpu':

            curdict = self.myres2D
            curtree = self.myitemsres2d

            if newobj is None:
                # with wx.BusyInfo(_('Importing 2D GPU model')):
                #     wait = wx.BusyCursor()
                #     newobj = wolfres2DGPU(filename, mapviewer=self)
                #     del wait
                newobj = wolfres2DGPU(filename, mapviewer=self)

            if newobj is None:
                logging.warning(_('Error while importing GPU results'))
                return -1

            newobj.get_nbresults(True)
            newobj.read_oneresult(-1)
            newobj.updatepalette()
            self.myres2D.append(newobj)
            self.active_res2d = newobj

        elif which.lower() == 'vector':
            curdict = self.myvectors
            curtree = self.myitemsvector
            if newobj is None:
                with wx.BusyInfo(_('Importing file')):
                    wait = wx.BusyCursor()
                    newobj = Zones(filename, parent=self)
                    del wait
            self.myvectors.append(newobj)

        elif which.lower() == 'cross_sections':

            curdict = self.myvectors
            curtree = self.myitemsvector

            if newobj is None:

                dlg = wx.MessageDialog(None, 'Load LAZ data?', style=wx.YES_NO | wx.NO_DEFAULT)
                ret = dlg.ShowModal()
                dlg.Destroy()
                dirlaz = ''

                if ret == wx.ID_YES:
                    if self.mylazgrid is not None:
                        dlg = wx.MessageDialog(None, 'Gridded LAZ data exist - use them ?', style=wx.YES_NO | wx.YES_DEFAULT)
                        ret = dlg.ShowModal()
                        dlg.Destroy()

                        if ret == wx.ID_YES:
                            dirlaz = self.mylazgrid
                        else:
                            dlg = wx.DirDialog(None, 'If exist, where are the LAZ data?')
                            ret = dlg.ShowModal()
                            if ret == wx.ID_OK:
                                dirlaz = dlg.GetPath()
                    else:
                        dlg = wx.DirDialog(None, 'If exist, where are the LAZ data?')
                        ret = dlg.ShowModal()
                        if ret == wx.ID_OK:
                            dirlaz = dlg.GetPath()

                with wx.BusyInfo(_('Importing cross sections')):
                    wait = wx.BusyCursor()
                    if curfilter == 1:  # txt 2022
                        newobj = crosssections(filename, format='2022', dirlaz=dirlaz, mapviewer=self)
                    if curfilter == 0:  # vecz
                        newobj = crosssections(filename, format='vecz', dirlaz=dirlaz, mapviewer=self)
                    elif curfilter == 2:  # sxy
                        newobj = crosssections(filename, format='sxy', dirlaz=dirlaz, mapviewer=self)
                    elif curfilter == 3:  # txt 2000
                        newobj = crosssections(filename, format='2000', dirlaz=dirlaz, mapviewer=self)
                    elif curfilter == 4:  # xlsx
                        newobj = crosssections(filename, format='2025_xlsx', dirlaz=dirlaz, mapviewer=self)
                    else:
                        newobj = crosssections(filename, format='2000', dirlaz=dirlaz, mapviewer=self)

                    del wait
            self.myvectors.append(newobj)
            newobj.mapviewer = self

        elif which.lower() =='laz':
            curdict = self.mylazdata
            curtree = self.myitemslaz

            if newobj is None:
                newobj = Wolf_LAZ_Data(mapviewer=self)
                newobj.from_file(filename)

            self.mylazdata.append(newobj)
            self.active_laz = newobj

            newobj.set_mapviewer(self)

        elif which.lower() == 'cloud':

            curdict = self.myclouds
            curtree = self.myitemscloud
            if newobj is None:

                loadhead = False
                if not filename.endswith('.dxf') and not filename.endswith('.shp'):
                    with open(filename,'r') as f:
                        text=f.read().splitlines()
                        tmphead=''
                        for i in range(min(4, len(text))):
                            tmphead += text[i].replace('\t','\\t') +'\n'

                    dlg = wx.MessageDialog(None,_('Is there a file header (one upper line containing column names)?') + '\n\n' + tmphead,style=wx.YES_NO|wx.NO_DEFAULT)
                    ret=dlg.ShowModal()

                    if ret == wx.ID_YES:
                        loadhead = True

                    newobj = cloud_vertices(filename, header=loadhead, mapviewer=self)

                elif filename.endswith('.dxf'):
                    types = ['POLYLINE','LWPOLYLINE','LINE', 'MTEXT', 'INSERT']

                    dlg = wx.MultiChoiceDialog(None, _('Choose the types of entities to import'), _('Choose entities'), types)
                    dlg.SetSelections = [3,4]

                    ret = dlg.ShowModal()
                    if ret == wx.ID_CANCEL:
                        dlg.Destroy()
                        return -1

                    types = [types[i] for i in dlg.GetSelections()]
                    dlg.Destroy()

                    newobj = cloud_vertices(filename, header=loadhead, mapviewer=self, dxf_imported_elts=types)

                elif filename.endswith('.shp'):

                    types = None

                    if Path(filename).stem == 'Vesdre_Bridges':
                        # We need to import the bridges from the clogging study/database
                        data = gpd.read_file(filename)

                        #filter 'Clogging' == Yes
                        clogged = data[data['Clogging'] == 'Yes']
                        unclogged = data[data['Clogging'] == 'No']
                        notsure = data[data['Clogging'] == 'No information']

                        from tempfile import TemporaryDirectory
                        with TemporaryDirectory() as tmpdirname:
                            clogged.to_file(tmpdirname + '/clogged.shp')
                            unclogged.to_file(tmpdirname + '/unclogged.shp')
                            notsure.to_file(tmpdirname + '/notsure.shp')

                            newobj = cloud_vertices(tmpdirname + '/unclogged.shp', header=loadhead, mapviewer=self, idx='unclogged')
                            self.myclouds.append(newobj)
                            newobj.set_mapviewer(self)
                            newobj.myprop.color = (0,255,0)
                            newobj.myprop.size = 10

                            myitem = self.treelist.AppendItem(curtree, newobj.idx, data=newobj)
                            self.treelist.CheckItem(myitem)
                            self.treelist.CheckItem(self.treelist.GetItemParent(myitem))
                            newobj.check_plot()

                            newobj = cloud_vertices(tmpdirname + '/notsure.shp', header=loadhead, mapviewer=self, idx='notsure')
                            self.myclouds.append(newobj)
                            newobj.set_mapviewer(self)
                            newobj.myprop.color = (0,0,255)
                            newobj.myprop.size = 10

                            myitem = self.treelist.AppendItem(curtree, newobj.idx, data=newobj)
                            self.treelist.CheckItem(myitem)
                            self.treelist.CheckItem(self.treelist.GetItemParent(myitem))
                            newobj.check_plot()

                            newobj = cloud_vertices(tmpdirname + '/clogged.shp', header=loadhead, mapviewer=self, idx='clogged')
                            newobj.myprop.color = (255,0,0)
                            newobj.myprop.size = 15
                            id = 'clogged'

                    else:
                        newobj = cloud_vertices(filename, header=loadhead, mapviewer=self)

            self.myclouds.append(newobj)
            self.active_cloud = newobj

            newobj.set_mapviewer(self)

            self.create_cloud_menu()

        elif which.lower() == 'triangulation':

            curdict = self.mytri
            curtree = self.myitemstri
            if newobj is None:
                with wx.BusyInfo(_('Importing triangulation')):
                    wait = wx.BusyCursor()
                    newobj = Triangulation(filename, mapviewer=self)
                    del wait

            self.mytri.append(newobj)
            self.active_tri = newobj

            self.create_triangles_menu()

        elif which.lower() == 'other':

            if not newobj is None:
                curdict = self.myothers
                curtree = self.myitemsothers
                self.myothers.append(newobj)
                newobj.mapviewer = self
            else:
                logging.warning(_('No object to add in "Other" category -- Please provide an object to add or check your code'))

        elif which.lower() == 'views':

            if newobj is None:
                newobj = WolfViews(plotted=ToCheck, mapviewer=self)
                newobj.read_from_file(filename)

            curdict = self.myviews
            curtree = self.myitemsviews
            self.myviews.append(newobj)

        elif which.lower() == 'wmsback':

            if not newobj is None:
                curdict = self.mywmsback
                curtree = self.myitemswmsback
                self.mywmsback.append(newobj)
            else:
                logging.warning(_('No object to add in "WMS background" category -- Please provide an object to add or check your code'))

        elif which.lower() == 'wmsfore':

            if not newobj is None:
                curdict = self.mywmsfore
                curtree = self.myitemswmsfore
                self.mywmsfore.append(newobj)
            else:
                logging.warning(_('No object to add in "WMS foreground" category -- Please provide an object to add or check your code'))

        elif which.lower() == 'particlesystem':

            curdict = self.mypartsystems
            curtree = self.myitemsps
            if newobj is None:
                    newobj = Particle_system(mapviewer=self)
                    newobj.load(filename)

            self.mypartsystems.append(newobj)
            self.active_particle_system = newobj

        elif which.lower() == 'drowning':

            curdict = self.mydrownings
            curtree = self.myitemsdrowning

            self.mydrownings.append(newobj)
            self.active_drowning = newobj


        elif which.lower() == 'dike':

            if not WOLFPYDIKE_AVAILABLE:
                logging.error('WolfPyDike module not available - cannot add dike')
                return -1

            curdict = self.mydikes
            curtree = self.myitemsdike

            self.mydikes.append(newobj)
            self.active_dike = newobj

        elif which.lower() == 'injector':

            curdict = self.myinjectors
            curtree = self.myitemsinjector

            self.myinjectors.append(newobj)
            self.active_injector = newobj

        # ID chooser
        if id == '':
            dlg = wx.TextEntryDialog(self, 'ID ? (case insensitive)', 'Choose an identifier', '')
            if filename != '':
                dlg.SetValue((Path(filename).with_suffix('')).name)
            else:
                dlg.SetValue('')

            endid = 1
            # ids = self.get_list_keys(None, checked_state=None) #[cur.idx for cur in curdict]
            while id.lower() in all_ids or id =='':
                if dlg.ShowModal() == wx.ID_OK:
                    id = dlg.GetValue()
                    if id =='':
                        id = str(endid).zfill(3)
                    endid += 1
            dlg.Destroy()

        # ids = [cur.idx for cur in curdict]
        # if id.lower() in ids:
        #     endid = 1
        #     while (id + str(endid).zfill(3)).lower() in ids:
        #         endid += 1
        #     id = id + str(endid).zfill(3)

        # all_ids = self.get_list_keys(None, checked_state=None)
        if id.lower() in all_ids:
            endid = 1
            while (id + str(endid).zfill(3)).lower() in all_ids:
                endid += 1
            id = id + str(endid).zfill(3)

        newobj.idx = id.lower()

        if curtree is not None:
            myitem = self.treelist.AppendItem(curtree, id, data=newobj)

            if ToCheck:
                self.treelist.CheckItem(myitem)
                self.treelist.CheckItem(self.treelist.GetItemParent(myitem))

                newobj.check_plot()
        else:
            logging.info(f'No tree item for this object {newobj.idx}')

        # curdict[id.lower()] = newobj
        if filename != '':
            newobj._filename_vector = Path(filename).name.lower() # FIXME useful ??
        newobj.checked = ToCheck

        if isinstance(newobj,crosssections):
            self.add_object('cloud',newobj=newobj.cloud    ,id=newobj.idx+'_intersect',ToCheck=False)
            self.add_object('cloud',newobj=newobj.cloud_all,id=newobj.idx+'_all',      ToCheck=False)

        elif type(newobj) == WolfArray:
            if self.active_cs is None:
                self.active_cs = self.get_cross_sections()

        return 0

    def replace_object(self, id: str, newobj, drawing_type: draw_type = None):
        """ Replace an object in the list of objects of type drawing_type """

        if drawing_type is None:
            for curdict in draw_type:
                keys = self.get_list_keys(curdict, checked_state=None)
                if id.lower() in keys:
                    # The object exists in the current dictionary
                    obj = self.get_obj_from_id(id, drawing_type=curdict)
                    obj.reset_listogl()
                    # Searching the object in all lists
                    if obj is not None:
                        curlist = self._get_list(drawing_type=curdict)
                        if obj in curlist:
                            pos = curlist.index(obj)
                            if isinstance(newobj, curlist[pos].__class__):
                                # Updating the tree item
                                self.treelist.SetItemData(self.get_treeitem_from_obj(obj), newobj)
                                curlist[pos] = newobj
                                newobj.idx = id.lower()
                            else:
                                logging.error(f'Cannot replace {id} with {newobj.idx} - Different type of object')
                        else:
                            logging.error(f'Object {id} not found in list')
                    else:
                        logging.error(f'Object {id} not found in list')
                else:
                    logging.error(f'Object {id} not found in dictionary {curdict}')
        else:
            keys = self.get_list_keys(drawing_type, checked_state=None)
            if id.lower() in keys:
                # The object exists in the current dictionary
                obj = self.get_obj_from_id(id, drawing_type=drawing_type)
                obj.reset_listogl()
                # Searching the object in all lists
                if obj is not None:
                    curlist = self._get_list(drawing_type=drawing_type)
                    if obj in curlist:
                        pos = curlist.index(obj)
                        if isinstance(newobj, curlist[pos].__class__):
                            # Updating the tree item
                            self.treelist.SetItemData(self.get_treeitem_from_obj(obj), newobj)
                            curlist[pos] = newobj
                            newobj.idx = id.lower()
                        else:
                            logging.error(f'Cannot replace {id} with {newobj.idx} - Different type of object')
                    else:
                        logging.error(f'Object {id} not found in list')
                else:
                    logging.error(f'Object {id} not found in list')
            else:
                logging.error(f'Object {id} not found in dictionary {drawing_type}')

        obj = self.get_obj_from_id(id, drawing_type=drawing_type)
        obj_from_tree = self.get_obj_from_treeitem(self.get_treeitem_from_id(id, drawing_type=drawing_type))
        if obj is not None and obj_from_tree is not None:
            if obj is obj_from_tree:
                logging.debug(f'Object {id} replaced successfully in the list and tree item')
            else:
                logging.error(f'Object {id} replaced in the list but not in the tree item - {obj} != {obj_from_tree}')
        else:
            logging.error(f'Object {id} not found in the list or tree item after replacement')

    def get_obj_from_treeitem(self, treeitem):
        """ Find the object associated with treeitem """

        return self.treelist.GetItemData(treeitem)


    def get_treeitem_from_id(self, id: str, drawing_type: draw_type = None):
        """ Find the tree item associated with id """

        obj = self.get_obj_from_id(id, drawing_type=drawing_type)
        if obj is not None:
            return self.get_treeitem_from_obj(obj)
        return None

    def get_treeitem_from_obj(self, obj):
        """ Find the tree item associated with obj.

        Alias for "gettreeitem".
        """

        return self.gettreeitem(obj)

    def getobj_from_id(self, id: str, drawing_type: draw_type = None):
        """ Find the object associated with id """

        if drawing_type is None:
            for curdict in draw_type:
                keys = self.get_list_keys(curdict, checked_state=None)
                if id.lower() in keys:
                    try:
                        idx = keys.index(id.lower())
                        return self.get_list_objects(curdict, checked_state=None)[idx]
                    except:
                        return None
        else:
            keys = self.get_list_keys(drawing_type, checked_state=None)
            if id.lower() in keys:
                try:
                    idx = keys.index(id.lower())
                    return self.get_list_objects(drawing_type, checked_state=None)[idx]
                except:
                    return None

    def get_obj_from_id(self, id: str, drawing_type: draw_type = None):
        """ Find the object associated with id in a specifid drawtype

        If you want to search in all drawtypes, use getobj_from_id instead.

        :param id: str : id of the object
        :param drawtype: draw_type : type of object to search

        """

        keys = self.get_list_keys(drawing_type, checked_state=None)
        if id.lower() in keys:
            try:
                idx = keys.index(id.lower())
                return self.get_list_objects(drawing_type, checked_state=None)[idx]
            except:
                return None

    def _get_list(self, drawing_type:draw_type = None):
        """ return the list of objects of type drawing_type """

        # ARRAYS = 'arrays'
        # BRIDGES= 'bridges'
        # WEIRS = 'weirs'
        # VECTORS = 'vectors'
        # CLOUD = 'clouds'
        # TRIANGULATION = 'triangulations'
        # PARTICLE_SYSTEM = 'particle systems'
        # CROSS_SECTIONS = 'cross_sections'
        # OTHER = 'others'
        # VIEWS = 'views'
        # RES2D = 'wolf2d'
        # WMSBACK = 'wms-background'
        # WMSFORE = 'wms-foreground'
        # PICTURE_COLLECTION = 'picture collections'

        if drawing_type is None:
            # return all_lists
            return self.myarrays + self.myvectors + self.myclouds + self.mytri + self.mypartsystems + self.myothers + self.myviews + self.myres2D + self.mydikes + self.mydrownings + self.myinjectors + self.mypicturecollections

        if drawing_type == draw_type.ARRAYS:
            return self.myarrays
        elif drawing_type == draw_type.VECTORS or drawing_type == draw_type.BRIDGES or drawing_type == draw_type.WEIRS or drawing_type == draw_type.CROSS_SECTIONS :
            return self.myvectors
        elif drawing_type == draw_type.TILES:
            return self.mytiles
        elif drawing_type == draw_type.CLOUD:
            return self.myclouds
        elif drawing_type == draw_type.TRIANGULATION:
            return self.mytri
        elif drawing_type == draw_type.RES2D:
            return self.myres2D
        elif drawing_type == draw_type.PARTICLE_SYSTEM:
            return self.mypartsystems
        elif drawing_type == draw_type.OTHER:
            return self.myothers
        elif drawing_type == draw_type.VIEWS:
            return self.myviews
        elif drawing_type == draw_type.WMSBACK:
            return self.mywmsback
        elif drawing_type == draw_type.WMSFORE:
            return self.mywmsfore
        elif drawing_type == draw_type.IMAGESTILES:
            return self.myimagestiles
        elif drawing_type == draw_type.LAZ:
            return self.mylazdata
        elif drawing_type == draw_type.DROWNING:
            return self.mydrownings
        elif drawing_type == draw_type.DIKE:
            return self.mydikes
        elif drawing_type == draw_type.INJECTOR:
            return self.myinjectors
        elif drawing_type == draw_type.PICTURECOLLECTION:
            return self.mypicturecollections
        else:
            logging.error('Unknown drawing type : ' + drawing_type)
            return None

    def get_list_keys(self, drawing_type:draw_type = None, checked_state:bool=True):
        """ Create a list of keys of type draw_type.

        Return a list of keys (idx) in LOWER CASE of objects of type draw_type.

        :param drawing_type: type of object to search - If None, return all objects
        :param checked_state: if True/False, return only keys of objects that are plotted or not. None return all objects.
        """

        if checked_state is None:
            return [curobj.idx for curobj in self._get_list(drawing_type)]
        else:
            return [curobj.idx for curobj in self._get_list(drawing_type) if curobj.plotted == checked_state]

    def get_list_ids(self, drawing_type:draw_type = None, checked_state:bool=True):
        """ Alias for get_list_keys """

        return self.get_list_keys(drawing_type, checked_state)

    def get_list_objects(self, drawing_type:draw_type = None, checked_state:bool=True):
        """ Create a list of objects of type draw_type.

        Return a list of keys (idx) in LOWER CASE of objects of type draw_type.

        :param drawing_type: type of object to search -- If None, return all objects.
        :param checked_state: if True/False, return only objects that are plotted or not. None return all objects.
        """

        if checked_state is None:
            return [curobj for curobj in self._get_list(drawing_type)]
        else:
            return [curobj for curobj in self._get_list(drawing_type) if curobj.plotted == checked_state]

    def single_choice_key(self, draw_type:draw_type, checked_state:bool=True, message:str=_('Make a choice'), title:str=_('Choice')):
        """ Create wx dialog to choose a key object of type draw_type """

        keys = self.get_list_keys(draw_type, checked_state)
        dlg = wx.SingleChoiceDialog(None, message, title, keys, style=wx.CHOICEDLG_STYLE)
        ret = dlg.ShowModal()

        if ret != wx.ID_OK:
            dlg.Destroy()
            return None

        idx = dlg.GetSelection()
        dlg.Destroy()

        return keys[idx]

    def single_choice_object(self, draw_type:draw_type, checked_state:bool=True, message:str=_('Make a choice'), title:str=_('Choice')):
        """ Create wx dialog to choose an object of type draw_type """

        keys = self.get_list_keys(draw_type, checked_state)
        obj = self.get_list_objects
        dlg = wx.SingleChoiceDialog(None, message, title, keys, style=wx.CHOICEDLG_STYLE)
        ret = dlg.ShowModal()

        if ret != wx.ID_OK:
            dlg.Destroy()
            return None

        idx = dlg.GetSelection()
        dlg.Destroy()

        return obj[idx]

    def multiple_choice_key(self, draw_type:draw_type, checked_state:bool=True, message:str=_('Make a choice'), title:str=_('Choice')):
        """ Create wx dialog to choose multiple keys object of type draw_type """

        keys = self.get_list_keys(draw_type, checked_state)
        dlg = wx.MultiChoiceDialog(None, message, title, keys, style=wx.CHOICEDLG_STYLE)
        ret = dlg.ShowModal()

        if ret != wx.ID_OK:
            dlg.Destroy()
            return None

        idx = dlg.GetSelections()
        dlg.Destroy()

        return [keys[i] for i in idx]

    def multiple_choice_object(self, draw_type:draw_type, checked_state:bool=True, message:str=_('Make a choice'), title:str=_('Choice')):
        """ Create wx dialog to choose multiple objects of type draw_type """

        keys = self.get_list_keys(draw_type, checked_state)
        obj = self.get_list_objects
        dlg = wx.MultiChoiceDialog(None, message, title, keys, style=wx.CHOICEDLG_STYLE)
        ret = dlg.ShowModal()

        if ret != wx.ID_OK:
            dlg.Destroy()
            return None

        idx = dlg.GetSelections()
        dlg.Destroy()

        return [obj[i] for i in idx]

    def iterator_over_objects(self, drawing_type:draw_type, checked_state:bool=True):
        """ Create iterator over objects of type draw_type """

        for obj in self.get_list_objects(drawing_type, checked_state):
            yield obj

    def gettreeitem(self, obj):
        """ Find the tree item associated with obj """

        up = self.treelist.GetFirstItem()
        updata = self.treelist.GetItemData(up)

        while updata is not obj:
            up = self.treelist.GetNextItem(up)
            updata = self.treelist.GetItemData(up)

        return up

    def removeobj(self):
        """Remove selected item from general tree"""

        if self.selected_treeitem is None:
            return

        id = self.treelist.GetItemText(self.selected_treeitem).lower()

        self.removeobj_from_id(id)

    def checkuncheckobj(self):
        """ Check/uncheck selected item from general tree """

        if self.selected_treeitem is None:
            return

        id = self.treelist.GetItemText(self.selected_treeitem).lower()

        current_check = self.treelist.GetCheckedState(self.selected_treeitem)
        myobj = self.getobj_from_id(id)
        if myobj is not None:
            if current_check == 0:
                self.treelist.CheckItem(self.selected_treeitem)
                myobj.check_plot()
            else:
                self.treelist.CheckItem(self.selected_treeitem, False)
                myobj.uncheck_plot()

    def removeobj_from_id(self, id:str, draw_type:draw_type = None):
        """ Remove object from id """

        myobj = self.getobj_from_id(id)
        if myobj is not None:
            self.treelist.DeleteItem(self.gettreeitem(myobj))

            for curlist in self.all_lists:
                if myobj in curlist:
                    curlist.pop(curlist.index(myobj))

            myobj.hide_properties()

            if myobj is self.active_array:
                self.active_array = None
                self.set_label_selecteditem('')
            elif myobj is self.active_res2d:
                self.active_res2d = None
                self.set_label_selecteditem('')
            elif myobj is self.active_tri:
                self.active_tri = None
                self.set_label_selecteditem('')
            elif myobj is self.active_particle_system:
                self.active_particle_system = None
                self.set_label_selecteditem('')
            elif myobj is self.active_cloud:
                self.active_cloud = None
                self.set_label_selecteditem('')
            elif myobj is self.active_cs:
                self.active_cs = None
                self.set_label_selecteditem('')
            elif myobj is self.active_tile:
                self.active_tile = None
                self.set_label_selecteditem('')
            elif myobj is self.active_bc:
                self.active_bc = None
                self.set_label_selecteditem('')
            elif myobj is self.active_vector:
                self.active_vector = None
                self.set_label_selecteditem('')
            elif myobj is self.active_view:
                self.active_view = None
                self.set_label_selecteditem('')
            elif myobj is self.active_zone:
                self.active_zone = None
                self.set_label_selecteditem('')
            elif myobj is self.active_zones:
                self.active_zones = None
                self.set_label_selecteditem('')
            elif myobj is self.active_drowning:
                self.active_drowning = None
                self.set_label_selecteditem('')
            elif myobj is self.active_dike:
                self.active_dike = None
                self.set_label_selecteditem('')
            elif myobj is self.active_injector:
                self.active_injector = None
                self.set_label_selecteditem('')

    def upobj(self):
        """Up selected item into general tree"""

        if self.selected_treeitem is None:
            return

        id:str
        id = self.treelist.GetItemText(self.selected_treeitem).lower()
        myobj = self.getobj_from_id(id)
        ischecked = self.treelist.GetCheckedState(self.selected_treeitem)

        assert self.selected_object is myobj, 'selected_object is not myobj'

        if myobj is not None:

            down = self.treelist.GetNextItem(self.selected_treeitem)
            up = self.treelist.GetFirstItem()
            up2 = up

            while self.treelist.GetNextItem(up) != self.selected_treeitem:
                up2= up
                up = self.treelist.GetNextItem(up)

            parent = self.treelist.GetItemParent(self.selected_treeitem)
            parentup = self.treelist.GetItemParent(up)
            parentup2 = self.treelist.GetItemParent(up2)

            if parent == parentup2:
                # up n'est pas le premier élément de la liste
                myitem = self.treelist.InsertItem(parent,up2,id,data=myobj)
                self.treelist.CheckItem(myitem,ischecked)
            elif parentup == parent:
                # up est le premier élément de la liste
                myitem = self.treelist.PrependItem(parent,id,data=myobj)
                self.treelist.CheckItem(myitem,ischecked)
            else:
                # nothing to do
                return

            self.treelist.DeleteItem(self.selected_treeitem)
            self.selected_treeitem = myitem

            # mouvement dans les listes pour garder l'ordre identique à l'arbre
            for curlist in self.all_lists:
                if myobj in curlist:
                    idx = curlist.index(myobj)
                    if idx>0:
                        curlist.pop(idx)
                        curlist.insert(idx-1,myobj)

            self.Refresh()

    def downobj(self):
        """Down selected item into general tree"""

        if self.selected_treeitem is None:
            return

        id = self.treelist.GetItemText(self.selected_treeitem).lower()
        myobj = self.getobj_from_id(id)
        ischecked = self.treelist.GetCheckedState(self.selected_treeitem)

        if myobj is not None:

            down = self.treelist.GetNextItem(self.selected_treeitem)
            down2 = self.treelist.GetNextItem(down)

            parent = self.treelist.GetItemParent(self.selected_treeitem)
            parentdown = self.treelist.GetItemParent(down)
            parentdown2 = self.treelist.GetItemParent(down2)

            if parent == parentdown:
                # on n'est pas sur le dernoier élément
                myitem = self.treelist.InsertItem(parent,down,id,data=myobj)
                self.treelist.CheckItem(myitem,ischecked)
            else:
                # nothing to do
                return

            self.treelist.DeleteItem(self.selected_treeitem)
            self.selected_treeitem = myitem

            for curlist in self.all_lists:
                if myobj in curlist:
                    if len(curlist)>1:
                        idx = curlist.index(myobj)
                        if idx == len(curlist)-1:
                            # dernier --> rien à faire
                            pass
                        elif idx==len(curlist)-2:
                            # avant-dernier --> passage en dernier
                            curlist.append(myobj)
                            curlist.pop(idx)
                        elif idx<len(curlist)-2:
                            curlist.insert(idx+2,myobj)
                            curlist.pop(idx)

            self.Refresh()

    def OnShowPopup(self, event):
        pos = event.GetPosition()
        if pos == (-1, -1):
            width, height = self.GetSize()
            pos = (width / 2., height / 2.)
        # else:
        #     pos = pos - self.GetPosition()

        self.PopupMenu(self.popupmenu, pos)

    def OnPopupItemSelected(self, event):
        """ Action to do when an item is selected in the popup menu """

        item = self.popupmenu.FindItemById(event.GetId())
        text = item.ItemLabel

        if text == _('Save'):
            if self.selected_object is not None:
                if issubclass(type(self.selected_object), WolfArray):
                    self.selected_object.write_all()
                elif type(self.selected_object) is Zones:
                    self.selected_object.saveas()
                elif type(self.selected_object) in [Bridge, Weir]:
                    self.selected_object.saveas()
                elif type(self.selected_object) is Triangulation:
                    self.selected_object.saveas()
                elif isinstance(self.selected_object, Particle_system):
                    self.selected_object.save()
                elif isinstance(self.selected_object, Drowning_victim_Viewer):
                    self.selected_object.save()
                elif isinstance(self.selected_object, DikeWolf):
                    self.selected_object.save()
                elif isinstance(self.selected_object, InjectorDike):
                    self.selected_object.save()

        elif text==_('Up'):
            self.upobj()

        elif text == _('Down'):
            self.downobj()

        elif text == _('Delete'):
            self.removeobj()

        elif text == ('Check/Uncheck'):
            self.checkuncheckobj()

        elif text == _('Rename'):
            #Modification du nom de l'objet sélectionné

            if self.selected_object is not None:
                #récupération de l'id courant
                all_ids = self.get_list_ids(checked_state=None)

                label = self.selected_object.idx

                all_ids.remove(label.lower())

                dlg = wx.TextEntryDialog(self, message=_('Chose a new label :'), value=label)
                ret=dlg.ShowModal()

                if ret == wx.ID_OK:
                    newlab = dlg.GetValue()

                    if newlab.lower() in all_ids:
                        wx.MessageBox(_('This label already exists. Please choose another one.'), _('Error'), wx.OK | wx.ICON_ERROR)
                        dlg.Destroy()
                        return

                    #MAJ de l'id dans l'objet
                    self.selected_object.idx = newlab
                    #MAJ de l'arbre
                    self.treelist.SetItemText(self.selected_treeitem, newlab)

                    if self.get_label_selecteditem() == _("Active : ") + label:
                        self.set_label_selecteditem(_("Active : ") + newlab)

                dlg.Destroy()

        elif text ==  _('Duplicate'):

            # Duplication de l'objet sélectionné
            if self.selected_object is not None:
                #récupération de l'id courant
                label = self.selected_object.idx + '_copy'
                dlg = wx.TextEntryDialog(self, message=_('Chose a label for the copy:'), value=label)
                ret=dlg.ShowModal()

                if ret != wx.ID_OK:
                    dlg.Destroy()
                    return

                newlab = dlg.GetValue()
                dlg.Destroy()

                if isinstance(self.selected_object, WolfArray) and (not type(self.selected_object) in [WolfArrayMB, WolfArrayMNAP]):

                    curtype = self.selected_object.dtype

                    if curtype == np.float64:
                        curtype = 'float64'
                    elif curtype == np.float32:
                        curtype = 'float32'
                    elif curtype == np.int32:
                        curtype = 'int32'
                    elif curtype == np.int16:
                        curtype = 'int16'
                    elif curtype == np.int8:
                        curtype = 'int8'

                    dlg = wx.MessageDialog(None, _('The type of the data is {}.\nDo you want to change this type?'.format(curtype)), style=wx.YES_NO | wx.NO_DEFAULT)
                    ret = dlg.ShowModal()
                    dlg.Destroy()

                    if ret == wx.ID_YES:
                        dlg = wx.SingleChoiceDialog(None, _('Choose a type'), _('Type'), ['float32','float64','int32','int16','int8'], style=wx.CHOICEDLG_STYLE)
                        ret = dlg.ShowModal()

                        if ret != wx.ID_OK:
                            dlg.Destroy()
                            return

                        idx = dlg.GetSelection()
                        dlg.Destroy()

                        if idx == 0:
                            curtype = WOLF_ARRAY_FULL_SINGLE # np.float32
                        elif idx == 1:
                            curtype = WOLF_ARRAY_FULL_DOUBLE # np.float64
                        elif idx == 2:
                            curtype = WOLF_ARRAY_FULL_INTEGER # np.int32
                        elif idx == 3:
                            curtype = WOLF_ARRAY_FULL_INTEGER16 #np.int16
                        elif idx == 4:
                            curtype = WOLF_ARRAY_FULL_INTEGER8 #np.int8

                        newarray = WolfArray(srcheader=self.selected_object.get_header(), whichtype=curtype, nullvalue=self.selected_object.nullvalue)

                        newarray.allocate_ressources()
                        asnewtype = self.selected_object.array.data.astype(newarray.dtype)
                        newarray.array.data[:,:] = asnewtype[:,:]
                        newarray.copy_mask(self.selected_object, forcenullvalue=True, link=False)
                    else:
                        newarray = WolfArray(mold=self.selected_object)

                    self.add_object('array', newobj=newarray, id=newlab)
                    self.Refresh()
                else:
                    logging.warning(_('Not yet implemented'))

        elif text == _('Save as'):
            # save objet to file, choosing the file name

            if self.selected_object is not None:
                if issubclass(type(self.selected_object), WolfArray):
                    filterArray = "bin (*.bin)|*.bin|Geotif (*.tif)|*.tif|Numpy (*.npy)|*.npy|all (*.*)|*.*"
                    fdlg = wx.FileDialog(self, "Choose file name for Array : " + self.selected_object.idx, wildcard=filterArray,
                                         style=wx.FD_SAVE)
                    ret = fdlg.ShowModal()
                    if ret == wx.ID_OK:
                        curfil = fdlg.GetFilterIndex()

                        self.selected_object.filename = fdlg.GetPath()

                        self.selected_object.write_all()
                    fdlg.Destroy()
                elif type(self.selected_object) in [Zones, Bridge, Weir]:
                    filterArray = "vec (*.vec)|*.vec|vecz (*.vecz)|*.vecz|Shapefile (*.shp)|*.shp|all (*.*)|*.*"
                    fdlg = wx.FileDialog(self, "Choose file name for Vector :" + self.selected_object.idx, wildcard=filterArray,
                                         style=wx.FD_SAVE)
                    ret = fdlg.ShowModal()
                    if ret == wx.ID_OK:
                        self.selected_object.saveas(fdlg.GetPath())
                    fdlg.Destroy()

                elif type(self.selected_object) in [PictureCollection, Particularites, Enquetes, Ouvrages, Profils]:
                    filterArray = "vec (*.vec)|*.vec|vecz (*.vecz)|*.vecz"
                    fdlg = wx.FileDialog(self, "Choose file name for Collection :" + self.selected_object.idx, wildcard=filterArray,
                                         style=wx.FD_SAVE)
                    ret = fdlg.ShowModal()
                    if ret == wx.ID_OK:
                        self.selected_object.saveas(fdlg.GetPath())
                    fdlg.Destroy()

                elif type(self.selected_object) is Triangulation:
                    filterArray = "tri (*.tri)|*.tri|all (*.*)|*.*"
                    fdlg = wx.FileDialog(self, "Choose file name for triangulation :" + self.selected_object.idx, wildcard=filterArray,
                                         style=wx.FD_SAVE)
                    ret = fdlg.ShowModal()
                    if ret == wx.ID_OK:
                        self.selected_object.saveas(fdlg.GetPath())

                    fdlg.Destroy()
                elif isinstance(self.selected_object, Particle_system):
                    filterArray = "json (*.json)|*.json|all (*.*)|*.*"
                    fdlg = wx.FileDialog(self, "Choose file name for particle system :" + self.selected_object.idx, wildcard=filterArray,
                                         style=wx.FD_SAVE)
                    ret = fdlg.ShowModal()
                    if ret == wx.ID_OK:
                        self.selected_object.save(fdlg.GetPath())

                    fdlg.Destroy()
                elif isinstance(self.selected_object, DikeWolf):
                    self.selected_object.save_as()

                elif isinstance(self.selected_object, InjectorDike):
                    self.selected_object.save_as()

                elif isinstance(self.selected_object, Wolf_LAZ_Data):
                    filterArray = "Dump (*.dump)|*.dmp|all (*.*)|*.*"
                    fdlg = wx.FileDialog(self, "Choose file name for LAZ data :" + self.selected_object.idx, wildcard=filterArray,
                                            style=wx.FD_SAVE)
                    ret = fdlg.ShowModal()
                    if ret == wx.ID_OK:
                        self.selected_object.saveas(fdlg.GetPath())

                    fdlg.Destroy()

                elif isinstance(self.selected_object, Particle_system):
                    self.selected_object.saveas()

                elif isinstance(self.selected_object, Drowning_victim_Viewer):
                    self.selected_object.saveas()

                elif isinstance(self.selected_object, crosssections):
                    filterArray = "vecz (*.vecz)|*.vecz|SXY (*.sxy)|*.sxy"
                    fdlg = wx.FileDialog(self, "Choose file name for Cross Sections :" + self.selected_object.idx, wildcard=filterArray,
                                         style=wx.FD_SAVE)
                    ret = fdlg.ShowModal()
                    if ret == wx.ID_OK:
                        self.selected_object.saveas(fdlg.GetPath())
                    fdlg.Destroy()

        elif text == _('Properties'):

            myobj = self.selected_object
            if type(myobj) in [WolfArray, WolfArrayMB, WolfArrayMNAP, Zones, Wolfresults_2D, wolfres2DGPU, Particle_system, Picc_data, Cadaster_data, hydrometry_wolfgui, Bridge, Weir, Wolf_LAZ_Data, DikeWolf, Drowning_victim_Viewer, InjectorDike]:
                myobj.show_properties()

            elif isinstance(myobj, cloud_vertices):
                myobj.show_properties()

        elif text == _('Boundary conditions'):
            bc = self.get_boundary_manager(self.selected_object)
            if bc is not None:
                bc.Show()

        elif text == _('Contours'):
            if isinstance(self.selected_object, WolfArray):
                cont = self.selected_object.contour()
                cont.prep_listogl()
                self.add_object('vector', newobj= cont, id= cont.idx)
                self.Paint()

        elif text == _('Rebin'):

            if isinstance(self.selected_object, WolfArray):
                dlg = wx.TextEntryDialog(self, _('Enter the rebin factor (>1 will decrease the resolution, <1 will increase the resolution) :'), _('Rebin'), '1')
                if dlg.ShowModal() == wx.ID_OK:
                    res = dlg.GetValue()
                    dlg.Destroy()
                    try:
                        res = float(res)

                        ops = ['Mean', 'Sum', 'Max', 'Min', 'Median']
                        dlg_op = wx.SingleChoiceDialog(None, _('Choose the operation'), _('Operation'), ops, style=wx.CHOICEDLG_STYLE)
                        ret = dlg_op.ShowModal()

                        if ret != wx.ID_OK:
                            dlg_op.Destroy()
                            logging.info(_('Rebin cancelled'))
                            return

                        op = dlg_op.GetSelection()
                        dlg_op.Destroy()

                        self.selected_object.rebin(res, ops[op].lower())
                    except:
                        logging.warning(_('Invalid value for rebin factor'))
                        return
                else:
                    dlg.Destroy()
            else:
                logging.warning(_('Rebin not yet implemented for this type of object'))

        elif text == _('Set NullValue'):

            if isinstance(self.selected_object, WolfArray):
                dlg = wx.TextEntryDialog(self, _('Enter the new null value :'), _('Set NullValue'), str(self.selected_object.array.data[0,0]))
                if dlg.ShowModal() == wx.ID_OK:
                    res = dlg.GetValue()
                    dlg.Destroy()

                    try:
                        res = float(res)
                        self.selected_object.nullvalue = res
                        self.selected_object.mask_data(res)
                        self.selected_object.reset_plot()
                        self.Refresh()
                    except:
                        logging.warning(_('Invalid value for null value'))
                        return

        elif _('Convert to mono-block') in text:

            if isinstance(self.selected_object, WolfArrayMB):
                mono = self.selected_object.as_WolfArray()
                self.add_object('array', newobj=mono, id=self.selected_object.idx + '_mono')
                logging.info(_('Mono-block created and added to the viewer'))

            elif isinstance(self.selected_object, Wolfresults_2D):
                mono = self.selected_object.as_WolfArray()

                if isinstance(mono, WolfArrayMB):
                    mono = mono.as_WolfArray()

                self.add_object('array', newobj=mono, id=self.selected_object.idx + '_mono')
                logging.info(_('Mono-block created and added to the viewer'))

            else:
                logging.warning(_('Convert to mono-blocks not yet implemented for this type of object'))

        elif _('Convert to multi-blocks') in text:

            if isinstance(self.selected_object, Wolfresults_2D):
                mb = self.selected_object.as_WolfArray(force_mb=True)

                if isinstance(mb, WolfArrayMB):
                    logging.info(_('Multi-blocks created and added to the viewer'))

                elif isinstance(mb, WolfArray):
                    logging.warning(_('Mono-blocks created and added to the viewer -- Instead of multi-blocks as only one block was found'))

                self.add_object('array', newobj=mb, id=self.selected_object.idx + '_mb')
            else:
                logging.warning(_('Convert to multi-blocks not yet implemented for this type of object'))

        elif _('Extract current step as IC') in text:

            if isinstance(self.selected_object, Wolfresults_2D):

                if self.selected_object.current_result is None or self.selected_object.current_result < 0 :
                    logging.warning(_('No current step defined'))
                    return

                import shutil

                logging.info(_('Extracting current step as IC'))

                dir = Path(self.selected_object.filename).parent
                files = ['h.npy', 'qx.npy', 'qy.npy']
                for curfile in files:
                    if (dir / curfile).exists():
                        # prepend the file with __
                        logging.info(f'File {curfile} already exists -- renaming it')
                        shutil.copyfile(dir / curfile, dir / ('__' + curfile))

                logging.info(_('Extracting current step as IC'))
                self.selected_object.export_as(Path(self.selected_object.filename).parent, [views_2D.WATERDEPTH, views_2D.QX, views_2D.QY], 'numpy', False)
                logging.info(_('Done !'))

        elif _('Export to Shape file') in text:

            if isinstance(self.selected_object, Zones | Bridge | Weir):
                filterArray = "Shapefile (*.shp)|*.shp"
                fdlg = wx.FileDialog(self, "Choose file name for Zones :" + self.selected_object.idx, wildcard=filterArray,
                                        style=wx.FD_SAVE)
                ret = fdlg.ShowModal()
                if ret == wx.ID_OK:
                    self.selected_object.export_to_shapefile(fdlg.GetPath())
                fdlg.Destroy()

        elif _('Export active zone to Shape file') in text:

            if isinstance(self.selected_object, Zones | Bridge | Weir):

                filterArray = "Shapefile (*.shp)|*.shp"
                fdlg = wx.FileDialog(self, "Choose file name for Vector :" + self.selected_object.idx, wildcard=filterArray,
                                        style=wx.FD_SAVE)
                ret = fdlg.ShowModal()
                if ret == wx.ID_OK:
                    self.selected_object.export_active_zone_to_shapefile(fdlg.GetPath())
                fdlg.Destroy()

        elif _('Set colormap') in text:

            if isinstance(self.selected_object, Wolf_LAZ_Data):
                self.selected_object.associated_color = self._choice_laz_colormap()

        elif _('Edit colormap') in text:

            if isinstance(self.selected_object, Wolf_LAZ_Data):
                self.selected_object.interactive_update_colors()

        elif _('Set classification') in text:

            if isinstance(self.selected_object, Wolf_LAZ_Data):
                self.selected_object.set_classification(self._choice_laz_classification())

        elif _('Edit selection') in text:

            if isinstance(self.selected_object, Wolf_LAZ_Data):
                self.selected_object._edit_selection()

        elif _('All to cloud') in text:

            if isinstance(self.selected_object, Wolf_LAZ_Data):

                if self.selected_object.num_points > 100000:
                    dlg = wx.MessageDialog(None, _('The number of points is high, it could take some time to convert to cloud.\nDo you want to continue ?'), _('Warning'), wx.YES_NO | wx.NO_DEFAULT)
                    ret = dlg.ShowModal()

                    if ret != wx.ID_YES:
                        dlg.Destroy()

                        return

                newcloud = cloud_vertices()
                newcloud.init_from_nparray(self.selected_object.xyz)
                self.add_object('cloud', newobj=newcloud, id=self.selected_object.idx + '_cloud')

        elif _('Selection to cloud') in text:

            if isinstance(self.selected_object, Wolf_LAZ_Data):

                xyz = self.selected_object.xyz_selected
                if xyz.shape[0] ==0:
                    logging.warning('No points selected')
                    return

                if xyz.shape[0] > 100000:
                    dlg = wx.MessageDialog(None, _('The number of points is high, it could take some time to convert to cloud.\nDo you want to continue ?'), _('Warning'), wx.YES_NO | wx.NO_DEFAULT)
                    ret = dlg.ShowModal()

                    if ret != wx.ID_YES:
                        dlg.Destroy()

                        return

                newcloud = cloud_vertices()
                newcloud.init_from_nparray(xyz)
                self.add_object('cloud', newobj=newcloud, id=self.selected_object.idx + '_cloud_sel')

        elif _('Selection to vector') in text:

            if isinstance(self.selected_object, Wolf_LAZ_Data):

                if self.active_zone is None:
                    logging.warning(_('No active zone selected'))
                    return

                xyz = self.selected_object.xyz_selected
                if xyz.shape[0] ==0:
                    logging.warning(_('No points selected'))

                    if self.selected_object._myprops[('Selection', 'Codes')] != '':
                        logging.info(_('You filtered the points with the codes : {}'.format(self.selected_object._myprops[('Selection', 'Codes')])))
                    return

                if xyz.shape[0] > 100000:
                    dlg = wx.MessageDialog(None, _('The number of points is high, it could take some time to convert to cloud.\nDo you want to continue ?'), _('Warning'), wx.YES_NO | wx.NO_DEFAULT)
                    ret = dlg.ShowModal()

                    if ret != wx.ID_YES:
                        dlg.Destroy()

                        return

                def approximate_vector(xyz):
                    """ Get a cloud of points and return a vector
                    based on the best approximated segment
                    and points projected on its trace
                    """

                    # best approximation of the segment
                    # based on the RANSAC algorithm from scikit-learn
                    model = linear_model.RANSACRegressor()
                    model.fit(xyz[:,0].reshape(-1,1), xyz[:,1])

                    # get the points projected on the segment
                    proj = model.predict(xyz[:,0].reshape(-1,1))

                    # get the coordinates of the projected points
                    xyz_proj = np.zeros((xyz.shape[0],3))
                    xyz_proj[:,0] = xyz[:,0]
                    xyz_proj[:,1] = proj
                    xyz_proj[:,2] = xyz[:,2]

                    #Sort the points
                    idx = np.argsort(xyz_proj[:,0])
                    xyz_proj = xyz_proj[idx]

                    return xyz_proj

                newvector = vector(name = self.selected_object.idx + '_vector_sel', fromnumpy= approximate_vector(xyz))
                self.active_zone.add_vector(newvector, forceparent=True)
                self.active_zone.parent.find_minmax(True)
                self.active_zone.parent.reset_listogl()

                self.active_zone.parent.fill_structure()
                self.Refresh()

        elif _('Play') in text:

            if isinstance(self.selected_object, Wolf_LAZ_Data):
                self.selected_object.play_flight()

        elif _('Add point') in text:

            if isinstance(self.selected_object, Wolf_LAZ_Data):
                self.selected_object.add_pose_in_memory()

        elif _('Record') in text:

            if isinstance(self.selected_object, Wolf_LAZ_Data):

                dlg = wx.DirDialog(self, _('Choose a directory to save the video'), style=wx.DD_DEFAULT_STYLE)
                if dlg.ShowModal() == wx.ID_OK:
                    self.selected_object.record_flight(dlg.GetPath())

                dlg.Destroy()

        elif _('Load flight') in text:

            if isinstance(self.selected_object, Wolf_LAZ_Data):
                dlg = wx.FileDialog(self, _('Choose a file to load the flight'), wildcard='JSON (*.json)|*.json|All (*.*)|*.*', style=wx.FD_OPEN)
                if dlg.ShowModal() == wx.ID_OK:
                    self.selected_object.load_flight(dlg.GetPath())

                dlg.Destroy()

        elif _('Save flight') in text:

            if isinstance(self.selected_object, Wolf_LAZ_Data):
                dlg = wx.FileDialog(self, _('Choose a file to save the flight'), wildcard='JSON (*.json)|*.json|All (*.*)|*.*', style=wx.FD_SAVE)
                if dlg.ShowModal() == wx.ID_OK:
                    self.selected_object.save_flight(dlg.GetPath())

                dlg.Destroy()

        elif _('Reload') in text:

            if isinstance(self.selected_object, WolfArray):
                if self.selected_object.filename is not None:

                    dlg = wx.MessageDialog(None, _('Do you want to reload the file ?'), _('Reload'), wx.YES_NO | wx.NO_DEFAULT)
                    ret = dlg.ShowModal()
                    if ret == wx.ID_YES:
                        self.selected_object.read_all()
                        self.selected_object.mask_data(self.selected_object.nullvalue)
                        self.selected_object.reset_plot()

                    dlg.Destroy()
            # elif isinstance(self.selected_object, Zones):
            #     if self.selected_object.filename is not None:
            #         dlg = wx.MessageDialog(None, _('Do you want to reload the file ?'), _('Reload'), wx.YES_NO | wx.NO_DEFAULT)
            #         ret = dlg.ShowModal()
            #         if ret == wx.ID_YES:
            #             self.selected_object.read()

            #         dlg.Destroy()

            else:
                logging.warning(_('Reload not yet implemented for this type of object'))

        elif _('Rasterize active zone') in text:

            if self.active_array is None and self.active_res2d is None:
                logging.warning(_('No active array selected'))
                return

            if self.active_array is not None and self.active_res2d is not None:
                # Show a dialog to choose between array or res2d
                dlg = wx.SingleChoiceDialog(None, _('Choose the type of rasterization'), _('Rasterization Type'), ['Array', 'Res2D'], style=wx.CHOICEDLG_STYLE)
                ret = dlg.ShowModal()
                if ret != wx.ID_OK:
                    dlg.Destroy()
                    return
                choice = dlg.GetSelection()
                dlg.Destroy()
                if choice == 0:
                    based_array = self.active_array
                else:
                    if self.active_res2d.nb_blocks == 1:
                        based_array = self.active_res2d.get_header_block(1)
                    elif self.active_res2d.nb_blocks > 1:
                        logging.error(_('Rasterization not implemented for multi-blocks res2d -- Convert to mono-blocks first'))
                        return

            elif self.active_array is not None:
                based_array = self.active_array
            elif self.active_res2d is not None:
                if self.active_res2d.nb_blocks == 1:
                    based_array = self.active_res2d.get_header_block(1)
                elif self.active_res2d.nb_blocks > 1:
                    logging.error(_('Rasterization not implemented for multi-blocks res2d -- Convert to mono-blocks first'))
                    return

            if isinstance(self.selected_object, Zones):
                if self.selected_object.active_zone is None:
                    logging.warning(_('No active zone selected'))
                    return
                if self.selected_object.active_zone.parent is None:
                    logging.warning(_('No parent object for the active zone'))
                else:
                    new_zone = based_array.rasterize_zone_along_grid(self.selected_object.active_zone, outformat= zone)
                    if new_zone is not None:
                        new_zone.myname = self.selected_object.active_zone.myname + '_rasterized'
                        self.selected_object.active_zone.parent.add_zone(new_zone, forceparent=True)
                        self.selected_object.active_zone.parent.fill_structure()
                        self.selected_object.active_zone.reset_listogl()

        elif _('Rasterize active vector') in text:

            if self.active_array is None and self.active_res2d is None:
                logging.warning(_('No active array selected'))
                return

            if self.active_array is not None and self.active_res2d is not None:
                # Show a dialog to choose between array or res2d
                dlg = wx.SingleChoiceDialog(None, _('Choose the type of rasterization'), _('Rasterization Type'), ['Array', 'Res2D'], style=wx.CHOICEDLG_STYLE)
                ret = dlg.ShowModal()
                if ret != wx.ID_OK:
                    dlg.Destroy()
                    return
                choice = dlg.GetSelection()
                dlg.Destroy()
                if choice == 0:
                    based_array = self.active_array
                else:
                    if self.active_res2d.nb_blocks == 1:
                        based_array = self.active_res2d.get_header_block(1)
                    elif self.active_res2d.nb_blocks > 1:
                        logging.error(_('Rasterization not implemented for multi-blocks res2d -- Convert to mono-blocks first'))
                        return

            elif self.active_array is not None:
                based_array = self.active_array
            elif self.active_res2d is not None:
                if self.active_res2d.nb_blocks == 1:
                    based_array = self.active_res2d.get_header_block(1)
                elif self.active_res2d.nb_blocks > 1:
                    logging.error(_('Rasterization not implemented for multi-blocks res2d -- Convert to mono-blocks first'))
                    return

            if isinstance(self.selected_object, Zones):
                if self.selected_object.active_vector is None:
                    logging.warning(_('No active vector selected'))
                    return

                vec_raster = based_array.rasterize_vector_along_grid(self.selected_object.active_vector)
                if vec_raster is not None:
                    vec_raster.myname = self.selected_object.active_vector.myname + '_rasterized'
                    self.selected_object.active_vector.parentzone.add_vector(vec_raster, forceparent=True, update_struct=True)
                    self.selected_object.active_vector.parentzone.reset_listogl()

        elif  _('Extrude on active array') in text:

            if isinstance(self.selected_object, Picc_data):

                if self.active_array is None:
                    logging.warning(_('No active array selected'))
                    return

                logging.info(_('Extruding polygons on active array'))
                logging.info(_('Please wait, it could take some time...'))
                self.selected_object.extrude_polygons(self.active_array)
                logging.info(_('Extrusion done !'))

                self.active_array.reset_plot()
                self.Refresh()

        elif  _('Interpolate on active array') in text:

            if isinstance(self.selected_object, Zones):

                if self.active_array is None:
                    logging.warning(_('No active array selected'))
                    return

                dlg = wx.SingleChoiceDialog(None, _('With respect to the array values, keep the extruded values which are ?'), _('Interpolate on active array'), ['Above', 'Below', 'All'], style=wx.CHOICEDLG_STYLE)
                ret = dlg.ShowModal()
                if ret != wx.ID_OK:
                    dlg.Destroy()
                    return
                choice = dlg.GetStringSelection().lower()
                dlg.Destroy()


                logging.info(_('Interpolating polygons on active array'))
                logging.info(_('Please wait, it could take some time...'))
                for curzone in self.selected_object.myzones:
                    self.active_array.interpolate_on_polygons(curzone, keep = choice)
                logging.info(_('Interpolation done !'))


    def OnClose(self, event):
        """ Close the application """

        nb = 0
        if self.linked:
            if self.linkedList is not None:
                if self in self.linkedList:
                    id = self.linkedList.index(self)
                    self.linkedList.pop(id)
                    nb = len(self.linkedList)

        if nb == 0:
            if self.wxlogging is not None:
                dlg = wx.MessageDialog(None, _('Do you want to quit Wolf ?'), _('Quit Wolf'), wx.YES_NO | wx.NO_DEFAULT)
                ret = dlg.ShowModal()
                dlg.Destroy()
                if ret == wx.ID_YES:
                    self.Destroy()
                    #FIXME : It is not a really proper way to quit the application
                    wx.Exit()
                    return
                else:
                    return
        self.Destroy()

    def OnSelectItem(self, event):
        """ Select the item in the tree list """

        ctrl = wx.GetKeyState(wx.WXK_CONTROL)
        alt = wx.GetKeyState(wx.WXK_ALT)

        myitem = event.GetItem()

        nameitem = self.treelist.GetItemText(myitem).lower()
        curobj = self.getobj_from_id(nameitem)
        myobj = self.treelist.GetItemData(myitem)

        if curobj is not myobj:
            logging.error(_('Bad association between object and tree item'))
            logging.error(_('Do you have 2 objects with the same id ?'))
            logging.error(_('It could be the case if you have drag/drop an object in the viewer...'))
            logging.error(_('I will continue but it is not normal...'))

        self.treelist.SetToolTip(self.treelist.GetItemText(myitem))

        # myparent = self.treelist.GetItemParent(myitem)
        # check = self.treelist.GetCheckedState(myitem)
        # if myparent is not None:
        #     nameparent = self.treelist.GetItemText(myparent).lower()

        self.selected_object = curobj
        self.selected_treeitem = myitem

    def OnCheckItem(self, event:TreeListEvent):
        """ Check the item in the tree list """

        myitem = event.GetItem()
        myparent = self.treelist.GetItemParent(myitem)
        check = self.treelist.GetCheckedState(myitem)
        nameparent = self.treelist.GetItemText(myparent).lower()
        nameitem = self.treelist.GetItemText(myitem).lower()

        ctrl = wx.GetKeyState(wx.WXK_CONTROL)
        shiftdown = wx.GetKeyState(wx.WXK_SHIFT)

        # ctrl = event.ControlDown()

        if nameparent != '':
            curobj = self.getobj_from_id(nameitem)
            if curobj is None:
                return

            if bool(check):
                try:
                    curobj.check_plot()

                    if isinstance(curobj, PlansTerrier):
                        if curobj.initialized:
                            self.menu_landmaps()
                            logging.info(_('Landmap initialized'))
                        else:
                            logging.warning(_('Landmap not initialized'))

                    elif isinstance(curobj, Ouvrages):
                        if curobj.initialized:
                            self.menu_pictcollection()
                            logging.info(_('Ouvrages collection initialized'))
                        else:
                            logging.warning(_('Ouvrages collection not initialized'))

                    elif isinstance(curobj, Particularites):
                        if curobj.initialized:
                            self.menu_pictcollection()
                            logging.info(_('Particularites collection initialized'))
                        else:
                            logging.warning(_('Particularites collection not initialized'))

                    elif isinstance(curobj, Enquetes):
                        if curobj.initialized:
                            self.menu_pictcollection()
                            logging.info(_('Enquetes collection initialized'))
                        else:
                            logging.warning(_('Enquetes collection not initialized'))

                    elif isinstance(curobj, Profils):
                        if curobj.initialized:
                            self.menu_pictcollection()
                            logging.info(_('Profils collection initialized'))
                        else:
                            logging.warning(_('Profils collection not initialized'))

                except Exception as ex:
                    wx.LogMessage(str(ex))
                    wx.MessageBox(str(ex), _("Error"), wx.ICON_ERROR)
            else:
                if issubclass(type(curobj), WolfArray):
                    curobj.uncheck_plot(not ctrl,ctrl)
                elif isinstance(curobj, Picc_data):
                    curobj.uncheck_plot(ctrl, shiftdown)
                else:
                    curobj.uncheck_plot()

            # if nameparent == 'vectors' or nameparent == 'cross_sections':
            #     if wx.GetKeyState(wx.WXK_CONTROL):
            #         curobj.showstructure(self)

            if curobj.idx == 'grid' and check:
                dlg = wx.TextEntryDialog(self, 'Size of the Grid ? (float)', 'Choose an size')
                dlg.SetValue('1000.')
                size = 1000.
                if dlg.ShowModal() == wx.ID_OK:
                    size = float(dlg.GetValue())
                curobj.creategrid(size, self.xmin, self.ymin, self.xmax, self.ymax)

            if 'alaro' in curobj.idx and check:
                if self.alaro_navigator is None:
                    self.alaro_navigator = Alaro_Navigator(self, curobj.idx, 'Alaro')
                    self.alaro_navigator.Show()

            self.Refresh()

    def _alaro_update_time(self):
        """ Update the time of the alaro navigator """

        objs = self.get_list_objects(drawing_type=draw_type.WMSFORE, checked_state=True)
        for obj in objs:
            obj.time = self.alaro_navigator.time_str
            obj.alpha = self.alaro_navigator.alpha
            obj.force_alpha = True

        self._update_foreground()
        self.Paint()

    def _alaro_legends(self):
        """ Show images of the checked alaro layers"""
        objs = self.get_list_objects(drawing_type=draw_type.WMSFORE, checked_state=True)
        for obj in objs:
            if obj.category == 'ALARO':
                img = Image.open(get_Alaro_legend(obj.idx.replace('alaro ', '')))
                img.show()

    def getXY(self, pospix):

        width, height = self.canvas.GetSize()
        X = float(pospix[0]) / self.sx + self.xmin
        Y = float(height - pospix[1]) / self.sy + self.ymin
        return X, Y

    def OnZoomGesture(self, e):
        pass

    def OnLeave(self, e):
        if e.ControlDown():
            self.mytooltip.Show(False)

    def get_cross_sections(self):
        """
        Récupération du premier objet crosssections disponible
        """
        for obj in self.iterator_over_objects(draw_type.VECTORS):
            if isinstance(obj,crosssections):
                return obj

        return None

    def set_active_profile(self, active_profile: profile):
        """
        This method sets the active profile in Pydraw (useful for interfaces communication).
        """
        self.active_profile = active_profile

    def set_active_vector(self, active_vector: vector):
        """
        This method sets the active vector in Pydraw (useful for interfaces communication).
        """
        self.active_vector = active_vector

    def get_active_profile(self):
        """
        This methods returns the  active profile in pydraw (useful for interfaces communication).
        """
        return self.active_profile

    def plot_cross(self, x:float, y:float):

        # Search for cross sections (List of profiles)
        if self.active_cs is None:
            self.active_cs = self.get_cross_sections()
            if self.active_cs is None:
                logging.warning(_('No cross sections available -- Please load a file or create data !'))
            return

        # Initialisation of the notebook where the active profile is plotted.
        if self.notebookprof is None:
            self.notebookprof = ProfileNotebook(mapviewer=self)
            # self.myfigprof = self.notebookprof.add('Figure 1', which= "all")
            self.myfigprof = self.notebookprof.add('Reference') # FIXME Updated add method

        else:
            try:
                self.notebookprof.Show()
            except:
                self.notebookprof = ProfileNotebook(mapviewer=self)
                # self.myfigprof = self.notebookprof.add('Figure 1', which= "all")
                self.myfigprof = self.notebookprof.add('Figure 1') # FIXME updated  add method

        # Initialisation of the active profile
        # 1. We uncolor the active profile in wolf GUI.
        self.active_profile: profile
        if self.active_profile is not None:
            self.active_profile.uncolor_active_profile()

        # 2. We select the closest profile corresponding to the user's right click in the GUI.
        self.active_profile = self.active_cs.select_profile(x, y)

        #Finally, we set the profile and the cross section (list of profiles) in the notebook.
        #FIXME Iden establishes the communications between pydraw and the notebook (to avoid circular information).
        self.myfigprof.cs_setter(mycross = self.active_cs, active_profile= self.active_profile, mapviewer = self)


    def On_Mouse_Right_Down(self, e: wx.MouseEvent):
        """
        Event when the right button of the mouse is pressed.

        We use this event to manage "action" set by others objects.

        """
        #self._user_activity_true()

        pos = e.GetPosition()
        x, y = self.getXY(pos)

        alt = e.AltDown()
        ctrl = e.ControlDown()
        shiftdown = e.ShiftDown()

        if self.action is None:
            if self.active_bc is not None:

                self.start_action('select bc', _('Select a boundary condition'))
                tmpvec = vector()

                self.last_active_vector = self.active_vector
                self.active_vector = tmpvec
                tmpvec.add_vertex(wolfvertex(x, y))

            self.rightdown = (x, y)

        elif self.action == 'plot alaro xy':

            if self.active_alaro._gdf is None:
                logging.warning(_('No Alaro run loaded -- Please load a run first'))
                return

            fig = self.active_alaro.plot_Rain_and_TotPrecip4XY(x, y)
            fig.show()

        elif self.action == 'distance along vector':

            # add a vertex to the vector
            self._tmp_vector_distance.add_vertex(wolfvertex(x, y))

        elif self.action == 'move triangles':

            if self.active_tri is None:
                logging.warning(_('No triangles selected -- Please select a triangulation first !'))
                return

            if self.active_tri._move_start is None:
                self.active_tri._move_start = (x, y)
                return

            delta_x = x - self.active_tri._move_start[0]
            delta_y = y - self.active_tri._move_start[1]

            if shiftdown:
                delta_y = 0.

            if alt:
                delta_x = 0.

            self.active_tri.move(delta_x, delta_y)
            self.active_tri.reset_plot()
            self.active_tri._move_start = None
            self.active_tri.clear_cache()
            self.end_action(_('End move triangulation'))

        elif self.action == 'rotate triangles':

            if self.active_tri is None:
                logging.warning(_('No vector selected -- Please select a triangulation first !'))
                return

            if self.active_tri._rotation_center is None:
                self.active_tri._rotation_center = (x,y)
                return

            if shiftdown:
                if ctrl:
                    self.active_tri._rotation_step = None
                else:
                    # Set the rotation step
                    self.active_tri._rotation_step = np.degrees(np.arctan2(y - self.active_tri._rotation_center[1], x - self.active_tri._rotation_center[0]))

            self.active_tri.rotate_xy(x, y)
            self.active_tri._rotation_center = None
            self.active_tri.clear_cache()
            self.active_tri.reset_plot()
            self.end_action(_('End rotate triangulation'))

        elif 'pick landmap' in self.action:
            # Pick a landmap if loaded

            if self.active_landmap is None:
                logging.warning(_('No landmap available -- Please activate the data and retry !'))
                return

            if 'full' in self.action:
                self.active_landmap.load_texture(x,y, which='full')
            else:
                self.active_landmap.load_texture(x,y, which='low')
            self.Refresh()

        elif self.action == 'pick municipality':
            # Pick a municipality if loaded
            if self.active_qdfidf is None:
                logging.warning(_('No municipality data available -- Please activate the data and retry !'))
                return

            self.active_qdfidf.pick_municipality(x, y, self.get_canvas_bounds())

        elif self.action == 'pick a picture':
            # Pick a picture

            if self.active_picturecollection is None:
                logging.warning(_('No picture collection available -- Please activate the data and retry !'))
                return

            vec = self.active_picturecollection.find_vector_containing_point(x, y)
            vec.myprop.imagevisible = not vec.myprop.imagevisible

            if shiftdown:
                # show/hide the legend
                vec.myprop.legendvisible = not vec.myprop.legendvisible

            vec.myprop.update_myprops()
            # vec.myprop.load_unload_image()
            vec.parentzone.reset_listogl()

            self.active_picturecollection.Activate_vector(vec)

            self.Refresh()

        elif self.action == 'pick bridge':
            self.pick_bridge(x, y)

        elif self.action == 'pick weir':
            self.pick_weir(x, y)

        elif self.action == 'bridge gltf':
            # Create a bridge in gltf format

            self.bridgepar = (x, y)

            dlg = wx.TextEntryDialog(self, 'Z maximum ?', 'Choose an elevation as top')
            dlg.SetValue('')

            zmax = 0.
            if dlg.ShowModal() == wx.ID_OK:
                zmax = float(dlg.GetValue())
            dlg.Destroy()

            dlg = wx.FileDialog(None, _('Choose filename'),
                                wildcard='glb (*.glb)|*.glb|gltf2 (*.gltf)|*.gltf|All (*.*)|*.*', style=wx.FD_SAVE)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            fn = dlg.GetPath()
            dlg.Destroy()

            points, triangles = self.active_vector.triangulation_ponts(self.bridgepar[0], self.bridgepar[1], zmax)
            self.active_cs.export_gltf_gen(points, triangles, fn)

            self.start_action('', 'None')

        elif self.action == 'Plot cross section':
            # Plot cross section
            self.plot_cross(x, y)

        elif self.action == 'Set 1D profile':
            # Set 1D profile

            if self.active_cs is None:
                self.active_cs = self.get_cross_sections()

            if self.active_cs is None:
                logging.warning(_('No cross sections available -- Please load a file or create data !'))
                return

            if self.notebookcs is None:
                self.notebookcs = PlotNotebook()
                self.myfigcs = self.notebookcs.add(_("Cross section"), "CS")
            else:
                try:
                    self.notebookcs.Show()
                except:
                    self.notebookcs = PlotNotebook()
                    self.myfigcs = self.notebookcs.add(_("Cross section"), "CS")

            self.active_profile: profile

            #on met l'ancien profil actif en noir
            if self.active_profile is not None:
                self.active_profile.uncolor_active_profile()
            if self.myfigcs.mycs is not None:
                self.myfigcs.mycs.uncolor_active_profile()

            # self.active_profile = self.active_cs.select_profile(x, y)
            self.active_profile = self.frame_create1Dfrom2D.active_profile

            self.myfigcs.set_linked_arrays(self.get_linked_arrays())
            self.myfigcs.set_cs(self.active_profile)

            #on met le profil en rouge et plus épais
            self.active_profile.color_active_profile()
            self.zoom_on_active_profile()

            self.Paint()

        elif self.action == 'Select nearest profile':
            # Select nearest profile

            if self.active_cs is None:
                self.active_cs = self.get_cross_sections()

            if self.active_cs is None:
                logging.warning(_('No cross sections available -- Please load a file or create data !'))
                return

            if self.notebookcs is None:
                self.notebookcs = PlotNotebook()
                self.myfigcs = self.notebookcs.add(_("Cross section"), "CS")
            else:
                try:
                    self.notebookcs.Show()
                except:
                    self.notebookcs = PlotNotebook()
                    self.myfigcs = self.notebookcs.add(_("Cross section"), "CS")

            self.active_profile: profile

            #on met l'ancien profil actif en noir
            if self.active_profile is not None:
                self.active_profile.uncolor_active_profile()
            if self.myfigcs.mycs is not None:
                self.myfigcs.mycs.uncolor_active_profile()

            self.active_profile = self.active_cs.select_profile(x, y)

            self.myfigcs.set_linked_arrays(self.get_linked_arrays())
            self.myfigcs.set_cs(self.active_profile)

            #on met le profil en rouge et plus épais
            self.active_profile.color_active_profile()

            self.Refresh()

        elif self.action == 'select active tile':
            # Select active tile from Lidar data

            self.active_tile.select_vectors_from_point(x, y, True)
            self.active_vector = self.active_tile.get_selected_vectors()

            tilearray = self.active_tile.get_array(self.active_vector)
            if tilearray is not None:
                if self.active_vector.myname =='':
                    bbox = self.active_vector.get_bounds()
                    id_label = '{}-{}'.format(bbox[0][0], bbox[1][1])
                else:
                    id_label = self.active_vector.myname

                self.add_object('array', newobj = tilearray, ToCheck=True, id=id_label)

        elif self.action == 'select active image tile':
            # select active image tile

            self.active_imagestiles.select_vectors_from_point(x, y, True)
            active_tile = self.active_imagestiles.get_selected_vectors()
            active_tile.myprop.imagevisible = not active_tile.myprop.imagevisible

        elif self.action.find('select active vector') > -1:
            # Select active vector

            inside = self.action.find('inside') > -1 # only polygons/closed polyline if 'inside' is in the action name
            onlyonezone = self.action.find('2') > -1 # only the active zone if '2' is in the action name, all zones otherwise

            if onlyonezone:
                self.active_zone.select_vectors_from_point(x, y, inside)
                self.active_vector = self.active_zone.get_selected_vectors()[0]

                if self.active_vector is not None:
                    self.active_zone.parent.Activate_vector(self.active_vector)
                    self.active_zone.active_vector = self.active_vector
                    self.active_zones.active_zone = self.active_vector.parentzone
            else:
                self.active_zones.select_vectors_from_point(x, y, inside)
                self.active_vector = self.active_zones.get_selected_vectors()

            if self.active_vector is not None:
                self.active_zones.Activate_vector(self.active_vector)
                self.active_zone = self.active_vector.parentzone
                self.active_zones.expand_tree(self.active_zone)

        elif 'select node by node' in self.action:
            # Select node by node

            if 'results' in self.action:
                curobj:Wolfresults_2D
                curobj = self.active_res2d.SelectionData
            else:
                curobj: WolfArray
                curobj = self.active_array.SelectionData

            if curobj.myselection == 'all':
                logging.warning(_('All nodes are selected !!'))
                logging.warning(_('Selecting node by node will force to reset the selection'))
                logging.warning(_('and start from scratch'))

            curobj.add_node_to_selection(x, y)
            curobj.update_nb_nodes_selection()
            self.Paint()

        elif 'select by tmp vector' in self.action or 'select by vector' in self.action:
            # Select nodes by vector or temporary vector
            self.active_vector.add_vertex(wolfvertex(x, y))

        elif 'laz tmp vector' == self.action:
            self.active_vector.add_vertex(wolfvertex(x, y))
            self.active_vector.find_minmax()

        elif self.action == 'create polygon - tiles':
            self.active_vector.add_vertex(wolfvertex(x, y))
            self.active_vector.find_minmax()

        elif self.action == 'capture vertices':

            if ctrl:
                if self.active_array is not None:
                    z = self.active_array.get_value(x, y)
                    self.active_vector.myvertices[-1].z = z
                else:
                    logging.warning(_('No array available and ctrl is pressed -- Please load a file or create data !'))

            self.active_vector.add_vertex(wolfvertex(x, y))
            self.active_vertex = self.active_vector.myvertices[-1]

            self.active_vector.find_minmax()
            self.active_zone.find_minmax()

            # Update the ogl list
            self.active_zone.reset_listogl()
            self.active_vector.reset_linestring()

        elif self.action == 'offset/scale image':

            if self.active_vector is None:
                logging.warning(_('No vector selected -- Please select a vector first !'))
                return

            if self.active_vector.myprop.textureimage is None:
                logging.warning(_('No image available -- Please load an image first !'))
                return

            if self.active_vector._move_start is None:
                self.active_vector._move_start = (x, y)
                return

            delta_x = x - self.active_vector._move_start[0]
            delta_y = y - self.active_vector._move_start[1]

            self.active_vector.myprop.offset_image(delta_x, delta_y)
            self.active_vector.myprop.update_myprops()
            self.active_vector._move_start = None
            self.end_action(_('End offset/scale image'))

        elif self.action == 'move vector':

            if self.active_vector is None:
                logging.warning(_('No vector selected -- Please select a vector first !'))
                return

            if self.active_vector._move_start is None:
                self.active_vector._move_start = (x, y)
                return

            delta_x = x - self.active_vector._move_start[0]
            delta_y = y - self.active_vector._move_start[1]

            if shiftdown:
                delta_y = 0.

            if alt:
                delta_x = 0.

            self.active_vector.move(delta_x, delta_y)
            self.active_vector.clear_cache()
            self.active_vector._move_start = None
            self.end_action(_('End move vector'))

        elif self.action == 'rotate vector':

            if self.active_vector is None:
                logging.warning(_('No vector selected -- Please select a vector first !'))
                return

            if self.active_vector._rotation_center is None:
                self.active_vector._rotation_center = (x,y)
                return

            if shiftdown:
                if ctrl:
                    self.active_vector._rotation_step = None
                else:
                    # Set the rotation step
                    self.active_vector._rotation_step = np.degrees(np.arctan2(y - self.active_vector._rotation_center[1], x - self.active_vector._rotation_center[0]))

            self.active_vector.rotate_xy(x, y)
            self.active_vector.clear_cache()
            self.active_vector._rotation_center = None
            self.end_action(_('End rotate vector'))

        elif self.action == 'move zone':

            if self.active_zone is None:
                logging.warning(_('No zone selected -- Please select a zone first !'))
                return

            if self.active_zone._move_start is None:
                self.active_zone._move_start = (x, y)
                return

            if shiftdown:
                delta_y = 0.

            if alt:
                delta_x = 0.

            self.active_zone.move(x - self.active_zone._move_start[0], y - self.active_zone._move_start[1])
            self.active_zone.clear_cache()
            self.active_zone._move_start = None
            self.end_action(_('End move zone'))

        elif self.action == 'rotate zone':

            if self.active_zone is None:
                logging.warning(_('No zone selected -- Please select a zone first !'))
                return

            if self.active_zone._rotation_center is None:
                self.active_zone._rotation_center = (x,y)
                return

            self.active_zone.rotate_xy(x, y)
            self.active_zone.clear_cache()
            self.active_zone._rotation_center = None
            self.end_action(_('End rotate zone'))

        elif self.action == 'dynamic parallel':
            # Create a dynamic parallel line
            if ctrl:
                if self.active_array is not None:
                    z = self.active_array.get_value(x, y)
                    self.active_vector.myvertices[-1].z = z
                else:
                    logging.warning(_('No array available and ctrl is pressed -- Please load a file or create data !'))

            self.active_vector.add_vertex(wolfvertex(x, y))

            self.active_zone.parallel_active(self.dynapar_dist)

        elif self.action == 'modify vertices':

            if self.active_vertex is None:
                self.active_vertex = self.active_vector.find_nearest_vert(x, y)
            else:
                self.active_vertex.limit2bounds(self.active_vector._mylimits)

                if ctrl:
                    if self.active_array is not None:
                        # Get the value of the array at the position of the vertex
                        z = self.active_array.get_value(x, y)
                        self.active_vertex.z = z
                    else:
                        logging.warning(_('No array available and ctrl is pressed -- Please load a file or create data !'))

                self.active_vertex = None

        elif self.action == 'insert vertices':

            if self.active_vertex is None:
                self.active_vertex = self.active_vector.insert_nearest_vert(x, y)
            else:

                if ctrl:
                    if self.active_array is not None:
                        z = self.active_array.get_value(x, y)
                        self.active_vertex.z = z
                    else:
                        logging.warning(_('No array available and ctrl is pressed -- Please load a file or create data !'))

                self.active_vertex = None

        else:
            self.rightdown = (x, y)


    def On_Mouse_Right_Up(self, e):

        #self._user_activity_true()

        pos = e.GetPosition()
        x, y = self.getXY(pos)

        if self.active_bc is not None:
            if self.action == 'select bc':
                try:
                    minx = min(self.rightdown[0], x)
                    miny = min(self.rightdown[1], y)
                    maxx = max(self.rightdown[0], x)
                    maxy = max(self.rightdown[1], y)

                    if minx != maxx and maxy != miny:
                        self.active_bc.ray_tracing_numpy([[minx, miny], [maxx, miny], [maxx, maxy], [minx, maxy]], 'X')
                        self.active_bc.ray_tracing_numpy([[minx, miny], [maxx, miny], [maxx, maxy], [minx, maxy]], 'Y')
                    else:
                        self.active_bc.query_kdtree((x, y))

                    self.active_bc.update_selection()
                    self.Refresh()

                    self.active_vector = self.last_active_vector

                    self.end_action(_('End selection BC'))

                except:
                    pass

    def On_Mouse_Button(self, e: wx.MouseEvent):

        #self._user_activity_true()

        d = e.GetWheelDelta()
        r = e.GetWheelRotation()
        a = e.GetWheelAxis()

        altdown = e.AltDown()
        ctrldown = e.ControlDown()
        shiftdown = e.ShiftDown()
        spacedown = wx.GetKeyState(wx.WXK_SPACE)

        if self.action == 'dynamic parallel' and shiftdown and not ctrldown:
            self.dynapar_dist *= (1 - .1 * (r / max(d, 1)))
            self.dynapar_dist = max(self.dynapar_dist, .01)

            self.active_zone.parallel_active(self.dynapar_dist)
            self.Refresh()
            return
        elif self.action == 'dynamic parallel' and shiftdown and ctrldown:
            dlg = wx.NumberEntryDialog(None,
                                       _('What is the desired size [cm] ?'),
                                       'ds',
                                       'ds size',
                                       int(self.dynapar_dist * 100.),
                                       1,
                                       100000)
            ret = dlg.ShowModal()
            if ret == wx.ID_CANCEL:
                dlg.Destroy()
                return

            self.dynapar_dist = float(dlg.GetValue()) / 100.
            self.dynapar_dist = max(self.dynapar_dist, .01)
            dlg.Destroy()

            self.active_zone.parallel_active(self.dynapar_dist)
            self.Refresh()
            return

        elif shiftdown:

            if self.active_vector is None:
                logging.warning(_('No vector selected -- Please select a vector first !'))
                return

            if self.active_vector.myprop.textureimage is None:
                logging.warning(_('No image available -- Please load an image first !'))
                return

            self.active_vector.myprop.image_scale /= (1 - .1 * (r / max(d, 1)))
            # limit to 1. or upper
            self.active_vector.myprop.image_scale = max(self.active_vector.myprop.image_scale, 1.)
            self.active_vector.myprop.update_myprops()
            self.active_vector.myprop.update_image_texture()
            self.Refresh()
            return

        # Allow the user to zoom onto the pixel where the
        # mouse cursor is

        # Step1: move the map so that the pixem under the mouse cursor
        # ends up right in the middle of the screen (this move is
        # not visible from the end user point of view, it's just
        # here to make computation seasier)
        if spacedown:
            self.center_view_on( *self.getXY( e.GetPosition()))

        # Zoom/dezoom, center pf the tranfromation is the center of the screen
        self.width = self.width * (1 - .1 * (r / max(d, 1)))
        self.height = self.height * (1 - .1 * (r / max(d, 1)))

        if spacedown:
            self.updatescalefactors() # not base on mousex

            # Translate back the pixel at the center of the screen to where the
            # mouse cursor is. For that we measure the delta in screen coordinates
            # and transform it to map space coordinates.
            x_mid, y_mid = self.canvas.GetSize()
            x_mid, y_mid = self.getXY((0.5*x_mid, y_mid*0.5))
            x, y = self.getXY( e.GetPosition())
            dx, dy = x_mid - x, y_mid - y
            self.mousex +=  dx
            self.mousey +=  dy

        # will translate and rescale the map view so that it fits the window.
        self.setbounds()

    def On_Right_Double_Clicks(self, e):

        pos = e.GetPosition()
        ctrldown = e.ControlDown()
        altdown = e.AltDown()
        shiftdown = e.ShiftDown()
        x, y = self.getXY(pos)

        self._endactions()

    def On_Left_Double_Clicks(self, e:wx.MouseEvent):
        pos = e.GetPosition()

        ctrldown = e.ControlDown()
        altdown = e.AltDown()
        shiftdown = e.ShiftDown()

        x, y = self.getXY(pos)

        if self.mousedown[0] == -99999:
            # Manage the case when the user double click rapidly on the map
            # self.mousedown can not be set in On_Mouse_Left_Down
            self.mousedown = (x, y)

        self.mousex, self.mousey = self.mousedown
        self.oneclick = False
        self.setbounds()

        if shiftdown:
            if self.active_array is not None:
                if self.active_viewer3d is not None:
                    self.active_viewer3d.force_view(self.mousex, self.mousey, self.active_array.get_value(self.mousex, self.mousey))
                    self.Refresh()

                if self.active_laz is not None:
                    if self.active_laz.viewer is not None:
                        self.active_laz.force_view(self.mousex, self.mousey, self.active_array.get_value(self.mousex, self.mousey))
            else:
                if self.active_viewer3d is not None:
                    self.active_viewer3d.force_view(self.mousex, self.mousey)
                    self.Refresh()

                if self.active_laz is not None:
                    if self.active_laz.viewer is not None:
                        self.active_laz.force_view(self.mousex, self.mousey)

    def On_Mouse_Left_Down(self, e):
        """ Event when the left button of the mouse is pressed """
        #self._user_activity_true()

        # if not self.move:
        pos = e.GetPosition()
        x, y = self.getXY(pos)
        self.mousedown = (x, y)
        # self.move = True

    def On_Mouse_Left_Up(self, e):
        """ Event when the left button of the mouse is released """

        pass
        # if self._thread_update_background is None:
        #     self._thread_update_background = threading.Thread(target=self.background_task, args=[self]) # thread to update the background
        #     self._thread_update_background.start() # start the thread

    def _set_active_bc(self):
        """Search and activate BCManager according to active_array"""
        if self.active_bc is not None:
            if self.active_array != self.active_bc.linked_array:
                # it is not the good one -> Hide
                self.active_bc.Hide()
            else:
                return
        # searching if bcmanager is attached to active_array
        self.active_bc = None
        for curbc in self.mybc:
            if self.active_array == curbc.linked_array:
                self.active_bc = curbc
                self.active_bc.Show()
                return

    def set_statusbar_text(self, txt:str):
        """ Set the status bar text """
        self.StatusBar.SetStatusText(txt)

    def set_label_selecteditem(self, nameitem:str):
        """ Set the label of the selected item in the tree list """
        self._lbl_selecteditem.SetLabel(nameitem)

    def get_label_selecteditem(self):
        """ Get the label of the selected item in the tree list """
        return self._lbl_selecteditem.GetLabel()

    def OnActivateTreeElem(self, e): #:dataview.TreeListEvent ):
        """ Activate the selected item in the tree list """
        curzones: Zones
        curzone: zone
        curvect: vector

        myitem = e.GetItem()
        ctrl = wx.GetKeyState(wx.WXK_CONTROL)
        alt = wx.GetKeyState(wx.WXK_ALT)

        myparent = self.treelist.GetItemParent(myitem)
        check = self.treelist.GetCheckedState(myitem)


        nameparent = self.treelist.GetItemText(myparent).lower()
        nameitem = self.treelist.GetItemText(myitem).lower()

        myobj = self.treelist.GetItemData(myitem)
        self.selected_object = myobj
        self.set_label_selecteditem(_('Active : ') + nameitem)

        #FIXME : To generalize using draw_type
        if type(myobj) == Zones:
            self.active_zones = myobj

            if ctrl:
                myobj.show_properties()

        elif type(myobj) == PictureCollection:
            self.active_picturecollection = myobj

            if ctrl:
                myobj.show_properties()

        elif type(myobj) == Wolf_LAZ_Data:

            self.active_laz = myobj

            if ctrl:
                myobj.show_properties()

        elif type(myobj) == Bridge:
            self.active_bridge = myobj

            if ctrl:
                myobj.show_properties()

        elif type(myobj) == Weir:
            self.active_weir = myobj

            if ctrl:
                myobj.show_properties()

        elif isinstance(myobj, PlansTerrier):
            self.active_landmap = myobj

        elif isinstance(myobj, Particularites | Enquetes | Ouvrages | Profils):
            self.active_picturecollection = myobj

        elif type(myobj) == hydrometry_wolfgui:
            if ctrl:
                myobj.show_properties()

        elif type(myobj) in [Picc_data, Cadaster_data]:
            if ctrl:
                myobj.show_properties()

        elif type(myobj) == Particle_system:
            if ctrl:
                myobj.show_properties()

        elif type(myobj) == Tiles:
            self.active_tile= myobj

        elif issubclass(type(myobj), WolfArray):
            if ctrl:
                myobj.show_properties()
                # myobj.myops.SetTitle(_('Operations on array: ')+myobj.idx)
                # myobj.myops.Show()

            logging.info(_('Activating array : ' + nameitem))
            self.active_array = myobj

            # If BC maneger is attached to the array, we activate it
            self._set_active_bc()

            #Print info in the status bar
            txt  = 'Dx : {:.4f} ; Dy : {:.4f}'.format(self.active_array.dx, self.active_array.dy)
            txt += ' ; Xmin : {:.4f} ; Ymin : {:.4f}'.format(self.active_array.origx, self.active_array.origy)
            txt += ' ; Xmax : {:.4f} ; Ymax : {:.4f}'.format(self.active_array.origx + self.active_array.dx * float(self.active_array.nbx),
                                                           self.active_array.origy + self.active_array.dy * float(self.active_array.nby))
            txt += ' ; Nx : {:d} ; Ny : {:d}'.format(self.active_array.nbx, self.active_array.nby)

            if self.active_array.nb_blocks > 0:
                txt += ' ; Nb blocks : {:d}'.format(self.active_array.nb_blocks)

            txt += ' ; Type : ' + self.active_array.dtype_str

            self.set_statusbar_text(txt)


        elif type(myobj) in [WolfViews]:
            logging.info(_('Activating view : ' + nameitem))
            self.active_view = myobj

        elif type(myobj) == cloud_vertices:
            self.active_cloud = myobj
            if ctrl:
                myobj.myprop.show()

        elif type(myobj) == crosssections:
            if ctrl:
                myobj.showstructure()
            logging.info(_('Activating cross sections : ' + nameitem))
            self.active_cs = myobj

        elif type(myobj) == Triangulation:
            self.active_tri = myobj

        elif type(myobj) == Wolfresults_2D:
            logging.info(_('Activating Wolf2d results : ' + nameitem))
            self.active_res2d = myobj

            if ctrl:
                myobj.show_properties()

            if alt:
                dlg = wx.MessageDialog(self,_('Do you want to open the 2D model?'),style=wx.YES_NO|wx.NO_DEFAULT)
                ret=dlg.ShowModal()
                dlg.Destroy()
                if ret == wx.ID_NO:
                    return

                from .PyGui import Wolf2DModel
                mywolf = Wolf2DModel(dir=os.path.dirname(self.active_res2d.filenamegen), splash=False)

        elif type(myobj) == wolfres2DGPU:
            logging.info(_('Activating Wolf2d results : ' + nameitem))
            self.active_res2d = myobj

            if ctrl:
                myobj.show_properties()

        elif type(myobj) == Drowning_victim_Viewer:
            logging.info(_('Activating Drowning victim event : ' + nameitem))
            self.active_drowning = myobj

        elif WOLFPYDIKE_AVAILABLE:
            if type(myobj) == DikeWolf:
                logging.info(_('Activating DikeWolf : ' + nameitem))
                self.active_dike = myobj
                if myobj.injector is not None:
                    self.active_injector = myobj.injector
                    logging.info(_('Activating InjectorDike : ' + nameitem))

                if ctrl:
                    myobj.show_properties()

            elif type(myobj) == InjectorDike:
                logging.info(_('Activating InjectorDike : ' + nameitem))
                self.active_injector = myobj

                if ctrl:
                    myobj.show_properties()

    def _update_mytooltip(self):
        """ Update the tooltip with the values of the active arrays and results at position x,y """

        x,y,pos = self._last_mouse_pos

        self.mytooltip.myparams.clear()


        curgroup = 'Position'
        self.mytooltip.myparams[curgroup] = {}

        curpar = _('Pixel (col,row)')
        self.mytooltip.add_param(groupname = curgroup,
                                 name = curpar,
                                 value = '(' + str(pos[0]) + ' ; ' + str(pos[1]) + ')',
                                 type = Type_Param.String,
                                 comment = '')

        curpar = _('Coordinate X [m]')
        self.mytooltip.add_param(groupname = curgroup,
                                 name = curpar,
                                 value = '{:3f}'.format(x),
                                 type = Type_Param.String,
                                 comment = '')

        curpar = _('Coordinate Y [m]')
        self.mytooltip.add_param(groupname = curgroup,
                                 name = curpar,
                                 value = '{:3f}'.format(y),
                                 type = Type_Param.String,
                                 comment = '')

        if self._tmp_vector_distance is not None:
            curgroup = _('Temporary vector')
            self.mytooltip.myparams[curgroup] = {}
            curpar = _('Length [m]')
            self._tmp_vector_distance.update_lengths()
            self.mytooltip.add_param(groupname = curgroup,
                                        name = curpar,
                                        value = '{:3f}'.format(self._tmp_vector_distance.length2D),
                                        type = Type_Param.Float,
                                        comment = '')

            if self._tmp_vector_distance.nbvertices > 4:
                _polygon = self._tmp_vector_distance.asshapely_pol()
                _area = _polygon.area
                curpar = _('Area [m2]')
                self.mytooltip.add_param(groupname = curgroup,
                                            name = curpar,
                                            value = '{:3f}'.format(_area),
                                            type = Type_Param.Float,
                                            comment = '')
                curpar = _('Area [ha]')
                self.mytooltip.add_param(groupname = curgroup,
                                            name = curpar,
                                            value = '{:3f}'.format(_area / 10000.),
                                            type = Type_Param.Float,
                                            comment = '')
                curpar = _('Area [km2]')
                self.mytooltip.add_param(groupname = curgroup,
                                            name = curpar,
                                            value = '{:3f}'.format(_area / 1e6),
                                            type = Type_Param.Float,
                                            comment = '')

        for locarray in self.myres2D:
            locarray:Wolfresults_2D
            curgroup = locarray.idx
            if locarray.checked:
                try:
                    vals,labs = locarray.get_values_labels(x,y)

                    i, j, curbloc = locarray.get_blockij_from_xy(x, y)

                    if i != '-' and i != -1:
                        curpar = 'Indices (i;j;bloc) (1-based)'

                        self.mytooltip.add_param(groupname = curgroup,
                                                 name = curpar,
                                                 value =  '(' + str(i) + ';' + str(j) + ';' + str(curbloc) + ')',
                                                 type = Type_Param.String,
                                                 comment = '')

                        for val,curpar in zip(vals,labs):

                            if val is np.nan:
                                self.mytooltip.add_param(groupname = curgroup,
                                                        name = 'Value',
                                                        value =  "Nan",
                                                        type = Type_Param.String,
                                                        comment = '')

                            elif np.ma.is_masked(val):

                                self.mytooltip.add_param(groupname = curgroup,
                                                        name = 'Value',
                                                        value =  "Masked",
                                                        type = Type_Param.String,
                                                        comment = '')

                            elif isinstance(val, str):
                                self.mytooltip.add_param(groupname = curgroup,
                                                        name = curpar,
                                                        value =  val,
                                                        type = Type_Param.String,
                                                        comment = '')

                            elif isinstance(val, int):
                                self.mytooltip.add_param(groupname = curgroup,
                                                        name = curpar,
                                                        value =  int(val),
                                                        type = Type_Param.Integer,
                                                        comment = '')

                            else:
                                self.mytooltip.add_param(groupname = curgroup,
                                                        name = curpar,
                                                        value =  float(val),
                                                        type = Type_Param.Float,
                                                        comment = '')

                except:
                    pass

        for locarray in self.myarrays:
            if locarray.checked:
                curgroup = locarray.idx

                try:
                    val = locarray.get_value(x, y)

                    if val != -99999.:

                        if locarray.wolftype in WOLF_ARRAY_MB:
                            i, j, curbloc = locarray.get_blockij_from_xy(x, y)
                            curpar = 'Indices (i;j;bloc) (1-based)'

                            self.mytooltip.add_param(groupname = curgroup,
                                                        name = curpar,
                                                        value =  '(' + str(i+1) + ';' + str(j+1) + ';' + str(curbloc) + ')',
                                                        type = Type_Param.String,
                                                        comment = '')

                        else:
                            i, j = locarray.get_ij_from_xy(x, y)
                            curpar = 'Indices (i;j) (1-based)'

                            self.mytooltip.add_param(groupname = curgroup,
                                                        name = curpar,
                                                        value =  '(' + str(i+1) + ';' + str(j+1) + ')',
                                                        type = Type_Param.String,
                                                        comment = '')

                        curpar = 'Value'

                        if val is np.nan:
                            self.mytooltip.add_param(groupname = curgroup,
                                                    name = 'Value',
                                                    value =  "Nan",
                                                    type = Type_Param.String,
                                                    comment = '')
                        elif np.ma.is_masked(val):

                            self.mytooltip.add_param(groupname = curgroup,
                                                    name = 'Value',
                                                    value =  "Masked",
                                                    type = Type_Param.String,
                                                    comment = '')

                        elif isinstance(val, str):
                            self.mytooltip.add_param(groupname = curgroup,
                                                    name = curpar,
                                                    value =  val,
                                                    type = Type_Param.String,
                                                    comment = '')

                        elif isinstance(val, int):
                            self.mytooltip.add_param(groupname = curgroup,
                                                    name = curpar,
                                                    value =  int(val),
                                                    type = Type_Param.Integer,
                                                    comment = '')

                        else:
                            self.mytooltip.add_param(groupname = curgroup,
                                                    name = curpar,
                                                    value =  float(val),
                                                    type = Type_Param.Float,
                                                    comment = '')

                except:
                    pass

        if self.linked:
            for curFrame in self.linkedList:
                if not curFrame is self:
                    title = curFrame.GetTitle() if curFrame.GetTitle() != 'Wolf - main data manager' else 'Main'

                    for locarray in curFrame.myarrays:
                        curgroup = title + ' -' + locarray.idx
                        if locarray.plotted:

                            try:
                                val = locarray.get_value(x, y)

                                if val != -99999.:
                                    if locarray.wolftype in WOLF_ARRAY_MB:
                                        i, j, curbloc = locarray.get_blockij_from_xy(x, y)
                                        curpar = 'Indices (i;j;bloc) (1-based)'

                                        self.mytooltip.add_param(groupname = curgroup,
                                                                    name = curpar,
                                                                    value =  '(' + str(i+1) + ';' + str(j+1) + ';' + str(curbloc) + ')',
                                                                    type = Type_Param.String,
                                                                    comment = '')

                                    else:
                                        i, j = locarray.get_ij_from_xy(x, y)
                                        curpar = 'Indices (i;j) (1-based)'

                                        self.mytooltip.add_param(groupname = curgroup,
                                                                    name = curpar,
                                                                    value =  '(' + str(i+1) + ';' + str(j+1) + ')',
                                                                    type = Type_Param.String,
                                                                    comment = '')

                                    curpar = 'Value'

                                    if val is np.nan:
                                        self.mytooltip.add_param(groupname = curgroup,
                                                                name = 'Value',
                                                                value =  "Nan",
                                                                type = Type_Param.String,
                                                                comment = '')

                                    elif np.ma.is_masked(val):

                                        self.mytooltip.add_param(groupname = curgroup,
                                                                name = 'Value',
                                                                value =  "Masked",
                                                                type = Type_Param.String,
                                                                comment = '')

                                    elif isinstance(val, str):
                                        self.mytooltip.add_param(groupname = curgroup,
                                                                name = curpar,
                                                                value =  val,
                                                                type = Type_Param.String,
                                                                comment = '')
                                    elif isinstance(val, int):
                                        self.mytooltip.add_param(groupname = curgroup,
                                                                name = curpar,
                                                                value =  int(val),
                                                                type = Type_Param.Integer,
                                                                comment = '')
                                    else:
                                        self.mytooltip.add_param(groupname = curgroup,
                                                                name = curpar,
                                                                value =  float(val),
                                                                type = Type_Param.Float,
                                                                comment = '')
                            except:
                                logging.warning(_('Error in linked frame Arrays -- Please check !'))

                    for locarray in curFrame.myres2D:
                        locarray:Wolfresults_2D
                        curgroup = title + ' - ' + locarray.idx
                        if locarray.checked:
                            try:
                                vals,labs = locarray.get_values_labels(x,y)

                                i, j, curbloc = locarray.get_blockij_from_xy(x, y)

                                if i != '-' and i != -1:
                                    curpar = 'Indices (i;j;bloc) (1-based)'

                                    self.mytooltip.add_param(groupname = curgroup,
                                                            name = curpar,
                                                            value =  '(' + str(i) + ';' + str(j) + ';' + str(curbloc) + ')',
                                                            type = Type_Param.String,
                                                            comment = '')

                                    for val,curpar in zip(vals,labs):

                                        if val is np.nan:
                                            self.mytooltip.add_param(groupname = curgroup,
                                                                    name = 'Value',
                                                                    value =  "Nan",
                                                                    type = Type_Param.String,
                                                                    comment = '')

                                        elif np.ma.is_masked(val):

                                            self.mytooltip.add_param(groupname = curgroup,
                                                                    name = 'Value',
                                                                    value =  "Masked",
                                                                    type = Type_Param.String,
                                                                    comment = '')

                                        elif isinstance(val, str):
                                            self.mytooltip.add_param(groupname = curgroup,
                                                                    name = curpar,
                                                                    value =  val,
                                                                    type = Type_Param.String,
                                                                    comment = '')

                                        elif isinstance(val, int):
                                            self.mytooltip.add_param(groupname = curgroup,
                                                                    name = curpar,
                                                                    value =  int(val),
                                                                    type = Type_Param.Integer,
                                                                    comment = '')

                                        else:
                                            self.mytooltip.add_param(groupname = curgroup,
                                                                    name = curpar,
                                                                    value =  float(val),
                                                                    type = Type_Param.Float,
                                                                    comment = '')

                            except:
                                logging.warning(_('Error in linked frame Results2D -- Please check !'))

        for loc_ps in self.mypartsystems:
            if loc_ps.checked:
                curgroup = loc_ps.idx
                try:
                    self.mytooltip.add_param(groupname = curgroup,
                                            name = _('Step [s]'),
                                            value =  loc_ps.current_step,
                                            type = Type_Param.Float,
                                            comment = 'Step in seconds')
                    self.mytooltip.add_param(groupname = curgroup,
                                            name = _('Step [-]'),
                                            value =  loc_ps.current_step_idx,
                                            type = Type_Param.Integer,
                                            comment = 'Step index')
                except:
                    pass

        for loc_drowning in self.mydrownings:
            if loc_drowning.checked:
                try:
                    if loc_drowning.bottom_cells is not None:
                        curgroup = loc_drowning.idx
                        try:
                            i_bottom, j_bottom = loc_drowning.bottom_cells.get_ij_from_xy(x, y)
                            i_surface, j_surface = loc_drowning.surface_cells.get_ij_from_xy(x, y)
                        except:
                            pass
                        value = loc_drowning.bottom_cells.array[i_bottom,j_bottom]
                        if not np.ma.is_masked(value) and value >0:
                            self.mytooltip.add_param(groupname = curgroup,
                                                        name = 'Bottom - percentage of the whole sample',
                                                        value =  value,
                                                        type = Type_Param.Float,
                                                        comment = '')
                        value = loc_drowning.surface_cells.array[i_surface,j_surface]
                        if not np.ma.is_masked(value) and value >0:
                            self.mytooltip.add_param(groupname = curgroup,
                                                        name = 'Surface - percentage of the whole sample',
                                                        value =  value,
                                                        type = Type_Param.Float,
                                                        comment = '')
                except:
                    pass


        self.mytooltip.PopulateOnePage()

    def On_Mouse_Motion(self, e: wx.MouseEvent):
        """ Mouse move event """


        #self._user_activity_true()

        # Déplacement de la souris sur le canvas OpenGL
        posframe = self.GetPosition() # Get the position of the frame -> useful to set the tooltip position
        pos = e.GetPosition() # Get the position of the mouse in the canvas
        x, y = self.getXY(pos) # Get the coordinates of the mouse in the map
        altdown = e.AltDown()       # Alt key pressed
        shiftdown = e.ShiftDown()   # Shift key pressed

        if e.LeftIsDown() or e.MiddleIsDown():
            # Left mouse button or middle mouse button is pressed
            #
            # Moving the map relative to the position where the mouse was clicked
            # the first time

            if self.mousedown[0] == -99999: # only if the mouse was clicked before
                self.mousedown = [x, y]

            if shiftdown:
                if self.active_vector is None:
                    logging.warning(_('Shift key pressed but no active vector -- Please select a vector first !'))
                    return
                if self.active_vector.myprop.textureimage is None:
                    logging.warning(_('Shift key pressed but no image texture -- Please select a vector with an image first !'))
                    return
                # We move the image texture
                delta_x = x - self.mousedown[0]
                delta_y = y - self.mousedown[1]
                self.active_vector.myprop._offset_image_texture(delta_x, delta_y)
                self.active_vector.myprop.update_myprops()
                self.active_vector.myprop.update_image_texture()
                self.Refresh()
                return

            self.mousex -= x - self.mousedown[0]
            self.mousey -= y - self.mousedown[1]

            self.setbounds(updatescale = False)
            return

        elif e.RightIsDown():
            # Right mouse button is pressed
            if self.action == 'select bc':
                if self.active_vector is None:

                    self.end_action(_('None because no active vector'))
                    return

                self.active_vector.myvertices=[wolfvertex(self.rightdown[0],self.rightdown[1]),
                                               wolfvertex(self.rightdown[0],y),
                                               wolfvertex(x,y),
                                               wolfvertex(x,self.rightdown[1]),
                                               wolfvertex(self.rightdown[0],self.rightdown[1])]
        else:
            # No mouse button is pressed
            # self.move=False
            self.mousedown = (-99999, -99999)
            # self.mouseup   = (-99999, -99999)

        if self.action is not None:

            if 'select by tmp vector' in self.action or \
               'select by vector' in self.action or \
               self.action == 'capture vertices' or \
               self.action == 'dynamic parallel' or \
               self.action == 'laz tmp vector' or \
               self.action == 'create polygon - tiles':

                if self.active_vector.nbvertices > 0:
                    self.active_vector.myvertices[-1].x = x
                    self.active_vector.myvertices[-1].y = y

                    self.active_vector.reset_linestring()
                    self.active_vector.parentzone.reset_listogl()

            if self.action == 'modify vertices' or \
               self.action == 'insert vertices':
                if self.active_vertex is not None:
                    if shiftdown:
                        # Shift key is pressed
                        # We move/Insert the vertex along the segment linking the first and last vertices of the active vector
                        ox = self.active_vector.myvertices[0].x
                        oy = self.active_vector.myvertices[0].y

                        dirx = self.active_vector.myvertices[-1].x - ox
                        diry = self.active_vector.myvertices[-1].y - oy
                        normdir = np.sqrt(dirx ** 2. + diry ** 2.)

                        dirx /= normdir
                        diry /= normdir

                        vecx = x - ox
                        vecy = y - oy

                        norm = np.sqrt(vecx ** 2. + vecy ** 2.)

                        self.active_vertex.x = ox + np.inner([dirx, diry], [vecx, vecy]) * dirx
                        self.active_vertex.y = oy + np.inner([dirx, diry], [vecx, vecy]) * diry

                    else:
                        self.active_vertex.x = x
                        self.active_vertex.y = y

                    self.active_vertex.limit2bounds(self.active_vector._mylimits)

            elif self.action == 'dynamic parallel':
                self.active_zone.parallel_active(self.dynapar_dist)

            elif self.action == 'move vector':
                if self.active_vector is not None:
                    if self.active_vector._move_start is not None:

                        delta_x = x - self.active_vector._move_start[0]
                        delta_y = y - self.active_vector._move_start[1]

                        if shiftdown:
                            delta_y = 0.

                        if altdown:
                            delta_x = 0.

                        self.active_vector.move(delta_x, delta_y)

            elif self.action == 'move triangles':
                if self.active_tri is not None:
                    if self.active_tri._move_start is not None:

                        delta_x = x - self.active_tri._move_start[0]
                        delta_y = y - self.active_tri._move_start[1]

                        if shiftdown:
                            delta_y = 0.

                        if altdown:
                            delta_x = 0.

                        self.active_tri.move(delta_x, delta_y)

                        self.active_tri.reset_plot()

            elif self.action == 'rotate vector':
                if self.active_vector is not None:
                    if self.active_vector._rotation_center is not None:
                        self.active_vector.rotate_xy(x, y)

            elif self.action == 'rotate triangles':
                if self.active_tri is not None:
                    if self.active_tri._rotation_center is not None:
                        self.active_tri.rotate_xy(x, y)
                        self.active_tri.reset_plot()

            elif self.action == 'move zone':
                if self.active_zone is not None:
                    if self.active_zone._move_start is not None:
                        delta_x = x - self.active_zone._move_start[0]
                        delta_y = y - self.active_zone._move_start[1]

                        if shiftdown:
                            delta_y = 0.

                        if altdown:
                            delta_x = 0.

                        self.active_zone.move(delta_x, delta_y)

            elif self.action == 'rotate zone':
                if self.active_zone is not None:
                    if self.active_zone._rotation_center is not None:
                        self.active_zone.rotate_xy(x, y)

            elif self.action == 'distance along vector':
                if self._tmp_vector_distance is not None:
                    self._tmp_vector_distance.myvertices[-1].x = x
                    self._tmp_vector_distance.myvertices[-1].y = y

                    if self._tmp_vector_distance.nbvertices ==2:
                        self._tmp_vector_distance.myvertices[0].x = x
                        self._tmp_vector_distance.myvertices[0].y = y

            self.Paint()

        # Store the position of the mouse as last known position
        self._last_mouse_pos = (x,y,pos)

        if self.active_vector is not None:
            if self.active_vector.myprop.textureimage is not None:
                self.active_vector.myprop._reset_cached_offset()

        # Update the tooltip with the values of the active arrays and results at position x,y
        self._update_mytooltip()

        if e.ControlDown():
            # Control key is pressed, so the tooltip must be near the mouse position
            if self._oldpos_tooltip is None:
                # Store the position of the tooltip
                # Useful to restore it after CTRL is released
                self._oldpos_tooltip = self.mytooltip.GetPosition()

            try:
                self.mytooltip.SetWindowStyle(wx.STAY_ON_TOP) # Just on top, without Title bar
                ttsize = self.mytooltip.GetSize()
                if ttsize[0] == 0 or ttsize[1] == 0:
                    pass
                self.mytooltip.position(pos + posframe + (ttsize[0] / 2. + 15, 15))
            except Exception as e:
                logging.warning(_('Error in tooltip positionning : ') + str(e))
                self.mytooltip.Show(True)
        else:

            width, height = self.GetSize()

            if self.IsMaximized():
                # Frame is maximized -> tooltip must be on the Screen
                self.mytooltip.SetWindowStyle(wx.STAY_ON_TOP | wx.DEFAULT_FRAME_STYLE)
            else:

                if self._oldpos_tooltip is None:
                    # No old position stored -> tooltip does not move
                    pos_tooltip = self.mytooltip.GetPosition()
                else:
                    # Restore the position of the tooltip
                    pos_tooltip = self._oldpos_tooltip

                # Reset the old position, so when CTRL is pressed again, the memory will be updated
                self._oldpos_tooltip = None

                if shiftdown or (pos_tooltip[0] == 0 and pos_tooltip[1] == 0):
                    # SHIFT is pressed or tooltip is at the top right corner of the Frame
                    # or it is the first time the tooltip is displayed
                    posframe[0] += width
                    posframe[1] -= 50
                    self.mytooltip.position(posframe)
                    w, h = self.mytooltip.GetSize()
                    self.mytooltip.SetSize((w, height))

                else:
                    # Force the position
                    self.mytooltip.SetPosition(pos_tooltip)

                # self.mytooltip.SetIcon(self.GetIcon()) # update icon
                self.mytooltip.SetWindowStyle(wx.DEFAULT_FRAME_STYLE) # | wx.STAY_ON_TOP) # on top, with Title bar

            #Force to show the tooltip --> useful when the tooltip was hidden by the user
            self.mytooltip.Show(True)

        e.Skip()

    def Autoscale(self, update_backfore=True):
        """ Redimensionnement de la fenêtre pour afficher tous les objets """

        self.findminmax()
        self.width = self.xmax - self.xmin
        self.height = self.ymax - self.ymin

        centerx = self.xmin + self.width / 2.
        centery = self.ymin + self.height / 2.

        iwidth = self.width * self.sx
        iheight = self.height * self.sy

        width, height = self.canvas.GetSize()

        if iwidth == 0 or iheight == 0:
            logging.warning(_('Width or height of the canvas is null -- Please check the "findminmax" routine in "Autoscale" !'))
            iwidth = 1
            iheight = 1

        sx = float(width) / float(iwidth)
        sy = float(height) / float(iheight)

        if sx == 0 or sy == 0:
            logging.error(_('At least one scale factor is null -- Please check the "Autoscale" routine !'))
            sx = 1.
            sy = 1.

        if sx > sy:
            self.xmax = self.xmin + self.width * sx / sy
            self.width = self.xmax - self.xmin
        else:
            self.ymax = self.ymin + self.height * sy / sx
            self.height = self.ymax - self.ymin

        self.mousex = centerx
        self.mousey = centery

        if update_backfore:
            # dessin du background
            for obj in self.iterator_over_objects(draw_type.WMSBACK):
                obj.reload()

            # dessin du foreground
            for obj in self.iterator_over_objects(draw_type.WMSFORE):
                obj.reload()

        self.setbounds()

    def _endactions(self):
        """
        End of actions

        Call when the user double click on the right button of the mouse or press return.

        Depending on the action, the method will call differnt routines and refresh the figure.

        Each action must call self.end_action() to nullify the action and print a message.
        """

        if self.action is not None:
            locaction = self.action

            if 'select by tmp vector' in locaction or 'select by vector' in locaction:
                inside_under = 'inside' in self.action
                outside_under = 'outside' in self.action

                self.end_action(_('End of vector selection'))

                self.active_vector.myvertices.pop(-1)

                if inside_under:
                    self.active_vector.close_force()
                    self.active_array.SelectionData.select_insidepoly(self.active_vector)
                elif outside_under:
                    self.active_vector.close_force()
                    self.active_array.SelectionData.select_outsidepoly(self.active_vector)
                else:
                    self.active_array.SelectionData.select_underpoly(self.active_vector)

                if 'tmp' in locaction:
                    # we must reset the temporary vector
                    self.active_vector.reset()

            elif locaction == 'distance along vector':

                dlg = wx.MessageDialog(None, _('Memorize the vector ?'), _('Confirm'), wx.YES_NO | wx.YES_DEFAULT | wx.ICON_QUESTION)
                ret = dlg.ShowModal()
                dlg.Destroy()

                if ret == wx.ID_YES:
                    self._distances[-1].add_vector(self._tmp_vector_distance, forceparent=True, update_struct=True)

                self._tmp_vector_distance = None

            elif locaction == 'pick landmap':

                self.end_action(_('End of landmap picking'))

            elif locaction == 'laz tmp vector':
                self.end_action(_('End of LAZ selection'))
                self.active_vector.myvertices.pop(-1)
                self.plot_laz_around_active_vec()
                self.active_vector.reset()

            elif locaction == 'create polygon - tiles':
                self.end_action(_('End of polygon creation'))
                self.active_vector.myvertices.pop(-1)
                self.active_vector.close_force()

                dlg = wx.MessageDialog(None, _('Do you want to align vertices on magnetic grid ?'), _('Confirm'), wx.YES_NO | wx.YES_DEFAULT | wx.ICON_QUESTION)
                ret = dlg.ShowModal()
                dlg.Destroy()

                if ret == wx.ID_YES:
                    dlg = wx.NumberEntryDialog(None, _('Which is the sptial step size [m] ?'), _('Size'),  _('Spatial grid size'), 50, 1, 10000)
                    ret = dlg.ShowModal()
                    if ret == wx.ID_CANCEL:
                        dlg.Destroy()
                        return
                    ds = dlg.GetValue()
                    dlg.Destroy()

                    vertices = self.active_vector.myvertices.copy()
                    self.active_vector.myvertices.clear()

                    x_aligned = np.asarray([(curvert.x // ds)*ds for curvert in vertices])
                    y_aligned = np.asarray([(curvert.y // ds)*ds for curvert in vertices])

                    if (x_aligned.min() == x_aligned.max()) and (y_aligned.min() == y_aligned.max()):
                        logging.error(_('All vertices are aligned on the same point -- Choose another step size'))
                        return

                    self.active_vector.add_vertex(wolfvertex(x_aligned.min(), y_aligned.min()))
                    self.active_vector.add_vertex(wolfvertex(x_aligned.max(), y_aligned.min()))
                    self.active_vector.add_vertex(wolfvertex(x_aligned.max(), y_aligned.max()))
                    self.active_vector.add_vertex(wolfvertex(x_aligned.min(), y_aligned.max()))
                    self.active_vector.add_vertex(wolfvertex(x_aligned.min(), y_aligned.min()))
                    self.active_vector.close_force()
                    self.active_vector.find_minmax()

                self._create_data_from_tiles_common()

            elif locaction == 'capture vertices':
                self.end_action(_('End of points capturing'))
                self.active_vector.myvertices.pop(-1)
                r = wx.MessageDialog(
                    None,
                    _('End of points capturing') + '\n' +
                    _('Force to close the vector ?'),
                    _('Confirm'),
                    wx.YES_NO | wx.YES_DEFAULT | wx.ICON_QUESTION
                ).ShowModal()
                if r == wx.ID_YES:
                    self.active_vector.close_force()

                # force to prepare OpenGL to accelerate the plot
                # Le test not(self in self.linkedList) permet de ne pas créer le liste OpenGL en cas de multi-viewers
                # car une liste OpenGL ne sera pas tracée sur les autres fenêtres
                # C'est donc plus lent mais plus sûr pour que l'affichage dynamique soit correct
                self.active_vector.parentzone.plot(prep = self.linkedList is None or not(self in self.linkedList))

            elif locaction == 'modify vertices':

                # end of vertices modification
                self.end_action(_('End of vertices modification'))

                # force to prepare OpenGL to accelerate the plot
                # Le test not(self in self.linkedList) permet de ne pas créer le liste OpenGL en cas de multi-viewers
                # car une liste OpenGL ne sera pas tracée sur les autres fenêtres
                # C'est donc plus lent mais plus sûr pour que l'affichage dynamique soit correct
                self.active_vector.parentzone.plot(prep = self.linkedList is None or not(self in self.linkedList))
                self.active_zones.find_minmax(True)

                self.active_vertex = None

            elif locaction == 'insert vertices':
                self.end_action(_('End of vertices insertion'))

                # force to prepare OpenGL to accelerate the plot
                # Le test not(self in self.linkedList) permet de ne pas créer le liste OpenGL en cas de multi-viewers
                # car une liste OpenGL ne sera pas tracée sur les autres fenêtres
                # C'est donc plus lent mais plus sûr pour que l'affichage dynamique soit correct
                self.active_vector.parentzone.plot(prep = self.linkedList is None or not(self in self.linkedList))
                self.active_zones.find_minmax(True)

                self.active_vertex = None

            elif locaction == 'dynamic parallel':
                self.active_vector.myvertices.pop(-1)
                self.active_zone.parallel_active(self.dynapar_dist)

                self.active_zones.fill_structure()
                self.active_zones.find_minmax(True)

                # force to prepare OpenGL to accelerate the plot
                # Le test not(self in self.linkedList) permet de ne pas créer le liste OpenGL en cas de multi-viewers
                # car une liste OpenGL ne sera pas tracée sur les autres fenêtres
                # C'est donc plus lent mais plus sûr pour que l'affichage dynamique soit correct
                self.active_vector.parentzone.plot(prep = self.linkedList is None or not(self in self.linkedList))

                self.active_vertex = None

                self.end_action(_('End of dynamic parallel'))

            elif 'select active vector' in locaction:

                self.end_action(_('End of vector selection'))

            elif 'select node by node' in locaction:
                self.end_action(_('End of node by node selection'))

        self.copyfrom = None

        self.Refresh()
        self.mimicme()

    def print_About(self):
        """ Print the About window """
        from .apps.version import WolfVersion

        version = WolfVersion()
        dlg = wx.MessageDialog(None, _('Wolf - Version {}\n\n'.format(str(version))) + _('Developed by : ') + 'HECE ULiège\n' + _('Contact : pierre.archambeau@uliege.be'), _('About'), wx.OK | wx.ICON_INFORMATION)
        dlg.ShowModal()
        dlg.Destroy()

    def check_for_updates(self):
        """ Check for updates """
        from .apps.version import WolfVersion
        import requests
        import importlib.metadata

        # check_gpu = False
        # try:
        #     import wolfgpu
        #     check_gpu = True
        # except:
        #     pass

        msg = ''

        current_version = str(WolfVersion())
        package_name = "wolfhece"

        try:
            available_version = requests.get(f"https://pypi.org/pypi/{package_name}/json").json()["info"]["version"]

            if available_version > current_version:
                msg += _("A new version is available: {}\n\nYour version is {}\n\nIf you want to upgrade, 'pip install wolfhece --upgrade' from your Python environment.").format(available_version, current_version)
            else:
                msg += _("You have the latest version.")

        except Exception as e:
            logging.error("Package not found on PyPI. -- {}".format(e))

        # if check_gpu:
        #     # find the version
        #     package_name = "wolfgpu"
        #     current_version = importlib.metadata.version(package_name)
        #     url_wolfgpu = "https://gitlab.uliege.be/api/v4/projects/4180/packages/pypi/simple/" #+ package_name
        #     try:
        #         response = requests.get(url_wolfgpu).json()
        #         update_wolfgpu = False
        #         #parcourir toutes le entrées et ne conserver que la version la plus récente
        #         for one in response:
        #             one["version"] > current_version
        #             update_wolfgpu = True

        #         if update_wolfgpu:
        #             msg += '\n'
        #             msg += _("A new version of WolfGPU is available: {}.{}.{}, please update it.").format(current_version[0], current_version[1], current_version[2])
        #         else:
        #             msg += '\n'
        #             msg+= _("You have the latest version of WolfGPU.")

        #     except Exception as e:
        #         logging.error("Package not found on PyPI. -- {}".format(e))

        msg+= '\n\n'
        msg+= _('If you use wolfgpu, please check the GPU version independently.')

        with wx.MessageDialog(None, msg, _("Upgrade"), wx.OK | wx.ICON_INFORMATION) as dlg:
            dlg.ShowModal()

    def print_shortcuts(self, inframe:bool = None):
        """ Print the list of shortcuts into logging """

        # shortcuts = "F1 : mise à jour du dernier pas de résultat\n \
        #     F2 : mise à jour du résultat pas suivant\n \
        #     F4 : mise à jour du particle system au pas suivant\n \
        #     Shift+F2 : mise à jour du résultat pas précédent\n \
        #     Shift+F4 : mise à jour du particle system au pas précédent\n \
        #     CTRL+F2 : choix du pas\n \
        #     CTRL+F4 : choix du pas (particle system)\n \
        #     CTRL+Shift+F2 : choix du pas sur base du temps\n \
        #     CTRL+Shift+F4 : choix du pas sur base du temps (particle system)\n \
        #     F5 : autoscale\n \
        #     F7 : refresh\n \
        #     F8 : Zoom on Whole Walonia\n \
        #     F9 : sélection de toutes les mailles dans la matrice courante\n \
        #     F11 : sélection sur matrice courante\n \
        #     F12 : opération sur matrice courante\n \
        #     \n \
        #     ESPACE : pause/resume animation\n \
        #     \n \
        #     Z  : zoom avant\n \
        #     z  : zoom artrière\n \
        #     Flèches : déplacements latéraux\n \
        #     P : sélection de profil\n \
        #     1,2 : Transfert de la sélection de la amtrice courante vers le dictionnaire\n \
        #     F, CTRL+F : recherche de la polyligne dans la zone courante ou dans toutes les zones\n \
        #     i : interpolation2D sur base de la sélection sur la matrice courante\n \
        #     +,- (numpad) : augmente ou diminue la taille des flèches de resultats 2D\n \
        #     \n \
        #     o, O : Gestion de la transparence de la matrice courante\n \
        #     CTRL+o, CTRL+O : Gestion de la transparence du résultat courant\n \
        #     \n \
        #     !! ACTIONs !!\n \
        #     N : sélection noeud par noeud de la matrice courante\n \
        #     b, B : sélection par vecteur de la matrice courante - trace du vecteur\n \
        #     v, V : sélection par vecteur de la matrice courante - zone intérieure\n \
        #     r : reset de la sélection de la matrice courante\n \
        #     R : reset de toutes les sélections de la matrice courante\n \
        #     P : sélection de la section transversale par click souris\n \
        #     D : calcule de distance le long d'un vecteur temporaire\n \
        #     \n \
        #     RETURN : end current action (cf aussi double clicks droit 'OnRDClick')\n \
        #     DELETE : remove item\n \
        #     \n \
        #     CTRL+Q : Quit application\n \
        #     CTRL+U : Import GLTF/GLB\n \
        #     CTRL+C : Set copy source\n \
        #     CTRL+V : Paste selected values\n \
        #     CTRL+ALT+V ou ALTGr+V : Paste/Recopy selection\n \
        #     CTRL+L : chargement d'une matrice sur base du nom de fichier de la tile\n \
        #     \n \
        #     ALT+C : Copy image"
        #     ALT+SHIFT+C : Copy images from multiviewers \
        #     CTRL+ALT+SHIFT+C : Copy images from all arrays as independent image \

        groups = ['Results', 'Particle system', 'Drawing', 'Arrays', 'Cross sections', 'Zones', 'Action', 'Tree', 'Tiles', 'GLTF/GLB', 'App']

        shortcuts = {'F1': _('Results : read the last step'),
                     'F2': _('Results : read the next step'),
                     'Shift+F2': _('Results : read the previous step'),
                     'CTRL+F2': _('Results : choose the step'),
                     'CTRL+Shift+F2': _('Results : choose the step based on time'),
                     '+,- (numpad)': _('Results : increase or decrease the size of 2D result arrows'),

                     'F4': _('Particle system : update to the next step'),
                     'Shift+F4': _('Particle system : update to the previous step'),
                     'CTRL+F4': _('Particle system : choose the step'),
                     'CTRL+Shift+F4': _('Particle system : choose the step based on time'),
                     'SPACE': _('Particle system : pause/resume animation'),

                     'LMB double clicks': _('Drawing : center the view on the clicked point -- future zoom will be centered on the point'),
                     'LMB and move': _('Drawing : translate the view'),
                     'Mouse wheel click and move': _('Drawing : translate the view'),
                     'Mouse wheel': _('Drawing : zoom in/out - centered on the middle of the canvas'),
                     'Mouse wheel + Space Bar': _('Drawing : zoom in/out - centered on the mouse position'),
                     'z, Z': _('Drawing : zoom out/in - centered on the middle of the canvas'),
                     'Touchpad 2 fingers': _('Drawing : zoom in/out - centered on the middle of the canvas'),
                     'CTRL + z': _('Drawing : Autoscale only on active array'),
                     'CTRL + Z': _('Drawing : Autoscale only on active vector'),

                     'F5': _('Drawing : autoscale'),
                     'F7': _('Drawing : refresh'),
                     'F8': _('Drawing : zoom on whole Walonia'),
                     'Arrows': _('Drawing : lateral movements'),
                     'c or C': _('Drawing : copy canvas to Clipboard wo axes'),
                     'ALT+C': _('Drawing : copy canvas to Clipboard as Matplotlib image'),
                     'ALT+SHIFT+C': _('Drawing : copy canvas to Clipboard as Matplotlib image with axes - multiviewers'),
                     'CTRL+ALT+SHIFT+C': _('Drawing : copy canvas to Clipboard as Matplotlib image with axes - all arrays one by one'),
                     'd or D': _('Drawing : calculate distance along a temporary vector'),

                     'CTRL+o': _('Results : increase transparency of the current result'),
                     'CTRL+O': _('Results : decrease transparency of the current result'),
                     'o': _('Arrays : increase transparency of the current array'),
                     'O': _('Arrays : decrease transparency of the current array'),

                     'F9': _('Arrays : select all cells'),
                     'F11': _('Arrays : select by criteria'),
                     'F12': _('Arrays : operations'),
                     'n or N': _('Arrays : node-by-node selection'),

                     'b or B': _('Arrays : temporary/active vector selection - along polyline'),
                     'v or V': _('Arrays : temporary/active vector selection - inside polygon'),

                     'r': _('Arrays : reset the selection'),
                     'R': _('Arrays : reset the selection and the associated dictionnary'),

                     '1,2...9': _('Arrays : transfer the selection to the associated dictionary - key 1 to 9'),
                     '>, <' : _('Arrays : dilate/erode the selection - cross-shaped neighbours'),
                     'CTRL+>, CTRL+<': _('Arrays : dilate/erode the selection unselecting the values inside the contour - cross-shaped neighbours'),

                     'i': _('Arrays : 2D interpolation based on the selection on the current matrix'),
                     'CTRL+C': _('Arrays : Set copy source and current selection to clipboard as string'),
                     'CTRL+X': _('Arrays : Crop the active array using the active vector and make a copy'),
                     'CTRL+V': _('Arrays : paste selected values'),
                     'CTRL+ALT+C or ALTGr+C': _('Arrays : Set copy source and current selection to clipboard as script'),
                     'CTRL+ALT+X or ALTGr+X': _('Arrays : Crop the active array using the active vector without masking the values outside the vector'),
                     'CTRL+ALT+V or ALTGr+V': _('Arrays : paste selection to active array'),

                     'p or P': _('Cross sections : Pick a profile/cross section'),

                     'f or F, CTRL+F': _('Zones : search for the polyline in the current zone or in all zones'),

                     'RETURN': _('Action : End the current action (see also right double-click -- OnRDClick)'),
                     'Press and Hold CTRL': _('Action : Data Frame follows the mouse cursor'),

                     'DELETE': _('Tree : Remove item'),

                     'CTRL+L': _('Tiles: Pick a tile by clicking on it'),

                     'CTRL+U': _('GLTF/GLB : import/update GLTF/GLB'),

                     'CTRL+Q': _('App : Quit application'),}

        def gettxt():
            txt = ''
            for curgroup in groups:
                txt += curgroup + '\n'
                for curkey, curval in shortcuts.items():
                    if curgroup in curval:
                        txt += '\t' + curkey + ' : ' + curval.split(':')[1] + '\n'
                txt += '\n'
            return txt

        logging.info(gettxt())

        if inframe :
            frame = wx.Frame(None, -1, _('Shortcuts'), size=(500, 800))
            # panel = wx.Panel(frame, -1)

            sizer = wx.BoxSizer(wx.VERTICAL)

            multiline = wx.TextCtrl(frame, -1, '', style=wx.TE_MULTILINE | wx.TE_READONLY | wx.TE_RICH2)

            multiline.SetValue(gettxt())

            sizer.Add(multiline, 1, wx.EXPAND)

            frame.SetSizer(sizer)

            icon = wx.Icon()
            icon_path = Path(__file__).parent / "apps/wolf.ico"
            icon.CopyFromBitmap(wx.Bitmap(str(icon_path), wx.BITMAP_TYPE_ANY))
            frame.SetIcon(icon)

            frame.SetAutoLayout(True)
            frame.Layout()

            frame.Show()

    def msg_action(self, which:int = 0):
        """ Message to end action """

        if which == 0:
            self.set_statusbar_text(_('Action in progress... -- To quit, press "RETURN" or "double clicks RIGHT" or press "ESC"'))
        else:
            self.set_statusbar_text('')

    def start_action(self, action:str, message:str=''):
        """ Message to start action """

        assert isinstance(action, str), 'action must be a string'
        if action == '':
            self.action = None
        else:
            self.action = action.lower()
        logging.info(_('ACTION : ') + _(message) if message != '' else _('ACTION : ') + _(action))
        self.msg_action(0)

    def end_action(self, message:str=''):
        """ Message to end action """

        self.action = None
        self.active_vertex = None
        logging.info(_('ACTION : ') + _(message) if message != '' else _('ACTION : End of action') )
        self.msg_action(1)

    def distance_by_multiple_clicks(self):
        """ Distance between multiple clicks """

        self.start_action('distance along vector', _('Distance by multiple clicks -- Select the points'))
        self._tmp_vector_distance = vector()
        self._tmp_vector_distance.add_vertex([wolfvertex(0., 0.),
                                              wolfvertex(0., 0.)])

    def OnHotKey(self, e: wx.KeyEvent):
        """
        Gestion des touches clavier -- see print_shortcuts for more details
        """

        key = e.GetKeyCode()
        ctrldown = e.ControlDown()
        altdown = e.AltDown()
        shiftdown = e.ShiftDown()

        myobj = e.EventObject

        logging.debug(_('You are pressing key code : ') + str(key))
        if ctrldown:
            logging.debug(_('Ctrl is down'))
        if altdown:
            logging.debug(_('Alt is down'))

        if ctrldown or altdown:
            if key == 60 and shiftdown: #'>'
                if self.active_array is not None:
                    if self.active_array.SelectionData is not None:
                        self.active_array.SelectionData.dilate_contour_selection(1)
                        self.active_array.reset_plot()
            elif key == 60 and not shiftdown: #'<'
                if self.active_array is not None:
                    if self.active_array.SelectionData is not None:
                        self.active_array.SelectionData.erode_contour_selection()
                        self.active_array.reset_plot()

            elif key == wx.WXK_F2 and ctrldown and altdown and shiftdown:

                if self.active_res2d is None:
                    logging.info(_('Please activate a simulation before search a specific result'))

                self._add_sim_explorer(self.active_res2d)

            elif key == wx.WXK_F2 and not shiftdown:

                if self.active_res2d is not None:
                    nb = self.active_res2d.get_nbresults()
                    dlg = wx.NumberEntryDialog(None,_('Please choose a step (1 -> {})'.format(nb)),'Step :', _('Select a specific step'), nb, min=1, max=nb)
                    ret = dlg.ShowModal()

                    nb = dlg.GetValue()
                    dlg.Destroy()

                    self.active_res2d.read_oneresult(nb-1)
                    self.active_res2d.set_currentview()
                    self.Refresh()

                    self._update_sim_explorer()

                else:
                    logging.info(_('Please activate a simulation before search a specific result'))

            elif key == wx.WXK_F2 and shiftdown:

                if self.active_res2d is not None:
                    nb = self.active_res2d.get_nbresults()

                    choices = ['{:.3f} [s] - {} [h:m:s]'.format(cur, timedelta(seconds=int(cur),
                                                                   milliseconds=int(cur-int(cur))*1000))
                               for cur in self.active_res2d.times]

                    dlg = wx.SingleChoiceDialog(None,
                                                _('Please choose a time step'),
                                                _('Select a specific step'),
                                                choices)
                    ret = dlg.ShowModal()
                    if ret == wx.ID_CANCEL:
                        dlg.Destroy()
                        return

                    keyvalue = dlg.GetStringSelection()
                    dlg.Destroy()

                    self.active_res2d.read_oneresult(choices.index(keyvalue))
                    self.active_res2d.set_currentview()
                    self.Refresh()

                    self._update_sim_explorer()

                else:
                    logging.info(_('Please activate a simulation before searching a specific result'))

            if key == wx.WXK_F4 and not shiftdown:

                if self.active_particle_system is not None:

                    nb = self.active_particle_system.nb_steps
                    dlg = wx.NumberEntryDialog(None,_('Please choose a step (1 -> {})'.format(nb)),'Step :', _('Select a specific step'), nb, min=1, max=nb)
                    ret = dlg.ShowModal()

                    nb = dlg.GetValue()
                    dlg.Destroy()

                    self.active_particle_system.current_step = nb-1
                    self.Refresh()
                    self._update_mytooltip()
                    self._update_sim_explorer()

                else:
                    logging.info(_('Please activate a particle system before searching a specific result'))

            elif key == wx.WXK_F4 and shiftdown:

                if self.active_particle_system is not None:

                    choices = ['{:.3f} [s] - {} [h:m:s]'.format(cur, timedelta(seconds=int(cur),
                                                                   milliseconds=int(cur-int(cur))*1000))
                               for cur in self.active_particle_system.get_times()]

                    dlg = wx.SingleChoiceDialog(None,
                                                _('Please choose a time step'),
                                                _('Select a specific step'),
                                                choices)
                    ret = dlg.ShowModal()
                    if ret == wx.ID_CANCEL:
                        dlg.Destroy()
                        return

                    keyvalue = dlg.GetStringSelection()
                    dlg.Destroy()

                    self.active_particle_system.current_step = choices.index(keyvalue)
                    self.Refresh()
                    self._update_mytooltip()
                    self._update_sim_explorer()

                else:
                    logging.info(_('Please activate a simulation before search a specific result'))

            elif key == wx.WXK_NUMPAD_ADD: #+ from numpad
                if self.active_res2d is not None:
                    self.active_res2d.update_zoom_2(1.1)
                    self.Refresh()

            elif key == wx.WXK_NUMPAD_SUBTRACT: #- from numpad
                if self.active_res2d is not None:
                    self.active_res2d.update_zoom_2(1./1.1)
                    self.Refresh()

            elif key == ord('X'):
                # Create a new array from the active array and the active vector
                # Node outside the vector are set to NullValue
                if self.active_array is not None and self.active_vector is not None:

                    bbox = self.active_vector.get_bounds_xx_yy()
                    newarray = self.active_array.crop_array(bbox)

                    if not altdown:
                        newarray.mask_outsidepoly(self.active_vector)

                    newarray.nullify_border(width=1)

                    #keys for arrays
                    keys = self.get_list_keys(draw_type.ARRAYS, checked_state=None)
                    new_key = self.active_array.idx + '_crop'

                    while new_key in keys:
                        new_key += '_'

                    self.add_object('array', newobj = newarray, id = new_key)

                    self.Refresh()

            elif key == ord('Q'):
                # If Ctrl-Q is hit, then we must *not* handle it
                # because it is tied to the Ctrl-Q accelerator
                # of the "quit" menu...
                e.Skip()
                return

            if key == ord('U'):
                # CTRL+U
                # Mise à jour des données par import du fichier gtlf2
                msg = ''
                if self.active_array is None:
                    msg += _('Active array is None\n')

                if msg != '':
                    msg += _('\n')
                    msg += _('Retry !\n')
                    wx.MessageBox(msg)
                    return

                self.set_fn_fnpos_gltf()
                self.update_blender_sculpting()

            elif key == ord('F'):
                if self.active_zones is not None:
                    self.start_action('select active vector all', _('Select active vector all'))

            elif key == ord('L'):
                if self.active_tile is not None:
                    self.start_action('select active tile', _('Select active tile'))

            elif key == wx.WXK_UP:
                self.upobj()

            elif key == wx.WXK_DOWN:
                self.downobj()

            elif key == ord('C') and altdown and not ctrldown and not shiftdown:
                # ALT+C
                #Copie du canvas dans le clipboard pour transfert vers autre application
                self.copy_canvasogl()

            elif key == ord('C') and altdown and not ctrldown and shiftdown:
                # ALT+SHIFT+C
                # Copie du canvas dans le clipboard pour transfert vers autre application
                # Copie des canvas liés

                if not self.linked:
                    logging.error(_('No linked canvas to copy -- calling ALT+C instead'))
                    self.copy_canvasogl()
                    return

                from tempfile import TemporaryDirectory

                logging.info(_('Creating images'))

                with TemporaryDirectory() as tmpdirname:
                    all_images = self.save_linked_canvas(Path(tmpdirname) / 'fig', mpl= True, ds= self.ticks_size, add_title= True)
                    if len(all_images) == 0:
                        logging.error(_('No image to combine -- aborting !'))
                        return
                    im_assembly = self.assembly_images(all_images, mode= self.assembly_mode)

                logging.info(_('Creating images - done'))

                # Copy image to clipboard
                if im_assembly is not None:
                    if wx.TheClipboard.Open():

                        #création d'un objet bitmap wx
                        wxbitmap = wx.Bitmap().FromBuffer(im_assembly.width, im_assembly.height, im_assembly.tobytes())

                        # objet wx exportable via le clipboard
                        dataobj = wx.BitmapDataObject()
                        dataobj.SetBitmap(wxbitmap)

                        wx.TheClipboard.SetData(dataobj)
                        wx.TheClipboard.Close()
                        logging.info(_('Image copied to clipboard'))
                    else:
                        logging.error(_('Cannot open the clipboard'))
                else:
                    logging.error(_('No image to copy to clipboard'))

            elif key == ord('C') and ctrldown and not altdown:
                # CTRL+C
                if self.active_array is None:
                    dlg = wx.MessageDialog(self,
                                        _('The active array is None - Please active an array from which to copy the values !'),
                                        style=wx.OK)
                    dlg.ShowModal()
                    dlg.Destroy()
                    return

                logging.info(_('Start copying values / Current selection to clipboard'))
                self.copyfrom = self.active_array
                self.mimicme_copyfrom()  # force le recopiage de copyfrom dans les autres matrices liées

                if len(self.active_array.SelectionData.myselection) > 5000:
                    dlg = wx.MessageDialog(self, _('The selection is large, copy to clipboard may be slow ! -- Continue?'), style=wx.OK | wx.CANCEL)
                    ret = dlg.ShowModal()

                    if ret == wx.ID_CANCEL:
                        logging.info(_('Copy to clipboard cancelled -- But source array is well defined !'))
                        dlg.Destroy()
                        return

                    dlg.Destroy()

                self.active_array.SelectionData.copy_to_clipboard()

                logging.info(_('Values copied to clipboard'))

            elif key == ord('C') and ctrldown and altdown and shiftdown:
                # CTRL+ALT+SHIFT+C
                # Copie du canvas dans le clipboard pour transfert vers autre application
                # Une matrice est associée à chaque canvas

                from tempfile import TemporaryDirectory

                logging.info(_('Creating images'))

                with TemporaryDirectory() as tmpdirname:
                    all_images = self.save_arrays_indep(Path(tmpdirname) / 'fig', mpl= True, ds= self.ticks_size, add_title= True)
                    if len(all_images) == 0:
                        logging.error(_('No image to combine -- aborting !'))
                        return
                    im_assembly = self.assembly_images(all_images, mode= self.assembly_mode)

                logging.info(_('Creating images - done'))

                # Copy image to clipboard
                if im_assembly is not None:
                    if wx.TheClipboard.Open():

                        #création d'un objet bitmap wx
                        wxbitmap = wx.Bitmap().FromBuffer(im_assembly.width, im_assembly.height, im_assembly.tobytes())

                        # objet wx exportable via le clipboard
                        dataobj = wx.BitmapDataObject()
                        dataobj.SetBitmap(wxbitmap)

                        wx.TheClipboard.SetData(dataobj)
                        wx.TheClipboard.Close()
                        logging.info(_('Image copied to clipboard'))
                    else:
                        logging.error(_('Cannot open the clipboard'))
                else:
                    logging.error(_('No image to copy to clipboard'))

            elif key == ord('C') and ctrldown and altdown and not shiftdown:
                if self.active_array is None:
                    dlg = wx.MessageDialog(self,
                                        _('The active array is None - Please active an array from which to copy the selection !'),
                                        style=wx.OK)
                    dlg.ShowModal()
                    dlg.Destroy()
                    return

                logging.info(_('Start copying selection / Current selection to clipboard as script (Python)'))
                self.copyfrom = self.active_array
                self.mimicme_copyfrom()  # force le recopiage de copyfrom dans les autres matrices liées

                if len(self.active_array.SelectionData.myselection) > 5000:
                    dlg = wx.MessageDialog(self, _('The selection is large, copy to clipboard may be slow ! -- Continue?'), style=wx.OK | wx.CANCEL)
                    ret = dlg.ShowModal()

                    if ret == wx.ID_CANCEL:
                        logging.info(_('Copy script to clipboard cancelled -- But source array is well defined !'))
                        dlg.Destroy()
                        return

                    dlg.Destroy()

                self.active_array.SelectionData.copy_to_clipboard(typestr='script')

                logging.info(_('Selection copied to clipboard as script (Python)'))

            elif key == ord('V') and ctrldown:
                # CTRL+V
                # CTRL+ALT+V ou Alt Gr + V

                if self.active_array is None:
                    if e.AltDown():
                        # CTRL+ALT+V
                        logging.warning(_('The active array is None - Please active an array into which to paste the selection !'))
                    else:
                        logging.warning(_('The active array is None - Please active an array into which to paste the values !'))

                    return

                fromarray = self.copyfrom
                if fromarray is None:
                    if self.linked:
                        if not self.linkedList is None:
                            for curFrame in self.linkedList:
                                if curFrame.copyfrom is not None:
                                    fromarray = curFrame.copyfrom
                                    break

                if fromarray is None:
                    logging.warning(_('No selection to be pasted !'))
                    return

                cursel = fromarray.SelectionData.myselection

                if e.AltDown():
                    logging.info(_('Paste selection position'))

                    if cursel == 'all':
                        self.active_array.SelectionData.myselection = 'all'
                    elif len(cursel) > 0:
                        self.active_array.SelectionData.myselection = cursel.copy()
                        # self.active_array.SelectionData.update_nb_nodes_selection()

                else:
                    logging.info(_('Paste selection values'))
                    if cursel == 'all':
                        self.active_array.paste_all(fromarray)

                    elif len(cursel) > 0:
                        z = fromarray.SelectionData.get_values_sel()
                        self.active_array.set_values_sel(cursel, z)

                self.Refresh()

                logging.info(_('Selection/Values pasted'))

            elif key == ord('Z'):

                if ctrldown:
                    if shiftdown:
                        if self.active_vector is not None:
                            self.zoom_on_vector(self.active_vector, canvas_height= self.canvas.GetSize()[1])
                        else:
                            logging.warning(_('No active vector to zoom on !'))
                    else:
                        if self.active_array is not None:
                            self.zoom_on_array(self.active_array, canvas_height= self.canvas.GetSize()[1])
                        else:
                            logging.warning(_('No active array to zoom on !'))
        else:
            if key == wx.WXK_DELETE:
                self.removeobj()

            elif key == 60 and shiftdown: #'>'
                if self.active_array is not None:
                    if self.active_array.SelectionData is not None:
                        self.active_array.SelectionData.dilate_selection(1)
                        self.active_array.reset_plot()
            elif key == 60 and not shiftdown: #'<'
                if self.active_array is not None:
                    if self.active_array.SelectionData is not None:
                        self.active_array.SelectionData.erode_selection(1)
                        self.active_array.reset_plot()

            elif key == wx.WXK_ESCAPE:

                logging.info(_('Escape key pressed -- Set all active objects and "action" to None'))

                self.action = None
                self.active_array = None
                self.active_vector = None
                self.active_zone = None
                self.active_zones = None
                self.active_res2d = None
                self.active_tile = None
                self.active_particle_system = None
                self.active_vertex = None
                self.active_cloud = None
                self.active_dike = None
                self.active_injector = None

                self.active_laz = None

                self.active_fig = None
                self.active_bridge  = None
                self.active_bridges = None
                self.active_bc = None
                self.active_cs = None
                self.active_imagestiles = None
                self.active_landmap = None
                self.active_profile = None
                self.active_tri = None
                self.active_viewer3d = None

                self.set_statusbar_text(_('Esc pressed - No more action in progress - No more active object'))
                self.set_label_selecteditem('')

            elif key == ord('C'):

                 self.copy_canvasogl(mpl = False)

            elif key == wx.WXK_SPACE:
                if self.timer_ps is not None and self.active_particle_system is not None :
                    if self.timer_ps.IsRunning():
                        self.timer_ps.Stop()
                    else:
                        if self.active_particle_system.current_step_idx == self.active_particle_system.nb_steps-1:
                            self.active_particle_system.current_step_idx = 0
                            self.active_particle_system.current_step = 0
                        self.timer_ps.Start(1000. / self.active_particle_system.fps)

            elif key == 388: #+ from numpad
                if self.active_res2d is not None:
                    self.active_res2d.update_arrowpixelsize_vectorfield(-1)
                    self.Refresh()

            elif key == 390: #- from numpad
                if self.active_res2d is not None:
                    self.active_res2d.update_arrowpixelsize_vectorfield(1)
                    self.Refresh()

            elif key == 13 or key==370 or key == wx.WXK_RETURN or key == wx.WXK_NUMPAD_ENTER:
                # 13 = RETURN classic keyboard
                # 370 = RETURN NUMPAD
                self._endactions()

            elif key == ord('I'):
                if self.active_array is not None :
                    self.active_array.interpolation2D()

            elif key == ord('F'):
                if self.active_zone is not None:
                    self.start_action('select active vector2 all', _('Select active vector2 all'))

            elif key in LIST_1TO9:

                if self.active_array is not None:

                    if self.active_array.SelectionData.myselection == 'all':
                        logging.warning(_('No selection to transfer to the dictionary !'))
                        logging.info(_('Please select some nodes before transfering to the dictionary, not ALL !'))
                        return

                    # colors = [(0, 0, 255, 255),
                    #           (0, 255, 0, 255),
                    #           (0, 128, 255, 255),
                    #           (255, 255, 0, 255),
                    #           (255, 165, 0, 255),
                    #           (128, 0, 128, 255),
                    #           (255, 192, 203, 255),
                    #           (165, 42, 42, 255),
                    #           (128, 128, 128, 255)]

                    idx = LIST_1TO9.index(key)
                    if idx > 8:
                        idx -= 9

                    self.active_array.SelectionData.move_selectionto(str(idx+1), self.colors1to9[idx])

            elif key == wx.WXK_F1:
                self.read_last_result()

            elif key == wx.WXK_F2 and shiftdown:
                self.simul_previous_step()

            elif key == wx.WXK_F4 and shiftdown:
                self.particle_previous_step()

            elif key == wx.WXK_F4:
                self.particle_next_step()

            elif key == wx.WXK_F2:
                self.simul_next_step()

            elif key == wx.WXK_F5:
                # Autoscale
                self.Autoscale()

            elif key == wx.WXK_F7:
                self.update()

            elif key == wx.WXK_F8:
                self.zoom_on_whole_walonia()

            elif key == wx.WXK_F12 or key == wx.WXK_F11:
                if self.active_array is not None:
                    self.active_array.myops.SetTitle(_('Operations on array: ')+self.active_array.idx)
                    self.active_array.myops.Show()
                    self.active_array.myops.array_ops.SetSelection(1)
                    self.active_array.myops.Center()

            elif key == wx.WXK_F9:

                if self.active_array is not None:
                    if self.active_array.SelectionData is not None:
                        self.active_array.SelectionData.myselection = 'all'
                        logging.info(_('Selecting all nodes in the active array !'))
                    else:
                        logging.warning(_('No selection manager for this array !'))

                    if self.active_array.myops is not None:
                        self.active_array.myops.nbselect.SetLabelText('All')
                    else:
                        logging.warning(_('No operations manager for this array !'))

            elif key == ord('N'):  # N
                if self.active_array is not None:
                    self.active_array.myops.select_node_by_node()

                if self.active_res2d is not None:
                    if self.active_array is not None:
                        msg = wx.MessageDialog(None, _('Do you want to select the nodes of the active result ?'), _('Select nodes'), wx.YES_NO | wx.ICON_QUESTION)
                        ret = msg.ShowModal()
                        if ret == wx.ID_YES:
                            self.active_res2d.properties.select_node_by_node()
                    else:
                        self.active_res2d.properties.select_node_by_node()

                if self.active_array is None and self.active_res2d is None:
                    logging.warning(_('No active array or result 2D to select node by node !'))

            elif key == ord('V'):  # V
                if self.active_array is not None:
                    if shiftdown:
                        self.active_array.myops.select_vector_inside_manager()
                    else:
                        self.active_array.myops.select_vector_inside_tmp()
                else:
                    logging.warning(_('No active array to select the vector inside !'))

            elif key == ord('B'):  # B
                if self.active_array is not None:
                    if shiftdown:
                        self.active_array.myops.select_vector_under_manager()
                    else:
                        self.active_array.myops.select_vector_under_tmp()
                else:
                    logging.warning(_('No active array to select the vector inside !'))

            elif key == ord('P'):  # P

                if self.active_cs is not None:
                    self.start_action('Select nearest profile', _('Select nearest profile'))
                else:
                    logging.warning(_('No active cross section to select the nearest profile !'))

            elif key == ord('Z') and shiftdown:  # Z
                self.width = self.width / 1.1
                self.height = self.height / 1.1
                self.setbounds()

            elif key == ord('Z'):  # z
                self.width = self.width * 1.1
                self.height = self.height * 1.1
                self.setbounds()

            elif key == ord('R') and shiftdown:  # R
                if self.active_array is not None:
                    self.active_array.myops.reset_all_selection()
                    self.Refresh()

                if self.active_res2d is not None:
                    self.active_res2d.SelectionData.reset_all()
                    self.Refresh()

            elif key == ord('R'):  # r
                if self.active_array is not None:
                    self.active_array.myops.reset_selection()
                    self.Refresh()

                if self.active_res2d is not None:
                    self.active_res2d.SelectionData.reset()
                    self.Refresh()

            elif key == ord('O'):
                # Active Opacity for the active array

                if ctrldown:
                    if self.active_res2d is None:
                        logging.warning(_('No active result 2D to change the opacity !'))
                        return

                    if shiftdown:
                        self.active_res2d.set_opacity(self.active_res2d.alpha + 0.25)
                    else:
                        self.active_res2d.set_opacity(self.active_res2d.alpha - 0.25)

                else:
                    if self.active_array is None:
                        logging.warning(_('No active array to change the opacity !'))
                        return

                    if shiftdown:
                        self.active_array.set_opacity(self.active_array.alpha + 0.25)
                    else:
                        self.active_array.set_opacity(self.active_array.alpha - 0.25)

            elif key == wx.WXK_UP:
                self.mousey = self.mousey + self.height / 10.
                self.setbounds()

            elif key == wx.WXK_DOWN:
                self.mousey = self.mousey - self.height / 10.
                self.setbounds()

            elif key == wx.WXK_LEFT:
                self.mousex = self.mousex - self.width / 10.
                self.setbounds()

            elif key == wx.WXK_RIGHT:
                self.mousex = self.mousex + self.width / 10.
                self.setbounds()

            elif key == ord('A'):
                if self.active_laz is not None:
                    self.active_laz.add_pose_in_memory()

            elif key == ord('D'):
                self.distance_by_multiple_clicks()

    def paste_values(self,fromarray:WolfArray):
        """ Paste selected values from a WolfArray to the active array """

        if self.active_array is None:
            logging.warning(_('The active array is None - Please active an array into which to paste the values !'))
            return

        logging.info(_('Paste selection values'))
        cursel = fromarray.SelectionData.myselection
        if cursel == 'all':
            self.active_array.paste_all(fromarray)
        elif len(cursel) > 0:
            z = fromarray.SelectionData.get_values_sel()
            self.active_array.set_values_sel(cursel, z)

    def paste_selxy(self,fromarray:WolfArray):
        """ Paste selected nodes from a WolfArray to the active array """

        if self.active_array is None:
            logging.warning(_('The active array is None - Please active an array into which to paste the selection !'))
            return

        logging.info(_('Paste selection position'))
        cursel = fromarray.SelectionData.myselection
        if cursel == 'all':
            self.active_array.SelectionData.OnAllSelect(0)
        elif len(cursel) > 0:
            self.active_array.SelectionData.myselection = cursel.copy()
            self.active_array.SelectionData.update_nb_nodes_selection()

    def OntreeRight(self, e: wx.MouseEvent):
        """ Gestion du menu contextuel sur l'arbre des objets """

        if self.selected_object is not None:

            # On va nettoyer le menu contextuel car certaines entrées ne sont
            # pas nécessairement pertinentes

            # Chaînes à supprimer
            tracks=[]
            tracks.append(_('Contours'))
            tracks.append(_('Boundary conditions'))
            tracks.append(_('Convert to mono-block'))
            tracks.append(_('Convert to mono-block (result)'))
            tracks.append(_('Convert to multi-blocks (result)'))
            tracks.append(_('Extract current step as IC (result)'))
            tracks.append(_('Export to Shape file'))
            tracks.append(_('Export active zone to Shape file'))
            tracks.append(_('Rebin'))
            tracks.append(_('Set NullValue'))
            tracks.append(_('Set colormap'))
            tracks.append(_('Edit colormap'))
            tracks.append(_('Set classification'))
            tracks.append(_('Convert to...'))
            tracks.append(_('Edit selection'))
            tracks.append(_('All to cloud'))
            tracks.append(_('Selection to cloud'))

            tracks.append(_('Colormap'))
            tracks.append(_('Movie'))
            tracks.append(_('Play'))
            tracks.append(_('Record'))
            tracks.append(_('Load flight'))
            tracks.append(_('Save flight'))

            tracks.append(_('Rasterize active zone'))
            tracks.append(_('Rasterize active vector'))

            tracks.append(_('Extrude on active array'))
            tracks.append(_('Interpolate on active array'))

            # Récupération des items du menu contextuel
            menuitems = self.popupmenu.GetMenuItems()
            text = [cur.GetItemLabelText() for cur in menuitems]

            # Liste des indices à supprimer
            # Pas possible de supprimer à la volée car cela modifie la liste
            to_delete = []
            for track in tracks:
                if track in text:
                    to_delete.append(text.index(track))

            # Suppression des items
            if len(to_delete) > 0:
                # Suppression en ordre décroissant pour ne pas décaler les indices
                to_delete.sort(reverse=True)
                for idx in to_delete:
                    self.popupmenu.Remove(menuitems[idx])

            # Add specific menu items for WolfArray
            if isinstance(self.selected_object, WolfArray):
                bc = self.get_boundary_manager(self.selected_object)
                if bc is not None:
                    self.popupmenu.Append(wx.ID_ANY, _('Boundary conditions'), _('Boundary conditions'))
                self.popupmenu.Append(wx.ID_ANY, _('Contours'))
                self.popupmenu.Append(wx.ID_ANY, _('Rebin'), _('Change the spatial resolution'))
                self.popupmenu.Append(wx.ID_ANY, _('Set NullValue'), _('Set NullValue'))

            # Add specific menu items for WolfArrayMB
            if isinstance(self.selected_object, WolfArrayMB):
                self.popupmenu.Append(wx.ID_ANY, _('Convert to mono-block'), _('Convert to mono-block'))

            # Add specific menu items for Wolfresults_2D
            if isinstance(self.selected_object, Wolfresults_2D):
                self.popupmenu.Append(wx.ID_ANY, _('Convert to mono-block (result)'), _('Convert to mono-block'))
                self.popupmenu.Append(wx.ID_ANY, _('Convert to multi-blocks (result)'), _('Convert to multi-blocks'))
                self.popupmenu.Append(wx.ID_ANY, _('Extract current step as IC (result)'), _('Extract current step as IC'))

            if isinstance(self.selected_object, Zones):
                self.popupmenu.Append(wx.ID_ANY, _('Rasterize active zone'), _('Rasterize active zone'))
                self.popupmenu.Append(wx.ID_ANY, _('Rasterize active vector'), _('Rasterize active vector'))
                self.popupmenu.Append(wx.ID_ANY, _('Interpolate on active array'), _('Interpolate Z-values on active array'))

            if isinstance(self.selected_object, Zones | Bridge | Weir):
                self.popupmenu.Append(wx.ID_ANY, _('Export to Shape file'), _('Export to Shape file'))
                self.popupmenu.Append(wx.ID_ANY, _('Export active zone to Shape file'), _('Export active zone to Shape file'))

            if isinstance(self.selected_object, Wolf_LAZ_Data):

                colrmapmenu = wx.Menu()
                self.popupmenu.AppendSubMenu(colrmapmenu, _('Colormap'))

                colrmapmenu.Append(wx.ID_ANY, _('Set colormap'), _('Change colormap'))
                colrmapmenu.Append(wx.ID_ANY, _('Edit colormap'), _('Edit colormap'))
                colrmapmenu.Append(wx.ID_ANY, _('Set classification'), _('Change classification'))

                converttomenu  = wx.Menu()
                self.popupmenu.AppendSubMenu(converttomenu, _('Convert to...'))

                converttomenu.Append(wx.ID_ANY, _('All to cloud'), _('Convert all to cloud'))
                converttomenu.Append(wx.ID_ANY, _('Selection to cloud'), _('Convert selection to cloud'))
                converttomenu.Append(wx.ID_ANY, _('Selection to vector'), _('Convert selection to vector'))

                self.popupmenu.Append(wx.ID_ANY, _('Edit selection'), _('Edit selection'))

                moviemenu = wx.Menu()
                self.popupmenu.AppendSubMenu(moviemenu, _('Movie'))

                moviemenu.Append(wx.ID_ANY, _('Add point'), _('Add point passage'))
                moviemenu.Append(wx.ID_ANY, _('Play'), _('Play'))
                # moviemenu.Append(wx.ID_ANY, _('Record'), _('Record'))
                moviemenu.Append(wx.ID_ANY, _('Load flight'), _('Load flight'))
                moviemenu.Append(wx.ID_ANY, _('Save flight'), _('Save flight'))

            if isinstance(self.selected_object, Picc_data):
                self.popupmenu.Append(wx.ID_ANY, _('Extrude on active array'), _('Extrude building elevation on active array'))

            self.treelist.PopupMenu(self.popupmenu)

    def zoom_on_whole_walonia(self):
        """ Zoom on the whole Walonia """

        xmin = 40_000
        xmax = 300_000
        ymin = 10_000
        ymax = 175_000

        self.mousex = (xmin + xmax) / 2.
        self.mousey = (ymin + ymax) / 2.

        self.width = xmax - xmin
        self.height = ymax - ymin

        self.setbounds()
        self.update()


    def _update_background(self):
        """
        Update background
        """

        # dessin du background
        for obj in self.iterator_over_objects(draw_type.WMSBACK):
            obj.reload()

    def _update_foreground(self):
        """
        Update foreground
        """
        # dessin du foreground
        for obj in self.iterator_over_objects(draw_type.WMSFORE):
            obj.reload()

    def update(self):
        """
        Update backgournd et foreground elements and arrays if local minmax is checked.

        """

        self._update_background()

        self._update_foreground()

        if self.locminmax.IsChecked() or self.update_absolute_minmax:
            for curarray in self.iterator_over_objects(draw_type.ARRAYS):
                curarray: WolfArray
                if self.update_absolute_minmax:
                    curarray.updatepalette()
                    self.update_absolute_minmax = False
                else:
                    curarray.updatepalette(onzoom=[self.xmin, self.xmax, self.ymin, self.ymax])
                curarray.delete_lists()

        self.Paint()

    def _plotting(self, drawing_type: draw_type, checked_state: bool = True):
        """ Drawing objets on canvas"""

        try:
            for curobj in self.iterator_over_objects(drawing_type, checked_state=checked_state):
                if not curobj.plotting:
                    curobj.plotting = True
                    curobj.plot(sx = self.sx, sy=self.sy, xmin=self.xmin, ymin=self.ymin, xmax=self.xmax, ymax=self.ymax, size = (self.xmax - self.xmin) / 100.)
                    curobj.plotting = False
        except Exception as ex:
            curobj.plotting = False
            logging.error(_('Error while plotting objects of type {}').format(drawing_type.name))

            traceback.print_exc()
            logging.error(ex)

    def get_MVP_Viewport_matrix(self):
        """ Get the modelview projection matrix """

        if self.SetCurrentContext():
            modelview = glGetFloatv(GL_MODELVIEW_MATRIX)
            projection = glGetFloatv(GL_PROJECTION_MATRIX)
            viewport = glGetIntegerv(GL_VIEWPORT)

            return modelview, projection, viewport
        else:

            return None, None, None

    def SetCurrentContext(self):
        """ Set the current OGL context if exists otherwise return False """

        if self.context is None:
            return False

        return self.canvas.SetCurrent(self.context)

    def _set_gl_projection_matrix(self):
        glMatrixMode(GL_PROJECTION)
        glLoadIdentity()
        glOrtho(self.xmin, self.xmax, self.ymin, self.ymax, -99999, 99999)

        glMatrixMode(GL_MODELVIEW)
        glLoadIdentity()

    def Paint(self):
        """ Dessin des éléments ajoutés au viewer """

        if self.currently_readresults:
            return

        width, height = self.canvas.GetSize()

        # C'est bien ici que la zone de dessin utile est calculée sur base du centre et de la zone en coordonnées réelles
        # Les commandes OpenGL sont donc traitées en coordonnées réelles puisque la commande glOrtho définit le cadre visible
        self.xmin = self.mousex - self.width / 2.
        self.ymin = self.mousey - self.height / 2.
        self.xmax = self.mousex + self.width / 2.
        self.ymax = self.mousey + self.height / 2.

        if self.SetCurrentContext():

            bkg_color = self.bkg_color

            glClearColor(bkg_color[0]/255., bkg_color[1]/255., bkg_color[2]/255., bkg_color[3]/255.)
            # glClearColor(0., 0., 1., 0)
            glClear(GL_COLOR_BUFFER_BIT | GL_DEPTH_BUFFER_BIT)
            glViewport(0, 0, int(width), int(height))

            self._set_gl_projection_matrix()

            # dessin du background
            self._plotting(draw_type.WMSBACK)

            # Dessin des matrices
            self._plotting(draw_type.ARRAYS)

            # Dessin des résultats 2D
            self._plotting(draw_type.RES2D)

            # Dessin des vecteurs
            self._plotting(draw_type.VECTORS)

            # Dessin des tuiles
            self._plotting(draw_type.TILES)
            self._plotting(draw_type.IMAGESTILES)

            if self.active_vector is not None:
                if self.active_vector.parentzone is None:
                    # we must plot this vector because it is a temporary vector outside any zone
                    self.active_vector.plot()

            # Dessin des triangulations
            self._plotting(draw_type.TRIANGULATION)

            # Dessin des nuages
            self._plotting(draw_type.CLOUD)

            # Dessin des vues
            self._plotting(draw_type.VIEWS)

            # Dessin des "particule systems"
            self._plotting(draw_type.PARTICLE_SYSTEM)

            # Dessin du reste
            self._plotting(draw_type.OTHER)

            # Dessin du noyé
            self._plotting(draw_type.DROWNING)

            # Dessin du Front
            self._plotting(draw_type.WMSFORE)

            # Dessin des images
            self._plotting(draw_type.PICTURECOLLECTION)

            # Dessin des QDF/IDF
            if self.active_qdfidf is not None:
                self.active_qdfidf.plot(sx = self.sx, sy=self.sy,
                                        xmin=self.xmin, ymin=self.ymin,
                                        xmax=self.xmax, ymax=self.ymax,
                                        size = (self.xmax - self.xmin) / 100.)

            # Gestion des BC (si actif)
            if self.active_bc is not None:
                self.active_bc.plot()

            if self._tmp_vector_distance is not None:
                self._tmp_vector_distance.plot()

            if self.active_vector is not None:
                if getIfromRGB(self.active_vector_color) != self.active_vector.myprop.color:
                    old = self.active_vector.myprop.color
                    self.active_vector.myprop.color = getIfromRGB(self.active_vector_color)
                    self.active_vector.plot()
                    self.active_vector._plot_square_at_vertices(size = self.active_vector_square_size)
                    self.active_vector.myprop.color = old
                else:
                    self.active_vector._plot_square_at_vertices(size = self.active_vector_square_size)

                if self.active_vector.myprop.plot_indices:
                    self.active_vector._plot_all_indices(sx = self.sx, sy=self.sy,
                                                            xmin=self.xmin, ymin=self.ymin,
                                                            xmax=self.xmax, ymax=self.ymax,
                                                            size = (self.xmax - self.xmin) / 100.)

                elif self.active_vertex is not None:
                    self.active_vector._plot_index_vertex(idx = self.active_vector.myvertices.index(self.active_vertex),
                                                          sx = self.sx, sy=self.sy,
                                                          xmin=self.xmin, ymin=self.ymin,
                                                          xmax=self.xmax, ymax=self.ymax,
                                                          size = (self.xmax - self.xmin) / 100.)

            self.canvas.SwapBuffers()
        else:
            raise NameError(
                'Opengl setcurrent -- maybe a conflict with an existing opengl32.dll file - please rename the opengl32.dll in the libs directory and retry')

    def OnPaint(self, e):
        """ event handler for paint event"""

        self.Paint()
        if e is not None:
            e.Skip()

    def findminmax(self, force=False):
        """ Find min/max of all objects """

        # FIXME : use iterator

        xmin = 1.e30
        ymin = 1.e30
        xmax = -1.e30
        ymax = -1.e30

        k = 0
        for locarray in self.myarrays:
            if locarray.plotted or force:
                [xmin_arr, xmax_arr], [ymin_arr, ymax_arr] = locarray.get_bounds()
                xmin = min(xmin, xmin_arr)
                xmax = max(xmax, xmax_arr)
                ymin = min(ymin, ymin_arr)
                ymax = max(ymax, ymax_arr)
                k += 1

        for locvector in self.myvectors:
            if locvector.plotted or force:
                if locvector.idx != 'grid':
                    locvector.find_minmax()
                    if isinstance(locvector,Zones):
                        xmin = min(locvector.xmin, xmin)
                        xmax = max(locvector.xmax, xmax)
                        ymin = min(locvector.ymin, ymin)
                        ymax = max(locvector.ymax, ymax)
                    elif isinstance(locvector,Bridges):
                        xmin = min(locvector.xmin, xmin)
                        xmax = max(locvector.xmax, xmax)
                        ymin = min(locvector.ymin, ymin)
                        ymax = max(locvector.ymax, ymax)
                    elif isinstance(locvector,crosssections):
                        xmin = min(locvector.xmin, xmin)
                        xmax = max(locvector.xmax, xmax)
                        ymin = min(locvector.ymin, ymin)
                        ymax = max(locvector.ymax, ymax)
                    k += 1

        for locvector in self.myimagestiles:
            if locvector.plotted or force:
                locvector.find_minmax()
                if isinstance(locvector,ImagesTiles):
                    xmin = min(locvector.xmin, xmin)
                    xmax = max(locvector.xmax, xmax)
                    ymin = min(locvector.ymin, ymin)
                    ymax = max(locvector.ymax, ymax)
                k += 1

        for locvector in self.mypicturecollections:
            if locvector.plotted or force:
                locvector.find_minmax()
                xmin = min(locvector.xmin, xmin)
                xmax = max(locvector.xmax, xmax)
                ymin = min(locvector.ymin, ymin)
                ymax = max(locvector.ymax, ymax)
                k += 1

        for locvector in self.mytiles:
            if locvector.plotted or force:
                if locvector.idx != 'grid':
                    locvector.find_minmax()
                    if isinstance(locvector,Zones):
                        xmin = min(locvector.xmin, xmin)
                        xmax = max(locvector.xmax, xmax)
                        ymin = min(locvector.ymin, ymin)
                        ymax = max(locvector.ymax, ymax)
                    elif isinstance(locvector,Bridges):
                        xmin = min(locvector.xmin, xmin)
                        xmax = max(locvector.xmax, xmax)
                        ymin = min(locvector.ymin, ymin)
                        ymax = max(locvector.ymax, ymax)
                    elif isinstance(locvector,crosssections):
                        xmin = min(locvector.xmin, xmin)
                        xmax = max(locvector.xmax, xmax)
                        ymin = min(locvector.ymin, ymin)
                        ymax = max(locvector.ymax, ymax)
                    k += 1

        for loccloud in self.myclouds:
            if loccloud.plotted or force:
                loccloud.find_minmax(force)
                xmin = min(loccloud.xbounds[0], xmin)
                xmax = max(loccloud.xbounds[1], xmax)
                ymin = min(loccloud.ybounds[0], ymin)
                ymax = max(loccloud.ybounds[1], ymax)
                k += 1

        for loctri in self.mytri:
            if loctri.plotted or force:
                loctri.find_minmax(force)
                xmin = min(loctri.xmin, xmin)
                xmax = max(loctri.xmax, xmax)
                ymin = min(loctri.ymin, ymin)
                ymax = max(loctri.ymax, ymax)
                k += 1

        for locres2d in self.myres2D:
            locres2d:Wolfresults_2D
            if locres2d.plotted or force:
                locres2d.find_minmax(force)
                xmin = min(locres2d.xmin, xmin)
                xmax = max(locres2d.xmax, xmax)
                ymin = min(locres2d.ymin, ymin)
                ymax = max(locres2d.ymax, ymax)
                k += 1

        for locps in self.mypartsystems:
            locps:Particle_system
            if locps.plotted or force:
                locps.find_minmax(force)
                xmin = min(locps.xmin, xmin)
                xmax = max(locps.xmax, xmax)
                ymin = min(locps.ymin, ymin)
                ymax = max(locps.ymax, ymax)
                k += 1

        for locview in self.myviews:
            locview.find_minmax(force)
            xmin = min(locview.xmin, xmin)
            xmax = max(locview.xmax, xmax)
            ymin = min(locview.ymin, ymin)
            ymax = max(locview.ymax, ymax)
            k += 1

        for locothers in self.myothers:
            if type(locothers) in [genericImagetexture]: #, hydrometry_wolfgui]:
                xmin = min(locothers.xmin, xmin)
                xmax = max(locothers.xmax, xmax)
                ymin = min(locothers.ymin, ymin)
                ymax = max(locothers.ymax, ymax)
                k += 1
            elif type(locothers) in [PlansTerrier]: #, hydrometry_wolfgui]:
                if locothers.initialized:
                    xmin = min(locothers.xmin, xmin)
                    xmax = max(locothers.xmax, xmax)
                    ymin = min(locothers.ymin, ymin)
                    ymax = max(locothers.ymax, ymax)
                    k += 1
            elif type(locothers) in [Particularites, Enquetes, Ouvrages, Profils]:
                if locothers.initialized:
                    xmin = min(locothers.xmin, xmin)
                    xmax = max(locothers.xmax, xmax)
                    ymin = min(locothers.ymin, ymin)
                    ymax = max(locothers.ymax, ymax)
                    k += 1

        for drown in self.mydrownings:
            if drown.plotted or force:
                drown.find_minmax(force)
                xmin = min(drown.xmin, xmin)
                xmax = max(drown.xmax, xmax)
                ymin = min(drown.ymin, ymin)
                ymax = max(drown.ymax, ymax)
                k += 1


        if k > 0:
            self.xmin = xmin
            self.xmax = xmax
            self.ymin = ymin
            self.ymax = ymax

    def resizeFrame(self, w:int, h:int):
        """ Resize the frame

        :param w: width in pixels
        :param h: height in pixels
        """

        self.SetClientSize(w, h)

    def mimicme(self):
        """
        Report des caractéristiques de la fenêtre sur les autres éléments liés
        """

        if self.linked and self.forcemimic:
            if not self.linkedList is None:
                width, height = self.GetClientSize()

                curFrame: WolfMapViewer
                for curFrame in self.linkedList:
                    curFrame.forcemimic = False

                for curFrame in self.linkedList:
                    if curFrame != self:
                        curFrame.resizeFrame(width, height)
                        curFrame.mousex = self.mousex
                        curFrame.mousey = self.mousey
                        curFrame.sx = self.sx
                        curFrame.sy = self.sy
                        curFrame.width = self.width
                        curFrame.height = self.height
                        curFrame.setbounds()

                        if curFrame.link_shareopsvect:
                            if curFrame.active_vector is not self.active_vector:
                                curFrame.Active_vector(self.active_vector)
                            if curFrame.active_array.myops.active_vector is not self.active_vector:
                                curFrame.active_array.myops.Active_vector(self.active_vector, False)
                            curFrame.action = self.action

                for curFrame in self.linkedList:
                    curFrame.forcemimic = True

    def mimicme_copyfrom(self):
        if self.linked and self.forcemimic:
            if not self.linkedList is None:
                width, height = self.GetClientSize()

                curFrame: WolfMapViewer
                for curFrame in self.linkedList:
                    curFrame.forcemimic = False

                for curFrame in self.linkedList:
                    if curFrame != self:
                        curFrame.copyfrom = self.copyfrom

                for curFrame in self.linkedList:
                    curFrame.forcemimic = True

    def Active_vector(self, vect):
        """ Active un vecteur et sa zone parent si existante """

        self.active_vector = vect

        if vect is not None:
            logging.info(_('Activating vector : ' + vect.myname))
            if vect.parentzone is not None:
                self.Active_zone(vect.parentzone)

        self.mimicme()
        self.Paint()

    def Active_zone(self, zone: zone):
        """ Active une zone et son parent si existant """

        self.active_zone = zone
        self.active_zones = zone.parent
        logging.info(_('Activating zone : ' + zone.myname))

    def list_background(self):
        return [cur.idx for cur in self.mywmsback]

    def list_foreground(self):
        return [cur.idx for cur in self.mywmsfore]

    def check_id(self, id=str, gridsize = 100.):
        """ Check an element from its id """

        curobj = self.getobj_from_id(id)

        if curobj is None:
            logging.warning('Bad id')
            return

        curobj.check_plot()

        curitem = self.gettreeitem(curobj)
        self.treelist.CheckItem(curitem, True)

        if id == 'grid':
            curobj.creategrid(gridsize, self.xmin, self.ymin, self.xmax, self.ymax)

    def uncheck_id(self, id=str, unload=True, forceresetOGL=True, askquestion=False):
        """ Uncheck an element from its id """

        curobj = self.getobj_from_id(id)

        if curobj is None:
            logging.warning('Bad id')
            return

        if issubclass(type(curobj), WolfArray):
            curobj.uncheck_plot(unload, forceresetOGL, askquestion)
        else:
            curobj.uncheck_plot()

        curitem = self.gettreeitem(curobj)
        self.treelist.UncheckItem(curitem)

    def get_current_zoom(self):
        """
        Get the current zoom

        :return: dict with keys 'center', 'xmin', 'xmax', 'ymin', 'ymax', 'width', 'height'

        """

        return {'center': (self.mousex, self.mousey),
                'xmin' :  self.xmin,
                'xmax' :  self.xmax,
                'ymin' :  self.ymin,
                'ymax' :  self.ymax,
                'width' : self.xmax-self.xmin,
                'height' : self.ymax-self.ymin}

    def save_current_zoom(self, filepath):
        """ Save the current zoom in a json file """

        zoom = self.get_current_zoom()
        with open(filepath, 'w') as fp:
            json.dump(zoom, fp)

    def read_current_zoom(self, filepath):
        """ Read the current zoom from a json file """

        if exists(filepath):
            with open(filepath, 'r') as fp:
                zoom = json.load(fp)

            self.zoom_on(zoom)


    def menu_bridges(self):

        if self.menu_bridge is None:
            self.menu_bridge = wx.Menu()
            self.menubar.Append(self.menu_bridge, _('&Bridges'))

            self._menu_add_bridge = self.menu_bridge.Append(wx.ID_ANY, _('New bridge'), _('Add a new bridge to the active collection...'))
            self._menu_find_bridge = self.menu_bridge.Append(wx.ID_ANY, _('Pick bridge'), _("Pick the mouse's nearest bridge"))
            self._menu_edit_bridge = self.menu_bridge.Append(wx.ID_ANY, _('Edit bridge'), _('Edit the active bridge'))

            self.Bind(wx.EVT_MENU, self.OnAddBridge, self._menu_add_bridge)
            self.Bind(wx.EVT_MENU, self.OnEditBridge, self._menu_edit_bridge)
            self.Bind(wx.EVT_MENU, self.OnFindBridge, self._menu_find_bridge)

    def OnAddBridge(self, e):
        """ Add a bridge """

        if self.active_bridges is None:
            logging.warning(_('No bridge collection !'))
            return

        newid = wx.TextEntryDialog(None, _('Enter the new bridge id'), _('New bridge id'), 'bridge')
        if newid.ShowModal() == wx.ID_OK:

            newid = newid.GetValue()
            while newid in self.get_list_keys(drawing_type=draw_type.VECTORS):
                newid = newid + '_'

            newbridge = self.active_bridges.addnew(newid)
            self.add_object('vector', newobj=newbridge, id=newid)

    def OnEditBridge(self, e):
        """ Edit a bridge """

        if self.active_bridge is None:
            logging.warning(_('No active bridge to edit !'))
            return

        keys = self.get_list_keys(drawing_type=draw_type.VECTORS)

        newid = self.active_bridge.idx
        while newid in keys:
            newid = newid + '_'

        self.add_object('vector', newobj=self.active_bridge, id= newid)

    def OnFindBridge(self, e):
        """ Find the nearest bridge """

        self.start_action('Pick bridge', _('Right click to pick the nearest bridge'))

    def pick_bridge(self, x:float, y:float):
        """ Find the nearest bridge """

        if self.active_bridges is None:
            logging.warning(_('No bridges to pick !'))
            return

        self.active_bridge = self.active_bridges.find_nearest(x,y)

    def menu_weirs(self):

        if self.menu_weir is None:
            self.menu_weir = wx.Menu()
            self.menubar.Append(self.menu_weir, _('&Weirs'))

            self._menu_add_weir = self.menu_weir.Append(wx.ID_ANY, _('New weir'), _('Add a new weir to the active collection'))
            self._menu_find_weir = self.menu_weir.Append(wx.ID_ANY, _('Pick weir'), _("Pick the mouse's nearest weir"))
            self._menu_edit_weir = self.menu_weir.Append(wx.ID_ANY, _('Edit weir'), _('Edit the active weir'))

            self.Bind(wx.EVT_MENU, self.OnAddWeir, self._menu_add_weir)
            self.Bind(wx.EVT_MENU, self.OnEditWeir, self._menu_edit_weir)
            self.Bind(wx.EVT_MENU, self.OnFindWeir, self._menu_find_weir)

    def OnAddWeir(self, e):
        """ Add a weir """

        logging.info(_('!! To be implemented !!'))
        pass

    def OnEditWeir(self, e):
        """ Edit a weir """

        if self.active_weir is None:
            logging.warning(_('No active weir to edit !'))
            return

        keys = self.get_list_keys(drawing_type=draw_type.VECTORS)

        newid = self.active_weir.idx
        while newid in keys:
            newid = newid + '_'

        self.add_object('vector', newobj=self.active_weir, id= newid)

    def OnFindWeir(self, e):
        """ Find the nearest weir """

        self.start_action('Pick weir', _('Right click to pick the nearest weir'))

    def pick_weir(self, x:float, y:float):
        """ Find the nearest weir """

        if self.active_weirs is None:
            logging.warning(_('No weirs to pick !'))
            return

        self.active_weir = self.active_weirs.find_nearest(x,y)



class Comp_Type(Enum):
    ARRAYS = 1
    ARRAYS_MB = 2
    RES2D = 3
    RES2D_GPU = 4
class Compare_Arrays_Results():

    def __init__(self, parent:WolfMapViewer = None, share_cmap_array:bool = False, share_cmap_diff:bool = False):

        self.parent = parent

        self.paths = []
        self.elements = []
        self.linked_elts = []
        self.diff = []
        self.mapviewers = []
        self.mapviewers_diff = []

        self.times = None

        self.share_cmap_array = share_cmap_array
        self.share_cmap_diff  = share_cmap_diff

        self.type = Comp_Type.ARRAYS

        self._initialized_viewers = False
        self.independent = True

    def _check_type(self, file:Path):
        """
        Check the type of the file/directory

        If it is a file and suffix is empty, it is considered as RES2D.
        If it is a directory and contains a simul_gpu_results, it is considered as RES2D_GPU.
        If it is a file and suffix is not empty, it is considered as ARRAYS. A check is done to see if it is a multi-block array.

        """

        file = Path(file)

        if file.suffix == '' and not file.is_dir():

            return Comp_Type.RES2D, file

        elif file.suffix in ('.bin', '.tif', '.tiff', '.npy', '.npz', '.top', '.frott', '.nap', '.hbin', '.hbinb', '.qxbin', '.qxbinb', '.qybin', '.qybinb', '.inf') :

            if file.suffix in ('.bin', '.top', '.frott', '.nap', '.hbin', '.hbinb', '.qxbin', '.qxbinb', '.qybin', '.qybinb', '.inf'):
                if file.with_suffix(file.suffix + '.txt').exists():

                    test = WolfArray(file, preload=False)
                    test.read_txt_header()

                    mb = test.nb_blocks > 0

                    if mb:
                        return Comp_Type.ARRAYS_MB, file

            return Comp_Type.ARRAYS, file

        elif (file.parent / 'simul_gpu_results').exists():
            file = file.parent / 'simul_gpu_results'
            return Comp_Type.RES2D_GPU, file

        elif (file.parent.parent / 'simul_gpu_results').exists():
            file = file.parent.parent / 'simul_gpu_results'
            return Comp_Type.RES2D_GPU, file
        else:
            return None, None

    def add(self, file_or_dir:Union[str, Path] = None):

        if file_or_dir is None:

            filterProject = "all (*.*)|*.*"
            file = wx.FileDialog(None, "Choose array/model", wildcard=filterProject)

            if file.ShowModal() == wx.ID_CANCEL:
                file.Destroy()
                return False
            else:
                filename = Path(file.GetPath())
                file.Destroy()

        self.paths.append(self._check_type(filename))

        if self.paths[-1][0] is None:
            logging.warning(_('File type not recognized -- Retry !'))
            self.paths.pop()
            return False

        return True

    def check(self):
        """ Check the consystency of the elements to compare """

        reftype = self.paths[0][0]
        for cur in self.paths:
            if cur[0] != reftype:
                logging.warning(_('Inconsistency in the type of the elements to compare'))
                return False

        return True

    def update_comp(self, idx=list[int]):
        """
        Update Arrays from 2D modellings

        :param idx: indexes of the time step to update --> steps to read

        """
        assert self.type in (Comp_Type.RES2D, Comp_Type.RES2D_GPU), 'This method is only for 2D results'

        self.linked_elts = []

        for curelt, curstep in zip(self.elements, idx):

            curelt.read_oneresult(curstep)

            self.linked_elts.append(curelt.as_WolfArray())

            for curelt, curlink in zip(self.elements, self.linked_elts):
                curlink.idx = curelt.idx + ' ' + curelt.get_currentview().value

        self.set_diff()

        if self._initialized_viewers:
            self.update_viewers()

    def update_type_result(self, newtype):
        """
        Update the result type for each element

        """
        assert newtype in views_2D, 'This type is not a 2D result'
        assert self.type in (Comp_Type.RES2D, Comp_Type.RES2D_GPU), 'This method is only for 2D results'

        for curelt in self.elements:
            curelt.set_currentview(newtype, force_updatepal = True)

        # remove elements
        for baselt, curelt, curmap in zip(self.elements, self.linked_elts, self.mapviewers):

            curmap.removeobj_from_id(curelt.idx)

        for curdiff, curmap in zip(self.diff, self.mapviewers_diff):

            curmap.removeobj_from_id(curdiff.idx)


        self.update_comp(self.times.get_times_idx())

    def set_elements(self):
        """ Set the elements to compare with the right type """
        from .ui.wolf_times_selection_comparison_models import Times_Selection

        if self.check():
            self.type = self.paths[0][0]

            if self.type == Comp_Type.RES2D_GPU:
                self.parent.menu_wolf2d()
                self.elements = [wolfres2DGPU(cur[1], plotted=False, idx = cur[1].name + '_' + str(idx)) for idx, cur in enumerate(self.paths)]

                times = [curmod.get_times_steps()[0] for curmod in self.elements]

                self.times = Times_Selection(self, wx.ID_ANY, _("Times"), size=(400,400), times = times, callback = self.update_comp)
                self.times.Show()

            elif self.type == Comp_Type.RES2D:
                self.parent.menu_wolf2d()
                self.elements = [Wolfresults_2D(cur[1], plotted=False, idx = cur[1].name + '_' + str(idx)) for idx, cur in enumerate(self.paths)]

                times = [curmod.get_times_steps()[0] for curmod in self.elements]

                self.times = Times_Selection(self, wx.ID_ANY, _("Times"), size=(400,400), times = times, callback = self.update_comp)
                self.times.Show()

            elif self.type == Comp_Type.ARRAYS:
                self.elements = [WolfArray(cur[1], plotted=False, idx = cur[1].name + '_' + str(idx)) for idx, cur in enumerate(self.paths)]

            elif self.type == Comp_Type.ARRAYS_MB:
                self.elements = [WolfArrayMB(cur[1], plotted=False, idx = cur[1].name + '_' + str(idx)) for idx, cur in enumerate(self.paths)]


    def set_diff(self):
        """ Set the differential between the elements and the first one, which is the reference """

        if self.type in (Comp_Type.ARRAYS, Comp_Type.ARRAYS_MB):

            ref = self.elements[0]

            # Recherche d'un masque union des masques partiels
            ref.mask_unions(self.elements[1:])

            # Création du différentiel -- Les opérateurs mathématiques sont surchargés
            self.diff = [cur - ref for cur in self.elements[1:]]

            for curdiff, cur in zip(self.diff, self.elements[1:]):
                curdiff.idx = _('Difference') + cur.idx +' - ' + ref.idx

        elif self.type in (Comp_Type.RES2D, Comp_Type.RES2D_GPU):

            if len(self.linked_elts) == 0:
                self.update_comp([-1] * len(self.elements))

            elif len(self.linked_elts) == len(self.elements):
                ref = self.linked_elts[0]

                self.diff = [cur - ref for cur in self.linked_elts[1:]]

                for curdiff, cur in zip(self.diff, self.linked_elts[1:]):
                    curdiff.idx = _('Difference') + cur.idx +' - ' + ref.idx

    def set_viewers(self, independent:bool = None):
        """
        Set viewers

        """

        if independent is None:
            dlg = wx.MessageDialog(None, _("Create a viewer for each element ?"), _("Viewers"), style = wx.YES_NO|wx.YES_DEFAULT)
            ret = dlg.ShowModal()

            self.independent = ret == wx.ID_YES
        else:
            self.independent = independent

        if not self.independent:
            self.mapviewers = [self.parent] * len(self.elements)
            self.mapviewers_diff = self.mapviewers
        else:
            # Création de plusieurs fenêtres de visualisation basées sur la classe "WolfMapViewer"
            self.mapviewers = []

            self.mapviewers.append(self.parent) # parent as viewer for first element

            for id, file in enumerate(self.elements[1:]):

                self.mapviewers.append(WolfMapViewer(None, file.idx, w=600, h=600, wxlogging=self.parent.wxlogging, wolfparent = self.parent.wolfparent))
                self.mapviewers_diff.append(WolfMapViewer(None, 'Difference' + file.idx, w=600, h=600, wxlogging=self.parent.wxlogging, wolfparent = self.parent.wolfparent))

            for curviewer in self.mapviewers[1:] + self.mapviewers_diff:
                curviewer.add_grid()
                curviewer.add_WMS()

        for curviewer in self.mapviewers + self.mapviewers_diff:
            curviewer.linked = True
            curviewer.linkedList = self.mapviewers + self.mapviewers_diff

        self._initialized_viewers = True

        self.update_viewers()

    def set_shields_param(self, diamsize:float = .001, graindensity:float = 2.65):
        """ Set the parameters for the shields diagram """

        for curelt in self.elements:
            curelt.sediment_diameter = diamsize
            curelt.sediment_density = graindensity
            curelt.load_default_colormap('shields_cst')

    def update_viewers(self):
        """ Update the viewers with the new elements """

        if self.type in (Comp_Type.ARRAYS, Comp_Type.ARRAYS_MB):
            elts = self.elements
        elif self.type in (Comp_Type.RES2D, Comp_Type.RES2D_GPU):
            elts = self.linked_elts

        # on attribue une matrice par interface graphique
        ref = elts[0]
        for baselt, curelt, curmap in zip(self.elements, elts, self.mapviewers):

            # if self.type in (Comp_Type.RES2D, Comp_Type.RES2D_GPU):
            #     curmap.active_res2d = baselt

            curmap.removeobj_from_id(curelt.idx)

            curelt.change_gui(curmap)
            curmap.active_array = curelt
            curelt.myops.myzones = ref.myops.myzones

        # diff = self.diff[0]
        for curdiff, curmap in zip(self.diff, self.mapviewers_diff):

            curmap.removeobj_from_id(curdiff.idx)

            curdiff.change_gui(curmap)
            curmap.active_array = curdiff
            curdiff.myops.myzones = ref.myops.myzones

        # on partage la palette de couleurs
        ref.mypal.automatic = False
        ref.myops.palauto.SetValue(0)

        if self.share_cmap_array:
            for curelt in elts[1:]:
                curelt.mypal.automatic = False
                curelt.myops.palauto.SetValue(0)
                ref.add_crosslinked_array(curelt)
                ref.share_palette()
        else:
            for curelt in elts[1:]:
                curelt.mypal.automatic = False
                curelt.myops.palauto.SetValue(0)
                curelt.mypal.updatefrompalette(ref.mypal)

        #palette de la différence
        diff = self.diff[0]
        diff.mypal = wolfpalette()
        if isinstance(diff, WolfArrayMB):
            diff.link_palette()

        path = os.path.dirname(__file__)
        fn = join(path, 'models\\diff16.pal')

        diff.mypal.readfile(fn)
        diff.mypal.automatic = False
        diff.myops.palauto.SetValue(0)

        if self.share_cmap_diff:
            for curelt in self.diff[1:]:
                curelt.mypal.automatic = False
                curelt.myops.palauto.SetValue(0)
                diff.add_crosslinked_array(curelt)
                diff.share_palette()
        else:
            for curelt in self.diff[1:]:
                curelt.mypal.automatic = False
                curelt.myops.palauto.SetValue(0)
                curelt.mypal.updatefrompalette(diff.mypal)

        # Ajout des matrices dans les fenêtres de visualisation
        for curelt, curmap in zip(elts, self.mapviewers):
            curmap.add_object('array', newobj = curelt, ToCheck = True, id = curelt.idx)

        for curdiff, curmap in zip(self.diff, self.mapviewers_diff):
            curmap.add_object('array', newobj = curdiff, ToCheck = True, id = curdiff.idx)

        if self.independent:
            for curmap in self.mapviewers + self.mapviewers_diff:
                curmap.Refresh()
        else:
            self.mapviewers[0].Refresh()

    def bake(self):

        self.set_elements()
        self.set_diff()
        self.set_viewers()

class InPaint_waterlevel(wx.Dialog):

    def __init__(self, parent, title:str = _('Inpainting'), size:tuple[int,int] = (400,400), mapviewer:WolfMapViewer=None, **kwargs):

        super().__init__(parent, title = title, size = size, **kwargs)

        self._array: WolfArray = None
        self._dem: WolfArray = None
        self._dtm: WolfArray = None

        self._mapviewer = mapviewer

        self._init_UI()

    def _init_UI(self):
        """ Create 2 listboxes for the arrays and the masks """
        import shutil

        if self._mapviewer is None:
            logging.warning(_('No mapviewer --> Nothing to do'))
            return

        self._sizer = wx.BoxSizer(wx.VERTICAL)

        self._sizer_lists = wx.BoxSizer(wx.HORIZONTAL)
        self._sizer_arrays = wx.BoxSizer(wx.VERTICAL)
        self._sizer_ignore = wx.BoxSizer(wx.VERTICAL)

        self._label_arrays = wx.StaticText(self, wx.ID_ANY, _('Array'))
        self._listbox_arrays = wx.ListBox(self, wx.ID_ANY, choices = self._mapviewer.get_list_keys(drawing_type = draw_type.ARRAYS), style = wx.LB_SINGLE)
        self._label_dem = wx.StaticText(self, wx.ID_ANY, _('DEM'))
        self._listbox_dems = wx.ListBox(self, wx.ID_ANY, choices = ['None'] + self._mapviewer.get_list_keys(drawing_type = draw_type.ARRAYS), style = wx.LB_SINGLE)
        self._label_dtm = wx.StaticText(self, wx.ID_ANY, _('DTM'))
        self._listbox_dtm = wx.ListBox(self, wx.ID_ANY, choices = ['None'] + self._mapviewer.get_list_keys(drawing_type = draw_type.ARRAYS), style = wx.LB_SINGLE)

        self._label_ignore = wx.StaticText(self, wx.ID_ANY, _('Ignore last holes'))
        self._listbox_ignore = wx.ListBox(self, wx.ID_ANY, choices = [str(i) for i in range(10)], style = wx.LB_SINGLE)

        self._sizer_ignore.Add(self._label_ignore, 0, wx.EXPAND)
        self._sizer_ignore.Add(self._listbox_ignore, 1, wx.EXPAND)

        self._sizer_arrays.Add(self._label_arrays, 0, wx.EXPAND)
        self._sizer_arrays.Add(self._listbox_arrays, 1, wx.EXPAND)
        self._sizer_arrays.Add(self._label_dem, 0, wx.EXPAND)
        self._sizer_arrays.Add(self._listbox_dems, 1, wx.EXPAND)
        self._sizer_arrays.Add(self._label_dtm, 0, wx.EXPAND)
        self._sizer_arrays.Add(self._listbox_dtm, 1, wx.EXPAND)

        self._sizer_lists.Add(self._sizer_arrays, 1, wx.EXPAND)
        self._sizer_lists.Add(self._sizer_ignore, 1, wx.EXPAND)

        self._sizer.Add(self._sizer_lists, 1, wx.EXPAND)

        self._sizer_btns = wx.BoxSizer(wx.HORIZONTAL)

        self._sizer_inpaint = wx.BoxSizer(wx.VERTICAL)
        self._btn_inpaint = wx.Button(self, wx.ID_ANY, _('Inpaint'))
        self._sizer_inpaint.Add(self._btn_inpaint, 1, wx.EXPAND)

        self._check_fortran = wx.CheckBox(self, wx.ID_ANY, _('Use Fortran'))
        self._check_fortran.SetValue(False)
        if shutil.which('holes.exe') is not None:
            self._sizer_inpaint.Add(self._check_fortran, 1, wx.EXPAND)

        self._btn_update_ids = wx.Button(self, wx.ID_ANY, _('Update IDs'))
        self._btn_select_holes = wx.Button(self, wx.ID_ANY, _('Select holes'))
        self._btn_create_mask = wx.Button(self, wx.ID_ANY, _('Create mask'))

        self._sizer_btns.Add(self._sizer_inpaint, 1, wx.EXPAND)
        self._sizer_btns.Add(self._btn_update_ids, 1, wx.EXPAND)
        self._sizer_btns.Add(self._btn_select_holes, 1, wx.EXPAND)
        self._sizer_btns.Add(self._btn_create_mask, 1, wx.EXPAND)

        self._sizer.Add(self._sizer_btns, 1, wx.EXPAND)

        self.SetSizer(self._sizer)

        self._listbox_arrays.Bind(wx.EVT_LISTBOX, self.OnSelectArray)
        self._listbox_dems.Bind(wx.EVT_LISTBOX, self.OnSelectMask)
        self._listbox_dtm.Bind(wx.EVT_LISTBOX, self.OnSelectDTM)
        self._btn_inpaint.Bind(wx.EVT_BUTTON, self.OnInpaint)
        self._btn_update_ids.Bind(wx.EVT_BUTTON, self.OnUpdateIDs)
        self._btn_select_holes.Bind(wx.EVT_BUTTON, self.OnSelectHoles)
        self._btn_create_mask.Bind(wx.EVT_BUTTON, self.OnCreateMask)

        self._listbox_dems.SetSelection(0)
        self._listbox_dtm.SetSelection(0)
        self._listbox_ignore.SetSelection(1)

        self.CenterOnScreen()
        self.Show()

    def OnUpdateIDs(self, e):
        """ Update the list of arrays/mask/dtm """

        self._listbox_arrays.Set(self._mapviewer.get_list_keys(drawing_type = draw_type.ARRAYS))
        self._listbox_dems.Set(['None'] + self._mapviewer.get_list_keys(drawing_type = draw_type.ARRAYS))
        self._listbox_dtm.Set(['None'] + self._mapviewer.get_list_keys(drawing_type = draw_type.ARRAYS))

    def OnSelectArray(self, e):
        """ Select an array """

        self._array = self._mapviewer.getobj_from_id(self._listbox_arrays.GetStringSelection())

    def OnSelectMask(self, e):
        """ Select a mask """

        mask_ = self._listbox_dems.GetStringSelection()
        if mask_ == 'None':
            self._dem = None
        else:
            self._dem = self._mapviewer.getobj_from_id(self._listbox_dems.GetStringSelection())

    def OnSelectDTM(self, e):
        """ Select a DTM """

        dtm_ = self._listbox_dtm.GetStringSelection()
        if dtm_ == 'None':
            self._dtm = None
        else:
            self._dtm = self._mapviewer.getobj_from_id(self._listbox_dtm.GetStringSelection())

    def OnInpaint(self, e):
        """ Inpaint the array with the mask """

        if self._array is None:
            logging.warning(_('Select an array, a mask and a DTM'))

        else:
            times, wl, wd = self._array._inpaint_waterlevel_dem_dtm(self._dem, self._dtm, ignore_last= self._listbox_ignore.GetSelection(), use_fortran= self._check_fortran.GetValue())
            logging.info(_('Inpainting done !'))

            dlg = wx.MessageDialog(None, _('Add water depth to the viewer ?'), _('Water depth'), style = wx.YES_NO|wx.YES_DEFAULT)
            ret = dlg.ShowModal()

            if ret == wx.ID_YES:
                self._mapviewer.add_object('array', newobj = wd, id = 'wd_' + self._array.idx)

            dlg.Destroy()

    def OnSelectHoles(self, e):
        """ Select the holes in the array """

        if self._array is None:
            logging.warning(_('Select an array'))
            return

        self._array.select_holes(ignore_last = self._listbox_ignore.GetSelection())
        self._mapviewer.Paint()

    def OnCreateMask(self, e):
        """ Create a mask from the array """

        if self._array is None:
            logging.warning(_('Select an array, a mask and a DTM'))
            return

        if self._dem is None:
            logging.warning(_('Select a dem'))
            return

        if self._dtm is None:
            logging.warning(_('Select a dtm'))
            return

        newmask = self._array._create_building_holes_dem_dtm(self._dem, self._dtm, ignore_last= self._listbox_ignore.GetSelection())
        self._mapviewer.add_object('array', newobj = newmask, id = 'mask_' + self._array.idx)

class InPaint_array(wx.Dialog):

    def __init__(self, parent, title:str = _('Inpainting'), size:tuple[int,int] = (400,400), mapviewer:WolfMapViewer=None, **kwargs):

        super().__init__(parent, title = title, size = size, **kwargs)

        self._array: WolfArray = None
        self._mask: WolfArray = None
        self._test: WolfArray = None

        self._mapviewer = mapviewer

        self._init_UI()

    def _init_UI(self):
        """ Create 2 listboxes for the arrays and the masks """
        import shutil

        if self._mapviewer is None:
            logging.warning(_('No mapviewer --> Nothing to do'))
            return

        self._sizer = wx.BoxSizer(wx.VERTICAL)

        self._sizer_lists = wx.BoxSizer(wx.HORIZONTAL)
        self._sizer_arrays = wx.BoxSizer(wx.VERTICAL)
        self._sizer_ignore = wx.BoxSizer(wx.VERTICAL)

        self._label_arrays = wx.StaticText(self, wx.ID_ANY, _('Array'))
        self._listbox_arrays = wx.ListBox(self, wx.ID_ANY, choices = self._mapviewer.get_list_keys(drawing_type = draw_type.ARRAYS), style = wx.LB_SINGLE)
        self._label_masks = wx.StaticText(self, wx.ID_ANY, _('Mask == where to inpaint'))
        self._listbox_masks = wx.ListBox(self, wx.ID_ANY, choices = ['None'] + self._mapviewer.get_list_keys(drawing_type = draw_type.ARRAYS), style = wx.LB_SINGLE)
        self._label_test = wx.StaticText(self, wx.ID_ANY, _('Test == local inpainted value must be greater than this value'))
        self._listbox_test = wx.ListBox(self, wx.ID_ANY, choices = ['None'] + self._mapviewer.get_list_keys(drawing_type = draw_type.ARRAYS), style = wx.LB_SINGLE)

        self._label_ignore = wx.StaticText(self, wx.ID_ANY, _('Ignore last holes'))
        self._listbox_ignore = wx.ListBox(self, wx.ID_ANY, choices = [str(i) for i in range(10)], style = wx.LB_SINGLE)

        self._sizer_ignore.Add(self._label_ignore, 0, wx.EXPAND)
        self._sizer_ignore.Add(self._listbox_ignore, 1, wx.EXPAND)

        self._sizer_arrays.Add(self._label_arrays, 0, wx.EXPAND)
        self._sizer_arrays.Add(self._listbox_arrays, 1, wx.EXPAND)
        self._sizer_arrays.Add(self._label_masks, 0, wx.EXPAND)
        self._sizer_arrays.Add(self._listbox_masks, 1, wx.EXPAND)
        self._sizer_arrays.Add(self._label_test, 0, wx.EXPAND)
        self._sizer_arrays.Add(self._listbox_test, 1, wx.EXPAND)

        self._sizer_lists.Add(self._sizer_arrays, 1, wx.EXPAND)
        self._sizer_lists.Add(self._sizer_ignore, 1, wx.EXPAND)

        self._sizer.Add(self._sizer_lists, 1, wx.EXPAND)

        self._sizer_btns = wx.BoxSizer(wx.HORIZONTAL)

        self._sizer_inpaint = wx.BoxSizer(wx.VERTICAL)
        self._btn_inpaint = wx.Button(self, wx.ID_ANY, _('Inpaint'))
        self._sizer_inpaint.Add(self._btn_inpaint, 1, wx.EXPAND)

        self._btn_update_ids = wx.Button(self, wx.ID_ANY, _('Update IDs'))
        self._btn_select_holes = wx.Button(self, wx.ID_ANY, _('Select holes'))
        self._btn_create_mask = wx.Button(self, wx.ID_ANY, _('Create mask'))

        self._sizer_btns.Add(self._sizer_inpaint, 1, wx.EXPAND)
        self._sizer_btns.Add(self._btn_update_ids, 1, wx.EXPAND)
        self._sizer_btns.Add(self._btn_select_holes, 1, wx.EXPAND)
        self._sizer_btns.Add(self._btn_create_mask, 1, wx.EXPAND)

        self._sizer.Add(self._sizer_btns, 1, wx.EXPAND)

        self.SetSizer(self._sizer)

        self._listbox_arrays.Bind(wx.EVT_LISTBOX, self.OnSelectArray)
        self._listbox_masks.Bind(wx.EVT_LISTBOX, self.OnSelectMask)
        self._listbox_test.Bind(wx.EVT_LISTBOX, self.OnSelectTest)
        self._btn_inpaint.Bind(wx.EVT_BUTTON, self.OnInpaint)
        self._btn_update_ids.Bind(wx.EVT_BUTTON, self.OnUpdateIDs)
        self._btn_select_holes.Bind(wx.EVT_BUTTON, self.OnSelectHoles)
        self._btn_create_mask.Bind(wx.EVT_BUTTON, self.OnCreateMask)

        self._listbox_masks.SetSelection(0)
        self._listbox_test.SetSelection(0)
        self._listbox_ignore.SetSelection(0)

        self.CenterOnScreen()
        self.Show()

    def OnUpdateIDs(self, e):
        """ Update the list of arrays/mask/dtm """

        self._listbox_arrays.Set(self._mapviewer.get_list_keys(drawing_type = draw_type.ARRAYS))
        self._listbox_masks.Set(['None'] + self._mapviewer.get_list_keys(drawing_type = draw_type.ARRAYS))
        self._listbox_test.Set(['None'] + self._mapviewer.get_list_keys(drawing_type = draw_type.ARRAYS))

    def OnSelectArray(self, e):
        """ Select an array """

        self._array = self._mapviewer.getobj_from_id(self._listbox_arrays.GetStringSelection())

    def OnSelectMask(self, e):
        """ Select a mask """

        mask_ = self._listbox_masks.GetStringSelection()
        if mask_ == 'None':
            self._mask = None
        else:
            self._mask = self._mapviewer.getobj_from_id(self._listbox_masks.GetStringSelection())

    def OnSelectTest(self, e):
        """ Select a DTM """

        dtm_ = self._listbox_test.GetStringSelection()
        if dtm_ == 'None':
            self._test = None
        else:
            self._test = self._mapviewer.getobj_from_id(self._listbox_test.GetStringSelection())

    def OnInpaint(self, e):
        """ Inpaint the array with the mask """

        if self._array is None:
            logging.warning(_('Select an array, a mask and a DTM'))

        else:

            times, wl, wd = self._array.inpaint(self._mask, self._test, ignore_last= self._listbox_ignore.GetSelection())
            logging.info(_('Inpainting done !'))

            dlg = wx.MessageDialog(None, _('Add extra array to the viewer ?'), _('Inerpolation - test data'), style = wx.YES_NO|wx.YES_DEFAULT)
            ret = dlg.ShowModal()

            if ret == wx.ID_YES:
                self._mapviewer.add_object('array', newobj = wd, id = 'extra_' + self._array.idx)

            dlg.Destroy()

    def OnSelectHoles(self, e):
        """ Select the holes in the array """

        if self._mask is None:
            logging.warning(_('Select a mask array'))
            return

        self._mask.select_holes(ignore_last = self._listbox_ignore.GetSelection())
        self._mapviewer.Paint()

    def OnCreateMask(self, e):
        """ Create a mask from the array """

        if self._array is None:
            logging.warning(_('Select an array'))
            return

        newmask = self._array.create_mask_holes(ignore_last= self._listbox_ignore.GetSelection())
        self._mapviewer.add_object('array', newobj = newmask, id = 'mask_' + self._array.idx)
