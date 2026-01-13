#!/usr/bin/env python3

"""pdf2ppt

Convert PDF Slides to PowerPoint Presentations (PPT)
Author: Teddy van Jerry (Wuqiong Zhao)
License: MIT
GitHub: https://github.com/Teddy-van-Jerry/pdf2ppt
"""

import argparse
import os
import sys
import shutil
import subprocess
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed

from pptx import Presentation
from pptx.util import Pt
from pypdf import PdfReader
from rich.console import Console
from rich.progress import Progress, SpinnerColumn, TextColumn, BarColumn, TaskProgressColumn
from rich.panel import Panel
from rich.table import Table

__version__ = '1.2.0'
TMP_DIR_NAME = '_pdf2ppt.tmp'
ERR_INPUT_NOT_FOUND = 1
ERR_DEPENDENCY_MISSING = 2
ERR_PDF2SVG = 101
ERR_SVG2EMF = 102

console = Console()


def check_dependency(cmd: str, name: str) -> bool:
    """Check if a command-line tool is available."""
    try:
        # pdf2svg doesn't have --version, just check if it exists
        if 'pdf2svg' in cmd:
            result = subprocess.run([cmd], capture_output=True, timeout=10)
            return True  # pdf2svg returns non-zero but exists
        else:
            result = subprocess.run([cmd, '--version'], capture_output=True, timeout=10)
            return True
    except FileNotFoundError:
        return False
    except subprocess.TimeoutExpired:
        return True  # If it times out, it probably exists


def check_dependencies(pdf2svg_path: str, inkscape_path: str) -> bool:
    """Check all required dependencies."""
    missing = []
    
    if not check_dependency(pdf2svg_path, 'pdf2svg'):
        missing.append(('pdf2svg', 'brew install pdf2svg', 'sudo apt install pdf2svg'))
    
    if not check_dependency(inkscape_path, 'inkscape'):
        missing.append(('Inkscape', 'brew install inkscape', 'sudo apt install inkscape'))
    
    if missing:
        console.print("\n[bold red]❌ Missing dependencies:[/bold red]\n")
        table = Table(show_header=True, header_style="bold magenta")
        table.add_column("Tool")
        table.add_column("macOS")
        table.add_column("Ubuntu/Debian")
        for tool, mac_cmd, linux_cmd in missing:
            table.add_row(tool, f"[cyan]{mac_cmd}[/cyan]", f"[cyan]{linux_cmd}[/cyan]")
        console.print(table)
        return False
    return True


def parse_page_range(page_str: str, total_pages: int) -> list:
    """Parse page range string like '1-5,7,9-11' into list of page numbers."""
    if not page_str:
        return list(range(1, total_pages + 1))
    
    pages = set()
    for part in page_str.split(','):
        part = part.strip()
        if '-' in part:
            start, end = part.split('-', 1)
            start = int(start) if start else 1
            end = int(end) if end else total_pages
            pages.update(range(max(1, start), min(total_pages, end) + 1))
        else:
            page = int(part)
            if 1 <= page <= total_pages:
                pages.add(page)
    return sorted(pages)


def pdf2svg(pdf_path: Path, pdf2svg_path: str, verbose: bool = False) -> bool:
    """Convert PDF to SVG using pdf2svg."""
    tmp_dir = pdf_path.parent / TMP_DIR_NAME
    tmp_dir.mkdir(parents=True, exist_ok=True)
    pdf_name = pdf_path.stem
    
    cmd = [pdf2svg_path, str(pdf_path), str(tmp_dir / f'{pdf_name}_%d.svg'), 'all']
    if verbose:
        console.print(f"[dim]Running: {' '.join(cmd)}[/dim]")
    
    result = subprocess.run(cmd, capture_output=True)
    return result.returncode == 0


def convert_single_svg(args: tuple) -> tuple:
    """Convert a single SVG to EMF. Returns (page_num, success, has_filter)."""
    svg_path, emf_path, inkscape_path, no_check = args
    page_num = int(svg_path.stem.split('_')[-1])
    
    cmd = [inkscape_path, '--export-type=emf', str(svg_path)]
    result = subprocess.run(cmd, capture_output=True)
    
    if result.returncode != 0:
        return (page_num, False, False)
    
    has_filter = False
    if not no_check:
        with open(svg_path, 'r', encoding='utf-8', errors='ignore') as f:
            if 'filter=' in f.read():
                has_filter = True
    
    return (page_num, True, has_filter)


def svg2emf(pdf_reader: PdfReader, pdf_path: Path, inkscape_path: str, 
            pages: list, verbose: bool = False, no_check: bool = False, 
            parallel: int = 1) -> tuple:
    """Convert SVG to EMF using inkscape."""
    tmp_dir = pdf_path.parent / TMP_DIR_NAME
    pdf_name = pdf_path.stem
    pages_with_filters = []
    pages_with_filters_svg_dir = pdf_path.parent / f'{pdf_name}_svg'
    
    # Prepare conversion tasks
    tasks = []
    for page in pages:
        svg_path = tmp_dir / f'{pdf_name}_{page}.svg'
        emf_path = tmp_dir / f'{pdf_name}_{page}.emf'
        tasks.append((svg_path, emf_path, inkscape_path, no_check))
    
    with Progress(
        SpinnerColumn(),
        TextColumn("[progress.description]{task.description}"),
        BarColumn(),
        TaskProgressColumn(),
        console=console,
        transient=not verbose
    ) as progress:
        task = progress.add_task("[cyan]Converting SVG to EMF...", total=len(tasks))
        
        if parallel > 1:
            with ThreadPoolExecutor(max_workers=parallel) as executor:
                futures = {executor.submit(convert_single_svg, t): t for t in tasks}
                for future in as_completed(futures):
                    page_num, success, has_filter = future.result()
                    if not success:
                        return (False, [])
                    if has_filter:
                        pages_with_filters.append(page_num)
                        if len(pages_with_filters) == 1:
                            shutil.rmtree(pages_with_filters_svg_dir, ignore_errors=True)
                            pages_with_filters_svg_dir.mkdir(parents=True, exist_ok=True)
                        svg_path = tmp_dir / f'{pdf_name}_{page_num}.svg'
                        shutil.copy(svg_path, pages_with_filters_svg_dir / f'{pdf_name}_{page_num}.svg')
                    progress.advance(task)
        else:
            for t in tasks:
                page_num, success, has_filter = convert_single_svg(t)
                if not success:
                    return (False, [])
                if has_filter:
                    pages_with_filters.append(page_num)
                    if len(pages_with_filters) == 1:
                        shutil.rmtree(pages_with_filters_svg_dir, ignore_errors=True)
                        pages_with_filters_svg_dir.mkdir(parents=True, exist_ok=True)
                    svg_path = tmp_dir / f'{pdf_name}_{page_num}.svg'
                    shutil.copy(svg_path, pages_with_filters_svg_dir / f'{pdf_name}_{page_num}.svg')
                progress.advance(task)
    
    return (True, sorted(pages_with_filters))


def emf2ppt(pdf_reader: PdfReader, pdf_path: Path, ppt_path: Path, 
            pages: list, verbose: bool = False):
    """Convert EMF files to PowerPoint presentation."""
    tmp_dir = pdf_path.parent / TMP_DIR_NAME
    pdf_name = pdf_path.stem

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    
    pdf_box = pdf_reader.pages[0].mediabox
    slide_width, slide_height = float(pdf_box[2]), float(pdf_box[3])
    prs.slide_width = Pt(slide_width)
    prs.slide_height = Pt(slide_height)

    # Add metadata
    if pdf_reader.metadata:
        if pdf_reader.metadata.title:
            prs.core_properties.title = pdf_reader.metadata.title
        if pdf_reader.metadata.author:
            prs.core_properties.author = pdf_reader.metadata.author
        if pdf_reader.metadata.subject:
            prs.core_properties.subject = pdf_reader.metadata.subject
        if pdf_reader.metadata.creation_date:
            prs.core_properties.created = pdf_reader.metadata.creation_date
    
    prs.core_properties.comments = 'Generated using pdf2ppt (https://github.com/neosun100/pdf2ppt)'
    prs.core_properties.last_modified_by = 'pdf2ppt'

    with Progress(
        SpinnerColumn(),
        TextColumn("[progress.description]{task.description}"),
        BarColumn(),
        TaskProgressColumn(),
        console=console,
        transient=not verbose
    ) as progress:
        task = progress.add_task("[cyan]Creating PowerPoint...", total=len(pages))
        
        for page in pages:
            slide = prs.slides.add_slide(blank_slide_layout)
            emf_path = tmp_dir / f'{pdf_name}_{page}.emf'
            slide.shapes.add_picture(str(emf_path), Pt(0), Pt(0), width=prs.slide_width)
            progress.advance(task)
    
    prs.save(ppt_path)


def clean_tmp(pdf_path: Path, verbose: bool = False):
    """Clean up temporary files."""
    tmp_dir = pdf_path.parent / TMP_DIR_NAME
    if tmp_dir.exists():
        shutil.rmtree(tmp_dir)
        if verbose:
            console.print("[dim]Cleaned up temporary files.[/dim]")


def main():
    parser = argparse.ArgumentParser(
        prog='pdf2ppt',
        description='Convert PDF Slides to PowerPoint Presentations with Vector Graphics'
    )
    parser.add_argument('-v', '--version', action='version', version=f'%(prog)s {__version__}')
    parser.add_argument('input', type=Path, help='Input PDF file')
    parser.add_argument('output', type=Path, nargs='?', help='Output PPTX file (default: input.pptx)')
    parser.add_argument('--verbose', action='store_true', help='Verbose output')
    parser.add_argument('--no-clean', action='store_true', help='Keep temporary files')
    parser.add_argument('--no-check', action='store_true', help='Skip SVG filter check')
    parser.add_argument('--force', '-f', action='store_true', help='Overwrite output file if exists')
    parser.add_argument('--pages', '-p', type=str, help='Page range (e.g., "1-5,7,9-11")')
    parser.add_argument('--parallel', '-j', type=int, default=1, help='Parallel workers (default: 1)')
    parser.add_argument('--pdf2svg-path', type=str, default='pdf2svg', help='Path to pdf2svg')
    parser.add_argument('--inkscape-path', type=str, default='inkscape', help='Path to inkscape')
    
    args = parser.parse_args()
    
    # Check input file
    if not args.input.exists():
        console.print(f"[bold red]❌ Error:[/bold red] Input file not found: {args.input}")
        sys.exit(ERR_INPUT_NOT_FOUND)
    
    # Check dependencies
    if not check_dependencies(args.pdf2svg_path, args.inkscape_path):
        sys.exit(ERR_DEPENDENCY_MISSING)
    
    # Determine output path
    ppt_path = args.output if args.output else args.input.with_suffix('.pptx')
    
    # Check output file
    if ppt_path.exists() and not args.force:
        console.print(f"[bold yellow]⚠️  Output file exists:[/bold yellow] {ppt_path}")
        console.print("Use [cyan]--force[/cyan] to overwrite.")
        sys.exit(1)
    
    # Show header
    console.print(Panel.fit(
        f"[bold blue]pdf2ppt[/bold blue] v{__version__}\n"
        f"[dim]Converting:[/dim] {args.input.name} → {ppt_path.name}",
        border_style="blue"
    ))
    
    # Read PDF
    pdf_reader = PdfReader(args.input)
    total_pages = len(pdf_reader.pages)
    
    # Parse page range
    pages = parse_page_range(args.pages, total_pages)
    if args.verbose:
        console.print(f"[dim]Processing {len(pages)} of {total_pages} pages[/dim]")
    
    # Step 1: PDF to SVG
    with console.status("[bold green]Converting PDF to SVG..."):
        if not pdf2svg(args.input, args.pdf2svg_path, args.verbose):
            console.print("[bold red]❌ Error:[/bold red] Failed to convert PDF to SVG")
            sys.exit(ERR_PDF2SVG)
    
    # Step 2: SVG to EMF
    success, pages_with_filters = svg2emf(
        pdf_reader, args.input, args.inkscape_path,
        pages, args.verbose, args.no_check, args.parallel
    )
    
    if not success:
        console.print("[bold red]❌ Error:[/bold red] Failed to convert SVG to EMF")
        sys.exit(ERR_SVG2EMF)
    
    if pages_with_filters:
        console.print(f"[bold yellow]⚠️  Warning:[/bold yellow] Pages {pages_with_filters} may have transparency issues.")
        console.print("[dim]   See: https://github.com/neosun100/pdf2ppt/issues/1[/dim]")
    
    # Step 3: EMF to PPT
    emf2ppt(pdf_reader, args.input, ppt_path, pages, args.verbose)
    
    # Cleanup
    if not args.no_clean:
        clean_tmp(args.input, args.verbose)
    
    # Success message
    console.print(f"\n[bold green]✅ Success![/bold green] Created: [cyan]{ppt_path}[/cyan]")
    
    # Show summary
    table = Table(show_header=False, box=None)
    table.add_row("Pages converted:", f"[cyan]{len(pages)}[/cyan]")
    table.add_row("Output file:", f"[cyan]{ppt_path}[/cyan]")
    console.print(table)


if __name__ == '__main__':
    main()
