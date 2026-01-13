"""
To modify slide layouts:
- open in WPS Presentation
- Design --> Edit slide Master
- Can right-click to rename layout
- make sure that no actual slides are created!
"""

import tempfile
from collections import Counter
import matplotlib.pyplot as plt
import pandas as pd
from slugify import slugify
import datetime as dt
from PIL import Image
import pptx
import pptx.util
from pptx.util import Pt
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from xdat import xsettings
from bidi.algorithm import get_display as fix_rtl_bidi
import re
import textwrap


DEFAULT_THEME = xsettings.XDAT_ROOT.joinpath('media', 'default_theme.pptx')
assert DEFAULT_THEME.exists(), DEFAULT_THEME


class Presentation:
    def __init__(self, title=None, theme=DEFAULT_THEME, print_layout=False, fake=False):
        self.theme = theme
        self.prs = pptx.Presentation(self.theme)
        self.fake = fake
        self.title = title

        if not fake:
            if print_layout:
                self.print_layout()

            if title:
                date_str = dt.datetime.now().strftime('%B %d, %Y')
                self.add_slide_title(title=title, subtitle=xsettings.PROJECT_NAME_PRETTY, note=date_str)

    def __bool__(self):
        return not self.fake

    def add_slide(self, layout_name, **kwargs):
        add_slide(self.prs, layout_name, **kwargs)

    def add_slide_title(self, title='', subtitle='', note=''):
        self.add_slide('title', title=title, subtitle=subtitle, text=note)

    def add_slide_h0(self, title=''):
        self.add_slide('main_point', title=title)

    def add_slide_h1(self, title='', subtitle=''):
        self.add_slide('section_header', title=title, subtitle=subtitle)

    def add_slide_h2(self, title='', subtitle='', desc=''):
        self.add_slide('section_title_and_description', title=title, subtitle=subtitle, text=desc)

    def add_slide_caption(self, title='', content=''):
        self.add_slide('caption', text=title, text_2=content)

    def _no_val(self, val):
        if val is None:
            return True
        if isinstance(val, str) and not val:
            return True
        return False

    def _has_val(self, val):
        return not self._no_val(val)

    def add_slide_content(self, title='', desc='', main_content='', sub_content='', sub_title='', slide_note=''):
        if self._no_val(desc) and self._no_val(sub_content) and self._no_val(sub_title):
            self.add_slide('title_and_body', title=title, text=main_content, slide_note=slide_note)
        elif self._no_val(sub_content) and self._no_val(sub_title):
            self.add_slide('left_column', title=title, text=desc, text_2=main_content, slide_note=slide_note)
        else:
            assert not (self._has_val(sub_content) and self._has_val(sub_title)), 'Can only specify one'
            if self._has_val(sub_title):
                self.add_slide('left_column_2', title=title, text=desc, text_2=main_content, text_3=sub_content, slide_note=slide_note)

            if self._has_val(sub_content):
                self.add_slide('left_column_3', title=title, text=desc, text_2=main_content, text_3=sub_content, slide_note=slide_note)

    def add_slide_content_2cols(self, title='', left='', right='', left_title='', right_title='', desc='', sub_content='', slide_note=''):
        if self._has_val(desc) or self._has_val(sub_content):
            self.add_slide('left_with_two_cols', title=title, text=desc, text_2=left, text_3=right, text_4=left_title, text_5=right_title, text_6=sub_content, slide_note=slide_note)
        elif self._no_val(left_title) and self._no_val(right_title):
            self.add_slide('two_columns', title=title, text=left, text_2=right, slide_note=slide_note)
        else:
            self.add_slide('two_columns_with_subtitles', title=title, text=left, text_2=right, text_3=left_title, text_4=right_title, slide_note=slide_note)

    def print_layout(self):
        print_layout(self.theme)

    def save(self, out_path=None):
        if not self.fake:
            if not out_path:
                title = self.title or 'unnamed'
                file_name = re.sub(r'[\/:*?"<>|]', ' ', title).strip()
                out_path = xsettings.OUTPUT_PATH.joinpath(f"{file_name}.pptx")

            self.prs.save(out_path)

    @classmethod
    def capture_image(cls):
        return Img()


class Img:
    """
    Everything's in inches...
    """
    DPI = 80

    def __init__(self, tight_layout=True):
        assert xsettings.CACHE_PATH is not None, "must set xsettings.CACHE_PATH"
        tmp_folder = xsettings.CACHE_PATH.joinpath('xpptx')
        tmp_folder.ensure_dir()

        self.img_path = tempfile.NamedTemporaryFile(suffix='.png', delete=False, dir=tmp_folder).name

        if tight_layout:
            plt.tight_layout()

        plt.savefig(self.img_path, pad_inches=0)
        plt.clf()
        plt.cla()
        plt.close('all')
        img = Image.open(self.img_path)
        w, h = img.size
        self.width = w/self.DPI
        self.height = h/self.DPI

    def box(self, width, height):
        rw = self.width / width
        rh = self.height / height
        rmax = max(rh, rw)
        return self.width/rmax, self.height/rmax

    def __str__(self):
        return self.img_path

    def __repr__(self):
        return self.img_path

    def __lt__(self, other):
        return False


def _slug(text):
    return slugify(text, separator="_")


def _slug_dict(input_dict):
    counts = Counter()
    out_dict = dict()
    for text, v in input_dict.items():
        try:
            parts = text.split()
            int(parts[-1])
            parts = parts[:-1]
            text = " ".join(parts)
        except:
            pass

        try:
            parts = text.split()
            if parts[-1].lower() == 'placeholder':
                parts = parts[:-1]
            text = " ".join(parts)
        except:
            pass

        text = _slug(text)
        counts[text] += 1

        if counts[text] > 1:
            text = f"{text}_{counts[text]}"

        out_dict[text] = v

    return out_dict


def _get_layout(prs):
    layouts = _slug_dict({l.name: l for l in prs.slide_layouts})
    return layouts


def _get_placeholders(slide):
    placeholders = _slug_dict({p.name: p for p in slide.placeholders})
    return placeholders


def add_slide(prs, layout_name, slide_note='', **kwargs):
    md_conv = MarkdownToTextFrame()
    layouts = _get_layout(prs)
    assert layout_name in layouts, f"{layout_name} not in: {sorted(layouts)}"

    layout = layouts[layout_name]
    slide = prs.slides.add_slide(layout)

    if slide_note:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        md_conv.auto_md(slide_note, text_frame)

    placeholders = _get_placeholders(slide)

    for k, v in kwargs.items():
        assert k in placeholders, f"{k} not in: {sorted(placeholders)}"
        p = placeholders[k]

        if isinstance(v, dict):
            v = pd.Series(v)

        if isinstance(v, pd.Series):
            v = pd.DataFrame(v).reset_index()
            v.columns = ['', '']

        if v is None:
            p.text = ' '

        elif isinstance(v, str):
            if len(v) == 0:
                p.text = ' '

            else:
                md_conv.auto_md(v, p.text_frame)

            nlines = len(v.split('\n'))
            if nlines < 10:
                fz = None
            elif nlines < 20:
                fz = 10
            else:
                fz = 8

            if fz:
                for prg in p.text_frame.paragraphs:
                    if prg.runs:
                        for prgr in prg.runs:
                            prgr.font.size = Pt(fz)
                    else:
                        # paragraph has no runs (blank spacer) → set paragraph-level font
                        prg.font.size = Pt(fz)

        elif isinstance(v, list):
            if len(v) == 0:
                p.text = ' '

            else:
                tf = p.text_frame

                tf.text = v[0]
                for item in v[1:]:
                    p = tf.add_paragraph()
                    p.text = str(item)
                    p.level = 0  # top‐level bullet

        elif isinstance(v, Img):
            w,h = v.box(p.width, p.height)
            w = int(w)
            h = int(h)
            slide.shapes.add_picture(str(v), p.left, p.top, height=h)
            p.text = ' '

        elif isinstance(v, pd.DataFrame):
            shape = v.shape
            p.text = ' '
            # (y, x)
            if len(v) < 10:
                fz = 10
            elif len(v) < 20:
                fz = 8
            else:
                fz = 6

            table = slide.shapes.add_table(shape[0]+1, shape[1], p.left, p.top, height=p.height, width=p.width).table
            for x in range(shape[1]):
                cell = table.cell(0, x)
                cell.text = fix_rtl_bidi(str(v.columns[x]))
                for p in cell.text_frame.paragraphs:
                    p.font.size = pptx.util.Pt(fz)

            for y in range(shape[0]):
                for x in range(shape[1]):
                    cell = table.cell(y+1, x)
                    cell.text = fix_rtl_bidi(str(v.iloc[y, x]))
                    for p in cell.text_frame.paragraphs:
                        p.font.size = pptx.util.Pt(fz)

        else:
            raise TypeError(type(v))

    for k in set(placeholders) - set(kwargs):
        p = placeholders[k]
        p.text = ' '

    return


def print_layout(template_path=None):
    prs = pptx.Presentation(template_path)
    layouts = _get_layout(prs)
    for layout_name in sorted(layouts):
        print(f"- {layout_name}")
        layout = layouts[layout_name]
        slide = prs.slides.add_slide(layout)
        placeholders = _get_placeholders(slide)
        for pl in sorted(placeholders):
            print(f"  + {pl}")

    return


class MarkdownToTextFrame:
    """
    Minimal Markdown -> python-pptx TextFrame renderer.

    Changes requested:
      - '_' emphasis works like '*' (italic).
      - #, ##, ### all rendered the same: bold + underlined (no size differences).
      - Plain (non-bulleted) paragraphs are forced to no wrap-indent (marL=0, indent=0).

    Supports:
      - headings: #, ##, ### (bold + underline)
      - unordered lists (- or *), nesting via leading spaces (2 per level)
      - ordered lists (1. or 1)), nesting via leading spaces (2 per level)
      - paragraphs
      - inline **bold**, *italic* / _italic_, `code`
      - links: [text](url) and bare URLs (http(s)://..., www....)
      - fenced code blocks: ``` ... ```
    """

    # ---------- regexes ----------
    _LIST_UN = re.compile(r"^(?P<indent>\s*)([-*])\s+(?P<text>.*)$")
    _LIST_OL = re.compile(r"^(?P<indent>\s*)(?P<num>\d+)[\.\)]\s+(?P<text>.*)$")
    _H1 = re.compile(r"^\s*#\s+(?P<text>.+)\s*$")
    _H2 = re.compile(r"^\s*##\s+(?P<text>.+)\s*$")
    _H3 = re.compile(r"^\s*###\s+(?P<text>.+)\s*$")

    # Token order matters:
    # link -> bold -> code -> italic(*) -> italic(_) -> bare URL
    _TOKENIZER = re.compile(
        r"("
        r"\[[^\]]+\]\([^)]+\)"  # [text](url)
        r"|\*\*[^*\n]+\*\*"  # **bold**
        r"|`[^`\n]+`"  # `code`
        r"|(?<!\w)\*(?!\s)[^*\n]+?\*(?!\w)"  # *italic* (not inside word)
        r"|(?<!\w)_(?!\s|_)[^_\n]+?_(?!\w)"  # _italic_ (not inside word)
        r"|https?://[^\s<>()\[\]{}\"']+[^\s<>()\[\]{}\"'.,;:!?]"  # URL
        r"|www\.[^\s<>()\[\]{}\"']+[^\s<>()\[\]{}\"'.,;:!?]"  # URL
        r")"
    )

    def __init__(
        self,
        bullet_char: str = "•",
        link_underline: bool = True,
        code_font_name: str = "Consolas",
        bullet_indent_per_level_emu: int = 1200,  # 1200 ≈ 0.167"
    ):
        self.bullet_char = bullet_char
        self.link_underline = link_underline
        self.code_font_name = code_font_name
        self._per_level = bullet_indent_per_level_emu

    def _looks_like_markdown(self, s: str) -> bool:
        """
        Heuristically detect 'real' Markdown without triggering on snake_case.
        - Headings, lists, fenced code blocks
        - Inline **bold**, *italic*, _italic_, `code`, [link](url)
        - Bare URLs alone do NOT trigger Markdown.
        """
        if not s:
            return False

        # Fenced code blocks
        if re.search(r"^\s*```", s, flags=re.M):
            return True

        # Line-level constructs
        for line in s.splitlines():
            if self._H1.match(line) or self._H2.match(line) or self._H3.match(line):
                return True
            if self._LIST_UN.match(line) or self._LIST_OL.match(line):
                return True

        # Inline constructs (strict)
        # links
        if re.search(r"\[[^\]]+\]\([^\n]+\)", s):
            return True
        # bold
        if re.search(r"\*\*[^*\n]+\*\*", s):
            return True
        # code
        if re.search(r"`[^`\n]+`", s):
            return True
        # italic with *  — must not be inside a word
        if re.search(r"(?<!\w)\*(?!\s)[^*\n]+?\*(?!\w)", s):
            return True
        # italic with _  — must not be inside a word (fix for perc_diff)
        if re.search(r"(?<!\w)_(?!\s|_)[^_\n]+?_(?!\w)", s):
            return True

        # Bare URLs do not trigger markdown mode
        return False

    # ---------- convenience API ----------
    def auto_md(self, original_text: str, text_frame):
        """
        If 'original_text' looks like Markdown, render with apply_markdown_to_text_frame().
        Otherwise, write the text verbatim to a single, plain paragraph with no bullets/indents.
        """
        if self._looks_like_markdown(original_text or ""):
            self.apply_markdown_to_text_frame(original_text, text_frame)
            return

        # Plain text fallback (no markdown processing)
        text_frame.clear()
        self._normalize_text_frame(text_frame)
        try:
            text_frame.margin_left = 0
        except Exception:
            pass

        p = self._new_para(text_frame)
        self._set_no_bullets(p)

        # Preserve original newlines; let python-pptx normalize CRLF
        try:
            p.text = (original_text or "").replace("\r\n", "\n")
        except Exception:
            # Fallback if .text assignment fails in some environments
            self._add_run(p, original_text or "")

        # Ensure plain paragraphs have no wrap-indent artifacts
        self._post_fix_plain_wrapping(text_frame)

    # ---------- safe paragraph creation ----------
    def _new_para(self, text_frame):
        # After clear(), many themes leave one empty paragraph behind.
        if len(text_frame.paragraphs) == 0:
            return text_frame.add_paragraph()
        if len(text_frame.paragraphs) == 1 and not text_frame.paragraphs[0].text:
            return text_frame.paragraphs[0]
        return text_frame.add_paragraph()

    # ---------- low-level helpers (oxml-safe) ----------
    def _p_el(self, paragraph):
        # works for both Paragraph and _Paragraph
        return getattr(paragraph, "_p", None) or getattr(paragraph, "_element", None) or paragraph

    def _pPr(self, paragraph):
        p_el = self._p_el(paragraph)
        return p_el.get_or_add_pPr()

    def _clear_bullet_elems(self, ppr):
        for tag in ("a:buAutoNum", "a:buChar", "a:buBlip", "a:buFont", "a:buNone"):
            el = ppr.find(qn(tag))
            if el is not None:
                ppr.remove(el)

    def _is_bulleted(self, paragraph) -> bool:
        ppr = self._pPr(paragraph)
        return (ppr.find(qn("a:buChar")) is not None) or (ppr.find(qn("a:buAutoNum")) is not None)

    def _set_no_bullets(self, paragraph):
        # plain paragraph: no bullets, no indent, no list level, no tabs
        try:
            if getattr(paragraph, "level", None) is not None:
                paragraph.level = 0
        except Exception:
            pass
        ppr = self._pPr(paragraph)
        self._clear_bullet_elems(ppr)

        # zero out margin + hanging indent
        ppr.set(qn("a:marL"), "0")
        ppr.set(qn("a:indent"), "0")

        # remove list level and tabs that cause wrap-indent
        lvl_attr = qn("a:lvl")
        if lvl_attr in ppr.attrib:
            del ppr.attrib[lvl_attr]
        tablst = ppr.find(qn("a:tabLst"))
        if tablst is not None:
            ppr.remove(tablst)

        # explicitly "no bullets"
        ppr.append(OxmlElement("a:buNone"))

    def _set_bullet_char(self, paragraph, level=0):
        try:
            if getattr(paragraph, "level", None) is not None:
                paragraph.level = level
        except Exception:
            pass
        ppr = self._pPr(paragraph)
        self._clear_bullet_elems(ppr)
        # control indent explicitly
        ppr.set(qn("a:marL"), str(level * self._per_level))
        ppr.set(qn("a:indent"), "0")
        buChar = OxmlElement("a:buChar")
        buChar.set("char", self.bullet_char)
        ppr.append(buChar)

    def _set_bullet_autonum(self, paragraph, level=0, num_type="arabicPeriod", start_at=None):
        try:
            if getattr(paragraph, "level", None) is not None:
                paragraph.level = level
        except Exception:
            pass
        ppr = self._pPr(paragraph)
        self._clear_bullet_elems(ppr)
        ppr.set(qn("a:marL"), str(level * self._per_level))
        ppr.set(qn("a:indent"), "0")
        bu = OxmlElement("a:buAutoNum")
        bu.set("type", num_type)
        if start_at is not None:
            bu.set("startAt", str(start_at))
        ppr.append(bu)

    # ---------- inline token handling ----------
    def _add_run(self, paragraph, text, *, bold=False, italic=False, code=False, link=None, underline=False):
        r = paragraph.add_run()  # works on both Paragraph and _Paragraph
        r.text = text.replace("\r\n", "\n")
        f = r.font
        try:
            f.bold = bold or None
            f.italic = italic or None
            if underline:
                f.underline = True
            if code:
                f.name = self.code_font_name
            if link:
                r.hyperlink.address = self._autonormalize_url(link)
                if self.link_underline:
                    f.underline = True
        except Exception:
            pass
        return r

    def _autonormalize_url(self, url: str) -> str:
        return "http://" + url if url.startswith("www.") else url

    def _add_inline_runs(self, paragraph, text, *, force_bold=False, force_underline=False):
        """
        Tokenizes and renders:
          [text](url) | **bold** | `code` | *italic* | _italic_ | bare URLs
        """
        pos = 0
        while True:
            m = self._TOKENIZER.search(text, pos)
            if not m:
                if pos < len(text):
                    self._add_run(paragraph, text[pos:], bold=force_bold, underline=force_underline)
                break

            if m.start() > pos:
                self._add_run(paragraph, text[pos:m.start()], bold=force_bold, underline=force_underline)

            tok = m.group(0)
            if tok.startswith("["):  # [text](url)
                close_br = tok.find("]")
                open_par = tok.find("(", close_br)
                close_par = tok.rfind(")")
                link_text = tok[1:close_br]
                link_url = tok[open_par + 1:close_par].strip()
                self._add_run(paragraph, link_text, link=link_url, bold=force_bold, underline=True or force_underline)
            elif tok.startswith("**"):
                self._add_run(paragraph, tok[2:-2], bold=True or force_bold, underline=force_underline)
            elif tok.startswith("`"):
                self._add_run(paragraph, tok[1:-1], code=True, underline=force_underline)
            elif tok.startswith("*"):
                self._add_run(paragraph, tok[1:-1], italic=True, bold=force_bold, underline=force_underline)
            elif tok.startswith("_") and re.fullmatch(r"(?<!\w)_(?!\s|_)[^_\n]+?_(?!\w)", tok):
                self._add_run(paragraph, tok[1:-1], italic=True, bold=force_bold, underline=force_underline)
            else:
                self._add_run(paragraph, tok, bold=force_bold, underline=force_underline)

            pos = m.end()

    # ---------- post-fix: remove wrap-indent from plain paragraphs ----------
    def _post_fix_plain_wrapping(self, text_frame):
        for prg in text_frame.paragraphs:
            if not self._is_bulleted(prg):
                ppr = self._pPr(prg)
                ppr.set(qn("a:marL"), "0")
                ppr.set(qn("a:indent"), "0")
                # also clear level + tab stops again, just in case
                lvl_attr = qn("a:lvl")
                if lvl_attr in ppr.attrib:
                    del ppr.attrib[lvl_attr]
                tablst = ppr.find(qn("a:tabLst"))
                if tablst is not None:
                    ppr.remove(tablst)

    def _normalize_text_frame(self, text_frame):
        """
        Kill frame-level left inset and theme list-style that cause wrap indents.
        Safe for python-pptx 1.0.x.
        """
        # Best-effort: API margin
        try:
            text_frame.margin_left = 0
        except Exception:
            pass

        # Access <a:txBody>
        tx = getattr(text_frame, "_txBody", None) or getattr(text_frame, "_element", None)
        if tx is None:
            return

        # <a:bodyPr lIns="0"> (also clear defTabSz if present)
        bodyPr = tx.find(qn("a:bodyPr"))
        if bodyPr is not None:
            bodyPr.set(qn("a:lIns"), "0")
            # optional: also zero these if you see extra padding:
            # bodyPr.set(qn("a:rIns"), "0"); bodyPr.set(qn("a:tIns"), "0"); bodyPr.set(qn("a:bIns"), "0")
            if qn("a:defTabSz") in bodyPr.attrib:
                del bodyPr.attrib[qn("a:defTabSz")]

        # Remove <a:lstStyle> so theme list levels don't reintroduce indents
        lstStyle = tx.find(qn("a:lstStyle"))
        if lstStyle is not None:
            tx.remove(lstStyle)

    # ---------- main API ----------
    def apply_markdown_to_text_frame(self, md: str, text_frame):
        """
        Render simple Markdown into a python-pptx TextFrame.
        Clears the frame first. Headings are bold + underlined (same style for #/##/###).
        Also normalizes plain paragraphs to avoid wrap-indent.
        """
        md = textwrap.dedent(md).strip()
        text_frame.clear()
        self._normalize_text_frame(text_frame)

        try:
            text_frame.margin_left = 0
        except Exception:
            pass

        def level_from_indent(s: str) -> int:
            return max(0, len(s) // 2)  # 2 spaces per level

        lines = md.splitlines()
        i = 0
        while i < len(lines):
            line = lines[i].rstrip()

            # fenced code block
            if line.strip().startswith("```"):
                j = i + 1
                code_lines = []
                while j < len(lines) and not lines[j].strip().startswith("```"):
                    code_lines.append(lines[j])
                    j += 1
                i = j + 1
                for cl in code_lines or [" "]:
                    p = self._new_para(text_frame)
                    self._set_no_bullets(p)
                    self._add_run(p, cl, code=True)
                continue

            if not line.strip():  # blank line
                p = self._new_para(text_frame)
                self._set_no_bullets(p)
                # Insert a zero-width space so later font sweeps can affect it if needed
                self._add_run(p, "\u200B")
                i += 1
                continue

            # headings (#, ##, ###) -> same styling (bold + underline)
            if self._H1.match(line) or self._H2.match(line) or self._H3.match(line):
                text = (self._H1.match(line) or self._H2.match(line) or self._H3.match(line)).group("text")
                p = self._new_para(text_frame)
                self._set_no_bullets(p)
                self._add_inline_runs(p, text, force_bold=True, force_underline=True)
                i += 1
                continue

            # unordered list
            m = self._LIST_UN.match(line)
            if m:
                lvl = level_from_indent(m.group("indent"))
                p = self._new_para(text_frame)
                self._set_bullet_char(p, level=lvl)
                self._add_inline_runs(p, m.group("text"))
                i += 1
                continue

            # ordered list
            m = self._LIST_OL.match(line)
            if m:
                lvl = level_from_indent(m.group("indent"))
                p = self._new_para(text_frame)
                self._set_bullet_autonum(p, level=lvl)
                self._add_inline_runs(p, m.group("text"))
                i += 1
                continue

            # normal paragraph (wrap consecutive plain lines)
            para_lines = [line]
            j = i + 1
            while j < len(lines):
                nxt = lines[j]
                if (self._LIST_UN.match(nxt) or self._LIST_OL.match(nxt) or
                        self._H1.match(nxt) or self._H2.match(nxt) or self._H3.match(nxt) or
                        not nxt.strip() or nxt.strip().startswith("```")):
                    break
                para_lines.append(nxt.rstrip())
                j += 1
            p = self._new_para(text_frame)
            self._set_no_bullets(p)
            # Use '\n' so single Markdown newlines render as line breaks (<br>-like)
            self._add_inline_runs(p, "\n".join(para_lines))
            i = j

        # Final sweep: remove wrap-indent from plain paragraphs, keep bullets as-is
        self._post_fix_plain_wrapping(text_frame)



if __name__ == "__main__":
    xsettings.PROJECT_NAME = 'xdat'
    xsettings.updated_config()
    plt.scatter([1, 2], [3, 4])
    i = Img()
    w, h = i.box(5, 5)

    prs = Presentation()

    df = pd.DataFrame({'hi': [1, 2, 3], 'there': ['a', 'b', 'c']})

    prs.add_slide('title_and_body', title='hi', text=i)

    prs.add_slide('title_and_body', title='hi', text=prs.capture_image())
    prs.add_slide('title_and_body', title='what', text=df)

    md_text = """
    
    # Hi there
    
    List 1:
    - world
    - hi
    - there
    
    # great!
    
    List2:
    1) Hi
    2) There
    
    Now _what_ **is** this?
    A link!  [the text](https://pypi.org/project/xdat/) ...
    
    ``` 
    code!
    lots of code
    ```    
    """

    prs.add_slide('title_and_body', title='what', text=md_text)

    prs.save('/tmp/xdat/test.pptx')
